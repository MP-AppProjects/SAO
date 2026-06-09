# -*- coding: utf-8 -*-
import streamlit as st
import pandas as pd
import numpy as np
import pyreadstat
import tempfile
import io
import json
import re
import copy
import string
import datetime
import plotly.express as px
import plotly.graph_objects as go
import scipy.stats as stats
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.cm as cm
from collections import defaultdict
from sklearn.cluster import KMeans
from sklearn.preprocessing import StandardScaler
from scipy.cluster.hierarchy import linkage, dendrogram, fcluster
import statsmodels.api as sm
from statsmodels.stats.outliers_influence import variance_inflation_factor
from factor_analyzer import FactorAnalyzer
from factor_analyzer.factor_analyzer import calculate_kmo, calculate_bartlett_sphericity
import sqlite3
import hashlib
import secrets
import uuid
import os
import ipaddress
import urllib.request
import urllib.error
import time

st.set_page_config(page_title="System Analiz Openfield (SAO)", layout="wide", page_icon="\U0001f4ca")

# =============================================================
# PANEL ADMINISTRACYJNY -- baza danych i funkcje auth
# =============================================================
# Backend auth / baza danych / panel admina -- wydzielone do sao_admin.py
from sao_admin import *  # noqa: F401,F403
# =============================================================
# LOGOWANIE
# =============================================================
# -------------------------------------------------------------
# CSS -- minimalistyczny, profesjonalny styl
# -------------------------------------------------------------
st.markdown("""
<style>
    /* \u2500\u2500 Extra top padding so content doesn\u2019t hide under Streamlit\u2019s fixed header \u2500\u2500 */
    .block-container { padding-top: 3.5rem !important; }

    /* \u2500\u2500 Sticky tab bar that stays below the Streamlit header (58px) \u2500\u2500 */
    .stTabs [data-baseweb="tab-list"] {
        position: sticky;
        top: 58px;
        z-index: 99;
        background: white;
        padding-top: 4px;
        padding-bottom: 3px;          /* room for the active-tab underline */
        overflow-x: auto;
        overflow-y: visible;          /* must be visible so underline isn\u2019t clipped */
        flex-wrap: nowrap;
        scrollbar-width: thin;
        scrollbar-color: #2E75B6 #f0f0f0;
        box-shadow: 0 1px 0 0 #e6e6e6; /* replicate the baseline separator */
    }

    /* \u2500\u2500 Compact tabs: smaller font + tight padding \u2500\u2500 */
    .stTabs [data-baseweb="tab"] {
        font-size: 0.78rem;
        padding: 5px 10px;
        white-space: nowrap;
    }

    /* Chrome / Edge / Safari scrollbar */
    .stTabs [data-baseweb="tab-list"]::-webkit-scrollbar { height: 4px; }
    .stTabs [data-baseweb="tab-list"]::-webkit-scrollbar-track {
        background: #f0f0f0; border-radius: 2px;
    }
    .stTabs [data-baseweb="tab-list"]::-webkit-scrollbar-thumb {
        background: #2E75B6; border-radius: 2px;
    }

    .metric-card {
        background: #f8f9fa; border-radius: 8px; padding: 14px 18px;
        border-left: 4px solid #2E75B6; margin-bottom: 8px;
    }
    .metric-card b { color: #2E75B6; font-size: 1.1rem; }
    .sig-green { color: #006100; font-weight: bold; }
    .sig-red   { color: #C00000; }
    .sig-orange{ color: #E36C09; }
    .section-header {
        background: linear-gradient(90deg, #2E75B6, #1F4E79);
        color: white; padding: 8px 16px; border-radius: 6px;
        font-size: 0.95rem; font-weight: bold; margin-bottom: 12px;
    }
    div[data-testid="stSidebarNav"] { display: none; }
    .stAlert { border-radius: 8px; }

</style>
""", unsafe_allow_html=True)

# =============================================================
# GATE LOGOWANIA
# =============================================================

# Idle timeout + walidacja tokenu w DB (wykrywa force-logout z panelu admina)
if st.session_state.get('authenticated'):
    _lg_token = st.session_state.get('session_token')
    _lg_db_sess = _validate_session(_lg_token)
    if _lg_db_sess is None:
        # Sesja zostala zakonieczona (force-logout lub invalid token)
        st.session_state.clear()
        st.warning("\U0001f512 Twoja sesja zosta\u0142a zako\u0144czona przez administratora. Zaloguj si\u0119 ponownie.")
    else:
        _lg_idle_timeout = _get_setting_int("idle_timeout_minutes", 60)
        _lg_last_ts = st.session_state.get('last_activity_ts', time.time())
        if time.time() - _lg_last_ts > _lg_idle_timeout * 60:
            _end_session(_lg_db_sess["id"], "idle_timeout")
            st.session_state.clear()
            st.warning("\u23f0 Sesja wygas\u0142a z powodu bezczynno\u015bci. Zaloguj si\u0119 ponownie.")
        else:
            st.session_state.last_activity_ts = time.time()
            _touch_session(_lg_db_sess["id"])

# Gate: blokuj nieuprawniony dostep
if not st.session_state.get('authenticated'):
    _lg_ip = _get_client_ip()
    _lg_ua = _get_user_agent()
    _lgc1, _lgc2, _lgc3 = st.columns([1, 2, 1])
    with _lgc2:
        st.markdown("""
<div style="text-align:center; padding:30px 0 10px 0;">
  <span style="font-size:2.5rem;">\U0001f4ca</span><br>
  <span style="font-size:1.4rem; font-weight:bold; color:#1F4E79;">System Analiz Openfield (SAO)</span><br>
  <span style="color:#666; font-size:0.9rem;">Zaloguj si\u0119, aby kontynuowa\u0107</span>
</div>""", unsafe_allow_html=True)
        with st.form("lg_login_form"):
            _lg_user_in = st.text_input("\U0001f464 Nazwa u\u017cytkownika", placeholder="login")
            _lg_pass_in = st.text_input("\U0001f511 Has\u0142o", type="password", placeholder="has\u0142o")
            _lg_btn = st.form_submit_button("\u25b6\ufe0f Zaloguj", type="primary",
                                            use_container_width=True)
        if _lg_btn:
            if not _lg_user_in or not _lg_pass_in:
                st.error("Podaj nazw\u0119 u\u017cytkownika i has\u0142o.")
            else:
                _lg_status, _lg_pl = _attempt_login(
                    _lg_user_in.strip(), _lg_pass_in, _lg_ip, _lg_ua)
                if _lg_status == "ok":
                    _lg_u = _lg_pl["user"]
                    st.session_state.authenticated         = True
                    st.session_state.current_user_id       = _lg_u["id"]
                    st.session_state.current_user_name     = _lg_u["username"]
                    st.session_state.current_user_role     = _lg_u["role"]
                    st.session_state.current_user_perms    = _load_user_perms(
                        _lg_u["id"], _lg_u["role"])
                    st.session_state.session_token         = _lg_pl["token"]
                    st.session_state.session_db_id         = _lg_pl["session_id"]
                    st.session_state.last_activity_ts      = time.time()
                    st.session_state.must_change_password  = bool(_lg_u["must_change_password"])
                    st.session_state.current_user_ip       = _lg_ip
                    _log_activity("system", "login",
                                  {"ip": _lg_ip, "country": _lg_pl.get("country")})
                    st.rerun()
                elif _lg_status == "locked":
                    _lk_until = str(_lg_pl or "")[:16].replace("T", " ")
                    st.error("\U0001f512 Konto zablokowane. Spr\u00f3buj ponownie po: "
                             + _lk_until + " UTC")
                elif _lg_status == "inactive":
                    st.error("\u26d4 To konto jest nieaktywne. Skontaktuj si\u0119 z administratorem.")
                elif _lg_status == "expired":
                    _exp_dt = str(_lg_pl or "")[:10]
                    st.error("\u23f3 Dost\u0119p wygas\u0142 (" + _exp_dt
                             + "). Skontaktuj si\u0119 z administratorem.")
                else:
                    st.error("\u274c Nieprawid\u0142owa nazwa u\u017cytkownika lub has\u0142o.")
    st.stop()

# Zmiana hasla przy pierwszym logowaniu (must_change_password)
if st.session_state.get('must_change_password'):
    _lgc1, _lgc2, _lgc3 = st.columns([1, 2, 1])
    with _lgc2:
        st.warning("\U0001f511 Przed kontynuowaniem musisz ustawi\u0107 nowe has\u0142o.")
        _cpw_uname = st.session_state.get("current_user_name", "")
        _cpw_min = str(_get_setting_int("min_pw_length", 10))
        with st.form("lg_change_pw_form"):
            _cpw_old  = st.text_input("Aktualne has\u0142o", type="password")
            _cpw_new1 = st.text_input("Nowe has\u0142o", type="password",
                                      help="Min. " + _cpw_min + " znak\u00f3w, cyfra i litera")
            _cpw_new2 = st.text_input("Powt\u00f3rz nowe has\u0142o", type="password")
            _cpw_btn  = st.form_submit_button("\u2705 Zmie\u0144 has\u0142o", type="primary",
                                              use_container_width=True)
        if _cpw_btn:
            _cpw_user = _get_user_by_name(_cpw_uname)
            if not _cpw_user:
                st.error("B\u0142\u0105d: nie znaleziono u\u017cytkownika.")
            elif not _verify_password(_cpw_old, _cpw_user["password_hash"],
                                      _cpw_user["password_salt"]):
                st.error("\u274c Nieprawid\u0142owe aktualne has\u0142o.")
            elif _cpw_new1 != _cpw_new2:
                st.error("\u274c Nowe has\u0142a nie s\u0105 identyczne.")
            else:
                _cpw_err = _validate_password_policy(_cpw_new1, _cpw_uname)
                if _cpw_err:
                    st.error("\u274c " + _cpw_err)
                else:
                    _cpw_h, _cpw_s = _hash_password(_cpw_new1)
                    get_db().execute(
                        "UPDATE users SET password_hash=?, password_salt=?,"
                        " must_change_password=0 WHERE id=?",
                        (_cpw_h, _cpw_s, _cpw_user["id"]))
                    get_db().execute(
                        "INSERT INTO audit_log(actor_user_id, target_user_id,"
                        " event_type, details_json, created_at) VALUES(?,?,?,?,?)",
                        (_cpw_user["id"], _cpw_user["id"], "change_password",
                         json.dumps({"reason": "first_login"}), _now_iso()))
                    st.session_state.must_change_password = False
                    _log_activity("system", "change_password",
                                  {"reason": "first_login"})
                    st.success("\u2705 Has\u0142o zmienione. Mo\u017cesz teraz u\u017cywa\u0107 aplikacji.")
                    st.rerun()
    st.stop()

# -------------------------------------------------------------
# FUNKCJE POMOCNICZE
# -------------------------------------------------------------



# Warstwa analityczna i eksporty -- wydzielone do sao_core.py
import sao_core
from sao_core import *  # noqa: F401,F403

# =============================================================
# SZABLONY WYKRESOW -- domyslny szablon + presety palet kolorow.
# Tworzone raz w module "Projekt i Slownik" (zakladka Szablony
# wykresow); przy eksporcie (PowerPoint / Word) sa tylko wybierane.
# =============================================================
CHART_TPL_DEFAULT = {
    "colors": ["#2E75B6", "#1F4E79", "#4472C4", "#A5A5A5", "#ED7D31",
               "#70AD47", "#FFC000", "#5B9BD5"],
    "title_color": "#1F4E79",
    "font_size_title": 14,
    "font_size_labels": 10,
    "font_size_data": 9,
    "show_data_labels": True,
    "data_label_format": "auto",
    "data_label_bold": True,
    "legend_position": "bottom",
    "show_gridlines": False,
    "show_y_axis": False,
    "show_x_axis": True,
    "bar_bold_labels": True,
}

CHART_TPL_PRESETS = {
    "Openfield":        ["#2E75B6", "#1F4E79", "#4472C4", "#A5A5A5",
                         "#ED7D31", "#70AD47", "#FFC000", "#5B9BD5"],
    "Korporacyjny":     ["#1F4E79", "#2E75B6", "#5B9BD5", "#9DC3E6",
                         "#404040", "#808080", "#BFBFBF", "#D6DCE4"],
    "\u017bywy":        ["#2E75B6", "#ED7D31", "#70AD47", "#FFC000",
                         "#7030A0", "#C00000", "#00B0F0", "#A5A5A5"],
    "Pastelowy":        ["#9DC3E6", "#F8CBAD", "#C5E0B4", "#FFE699",
                         "#D9C2E9", "#F4B6B6", "#B6E1E0", "#D9D9D9"],
    "Stonowany":        ["#4C72B0", "#DD8452", "#55A868", "#C44E52",
                         "#8172B3", "#937860", "#DA8BC3", "#8C8C8C"],
    "Monochromatyczny": ["#1F4E79", "#2E75B6", "#4A90D9", "#7BAFE0",
                         "#A9CCE8", "#1A3A57", "#3A6AA0", "#9DC3E6"],
}

# =============================================================
# Inicjalizacja st.session_state -- raz, zaraz po zalogowaniu.
# Klucze danych MUSZA istniec, zanim wykona sie kod przed dispatchem
# menu (memoizacja pipeline'u, dashboard, pasek boczny, selektor arkusza)
# -- inaczej leci AttributeError "has no attribute ...".
# =============================================================
_SS_DEFAULTS = {
    'anova_results': [],
    'chi_results': {},
    'cleaning_ops': [],
    'conjoint_results': [],
    'custom_missing': {},
    'custom_val_labels': {},
    'custom_var_labels': {},
    'excel_col_types': {},
    'excel_sheet': "",
    'factor_results': [],
    'hclust_results': [],
    'logistic_results': [],
    'matrix_results': [],
    'matrix_sets': {},
    'maxdiff_pairs': [('', '')],
    'maxdiff_results': [],
    'mrs_sets': {},
    'normality_results': {},
    'ppt_chart_templates': {},
    'recodings': [],
    'reg_blocks': [[]],
    'regression_results': [],
    'segmentations': [],
    'split_var': None,
    'treat_empty_as_miss': False,
    'value_orders': {},
    'weight_targets': {},
    'weights': None,
    'wordcloud_results': [],
}
for _ssk, _ssv in _SS_DEFAULTS.items():
    if _ssk not in st.session_state:
        st.session_state[_ssk] = _ssv
# struktury wymagajace specjalnego ksztaltu/typu:
if 'results' not in st.session_state:
    st.session_state['results'] = {
        'czestosci': {}, 'krzyzowe': {}, 'srednie': {},
        'opisowe': {}, 'korelacje': {}, 'banner': {},
    }
# Doposaz istniejacy slownik (np. po hot-reload lub wczytaniu starego projektu)
for _rk_def in ('czestosci', 'krzyzowe', 'srednie', 'opisowe', 'korelacje', 'banner'):
    st.session_state['results'].setdefault(_rk_def, {})
if 'box_sets' not in st.session_state:
    st.session_state['box_sets'] = defaultdict(dict)
if 'user_cleared_val_labels' not in st.session_state:
    st.session_state['user_cleared_val_labels'] = set()


# =============================================================
# Szczegolowy wybor wynikow do eksportu (wspoldzielony: Excel / PPT).
# =============================================================
def _result_key_label(key, var_labels):
    """Zamien klucz wyniku (np. 'Q1 | plec=Kobieta' albo 'Q1 x Q2 [N]')
    na czytelna etykiete zmiennej, jesli istnieje."""
    s = str(key)
    _m = re.search(r'\s*\[[^\]]*\]\s*$', s)
    if _m:
        s = s[:_m.start()]
    _base, _grp = _extract_split_from_title(s)
    if ' x ' in _base:
        _lbl = ' \u00d7 '.join(get_var_display_name(_p.strip(), var_labels)
                          for _p in _base.split(' x '))
    else:
        _lbl = get_var_display_name(_base, var_labels)
    return _lbl + (' \u2014 ' + _grp if _grp else '')


def render_granular_selector(prefix, cats, var_labels, intro=None):
    """Renderuj szczegolowy wybor wynikow do eksportu (jak w module Word).

    cats: lista krotek (cat_key, tytul, opcje_list, labelize_bool).
    Zapisuje zatwierdzony wybor w st.session_state[prefix + '_sel_confirmed'].
    Zwraca biezacy (live) slownik wyboru {cat_key: [zaznaczone]}.
    """
    live = {}
    _saved_conf = st.session_state.get(prefix + "_sel_confirmed", {})
    if intro:
        st.caption(intro)
    for cat_key, title, opts, labelize in cats:
        if not opts:
            continue
        if labelize:
            _ff = lambda k, _vl=var_labels: _result_key_label(k, _vl)
        else:
            _ff = lambda k: str(k)
        # Uzyj zapisanego wyboru jako domyslnego (gdy klucz widgetu zniknal po
        # nawigacji), zachowujac tylko klucze ktore nadal istnieja w opts.
        _opts_set = set(opts)
        _default = [k for k in _saved_conf.get(cat_key, opts) if k in _opts_set]
        _n_sel = len(_default)
        _n_tot = len(opts)
        _cnt_str = (str(_n_sel) + "/" + str(_n_tot)) if _n_sel < _n_tot else str(_n_tot)
        with st.expander(title + " (" + _cnt_str + ")", expanded=False):
            live[cat_key] = st.multiselect(
                "Wybierz:", opts, default=_default,
                key=prefix + "_sel_" + cat_key,
                format_func=_ff, label_visibility="collapsed")
    st.markdown("")
    if st.button("\u2705 Zatwierd\u017a wyb\u00f3r wynik\u00f3w do eksportu",
                 type="primary", key=prefix + "_confirm_sel"):
        st.session_state[prefix + "_sel_confirmed"] = {
            _k: list(_v) for _k, _v in live.items()}
        st.success("Wyb\u00f3r zatwierdzony \u2014 zostanie u\u017cyty przy eksporcie.")
    _conf = st.session_state.get(prefix + "_sel_confirmed")
    if _conf is not None:
        _n = sum(len(_v) for _v in _conf.values())
        st.caption("\u2705 Zatwierdzono do eksportu: " + str(_n) + " wynik\u00f3w. "
                   "Po zmianie zaznacze\u0144 kliknij ponownie **Zatwierd\u017a wyb\u00f3r**.")
    else:
        st.caption("\u2139\ufe0f Wyb\u00f3r niezatwierdzony \u2014 domy\u015blnie eksportowane s\u0105 wszystkie wyniki.")
    return live


def export_selection_filter(prefix):
    """Zwroc funkcje keep(cat_key, result_key) wg zatwierdzonego wyboru
    (lub przepuszczajaca wszystko, gdy nic nie zatwierdzono)."""
    _conf = st.session_state.get(prefix + "_sel_confirmed")
    if _conf is None:
        return lambda cat, key: True
    _sets = {k: set(v) for k, v in _conf.items()}

    def _keep(cat, key):
        if cat not in _sets:
            return True
        return key in _sets[cat]
    return _keep

def _weights_ignored_note(active):
    """Adnotacja dla modulow, ktore nie stosuja wag (Conjoint, MaxDiff,
    Chmura slow, Skupienia). 'active' = czy globalne wazenie jest wlaczone."""
    if active:
        st.info("\u2696\ufe0f Uwaga: ten modu\u0142 liczy wyniki **bez wag** \u2014 "
                "globalne wa\u017cenie nie jest tu stosowane (metoda nie korzysta "
                "z wag respondent\u00f3w).")

def render_reorder_ui(prefix, categories, var_labels):
    """Panel zmiany kolejnosci wynikow w module eksportu.

    prefix     : 'wd' | 'excel' | 'ppt'
    categories : lista krotek (cat_key, display_name, selected_keys)
                 selected_keys = aktualna lista zaznaczonych kluczy
    var_labels : slownik etykiet zmiennych
    Zwraca     : dict {cat_key: [klucze w ustalonej kolejnosci]}
                 uzyj przy generowaniu zamiast oryginalnego slownika

    Kolejnosc zapisywana w st.session_state[prefix + '_result_order'].
    Przy pierwszym wywolaniu lub gdy zestaw kluczy sie zmienil, kolejnosc
    jest inicjalizowana zachowujac poprzednie ustawienia i doklejajac nowe
    na koniec.
    """
    _order_key = prefix + '_result_order'
    if _order_key not in st.session_state:
        st.session_state[_order_key] = {}
    _stored = st.session_state[_order_key]

    # Zsynchronizuj kazda kategorie z aktualnym zestawem kluczy
    _result = {}
    for cat_key, _dname, sel_keys in categories:
        if not sel_keys:
            continue
        _sel_set = set(sel_keys)
        _prev = _stored.get(cat_key, [])
        # Zachowaj istniejacy porzadek, usuniete klucze odpadaja, nowe na koniec
        _kept = [k for k in _prev if k in _sel_set]
        _new  = [k for k in sel_keys if k not in set(_prev)]
        _synced = _kept + _new
        _stored[cat_key] = _synced
        _result[cat_key] = _synced

    # Pokaz UI gdy jest co najmniej jedna niepusta kategoria
    if not any(_result.values()):
        return _result

    st.markdown("---")
    st.markdown("#### Kolejnosc wynikow w eksporcie")
    st.caption(
        "Kliknij \u2191 / \u2193 aby zmienic kolejnosc wynikow. "
        "Wyniki zostana wyeksportowane w podanej kolejnosci. "
        "Kolejnosc jest zachowywana miedzy sesjami w pliku projektu.")

    for cat_key, disp_name, sel_keys in categories:
        _ordered = _stored.get(cat_key, [])
        if not _ordered:
            continue
        _multi = len(_ordered) >= 2
        with st.expander(disp_name, expanded=False):
            for _i, _key in enumerate(_ordered):
                _lbl = _result_key_label(_key, var_labels)
                if _multi:
                    _cn, _cu, _cd, _cl = st.columns([0.5, 0.5, 0.5, 8.5])
                    _cn.markdown(
                        "<div style='text-align:right;padding-top:6px;"
                        "color:#999;font-size:.82rem;'>" + str(_i + 1) + ".</div>",
                        unsafe_allow_html=True)
                    if _cu.button(
                        "\u2191",
                        key=prefix + "_ro_up_" + cat_key + "_" + str(_i),
                        disabled=(_i == 0),
                        use_container_width=True
                    ):
                        _ordered[_i], _ordered[_i - 1] = _ordered[_i - 1], _ordered[_i]
                        st.rerun()
                    if _cd.button(
                        "\u2193",
                        key=prefix + "_ro_dn_" + cat_key + "_" + str(_i),
                        disabled=(_i == len(_ordered) - 1),
                        use_container_width=True
                    ):
                        _ordered[_i], _ordered[_i + 1] = _ordered[_i + 1], _ordered[_i]
                        st.rerun()
                    _cl.markdown(
                        "<div style='padding-top:6px;font-size:.9rem;line-height:1.3;'>"
                        + _lbl + "</div>",
                        unsafe_allow_html=True)
                else:
                    # Jedna pozycja \u2014 pokaz bez przyciskow (potwierdzenie co zostanie wyeksportowane)
                    st.markdown(
                        "<div style='padding:5px 4px;font-size:.9rem;color:#555;'>"
                        "1. " + _lbl + "</div>",
                        unsafe_allow_html=True)
    return _result

# -------------------------------------------------------------
# Admin (rola admin) nigdy nie potrzebuje pliku
_is_admin_role = (st.session_state.get("current_user_role") == "admin")

# \u0179r\u00f3d\u0142o danych pochodzi z kreatora importu (modu\u0142 "\U0001f4e5 Import danych").
# st.session_state.imported_file = {kind, name, bytes, sheet?, csv?}
_imp   = st.session_state.get('imported_file')
_kind  = _imp.get('kind') if _imp else None
is_spss    = (_kind == 'spss')
is_excel   = (_kind == 'excel')
is_csv     = (_kind == 'csv')
is_tabular = is_excel or is_csv
selected_sheet = _imp.get('sheet') if _imp else None

# Admin nie pracuje na danych; brak importu => aplikacja dzia\u0142a bez danych
_has_import   = (_imp is not None) and not _is_admin_role
_no_data_file = not _has_import

# Domy\u015blne (puste) warto\u015bci \u2014 gdy brak wczytanego pliku
df_orig_raw = pd.DataFrame()
df_orig     = pd.DataFrame()
meta_orig   = ExcelMeta([])
loaded_name = ""
_load_error = None

# Wczytanie danych z bajt\u00f3w zapisanych w sesji. BRAK twardej bramki \u2014
# aplikacja renderuje si\u0119 bez danych, a modu\u0142y analityczne wymagaj\u0105 ich
# przez _require_data() (komunikat zamiast zatrzymywania ca\u0142ej aplikacji).
if _has_import:
    try:
        with st.spinner("Wczytywanie i strukturyzowanie bazy..."):
            _fb = _imp['bytes']
            if is_spss:
                df_orig_raw, df_orig, meta_orig = load_spss_data(_fb)
            elif is_excel:
                _ov = json.dumps(st.session_state.excel_col_types, sort_keys=True)
                _ms = json.dumps(st.session_state.custom_missing, sort_keys=True)
                df_orig_raw, df_orig, meta_orig = load_excel_data(_fb, selected_sheet, _ov, _ms)
            else:  # csv
                _ov = json.dumps(st.session_state.excel_col_types, sort_keys=True)
                _ms = json.dumps(st.session_state.custom_missing, sort_keys=True)
                _co = _imp.get('csv', {})
                df_orig_raw, df_orig, meta_orig = load_csv_data(
                    _fb, sep=_co.get('sep', ';'), decimal=_co.get('decimal', ','),
                    encoding=_co.get('encoding', 'utf-8'), header_row=_co.get('header', 0),
                    col_type_overrides_json=_ov, custom_missing_json=_ms)
            loaded_name = _imp.get('name', '')
            # Excel/CSV: zapisz mapy tekst->liczba jako etykiety warto\u015bci
            if is_tabular:
                _tnm = getattr(meta_orig, '_text_to_num_maps', {})
                for _col, _lmap in _tnm.items():
                    if (_col not in st.session_state.custom_val_labels
                            and _col not in st.session_state.user_cleared_val_labels):
                        st.session_state.custom_val_labels[_col] = _lmap
    except Exception as _le:
        _load_error = str(_le)
        df_orig_raw = pd.DataFrame()
        df_orig     = pd.DataFrame()
        meta_orig   = ExcelMeta([])
        loaded_name = ""

original_cols = set(df_orig_raw.columns)
sao_core.original_cols = original_cols  # wstrzyknij dla get_var_display_name

# ---- Memoizacja roboczych ramek danych (df_raw / df / var_labels) ----
# Caly pipeline przygotowania (kopie, segmentacje/KMeans, rekodowania,
# czyszczenie, braki) przelicza sie ponownie TYLKO gdy zmieni sie
# ktorekolwiek z jego wejsc. Przy niezmienionych ustawieniach zwracamy
# gotowy wynik z pamieci sesji -- to eliminuje m.in. ponowne liczenie
# KMeans oraz regex-replace przy kazdym przeladowaniu strony.
_wd_payload = {
    'src':     [loaded_name, id(df_orig_raw), id(df_orig),
                list(df_orig_raw.shape), [str(_wc) for _wc in df_orig_raw.columns]],
    'is_spss': bool(is_spss),
    'cvar':    st.session_state.get('custom_var_labels', {}),
    'seg':     st.session_state.get('segmentations', []),
    'rec':     st.session_state.get('recodings', []),
    'hclust':  [[_wr.get('var_name'), _wr.get('n_clusters'), _wr.get('method'),
                 hash(tuple(_wr.get('labels_data', {}).values()))]
                for _wr in st.session_state.get('hclust_results', [])],
    'clean':   st.session_state.get('cleaning_ops', []),
    'empty':   st.session_state.get('treat_empty_as_miss', False),
    'miss':    st.session_state.get('custom_missing', {}),
}
_wd_key = json.dumps(_wd_payload, sort_keys=True, default=str)

if (st.session_state.get('_wd_key') == _wd_key
        and st.session_state.get('_wd_cache') is not None):
    _wd_c = st.session_state['_wd_cache']
    df_raw     = _wd_c[0].copy()
    df         = _wd_c[1].copy()
    var_labels = dict(_wd_c[2])
else:
    df_raw = df_orig_raw.copy()
    df     = df_orig.copy()
    var_labels = meta_orig.column_names_to_labels.copy()


    # Apply user-edited variable labels (from S\u0142ownik tab)
    for _col, _lbl in st.session_state.get('custom_var_labels', {}).items():
        var_labels[_col] = _lbl


    apply_segmentations(df_raw, df, var_labels, st.session_state.get('segmentations', []))
    apply_recodings(df_raw, df, var_labels, st.session_state.get('recodings', []))
    apply_hclust_columns(df_raw, df, var_labels, st.session_state.get('hclust_results', []))

    # Apply in-place text cleaning
    apply_cleaning_ops(df_raw, df, st.session_state.get('cleaning_ops', []))

    # Apply empty-as-missing
    if st.session_state.get('treat_empty_as_miss', False):
        for c in df_raw.columns:
            if df_raw[c].dtype == object:
                df_raw[c] = df_raw[c].replace(r'^\s*$', np.nan, regex=True)
                df[c]     = df[c].replace(r'^\s*$', np.nan, regex=True)

    # Apply custom missing values
    for c, m_vals in st.session_state.get('custom_missing', {}).items():
        if c in df_raw.columns:
            # Build replace list \u2014 include value as-is plus numeric variants
            replace_vals = []
            for v in m_vals:
                replace_vals.append(v)
                try:
                    iv = int(float(v))
                    replace_vals.append(iv)
                    replace_vals.append(str(iv))
                except (ValueError, TypeError):
                    pass
                try:
                    replace_vals.append(str(v))
                    replace_vals.append(str(float(v)))
                except (ValueError, TypeError):
                    pass
            replace_vals = list(dict.fromkeys(replace_vals))

            # Convert Categorical to object before replace (required for text columns)
            if hasattr(df_raw[c], 'cat'):
                df_raw[c] = df_raw[c].astype(object)
            df_raw[c] = df_raw[c].replace(replace_vals, np.nan)

            if is_spss:
                # For SPSS: replace both numeric codes and their labels in df
                label_vals = []
                _vvl = meta_orig.variable_value_labels.get(c, {})
                for v in m_vals:
                    lbl = _vvl.get(v, _vvl.get(str(v), None))
                    if lbl is not None:
                        label_vals.append(lbl)
                if hasattr(df[c], 'cat'):
                    df[c] = df[c].astype(object)
                df[c] = df[c].replace(label_vals + replace_vals, np.nan)
            else:
                if hasattr(df[c], 'cat'):
                    df[c] = df[c].astype(object)
                df[c] = df[c].replace(replace_vals, np.nan)
    st.session_state['_wd_cache'] = (df_raw.copy(), df.copy(), dict(var_labels))
    st.session_state['_wd_key']   = _wd_key

hidden_cols = set()
for set_data in st.session_state.get('mrs_sets', {}).values():
    cols = set_data if isinstance(set_data, list) else set_data.get('cols', [])
    hidden_cols.update(cols)
visible_columns = [c for c in df.columns if c not in hidden_cols]
all_options     = visible_columns + list(st.session_state.get('mrs_sets', {}).keys()) + list(st.session_state.get('matrix_sets', {}).keys())
# Wariant bez matrix_sets: do tabel krzyzowych i srednich (baterie sa tylko dla czestosci/matrycowych)
all_options_no_matrix = visible_columns + list(st.session_state.get('mrs_sets', {}).keys())
numeric_cols_raw = df_raw.select_dtypes(include=[np.number]).columns.tolist()
numeric_cols     = [c for c in numeric_cols_raw if c not in hidden_cols and c in visible_columns]

# Sidebar status (tylko gdy plik wczytany; dla admina ukryte)
n_rows, n_cols = len(df_raw), len(df_raw.columns)
use_weights = False
if not _is_admin_role:
    st.sidebar.markdown("---")
    if loaded_name:
        src_icon = "\U0001f4ca" if is_spss else ("\U0001f4d1" if is_csv else "\U0001f4c8")
        st.sidebar.success(f"{src_icon} **{loaded_name}**\n\n{n_rows:,} respondent\u00f3w \u00b7 {n_cols} zmiennych")
    else:
        st.sidebar.info("\U0001f4e5 Brak danych \u2014 wczytaj plik w module **Import danych**.")
        if _load_error:
            st.sidebar.error("\u274c Ostatni import nie powi\u00f3d\u0142 si\u0119.")

    if st.session_state.get('weights') is not None:
        st.sidebar.markdown("---")
        use_weights = st.sidebar.checkbox("\u2696\ufe0f Zastosuj wagi w analizach", value=True)

st.sidebar.markdown("---")

# Sidebar: info o zalogowanym uzytkowniku + wylogowanie
_sw_name  = st.session_state.get("current_user_name", "")
_sw_role  = st.session_state.get("current_user_role", "")
_sw_label = "Administrator" if _sw_role == "admin" else "U\u017cytkownik"
st.sidebar.markdown(
    "<div style='background:#E8F0FB;border-radius:6px;padding:6px 10px;"
    "margin-bottom:6px;font-size:0.82rem;'>"
    "<b>\U0001f464 " + _sw_name + "</b><br>"
    "<span style='color:#555;'>" + _sw_label + "</span></div>",
    unsafe_allow_html=True)
if st.sidebar.button("\u21a9\ufe0f Wyloguj", key="sidebar_logout",
                     use_container_width=True):
    _sw_sid = st.session_state.get('session_db_id')
    if _sw_sid:
        _end_session(_sw_sid, "logout")
    _log_activity("system", "logout", {"username": _sw_name})
    st.session_state.clear()
    st.rerun()

st.sidebar.markdown("## \U0001f4cc Nawigacja")

def _strip_emoji(label):
    i = 0
    while i < len(label) and ord(label[i]) > 127:
        i += 1
    return label[i:].lstrip(' ')

# Menu filtrowane wg uprawnien uzytkownika
_MENU_ITEMS = [lbl for key, lbl in _MODULE_KEYS.items() if _user_can_access(key)]

# Allow tile clicks to navigate by writing to session state
if 'nav_to' not in st.session_state:
    st.session_state.nav_to = None
if 'current_menu' not in st.session_state or st.session_state.current_menu not in _MENU_ITEMS:
    st.session_state.current_menu = _MENU_ITEMS[0] if _MENU_ITEMS else ""

# If a tile was clicked, update current_menu then clear nav_to
if st.session_state.nav_to and st.session_state.nav_to in _MENU_ITEMS:
    st.session_state.current_menu = st.session_state.nav_to
    st.session_state.nav_to = None

# Use key="current_menu" so Streamlit syncs radio <-> session_state automatically.
# This avoids the double-click bug caused by manually setting index= every rerun.
menu = st.sidebar.radio("", _MENU_ITEMS,
                         key="current_menu",
                         label_visibility="collapsed",
                         format_func=_strip_emoji)


def _require_data():
    """Wymusza wczytane dane w module analitycznym. Gdy brak danych \u2014 pokazuje
    komunikat z przyciskiem do kreatora importu i zatrzymuje TYLKO ten modul
    (nie cala aplikacje)."""
    if df_raw is None or df_raw.empty:
        st.warning("\U0001f4e5 Ten modu\u0142 wymaga danych. Najpierw wczytaj plik "
                   "w module **Import danych**.")
        if st.button("\U0001f4e5 Przejd\u017a do importu danych", type="primary",
                     key=f"goimport_{menu}"):
            st.session_state.nav_to = "\U0001f4e5 Import danych"
            st.rerun()
        st.stop()


# =============================================================
# DASHBOARD
# =============================================================
if menu == "\U0001f3e0 Dashboard":
    _require_module_access("dashboard")
    n_rows, n_cols_db = len(df_raw), len(df_raw.columns)
    _has_data = bool(loaded_name) and n_rows > 0
    if _has_data:
        _hero_sub = (f"\U0001f4c1 <b>{loaded_name}</b> &nbsp;\u00b7&nbsp; "
                     f"\U0001f465 {n_rows:,} respondent\u00f3w &nbsp;\u00b7&nbsp; "
                     f"\U0001f4ca {n_cols_db} zmiennych")
    else:
        _hero_sub = ("\u26a0\ufe0f Brak wczytanych danych \u2014 przejd\u017a do modu\u0142u "
                     "<b>\U0001f4e5 Import danych</b>, aby odblokowa\u0107 analizy.")

    st.markdown(f"""
<div style="background:linear-gradient(90deg,#1F4E79,#2E75B6);
     padding:28px 36px;border-radius:12px;margin-bottom:24px;color:white;">
  <h2 style="margin:0;font-size:1.8rem;">System Analiz Openfield (SAO)</h2>
  <p style="margin:6px 0 0;opacity:.85;font-size:1rem;">{_hero_sub}</p>
</div>
""", unsafe_allow_html=True)

    if not _has_data:
        if st.button("\U0001f4e5 Wczytaj dane, aby rozpocz\u0105\u0107", type="primary",
                     key="dash_goimport"):
            st.session_state.nav_to = "\U0001f4e5 Import danych"
            st.rerun()

    # Tile definitions: (emoji, label, description, menu_key)
    _TILES = [
        ("\U0001f4e5", "Import danych",
         "Wczytaj baz\u0119 z pliku SPSS (.sav), Excel (.xlsx) lub CSV (.csv)",
         "\U0001f4e5 Import danych"),
        ("\U0001f4c1", "Projekt i S\u0142ownik",
         "Zapisz/wczytaj projekt, przegl\u0105daj s\u0142ownik zmiennych i edytuj etykiety",
         "\U0001f4c1 Projekt i S\u0142ownik"),
        ("\U0001f6e0\ufe0f", "Przygotowanie Danych",
         "Braki, rekodowanie, czyszczenie, zestawy MRS, pytania matrycowe, wa\u017cenie",
         "\U0001f6e0\ufe0f Przygotowanie Danych"),
        ("\U0001f4c8", "Analizy i Tabele",
         "Tablice cz\u0119sto\u015bci, tabele krzy\u017cowe, \u015brednie, statystyki opisowe, korelacje",
         "\U0001f4c8 Analizy i Tabele"),
        ("\U0001f4d0", "Testy Normalno\u015bci",
         "Shapiro-Wilk, Kolmogorov-Smirnov, Lilliefors, D\u2019Agostino \u2014 wykresy Q-Q i histogramy",
         "\U0001f4d0 Testy Normalno\u015bci"),
        ("\U0001f4c9", "Regresja",
         "OLS (liniowa) i logistyczna (binarna/wielomianowa) w jednym module",
         "\U0001f4c9 Regresja"),
        ("\U0001f4ca", "ANOVA",
         "Jednoczynnikowa ANOVA, test Levene\u2019a, eta\u00b2, Tukey HSD post-hoc",
         "\U0001f4ca ANOVA"),
        ("\U0001f52c", "Analiza Czynnikowa",
         "EFA, KMO, test Bartletta, wykres osypiska, macierz \u0142adunk\u00f3w",
         "\U0001f52c Analiza Czynnikowa"),
        ("\U0001f3af", "Skupienia i Segmentacja",
         "Hierarchiczne (dendrogram) + K-Means: dobr\u00f3r skupie\u0144, profile grup",
         "\U0001f3af Skupienia i Segmentacja"),
        ("\U0001f4ca", "Conjoint",
         "Rating-based (OLS) i CBC (Logit) \u2014 u\u017cyteczno\u015bci, wa\u017cno\u015b\u0107 atrybut\u00f3w",
         "\U0001f4ca Conjoint"),
        ("\U0001f522", "MaxDiff",
         "Best-Worst Scaling, B-W scores, ranking wa\u017cno\u015bci, wynik standaryzowany 0-100",
         "\U0001f522 MaxDiff"),
        ("\u2601\ufe0f", "Chmura S\u0142\u00f3w",
         "Wizualizacja pyta\u0144 otwartych, stop words, palety kolor\u00f3w, eksport PNG/JPG",
         "\u2601\ufe0f Chmura S\u0142\u00f3w"),
        ("\U0001f30a", "Por\u00f3wnanie fal",
         "Zestaw cz\u0119sto\u015bci z r\u00f3\u017cnych fal badania, delty (pp) i test istotno\u015bci zmian",
         "\U0001f30a Por\u00f3wnanie fal"),
        ("\U0001f4be", "Eksport do Excela",
         "Raport analityczny, wykresy, baza danych, spis tre\u015bci z hiperlink\u00f3w",
         "\U0001f4be Eksport do Excela"),
        ("\U0001f4ca", "Eksport do PowerPoint",
         "Edytowalne wykresy kolumnowe z cz\u0119sto\u015bci i tabel krzy\u017cowych \u2014 ka\u017cdy na osobnym slajdzie",
         "\U0001f4ca Eksport do PowerPoint"),
        ("\U0001f4c4", "Eksport do Worda",
         "Raport .docx z tabelami i interaktywnymi wykresami HTML w archiwum ZIP",
         "\U0001f4c4 Eksport do Worda"),
    ]

    # Status badges for tiles that have results
    def _has_results(menu_key):
        if "Cz\u0119sto\u015bci" in menu_key or "Analizy" in menu_key:
            return any(st.session_state.get('results', {}).get(g)
                       for g in ['czestosci','krzyzowe','srednie','opisowe','korelacje']) \
                   or bool(st.session_state.get('matrix_results'))
        if "Regresja" in menu_key and "Logistyczna" not in menu_key:
            return bool(st.session_state.get('regression_results'))
        if "Regresja" in menu_key and "Logistyczna" in menu_key:
            return bool(st.session_state.get('logistic_results'))
        if "ANOVA" in menu_key:
            return bool(st.session_state.get('anova_results'))
        if "Czynnikowa" in menu_key:
            return bool(st.session_state.get('factor_results'))
        if "Skupienia" in menu_key or "Segmentacja" in menu_key:
            return bool(st.session_state.get('hclust_results'))
        if "Conjoint" in menu_key:
            return bool(st.session_state.get('conjoint_results'))
        if "MaxDiff" in menu_key:
            return bool(st.session_state.get('maxdiff_results'))
        return False

    st.markdown("### Modu\u0142y")

    # Filter tiles to only those the current user can access
    _TILES = [t for t in _TILES if _user_can_access(_LABEL_TO_KEY.get(t[3], ""))]

    # Render tiles in rows of 3
    for row_start in range(0, len(_TILES), 3):
        cols = st.columns(3, gap="medium")
        for ci, tile in enumerate(_TILES[row_start:row_start+3]):
            icon, title, desc, key = tile
            has_res = _has_results(key)
            badge = " \u2705" if has_res else ""
            with cols[ci]:
                st.markdown(f"""
<div style="
    background:#fff;border-radius:10px;
    border:1.5px solid {'#2E75B6' if has_res else '#e0e0e0'};
    padding:20px 22px 14px;height:160px;
    box-shadow:0 2px 8px rgba(0,0,0,0.06);
    transition:border-color .2s;
    display:flex;flex-direction:column;justify-content:space-between;
">
  <div>
    <span style="font-size:1.6rem">{icon}</span>
    <span style="font-size:1rem;font-weight:700;color:#1F4E79;margin-left:8px">{title}{badge}</span>
    <p style="font-size:0.78rem;color:#595959;margin:8px 0 0;line-height:1.4">{desc}</p>
  </div>
</div>
""", unsafe_allow_html=True)
                if st.button(f"Przejd\u017a \u2192", key=f"tile_{key}",
                              use_container_width=True):
                    st.session_state.nav_to = key
                    st.rerun()
        st.write("")  # vertical gap between rows

    # Quick stats row
    st.divider()
    st.markdown("### Podsumowanie sesji")
    qs = st.columns(4)
    _split_active = st.session_state.get('split_var') is not None
    _split_label = (
        "\u2705 " + str(st.session_state.split_var)
        if _split_active else "\u274c Wy\u0142\u0105czony"
    )
    qs[0].metric("\U0001f500 Podzia\u0142 na podzbiory", _split_label)
    qs[1].metric("\u2696\ufe0f Wagi",
                 "\u2705 Tak" if st.session_state.weights is not None else "\u274c Nie")
    qs[2].metric("\U0001f9f9 Regu\u0142y czyszczenia", len(st.session_state.cleaning_ops))
    n_analyses = (
        sum(1 for g in ['czestosci','krzyzowe','srednie','opisowe','korelacje']
            if st.session_state.results.get(g))
        + bool(st.session_state.regression_results)
        + bool(st.session_state.anova_results)
        + bool(st.session_state.factor_results)
        + bool(st.session_state.conjoint_results)
        + bool(st.session_state.maxdiff_results)
    )
    qs[3].metric("\U0001f4ca Wykonane analizy", n_analyses)

# -------------------------------------------------------------
# MODUL: IMPORT DANYCH (kreator importu SPSS / Excel / CSV)
# -------------------------------------------------------------
elif menu == "\U0001f4e5 Import danych":
    _require_module_access("import")
    module_header("\U0001f4e5", "Import danych",
                  "Wczytaj plik z danymi (SPSS, Excel lub CSV), aby odblokowa\u0107 analizy")

    # \u2014 aktualny stan \u2014
    _cur = st.session_state.get('imported_file')
    if _cur:
        _ck = (_cur.get('kind') or '').upper()
        st.success(f"\u2705 Aktualnie wczytany plik: **{_cur.get('name','')}** ({_ck}) "
                   f"\u2014 {len(df_raw):,} wierszy \u00d7 {len(df_raw.columns)} kolumn")
        if st.button("\U0001f5d1\ufe0f Usu\u0144 wczytane dane", key="imp_clear"):
            st.session_state.pop('imported_file', None)
            st.session_state.pop('_wd_cache', None)
            st.session_state.pop('_wd_key', None)
            st.rerun()
        st.divider()

    if _load_error:
        st.error(f"\u274c B\u0142\u0105d wczytywania ostatniego pliku: {_load_error}")

    with st.expander("Instrukcja \u2014 jak wczyta\u0107 dane", expanded=False):
        st.markdown("""
##### \U0001f3af Do czego s\u0142u\u017cy ten modu\u0142?
To **kreator importu danych**. Bez wczytanego pliku zobaczysz Dashboard i menu,
ale modu\u0142y analityczne pozostan\u0105 zablokowane. Tutaj wczytujesz baz\u0119 w jednym
z trzech format\u00f3w: **SPSS (.sav)**, **Excel (.xlsx/.xls)** lub **CSV (.csv)**.

##### \U0001f527 Jak korzysta\u0107
1. Wybierz **format pliku**.
2. Wgraj plik (dla Excela wska\u017c **arkusz**, dla CSV ustaw **separator/kodowanie**).
3. Sprawd\u017a **podgl\u0105d** pierwszych wierszy.
4. Kliknij **Zatwierd\u017a import** \u2014 dane b\u0119d\u0105 dost\u0119pne we wszystkich modu\u0142ach.

> CSV: separator kolumn i dziesi\u0119tny oraz kodowanie s\u0105 wykrywane automatycznie,
> ale mo\u017cesz je poprawi\u0107, je\u015bli podgl\u0105d wygl\u0105da niepoprawnie (polskie pliki
> z Excela to zwykle separator **\u015brednik `;`** i przecinek dziesi\u0119tny).
        """)

    # \u2014 Krok 1: format \u2014
    _fmt = st.radio("Format pliku:",
                    ["SPSS (.sav)", "Excel (.xlsx / .xls)", "CSV (.csv)"],
                    horizontal=True, key="imp_fmt")
    _kind_sel = {"SPSS (.sav)": "spss",
                 "Excel (.xlsx / .xls)": "excel",
                 "CSV (.csv)": "csv"}[_fmt]

    # \u2014 Krok 2: upload \u2014
    _types = {"spss": ["sav"], "excel": ["xlsx", "xls"], "csv": ["csv", "txt"]}[_kind_sel]
    _up = st.file_uploader(f"Wybierz plik \u2014 {_fmt}", type=_types, key="imp_uploader")

    if _up is not None:
        _bytes = _up.getvalue()
        _opts = {}
        _prev = None

        if _kind_sel == "excel":
            try:
                _sheets = pd.ExcelFile(io.BytesIO(_bytes)).sheet_names
            except Exception as _e:
                _sheets = []
                st.error(f"Nie mo\u017cna otworzy\u0107 pliku Excel: {_e}")
            if _sheets:
                _sheet = (_sheets[0] if len(_sheets) == 1
                          else st.selectbox("Arkusz:", _sheets, key="imp_sheet"))
                _opts["sheet"] = _sheet
                try:
                    _prev = pd.read_excel(io.BytesIO(_bytes), sheet_name=_sheet,
                                          header=0, nrows=12)
                except Exception as _e:
                    st.error(f"Nie mo\u017cna odczyta\u0107 arkusza: {_e}")

        elif _kind_sel == "csv":
            _sn = sniff_csv_dialect(_bytes)
            _enc_opts = ["utf-8", "cp1250", "iso-8859-2", "latin1"]
            _sep_opts = {"\u015arednik  ;": ";", "Przecinek  ,": ",",
                         "Tabulator": "\t", "Pionowa  |": "|"}
            _dec_opts = {"Przecinek  ,": ",", "Kropka  .": "."}
            _sep_rev = {v: k for k, v in _sep_opts.items()}
            _dec_rev = {v: k for k, v in _dec_opts.items()}
            _c1, _c2, _c3 = st.columns(3)
            _enc = _c1.selectbox(
                "Kodowanie:", _enc_opts,
                index=_enc_opts.index(_sn["encoding"]) if _sn["encoding"] in _enc_opts else 0,
                key="imp_enc")
            _sep_lbl = _c2.selectbox(
                "Separator kolumn:", list(_sep_opts.keys()),
                index=list(_sep_opts.keys()).index(_sep_rev.get(_sn["sep"], "\u015arednik  ;")),
                key="imp_sep")
            _dec_lbl = _c3.selectbox(
                "Separator dziesi\u0119tny:", list(_dec_opts.keys()),
                index=list(_dec_opts.keys()).index(_dec_rev.get(_sn["decimal"], "Przecinek  ,")),
                key="imp_dec")
            _sep = _sep_opts[_sep_lbl]
            _decv = _dec_opts[_dec_lbl]
            _opts["csv"] = {"sep": _sep, "decimal": _decv, "encoding": _enc, "header": 0}
            try:
                _prev = pd.read_csv(io.BytesIO(_bytes), sep=_sep, decimal=_decv,
                                    encoding=_enc, header=0, nrows=12,
                                    engine="python", skipinitialspace=True)
            except Exception as _e:
                st.error(f"Nie mo\u017cna sparsowa\u0107 CSV z tymi ustawieniami: {_e}")

        else:  # spss
            try:
                _pr_raw, _pr_lab, _pr_meta = load_spss_data(_bytes)
                _prev = _pr_lab.head(12)
            except Exception as _e:
                st.error(f"Nie mo\u017cna wczyta\u0107 pliku SPSS: {_e}")

        # \u2014 Krok 3: podgl\u0105d \u2014
        if _prev is not None:
            st.markdown("**Podgl\u0105d (pierwsze wiersze):**")
            # Arrow nie serializuje kolumn 'category'/'object' o mieszanych typach
            # (np. etykiety SPSS jako tekst + niezmapowane kody jako float) \u2014
            # na potrzeby podgl\u0105du rzutujemy takie kolumny na tekst.
            _prev_disp = _prev.copy()
            for _pc in _prev_disp.columns:
                if (isinstance(_prev_disp[_pc].dtype, pd.CategoricalDtype)
                        or _prev_disp[_pc].dtype == object):
                    _prev_disp[_pc] = _prev_disp[_pc].astype(str)
            try:
                st.dataframe(_prev_disp, use_container_width=True)
            except Exception:
                st.dataframe(_prev.astype(str), use_container_width=True)
            st.caption(f"Wykrytych kolumn: {_prev.shape[1]}")

            # \u2014 Krok 4: zatwierdzenie \u2014
            if st.button("\u2705 Zatwierd\u017a import", type="primary", key="imp_confirm"):
                st.session_state.imported_file = {
                    "kind": _kind_sel, "name": _up.name, "bytes": _bytes, **_opts,
                }
                # przelicz pipeline od nowa dla nowych danych
                st.session_state.pop('_wd_cache', None)
                st.session_state.pop('_wd_key', None)
                _log_activity("import", "data_import",
                              {"kind": _kind_sel, "name": _up.name})
                st.success("\u2705 Dane wczytane! Mo\u017cesz przej\u015b\u0107 do Dashboardu lub modu\u0142\u00f3w analiz.")
                st.rerun()

# -------------------------------------------------------------
# MODUL 1: PROJEKT I SLOWNIK
# -------------------------------------------------------------
elif menu == "\U0001f4c1 Projekt i S\u0142ownik":
    _require_module_access("project")
    module_header("\U0001f4c1", "Projekt i S\u0142ownik")
    if st.session_state.pop("_proj_loaded_ok", False):
        st.success("\u2705 Projekt zosta\u0142 wczytany pomy\u015blnie!")
    tab_proj, tab_summary, tab_dict, tab_charts = st.tabs(["Projekt", "Podsumowanie Bazy", "S\u0142ownik Zmiennych", "Szablony wykres\u00f3w"])

    with tab_proj:
        st.markdown("#### Zarz\u0105dzanie projektem")

        # \u2500\u2500 Helpers for serialising / deserialising DataFrames \u2500
        def _df_to_dict(df):
            """Serialize a DataFrame to a JSON-safe dict."""
            if df is None or not isinstance(df, pd.DataFrame):
                return None
            return {'__df__': True, 'data': df.to_json(orient='split')}

        def _dict_to_df(d):
            """Deserialize a DataFrame from saved dict."""
            if not isinstance(d, dict) or not d.get('__df__'):
                return None
            import io as _io
            return pd.read_json(_io.StringIO(d['data']), orient='split')

        def _ser_results(res_dict):
            """Convert {title: DataFrame} dict to JSON-safe format."""
            return {k: _df_to_dict(v) for k, v in res_dict.items()}

        def _deser_results(raw):
            """Restore {title: DataFrame} from saved format."""
            if not isinstance(raw, dict):
                return {}
            return {k: _dict_to_df(v) for k, v in raw.items()
                    if _dict_to_df(v) is not None}

        def _ser_matrix_results(lst):
            """Serialize matrix_results list."""
            out = []
            for entry in lst:
                e = {}
                for k, v in entry.items():
                    e[k] = _df_to_dict(v) if isinstance(v, pd.DataFrame) else v
                out.append(e)
            return out

        def _deser_matrix_results(lst):
            if not isinstance(lst, list):
                return []
            out = []
            for entry in lst:
                e = {}
                for k, v in entry.items():
                    e[k] = _dict_to_df(v) if isinstance(v, dict) and v.get('__df__') else v
                out.append(e)
            return out

        def _safe_val(v):
            """Convert any value to a JSON-safe type, silently dropping non-serializable objects."""
            if v is None or isinstance(v, (bool, int, float, str)):
                return v
            if isinstance(v, pd.DataFrame):
                return _df_to_dict(v)
            if isinstance(v, np.ndarray):
                return v.tolist()
            if isinstance(v, (np.integer,)):
                return int(v)
            if isinstance(v, (np.floating,)):
                return float(v)
            if isinstance(v, (np.bool_,)):
                return bool(v)
            if isinstance(v, dict):
                return {kk: _safe_val(vv) for kk, vv in v.items()}
            if isinstance(v, (list, tuple)):
                return [_safe_val(i) for i in v]
            # Try basic JSON round-trip; if it fails, drop the value
            try:
                import json as _json
                _json.dumps(v)
                return v
            except (TypeError, ValueError):
                return None   # drop silently

        def _ser_reg_results(lst):
            """Serialize regression results \u2014 drop unserialisable model objects."""
            out = []
            for res in lst:
                e = {}
                for k, v in res.items():
                    if k == 'model':
                        continue   # statsmodels object \u2014 always skip
                    e[k] = _safe_val(v)
                out.append(e)
            return out

        def _deser_reg_results(lst):
            if not isinstance(lst, list):
                return []
            out = []
            for entry in lst:
                e = {}
                for k, v in entry.items():
                    if isinstance(v, dict) and v.get('__df__'):
                        e[k] = _dict_to_df(v)
                    else:
                        e[k] = v
                e['model'] = None   # model not restored \u2014 display from coef_df
                out.append(e)
            return out

        def _ser_factor_results(lst):
            out = []
            for res in lst:
                e = {k: _safe_val(v) for k, v in res.items()}
                out.append(e)
            return out

        def _deser_factor_results(lst):
            if not isinstance(lst, list):
                return []
            out = []
            for entry in lst:
                e = {}
                for k, v in entry.items():
                    e[k] = _dict_to_df(v) if isinstance(v, dict) and v.get('__df__') else v
                out.append(e)
            return out

        def _ser_hclust(lst):
            out = []
            for res in lst:
                e = {k: _safe_val(v) for k, v in res.items()}
                out.append(e)
            return out

        def _deser_hclust(lst):
            if not isinstance(lst, list):
                return []
            out = []
            for entry in lst:
                e = {}
                for k, v in entry.items():
                    if isinstance(v, dict) and v.get('__df__'):
                        e[k] = _dict_to_df(v)
                    elif isinstance(v, list) and k == 'Z':
                        e[k] = v   # keep as list
                    else:
                        e[k] = v
                out.append(e)
            return out

        def _build_project_data(include_results=True, include_data=True):
            data = {
                "_version": "3.0",
                "_saved_at": datetime.datetime.now().isoformat(timespec="seconds"),
                "_source": "spss" if is_spss else ("csv" if is_csv else "excel"),
                "meta": {
                    "name":   st.session_state.get("proj_name_inp", ""),
                    "author": st.session_state.get("proj_author_inp", ""),
                    "desc":   st.session_state.get("proj_desc_inp", ""),
                },
                # Config
                "treat_empty_as_miss": st.session_state.treat_empty_as_miss,
                "custom_missing":      st.session_state.custom_missing,
                "excel_col_types":     st.session_state.excel_col_types,
                "mrs_sets":            st.session_state.mrs_sets,
                "matrix_sets":         st.session_state.matrix_sets,
                "box_sets":            dict(st.session_state.box_sets),
                "value_orders":        dict(st.session_state.value_orders),
                "custom_var_labels":   st.session_state.custom_var_labels,
                "custom_val_labels":   st.session_state.custom_val_labels,
                "user_cleared_val_labels": list(st.session_state.user_cleared_val_labels),
                "ppt_chart_templates": st.session_state.ppt_chart_templates,
                "recodings":           st.session_state.recodings,
                "cleaning_ops":        st.session_state.cleaning_ops,
                "segmentations":       st.session_state.segmentations,
                "weight_targets":      st.session_state.weight_targets,
                "weights":             list(st.session_state.weights)
                                       if st.session_state.weights is not None else None,
                "maxdiff_pairs":       st.session_state.maxdiff_pairs,
                "reg_blocks":          st.session_state.reg_blocks,
                "split_var":           st.session_state.split_var,
            }
            # Osadzenie pliku zrodlowego (base64) \u2014 projekt samowystarczalny:
            # po wczytaniu nie trzeba osobno wgrywac pliku z danymi.
            _imp_src = st.session_state.get("imported_file")
            if include_data and _imp_src and _imp_src.get("bytes") is not None:
                import base64 as _b64
                _emb = {_k: _v for _k, _v in _imp_src.items() if _k != "bytes"}
                _emb["bytes_b64"] = _b64.b64encode(_imp_src["bytes"]).decode("ascii")
                data["imported_file"] = _emb
                data["_data_columns"] = list(df_raw.columns)
            if include_results:
                data["results"] = {
                    "czestosci": _ser_results(st.session_state.results.get('czestosci', {})),
                    "krzyzowe":  _ser_results(st.session_state.results.get('krzyzowe', {})),
                    "srednie":   _ser_results(st.session_state.results.get('srednie', {})),
                    "opisowe":   _ser_results(st.session_state.results.get('opisowe', {})),
                    "korelacje": _ser_results(st.session_state.results.get('korelacje', {})),
                    "banner":    _ser_results(st.session_state.results.get('banner', {})),
                }
                data["chi_results"]       = st.session_state.chi_results
                data["matrix_results"]    = _ser_matrix_results(st.session_state.matrix_results)
                data["regression_results"]= _ser_reg_results(st.session_state.regression_results)
                data["logistic_results"]  = _ser_reg_results(st.session_state.logistic_results)
                data["anova_results"]     = _ser_reg_results(st.session_state.anova_results)
                data["factor_results"]    = _ser_factor_results(st.session_state.factor_results)
                data["conjoint_results"]  = _ser_factor_results(st.session_state.conjoint_results)
                data["maxdiff_results"]   = _ser_factor_results(st.session_state.maxdiff_results)
                data["hclust_results"]    = _ser_hclust(st.session_state.hclust_results)
            return data

        def _validate_project_data(raw_data, df_raw):
            """Validate project JSON against current df_raw; return (cleaned, report).
            report is a list of dicts: severity in ('critical','warning','info'),
            category (str), message (str), items (list)."""
            import copy as _copy
            cleaned = _copy.deepcopy(raw_data) if raw_data else {}
            report = []
            _cols = set(df_raw.columns) if df_raw is not None else set()
            _n_rows = len(df_raw) if df_raw is not None else 0

            # 1. split_var
            _sv = cleaned.get("split_var")
            if isinstance(_sv, str) and _sv and _sv not in _cols:
                report.append({
                    'severity': 'critical',
                    'category': 'Podzia\u0142 na podzbiory',
                    'message': f"Zmienna splita `{_sv}` nie istnieje w aktualnych danych \u2014 split wy\u0142\u0105czony.",
                    'items': [_sv],
                })
                cleaned["split_var"] = None

            # 2. mrs_sets
            _mrs = cleaned.get("mrs_sets", {}) or {}
            _mrs_drop, _mrs_trim, _mrs_new = [], [], {}
            for _name, _clist in _mrs.items():
                if not isinstance(_clist, list):
                    continue
                _kept    = [c for c in _clist if c in _cols]
                _missing = [c for c in _clist if c not in _cols]
                if len(_kept) < 2:
                    _mrs_drop.append(_name)
                else:
                    _mrs_new[_name] = _kept
                    if _missing:
                        _mrs_trim.append((_name, _missing))
            if _mrs_drop:
                report.append({
                    'severity': 'warning',
                    'category': 'Zestawy MRS',
                    'message': f"Usuni\u0119to {len(_mrs_drop)} zestaw(\u00f3w) MRS (za ma\u0142o istniej\u0105cych kolumn): " + ", ".join(_mrs_drop),
                    'items': _mrs_drop,
                })
            if _mrs_trim:
                _trim_txt = "; ".join(f"{n} (brak: {', '.join(m)})" for n, m in _mrs_trim)
                report.append({
                    'severity': 'warning',
                    'category': 'Zestawy MRS',
                    'message': "Skr\u00f3cono zestawy MRS \u2014 u\u017cyta b\u0119dzie wersja tylko dla istniej\u0105cych kolumn: " + _trim_txt,
                    'items': [n for n, _ in _mrs_trim],
                })
            cleaned["mrs_sets"] = _mrs_new

            # 3. matrix_sets
            _ms = cleaned.get("matrix_sets", {}) or {}
            _ms_drop, _ms_trim, _ms_new = [], [], {}
            for _name, _clist in _ms.items():
                if not isinstance(_clist, list):
                    continue
                _kept    = [c for c in _clist if c in _cols]
                _missing = [c for c in _clist if c not in _cols]
                if len(_kept) < 2:
                    _ms_drop.append(_name)
                else:
                    _ms_new[_name] = _kept
                    if _missing:
                        _ms_trim.append((_name, _missing))
            if _ms_drop:
                report.append({
                    'severity': 'warning',
                    'category': 'Pytania matrycowe',
                    'message': f"Usuni\u0119to {len(_ms_drop)} pyta\u0144 matrycowych (za ma\u0142o istniej\u0105cych kolumn): " + ", ".join(_ms_drop),
                    'items': _ms_drop,
                })
            if _ms_trim:
                _trim_txt = "; ".join(f"{n} (brak: {', '.join(m)})" for n, m in _ms_trim)
                report.append({
                    'severity': 'warning',
                    'category': 'Pytania matrycowe',
                    'message': "Skr\u00f3cono pytania matrycowe: " + _trim_txt,
                    'items': [n for n, _ in _ms_trim],
                })
            cleaned["matrix_sets"] = _ms_new

            # 4. box_sets
            _box = cleaned.get("box_sets", {}) or {}
            _box_drop = [c for c in list(_box.keys()) if c not in _cols]
            if _box_drop:
                for _c in _box_drop:
                    _box.pop(_c, None)
                report.append({
                    'severity': 'warning',
                    'category': 'Grupowanie odpowiedzi',
                    'message': f"Usuni\u0119to grupowania dla {len(_box_drop)} nieistniej\u0105cych kolumn: " + ", ".join(_box_drop),
                    'items': _box_drop,
                })
            cleaned["box_sets"] = _box

            # 5. weight_targets
            _wt = cleaned.get("weight_targets", {}) or {}
            _wt_drop = [c for c in list(_wt.keys()) if c not in _cols]
            _weights_reset_due_wt = False
            if _wt_drop:
                for _c in _wt_drop:
                    _wt.pop(_c, None)
                _weights_reset_due_wt = True
                report.append({
                    'severity': 'warning',
                    'category': 'Cele wa\u017cenia',
                    'message': f"Usuni\u0119to cele wa\u017cenia dla {len(_wt_drop)} kolumn: " + ", ".join(_wt_drop) + ". Wagi zostan\u0105 zresetowane.",
                    'items': _wt_drop,
                })
            cleaned["weight_targets"] = _wt

            # 6. weights
            _w = cleaned.get("weights")
            if _w is not None:
                try:
                    _wlen = len(_w)
                except TypeError:
                    _wlen = -1
                if _weights_reset_due_wt or _wlen != _n_rows:
                    cleaned["weights"] = None
                    if not _weights_reset_due_wt:
                        report.append({
                            'severity': 'critical',
                            'category': 'Wagi',
                            'message': f"Wagi w projekcie maj\u0105 {_wlen} element\u00f3w, a aktualne dane {_n_rows} wierszy \u2014 wagi zresetowane.",
                            'items': [],
                        })

            # 7. Metadata dicts keyed by column name
            _meta_specs = [
                ("custom_var_labels", "Etykiety zmiennych"),
                ("custom_val_labels", "Etykiety warto\u015bci"),
                ("custom_missing",    "Braki danych"),
                ("excel_col_types",   "Typy kolumn"),
            ]
            for _key, _label in _meta_specs:
                _d = cleaned.get(_key, {}) or {}
                _drop = [c for c in list(_d.keys()) if c not in _cols]
                if _drop:
                    for _c in _drop:
                        _d.pop(_c, None)
                    report.append({
                        'severity': 'info',
                        'category': _label,
                        'message': f"Usuni\u0119to wpisy dla {len(_drop)} nieistniej\u0105cych kolumn.",
                        'items': _drop,
                    })
                cleaned[_key] = _d
            # user_cleared_val_labels is a list
            _ucl = cleaned.get("user_cleared_val_labels", []) or []
            if isinstance(_ucl, list):
                _new_ucl  = [c for c in _ucl if c in _cols]
                _drop_ucl = [c for c in _ucl if c not in _cols]
                if _drop_ucl:
                    report.append({
                        'severity': 'info',
                        'category': 'Etykiety warto\u015bci (wyczyszczone)',
                        'message': f"Usuni\u0119to flag\u0119 dla {len(_drop_ucl)} nieistniej\u0105cych kolumn.",
                        'items': _drop_ucl,
                    })
                cleaned["user_cleared_val_labels"] = _new_ucl

            # 8. recodings
            _rec = cleaned.get("recodings", []) or []
            _rec_bad, _rec_new = [], []
            for _r in _rec:
                if not isinstance(_r, dict):
                    continue
                _src = _r.get("source") or _r.get("src") or _r.get("col")
                if _src and _src not in _cols:
                    _rec_bad.append(f"{_src} \u2192 {_r.get('new_name', '?')}")
                else:
                    _rec_new.append(_r)
            if _rec_bad:
                report.append({
                    'severity': 'warning',
                    'category': 'Rekodowania',
                    'message': f"Usuni\u0119to {len(_rec_bad)} rekodowa\u0144 (brak kolumny \u017ar\u00f3d\u0142owej): " + ", ".join(_rec_bad),
                    'items': _rec_bad,
                })
            cleaned["recodings"] = _rec_new

            # 9. cleaning_ops
            _cops = cleaned.get("cleaning_ops", []) or []
            _cops_new, _cops_trim, _cops_drop = [], 0, 0
            for _op in _cops:
                if not isinstance(_op, dict):
                    continue
                _opcols_key = 'cols' if 'cols' in _op else ('vars' if 'vars' in _op else None)
                if _opcols_key is None:
                    _cops_new.append(_op)
                    continue
                _opcols = _op.get(_opcols_key, []) or []
                _kept = [c for c in _opcols if c in _cols]
                if not _kept:
                    _cops_drop += 1
                    continue
                if len(_kept) != len(_opcols):
                    _cops_trim += 1
                    _new_op = dict(_op)
                    _new_op[_opcols_key] = _kept
                    _cops_new.append(_new_op)
                else:
                    _cops_new.append(_op)
            if _cops_drop:
                report.append({
                    'severity': 'warning',
                    'category': 'Operacje czyszczenia',
                    'message': f"Usuni\u0119to {_cops_drop} operacji czyszczenia (wszystkie kolumny brakuj\u0105).",
                    'items': [],
                })
            if _cops_trim:
                report.append({
                    'severity': 'info',
                    'category': 'Operacje czyszczenia',
                    'message': f"Skr\u00f3cono {_cops_trim} operacji czyszczenia (cz\u0119\u015b\u0107 kolumn brakuje).",
                    'items': [],
                })
            cleaned["cleaning_ops"] = _cops_new

            # 10. segmentations
            _segs = cleaned.get("segmentations", []) or []
            _segs_new, _segs_bad = [], []
            for _s in _segs:
                if not isinstance(_s, dict):
                    continue
                _svars = _s.get('vars', []) or []
                _kept = [c for c in _svars if c in _cols]
                if _kept and len(_kept) == len(_svars):
                    _segs_new.append(_s)
                else:
                    _segs_bad.append(_s.get('name', '?'))
            if _segs_bad:
                report.append({
                    'severity': 'warning',
                    'category': 'Segmentacje',
                    'message': f"Usuni\u0119to {len(_segs_bad)} segmentacji (brak wymaganych kolumn): " + ", ".join(_segs_bad),
                    'items': _segs_bad,
                })
            cleaned["segmentations"] = _segs_new

            # 11. reg_blocks
            _rb = cleaned.get("reg_blocks", []) or []
            _rb_new, _rb_trim, _rb_drop = [], 0, 0
            for _blk in _rb:
                if not isinstance(_blk, list):
                    _rb_new.append(_blk)
                    continue
                _kept = [c for c in _blk if c in _cols]
                if len(_kept) != len(_blk):
                    if _kept:
                        _rb_new.append(_kept)
                        _rb_trim += 1
                    else:
                        _rb_drop += 1
                else:
                    _rb_new.append(_blk)
            if not _rb_new:
                _rb_new = [[]]
            if _rb_trim or _rb_drop:
                report.append({
                    'severity': 'info',
                    'category': 'Bloki regresji',
                    'message': f"Zaktualizowano bloki regresji (skr\u00f3cone: {_rb_trim}, usuni\u0119te: {_rb_drop}).",
                    'items': [],
                })
            cleaned["reg_blocks"] = _rb_new

            # 12. maxdiff_pairs
            _mp = cleaned.get("maxdiff_pairs", []) or []
            _mp_new, _mp_bad = [], 0
            for _pair in _mp:
                if isinstance(_pair, (list, tuple)) and len(_pair) == 2:
                    _a, _b = _pair
                    _a_ok = (not _a) or (_a in _cols)
                    _b_ok = (not _b) or (_b in _cols)
                    if _a_ok and _b_ok:
                        _mp_new.append(list(_pair))
                    else:
                        _mp_bad += 1
                elif isinstance(_pair, (list, tuple)):
                    _mp_new.append(list(_pair))
            if _mp_bad:
                report.append({
                    'severity': 'info',
                    'category': 'Pary MaxDiff',
                    'message': f"Usuni\u0119to {_mp_bad} par MaxDiff z nieistniej\u0105cymi kolumnami.",
                    'items': [],
                })
            if not _mp_new:
                _mp_new = [('', '')]
            cleaned["maxdiff_pairs"] = _mp_new

            # 13. Results cache \u2014 detect references to missing columns/sets
            _all_known = _cols | set(cleaned.get("mrs_sets", {}).keys()) | set(cleaned.get("matrix_sets", {}).keys())
            _stale = set()
            _groups_res = (raw_data or {}).get("results", {}) or {}
            for _grp in ("czestosci", "krzyzowe", "srednie", "opisowe", "korelacje", "banner"):
                _g = _groups_res.get(_grp) or {}
                if not isinstance(_g, dict):
                    continue
                for _rkey in _g.keys():
                    _base = _rkey.split(" | ")[0].strip() if " | " in _rkey else _rkey
                    if " [" in _base:
                        _base = _base.rsplit(" [", 1)[0].strip()
                    if " x " in _base:
                        for _p in _base.split(" x "):
                            _p = _p.strip()
                            if _p and _p not in _all_known:
                                _stale.add(_p)
                    else:
                        if _base and _base not in _all_known:
                            _stale.add(_base)
            if _stale:
                _stale_list = sorted(_stale)
                report.append({
                    'severity': 'warning',
                    'category': 'Zapisane wyniki',
                    'message': f"Zapisane wyniki referuj\u0105 do {len(_stale_list)} nieistniej\u0105cych kolumn/zestaw\u00f3w: " + ", ".join(_stale_list) + ". Rozwa\u017c usuni\u0119cie starych wynik\u00f3w po wczytaniu.",
                    'items': _stale_list,
                })

            return cleaned, report

        col1, col2 = st.columns(2)

        # \u2500\u2500 SAVE \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
        with col1:
            st.markdown("**\U0001f4e5 Zapisz projekt**")

            proj_name   = st.text_input("Nazwa projektu:", key="proj_name_inp",
                                         placeholder="np. Badanie satysfakcji klient\u00f3w 2025")
            proj_author = st.text_input("Autor:", key="proj_author_inp",
                                         placeholder="np. Jan Kowalski")
            proj_desc   = st.text_area("Opis / notatki:", key="proj_desc_inp",
                                        height=70,
                                        placeholder="Opcjonalny opis badania, wersji danych itp.")

            include_res = st.checkbox(
                "Do\u0142\u0105cz wyniki analiz do pliku",
                value=True, key="proj_save_results",
                help="Tablice cz\u0119sto\u015bci, tabele krzy\u017cowe, regresje, ANOVA, EFA itp. "
                     "Plik b\u0119dzie wi\u0119kszy, ale po wczytaniu nie trzeba b\u0119dzie "
                     "ponownie wyklikywa\u0107 analiz."
            )

            _data_present = bool(st.session_state.get("imported_file"))
            include_data = st.checkbox(
                "Do\u0142\u0105cz dane \u017ar\u00f3d\u0142owe (wgrany plik) do projektu",
                value=_data_present, disabled=not _data_present,
                key="proj_save_data",
                help="Zapisuje wgrany plik (SPSS/Excel/CSV) wewn\u0105trz projektu. "
                     "Dzi\u0119ki temu po wczytaniu projektu dane s\u0105 od razu dost\u0119pne \u2014 "
                     "nie trzeba osobno wgrywa\u0107 pliku. Plik projektu b\u0119dzie wi\u0119kszy."
            )
            if not _data_present:
                st.caption("Brak wczytanych danych \u2014 najpierw wczytaj plik w module Import danych.")

            try:
                proj_json = json.dumps(
                    _build_project_data(include_results=include_res,
                                        include_data=include_data),
                    ensure_ascii=False, indent=2
                )
            except Exception as _se:
                proj_json = None
                st.error(f"B\u0142\u0105d serializacji: {_se}")

            safe_name = (proj_name or "Projekt").replace(" ", "_")[:40]
            if proj_json:
                _sz_mb = len(proj_json.encode("utf-8")) / (1024 * 1024)
                st.download_button(
                    "\U0001f4e5 Zapisz projekt (.json)",
                    data=proj_json,
                    file_name=f"{safe_name}.json",
                    mime="application/json",
                    type="primary",
                    use_container_width=True
                )
                _sz_note = " (z danymi \u017ar\u00f3d\u0142owymi)" if (include_data and _data_present) else ""
                st.caption(f"Rozmiar pliku projektu: ~{_sz_mb:.1f} MB{_sz_note}")

            n_analyses = (
                sum(1 for g in ['czestosci','krzyzowe','srednie','opisowe','korelacje']
                    if st.session_state.results.get(g)) +
                bool(st.session_state.regression_results) +
                bool(st.session_state.logistic_results) +
                bool(st.session_state.anova_results) +
                bool(st.session_state.factor_results) +
                bool(st.session_state.conjoint_results) +
                bool(st.session_state.maxdiff_results)
            )
            n_config = (len(st.session_state.mrs_sets) + len(st.session_state.matrix_sets) +
                        len(st.session_state.recodings) + len(st.session_state.segmentations) +
                        len(st.session_state.custom_var_labels) + len(st.session_state.custom_missing) +
                        sum(len(v) for v in st.session_state.box_sets.values()))
            st.caption(f"Konfiguracja: {n_config} element\u00f3w | Analizy: {n_analyses} modu\u0142\u00f3w z wynikami")

        # \u2500\u2500 LOAD \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
        with col2:
            st.markdown("**\U0001f504 Wczytaj projekt**")
            proj_file = st.file_uploader("Wgraj plik projektu (.json)", type="json",
                                          key="proj_uploader")
            if proj_file is not None:
                try:
                    raw_data = json.loads(proj_file.getvalue())
                except Exception:
                    st.error("Nieprawid\u0142owy plik JSON.")
                    raw_data = None

                if raw_data is not None:
                    ver      = raw_data.get("_version", "1.0")
                    meta     = raw_data.get("meta", {})
                    saved_at = raw_data.get("_saved_at", "")
                    src_lbl  = raw_data.get("_source", "?")
                    has_res  = "results" in raw_data
                    _emb_data = raw_data.get("imported_file")
                    _has_emb  = bool(_emb_data and _emb_data.get("bytes_b64"))

                    # Gdy projekt zawiera wlasne dane zrodlowe \u2014 waliduj wzgledem
                    # kolumn ZAPISANYCH w projekcie (a nie aktualnie wczytanych,
                    # ktore moga byc puste przy wczytywaniu projektu bez danych).
                    if _has_emb and raw_data.get("_data_columns"):
                        _val_target = pd.DataFrame(columns=raw_data["_data_columns"])
                    else:
                        _val_target = df_raw
                    cleaned_data, _val_report = _validate_project_data(raw_data, _val_target)
                    _n_crit = sum(1 for r in _val_report if r['severity'] == 'critical')
                    _n_warn = sum(1 for r in _val_report if r['severity'] == 'warning')
                    _n_info = sum(1 for r in _val_report if r['severity'] == 'info')

                    with st.expander("Podgl\u0105d wczytywanego projektu", expanded=True):
                        if meta.get("name"):
                            st.markdown(f"**Nazwa:** {meta['name']}")
                        if meta.get("author"):
                            st.markdown(f"**Autor:** {meta['author']}")
                        if meta.get("desc"):
                            st.markdown(f"**Opis:** {meta['desc']}")
                        st.caption(
                            f"Wersja: {ver} \u00b7 "
                            f"\u0179r\u00f3d\u0142o: {'SPSS' if src_lbl=='spss' else 'Excel' if src_lbl=='excel' else src_lbl} \u00b7 "
                            f"Zapisano: {saved_at or 'nieznane'}"
                        )
                        if has_res:
                            n_saved = sum(
                                len(raw_data.get("results", {}).get(g, {}))
                                for g in ['czestosci','krzyzowe','srednie','opisowe','korelacje']
                            )
                            st.success(f"\U0001f4ca Plik zawiera wyniki analiz ({n_saved} tabel).")
                        else:
                            st.info("Plik zawiera tylko konfiguracj\u0119 (bez wynik\u00f3w analiz).")
                        if _has_emb:
                            _emb_name = _emb_data.get("name", "plik")
                            _emb_kind = (_emb_data.get("kind") or "").upper()
                            st.success(f"\U0001f4e6 Projekt zawiera dane \u017ar\u00f3d\u0142owe: **{_emb_name}** "
                                       f"({_emb_kind}) \u2014 zostan\u0105 wczytane wraz z projektem.")

                        summary_parts = []
                        for key, label in [("mrs_sets","zestawy MRS"),("matrix_sets","pytania matrycowe"),
                                           ("recodings","rekodowania"),("segmentations","segmentacje")]:
                            n = len(raw_data.get(key, {}))
                            if n: summary_parts.append(f"{n} {label}")
                        if summary_parts:
                            st.info(", ".join(summary_parts))

                        st.markdown("---")
                        st.markdown("**\U0001f50d Walidacja zgodnosci z aktualnymi danymi:**")
                        if _n_crit == 0 and _n_warn == 0 and _n_info == 0:
                            st.success("\u2705 Projekt w pelni zgodny z aktualnymi danymi.")
                        else:
                            _crit_list = [r for r in _val_report if r['severity'] == 'critical']
                            _warn_list = [r for r in _val_report if r['severity'] == 'warning']
                            _info_list = [r for r in _val_report if r['severity'] == 'info']
                            if _crit_list:
                                st.error("\u274c Niezgodnosci krytyczne (wplywaja na dzialanie modulow):")
                                for _r in _crit_list:
                                    st.markdown(f"- **{_r['category']}**: {_r['message']}")
                            if _warn_list:
                                st.warning("\u26a0\ufe0f Wpisy zostana zmodyfikowane:")
                                for _r in _warn_list:
                                    st.markdown(f"- **{_r['category']}**: {_r['message']}")
                            if _info_list:
                                st.info("\u2139\ufe0f Wpisy zostana oczyszczone (tylko metadane):")
                                for _r in _info_list:
                                    st.markdown(f"- **{_r['category']}**: {_r['message']}")
                            if has_res and (_n_crit or _n_warn):
                                st.info(
                                    "\u2139\ufe0f Wyniki analiz pozostana zachowane, ale moga pochodzic ze starej wersji danych "
                                    "\u2014 rozwaz ich usuniecie w poszczegolnych modulach."
                                )

                    _restore_data = False
                    if _has_emb:
                        _restore_data = st.checkbox(
                            "Wczytaj dane \u017ar\u00f3d\u0142owe zapisane w projekcie",
                            value=True, key="proj_load_data",
                            help="Odtwarza wgrany plik (SPSS/Excel/CSV) z projektu \u2014 "
                                 "nie trzeba osobno wgrywa\u0107 danych. Odznacz, je\u015bli chcesz "
                                 "zachowa\u0107 aktualnie wczytane dane.")

                    _btn_caption = "\u2705 Przywr\u00f3\u0107 z pliku"
                    if _n_crit or _n_warn:
                        _btn_caption = "\u2705 Przywr\u00f3\u0107 z pliku (z automatyczn\u0105 korekt\u0105 niezgodno\u015bci)"
                    if st.button(_btn_caption,
                                  type="primary", use_container_width=True):
                        # Dane zrodlowe osadzone w projekcie (przed konfiguracja,
                        # by pipeline przeliczyl sie na nowych danych po rerun)
                        if _has_emb and _restore_data:
                            import base64 as _b64
                            _imp_new = {_k: _v for _k, _v in _emb_data.items()
                                        if _k != "bytes_b64"}
                            _imp_new["bytes"] = _b64.b64decode(_emb_data["bytes_b64"])
                            st.session_state.imported_file = _imp_new
                            st.session_state.pop('_wd_cache', None)
                            st.session_state.pop('_wd_key', None)
                        # Config
                        st.session_state.mrs_sets          = cleaned_data.get("mrs_sets", {})
                        st.session_state.matrix_sets       = cleaned_data.get("matrix_sets", {})
                        st.session_state.custom_var_labels = cleaned_data.get("custom_var_labels", {})
                        st.session_state.custom_val_labels = cleaned_data.get("custom_val_labels", {})
                        st.session_state.ppt_chart_templates = cleaned_data.get("ppt_chart_templates", {})
                        st.session_state.box_sets          = defaultdict(dict, cleaned_data.get("box_sets", {}))
                        st.session_state.value_orders      = cleaned_data.get("value_orders", {})
                        st.session_state.segmentations     = cleaned_data.get("segmentations", [])
                        st.session_state.recodings         = cleaned_data.get("recodings", [])
                        st.session_state.cleaning_ops      = cleaned_data.get("cleaning_ops", [])
                        st.session_state.custom_missing    = cleaned_data.get("custom_missing", {})
                        st.session_state.user_cleared_val_labels = set(cleaned_data.get("user_cleared_val_labels", []))
                        st.session_state.weight_targets    = cleaned_data.get("weight_targets", {})
                        st.session_state.treat_empty_as_miss = cleaned_data.get("treat_empty_as_miss", False)
                        st.session_state.excel_col_types   = cleaned_data.get("excel_col_types", {})
                        st.session_state.maxdiff_pairs     = cleaned_data.get("maxdiff_pairs", [('', '')])
                        st.session_state.reg_blocks        = cleaned_data.get("reg_blocks", [[]])
                        st.session_state.split_var         = cleaned_data.get("split_var", None)
                        w = cleaned_data.get("weights")
                        st.session_state.weights = np.array(w) if w else None
                        # Results
                        if has_res:
                            raw_res = cleaned_data.get("results", {})
                            st.session_state.results = {
                                'czestosci': _deser_results(raw_res.get('czestosci', {})),
                                'krzyzowe':  _deser_results(raw_res.get('krzyzowe', {})),
                                'srednie':   _deser_results(raw_res.get('srednie', {})),
                                'opisowe':   _deser_results(raw_res.get('opisowe', {})),
                                'korelacje': _deser_results(raw_res.get('korelacje', {})),
                                'banner':    _deser_results(raw_res.get('banner', {})),
                            }
                            st.session_state.chi_results        = cleaned_data.get("chi_results", {})
                            st.session_state.matrix_results     = _deser_matrix_results(cleaned_data.get("matrix_results", []))
                            st.session_state.regression_results = _deser_reg_results(cleaned_data.get("regression_results", []))
                            st.session_state.logistic_results   = _deser_reg_results(cleaned_data.get("logistic_results", []))
                            st.session_state.anova_results      = _deser_reg_results(cleaned_data.get("anova_results", []))
                            st.session_state.factor_results     = _deser_factor_results(cleaned_data.get("factor_results", []))
                            st.session_state.conjoint_results   = _deser_factor_results(cleaned_data.get("conjoint_results", []))
                            st.session_state.maxdiff_results    = _deser_factor_results(cleaned_data.get("maxdiff_results", []))
                            st.session_state.hclust_results     = _deser_hclust(cleaned_data.get("hclust_results", []))
                        st.session_state._proj_loaded_ok = True
                        st.rerun()

    with tab_summary:
        st.markdown("#### Podsumowanie bazy danych")
        total_rows = len(df_raw)
        total_cols_n = len(df_raw.columns)
        num_c = len(df_raw.select_dtypes(include=[np.number]).columns)
        cat_c = total_cols_n - num_c
        complete = int(df_raw.dropna().shape[0])
        miss_cells = int(df_raw.isna().sum().sum())
        total_cells = total_rows * total_cols_n
        miss_pct = miss_cells / total_cells * 100 if total_cells > 0 else 0

        m1, m2, m3, m4 = st.columns(4)
        m1.metric("\U0001f465 Respondenci", f"{total_rows:,}")
        m2.metric("\U0001f4ca Zmienne", f"{total_cols_n:,}")
        m3.metric("\U0001f522 Numeryczne", f"{num_c:,}")
        m4.metric("\U0001f524 Kategoryczne", f"{cat_c:,}")
        m5, m6, m7, m8 = st.columns(4)
        m5.metric("\u2705 Kompletne wiersze", f"{complete:,}")
        m6.metric("\u26a0\ufe0f Wiersze z brakami", f"{total_rows - complete:,}")
        m7.metric("\U0001f573\ufe0f Kom\u00f3rki z NaN", f"{miss_cells:,}")
        m8.metric("\U0001f4c9 % brak\u00f3w w bazie", f"{miss_pct:.1f}%")

        st.divider()
        st.markdown("**Braki danych i statystyki wg zmiennej:**")
        summ_rows = []
        for c in df_raw.columns:
            n_miss = df_raw[c].isna().sum()
            summ_rows.append({
                'Zmienna': c, 'Etykieta': var_labels.get(c, ''),
                'Typ danych': str(df_raw[c].dtype),
                'Braki [N]': n_miss,
                'Braki [%]': round(n_miss / total_rows * 100, 1),
                'Unikalne warto\u015bci': df_raw[c].nunique(),
                'Min': df_raw[c].min() if pd.api.types.is_numeric_dtype(df_raw[c]) else '--',
                'Max': df_raw[c].max() if pd.api.types.is_numeric_dtype(df_raw[c]) else '--',
            })
        st.dataframe(pd.DataFrame(summ_rows), use_container_width=True, height=400)

    with tab_dict:
        st.markdown("#### S\u0142ownik zmiennych")
        st.info("Edycja etykiet zmiennych i warto\u015bci dost\u0119pna w module "
                "**Przygotowanie Danych \u2192 Etykiety**.")
        dict_rows = []
        for col in df_raw.columns:
            orig_lbl = meta_orig.column_names_to_labels.get(col, "")
            curr_lbl = var_labels.get(col, orig_lbl)
            has_vl   = (bool(meta_orig.variable_value_labels.get(col))
                        or bool(st.session_state.custom_val_labels.get(col)))
            is_custom_var = col in st.session_state.custom_var_labels
            is_custom_val = col in st.session_state.custom_val_labels
            dict_rows.append({
                "Zmienna": col,
                "Etykieta bie\u017c\u0105ca": curr_lbl,
                "Etykieta oryginalna": orig_lbl,
                "Typ": str(df_raw[col].dtype),
                "Et. warto\u015bci": "\u2705" if has_vl else "--",
                "Zmodyfikowana": ("\U0001f3f7\ufe0f+\u270f\ufe0f" if is_custom_var and is_custom_val
                                  else "\U0001f3f7\ufe0f" if is_custom_val
                                  else "\u270f\ufe0f" if is_custom_var else ""),
            })
        st.dataframe(pd.DataFrame(dict_rows), use_container_width=True, height=420)
        st.caption("\u270f\ufe0f = zmieniona etykieta zmiennej | \U0001f3f7\ufe0f = zmienione etykiety warto\u015bci")


    # -- SZABLONY WYKRESOW (tworzone raz, wybierane przy eksporcie) --
    with tab_charts:
        st.session_state.setdefault('ppt_chart_templates', {})
        _cts_store = st.session_state.ppt_chart_templates

        with st.expander("Instrukcja \u2014 jak dzia\u0142aj\u0105 szablony wykres\u00f3w", expanded=False):
            st.markdown("""
##### \U0001f3af Do czego s\u0142u\u017cy ta zak\u0142adka?
Tworzysz tu **szablony wygl\u0105du wykres\u00f3w** raz, a nast\u0119pnie u\u017cywasz ich wielokrotnie
w eksporcie do **PowerPoint** i **Worda**. Dzi\u0119ki temu wszystkie wykresy w raporcie
maj\u0105 sp\u00f3jn\u0105, profesjonaln\u0105 opraw\u0119.

##### \U0001f527 Jak korzysta\u0107
1. **Wybierz gotow\u0105 palet\u0119** lub ustaw w\u0142asne kolory (8 serii).
2. Dostosuj **etykiety danych, typografi\u0119, legend\u0119 i osie**.
3. Sprawd\u017a **podgl\u0105d na \u017cywo** po prawej stronie.
4. Nadaj **nazw\u0119** i kliknij **Zapisz szablon**.
5. W module eksportu wybierzesz zapisany szablon z listy.

##### \U0001f4a1 Wskaz\u00f3wki
- Aby **edytowa\u0107** istniej\u0105cy szablon, kliknij **Edytuj** na jego kafelku w galerii.
- Szablony zapisuj\u0105 si\u0119 razem z projektem (plik `.json`).
""")

        # -- wartosci poczatkowe edytora (raz na sesje) --
        _CTS_DEFAULTS = {
            "cts_c_0": CHART_TPL_DEFAULT["colors"][0],
            "cts_c_1": CHART_TPL_DEFAULT["colors"][1],
            "cts_c_2": CHART_TPL_DEFAULT["colors"][2],
            "cts_c_3": CHART_TPL_DEFAULT["colors"][3],
            "cts_c_4": CHART_TPL_DEFAULT["colors"][4],
            "cts_c_5": CHART_TPL_DEFAULT["colors"][5],
            "cts_c_6": CHART_TPL_DEFAULT["colors"][6],
            "cts_c_7": CHART_TPL_DEFAULT["colors"][7],
            "cts_title_color": CHART_TPL_DEFAULT["title_color"],
            "cts_font_title":  CHART_TPL_DEFAULT["font_size_title"],
            "cts_font_labels": CHART_TPL_DEFAULT["font_size_labels"],
            "cts_font_data":   CHART_TPL_DEFAULT["font_size_data"],
            "cts_show_dls":    CHART_TPL_DEFAULT["show_data_labels"],
            "cts_dls_bold":    CHART_TPL_DEFAULT["data_label_bold"],
            "cts_dls_fmt":     CHART_TPL_DEFAULT["data_label_format"],
            "cts_legend":      CHART_TPL_DEFAULT["legend_position"],
            "cts_show_y":      CHART_TPL_DEFAULT["show_y_axis"],
            "cts_show_grid":   CHART_TPL_DEFAULT["show_gridlines"],
            "cts_show_x":      CHART_TPL_DEFAULT["show_x_axis"],
            "cts_chart_type":  "czestosci",
        }
        for _ck, _cv in _CTS_DEFAULTS.items():
            st.session_state.setdefault(_ck, _cv)

        # -- zastosuj odroczone zadania (preset / wczytanie) PRZED widgetami --
        _cts_req_preset = st.session_state.pop("_cts_preset_req", None)
        if _cts_req_preset and _cts_req_preset in CHART_TPL_PRESETS:
            for _pi2, _pc2 in enumerate(CHART_TPL_PRESETS[_cts_req_preset]):
                st.session_state["cts_c_" + str(_pi2)] = _pc2
        _cts_req_load = st.session_state.pop("_cts_load_req", None)
        if _cts_req_load and _cts_req_load in _cts_store:
            _lt = _cts_store[_cts_req_load]
            _ltc = _lt.get("colors", CHART_TPL_DEFAULT["colors"])
            for _li in range(8):
                st.session_state["cts_c_" + str(_li)] = (
                    _ltc[_li] if _li < len(_ltc) else CHART_TPL_DEFAULT["colors"][_li])
            st.session_state["cts_title_color"] = _lt.get("title_color", CHART_TPL_DEFAULT["title_color"])
            st.session_state["cts_font_title"]  = int(_lt.get("font_size_title", 14))
            st.session_state["cts_font_labels"] = int(_lt.get("font_size_labels", 10))
            st.session_state["cts_font_data"]   = int(_lt.get("font_size_data", 9))
            st.session_state["cts_show_dls"]    = bool(_lt.get("show_data_labels", True))
            st.session_state["cts_dls_bold"]    = bool(_lt.get("data_label_bold", True))
            st.session_state["cts_dls_fmt"]     = _lt.get("data_label_format", "auto")
            st.session_state["cts_legend"]      = _lt.get("legend_position", "bottom")
            st.session_state["cts_show_y"]      = bool(_lt.get("show_y_axis", False))
            st.session_state["cts_show_grid"]   = bool(_lt.get("show_gridlines", False))
            st.session_state["cts_show_x"]      = bool(_lt.get("show_x_axis", True))
            st.session_state["cts_chart_type"]  = _lt.get("chart_type", "czestosci")
            st.session_state["cts_name"]        = _cts_req_load

        def _cts_collect():
            return {
                "colors": [st.session_state["cts_c_" + str(_i)] for _i in range(8)],
                "title_color": st.session_state["cts_title_color"],
                "font_size_title": int(st.session_state["cts_font_title"]),
                "font_size_labels": int(st.session_state["cts_font_labels"]),
                "font_size_data": int(st.session_state["cts_font_data"]),
                "show_data_labels": bool(st.session_state["cts_show_dls"]),
                "data_label_format": st.session_state["cts_dls_fmt"],
                "data_label_bold": bool(st.session_state["cts_dls_bold"]),
                "legend_position": st.session_state["cts_legend"],
                "show_gridlines": bool(st.session_state["cts_show_grid"]),
                "show_y_axis": bool(st.session_state["cts_show_y"]),
                "show_x_axis": bool(st.session_state["cts_show_x"]),
                "bar_bold_labels": True,
                "chart_type": st.session_state.get("cts_chart_type", "czestosci"),
            }

        # \u2500\u2500 Typ wykresu \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
        st.markdown("")
        st.radio(
            "Typ wykresu:",
            ["czestosci", "krzyzowe"],
            format_func=lambda x: {
                "czestosci": "\U0001f4ca Tablice cz\u0119sto\u015bci \u2014 jeden s\u0142upek na kategori\u0119",
                "krzyzowe":  "\U0001f500 Tabele krzy\u017cowe \u2014 grupowane s\u0142upki (wiele serii)",
            }[x],
            key="cts_chart_type", horizontal=True)

        st.markdown("---")

        # \u2500\u2500 Kontrolki w 3 kolumnach \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
        _cts_col1, _cts_col2, _cts_col3 = st.columns([1.2, 1.2, 1], gap="large")

        with _cts_col1:
            _cts_is_freq = st.session_state.get("cts_chart_type", "czestosci") == "czestosci"
            if _cts_is_freq:
                st.markdown("**Kolor s\u0142upk\u00f3w**")
                st.color_picker("Kolor s\u0142upk\u00f3w:", key="cts_c_0",
                                label_visibility="collapsed")
                for _cci in range(1, 8):
                    st.session_state.setdefault("cts_c_" + str(_cci),
                                                CHART_TPL_DEFAULT["colors"][_cci])
            else:
                st.markdown("**Kolory serii (8)**")
                _cts_ccols = st.columns(4)
                for _cci in range(8):
                    _cts_ccols[_cci % 4].color_picker(
                        str(_cci + 1), key="cts_c_" + str(_cci))

        with _cts_col2:
            st.markdown("**Etykiety danych i typografia**")
            st.color_picker("Kolor etykiet / tytu\u0142u:", key="cts_title_color")
            _fc1, _fc2 = st.columns(2)
            _fc1.number_input("Tytu\u0142 (pt)", min_value=8, max_value=32,
                              step=1, key="cts_font_title")
            _fc2.number_input("Etykiety osi (pt)", min_value=6, max_value=20,
                              step=1, key="cts_font_labels")
            _fc3, _fc4 = st.columns(2)
            _fc3.number_input("Etykiety danych (pt)", min_value=6, max_value=20,
                              step=1, key="cts_font_data")
            _fc4.markdown("<div style='height:8px'></div>", unsafe_allow_html=True)
            _fc4.checkbox("Etykiety danych", key="cts_show_dls")
            _fc4.checkbox("Pogrubione", key="cts_dls_bold")
            st.radio("Format etykiet:", ["auto", "percent", "number"],
                     format_func=lambda x: {"auto": "Auto",
                                            "percent": "Procent (%)",
                                            "number": "Liczba"}[x],
                     horizontal=True, key="cts_dls_fmt")

        with _cts_col3:
            st.markdown("**Legenda, osie i siatka**")
            st.selectbox("Pozycja legendy",
                         ["bottom", "top", "right", "left", "none"],
                         format_func=lambda x: {"bottom": "D\u00f3\u0142", "top": "G\u00f3ra",
                                                "right": "Prawo", "left": "Lewo",
                                                "none": "Ukryj"}[x],
                         key="cts_legend")
            st.checkbox("Poka\u017c o\u015b X", key="cts_show_x")
            st.checkbox("Poka\u017c o\u015b Y", key="cts_show_y")
            st.checkbox("Linie siatki", key="cts_show_grid")

        st.markdown("---")

        # \u2500\u2500 Nazwa + przyciski zapisu \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
        _cts_nb1, _cts_nb2, _cts_nb3 = st.columns([3, 1.2, 1.2])
        _cts_name = _cts_nb1.text_input(
            "Nazwa szablonu", key="cts_name",
            placeholder="np. Korporacyjny / Prosty",
            label_visibility="collapsed")
        _cts_nm = _cts_name.strip()
        _cts_nb1.caption("Wpisz nazw\u0119 szablonu \u2014 zostanie zapisany po klikni\u0119ciu przycisku.")
        with _cts_nb2:
            st.markdown("<div style='height:4px'></div>", unsafe_allow_html=True)
            if st.button("\U0001f4be Zapisz szablon", type="primary",
                         use_container_width=True, key="cts_save",
                         disabled=(not _cts_nm)):
                _cts_store[_cts_nm] = _cts_collect()
                st.success("\u2705 Zapisano: " + _cts_nm)
                st.rerun()
        with _cts_nb3:
            st.markdown("<div style='height:4px'></div>", unsafe_allow_html=True)
            if _cts_nm in _cts_store:
                if st.button("\U0001f5d1\ufe0f Usu\u0144", use_container_width=True,
                             key="cts_del_cur"):
                    _cts_store.pop(_cts_nm, None)
                    st.rerun()

        # \u2500\u2500 Podgl\u0105d pe\u0142n\u0105 szeroko\u015bci\u0105 \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
        st.markdown("**Podgl\u0105d wykresu:**")
        _cts_cur = _cts_collect()
        _cts_typ = st.session_state.get("cts_chart_type", "czestosci")
        if _cts_typ == "krzyzowe":
            _cts_pdf = pd.DataFrame({
                "Kategoria": ["Kat. A", "Kat. B", "Kat. C"],
                "Seria 1": [42.0, 28.0, 18.0],
                "Seria 2": [35.0, 38.0, 24.0],
                "Seria 3": [23.0, 34.0, 58.0],
            })
            _cts_fig = px.bar(
                _cts_pdf, x="Kategoria",
                y=["Seria 1", "Seria 2", "Seria 3"],
                barmode="group",
                color_discrete_sequence=_cts_cur["colors"])
            _cts_ymax = 70.0
        else:
            _cts_pdf = pd.DataFrame({
                "Kategoria": ["Bardzo dobrze", "Dobrze", "\u015arednio", "\u0179le"],
                "Warto\u015b\u0107": [38.0, 34.0, 19.0, 9.0],
            })
            _cts_fig = px.bar(_cts_pdf, x="Kategoria", y="Warto\u015b\u0107")
            for _cts_tr in _cts_fig.data:
                _cts_tr.marker.color = _cts_cur["colors"][0]
            _cts_ymax = float(_cts_pdf["Warto\u015b\u0107"].max()) * 1.25
        _cts_fig.update_layout(
            height=380, margin=dict(l=20, r=20, t=30, b=50),
            showlegend=(_cts_cur["legend_position"] != "none"),
            xaxis=dict(visible=_cts_cur["show_x_axis"],
                       tickfont=dict(size=_cts_cur["font_size_labels"])),
            yaxis=dict(visible=_cts_cur["show_y_axis"],
                       showgrid=_cts_cur["show_gridlines"],
                       range=[0, _cts_ymax],
                       tickfont=dict(size=_cts_cur["font_size_labels"])),
            plot_bgcolor="white", paper_bgcolor="#fafafa",
            legend=dict(
                orientation="h" if _cts_cur["legend_position"] in ("bottom", "top") else "v",
                y=-0.2 if _cts_cur["legend_position"] == "bottom" else
                  1.05 if _cts_cur["legend_position"] == "top" else 0.5,
                x=0.5 if _cts_cur["legend_position"] in ("bottom", "top") else
                  1.02 if _cts_cur["legend_position"] == "right" else -0.15,
                xanchor="center" if _cts_cur["legend_position"] in ("bottom", "top") else "left",
            ) if _cts_cur["legend_position"] != "none" else dict(visible=False),
        )
        if _cts_cur["show_data_labels"]:
            if _cts_typ == "krzyzowe":
                _cts_fig.update_traces(texttemplate="%{y:.1f}%", textposition="outside",
                                       textfont=dict(size=_cts_cur["font_size_data"],
                                                     color=_cts_cur["title_color"]))
            else:
                _cts_yvals = _cts_pdf["Warto\u015b\u0107"].tolist()
                if _cts_cur["data_label_format"] == "number":
                    _cts_txt = [f"{_v:.0f}" for _v in _cts_yvals]
                else:
                    _cts_txt = [f"{_v:.1f}%" for _v in _cts_yvals]
                _cts_fig.update_traces(text=_cts_txt, textposition="outside",
                                       textfont=dict(size=_cts_cur["font_size_data"],
                                                     color=_cts_cur["title_color"]))
        # Pasek kolorow pod wykresem
        _cts_sw = "".join(
            "<span style='display:inline-block;width:22px;height:22px;background:"
            + _sc + ";border-radius:4px;margin:2px;border:1px solid #e0e0e0;' "
            + "title='" + _sc + "'></span>"
            for _sc in _cts_cur["colors"])
        st.plotly_chart(_cts_fig, use_container_width=True, key="cts_preview")
        st.markdown(
            "<div style='margin:2px 0 8px;display:flex;align-items:center;gap:6px;'>"
            "<span style='font-size:.78rem;color:#888;'>Paleta:</span>"
            + _cts_sw + "</div>",
            unsafe_allow_html=True)

        # -- galeria zapisanych szablonow --
        st.divider()
        st.markdown("#### Zapisane szablony")
        if not _cts_store:
            st.info("Nie masz jeszcze \u017cadnego zapisanego szablonu. "
                    "Ustaw wygl\u0105d powy\u017cej i kliknij **Zapisz szablon**.")
        else:
            _cts_gnames = list(_cts_store.keys())
            _cts_legmap = {"bottom": "d\u00f3\u0142", "top": "g\u00f3ra", "right": "prawo",
                           "left": "lewo", "none": "brak"}
            for _cts_gi in range(0, len(_cts_gnames), 3):
                _cts_gcols = st.columns(3)
                for _cts_gj, _cts_gn in enumerate(_cts_gnames[_cts_gi:_cts_gi + 3]):
                    _cts_td = _cts_store[_cts_gn]
                    with _cts_gcols[_cts_gj]:
                        with st.container(border=True):
                            _cts_gsw = "".join(
                                "<span style='display:inline-block;width:20px;height:20px;"
                                "background:" + _gc + ";border-radius:4px;margin:1px;"
                                "border:1px solid #e0e0e0;'></span>"
                                for _gc in _cts_td.get("colors", [])[:8])
                            _cts_leg = _cts_legmap.get(
                                _cts_td.get("legend_position", "bottom"), "d\u00f3\u0142")
                            _cts_dl = "tak" if _cts_td.get("show_data_labels") else "nie"
                            _cts_gr = "tak" if _cts_td.get("show_gridlines") else "nie"
                            _cts_gtype = _cts_td.get("chart_type", "czestosci")
                            _cts_gtype_lbl = ("Tablice cz\u0119sto\u015bci"
                                               if _cts_gtype == "czestosci"
                                               else "Tabele krzy\u017cowe")
                            _cts_gtype_col = "#1F4E79" if _cts_gtype == "czestosci" else "#7030A0"
                            st.markdown(
                                "<div style='font-weight:700;color:#1F4E79;font-size:.95rem;"
                                "margin-bottom:2px;overflow:hidden;text-overflow:ellipsis;"
                                "white-space:nowrap;'>\U0001f516 " + _cts_gn + "</div>"
                                "<div style='display:inline-block;background:" + _cts_gtype_col + ";"
                                "color:white;font-size:.7rem;padding:1px 6px;border-radius:3px;"
                                "margin-bottom:5px;'>" + _cts_gtype_lbl + "</div>"
                                "<div style='margin-bottom:6px;'>" + _cts_gsw + "</div>"
                                "<div style='font-size:.74rem;color:#666;'>"
                                "Legenda: <b>" + _cts_leg + "</b> \u00b7 "
                                "Etykiety: <b>" + _cts_dl + "</b> \u00b7 "
                                "Siatka: <b>" + _cts_gr + "</b></div>",
                                unsafe_allow_html=True)
                            _cts_ga, _cts_gb = st.columns(2)
                            if _cts_ga.button("\u270f\ufe0f Edytuj",
                                              key="cts_g_edit_" + _cts_gn,
                                              use_container_width=True):
                                st.session_state["_cts_load_req"] = _cts_gn
                                st.rerun()
                            if _cts_gb.button("\U0001f5d1\ufe0f Usu\u0144",
                                              key="cts_g_del_" + _cts_gn,
                                              use_container_width=True):
                                _cts_store.pop(_cts_gn, None)
                                st.rerun()

# -------------------------------------------------------------
# MODU? 2: PRZYGOTOWANIE DANYCH
# -------------------------------------------------------------
elif menu == "\U0001f6e0\ufe0f Przygotowanie Danych":
    _require_module_access("prep")
    _require_data()
    module_header("\U0001f6e0\ufe0f", "Przygotowanie Danych")
    # For Excel: add an extra tab for type overrides
    if is_tabular:
        tab_miss, tab_labels, tab_types, tab_clean, tab_mrs, tab_matrix, tab_weight, tab_recode, tab_box, tab_order, tab_split = st.tabs([
            'Braki', 'Etykiety', 'Typy', 'Czyszczenie',
            'Wielokrotne odp.', 'Matrycowe',
            'Wa\u017cenie', 'Rekodowanie', 'Grupowanie odpowiedzi',
            'Kolejno\u015b\u0107 warto\u015bci',
            'Podzia\u0142 na podzbiory'
        ])
    else:
        tab_miss, tab_labels, tab_clean, tab_mrs, tab_matrix, tab_weight, tab_recode, tab_box, tab_order, tab_split = st.tabs([
            'Braki', 'Etykiety', 'Czyszczenie',
            'Wielokrotne odp.', 'Matrycowe',
            'Wa\u017cenie', 'Rekodowanie', 'Grupowanie odpowiedzi',
            'Kolejno\u015b\u0107 warto\u015bci',
            'Podzia\u0142 na podzbiory'
        ])
        tab_types = None

    # -- BRAKI DANYCH ---------------------------------
    with tab_miss:
        st.markdown("#### Definiowanie brak\u00f3w danych")

        with st.expander("Instrukcja \u2014 jak definiowa\u0107 braki danych", expanded=False):
            st.markdown("""
##### \U0001f3af Do czego s\u0142u\u017cy ta zak\u0142adka?
Pozwala okre\u015bli\u0107, kt\u00f3re warto\u015bci w bazie danych maj\u0105 by\u0107 traktowane
jako **braki danych** (NaN / missing). To kluczowy krok \u2014 bez niego narz\u0119dzie
mo\u017ce w\u0142\u0105cza\u0107 do analiz niepoprawne lub puste warto\u015bci.

##### \U0001f527 Dost\u0119pne opcje
1. **Puste ci\u0105gi tekstowe jako NaN** \u2014 zaznacz checkbox, je\u015bli chcesz,
   aby puste kom\u00f3rki tekstowe (`""`, `" "`) by\u0142y traktowane jak brak odpowiedzi.
   Zalecane przy danych z ankiet tekstowych.
2. **Domy\u015blne braki z pliku SPSS** *(tylko .sav)* \u2014 plik SPSS mo\u017ce
   zawiera\u0107 zakodowane warto\u015bci brak\u00f3w (np. 99 = "brak odpowiedzi"). Zaznacz
   t\u0119 opcj\u0119, aby automatycznie je rozpozna\u0107.
3. **Niestandardowe warto\u015bci brak\u00f3w** \u2014 dla danych Excel lub gdy SPSS nie
   ma zakodowanych brak\u00f3w: wybierz zmienn\u0105 i podaj warto\u015bci, kt\u00f3re maj\u0105
   by\u0107 traktowane jako braki (np. 999, -1, "nie dotyczy").

##### \U0001f4a1 Wa\u017cne
- Braki danych s\u0105 stosowane globalnie \u2014 wp\u0142ywaj\u0105 na **wszystkie modu\u0142y** analiz.
- Warto zdefiniowa\u0107 braki **przed** tworzeniem wag i wykonaniem analiz.
- Ustawienia s\u0105 zapisywane w pliku projektu JSON.
"""
            )

        st.session_state.treat_empty_as_miss = st.checkbox(
            "Traktuj puste ci\u0105gi tekstowe jako braki danych",
            value=st.session_state.treat_empty_as_miss,
            help="Zaznacz, je\u015bli puste kom\u00f3rki tekstowe maj\u0105 by\u0107 traktowane jako brak odpowiedzi."
        )

        # SPSS-specific: default missing values from file
        if is_spss:
            use_spss_missing = st.checkbox("U\u017cywaj domy\u015blnych brak\u00f3w danych z pliku SPSS (zalecane)", value=True)
        else:
            use_spss_missing = False
            st.info("\U0001f4c8 Dane Excel nie maj\u0105 wbudowanych brak\u00f3w danych. Zdefiniuj warto\u015bci brak\u00f3w poni\u017cej.")

        if not use_spss_missing:
            col1, col2 = st.columns(2)
            with col1:
                st.markdown("**Globalne braki danych**")
                global_missing_str = st.text_input("Warto\u015bci oddzielone przecinkiem (np. 98, 99):")
                if st.button("Zastosuj do wszystkich zmiennych", use_container_width=True):
                    try:
                        vals = [float(x.strip()) for x in global_missing_str.split(',') if x.strip()]
                        for c in df_raw.columns:
                            st.session_state.custom_missing[c] = vals
                        st.success(f"Zastosowano braki {vals} do wszystkich zmiennych.")
                    except Exception:
                        st.error("Wprowad\u017a poprawne liczby.")
            with col2:
                st.markdown("**Braki dla konkretnej zmiennej**")
                var_missing = st.selectbox("Zmienna:", df_raw.columns,
                    format_func=lambda x: get_var_display_name(x, var_labels),
                    key="missing_var_select")
                var_missing_vals = st.multiselect("Warto\u015bci traktowane jako braki:",
                    sorted(df_raw[var_missing].dropna().unique(), key=lambda x: str(x)))
                if st.button("Ustaw braki dla tej zmiennej", use_container_width=True):
                    st.session_state.custom_missing[var_missing] = var_missing_vals
                    st.success("Zapisano!")

            if st.session_state.custom_missing:
                st.divider()
                st.markdown(f"**\U0001f4cb Zapisane braki danych ({len(st.session_state.custom_missing)} zmiennych):**")

                for _col, _vals in list(st.session_state.custom_missing.items()):
                    mc1, mc2, mc3, mc4 = st.columns([3, 3, 2, 1])
                    mc1.markdown(f"`{_col}` {var_labels.get(_col, '')}")
                    new_miss_str = mc2.text_input(
                        "", value=", ".join(str(v) for v in _vals),
                        key=f"miss_edit_{_col}", label_visibility="collapsed"
                    )
                    with mc3:
                        if st.button("\U0001f4be Zapisz", key=f"miss_upd_{_col}",
                                     use_container_width=True):
                            try:
                                new_vals = [float(x.strip())
                                            for x in new_miss_str.split(',') if x.strip()]
                                st.session_state.custom_missing[_col] = new_vals
                                st.rerun()
                            except Exception:
                                st.error("Nieprawid\u0142owe warto\u015bci.")
                    with mc4:
                        if st.button("\U0001f5d1\ufe0f", key=f"miss_del_{_col}",
                                     help=f"Usu\u0144 braki dla {_col}"):
                            st.session_state.custom_missing.pop(_col, None)
                            st.rerun()

                if st.button("\U0001f5d1\ufe0f Usu\u0144 wszystkie braki danych", type="secondary",
                              use_container_width=True, key="miss_clear_all"):
                    st.session_state.custom_missing = {}
                    st.rerun()

    # -- REKODOWANIE -----------------------------------
    with tab_recode:
        st.markdown("#### Rekodowanie zmiennych")

        with st.expander("Instrukcja \u2014 jak rekodowa\u0107 zmienne", expanded=False):
            st.markdown("""
##### \U0001f3af Do czego s\u0142u\u017cy ta zak\u0142adka?
Rekodowanie pozwala tworzy\u0107 **now\u0105 zmienn\u0105** na podstawie istniej\u0105cej,
przypisuj\u0105c nowe warto\u015bci do wybranych kod\u00f3w. Przyk\u0142ady:
- Grupowanie wieku: 18\u201324 \u2192 "M\u0142odzi", 25\u201344 \u2192 "\u015arednio", 45+ \u2192 "Starsi"
- Odwracanie skali: 1 \u2192 5, 2 \u2192 4, 3 \u2192 3, 4 \u2192 2, 5 \u2192 1
- Tworzenie flag: je\u015bli zmienna=1 to "Tak", inaczej "Nie"
- Ujednolicanie tekst\u00f3w: "tak" i "Tak" \u2192 "tak"

##### \U0001f527 Jak u\u017cywa\u0107 \u2014 krok po kroku
1. **Wybierz zmienn\u0105 \u017ar\u00f3d\u0142ow\u0105** \u2014 kolumn\u0119, na podstawie kt\u00f3rej tworzysz now\u0105.
2. **Podaj nazw\u0119 nowej zmiennej** \u2014 b\u0119dzie dodana do bazy jako nowa kolumna.
3. **Zdefiniuj regu\u0142y mapowania** \u2014 dla ka\u017cdej warto\u015bci (lub zakresu) podaj
   docelow\u0105 warto\u015b\u0107 nowej zmiennej. Mo\u017cesz te\u017c ustawi\u0107 warto\u015b\u0107 domy\u015bln\u0105
   (dla warto\u015bci, kt\u00f3re nie pasuj\u0105 do \u017cadnej regu\u0142y).
4. Kliknij **"\u25b6\ufe0f Zastosuj rekodowanie"** \u2014 nowa zmienna zostanie dodana do bazy.

##### \U0001f4a1 Wa\u017cne
- Nowa zmienna jest widoczna natychmiast we wszystkich modu\u0142ach analitycznych.
- Rekodowania s\u0105 zapisywane w projekcie JSON \u2014 po wczytaniu projektu zostan\u0105 ponownie zastosowane.
- Mo\u017cesz usuwa\u0107 istniej\u0105ce rekodowania z listy "Zapisane rekodowania".
- Zarwono zmienne numeryczne jak i tekstowe mog\u0105 by\u0107 rekodowane.
"""
            )

        col_r1, col_r2 = st.columns(2)
        with col_r1:
            # All visible columns -- numeric AND text
            all_recode_candidates = visible_columns
            src_var = st.selectbox(
                "Zmienna \u017ar\u00f3d\u0142owa:",
                all_recode_candidates,
                format_func=lambda x: get_var_display_name(x, var_labels),
                key="rec_src"
            )
            new_var_name = st.text_input("Nazwa nowej zmiennej:", value=f"{src_var}_r", key="rec_name")
            new_var_label = st.text_input("Etykieta nowej zmiennej:", key="rec_label")

            # Detect source type
            is_numeric_src = src_var in numeric_cols
            src_type_label = "numeryczna" if is_numeric_src else "tekstowa"
            st.caption(f"Typ zmiennej \u017ar\u00f3d\u0142owej: **{src_type_label}**")

            output_type = st.radio(
                "Typ warto\u015bci wyj\u015bciowych:",
                ["auto", "numeric", "text"],
                index=0,
                horizontal=True,
                key="rec_output_type",
                help="'auto' = automatyczne wykrycie (liczby je\u015bli mo\u017cliwe, tekst w p.p.)\n'numeric' = wymu\u015b liczby\n'text' = wymu\u015b tekst"
            )

        with col_r2:
            # Use df_raw for numeric, df (labeled) for categorical/text
            if is_numeric_src:
                unique_vals = sorted(df_raw[src_var].dropna().unique()) if src_var in df_raw.columns else []
            else:
                unique_vals = sorted(df[src_var].dropna().unique().astype(str)) if src_var in df.columns else []

            st.markdown(f"**Mapowanie warto\u015bci** (`{src_var}` -- {len(unique_vals)} unikalnych warto\u015bci):")

            mapping = {}
            if unique_vals:
                if len(unique_vals) > 50:
                    st.warning(f"Zmienna ma {len(unique_vals)} warto\u015bci -- wy\u015bwietlono pierwsze 50.")
                    unique_vals = unique_vals[:50]

                for val in unique_vals:
                    val_str = str(val)
                    # Show SPSS label in the field name if available (numeric vars)
                    if is_numeric_src:
                        try:
                            orig_label = meta_orig.variable_value_labels.get(src_var, {}).get(float(val_str), '')
                        except (ValueError, TypeError):
                            orig_label = ''
                    else:
                        orig_label = ''
                    display = f"{val_str}" + (f"  ({orig_label})" if orig_label else "")

                    # Always use text_input so both numeric and text targets are possible
                    new_val = st.text_input(
                        f"{display}  \u2192",
                        value=val_str,
                        key=f"rec_{src_var}_{val_str}"
                    )
                    mapping[val_str] = new_val
            else:
                st.info("Brak warto\u015bci do rekodowania (zmienna pusta lub brak danych).")

        if unique_vals and _tracked_button("\u2705 Utw\u00f3rz rekodowan\u0105 zmienn\u0105", "prep", "create_recoding", type="primary", use_container_width=True):
            if not new_var_name.strip():
                st.error("Podaj nazw\u0119 nowej zmiennej.")
            elif new_var_name.strip() in df_raw.columns:
                st.error(f"Zmienna `{new_var_name.strip()}` ju\u017c istnieje. Wybierz inn\u0105 nazw\u0119.")
            else:
                rec_entry = {
                    'source': src_var,
                    'new_name': new_var_name.strip(),
                    'label': new_var_label.strip() or f"Rekodowanie: {src_var}",
                    'mapping': mapping,
                    'output_type': output_type,
                }
                st.session_state.recodings.append(rec_entry)
                st.success(f"\u2705 Zmienna `{new_var_name.strip()}` zostanie dodana po od\u015bwie\u017ceniu.")
                st.rerun()

        if st.session_state.recodings:
            st.divider()
            st.markdown(f"**Zapisane rekodowania ({len(st.session_state.recodings)}):**")
            to_del = None
            for i, rec in enumerate(st.session_state.recodings):
                src_v   = rec['source']
                new_v   = rec['new_name']
                lbl_v   = rec.get('label', '')
                out_t   = rec.get('output_type', 'auto')
                mapping = rec.get('mapping', {})
                with st.expander(
                    f"`{new_v}` \u2190 `{src_v}` \u00b7 {lbl_v or '(brak etykiety)'}",
                    expanded=False
                ):
                    # Editable fields
                    e1, e2 = st.columns(2)
                    new_lbl_r = e1.text_input("Etykieta:", value=lbl_v, key=f"rec_lbl_{i}")
                    new_name_r = e2.text_input("Nazwa zmiennej:", value=new_v, key=f"rec_nm_{i}")

                    if mapping:
                        st.markdown("**Mapa kodowania:**")
                        for old_k, new_k in mapping.items():
                            mc1, mc2 = st.columns(2)
                            mc1.markdown(f"`{old_k}` \u2192")
                            mc2.text_input("", value=str(new_k),
                                key=f"rec_map_{i}_{old_k}",
                                label_visibility="collapsed")

                    bc1, bc2 = st.columns(2)
                    with bc1:
                        if st.button("\U0001f4be Zapisz zmiany", key=f"rec_save_{i}",
                                     use_container_width=True):
                            st.session_state.recodings[i]['label']    = new_lbl_r
                            st.session_state.recodings[i]['new_name'] = new_name_r
                            # Update mapping if editable
                            if mapping:
                                updated_map = {}
                                for old_k in mapping:
                                    raw_new = st.session_state.get(f"rec_map_{i}_{old_k}", str(mapping[old_k]))
                                    # Try to cast to original type
                                    try:
                                        raw_new = float(raw_new) if '.' in str(raw_new) else int(raw_new)
                                    except (ValueError, TypeError):
                                        pass
                                    updated_map[old_k] = raw_new
                                st.session_state.recodings[i]['mapping'] = updated_map
                            st.success("Zapisano.")
                            st.rerun()
                    with bc2:
                        if st.button("\U0001f5d1\ufe0f Usu\u0144 rekodowanie", key=f"rec_del_{i}",
                                     use_container_width=True):
                            to_del = i
            if to_del is not None:
                st.session_state.recodings.pop(to_del)
                st.rerun()

            if st.button("\U0001f5d1\ufe0f Usu\u0144 wszystkie rekodowania", type="secondary",
                         use_container_width=True, key="rec_del_all"):
                st.session_state.recodings = []
                st.rerun()
    # -- ETYKIETY ZMIENNYCH I WARTOSCI ----------------
    with tab_labels:
        st.markdown("#### Etykiety zmiennych i warto\u015bci")


        with st.expander("Instrukcja \u2014 jak nadawa\u0107 etykiety zmiennym i warto\u015bciom", expanded=False):
            st.markdown("""
##### \U0001f3af Do czego s\u0142u\u017cy ta zak\u0142adka?
Etykiety zast\u0119puj\u0105 techniczne nazwy kolumn i kody liczbowe czytelnymi opisami.
Na przyk\u0142ad kolumna `p1` mo\u017ce mie\u0107 etykiet\u0119 *"P1. Jak ocenia Pan/i..."*, a kod `1`
etykiet\u0119 *"Zdecydowanie tak"*. Etykiety wy\u015bwietlane s\u0105 we wszystkich tabelach i wykresach.

##### \U0001f527 Jak u\u017cywa\u0107 \u2014 krok po kroku

**Zak\u0142adka "\u270f\ufe0f Etykieta zmiennej":**
1. Wybierz zmienn\u0105 z listy.
2. Wpisz now\u0105 etykiet\u0119 (np. pe\u0142n\u0105 tre\u015b\u0107 pytania ankietowego).
3. Kliknij **Zapisz etykiet\u0119**.

**Zak\u0142adka "\U0001f3f7\ufe0f Etykiety warto\u015bci (kody)":**
1. Wybierz zmienn\u0105 do edycji.
2. Dla ka\u017cdej warto\u015bci liczbowej (np. 1, 2, 3) wpisz opis s\u0142owny.
3. Kliknij **Zapisz etykiety warto\u015bci**.
4. Opcjonalnie: **Wyczy\u015b\u0107 etykiety** \u2014 usuwa wszystkie etykiety warto\u015bci dla tej zmiennej.

**Zak\u0142adka "\U0001f4cb Wszystkie moje zmiany":**
- Podgl\u0105d wszystkich wprowadzonych zmian \u2014 etykiet zmiennych i warto\u015bci.
- Mo\u017cesz tu usun\u0105\u0107 wybran\u0105 etykiet\u0119 lub wyczy\u015bci\u0107 wszystkie.

##### \U0001f4a1 Wa\u017cne
- Pliki SPSS (.sav) zazwyczaj zawieraj\u0105 gotowe etykiety \u2014 s\u0105 automatycznie wczytywane.
- Zmiany etykiet **nie modyfikuj\u0105** oryginalnych danych \u2014 dzia\u0142aj\u0105 tylko w tej sesji.
- Etykiety s\u0105 zapisywane w pliku projektu JSON.
"""
            )

        lab_sub_var, lab_sub_val, lab_sub_all = st.tabs([
            "Etykieta zmiennej",
            "Etykiety warto\u015bci (kody)",
            "Wszystkie moje zmiany"
        ])

        # \u2500\u2500 Sub-tab 1: Variable label \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
        with lab_sub_var:
            col_ev1, col_ev2 = st.columns(2)
            with col_ev1:
                edit_var = st.selectbox("Zmienna:", df_raw.columns,
                    format_func=lambda x: f"[{x}] {var_labels.get(x, x)}",
                    key="lab_edit_var")
                new_lbl = st.text_input("Nowa etykieta:",
                    value=var_labels.get(edit_var, ''), key="lab_new_lbl")
                bcol1, bcol2 = st.columns(2)
                with bcol1:
                    if st.button("\U0001f4be Zapisz", key="lab_save_var_lbl",
                                 use_container_width=True, type="primary"):
                        st.session_state.custom_var_labels[edit_var] = new_lbl
                        var_labels[edit_var] = new_lbl
                        st.success(f"Zapisano etykiet\u0119 `{edit_var}`.")
                        st.rerun()
                with bcol2:
                    if edit_var in st.session_state.custom_var_labels:
                        if st.button("\U0001f5d1\ufe0f Usu\u0144", key="lab_del_var_lbl",
                                     use_container_width=True):
                            st.session_state.custom_var_labels.pop(edit_var, None)
                            st.rerun()
            with col_ev2:
                st.markdown("**Etykiety warto\u015bci tej zmiennej:**")
                spss_vvl_preview = meta_orig.variable_value_labels.get(edit_var, {})
                cust_vl_preview  = st.session_state.custom_val_labels.get(edit_var, {})
                if spss_vvl_preview or cust_vl_preview:
                    preview_rows = []
                    for k, v in sorted(spss_vvl_preview.items()):
                        preview_rows.append({'Kod': k, 'Etykieta \u017ar\u00f3d\u0142owa': v,
                            'Niestandardowa': cust_vl_preview.get(str(k), '--')})
                    for rk, cv in cust_vl_preview.items():
                        if not any(str(r['Kod']) == rk for r in preview_rows):
                            preview_rows.append({'Kod': rk, 'Etykieta \u017ar\u00f3d\u0142owa': '--',
                                'Niestandardowa': cv})
                    st.dataframe(pd.DataFrame(preview_rows), use_container_width=True, hide_index=True)
                else:
                    uniq = sorted(df_raw[edit_var].dropna().unique())
                    st.info(f"Brak etykiet warto\u015bci. Unikalne warto\u015bci: "
                            f"{', '.join(str(v) for v in uniq[:20])}"
                            + (" ..." if len(uniq) > 20 else ""))

        # \u2500\u2500 Sub-tab 2: Value labels editor \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
        with lab_sub_val:
            st.caption("Zmie\u0144 wy\u015bwietlane etykiety warto\u015bci dla wybranej zmiennej.")
            col_vv1, col_vv2 = st.columns([1, 1])
            with col_vv1:
                val_edit_var = st.selectbox("Zmienna:", df_raw.columns,
                    format_func=lambda x: f"[{x}] {var_labels.get(x, x)}",
                    key="lab_val_edit_var")
                spss_codes = meta_orig.variable_value_labels.get(val_edit_var, {})
                raw_vals   = sorted(df_raw[val_edit_var].dropna().unique())
                st.caption(f"{len(raw_vals)} unikalnych warto\u015bci \u00b7 {len(spss_codes)} etykiet \u017ar\u00f3d\u0142owych")
                display_items = {}
                for v in raw_vals:
                    raw_str = str(v)
                    spss_lbl = spss_codes.get(v, spss_codes.get(raw_str, ''))
                    display_items[raw_str] = spss_lbl if spss_lbl else raw_str
                if len(display_items) > 50:
                    st.warning(f"Wy\u015bwietlono pierwsze 50 z {len(display_items)} warto\u015bci.")
                    display_items = dict(list(display_items.items())[:50])
            with col_vv2:
                st.markdown("**Nowe etykiety wy\u015bwietlania:**")
                existing_custom = st.session_state.custom_val_labels.get(val_edit_var, {})
                new_val_map = {}
                for raw_str, spss_lbl in display_items.items():
                    current = existing_custom.get(raw_str, spss_lbl)
                    hint = f" [{spss_lbl}]" if spss_lbl and spss_lbl != raw_str else ""
                    new_label = st.text_input(f"Kod {raw_str}{hint}:", value=current,
                        key=f"lab_vl_{val_edit_var}_{raw_str}")
                    new_val_map[raw_str] = new_label
                col_bsave, col_bclear = st.columns(2)
                with col_bsave:
                    if st.button("\U0001f4be Zapisz etykiety warto\u015bci", key="lab_save_val_lbls",
                                 use_container_width=True, type="primary"):
                        filtered = {k: v for k, v in new_val_map.items()
                                    if v.strip() and v.strip() != k}
                        st.session_state.custom_val_labels[val_edit_var] = filtered
                        st.session_state.user_cleared_val_labels.discard(val_edit_var)
                        st.success(f"Zapisano {len(filtered)} etykiet dla `{val_edit_var}`.")
                        st.rerun()
                with col_bclear:
                    if existing_custom and st.button("\U0001f5d1\ufe0f Usu\u0144",
                                                     key="lab_clear_val_lbls",
                                                     use_container_width=True):
                        st.session_state.custom_val_labels.pop(val_edit_var, None)
                        st.session_state.user_cleared_val_labels.add(val_edit_var)
                        st.rerun()

            # \u2500\u2500 Saved labels for current variable \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
            all_saved = st.session_state.custom_val_labels
            if all_saved:
                st.divider()
                st.markdown("**\U0001f4cb Zapisane etykiety warto\u015bci (wszystkie zmienne):**")
                for sv_col, sv_map in list(all_saved.items()):
                    with st.expander(
                        f"`{sv_col}` \u2014 {var_labels.get(sv_col, sv_col)} ({len(sv_map)} etykiet)",
                        expanded=(sv_col == val_edit_var)
                    ):
                        sv_edited = {}
                        for code, lbl in sorted(sv_map.items(), key=lambda x: x[0]):
                            ec1, ec2, ec3 = st.columns([1, 4, 1])
                            ec1.markdown(f"**{code}**")
                            new_v = ec2.text_input("", value=str(lbl),
                                key=f"sv_val_{sv_col}_{code}",
                                label_visibility="collapsed")
                            sv_edited[code] = new_v
                            with ec3:
                                if st.button("\U0001f5d1\ufe0f", key=f"sv_del_row_{sv_col}_{code}"):
                                    upd = {k: v for k, v in sv_map.items() if str(k) != str(code)}
                                    if upd:
                                        st.session_state.custom_val_labels[sv_col] = upd
                                    else:
                                        st.session_state.custom_val_labels.pop(sv_col, None)
                                        st.session_state.user_cleared_val_labels.add(sv_col)
                                    st.rerun()
                        sc1, sc2 = st.columns(2)
                        with sc1:
                            if st.button("\U0001f4be Zapisz", key=f"sv_save_{sv_col}",
                                         use_container_width=True):
                                st.session_state.custom_val_labels[sv_col] = sv_edited
                                st.success("Zapisano.")
                                st.rerun()
                        with sc2:
                            if st.button("\U0001f5d1\ufe0f Usu\u0144 ca\u0142\u0105 zmienn\u0105",
                                         key=f"sv_del_{sv_col}", use_container_width=True):
                                st.session_state.custom_val_labels.pop(sv_col, None)
                                st.session_state.user_cleared_val_labels.add(sv_col)
                                st.rerun()

        # \u2500\u2500 Sub-tab 3: All custom labels overview \u2500\u2500\u2500\u2500\u2500
        with lab_sub_all:
            has_any = st.session_state.custom_var_labels or st.session_state.custom_val_labels
            if not has_any:
                st.info("Nie wprowadzono \u017cadnych niestandardowych etykiet.")
            else:
                if st.session_state.custom_var_labels:
                    st.markdown("**\u270f\ufe0f Zmienione etykiety zmiennych:**")
                    changed_var = False
                    for col_cv, lbl_cv in list(st.session_state.custom_var_labels.items()):
                        rc1, rc2, rc3 = st.columns([3, 3, 1])
                        rc1.markdown(f"`{col_cv}`")
                        new_cv = rc2.text_input("", value=lbl_cv,
                            key=f"all_var_lbl_{col_cv}", label_visibility="collapsed")
                        with rc3:
                            if st.button("\U0001f5d1\ufe0f", key=f"all_del_var_{col_cv}"):
                                st.session_state.custom_var_labels.pop(col_cv, None)
                                st.rerun()
                        if new_cv != lbl_cv:
                            st.session_state.custom_var_labels[col_cv] = new_cv
                            changed_var = True
                    st.divider()

                if st.session_state.custom_val_labels:
                    st.markdown("**\U0001f3f7\ufe0f Zmienione etykiety warto\u015bci:**")
                    for col_cv, val_map in list(st.session_state.custom_val_labels.items()):
                        with st.expander(
                            f"`{col_cv}` \u2014 {var_labels.get(col_cv, col_cv)}"
                            f" ({len(val_map)} etykiet)", expanded=False
                        ):
                            edited_vmap = {}
                            for code, lbl in sorted(val_map.items(), key=lambda x: x[0]):
                                ec1, ec2 = st.columns([1, 3])
                                ec1.markdown(f"Kod **{code}**")
                                new_v = ec2.text_input("", value=str(lbl),
                                    key=f"all_val_{col_cv}_{code}",
                                    label_visibility="collapsed")
                                edited_vmap[code] = new_v
                            sc1, sc2 = st.columns(2)
                            with sc1:
                                if st.button("\U0001f4be Zapisz", key=f"all_save_val_{col_cv}",
                                             use_container_width=True):
                                    st.session_state.custom_val_labels[col_cv] = edited_vmap
                                    st.session_state.user_cleared_val_labels.discard(col_cv)
                                    st.success("Zapisano.")
                                    st.rerun()
                            with sc2:
                                if st.button("\U0001f5d1\ufe0f Usu\u0144",
                                             key=f"all_del_val_{col_cv}",
                                             use_container_width=True):
                                    st.session_state.custom_val_labels.pop(col_cv, None)
                                    st.session_state.user_cleared_val_labels.add(col_cv)
                                    st.rerun()

            st.divider()
            if st.button("\U0001f5d1\ufe0f Usu\u0144 WSZYSTKIE moje etykiety",
                         type="secondary", use_container_width=True, key="all_del_all_labels"):
                _all_val_cols = set(st.session_state.custom_val_labels.keys())
                st.session_state.custom_var_labels = {}
                st.session_state.custom_val_labels = {}
                st.session_state.user_cleared_val_labels = _all_val_cols
                st.rerun()

    # -- CZYSZCZENIE DANYCH -------------------------
    with tab_clean:
        st.markdown("#### Czyszczenie danych tekstowych")

        with st.expander("Instrukcja \u2014 jak czy\u015bci\u0107 dane tekstowe", expanded=False):
            st.markdown("""
##### \U0001f3af Do czego s\u0142u\u017cy ta zak\u0142adka?
Dane tekstowe z ankiet cz\u0119sto zawieraj\u0105 niepo\u017c\u0105dane bia\u0142e znaki, r\u00f3\u017cne wielko\u015bci liter
lub b\u0142\u0105d formatowania. Ta zak\u0142adka pozwala je automatycznie poprawi\u0107.

##### \U0001f527 Dost\u0119pne operacje
- **Trimming** \u2014 usuwanie spacji z pocz\u0105tku i ko\u0144ca warto\u015bci (np. `" tak "` \u2192 `"tak"`).
- **Lowercase / Uppercase / Title Case** \u2014 ujednolicenie wielko\u015bci liter.
- **Usuwanie wielokrotnych spacji** \u2014 zast\u0105pienie podw\u00f3jnych spacji pojedynczymi.
- **Zast\u0119powanie warto\u015bci** \u2014 podmie\u0144 konkretn\u0105 warto\u015b\u0107 inn\u0105 (np. `"tak"` \u2192 `"Tak"`).

##### \U0001f527 Jak u\u017cywa\u0107 \u2014 krok po kroku
1. Wybierz zakres: **globalnie** (wszystkie zmienne tekstowe) lub **wybrana zmienna**.
2. Zaznacz operacje do wykonania.
3. Kliknij **Zastosuj**.
4. Podgl\u0105d zmian pojawi si\u0119 poni\u017cej \u2014 zmiany s\u0105 zapisywane jako operacje czyszczenia.

##### \U0001f4a1 Wa\u017cne
- Operacje czyszczenia s\u0105 tymczasowe \u2014 dzia\u0142aj\u0105 w sesji i zapisywane w projekcie JSON.
- Mo\u017cna je cofn\u0105\u0107 usuwaj\u0105c z listy zapisanych operacji.
- Czyszczenie wp\u0142ywa na analizy tekstowe, chmury s\u0142\u00f3w i cz\u0119sto\u015bci tekst\u00f3w otwartych.
"""
            )

        # \u2500\u2500 Helper: text columns \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
        text_cols = [c for c in visible_columns if df_raw[c].dtype == object]

        if not text_cols:
            st.warning("Brak zmiennych tekstowych w bazie danych.")
        else:
            # \u2500\u2500 SECTION A: Globally on all text columns \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
            st.markdown("### Globalnie \u2014 wszystkie zmienne tekstowe")
            st.caption(f"Dzia\u0142a na {len(text_cols)} zmiennych tekstowych jednocze\u015bnie.")

            clean_col1, clean_col2 = st.columns(2)
            with clean_col1:
                st.markdown("**Czyszczenie spacji i bia\u0142ych znak\u00f3w:**")
                do_strip    = st.checkbox("Usu\u0144 spacje na pocz\u0105tku i ko\u0144cu", key="clean_g_strip", value=True)
                do_dbl_sp   = st.checkbox("Usu\u0144 podw\u00f3jne spacje",               key="clean_g_dbl")
                do_tabs     = st.checkbox("Usu\u0144 tabulatory",                   key="clean_g_tabs")
                do_newlines = st.checkbox("Usu\u0144 znaki nowej linii",       key="clean_g_nl")

                st.markdown("**Standaryzacja cudzys\u0142ow\u00f3w:**")
                do_quotes   = st.checkbox("Zamie\u0144 cudzys\u0142owy na standardowe ASCII (\u201c\u201d\u2018\u2019 \u2192 \"')", key="clean_g_quotes")

            with clean_col2:
                st.markdown("**Wielko\u015b\u0107 liter:**")
                case_mode = st.radio(
                    "Zastosuj:",
                    ["Bez zmian", "WIELKIE LITERY", "ma\u0142e litery", "Pierwsza Wielka (Title Case)"],
                    key="clean_g_case", horizontal=False
                )

                st.markdown("**Znaki specjalne:**")
                do_special = st.checkbox(
                    "Usu\u0144 znaki specjalne (pozostaw litery, cyfry i spacje)",
                    key="clean_g_special",
                    help="Usuwa wszystko poza literami (w tym polskimi), cyframi i spacjami."
                )

            if st.button("\u25b6\ufe0f Zastosuj globalnie do wszystkich zmiennych tekstowych",
                         type="primary", key="clean_global_apply", use_container_width=True):
                ops = {
                    'strip':    do_strip,
                    'dbl_sp':   do_dbl_sp,
                    'tabs':     do_tabs,
                    'newlines': do_newlines,
                    'quotes':   do_quotes,
                    'case':     {'Bez zmian': 'none', 'WIELKIE LITERY': 'upper',
                                 'ma\u0142e litery': 'lower', 'Pierwsza Wielka (Title Case)': 'title'}.get(case_mode, 'none'),
                    'special':  do_special,
                }
                if any(v for k, v in ops.items() if k != 'case') or ops['case'] != 'none':
                    st.session_state.cleaning_ops.append({'cols': list(text_cols), 'ops': ops})
                    st.success(f"\u2705 Czyszczenie zapisane dla {len(text_cols)} zmiennych i b\u0119dzie stosowane przy ka\u017cdym wczytaniu danych.")
                    st.rerun()
                else:
                    st.warning("Nie wybrano \u017cadnej operacji czyszczenia.")

            st.divider()

            # \u2500\u2500 SECTION B: Per-column \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
            st.markdown("### Dla wybranych zmiennych")

            sel_cols = st.multiselect(
                "Wybierz zmienne do czyszczenia:",
                text_cols,
                format_func=lambda x: get_var_display_name(x, var_labels),
                key="clean_sel_cols"
            )

            if sel_cols:
                c1, c2 = st.columns(2)
                with c1:
                    st.markdown("**Spacje i bia\u0142e znaki:**")
                    p_strip    = st.checkbox("Usu\u0144 spacje na pocz\u0105tku i ko\u0144cu", key="clean_p_strip",  value=True)
                    p_dbl      = st.checkbox("Usu\u0144 podw\u00f3jne spacje",               key="clean_p_dbl")
                    p_tabs     = st.checkbox("Usu\u0144 tabulatory",                         key="clean_p_tabs")
                    p_nl       = st.checkbox("Usu\u0144 znaki nowej linii",                  key="clean_p_nl")
                    p_quotes   = st.checkbox("Zamie\u0144 cudzys\u0142owy na ASCII",         key="clean_p_quotes")
                with c2:
                    st.markdown("**Wielko\u015b\u0107 liter:**")
                    p_case = st.radio(
                        "Zastosuj:",
                        ["Bez zmian", "WIELKIE LITERY", "ma\u0142e litery", "Pierwsza Wielka"],
                        key="clean_p_case", horizontal=False
                    )
                    st.markdown("**Znaki specjalne:**")
                    p_special = st.checkbox(
                        "Usu\u0144 znaki specjalne",
                        key="clean_p_special"
                    )

                # Preview
                if sel_cols:
                    st.markdown("**Podgl\u0105d pierwszych 5 wierszy (przed / po):**")
                    preview_col = sel_cols[0]
                    orig = df_raw[preview_col].dropna().astype(str).head(5)
                    prev = orig.copy()
                    if p_strip:   prev = prev.str.strip()
                    if p_dbl:     prev = prev.str.replace(r' {2,}', ' ', regex=True)
                    if p_tabs:    prev = prev.str.replace('\t', ' ', regex=False)
                    if p_nl:      prev = prev.str.replace(r'[\n\r]', ' ', regex=True)
                    if p_quotes:
                        for old_q, new_q in [('\u201c','"'), ('\u201d','"'), ('\u201e','"'),
                                              ('\u2018',"'"), ('\u2019',"'"), ('\u201a',"'")]:
                            prev = prev.str.replace(old_q, new_q, regex=False)
                    if p_case == "WIELKIE LITERY":        prev = prev.str.upper()
                    elif p_case == "ma\u0142e litery":    prev = prev.str.lower()
                    elif p_case == "Pierwsza Wielka":     prev = prev.str.title()
                    if p_special:
                        prev = prev.str.replace(r'[^\w\s]', '', regex=True).str.replace('_','', regex=False)

                    preview_df = pd.DataFrame({'Przed': orig.values, 'Po': prev.values})
                    st.dataframe(preview_df, use_container_width=True, hide_index=True)

                if st.button("\u25b6\ufe0f Zastosuj dla wybranych zmiennych",
                             type="primary", key="clean_per_apply", use_container_width=True):
                    ops = {
                        'strip':    p_strip,
                        'dbl_sp':   p_dbl,
                        'tabs':     p_tabs,
                        'newlines': p_nl,
                        'quotes':   p_quotes,
                        'case':     {'Bez zmian': 'none', 'WIELKIE LITERY': 'upper',
                                     'ma\u0142e litery': 'lower', 'Pierwsza Wielka': 'title'}.get(p_case, 'none'),
                        'special':  p_special,
                    }
                    if any(v for k, v in ops.items() if k != 'case') or ops['case'] != 'none':
                        st.session_state.cleaning_ops.append({'cols': list(sel_cols), 'ops': ops})
                        st.success(f"\u2705 Czyszczenie zapisane dla {len(sel_cols)} zmiennych: {', '.join(sel_cols[:3])}{'...' if len(sel_cols) > 3 else ''}.")
                        st.rerun()
                    else:
                        st.warning("Nie wybrano \u017cadnej operacji czyszczenia.")

            # \u2500\u2500 Active cleaning rules panel \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
            if st.session_state.cleaning_ops:
                st.divider()
                st.markdown("### Aktywne regu\u0142y czyszczenia")
                st.caption("Regu\u0142y s\u0105 stosowane automatycznie przy ka\u017cdym wczytaniu danych. Kolejno\u015b\u0107 ma znaczenie.")
                to_remove = None
                for i, entry in enumerate(st.session_state.cleaning_ops):
                    ops  = entry['ops']
                    cols = entry['cols']
                    # Build human-readable summary
                    op_labels = []
                    if ops.get('strip'):    op_labels.append("trim spacji")
                    if ops.get('dbl_sp'):   op_labels.append("podw\u00f3jne spacje")
                    if ops.get('tabs'):     op_labels.append("tabulatory")
                    if ops.get('newlines'): op_labels.append("nowe linie")
                    if ops.get('quotes'):   op_labels.append("cudzys\u0142owy")
                    case = ops.get('case', 'none')
                    if case != 'none': op_labels.append({'upper': 'WIELKIE', 'lower': 'ma\u0142e', 'title': 'Pierwsze Wielkie'}.get(case, case))
                    if ops.get('special'):  op_labels.append("znaki specjalne")
                    col_desc = f"{len(cols)} zmiennych" if len(cols) > 3 else ", ".join(cols)
                    row_c1, row_c2 = st.columns([5, 1])
                    row_c1.markdown(f"**{i+1}.** `{col_desc}` \u2192 {', '.join(op_labels)}")
                    if row_c2.button("\U0001f5d1\ufe0f", key=f"del_clean_{i}", help="Usu\u0144 t\u0119 regu\u0142\u0119"):
                        to_remove = i
                if to_remove is not None:
                    st.session_state.cleaning_ops.pop(to_remove)
                    st.rerun()
                if st.button("\U0001f5d1\ufe0f Usu\u0144 wszystkie regu\u0142y czyszczenia", type="secondary", key="clean_clear_all"):
                    st.session_state.cleaning_ops = []
                    st.rerun()

    # -- TYPY ZMIENNYCH (tylko Excel) ------------------
    if is_tabular and tab_types is not None:
        with tab_types:
            st.markdown("#### Typy zmiennych (auto-detekcja + korekta)")

            with st.expander("Instrukcja \u2014 jak korygowa\u0107 typy zmiennych", expanded=False):
                st.markdown("""
##### \U0001f3af Do czego s\u0142u\u017cy ta zak\u0142adka?
Narz\u0119dzie automatycznie przypisuje typ **numeryczny** lub **kategoryczny** ka\u017cdej kolumnie
na podstawie heurystyki: je\u015bli kolumna zawiera \u226510 unikalnych warto\u015bci
i same liczby \u2192 numeryczna; w przeciwnym razie \u2192 kategoryczna. Ta zak\u0142adka
pozwala r\u0119cznie poprawi\u0107 b\u0142\u0119dnie wykryte typy.

##### \U0001f527 Jak u\u017cywa\u0107
1. W tabeli widocznej poni\u017cej znajd\u017a kolumn\u0119, kt\u00f3r\u0105 chcesz zmieni\u0107.
2. W kolumnie **"Typ (r\u0119czny)"** zmie\u0144 warto\u015b\u0107 z `auto` na `numeryczna` lub `kategoryczna`.
3. Kliknij **Zapisz zmiany typ\u00f3w**.

##### \U0001f4a1 Przyk\u0142ady kiedy warto zmieni\u0107 typ
- Kolumna zawiera kody liczbowe (np. 1, 2, 3 dla odpowiedzi) \u2014 auto-detekcja
  mo\u017ce j\u0105 oznaczy\u0107 jako **numeryczn\u0105**, ale powinna by\u0107 **kategoryczna**.
- Kolumna zawiera rok (np. 2020, 2021) \u2014 mo\u017ce by\u0107 sensowna jako numeryczna
  (do \u015brednich/regresji) lub kategoryczna (do tabel cz\u0119sto\u015bci).
- Skale Likerta (1\u20135, 1\u20137) \u2014 zale\u017cnie od analizy warto mie\u0107 obie opcje.

##### \U0001f4a1 Wa\u017cne
- Ta zak\u0142adka jest dost\u0119pna **tylko dla plik\u00f3w Excel** (pliki SPSS maj\u0105 typy w metadata).
- Zmiana typu nie modyfikuje danych \u2014 tylko instruuje narz\u0119dzie jak je traktowa\u0107.
"""
                )

            type_rows = []
            _type_label = {"numeric": "numeryczna", "categorical": "kategoryczna", "auto": "auto"}
            for col in df_orig_raw.columns:
                auto_type = "numeryczna" if col in numeric_cols_raw else "kategoryczna"
                override  = st.session_state.excel_col_types.get(col, "auto")
                effective = auto_type if override == "auto" else _type_label.get(override, override)
                type_rows.append({
                    "Zmienna": col,
                    "Etykieta": var_labels.get(col, col),
                    "Auto-detekcja": auto_type,
                    "Korekta u\u017cytkownika": _type_label.get(override, override),
                    "Efektywny typ": effective,
                    "Unikalnych warto\u015bci": df_orig_raw[col].nunique(),
                    "Brak\u00f3w [N]": df_orig_raw[col].isna().sum(),
                })
            st.dataframe(pd.DataFrame(type_rows), use_container_width=True, height=280)

            st.divider()
            st.markdown("**Zmie\u0144 typ wybranej zmiennej:**")
            col_t1, col_t2, col_t3 = st.columns([3, 2, 2])
            with col_t1:
                type_edit_var = st.selectbox(
                    "Zmienna:", df_orig_raw.columns,
                    format_func=lambda x: f"[{x}] {var_labels.get(x, x)}",
                    key="type_edit_var"
                )
            with col_t2:
                cur_override = st.session_state.excel_col_types.get(type_edit_var, "auto")
                new_type = st.selectbox(
                    "Typ:", ["auto", "numeric", "categorical"],
                    format_func=lambda x: {"auto": "auto", "numeric": "numeryczna", "categorical": "kategoryczna"}[x],
                    index=["auto", "numeric", "categorical"].index(cur_override) if cur_override in ["auto", "numeric", "categorical"] else 0,
                    key="type_edit_val"
                )
            with col_t3:
                st.write("")
                st.write("")
                if st.button("\U0001f4be Zastosuj", key="type_apply", use_container_width=True):
                    if new_type == "auto":
                        st.session_state.excel_col_types.pop(type_edit_var, None)
                    else:
                        st.session_state.excel_col_types[type_edit_var] = new_type
                    load_excel_data.clear()
                    st.success(f"Typ `{type_edit_var}` zmieniony na **{new_type}**. Strona zostanie od\u015bwie\u017cona.")
                    st.rerun()

            if st.session_state.excel_col_types:
                if st.button("\U0001f5d1\ufe0f Zresetuj wszystkie korekty typ\u00f3w", type="secondary", key="type_reset"):
                    st.session_state.excel_col_types = {}
                    load_excel_data.clear()
                    st.rerun()

            # \u2500\u2500 Summary of converted variables \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
            numeric_overrides = {c: v for c, v in st.session_state.excel_col_types.items()
                                 if v == "numeric"}
            if numeric_overrides:
                st.divider()
                st.markdown("**\u2705 Zmienne zmienione na numeryczne:**")
                for col_n in numeric_overrides:
                    val_lbl = st.session_state.custom_val_labels.get(col_n, {})
                    tag = " \u2014 zakodowana tekstowo" if val_lbl else " \u2014 by\u0142a ju\u017c liczbowa"
                    with st.expander(f"**{var_labels.get(col_n, col_n)}** (`{col_n}`){tag}", expanded=False):
                        if val_lbl:
                            st.markdown("**Etykiety warto\u015bci** (mo\u017cesz edytowa\u0107):")
                            edited_labels = {}
                            for code, lbl in sorted(val_lbl.items(), key=lambda x: x[0]):
                                new_lbl = st.text_input(
                                    f"Kod {code}:",
                                    value=str(lbl),
                                    key=f"lbl_edit_{col_n}_{code}"
                                )
                                edited_labels[code] = new_lbl
                            if st.button("\U0001f4be Zapisz etykiety", key=f"lbl_save_{col_n}",
                                         use_container_width=True):
                                st.session_state.custom_val_labels[col_n] = edited_labels
                                st.success("Etykiety zapisane.")
                                st.rerun()
                        else:
                            st.info("Warto\u015bci by\u0142y ju\u017c liczbowe \u2014 brak etykiet do edycji.")
                        if st.button("\u21a9\ufe0f Cofnij zmian\u0119 typu", key=f"type_revert_{col_n}",
                                     use_container_width=True):
                            st.session_state.excel_col_types.pop(col_n, None)
                            st.session_state.custom_val_labels.pop(col_n, None)
                            load_excel_data.clear()
                            st.rerun()

    # -- MRS ------------------------------------------
    with tab_mrs:
        with st.expander("Instrukcja \u2014 jak definiowa\u0107 pytania wielokrotnego wyboru (MRS)", expanded=False):
            st.markdown("""
##### \U0001f3af Do czego s\u0142u\u017cy ta zak\u0142adka?
Pytania wielokrotnego wyboru (MRS \u2014 Multiple Response Sets) to pytania, w kt\u00f3rych
respondent m\u00f3g\u0142 zaznaczy\u0107 kilka odpowiedzi jednocze\u015bnie. W SPSS s\u0105 zakodowane
jako kilka kolumn binarnych (0/1), np. `q3_1`, `q3_2`, `q3_3`. Zestaw MRS \u0142\u0105czy je
w jedno pytanie i umo\u017cliwia prawid\u0142owe obliczanie procent\u00f3w (od podstawy respondent\u00f3w,
nie od liczby wskaza\u0144).

##### \U0001f527 Jak u\u017cywa\u0107 \u2014 krok po kroku
1. **Autowykrywanie** \u2014 kliknij "\U0001f50d Autowykrywanie (zmienne 0/1)" \u2014 narz\u0119dzie
   automatycznie znajdzie grupy kolumn, kt\u00f3re wygl\u0105daj\u0105 jak zestaw MRS (maj\u0105 wsp\u00f3lny
   prefiks i zawieraj\u0105 tylko 0 i 1).
2. **R\u0119czne dodanie** \u2014 wpisz nazw\u0119 zestawu, wybierz kolumny z listy (min. 2),
   kliknij "Dodaj zestaw".
3. Zestawy MRS pojawiaj\u0105 si\u0105 w prawej kolumnie jako lista. Mo\u017cesz je usuwa\u0107.

##### \U0001f4ca Gdzie s\u0105 widoczne zestawy MRS?
- W module **Analizy \u2192 Matrycowe** \u2014 zestawy MRS tworz\u0105 tabele wielokrotnych odpowiedzi
  z procentami od bazy respondent\u00f3w.
- W module **Eksport do Excela** \u2014 zestawy MRS mog\u0105 by\u0107 eksportowane jako bloki tabelaryczne.

##### \U0001f4a1 Wa\u017cne
- Kolumny w zestawie MRS **musz\u0105 by\u0107 numeryczne** (0/1 lub 0/wi\u0119cej ni\u017c 0).
- Zestawy s\u0105 zapisywane w pliku projektu JSON.
"""
            )

        col1, col2 = st.columns([1, 1])
        with col1:
            st.markdown("#### Dodaj zestaw wielokrotnego wyboru")

            # Numeric columns only
            numeric_mrs_cols = [c for c in df_raw.columns
                                if pd.api.types.is_numeric_dtype(df_raw[c])]

            if st.button("\U0001f50d Autowykrywanie (zmienne 0/1)", use_container_width=True):
                detected = auto_detect_mrs(df_raw)
                added = 0
                skipped = []
                for name, cols in detected.items():
                    key = f"Auto_{name}"
                    if key not in st.session_state.mrs_sets:
                        st.session_state.mrs_sets[key] = {'cols': cols, 'count_val': 1}
                        added += 1
                    else:
                        skipped.append(key)
                if added:
                    st.success(f"Wykryto i dodano {added} zestaw\u00f3w wielokrotnych odpowiedzi.")
                elif not detected:
                    st.info("Nie znaleziono zmiennych binarnych (0/1) nadaj\u0105cych si\u0119 do zestawu MRS.")
                else:
                    st.info("Nie dodano nowych zestaw\u00f3w \u2014 wszystkie wykryte ju\u017c istniej\u0105.")
                if skipped:
                    _skip_list = ", ".join(f"`{s}`" for s in skipped[:5])
                    if len(skipped) > 5:
                        _skip_list += f" i {len(skipped) - 5} wi\u0119cej"
                    st.warning(
                        f"Pomini\u0119to {len(skipped)} zestaw\u00f3w \u2014 nazwy ju\u017c zaj\u0119te: {_skip_list}. "
                        "Edytuj lub usu\u0144 istniej\u0105ce zestawy, a nast\u0119pnie uruchom autowykrywanie ponownie.")
                st.rerun()

            mrs_name = st.text_input("Nazwa nowego zestawu:", key="mrs_new_name")
            mrs_cols = st.multiselect(
                "Zmienne (numeryczne):",
                options=numeric_mrs_cols,
                format_func=lambda x: get_var_display_name(x, var_labels),
                key="mrs_new_cols"
            )

            # Show possible values from selected columns
            if mrs_cols:
                all_vals = sorted(set(
                    v for c in mrs_cols
                    for v in df_raw[c].dropna().unique()
                ))
                val_options = [int(v) if float(v) == int(v) else float(v) for v in all_vals]
            else:
                val_options = [1]

            mrs_count_val = st.selectbox(
                "Warto\u015b\u0107 do zliczania:",
                options=val_options if val_options else [1],
                index=val_options.index(1) if 1 in val_options else 0,
                key="mrs_new_val",
                help="Kt\u00f3ra warto\u015b\u0107 oznacza zaznaczenie odpowiedzi? Najcz\u0119\u015bciej: 1."
            )

            if st.button("\u2795 Dodaj zestaw", use_container_width=True):
                if mrs_name and mrs_cols:
                    if mrs_name in st.session_state.mrs_sets:
                        st.error(
                            f"Zestaw o nazwie `{mrs_name}` ju\u017c istnieje. "
                            "Zmie\u0144 nazw\u0119 lub edytuj istniej\u0105cy zestaw w panelu poni\u017cej.")
                    else:
                        st.session_state.mrs_sets[mrs_name] = {
                            'cols': mrs_cols,
                            'count_val': mrs_count_val
                        }
                        st.rerun()
                else:
                    st.warning("Podaj nazw\u0119 i wybierz co najmniej jedn\u0105 zmienn\u0105.")

        with col2:
            st.markdown("#### Zdefiniowane zestawy -- kliknij aby edytowa\u0107")
            if not st.session_state.mrs_sets:
                st.info("Brak zdefiniowanych zestaw\u00f3w wielokrotnych odpowiedzi.")
            for set_name, set_data in list(st.session_state.mrs_sets.items()):
                # Support both old format (list) and new format (dict)
                if isinstance(set_data, list):
                    set_cols = set_data
                    set_val  = 1
                else:
                    set_cols = set_data.get('cols', [])
                    set_val  = set_data.get('count_val', 1)

                with st.expander(f"**{set_name}** ({len(set_cols)} zmiennych, zliczana: {set_val})", expanded=False):
                    new_mrs_name = st.text_input("Nazwa zestawu:", value=set_name, key=f"mrs_rename_{set_name}")
                    st.caption("Wybierz zmienne numeryczne zestawu:")
                    edited_cols = st.multiselect(
                        "Zmienne zestawu:",
                        options=numeric_mrs_cols,
                        default=[c for c in set_cols if c in numeric_mrs_cols],
                        format_func=lambda x: get_var_display_name(x, var_labels),
                        key=f"mrs_edit_{set_name}"
                    )
                    if edited_cols:
                        edit_vals = sorted(set(
                            v for c in edited_cols
                            for v in df_raw[c].dropna().unique()
                        ))
                        edit_val_options = [int(v) if float(v) == int(v) else float(v) for v in edit_vals]
                    else:
                        edit_val_options = [1]

                    edited_val = st.selectbox(
                        "Warto\u015b\u0107 do zliczania:",
                        options=edit_val_options if edit_val_options else [1],
                        index=edit_val_options.index(set_val) if set_val in edit_val_options else 0,
                        key=f"mrs_val_{set_name}"
                    )

                    col_save, col_del = st.columns([3, 1])
                    with col_save:
                        if st.button("\U0001f4be Zapisz zmiany", key=f"save_mrs_{set_name}", use_container_width=True):
                            if not edited_cols:
                                st.error("Zestaw nie mo\u017ce by\u0107 pusty.")
                            else:
                                del st.session_state.mrs_sets[set_name]
                                final_name = new_mrs_name.strip() or set_name
                                st.session_state.mrs_sets[final_name] = {
                                    'cols': edited_cols,
                                    'count_val': edited_val
                                }
                                st.success(f"Zapisano '{final_name}'.")
                                st.rerun()
                    with col_del:
                        if st.button("\U0001f5d1\ufe0f Usu\u0144", key=f"del_mrs_{set_name}", use_container_width=True):
                            del st.session_state.mrs_sets[set_name]
                            st.rerun()

    # -- BATERIE MATRYCOWE -----------------------------
    with tab_matrix:
        st.markdown("#### Pytania matrycowe (pytania Likerta)")

        with st.expander("Instrukcja \u2014 jak definiowa\u0107 pytania matrycowe (baterie Likerta)", expanded=False):
            st.markdown("""
##### \U0001f3af Do czego s\u0142u\u017cy ta zak\u0142adka?
Pytania matrycowe (baterie) to grupy subpyta\u0144, kt\u00f3re wsp\u00f3\u0142dziel\u0105 t\u0119 sam\u0105 skal\u0119 odpowiedzi,
np. "Q5. Oce\u0144 ka\u017cdy z poni\u017cszych aspekt\u00f3w (1\u20135): Q5_1 obs\u0142uga, Q5_2 cena, Q5_3 jako\u015b\u0107".
Grupuj\u0105c je w zestaw matrycowy, umo\u017cliwiasz tworzenie przejrzystych tabel
ze \u015brednymi lub rozk\u0142adami dla ca\u0142ej baterii jednocze\u015bnie.

##### \U0001f527 Jak u\u017cywa\u0107 \u2014 krok po kroku
1. **Autowykrywanie** \u2014 kliknij "\U0001f50d Autowykrywanie (wsp\u00f3lny prefiks)" \u2014 narz\u0119dzie
   znajdzie grupy kolumn z wsp\u00f3lnym prefiksem (np. `q5_1`, `q5_2`, `q5_3` \u2192 bateria "q5").
2. **R\u0119czne dodanie** \u2014 wpisz nazw\u0119 baterii, wybierz kolumny z listy (min. 2),
   kliknij "Dodaj bateri\u0119".
3. Baterie pojawi\u0105 si\u0119 w prawej kolumnie. Mo\u017cesz je usuwa\u0107 jednym klikni\u0119ciem.

##### \U0001f4ca Gdzie s\u0105 widoczne baterie?
- W module **Analizy \u2192 Matrycowe** \u2014 ka\u017cda bateria tworzy tabel\u0119 matrycow\u0105
  ze \u015bredniali lub procentami dla wszystkich subpyta\u0144.
- W module **Eksport do Excela** \u2014 baterie mog\u0105 by\u0107 eksportowane jako bloki.

##### \U0001f4a1 Wa\u017cne
- Subpytania powinny mie\u0107 t\u0119 sam\u0105 skal\u0119 odpowiedzi (np. 1\u20135 lub 1\u20137).
- Zaleca si\u0119 nadanie etykiet zmiennych (w zak\u0142adce Etykiety) przed definiowaniem baterii
  \u2014 etykiety s\u0105 wtedy wy\u015bwietlane w tabeli matrycowej jako nag\u0142\u00f3wki wierszy.
- Zestawy s\u0105 zapisywane w pliku projektu JSON.
"""
            )

        col1, col2 = st.columns([1, 1])
        with col1:
            st.markdown("**Dodaj bateri\u0119**")
            if st.button("\U0001f50d Autowykrywanie (wsp\u00f3lny prefiks)", use_container_width=True, key="auto_matrix"):
                detected_m = auto_detect_matrix(df_raw)
                added_m = 0
                skipped_m = []
                for name, cols in detected_m.items():
                    key = f"Auto_{name}"
                    if key not in st.session_state.matrix_sets:
                        st.session_state.matrix_sets[key] = cols
                        added_m += 1
                    else:
                        skipped_m.append(key)
                if added_m:
                    st.success(f"Wykryto i dodano {added_m} pyta\u0144 matrycowych.")
                elif not detected_m:
                    st.info("Nie znaleziono kolumn ze wsp\u00f3lnym prefiksem lub wsp\u00f3ln\u0105 cz\u0119\u015bci\u0105 nazwy.")
                else:
                    st.info("Nie dodano nowych baterii \u2014 wszystkie wykryte ju\u017c istniej\u0105.")
                if skipped_m:
                    _skip_m = ", ".join(f"`{s}`" for s in skipped_m[:5])
                    if len(skipped_m) > 5:
                        _skip_m += f" i {len(skipped_m) - 5} wi\u0119cej"
                    st.warning(
                        f"Pomini\u0119to {len(skipped_m)} baterii \u2014 nazwy ju\u017c zaj\u0119te: {_skip_m}. "
                        "Edytuj lub usu\u0144 istniej\u0105ce pytania matrycowe, a nast\u0119pnie uruchom autowykrywanie ponownie.")
                st.rerun()

            matrix_name = st.text_input("Nazwa baterii (np. 'Satysfakcja z produktu'):", key="matrix_new_name")
            matrix_cols_sel = st.multiselect(
                "Subpytania (zmienne numeryczne lub tekstowe):",
                options=[c for c in visible_columns if c not in hidden_cols],
                format_func=lambda x: get_var_display_name(x, var_labels),
                key="matrix_new_cols"
            )
            if st.button("\u2795 Dodaj bateri\u0119", use_container_width=True, key="add_matrix"):
                if matrix_name and len(matrix_cols_sel) >= 2:
                    if matrix_name in st.session_state.matrix_sets:
                        st.error(
                            f"Bateria o nazwie `{matrix_name}` ju\u017c istnieje. "
                            "Zmie\u0144 nazw\u0119 lub edytuj istniej\u0105c\u0105 bateri\u0119 w panelu poni\u017cej.")
                    else:
                        st.session_state.matrix_sets[matrix_name] = matrix_cols_sel
                        st.success(f"Dodano bateri\u0119 '{matrix_name}'.")
                        st.rerun()
                else:
                    st.warning("Podaj nazw\u0119 i wybierz co najmniej 2 subpytania.")

        with col2:
            st.markdown("**Zdefiniowane pytania -- kliknij aby edytowa\u0107**")
            if not st.session_state.matrix_sets:
                st.info("Brak zdefiniowanych pyta\u0144 matrycowych.")
            to_del_mat = None
            for mat_name, mat_cols in list(st.session_state.matrix_sets.items()):
                with st.expander(f"**{mat_name}** ({len(mat_cols)} subpyta\u0144)", expanded=False):
                    new_mat_name = st.text_input("Nazwa pytania matrycowego:", value=mat_name, key=f"matrix_rename_{mat_name}")
                    edited_mat_cols = st.multiselect(
                        "Subpytania:",
                        options=[c for c in visible_columns if c not in hidden_cols],
                        default=[c for c in mat_cols if c in visible_columns],
                        format_func=lambda x: get_var_display_name(x, var_labels),
                        key=f"matrix_edit_{mat_name}"
                    )
                    for c in edited_mat_cols:
                        st.caption(f"  `{c}` -- {var_labels.get(c, '')}")
                    col_ms, col_md = st.columns([3, 1])
                    with col_ms:
                        if st.button("\U0001f4be Zapisz zmiany", key=f"save_mat_{mat_name}", use_container_width=True):
                            if len(edited_mat_cols) >= 2:
                                final_name = new_mat_name.strip() or mat_name
                                del st.session_state.matrix_sets[mat_name]
                                st.session_state.matrix_sets[final_name] = edited_mat_cols
                                st.success(f"Zapisano '{final_name}'.")
                                st.rerun()
                            else:
                                st.error("Pytanie musi mie\u0107 co najmniej 2 subpytania.")
                    with col_md:
                        if st.button("\U0001f5d1\ufe0f Usu\u0144", key=f"del_mat_{mat_name}", use_container_width=True):
                            del st.session_state.matrix_sets[mat_name]
                            st.rerun()

    # -- TOP/BOTTOM BOX --------------------------------
    with tab_box:
        with st.expander("Instrukcja \u2014 jak definiowa\u0107 grupy odpowiedzi (Top/Bottom Box)", expanded=False):
            st.markdown("""
##### \U0001f3af Do czego s\u0142u\u017cy ta zak\u0142adka?
Grupowanie odpowiedzi pozwala \u0142\u0105czy\u0107 kilka warto\u015bci skali w jedn\u0105 grup\u0119.
Typowe zastosowania:
- **Top 2 Box** \u2014 \u0142\u0105czenie ocen 4 i 5 na skali 1\u20135 jako "pozytywna ocena"
- **Bottom 2 Box** \u2014 \u0142\u0105czenie ocen 1 i 2 jako "negatywna ocena"
- **Grupy wiekowe** \u2014 np. 18\u201324, 25\u201334, 35+ jako "m\u0142odzi", "\u015brednio zaawansowani", "starsi"

Grupy s\u0105 widoczne w tabelach cz\u0119sto\u015bci i krzy\u017cowych obok normalnych warto\u015bci.

##### \U0001f527 Jak u\u017cywa\u0107 \u2014 krok po kroku
1. Wybierz **zmienn\u0105** do grupowania.
2. Wpisz **nazw\u0119 grupy** (np. "Top 2 Box", "Oceny negatywne").
3. Wybierz **odpowiedzi**, kt\u00f3re wchodz\u0105 do grupy (multiselect).
4. Kliknij **"\u2795 Dodaj grup\u0119"**.
5. Grupy wy\u015bwietlaj\u0105 si\u0119 po prawej \u2014 mo\u017cesz je usuwa\u0107 per zmiennej.

##### \U0001f4ca Gdzie grupy s\u0105 widoczne?
- W module **Analizy \u2192 Cz\u0119sto\u015bci** \u2014 grupy s\u0105 wy\u015bwietlane jako dodatkowe wiersze
  (zaznaczone nawiasami kwadratowymi, np. `[Top 2 Box]`).
- W module **Analizy \u2192 Krzy\u017cowe** \u2014 grupy pojawiaj\u0105 si\u0119 w rozk\u0142adach.
- W **Eksporcie do Excela** \u2014 grupy mog\u0105 by\u0107 eksportowane jako osobne wiersze.

##### \U0001f4a1 Wa\u017cne
- Mo\u017cesz zdefiniowa\u0107 wiele grup dla tej samej zmiennej.
- Grupy s\u0105 zapisywane w pliku projektu JSON.
"""
            )

        col1, col2 = st.columns([1, 1])
        with col1:
            st.markdown("#### Dodaj grup\u0119 (np. Top/Bottom Box)")
            box_var = st.selectbox("Zmienna:", visible_columns, format_func=lambda x: get_var_display_name(x, var_labels), key="box_var_select")
            if box_var:
                box_name = st.text_input("Nazwa grupy (np. Top 2 Box):")
                box_cats = st.multiselect("Odpowiedzi w grupie:", df[box_var].dropna().unique())
                if st.button("\u2795 Dodaj grup\u0119", use_container_width=True):
                    if box_name and box_cats:
                        st.session_state.box_sets[box_var][f"[{box_name}]"] = box_cats
                        st.success(f"Dodano: {box_name}")
                        st.rerun()
        with col2:
            st.markdown("#### Zdefiniowane grupy")
            if not st.session_state.box_sets:
                st.info("Brak grup Box.")
            for var, boxes in list(st.session_state.box_sets.items()):
                with st.expander(
                    f"**{var_labels.get(var, var)}** (`{var}`) \u2014 {len(boxes)} grup",
                    expanded=False
                ):
                    all_cats = list(df[var].dropna().unique())
                    to_del_b = None
                    for b_name, b_cats in list(boxes.items()):
                        st.markdown(f"**{b_name}**")
                        new_cats = st.multiselect(
                            f"Odpowiedzi w grupie {b_name}:",
                            options=all_cats,
                            default=[c for c in b_cats if c in all_cats],
                            key=f"box_edit_{var}_{b_name}"
                        )
                        sc1, sc2 = st.columns(2)
                        with sc1:
                            if st.button(f"\U0001f4be Zapisz {b_name}",
                                         key=f"save_box_{var}_{b_name}",
                                         use_container_width=True):
                                if new_cats:
                                    st.session_state.box_sets[var][b_name] = new_cats
                                    st.success("Zapisano.")
                                    st.rerun()
                                else:
                                    st.error("Wybierz co najmniej jedn\u0105 odpowied\u017a.")
                        with sc2:
                            if st.button(f"\U0001f5d1\ufe0f Usu\u0144 {b_name}",
                                         key=f"del_box_{var}_{b_name}",
                                         use_container_width=True):
                                to_del_b = b_name

                    if to_del_b:
                        del st.session_state.box_sets[var][to_del_b]
                        if not st.session_state.box_sets[var]:
                            del st.session_state.box_sets[var]
                        st.rerun()

                    st.divider()
                    if st.button(f"\U0001f5d1\ufe0f Usu\u0144 wszystkie grupy dla tej zmiennej",
                                 key=f"del_box_var_{var}", use_container_width=True):
                        del st.session_state.box_sets[var]
                        st.rerun()

    # -- SEGMENTACJA -----------------------------------
    # -- WAZENIE ---------------------------------------
    with tab_weight:
        st.markdown("#### Wa\u017cenie RIM (iteracyjne dopasowanie proporcjonalne)")

        with st.expander("Instrukcja \u2014 jak dzia\u0142aj\u0105 wagi i kiedy ich u\u017cywa\u0107", expanded=False):
            st.markdown("""
##### \U0001f3af Po co s\u0105 wagi?
Wagi (ang. *weights* / *post-stratification weights*) to metoda **korygowania pr\u00f3by**, aby odzwierciedla\u0142a rzeczywiste proporcje w populacji. Je\u015bli w pr\u00f3bie jest np. 60% kobiet, a w populacji 50% \u2014 wyniki b\u0119d\u0105 zaburzone (\u201eprze-reprezentuj\u0105\u201d opinie kobiet). Wagi sprawiaj\u0105, \u017ce ka\u017cda obserwacja \"liczy si\u0119\" odpowiednio mniej lub wi\u0119cej, tak by odtworzy\u0107 za\u0142o\u017cone proporcje.

##### \u2699\ufe0f Algorytm RIM (raking, iterative proportional fitting)
W narz\u0119dziu u\u017cywana jest metoda **RIM** (ta sama co w SPSS). Dzia\u0142a tak:
1. Ka\u017cda obserwacja dostaje startow\u0105 wag\u0119 = 1
2. Dla ka\u017cdej zmiennej po kolei wagi s\u0105 korygowane tak, \u017ceby rozk\u0142ad tej zmiennej (wa\u017cony) pasowa\u0142 do celu
3. Ta korekta psuje poprzednie dopasowania \u2014 wi\u0119c proces powtarza si\u0119 iteracyjnie a\u017c **wszystkie rozk\u0142ady jednocze\u015bnie** s\u0105 zgodne z celami
4. Na ko\u0144cu wagi s\u0105 **normalizowane do sumy = N** (liczba obserwacji) \u2014 zgodnie z konwencj\u0105 SPSS

##### \U0001f527 Jak ustawi\u0107 wagi \u2014 krok po kroku
1. **Wybierz zmienne** po kt\u00f3rych chcesz wa\u017cy\u0107 (np. p\u0142e\u0107, wiek_grupa, region). Tylko zmienne z 2\u201310 kategoriami.
2. Dla ka\u017cdej kategorii **wpisz docelowy odsetek** w populacji (np. Kobieta 50%, M\u0119\u017cczyzna 50%). **Suma musi wynosi\u0107 100%**.
3. Kliknij **\u2696\ufe0f Oblicz wagi**. Narz\u0119dzie wykona iteracje RIM i zapisze wagi.
4. Od tej chwili wszystkie analizy (cz\u0119sto\u015bci, \u015brednie, krzy\u017cowe, regresja, ANOVA\u2026) b\u0119d\u0105 **automatycznie liczone z uwzgl\u0119dnieniem wag**.

##### \U0001f4ca Jak rozpozna\u0107 \u017ce wagi dzia\u0142aj\u0105
- W nag\u0142\u00f3wku ka\u017cdego modu\u0142u pojawia si\u0119 **zielony pasek**: `\u2696\ufe0f Aktywne wagi: N=..., min=..., max=..., zmienne: ...`
- W sidebarze checkbox **\"U\u017cyj wag w analizach\"** jest zaznaczony (mo\u017cesz go odznaczy\u0107 aby chwilowo zobaczy\u0107 niewa\u017cone wyniki bez usuwania wag)
- Kolumny `N` w tabelach pokazuj\u0105 **wa\u017cone liczebno\u015bci** (zaokr\u0105glone do liczb ca\u0142kowitych)
- \u015arednie, procenty, testy istotno\u015bci \u2014 wszystko uwzgl\u0119dnia wagi

##### \U0001f4a1 Kluczowe fakty
- **Suma wag = N** \u2014 dzi\u0119ki normalizacji SPSS-compatible wa\u017cona baza ma t\u0119 sam\u0105 wielko\u015b\u0107 co niewa\u017cona. Nie musisz si\u0119 martwi\u0107 o \u201esztucznie zawy\u017cone\u201d liczebno\u015bci.
- **Min/Max wag** \u2014 warto\u015bci bliskie 1 (np. 0.7\u20131.3) oznaczaj\u0105 \u017ce pr\u00f3ba jest blisko populacji. Skrajne wagi (np. 0.2 lub 5.0) oznaczaj\u0105 znaczne skrzywienie pr\u00f3by \u2014 wyniki staj\u0105 si\u0119 mniej stabilne. Je\u015bli wagi s\u0105 bardzo rozstrzelone, rozwa\u017c **trimming** (obci\u0119cie skrajnych) lub sprawdzenie czy pr\u00f3ba w og\u00f3le nadaje si\u0119 do waszej populacji.
- **Wagi + podzia\u0142 na podzbiory** \u2014 dzia\u0142aj\u0105 razem. W ka\u017cdej grupie podzia\u0142u wagi s\u0105 przycinane do obserwacji tej grupy; suma wag w grupie \u2248 liczba respondent\u00f3w w grupie.
- **Testy istotno\u015bci przy wa\u017conych danych** \u2014 narz\u0119dzie u\u017cywa **efektywnej wielko\u015bci pr\u00f3by (ESS)** dla test\u00f3w t/F. Je\u015bli wagi s\u0105 bardzo nier\u00f3wne, ESS b\u0119dzie znacznie mniejsze ni\u017c N, co poprawnie koryguje obliczenia p-value.

##### \U0001f504 Jak zmieni\u0107 / przeliczy\u0107 / usun\u0105\u0107 wagi
- **Zmieni\u0107 cele** \u2014 rozwi\u0144 \"Aktualne cele wa\u017cenia\", edytuj warto\u015bci, kliknij \"Oblicz wagi\"
- **Usun\u0105\u0107 wagi** \u2014 przycisk \"\U0001f5d1\ufe0f Usu\u0144 wagi\" wy\u015bwietlony obok podsumowania
- **Chwilowo wy\u0142\u0105czy\u0107** \u2014 odznacz checkbox \"U\u017cyj wag w analizach\" w sidebarze; wagi zostaj\u0105 zapisane, ale nie s\u0105 stosowane

##### \U0001f4be Zapis w projekcie
Wagi i cele wa\u017cenia s\u0105 zapisywane w pliku projektu JSON \u2014 przy ponownym wczytaniu projektu zostan\u0105 przywr\u00f3cone bez potrzeby ponownego przeliczania.

##### \u26a0\ufe0f Kiedy NIE nale\u017cy wa\u017cy\u0107?
- Gdy pr\u00f3ba jest ju\u017c reprezentatywna (badania probabilistyczne, pe\u0142na populacja)
- Gdy cele wa\u017cenia nie s\u0105 oparte na rzetelnych danych populacyjnych (wtedy wagi \u201enaprawiaj\u0105\u201d pr\u00f3b\u0119 do b\u0142\u0119dnego wzorca)
- Gdy wagi b\u0119d\u0105 ekstremalnie nier\u00f3wne (> 5.0 albo < 0.2) \u2014 to znak \u017ce pr\u00f3ba jest za bardzo odbiegaj\u0105ca od populacji
            """)

        st.divider()

        # \u2500\u2500 Active weights summary \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
        if st.session_state.weights is not None:
            w_arr = st.session_state.weights
            wc1, wc2, wc3 = st.columns(3)
            wc1.metric("N obserwacji", len(w_arr))
            wc2.metric("Min waga", f"{w_arr.min():.4f}")
            wc3.metric("Max waga", f"{w_arr.max():.4f}")

            if st.session_state.weight_targets:
                with st.expander("Aktualne cele wa\u017cenia (kliknij aby zobaczy\u0107 / zmodyfikowa\u0107)", expanded=False):
                    st.caption("Poni\u017cej mo\u017cesz zmieni\u0107 docelowe odsetki i przelicy\u0107 wagi.")
                    mod_targets = {}
                    mod_valid = True
                    for wv, wt in st.session_state.weight_targets.items():
                        st.markdown(f"**{get_var_display_name(wv, var_labels)}**")
                        mod_targets[wv] = {}
                        cats_w = list(wt.keys())
                        w_cols = st.columns(min(len(cats_w), 4))
                        sum_w = 0.0
                        for i, cat in enumerate(cats_w):
                            cur_pct = round(wt[cat] * 100, 1)
                            new_pct = w_cols[i % 4].number_input(
                                f"{cat}", 0.0, 100.0, cur_pct,
                                key=f"wmod_{wv}_{cat}"
                            )
                            mod_targets[wv][cat] = new_pct / 100.0
                            sum_w += new_pct
                        if not np.isclose(sum_w, 100.0, atol=0.1):
                            st.error(f"{get_var_display_name(wv, var_labels)}: suma = {sum_w:.1f}% (wymagane 100%)")
                            mod_valid = False
                    if mod_valid and _tracked_button("\u2696\ufe0f Przelicz wagi z nowymi celami", "prep", "recalculate_weights", type="primary",
                                                    key="reweight_btn"):
                        st.session_state.weights = calculate_rim_weights(df, mod_targets)
                        st.session_state.weight_targets = mod_targets
                        st.success("\u2705 Wagi przeliczone!")
                        st.rerun()

            if st.button("\U0001f5d1\ufe0f Usu\u0144 wagi", type="secondary", key="del_weights",
                         use_container_width=False):
                st.session_state.weights = None
                st.session_state.weight_targets = {}
                st.success("\u2705 Wagi usuni\u0119te.")
                st.rerun()

            st.divider()
            st.markdown("**Przelicz wagi od nowa (nowe zmienne / cele):**")

        weight_vars = st.multiselect("Zmienne do wa\u017cenia:", visible_columns,
                                      format_func=lambda x: get_var_display_name(x, var_labels),
                                      key="weight_vars_sel")
        if weight_vars:
            st.write("Wprowad\u017a docelowe odsetki (suma musi wynosi\u0107 100%):")
            targets, valid_targets = {}, True
            for w_var in weight_vars:
                st.markdown(f"**{get_var_display_name(w_var, var_labels)}**")
                categories = df[w_var].dropna().unique()
                targets[w_var] = {}
                cols = st.columns(min(len(categories), 4))
                sum_pct = 0
                for i, cat in enumerate(categories):
                    val = cols[i % 4].number_input(f"{cat}", 0.0, 100.0, 0.0, key=f"w_{w_var}_{cat}")
                    targets[w_var][cat] = val / 100.0
                    sum_pct += val
                if not np.isclose(sum_pct, 100.0, atol=0.1):
                    st.error(f"Suma = {sum_pct:.1f}%. Musi wynosi\u0107 100%!")
                    valid_targets = False
            if valid_targets and _tracked_button("\u2696\ufe0f Oblicz wagi", "prep", "calculate_weights", type="primary", key="calc_weights"):
                st.session_state.weights = calculate_rim_weights(df, targets)
                st.session_state.weight_targets = targets
                st.success("\u2705 Wagi obliczone!")
                st.rerun()


    # -- KOLEJNOSC WARTOSCI ------------------------------
    with tab_order:
        st.markdown("### Kolejno\u015b\u0107 wy\u015bwietlania warto\u015bci")

        with st.expander("Instrukcja \u2014 jak definiowa\u0107 kolejno\u015b\u0107 warto\u015bci", expanded=False):
            st.markdown("""
##### \U0001f3af Do czego s\u0142u\u017cy ta zak\u0142adka?
Domy\u015blnie warto\u015bci w tabelach sortowane s\u0105 alfabetycznie lub numerycznie.
Je\u015bli masz skal\u0119 Likerta (np. "zdecydowanie nie", "raczej nie", "nie wiem",
"raczej tak", "zdecydowanie tak"), kolejno\u015b\u0107 alfabetyczna b\u0119dzie b\u0142\u0119dna.
Ta zak\u0142adka pozwala ustawi\u0107 dowolny porz\u0105dek warto\u015bci dla wybranej zmiennej.

##### \U0001f527 Jak u\u017cywa\u0107 \u2014 krok po kroku
1. **Wybierz zmienn\u0105** z listy rozwijanej.
2. W tabeli po lewej pojawi\u0105 si\u0119 aktualne warto\u015bci zmiennej.
   Zmie\u0144 liczby w kolumnie **"Pozycja"** (1 = pierwsza, 2 = druga itd.),
   aby ustawi\u0107 \u017c\u0105dan\u0105 kolejno\u015b\u0107.
3. Kliknij **"\U0001f4be Zapisz kolejno\u015b\u0107"**.
4. Zmienna pojawi si\u0119 na li\u015bcie "Zmienne z niestandardow\u0105 kolejno\u015bci\u0105".
5. Kliknij **"\u270f\ufe0f Edytuj"** aby zmieni\u0107 kolejno\u015b\u0107, lub **"\U0001f5d1\ufe0f Usu\u0144"** aby wr\u00f3ci\u0107
   do domy\u015blnego sortowania.

##### \U0001f4ca Gdzie niestandardowa kolejno\u015b\u0107 jest stosowana?
- Tabele **Cz\u0119sto\u015bci** \u2014 wiersze w kolejno\u015bci zdefiniowanej przez Ciebie.
- Tabele **Krzy\u017cowe** \u2014 zar\u00f3wno wiersze jak i kolumny.
- Tabele **Matrycowe** \u2014 etykiety odpowiedzi.

##### \U0001f4a1 Wa\u017cne
- Warto\u015bci spoza zdefiniowanej kolejno\u015bci b\u0119d\u0105 dodane **na ko\u0144cu** tabeli.
- Zmiana kolejno\u015bci dzia\u0142a natychmiastowo \u2014 aby zobaczy\u0107 efekt, ponownie uruchom analiz\u0119.
- Kolejno\u015bci s\u0105 zapisywane w pliku projektu JSON.
"""
            )

        _vo_all_vars = list(df_raw.columns)
        if not _vo_all_vars:
            st.info("Brak zmiennych do skonfigurowania.")
        else:
            _vo_jump = st.session_state.pop('_vo_jump_to', None)
            if _vo_jump and _vo_jump in _vo_all_vars:
                if 'vo_var_select' in st.session_state:
                    del st.session_state['vo_var_select']
                _vo_default_idx = _vo_all_vars.index(_vo_jump)
            else:
                _vo_default_idx = 0
            _vo_var = st.selectbox(
                "Wybierz zmienn\u0105:",
                _vo_all_vars,
                index=_vo_default_idx,
                format_func=lambda x: get_var_display_name(x, var_labels),
                key="vo_var_select",
            )
            if _vo_var:
                _vo_series = df_raw[_vo_var].dropna()
                _vo_unique = []
                _vo_seen = set()
                if hasattr(df_raw[_vo_var], 'cat'):
                    for _c in df_raw[_vo_var].cat.categories:
                        _cs = str(_c)
                        if _cs not in _vo_seen:
                            _vo_unique.append(_cs)
                            _vo_seen.add(_cs)
                if not _vo_unique:
                    for _v in _vo_series.unique():
                        _vs = str(_v)
                        if _vs not in _vo_seen:
                            _vo_unique.append(_vs)
                            _vo_seen.add(_vs)
                    try:
                        _vo_unique = sorted(_vo_unique, key=lambda x: float(x))
                    except (ValueError, TypeError):
                        _vo_unique = sorted(_vo_unique)

                if not _vo_unique:
                    st.warning("Ta zmienna nie zawiera \u017cadnych warto\u015bci (po odfiltrowaniu brak\u00f3w).")
                else:
                    _vo_existing = st.session_state.value_orders.get(_vo_var, [])
                    _vo_initial = []
                    _seen_init = set()
                    for _x in _vo_existing:
                        _xs = str(_x)
                        if _xs in _vo_seen and _xs not in _seen_init:
                            _vo_initial.append(_xs)
                            _seen_init.add(_xs)
                    for _x in _vo_unique:
                        if _x not in _seen_init:
                            _vo_initial.append(_x)
                            _seen_init.add(_x)

                    _vo_df = pd.DataFrame({
                        'Pozycja':   list(range(1, len(_vo_initial) + 1)),
                        'Warto\u015b\u0107':   _vo_initial,
                    })
                    st.markdown("**Wpisz numer pozycji w kolumnie \u201ePozycja\u201d. Po zapisaniu kolejno\u015b\u0107 zostanie posortowana wg tych numer\u00f3w.**")
                    _vo_edited = st.data_editor(
                        _vo_df,
                        use_container_width=True,
                        hide_index=True,
                        column_config={
                            'Pozycja': st.column_config.NumberColumn(
                                'Pozycja', min_value=1, step=1,
                                help='Ni\u017csza liczba = wcze\u015bniej w wynikach'),
                            'Warto\u015b\u0107': st.column_config.TextColumn('Warto\u015b\u0107', disabled=True),
                        },
                        key=f"vo_editor_{_vo_var}",
                    )
                    _vo_c1, _vo_c2, _vo_c3 = st.columns([2, 2, 6])
                    with _vo_c1:
                        if st.button("\U0001f4be Zapisz kolejno\u015b\u0107", type="primary",
                                     key=f"vo_save_{_vo_var}", use_container_width=True):
                            _vo_sorted = _vo_edited.sort_values('Pozycja', kind='stable')
                            _vo_order = [str(v) for v in _vo_sorted['Warto\u015b\u0107'].tolist()]
                            st.session_state.value_orders[_vo_var] = _vo_order
                            st.session_state._vo_saved = _vo_var
                            st.rerun()
                    with _vo_c2:
                        if _vo_var in st.session_state.value_orders:
                            if st.button("\U0001f5d1\ufe0f Usu\u0144 kolejno\u015b\u0107",
                                         key=f"vo_clear_{_vo_var}", use_container_width=True):
                                del st.session_state.value_orders[_vo_var]
                                st.session_state._vo_cleared = _vo_var
                                st.rerun()
                    if st.session_state.pop('_vo_saved', None) == _vo_var:
                        st.success(f"\u2705 Kolejno\u015b\u0107 dla zmiennej `{_vo_var}` zosta\u0142a zapisana. Uruchom analizy ponownie, aby zobaczy\u0107 efekt.")
                    if st.session_state.pop('_vo_cleared', None) == _vo_var:
                        st.success(f"\u2705 Niestandardowa kolejno\u015b\u0107 dla `{_vo_var}` zosta\u0142a usuni\u0119ta.")

                    if _vo_var in st.session_state.value_orders:
                        _vo_order_disp = ' \u2192 '.join(str(v) for v in st.session_state.value_orders[_vo_var])
                        st.caption(f"Aktualna kolejno\u015b\u0107: {_vo_order_disp}")

            if st.session_state.value_orders:
                st.markdown("---")
                st.markdown("**Zmienne z niestandardow\u0105 kolejno\u015bci\u0105:**")
                if st.session_state.pop('_vo_deleted_global', None):
                    st.success(f"\u2705 Niestandardowa kolejno\u015b\u0107 dla `{st.session_state.pop('_vo_deleted_name', '')}` zosta\u0142a usuni\u0119ta.")
                _vo_items = list(st.session_state.value_orders.items())
                for _vi, (_v_name, _v_ord) in enumerate(_vo_items):
                    _v_disp = get_var_display_name(_v_name, var_labels)
                    _vrow_c1, _vrow_c2, _vrow_c3 = st.columns([8, 1, 1])
                    with _vrow_c1:
                        st.markdown(f"\u2022 `{_v_name}` ({_v_disp}): {' \u2192 '.join(_v_ord)}")
                    with _vrow_c2:
                        if st.button("\u270f\ufe0f", key=f"vo_edit_{_vi}_{_v_name}",
                                     help="Edytuj kolejno\u015b\u0107 tej zmiennej"):
                            st.session_state._vo_jump_to = _v_name
                            st.rerun()
                    with _vrow_c3:
                        if st.button("\U0001f5d1\ufe0f", key=f"vo_del_{_vi}_{_v_name}",
                                     help="Usu\u0144 niestandardow\u0105 kolejno\u015b\u0107 dla tej zmiennej"):
                            del st.session_state.value_orders[_v_name]
                            st.session_state._vo_deleted_global = True
                            st.session_state._vo_deleted_name = _v_name
                            st.rerun()

    # -- PODZIAL NA PODZBIORY (SPSS Split File) --------
    with tab_split:
        st.markdown("### Podzia\u0142 analiz na podzbiory")

        with st.expander("Instrukcja \u2014 jak korzysta\u0107 z podzia\u0142u na podzbiory", expanded=False):
            st.markdown("""
##### \U0001f3af Do czego to s\u0142u\u017cy?
Podzia\u0142 na podzbiory (odpowiednik **SPSS Split File**) pozwala automatycznie wykona\u0107 t\u0119 sam\u0105 analiz\u0119 **osobno dla ka\u017cdej grupy respondent\u00f3w**. Np. chcesz zobaczy\u0107 czy kobiety i m\u0119\u017cczy\u017ani r\u00f3\u017cni\u0105 si\u0119 w odpowiedziach \u2014 zamiast r\u0119cznie filtrowa\u0107 baz\u0119 dwa razy, jedno klikni\u0119cie daje osobne wyniki dla obu grup.

##### \U0001f527 Jak w\u0142\u0105czy\u0107 podzia\u0142 \u2014 krok po kroku
1. **Wybierz zmienn\u0105 grupuj\u0105c\u0105** z listy rozwijanej powy\u017cej. Pokazuj\u0105 si\u0119 tylko zmienne kategoryczne (2-20 unikalnych warto\u015bci), np. `p\u0142e\u0107`, `wiek_grupa`, `region`.
2. Kliknij **\u2705 Zastosuj**. Na g\u00f3rze ekranu pojawi si\u0119 \u017c\u00f3\u0142ty pasek: `\U0001f500 Aktywny podzia\u0142 na podzbiory: [twoja zmienna]` \u2014 widoczny w ka\u017cdym module.
3. **Przejd\u017a do dowolnej analizy** (Cz\u0119sto\u015bci, Krzy\u017cowe, Regresja, ANOVA itd.) i wygeneruj wyniki jak zwykle. System sam podzieli baz\u0119 na grupy.

##### \U0001f4ca Jak rozpozna\u0107 wyniki z podzia\u0142em
Ka\u017cdy wynik ma teraz w tytule etykiet\u0119 grupy oddzielon\u0105 znakiem `|`:
- **Bez podzia\u0142u:** `Q1 \u2014 Ocena produktu`
- **Z podzia\u0142em po `p\u0142e\u0107`:** pojawi\u0105 si\u0119 dwa wyniki:
  - `Q1 \u2014 Ocena produktu | p\u0142e\u0107=Kobieta`
  - `Q1 \u2014 Ocena produktu | p\u0142e\u0107=M\u0119\u017cczyzna`

Wyniki kumuluj\u0105 si\u0119 \u2014 mo\u017cesz wykona\u0107 tak\u017ce analiz\u0119 bez podzia\u0142u, potem w\u0142\u0105czy\u0107 podzia\u0142 i zobaczy\u0107 obie wersje obok siebie.

##### \U0001f501 Jak wy\u0142\u0105czy\u0107 podzia\u0142
Wr\u00f3\u0107 do tej zak\u0142adki, wybierz z listy **`(brak - pe\u0142na baza)`** i kliknij **\u2705 Zastosuj**. Nast\u0119pne analizy b\u0119d\u0105 wykonywane na pe\u0142nej bazie.

##### \U0001f4a1 Wa\u017cne
- **Poprzednie wyniki nie s\u0105 kasowane.** W\u0142\u0105czenie / wy\u0142\u0105czenie podzia\u0142u nie zmienia ju\u017c wygenerowanych tabel \u2014 wp\u0142ywa tylko na kolejne analizy, kt\u00f3re uruchomisz.
- **Ka\u017cdy modu\u0142 pami\u0119ta wyniki osobno** \u2014 np. mo\u017cesz mie\u0107 trzy tabele cz\u0119sto\u015bci tej samej zmiennej: jedn\u0105 dla pe\u0142nej bazy, jedn\u0105 dla Kobiet, jedn\u0105 dla M\u0119\u017cczyzn.
- **Skupienia hierarchiczne** dzia\u0142aj\u0105 inaczej \u2014 dla ka\u017cdej grupy tworzona jest osobna zmienna w bazie (z sufiksem, np. `Skupienie_H_plec_Kobieta`), bo ka\u017cda grupa ma w\u0142asne klastry.
- **Wa\u017cenie** (je\u015bli w\u0142\u0105czone) dzia\u0142a wewn\u0105trz ka\u017cdej grupy \u2014 wagi zostaj\u0105 odpowiednio przypisane do obserwacji tej grupy.
- **Ma\u0142e grupy** \u2014 je\u015bli grupa ma za ma\u0142o obserwacji do danej analizy (np. <20 dla regresji), w tej grupie wynik zostanie pomini\u0119ty z komunikatem, ale inne grupy b\u0119d\u0105 policzone.

##### \U0001f4be Zapis w projekcie
Aktywny podzia\u0142 jest zapisywany w pliku projektu JSON \u2014 przy ponownym wczytaniu projektu podzia\u0142 zostanie przywr\u00f3cony.
            """)

        #st.divider()

        # Candidate variables: categorical with 2-20 unique values
        _split_candidates = []
        for c in df.columns:
            try:
                n_unique = df[c].dropna().nunique()
                if 2 <= n_unique <= 20:
                    _split_candidates.append(c)
            except Exception:
                pass

        _current = st.session_state.split_var
        _opts = ["(brak - pe\u0142na baza)"] + _split_candidates
        _default_idx = _opts.index(_current) if _current in _opts else 0

        split_choice = st.selectbox(
            "Zmienna grupuj\u0105ca:",
            _opts,
            index=_default_idx,
            format_func=lambda x: x if x == "(brak - pe\u0142na baza)"
                                    else f"[{x}] {var_labels.get(x, x)}",
            key="split_var_select",
            help="Wybierz zmienn\u0105 kategoryczn\u0105 (2-20 kategorii). "
                 "Wszystkie analizy b\u0119d\u0105 podzielone na grupy wg jej warto\u015bci."
        )

        if _tracked_button("\u2705 Zastosuj", "prep", "apply_split", type="primary", key="split_apply"):
            new_val = None if split_choice == "(brak - pe\u0142na baza)" else split_choice
            st.session_state.split_var = new_val
            if new_val:
                st.success(f"\u2705 Podzia\u0142 aktywny: `{new_val}`. Wszystkie analizy b\u0119d\u0105 liczone per grupa.")
            else:
                st.success("\u2705 Podzia\u0142 wy\u0142\u0105czony. Analizy wykonywane na pe\u0142nej bazie.")
            st.rerun()

        # Status
        st.divider()
        if st.session_state.split_var:
            lbl = var_labels.get(st.session_state.split_var, st.session_state.split_var)
            groups_preview = df[st.session_state.split_var].dropna().unique()
            try:
                groups_preview = sorted(groups_preview, key=lambda x: str(x))
            except Exception:
                pass
            st.success(
                f"\U0001f500 **Aktywny podzia\u0142:** `{st.session_state.split_var}` \u2014 {lbl}  \n"
                f"**Liczba grup:** {len(groups_preview)}  \n"
                f"**Grupy:** {', '.join(str(g) for g in groups_preview[:10])}"
                + ("..." if len(groups_preview) > 10 else "")
            )
        else:
            st.info("Brak aktywnego podzia\u0142u \u2014 analizy wykonywane na pe\u0142nej bazie.")



# -------------------------------------------------------------
# MODU? 3: ANALIZY I TABELE
# -------------------------------------------------------------
elif menu == "\U0001f4c8 Analizy i Tabele":
    _require_module_access("analyses")
    _require_data()
    module_header("\U0001f4c8", "Analizy i Tabele")
    tab_freq, tab_matrix_an, tab_cross, tab_banner, tab_means, tab_desc, tab_corr = st.tabs([
        "Cz\u0119sto\u015bci", "Pytania Matrycowe", "Tabele Krzy\u017cowe", "Tabele zbiorcze (Banner)", "\u015arednie (T-Test)", "Statystyki Opisowe", "Korelacje"
    ])

    with tab_freq:
        freq_vars = st.multiselect("Wybierz zmienne:", all_options, format_func=lambda x: get_var_display_name(x, var_labels))
        show_charts_freq = st.checkbox("Wy\u015bwietlaj wykresy", key="charts_freq")
        if _tracked_button("\u25b6\ufe0f Generuj tablice cz\u0119sto\u015bci", "analyses", "freq_table", type="primary") and freq_vars:
            _w_full = st.session_state.weights if use_weights else np.ones(len(df_raw))
            for freq_var in freq_vars:
                # Iterate over split groups (single iteration if no split)
                for _grp_lbl, _df_s, _df_raw_s, _w_s in _iter_split_groups(
                        df, df_raw, var_labels, st.session_state.split_var,
                        weights=_w_full):
                    w = _w_s if _w_s is not None else np.ones(len(_df_raw_s))
                    # Key in results dict includes group label to avoid overwrites
                    result_key = f"{freq_var} | {_grp_lbl}" if _grp_lbl else freq_var
                    expander_title = (f"{get_var_display_name(freq_var, var_labels)} \u2014 {_grp_lbl}"
                                      if _grp_lbl else get_var_display_name(freq_var, var_labels))
                    if freq_var in st.session_state.matrix_sets:
                        # Matrix variable -- show as transposed freq table
                        mat_cols = st.session_state.matrix_sets[freq_var]
                        try:
                            mat_df, cats, sub_lbls = build_matrix_table(_df_s, _df_raw_s, mat_cols, var_labels, w, meta_orig.variable_value_labels, st.session_state.custom_val_labels)
                            st.session_state.results['czestosci'][result_key] = mat_df
                            with st.expander(f"[Pytanie matrycowe] {expander_title}", expanded=True):
                                pct_cols_m = [f"{s} [%]" for s in sub_lbls]
                                n_cols_m   = [f"{s} [N]"  for s in sub_lbls]
                                st.dataframe(
                                    mat_df.style
                                        .format(lambda x: f"{x:.0f}" if isinstance(x, (int, float)) else x, subset=n_cols_m)
                                        .format(lambda x: f"{x:.1f}%" if isinstance(x, (int, float)) else x, subset=pct_cols_m),
                                    use_container_width=True
                                )
                        except Exception as e:
                            st.error(f"B\u0142\u0105d dla baterii {freq_var}: {e}")
                        continue

                    if freq_var in st.session_state.mrs_sets:
                        set_data = st.session_state.mrs_sets[freq_var]
                        if isinstance(set_data, list):
                            cols = set_data
                            count_val = 1
                        else:
                            cols = set_data.get('cols', [])
                            count_val = set_data.get('count_val', 1)
                        missing_mask = _df_raw_s[cols].isna().all(axis=1)
                        missing_count = w[missing_mask].sum()
                        counts = pd.Series({var_labels.get(c, c): w[(_df_raw_s[c] == count_val).values].sum() for c in cols})
                        total_respondents = w[~missing_mask].sum()
                        pcts = (counts / total_respondents) * 100 if total_respondents > 0 else counts * 0
                        res_df = pd.DataFrame({'Liczebnosc [N]': counts, 'Procent [%]': pcts})
                        res_df.loc['Og\u00f3\u0142em (Wa\u017cne)'] = [total_respondents, pcts.sum()]
                        res_df.loc['Braki danych'] = [missing_count, np.nan]
                    else:
                        missing_mask = _df_s[freq_var].isna()
                        valid_df = pd.DataFrame({'val': _df_s[freq_var], 'w': w}).dropna()
                        counts = valid_df.groupby('val', observed=False, sort=False)['w'].sum()
                        # Restore SPSS-defined category order when column is Categorical
                        if hasattr(_df_s[freq_var], 'cat'):
                            _freq_order = [str(c) for c in _df_s[freq_var].cat.categories]
                            _present = [c for c in _freq_order if c in counts.index]
                            _extra   = [c for c in counts.index if c not in _freq_order]
                            counts = counts.reindex(_present + _extra)
                        # Apply user-defined value order (overrides SPSS / alphabetical)
                        if freq_var in st.session_state.value_orders:
                            _vo_idx = _apply_value_order(counts.index, freq_var)
                            counts = counts.reindex(_vo_idx)

                        # Apply custom value labels to rename index (e.g. 1\u2192Kobieta)
                        _cvl_freq = st.session_state.custom_val_labels.get(freq_var, {})
                        if _cvl_freq:
                            counts.index = counts.index.map(
                                lambda x: _cvl_freq.get(str(x), _cvl_freq.get(x, x))
                            )
                        if freq_var in st.session_state.box_sets:
                            for box_name, b_cats in st.session_state.box_sets[freq_var].items():
                                box_val = counts[counts.index.isin(b_cats)].sum()
                                counts.loc[box_name] = box_val
                        pcts = (counts / valid_df['w'].sum()) * 100 if valid_df['w'].sum() > 0 else counts * 0
                        res_df = pd.DataFrame({'Liczebnosc [N]': counts, 'Procent [%]': pcts})
                        sum_n = counts[~counts.index.astype(str).str.startswith('[')].sum()
                        res_df.loc['Suma'] = [sum_n, 100.0 if sum_n > 0 else 0]

                        missing_count = w[missing_mask.values].sum()
                        res_df.loc['Braki danych'] = [missing_count, np.nan]
                    st.session_state.results['czestosci'][result_key] = res_df
                    with st.expander(expander_title, expanded=True):
                        st.dataframe(res_df.style.format(get_streamlit_format(res_df)), use_container_width=True)
                        if show_charts_freq:
                            plot_df = res_df.drop(index=['Suma', 'Og\u00f3\u0142em (Wa\u017cne)', 'Braki danych'], errors='ignore')
                            plot_df = plot_df[~plot_df.index.astype(str).str.startswith('[')]
                            if not plot_df.empty and 'Procent [%]' in plot_df.columns:
                                plot_df = plot_df.dropna(subset=['Procent [%]'])
                            if not plot_df.empty and 'Procent [%]' in plot_df.columns:
                                fig = px.bar(plot_df, x='Procent [%]', y=plot_df.index, orientation='h',
                                             title=expander_title, color_discrete_sequence=['#2E75B6'])
                                fig.update_layout(yaxis={'categoryorder': 'total ascending'}, height=max(300, len(plot_df) * 30 + 80))
                                st.plotly_chart(fig, use_container_width=True, key=f"pc_freq_gen_{result_key}")
            st.success("\u2705 Tablice cz\u0119sto\u015bci wygenerowane!")

        # \u2500\u2500 Persistent display of stored results \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
        if st.session_state.results.get('czestosci'):
            st.divider()
            _fc1, _fc2 = st.columns([5, 1])
            _fc1.markdown(f"**Zapisane wyniki ({len(st.session_state.results['czestosci'])} tablic):**")
            with _fc2:
                if st.button("\U0001f5d1\ufe0f Usu\u0144 wszystkie", key="del_all_freq",
                             use_container_width=True):
                    st.session_state.results['czestosci'] = {}
                    st.rerun()
            for freq_var, res_df in list(st.session_state.results['czestosci'].items()):
                _base_var, _grp_lbl = _extract_split_from_title(freq_var)
                _ec1, _ec2 = st.columns([6, 1])
                with _ec1:
                    _display_title = get_var_display_name(_base_var, var_labels)
                    if _grp_lbl:
                        _display_title += f" \u2014 \U0001f500 {_grp_lbl}"
                    _exp = st.expander(_display_title, expanded=False)
                with _ec2:
                    if st.button("\U0001f5d1\ufe0f", key=f"del_freq_{freq_var}",
                                 help=f"Usu\u0144 wynik dla {freq_var}"):
                        st.session_state.results['czestosci'].pop(freq_var, None)
                        st.rerun()
                with _exp:
                    _split_badge(_grp_lbl)
                    st.dataframe(res_df.style.format(get_streamlit_format(res_df)),
                                 use_container_width=True)
                    if show_charts_freq:
                        plot_df = res_df.drop(
                            index=['Suma','Og\u00f3\u0142em (Wa\u017cne)','Braki danych'], errors='ignore')
                        plot_df = plot_df[~plot_df.index.astype(str).str.startswith('[')]
                        if not plot_df.empty and 'Procent [%]' in plot_df.columns:
                            fig = px.bar(plot_df, x='Procent [%]', y=plot_df.index,
                                         orientation='h',
                                         title=var_labels.get(freq_var, freq_var),
                                         color_discrete_sequence=['#2E75B6'])
                            fig.update_layout(yaxis={'categoryorder':'total ascending'}, height=350)
                            st.plotly_chart(fig, use_container_width=True, key=f"pc_freq_saved_{freq_var}")

    # -- BATERIE MATRYCOWE -- dedykowana zak?adka analityczna --
    with tab_matrix_an:
        st.markdown("#### Tabele matrycowe (pytania matrycowe (Likert))")
        st.info("Wiersze = warto\u015bci skali, Kolumny = subpytania. Dost\u0119pne tryby: tylko N, tylko %, lub N + %.")

        if not st.session_state.matrix_sets:
            st.warning("Brak zdefiniowanych pyta\u0144 matrycowych. Przejd\u017a do **Przygotowanie Danych \u2192 Pytania Matrycowe** i dodaj baterie.")
        else:
            mat_sel = st.multiselect(
                "Wybierz pytania matrycowe:",
                list(st.session_state.matrix_sets.keys()),
                default=list(st.session_state.matrix_sets.keys()),
                key="matrix_an_sel"
            )

            col_mode, col_chart = st.columns([2, 1])
            with col_mode:
                mat_display_mode = st.radio(
                    "Prezentuj warto\u015bci:",
                    ["N + %", "Tylko N", "Tylko %"],
                    index=0, horizontal=True, key="mat_display_mode"
                )
            with col_chart:
                show_chart_mat = st.checkbox("\U0001f4ca Wy\u015bwietl wykres", key="chart_mat")

            if _tracked_button("\u25b6\ufe0f Generuj tabele matrycowe", "analyses", "matrix_table", type="primary", key="gen_matrix"):
                _w_full = st.session_state.weights if use_weights else np.ones(len(df_raw))
                for _grp_lbl, _df_s, _df_raw_s, _w_s in _iter_split_groups(
                        df, df_raw, var_labels, st.session_state.split_var, weights=_w_full):
                    w = _w_s if _w_s is not None else np.ones(len(_df_raw_s))
                    for mat_name in mat_sel:
                        mat_cols = st.session_state.matrix_sets[mat_name]
                        try:
                            mat_df, cats, sub_lbls = build_matrix_table(_df_s, _df_raw_s, mat_cols, var_labels, w, meta_orig.variable_value_labels, st.session_state.custom_val_labels)
                            entry_name = f"{mat_name} | {_grp_lbl}" if _grp_lbl else mat_name
                            _mat_mode_map = {"N + %": "N+%", "Tylko N": "N", "Tylko %": "%"}
                            _mat_sfx = _mat_mode_map.get(mat_display_mode, mat_display_mode)
                            entry_name = f"{entry_name} [{_mat_sfx}]"
                            _merge_result(st.session_state.matrix_results, {
                                'name': entry_name, 'df': mat_df, 'cats': cats,
                                'sub_labels': sub_lbls, 'cols': mat_cols,
                                'display_mode': mat_display_mode,
                            }, key_fn=lambda r: r['name'])
                        except Exception as e:
                            st.error(f"B\u0142\u0105d dla '{mat_name}': {e}")
                st.success(f"\u2705 Wygenerowano {len(st.session_state.matrix_results)} tabel matrycowych.")

            if st.session_state.matrix_results:
                _mrc1, _mrc2 = st.columns([5, 1])
                _mrc1.markdown(f"**Zapisane tabele matrycowe ({len(st.session_state.matrix_results)}):**")
                with _mrc2:
                    if st.button("\U0001f5d1\ufe0f Usu\u0144 wszystkie", key="del_all_matrix",
                                 use_container_width=True):
                        st.session_state.matrix_results = []
                        st.rerun()

            for _mi, entry in enumerate(list(st.session_state.matrix_results)):
                _base_n, _grp_n = _extract_split_from_title(entry['name'])
                _mtc1, _mtc2 = st.columns([6, 1])
                with _mtc1:
                    _mtexp = st.expander(
                        f"\U0001f522 {_base_n}" + (f" \u2014 \U0001f500 {_grp_n}" if _grp_n else ""),
                        expanded=True
                    )
                with _mtc2:
                    if st.button("\U0001f5d1\ufe0f", key=f"del_matrix_{_mi}",
                                 help=f"Usu\u0144 {entry['name']}"):
                        st.session_state.matrix_results.pop(_mi)
                        st.rerun()
                with _mtexp:
                    _split_badge(_grp_n)
                    mat_df   = entry['df']
                    cats     = entry['cats']
                    sub_lbls = entry.get('sub_labels', [])

                    pct_cols_m = [f"{s} [%]" for s in sub_lbls]
                    n_cols_m   = [f"{s} [N]"  for s in sub_lbls]

                    # Build view according to display mode
                    SUMROW = "Baza (N) / Suma (%)"
                    if mat_display_mode == "Tylko N":
                        view_cols = n_cols_m
                        view_df   = mat_df[view_cols].copy()
                    elif mat_display_mode == "Tylko %":
                        view_cols = pct_cols_m
                        view_df   = mat_df[view_cols].copy()
                    else:  # N + %
                        view_cols = []
                        for s in sub_lbls:
                            view_cols += [f"{s} [N]", f"{s} [%]"]
                        view_df = mat_df[view_cols].copy()

                    _style_matrix_row = _make_style_matrix_row(SUMROW)

                    # Deduplicate index/columns before styling \u2014 Styler crashes on duplicates
                    _sdf = view_df.copy()
                    if not _sdf.index.is_unique:
                        _seen = {}
                        _new_idx = []
                        for v in _sdf.index:
                            k = str(v)
                            if k in _seen:
                                _seen[k] += 1
                                _new_idx.append(f"{v} ({_seen[k]})")
                            else:
                                _seen[k] = 1
                                _new_idx.append(v)
                        _sdf.index = _new_idx
                    if not _sdf.columns.is_unique:
                        _seen = {}
                        _new_cols = []
                        for v in _sdf.columns:
                            k = str(v)
                            if k in _seen:
                                _seen[k] += 1
                                _new_cols.append(f"{v} ({_seen[k]})")
                            else:
                                _seen[k] = 1
                                _new_cols.append(v)
                        _sdf.columns = _new_cols

                    # Per-column format: N columns as integer, % columns as 1 decimal
                    _fmt_map = {}
                    for _c in _sdf.columns:
                        _cstr = str(_c)
                        if "[N]" in _cstr:
                            _fmt_map[_c] = lambda v: ("" if (v == "" or (isinstance(v, float) and np.isnan(v)))
                                                       else (f"{int(round(float(v))):,}"
                                                             if isinstance(v, (int, float)) and not (isinstance(v, float) and np.isnan(v))
                                                             else str(v)))
                        elif "[%]" in _cstr:
                            _fmt_map[_c] = _fmt_cell
                        else:
                            _fmt_map[_c] = _fmt_cell

                    styled = _sdf.style.apply(_style_matrix_row, axis=1).format(_fmt_map)
                    st.dataframe(styled, use_container_width=True)

                    if show_chart_mat and cats and sub_lbls:
                        chart_rows = [c for c in cats if c in mat_df.index]
                        chart_data = pd.DataFrame(index=chart_rows, columns=sub_lbls, dtype=float)
                        for sub_l in sub_lbls:
                            col_key = f"{sub_l} [%]"
                            if col_key in mat_df.columns:
                                for cat_v in chart_rows:
                                    try:
                                        chart_data.loc[cat_v, sub_l] = float(mat_df.loc[cat_v, col_key])
                                    except Exception:
                                        chart_data.loc[cat_v, sub_l] = 0
                        chart_data = chart_data.reset_index().rename(columns={'index': 'Wartosc'})
                        chart_long = chart_data.melt(id_vars='Wartosc', var_name='Subpytanie', value_name='%')
                        fig_mat = px.bar(
                            chart_long, x='%', y='Wartosc', color='Subpytanie',
                            barmode='group', orientation='h', title=entry['name'],
                            color_discrete_sequence=px.colors.qualitative.Set2
                        )
                        fig_mat.update_layout(yaxis={'categoryorder': 'category ascending'},
                                              height=max(300, len(cats) * 40 + 80))
                        st.plotly_chart(fig_mat, use_container_width=True, key=f"pc_mat_{entry.get('name','m')}")

    with tab_cross:
        with st.expander("Jak wykona\u0107 i interpretowa\u0107 tabel\u0119 krzy\u017cow\u0105", expanded=False):
            st.markdown("""
### Czym jest tabela krzy\u017cowa?

Tabela krzy\u017cowa (tabulacja krzy\u017cowa) pokazuje jak rozk\u0142ada si\u0119 jedna zmienna wzgl\u0119dem innej.
Pozwala sprawdzi\u0107 czy istnieje zwi\u0105zek mi\u0119dzy dwiema zmiennymi kategorycznymi.

---

### Jak wype\u0142ni\u0107 pola?

- **Zmienne w wierszach** \u2014 zmienna kt\u00f3r\u0105 chcesz analizowa\u0107 (np. opinia, zachowanie)
- **Zmienne w kolumnach** \u2014 zmienna grupuj\u0105ca (np. p\u0142e\u0107, wiek, region)

> Wskaz\u00f3wka: w kolumnach najcz\u0119\u015bciej umieszcza si\u0119 metryczk\u0119 (p\u0142e\u0107, wykszta\u0142cenie), a w wierszach pytanie badawcze.

---

### Spos\u00f3b prezentacji

| Opcja | Kiedy u\u017cywa\u0107 |
|---|---|
| **Liczebno\u015bci** | Por\u00f3wnanie surowych N, ma\u0142e pr\u00f3by |
| **Kolumnowe (%)** | Najcz\u0119\u015bcej stosowane \u2014 % w ramach ka\u017cdej grupy z kolumny (np. % kobiet kt\u00f3re powiedzia\u0142y TAK) |
| **Wierszowe (%)** | % w ramach ka\u017cdego wiersza \u2014 przydatne gdy zmienna wierszowa jest pytaniem wielokrotnego wyboru |
| **Liczebno\u015bci + %** | Pe\u0142ny obraz: N i % razem |

---

### Testy statystyczne

- **Chi-kwadrat (\u03c7\u00b2)** \u2014 sprawdza czy rozk\u0142ad zmiennej jest niezale\u017cny od grupy
  - p < 0.05 \u2192 istnieje statystycznie istotny zwi\u0105zek
  - p \u2265 0.05 \u2192 brak podstaw do odrzucenia niezale\u017cno\u015bci
- **V Kramera** \u2014 si\u0142a zwi\u0105zku niezale\u017cnie od rozmiaru tabeli
  - < 0.1 brak | 0.1\u20130.3 s\u0142aby | 0.3\u20130.5 umiarkowany | > 0.5 silny
- **Testy Z (95%)** \u2014 por\u00f3wnanie par kolumn \u2014 oznaczenia literowe (A, B...) wskazuj\u0105 istotne r\u00f3\u017cnice mi\u0119dzy grupami

---

### Uwagi praktyczne

- Tabela krzy\u017cowa wymaga zmiennych **kategorycznych** \u2014 unikaj zmiennych z du\u017c\u0105 liczb\u0105 unikalnych warto\u015bci
- Przy ma\u0142ych liczebno\u015bciach (N < 5 w kom\u00f3rce) wyniki chi-kwadrat s\u0105 nierzetelne
- Przy wa\u017ceniu \u2014 N i % s\u0105 wa\u017cone, ale chi-kwadrat obliczany jest na wa\u017conych liczebno\u015bciach
            """)

        col1, col2 = st.columns(2)
        with col1: row_vars = st.multiselect("Zmienne w wierszach:", all_options_no_matrix, format_func=lambda x: get_var_display_name(x, var_labels))
        with col2: col_vars = st.multiselect("Zmienne w kolumnach:", all_options_no_matrix, format_func=lambda x: get_var_display_name(x, var_labels))
        pct_type = st.radio("Spos\u00f3b prezentacji:", ["Liczebno\u015bci", "Kolumnowe (%)", "Wierszowe (%)", "Liczebno\u015bci + Kolumnowe (%)", "Liczebno\u015bci + Wierszowe (%)"], horizontal=True)
        c1, c2, c3, c4 = st.columns(4)
        do_sig_test = c1.checkbox("\U0001f520 Testy Z (95%)")
        do_chi_square = c2.checkbox("\U0001f9ee Chi-kwadrat")
        do_cramer = c3.checkbox("\U0001f4cf V Kramera", help="Si\u0142a zwi\u0105zku: 0=brak, 0.1=s\u0142aby, 0.3=umiarkowany, 0.5+=silny")
        show_charts_cross = c4.checkbox("\U0001f4ca Wykresy")

        if _tracked_button("\u25b6\ufe0f Generuj tabele krzy\u017cowe", "analyses", "crosstab", type="primary") and row_vars and col_vars:
            _w_full = st.session_state.weights if use_weights else np.ones(len(df_raw))
            for _grp_lbl, _df_s, _df_raw_s, _w_s in _iter_split_groups(
                    df, df_raw, var_labels, st.session_state.split_var, weights=_w_full):
                w = _w_s if _w_s is not None else np.ones(len(_df_raw_s))
                for row_var in row_vars:
                    for col_var in col_vars:
                        is_row_mrs = row_var in st.session_state.mrs_sets
                        is_col_mrs = col_var in st.session_state.mrs_sets
                        try:
                            tmp_df = pd.DataFrame({'w': w})
                            if is_row_mrs: tmp_df['R_miss'] = _df_raw_s[st.session_state.mrs_sets[row_var]].isna().all(axis=1).values
                            else: tmp_df['R_miss'] = _df_s[row_var].isna().values
                            if is_col_mrs: tmp_df['C_miss'] = _df_raw_s[st.session_state.mrs_sets[col_var]].isna().all(axis=1).values
                            else: tmp_df['C_miss'] = _df_s[col_var].isna().values
                            missing_count = tmp_df.loc[tmp_df['R_miss'] | tmp_df['C_miss'], 'w'].sum()

                            if not is_row_mrs and not is_col_mrs:
                                df_n = pd.crosstab(_df_s[row_var], _df_s[col_var], values=w, aggfunc='sum', dropna=False).fillna(0)
                            elif is_row_mrs and not is_col_mrs:
                                cols = st.session_state.mrs_sets[row_var]
                                mrs_w = _df_raw_s[cols].replace(np.nan, 0).multiply(w, axis=0)
                                df_n = mrs_w.groupby(_df_s[col_var].values, observed=False).sum().T
                                df_n.index = [var_labels.get(c, c) for c in df_n.index]
                            elif not is_row_mrs and is_col_mrs:
                                cols = st.session_state.mrs_sets[col_var]
                                mrs_w = _df_raw_s[cols].replace(np.nan, 0).multiply(w, axis=0)
                                df_n = mrs_w.groupby(_df_s[row_var].values, observed=False).sum()
                                df_n.columns = [var_labels.get(c, c) for c in df_n.columns]
                            # Apply user-defined value order to rows (when not MRS)
                            if not is_row_mrs and row_var in st.session_state.value_orders:
                                _vo_ridx = _apply_value_order(df_n.index, row_var)
                                df_n = df_n.reindex(index=_vo_ridx, fill_value=0)
                            if not is_col_mrs and col_var in st.session_state.value_orders:
                                _vo_cidx = _apply_value_order(df_n.columns, col_var)
                                df_n = df_n.reindex(columns=_vo_cidx, fill_value=0)

                            if row_var in st.session_state.box_sets and not is_row_mrs:
                                for box_name, b_cats in st.session_state.box_sets[row_var].items():
                                    df_n.loc[box_name] = df_n.loc[df_n.index.intersection(b_cats)].sum(axis=0)

                            title_base = f"{row_var} x {col_var}"
                            title = f"{title_base} | {_grp_lbl}" if _grp_lbl else title_base
                            _pct_sfx_map = {
                                "Liczebno\u015bci": "N",
                                "Kolumnowe (%)": "kol.%",
                                "Wierszowe (%)": "wier.%",
                                "Liczebno\u015bci + Kolumnowe (%)": "N+kol.%",
                                "Liczebno\u015bci + Wierszowe (%)": "N+wier.%",
                            }
                            _pct_sfx = _pct_sfx_map.get(pct_type, pct_type)
                            title = f"{title} [{_pct_sfx}]"

                            if (do_chi_square or do_cramer) and not is_row_mrs and not is_col_mrs:
                                obs = df_n.loc[~df_n.index.astype(str).str.startswith('[')]
                                obs = obs[[c for c in obs.columns if not str(c).startswith('[') and c != 'Suma']]

                                # Pre-check: identify empty rows/columns BEFORE calling scipy
                                _row_sums = obs.sum(axis=1)
                                _col_sums = obs.sum(axis=0)
                                _empty_rows = _row_sums[_row_sums == 0].index.tolist()
                                _empty_cols = _col_sums[_col_sums == 0].index.tolist()

                                if obs.shape[0] <= 1 or obs.shape[1] <= 1:
                                    st.session_state.chi_results[title] = (
                                        "\u26a0\ufe0f Test nie mo\u017ce by\u0107 wykonany: tabela musi mie\u0107 co najmniej 2 wiersze i 2 kolumny "
                                        "(po usuni\u0119ciu zagregowanych wierszy/kolumn)."
                                    )
                                elif obs.sum().sum() == 0:
                                    st.session_state.chi_results[title] = (
                                        "\u26a0\ufe0f Test nie mo\u017ce by\u0107 wykonany: tabela jest ca\u0142kowicie pusta (wszystkie liczebno\u015bci = 0)."
                                    )
                                elif _empty_rows or _empty_cols:
                                    _msg_parts = []
                                    if _empty_rows:
                                        _msg_parts.append(f"wiersze bez obserwacji: {', '.join(str(r) for r in _empty_rows)}")
                                    if _empty_cols:
                                        _msg_parts.append(f"kolumny bez obserwacji: {', '.join(str(c) for c in _empty_cols)}")
                                    st.session_state.chi_results[title] = (
                                        "\u26a0\ufe0f Test Chi\u00b2 niedost\u0119pny \u2014 " + "; ".join(_msg_parts) + ". "
                                        "Przyczyna: te kategorie istniej\u0105 w s\u0142owniku warto\u015bci, ale nikt ich nie wybra\u0142 "
                                        "(lub wszystkie odpowiedzi w tej kategorii s\u0105 zakodowane jako braki). "
                                        "Aby wykona\u0107 test: usu\u0144 nieu\u017cywane kategorie w zak\u0142adce Etykiety warto\u015bci "
                                        "lub dodaj je do brak\u00f3w danych."
                                    )
                                else:
                                    try:
                                        chi2, p, dof, ex = stats.chi2_contingency(obs)
                                        n_total = obs.sum().sum()
                                        k = min(obs.shape[0], obs.shape[1])
                                        cramer_v = float(np.sqrt(chi2 / (n_total * (k - 1)))) if n_total > 0 and k > 1 else 0.0

                                        if cramer_v < 0.1: v_interp = "brak/zaniedbywalny"
                                        elif cramer_v < 0.3: v_interp = "s\u0142aby"
                                        elif cramer_v < 0.5: v_interp = "umiarkowany"
                                        else: v_interp = "silny"

                                        parts = []
                                        if do_chi_square:
                                            parts.append(f"Chi\u00b2={chi2:.2f}, df={dof}, p={p:.3f}")
                                        if do_cramer:
                                            parts.append(f"V Kramera={cramer_v:.3f} ({v_interp})")

                                        # Warn if >20% of cells have expected < 5 (Cochran's rule)
                                        _low_expected = int((ex < 5).sum())
                                        _total_cells = ex.size
                                        if _total_cells > 0 and _low_expected / _total_cells > 0.20:
                                            parts.append(
                                                f"\u26a0\ufe0f {_low_expected}/{_total_cells} kom\u00f3rek ma oczekiwan\u0105 liczebno\u015b\u0107 < 5 "
                                                "(test mo\u017ce by\u0107 niestabilny; rozwa\u017c Fisher Exact lub \u0142\u0105czenie kategorii)"
                                            )
                                        st.session_state.chi_results[title] = " | ".join(parts)
                                    except Exception as _chi_err:
                                        st.session_state.chi_results[title] = (
                                            f"\u26a0\ufe0f Test Chi\u00b2 nie powi\u00f3d\u0142 si\u0119: {_chi_err}. "
                                            "Sprawd\u017a czy w tabeli nie ma pustych wierszy/kolumn."
                                        )

                            if not is_row_mrs and not is_col_mrs:
                                df_n['Suma'] = df_n.loc[~df_n.index.astype(str).str.startswith('[')].sum(axis=1)
                                df_n.loc['Suma'] = df_n.loc[~df_n.index.astype(str).str.startswith('[')].sum(axis=0)
                            elif is_row_mrs and not is_col_mrs:
                                df_n['Suma'] = mrs_w.sum().values
                                df_n.loc['Suma'] = df_n.sum(axis=0)
                            elif not is_row_mrs and is_col_mrs:
                                df_n['Suma'] = df_n.loc[~df_n.index.astype(str).str.startswith('[')].sum(axis=1)
                                df_n.loc['Suma'] = mrs_w.sum().values.tolist() + [mrs_w.sum().sum()]

                            df_pct = pd.DataFrame(np.nan, index=df_n.index, columns=df_n.columns)
                            if "Kolumnowe" in pct_type:
                                if not is_row_mrs and not is_col_mrs: df_pct = df_n.div(df_n.loc['Suma'].replace(0, np.nan), axis=1) * 100
                                elif is_row_mrs and not is_col_mrs:
                                    base = tmp_df.loc[~tmp_df['C_miss']].groupby(_df_s[col_var].values, observed=False)['w'].sum()
                                    base['Suma'] = base.sum()
                                    df_pct = df_n.div(base.replace(0, np.nan), axis=1) * 100
                                elif not is_row_mrs and is_col_mrs:
                                    df_pct = df_n.div(df_n.loc['Suma'].replace(0, np.nan), axis=1) * 100
                            elif "Wierszowe" in pct_type:
                                if not is_row_mrs and not is_col_mrs: df_pct = df_n.div(df_n['Suma'].replace(0, np.nan), axis=0) * 100
                                elif is_row_mrs and not is_col_mrs: df_pct = df_n.div(df_n['Suma'].replace(0, np.nan), axis=0) * 100
                                elif not is_row_mrs and is_col_mrs:
                                    base = tmp_df.loc[~tmp_df['R_miss']].groupby(_df_s[row_var].values, observed=False)['w'].sum()
                                    for box_name, b_cats in st.session_state.box_sets.get(row_var, {}).items():
                                        base[box_name] = base[base.index.intersection(b_cats)].sum()
                                    base['Suma'] = base.loc[~base.index.astype(str).str.startswith('[')].sum()
                                    df_pct = df_n.div(base.replace(0, np.nan), axis=0) * 100

                            if "Kolumnowe" in pct_type or "Wierszowe" in pct_type:
                                df_pct = df_pct.fillna(0)

                            _sig_run = do_sig_test and (
                                "Kolumnowe" in pct_type
                                or pct_type == "Liczebno\u015bci"
                            )
                            if _sig_run:
                                try:
                                    # Dla trybu "Liczebnosci" oblicz procenty wewnetrznie
                                    # tylko na potrzeby testu (df_pct jest wtedy NaN)
                                    if "Kolumnowe" not in pct_type:
                                        _col_base = df_n.loc['Suma'].replace(0, np.nan)
                                        _pct_for_test = df_n.div(_col_base, axis=1) * 100
                                    else:
                                        _pct_for_test = df_pct
                                    # Przy wagach: efektywna baza (ESS) zamiast wazonego N,
                                    # by test Z nie byl przeszacowany. Deflacja design effect
                                    # (Sigma_w / Sigma_w^2) -- spojnie z testem srednich (ESS).
                                    _ess_bases = None
                                    if use_weights and _w_s is not None:
                                        _wv = np.asarray(w, dtype=float)
                                        _sw = float(_wv.sum()); _sw2 = float((_wv ** 2).sum())
                                        if _sw2 > 0:
                                            _ess_bases = df_n.loc['Suma'] * (_sw / _sw2)
                                    sig_df, col_letters = apply_sig_testing(
                                        _pct_for_test, df_n, bases=_ess_bases)
                                    rename_dict = {c: f"{c} [{col_letters.get(c, '')}]"
                                                   for c in df_n.columns if c != 'Suma'}
                                    df_n.rename(columns=rename_dict, inplace=True)
                                    df_pct.rename(columns=rename_dict, inplace=True)
                                    sig_df.rename(columns=rename_dict, inplace=True)
                                    if "Kolumnowe" in pct_type:
                                        # dtype=object zapobiega blendowi StringDtype w pandas 2.x
                                        df_pct_str = pd.DataFrame(
                                            "", index=df_pct.index,
                                            columns=df_pct.columns, dtype=object)
                                        for c in df_pct.columns:
                                            df_pct_str[c] = df_pct[c].apply(
                                                lambda x: f"{x:.1f}%" if pd.notna(x) else "")
                                        # Bezpieczne laczenie stringow kolumna po kolumnie
                                        for c in df_pct.columns:
                                            _sv = sig_df[c] if c in sig_df.columns else pd.Series("", index=df_pct.index)
                                            df_pct_str[c] = df_pct_str[c].astype(str) + _sv.astype(str)
                                        df_pct = df_pct_str
                                    elif pct_type == "Liczebno\u015bci":
                                        # Tryb tylko liczebnosci: dolacz litery do wartosci N
                                        # (N zaokraglone do liczby calkowitej + np. " E F")
                                        df_n_str = pd.DataFrame(
                                            "", index=df_n.index,
                                            columns=df_n.columns, dtype=object)
                                        for c in df_n.columns:
                                            _scol = sig_df[c] if c in sig_df.columns else pd.Series("", index=df_n.index)
                                            for _idx in df_n.index:
                                                _v = df_n.loc[_idx, c]
                                                if pd.notna(_v):
                                                    try:
                                                        _base = str(int(round(float(_v))))
                                                    except (TypeError, ValueError):
                                                        _base = str(_v)
                                                else:
                                                    _base = ""
                                                _let = str(_scol.loc[_idx]) if _idx in _scol.index else ""
                                                df_n_str.loc[_idx, c] = _base + _let
                                        df_n = df_n_str
                                except Exception as _sig_z_err:
                                    st.warning(f"Nie uda\u0142o si\u0119 oznaczy\u0107 istotno\u015bci Z: {_sig_z_err}")

                            if pct_type == "Liczebno\u015bci":
                                cross_df = df_n.add_suffix(' [N]')
                            elif pct_type in ["Kolumnowe (%)", "Wierszowe (%)"]:
                                cross_df = df_pct.add_suffix(' [%]')
                            else:
                                p_lbl = "[% Kolumnowe]" if "Kolumnowe" in pct_type else "[% Wierszowe]"
                                cross_cols = []
                                for c in df_n.columns:
                                    cross_cols.extend([f"{c} [N]", f"{c} {p_lbl}"])
                                cross_df = pd.DataFrame(index=df_n.index, columns=cross_cols)
                                for c in df_n.columns:
                                    cross_df[f"{c} [N]"] = df_n[c]
                                    cross_df[f"{c} {p_lbl}"] = df_pct[c]

                            cross_df.loc['Braki danych (wykluczone z tabeli)'] = [missing_count] + [np.nan] * (len(cross_df.columns) - 1)
                            st.session_state.results['krzyzowe'][title] = cross_df
                            with st.expander(title):
                                st.dataframe(safe_style(cross_df), use_container_width=True)
                                if title in st.session_state.chi_results:
                                    st.caption(f"\U0001f9ee {st.session_state.chi_results[title]}")
                                if show_charts_cross:
                                    plot_df = cross_df.drop(index=['Suma', 'Braki danych', 'Braki danych (wykluczone z tabeli)'], errors='ignore')
                                    plot_df = plot_df[~plot_df.index.astype(str).str.startswith('[')]
                                    if "Kolumnowe" in pct_type or "Wierszowe" in pct_type:
                                        p_cols = [c for c in plot_df.columns if '[%]' in c or '[% Kolumnowe]' in c or '[% Wierszowe]' in c]
                                    else:
                                        p_cols = [c for c in plot_df.columns if '[N]' in c]
                                    p_cols = [c for c in p_cols if 'Suma' not in c]
                                    if p_cols and not plot_df.empty:
                                        temp_plot = plot_df[p_cols].copy()
                                        for col in temp_plot.columns:
                                            temp_plot[col] = pd.to_numeric(temp_plot[col].apply(_to_float_pct), errors='coerce')
                                        temp_plot.columns = [c.split(' [')[0] for c in temp_plot.columns]
                                        fig = px.bar(temp_plot, barmode='group', orientation='h',
                                                     title=title)
                                        fig.update_layout(yaxis={'categoryorder': 'total ascending'}, height=350)
                                        st.plotly_chart(fig, use_container_width=True, key=f"pc_cross_{title}")
                        except Exception as e:
                            st.error(f"B\u0142\u0105d dla {row_var} \u00d7 {col_var}: {e}")
            st.success("\u2705 Tabele krzy\u017cowe wygenerowane!")

        # \u2500\u2500 Persistent display of stored cross-tab results \u2500\u2500\u2500\u2500
        if st.session_state.results.get('krzyzowe'):
            st.divider()
            st.markdown(f"**Zapisane wyniki ({len(st.session_state.results['krzyzowe'])} tabel):**")
            _cc1, _cc2 = st.columns([5, 1])
            _cc1.markdown(f"**Zapisane tabele krzy\u017cowe ({len(st.session_state.results['krzyzowe'])}):**")
            with _cc2:
                if st.button("\U0001f5d1\ufe0f Usu\u0144 wszystkie", key="del_all_cross",
                             use_container_width=True):
                    st.session_state.results['krzyzowe'] = {}
                    st.session_state.chi_results = {}
                    st.rerun()
            for title, cross_df in list(st.session_state.results['krzyzowe'].items()):
                _xc1, _xc2 = st.columns([6, 1])
                with _xc1:
                    _base_x, _grp_x = _extract_split_from_title(title)
                    _title_display = _base_x + (f" \u2014 \U0001f500 {_grp_x}" if _grp_x else "")
                    _xexp = st.expander(_title_display, expanded=False)
                with _xc2:
                    if st.button("\U0001f5d1\ufe0f", key=f"del_cross_{title}",
                                 help=f"Usu\u0144 {title}"):
                        st.session_state.results['krzyzowe'].pop(title, None)
                        st.session_state.chi_results.pop(title, None)
                        st.rerun()
                with _xexp:
                    _split_badge(_grp_x)
                    st.dataframe(safe_style(cross_df), use_container_width=True)
                    if title in st.session_state.chi_results:
                        st.caption(f"\U0001f9ee {st.session_state.chi_results[title]}")

    with tab_banner:
        with st.expander("Jak tworzy\u0107 i interpretowa\u0107 tabele zbiorcze (Banner)", expanded=False):
            st.markdown("""
### Czym jest tabela zbiorcza (Banner)?

Tabela zbiorcza (banner) to zaawansowana forma tabelacji \u2014 jedno pytanie w wierszach i wiele zmiennych grupuj\u0105cych obok siebie w kolumnach.
Pozwala por\u00f3wna\u0107 rozk\u0142ad jednej zmiennej wzgl\u0119dem wielu grup jednocze\u015bnie \u2014 jak SPSS Custom Tables.

---

### Jak wype\u0142ni\u0107 pola?

- **Wiersze (pytania)** \u2014 zmienna, kt\u00f3rej rozk\u0142ad chcesz zbada\u0107 (np. NPS, opinia, zachowanie). Obs\u0142ugiwane: zmienne kategoryczne i pytania wielokrotnego wyboru (MRS).
- **Banner (zmienne w kolumnach)** \u2014 zmienne grupuj\u0105ce wy\u015bwietlane obok siebie (np. p\u0142e\u0107, wiek, region). Ka\u017cda zmienna tworzy oddzielny blok kolumn.

> Wskaz\u00f3wka: kolumna **Og\u00f3\u0142em** (opcjonalna) pokazuje rozk\u0142ad brze\u017cowy zmiennej wierszowej \u2014 bez podzia\u0142u na grupy.

---

### Spos\u00f3b prezentacji

| Opcja | Kiedy u\u017cywa\u0107 |
|---|---|
| **N + %** | Pe\u0142ny obraz: liczebno\u015bci i procenty razem |
| **Tylko %** | Czytelne por\u00f3wnanie udzia\u0142\u00f3w mi\u0119dzy grupami |
| **Tylko N** | Liczebno\u015bci surowe \u2014 przydatne przy ma\u0142ych pr\u00f3bach |

Procent zawsze **kolumnowy** \u2014 obliczany w ramach ka\u017cdej grupy (kolumny) bloku bannera.

---

### Test Z (95%) per blok

Por\u00f3wnuje pary kolumn w ka\u017cdym bloku bannera niezale\u017cnie. Oznaczenia literowe (A, B, C...) resetuj\u0105 si\u0119 dla ka\u017cdego bloku bannera.

- Kolumna z liter\u0105 "A" oznacza grup\u0119 istotnie r\u00f3\u017cni\u0105c\u0105 si\u0119 od grupy "B" i vice versa
- Przy wagach test uwzgl\u0119dnia efektywn\u0105 wielko\u015b\u0107 pr\u00f3by (ESS)

---

### Uwagi praktyczne

- Zmienna wierszowa powinna by\u0107 kategoryczna (lub MRS) \u2014 unikaj zmiennych ci\u0105g\u0142ych z wieloma unikalnymi warto\u015bciami
- Przy MRS w kolumnie bannera baza % = respondenci (nie odpowiedzi)
- Przy wa\u017ceniu \u2014 N i % s\u0105 wa\u017cone, test Z uwzgl\u0119dnia ESS
            """)
        def _render_banner_html(df):
            import math as _math
            from html import escape as _esc
            _blocks = parse_banner_blocks(df.columns)
            _SUM = {
                'Baza (N)', 'Suma', 'Baza (N) / Suma (%)',
                'Braki danych', 'Braki danych (wykluczone z tabeli)',
                'Og\u00f3\u0142em (Wa\u017cne)'
            }
            def _fmt_val(val, col_s):
                if val is None:
                    return ''
                if isinstance(val, str):
                    return _esc(val)
                try:
                    f = float(val)
                except (TypeError, ValueError):
                    return _esc(str(val))
                if _math.isnan(f):
                    return ''
                col_low = col_s.lower()
                if '%' in col_low:
                    return f'{f:.0f}%'
                if '[n]' in col_low:
                    return str(int(round(f)))
                return f'{f:.1f}'
            style = (
                '<style>'
                '.bn-tbl{border-collapse:collapse;font-size:12px;width:100%;}'
                '.bn-tbl th,.bn-tbl td{border:1px solid #B0C4DE;padding:4px 8px;white-space:nowrap;}'
                '.bn-tbl .blk{background:#1F4E79;color:#fff;font-weight:bold;text-align:center;}'
                '.bn-tbl .cat{background:#D6E4F0;font-weight:bold;text-align:center;color:#1F4E79;}'
                '.bn-tbl .idx{text-align:left;background:#F2F2F2;font-weight:bold;min-width:140px;}'
                '.bn-tbl .dat{text-align:center;}'
                '.bn-tbl tr.even td.dat{background:#FAFAFA;}'
                '.bn-tbl tr.sum td{background:#D6E4F0;font-weight:bold;}'
                '.bn-tbl tr.qhdr td{background:#2E75B6;color:#fff;font-weight:bold;text-align:left;}'
                '</style>'
            )
            hdr1 = '<tr><th class="blk idx"></th>'
            for _blk_lbl, _blk_cols in _blocks:
                hdr1 += '<th class="blk" colspan="' + str(len(_blk_cols)) + '">' + _esc(_blk_lbl) + '</th>'
            hdr1 += '</tr>'
            hdr2 = '<tr><th class="cat idx">Kategorie</th>'
            for _blk_lbl, _blk_cols in _blocks:
                for _cn, _cd in _blk_cols:
                    hdr2 += '<th class="cat">' + _esc(_cd) + '</th>'
            hdr2 += '</tr>'
            body = ''
            _ncol = len(df.columns)
            for i, (idx, row) in enumerate(df.iterrows()):
                _is_qhdr = all((v is None) or (isinstance(v, float) and _math.isnan(v))
                               for v in row)
                if _is_qhdr:
                    # wiersz-naglowek pytania: scalony pasek na calej szerokosci
                    body += ('<tr class="qhdr"><td colspan="' + str(_ncol + 1) + '">'
                             + _esc(str(idx)) + '</td></tr>')
                    continue
                is_sum = str(idx) in _SUM
                tr_cls = 'sum' if is_sum else ('even' if i % 2 == 0 else 'odd')
                body += '<tr class="' + tr_cls + '"><td class="idx">' + _esc(str(idx)) + '</td>'
                for col in df.columns:
                    body += '<td class="dat">' + _fmt_val(row[col], str(col)) + '</td>'
                body += '</tr>'
            return (style
                    + '<div style="overflow-x:auto"><table class="bn-tbl">'
                    + '<thead>' + hdr1 + hdr2 + '</thead><tbody>' + body
                    + '</tbody></table></div>')

        _bn1, _bn2 = st.columns(2)
        with _bn1:
            banner_rows = st.multiselect(
                "Wiersze (pytania):", all_options_no_matrix,
                format_func=lambda x: get_var_display_name(x, var_labels),
                key="banner_rows")
        with _bn2:
            banner_cols = st.multiselect(
                "Banner (zmienne w kolumnach, obok siebie):", all_options_no_matrix,
                format_func=lambda x: get_var_display_name(x, var_labels),
                key="banner_cols")
        _bm1, _bm2, _bm3 = st.columns([2, 1, 1])
        with _bm1:
            banner_measure = st.radio(
                "Pomiar:", ["N + %", "Tylko %", "Tylko N"],
                horizontal=True, key="banner_measure")
        with _bm2:
            banner_sig = st.checkbox("\U0001f520 Test Z (95%) per blok", key="banner_sig")
        with _bm3:
            banner_total = st.checkbox("Do\u0142\u0105cz kolumn\u0119 Og\u00f3\u0142em", value=True,
                                       key="banner_total",
                                       help="Kolumna Og\u00f3\u0142em (rozk\u0142ad brzegowy) na ko\u0144cu tabeli.")

        if _tracked_button("\u25b6\ufe0f Generuj tabele zbiorcze", "analyses",
                           "banner_table", type="primary", key="gen_banner") and banner_rows and banner_cols:
            _w_full = st.session_state.weights if use_weights else np.ones(len(df_raw))
            _meas_sfx = {"N + %": "N+%", "Tylko %": "%", "Tylko N": "N"}.get(banner_measure, banner_measure)
            for _grp_lbl, _df_s, _df_raw_s, _w_s in _iter_split_groups(
                    df, df_raw, var_labels, st.session_state.split_var, weights=_w_full):
                w = _w_s if _w_s is not None else np.ones(len(_df_raw_s))
                try:
                    bdf, bmeta = build_banner_table_multi(
                        banner_rows, banner_cols, _df_s, _df_raw_s, w, var_labels,
                        mrs_sets=st.session_state.mrs_sets,
                        value_orders=st.session_state.value_orders,
                        box_sets=st.session_state.box_sets,
                        measure=banner_measure, do_sig=banner_sig,
                        include_total=banner_total)
                    _bn_base = " + ".join(banner_rows) + " [banner]"
                    _bn_title = (f"{_bn_base} | {_grp_lbl}" if _grp_lbl else _bn_base) + f" [{_meas_sfx}]"
                    st.session_state.results['banner'][_bn_title] = bdf
                    _bn_disp = (", ".join(get_var_display_name(rv, var_labels) for rv in banner_rows)
                                + (f" \u2014 {_grp_lbl}" if _grp_lbl else ""))
                    with st.expander(_bn_disp, expanded=True):
                        _split_badge(_grp_lbl)
                        st.markdown(_render_banner_html(bdf), unsafe_allow_html=True)
                        if bmeta:
                            _leg = "  |  ".join(
                                f"**{_blk}**: " + ", ".join(f"{_cat} ({_lt})" for _cat, _lt in _d.items())
                                for _blk, _d in bmeta.items())
                            st.caption("Legenda liter (test Z 95%, por\u00f3wnania w obr\u0119bie bloku): " + _leg)
                except Exception as _be:
                    _bn_gd = _grp_lbl or 'pe\u0142na baza'
                    st.error(f"B\u0142\u0105d bannera (grupa `{_bn_gd}`): {_be}")
            st.success("\u2705 Tabele zbiorcze wygenerowane!")

        # \u2500\u2500 Persistent display of stored banner results \u2500\u2500\u2500\u2500
        if st.session_state.results.get('banner'):
            st.divider()
            _bnp1, _bnp2 = st.columns([5, 1])
            _bnp1.markdown(f"**Zapisane tabele zbiorcze ({len(st.session_state.results['banner'])}):**")
            with _bnp2:
                if st.button("\U0001f5d1\ufe0f Usu\u0144 wszystkie", key="del_all_banner",
                             use_container_width=True):
                    st.session_state.results['banner'] = {}
                    st.rerun()
            for _bn_title, _bn_df in list(st.session_state.results['banner'].items()):
                _bx1, _bx2 = st.columns([6, 1])
                with _bx1:
                    _bn_b, _bn_g = _extract_split_from_title(_bn_title)
                    _bn_td = _bn_b + (f" \u2014 \U0001f500 {_bn_g}" if _bn_g else "")
                    _bn_exp = st.expander(_bn_td, expanded=False)
                with _bx2:
                    if st.button("\U0001f5d1\ufe0f", key=f"del_banner_{_bn_title}",
                                 help=f"Usu\u0144 {_bn_title}"):
                        st.session_state.results['banner'].pop(_bn_title, None)
                        st.rerun()
                with _bn_exp:
                    _split_badge(_bn_g)
                    st.markdown(_render_banner_html(_bn_df), unsafe_allow_html=True)

    with tab_means:
        col1, col2 = st.columns(2)
        with col1: mean_rows = st.multiselect("Zmienne ci\u0105g\u0142e (wiersze):", numeric_cols, format_func=lambda x: get_var_display_name(x, var_labels))
        with col2: mean_cols_sel = st.multiselect("Metryczka (kolumny):", all_options_no_matrix, format_func=lambda x: get_var_display_name(x, var_labels))
        do_means_sig = st.checkbox("\U0001f520 Oznacz istotne r\u00f3\u017cnice \u015brednich (T-Test 95%)")
        if _tracked_button("\u25b6\ufe0f Generuj tabele \u015brednich", "analyses", "means_table", type="primary") and mean_rows and mean_cols_sel:
            _w_full = st.session_state.weights if use_weights else np.ones(len(df_raw))
            for _grp_lbl, _df_s, _df_raw_s, _w_s in _iter_split_groups(
                    df, df_raw, var_labels, st.session_state.split_var, weights=_w_full):
                w = _w_s if _w_s is not None else np.ones(len(_df_raw_s))
                for row_var in mean_rows:
                    for col_var in mean_cols_sel:
                        try:
                            x = pd.to_numeric(_df_raw_s[row_var], errors='coerce')
                            c_series = _df_s[col_var]
                            cats = c_series.dropna().unique()
                            df_means = pd.DataFrame(index=['Srednia', 'Odchylenie Std.', 'Baza (N)'], columns=cats)
                            df_vars = pd.DataFrame(index=['Srednia'], columns=cats)
                            df_ess = pd.DataFrame(index=['Srednia'], columns=cats)
                            for cat in cats:
                                mask = (c_series == cat).values
                                mean, var, ess = get_weighted_stats(x[mask].values, w[mask])
                                std = np.sqrt(var) if pd.notna(var) and var >= 0 else np.nan
                                df_means.loc['Srednia', cat] = mean
                                df_means.loc['Odchylenie Std.', cat] = std
                                df_means.loc['Baza (N)', cat] = w[mask & ~np.isnan(x.values)].sum()
                                df_vars.loc['Srednia', cat] = var
                                df_ess.loc['Srednia', cat] = ess
                            mean, var, ess = get_weighted_stats(x.values, w)
                            df_means['Og\u00f3\u0142em'] = [mean, np.sqrt(var) if pd.notna(var) and var >= 0 else np.nan, w[~np.isnan(x.values)].sum()]
                            df_vars['Og\u00f3\u0142em'] = var
                            df_ess['Og\u00f3\u0142em'] = ess
                            if do_means_sig:
                                try:
                                    sig_df, col_letters = apply_means_sig_testing(
                                        df_means.loc[['Srednia']],
                                        df_vars.loc[['Srednia']],
                                        df_ess.loc[['Srednia']])
                                    # Przemianuj tylko kolumny grup (nie 'Ogolniem')
                                    rename_dict = {c: f"{c} [{col_letters.get(c, '')}]"
                                                   for c in cats if c in col_letters}
                                    df_means.rename(columns=rename_dict, inplace=True)
                                    sig_df.rename(columns=rename_dict, inplace=True)
                                    # dtype=object: pandas 2.x StringDtype odrzuca float
                                    df_str = pd.DataFrame(
                                        "", index=df_means.index,
                                        columns=df_means.columns, dtype=object)
                                    for c in df_means.columns:
                                        _raw = df_means.loc['Srednia', c]
                                        try:
                                            _fv = float(_raw)
                                            df_str.loc['Srednia', c] = f"{_fv:.2f}" if np.isfinite(_fv) else ""
                                        except (TypeError, ValueError):
                                            df_str.loc['Srednia', c] = str(_raw) if _raw is not None else ""
                                        # Odchylenie i Baza zostawiamy jako float (formatowane pozniej)
                                        df_str.loc['Odchylenie Std.', c] = df_means.loc['Odchylenie Std.', c]
                                        df_str.loc['Baza (N)', c] = df_means.loc['Baza (N)', c]
                                    # Dolacz litery - dla kazdej kol. osobno (bezpieczne)
                                    for c in df_means.columns:
                                        _sig_val = str(sig_df.loc['Srednia', c]) if c in sig_df.columns else ""
                                        df_str.loc['Srednia', c] = str(df_str.loc['Srednia', c]) + _sig_val
                                    df_means = df_str
                                except Exception as _sig_err:
                                    st.warning(f"Nie uda\u0142o si\u0119 oznaczy\u0107 istotno\u015bci: {_sig_err}")
                            title_base = f"{row_var} x {col_var}"
                            title = f"{title_base} | {_grp_lbl}" if _grp_lbl else title_base
                            st.session_state.results['srednie'][title] = df_means
                            with st.expander(title):
                                st.dataframe(_format_means_table(df_means), use_container_width=True)
                        except Exception as e:
                            st.error(f"B\u0142\u0105d: {e}")
            st.success("\u2705 Tabele \u015brednich wygenerowane!")

        # Persistent display of means with delete
        if st.session_state.results.get('srednie'):
            st.divider()
            _mc1, _mc2 = st.columns([5, 1])
            _mc1.markdown(f"**Zapisane tabele \u015brednich ({len(st.session_state.results['srednie'])}):**")
            with _mc2:
                if st.button("\U0001f5d1\ufe0f Usu\u0144 wszystkie", key="del_all_means",
                             use_container_width=True):
                    st.session_state.results['srednie'] = {}
                    st.rerun()
            for title, df_means in list(st.session_state.results['srednie'].items()):
                _base_m, _grp_m = _extract_split_from_title(title)
                _mec1, _mec2 = st.columns([6, 1])
                with _mec1:
                    _mexp = st.expander(_base_m + (f" \u2014 \U0001f500 {_grp_m}" if _grp_m else ""), expanded=False)
                with _mec2:
                    if st.button("\U0001f5d1\ufe0f", key=f"del_means_{title}",
                                 help=f"Usu\u0144 {title}"):
                        st.session_state.results['srednie'].pop(title, None)
                        st.rerun()
                with _mexp:
                    _split_badge(_grp_m)
                    st.dataframe(_format_means_table(df_means), use_container_width=True)

    with tab_desc:
        desc_vars = st.multiselect("Zmienne numeryczne:", numeric_cols,
                                    format_func=lambda x: get_var_display_name(x, var_labels))

        st.markdown("**Wybierz statystyki do prezentacji:**")

        # Group checkboxes into 3 columns for a clean layout
        d1, d2, d3 = st.columns(3)
        with d1:
            st.markdown("*Tendencja centralna*")
            d_mean     = st.checkbox("Srednia",              value=True,  key="ds_mean")
            d_median   = st.checkbox("Mediana (Q2 / 50%)",   value=True,  key="ds_median")
            d_mode     = st.checkbox("Dominanta (moda)",     value=False, key="ds_mode")
            d_trimmed  = st.checkbox("Srednia obci\u0119ta (5%)", value=False, key="ds_trimmed",
                                      help="Srednia po odci\u0119ciu 5% obserwacji z ka\u017cdego ko\u0144ca")
        with d2:
            st.markdown("*Rozrzut*")
            d_std      = st.checkbox("Odchylenie std.",      value=True,  key="ds_std")
            d_var      = st.checkbox("Wariancja",            value=False, key="ds_var")
            d_se       = st.checkbox("B\u0142\u0105d std. sredniej (SE)", value=False, key="ds_se")
            d_range    = st.checkbox("Rozst\u0119p (max-min)", value=False, key="ds_range")
            d_iqr      = st.checkbox("IQR (Q3-Q1)",          value=False, key="ds_iqr")
            d_cv       = st.checkbox("Wsp. zmienno\u015bci (%)", value=False, key="ds_cv",
                                      help="CV = (Odch. std / Srednia) * 100%")
        with d3:
            st.markdown("*Kszta\u0142t rozk\u0142adu*")
            d_skew     = st.checkbox("Sko\u015bno\u015b\u0107",  value=False, key="ds_skew")
            d_kurt     = st.checkbox("Kurtoza",               value=False, key="ds_kurt")
            d_min      = st.checkbox("Min",                   value=True,  key="ds_min")
            d_max      = st.checkbox("Max",                   value=True,  key="ds_max")
            d_q1       = st.checkbox("Q1 (25. percentyl)",    value=False, key="ds_q1")
            d_q3       = st.checkbox("Q3 (75. percentyl)",    value=False, key="ds_q3")
            st.markdown("*Obserwacje*")
            d_n_valid  = st.checkbox("N wa\u017cnych",         value=True,  key="ds_nvalid")
            d_n_miss   = st.checkbox("N brak\u00f3w",          value=True,  key="ds_nmiss")

        if _tracked_button("\u25b6\ufe0f Generuj statystyki opisowe", "analyses", "descriptive_stats", type="primary") and desc_vars:
            try:
                _w_full_desc = (st.session_state.weights if use_weights and st.session_state.weights is not None
                                else None)

                for _grp_lbl, _df_s, _df_raw_s, _w_s in _iter_split_groups(
                        df, df_raw, var_labels, st.session_state.split_var,
                        weights=_w_full_desc):

                    def _wstat(col):
                        """Return (values, weights) arrays after dropping NaN (from slice)."""
                        mask = _df_raw_s[col].notna()
                        x = _df_raw_s.loc[mask, col].values.astype(float)
                        if _w_s is not None:
                            w = pd.Series(_w_s, index=_df_raw_s.index).loc[mask].values.clip(min=0)
                        else:
                            w = np.ones(len(x))
                        return x, w

                    def _wmean(x, w):
                        return float((x * w).sum() / w.sum()) if w.sum() > 0 else np.nan

                    def _wvar(x, w):
                        m = _wmean(x, w)
                        n_eff = w.sum()
                        return float((w * (x - m) ** 2).sum() / max(n_eff - 1, 1))

                    def _wquantile(x, w, q):
                        idx = np.argsort(x)
                        xs, ws = x[idx], w[idx]
                        cumw = np.cumsum(ws)
                        total = cumw[-1]
                        target = q * total
                        i = np.searchsorted(cumw, target)
                        if i == 0:
                            return float(xs[0])
                        if i >= len(xs):
                            return float(xs[-1])
                        frac = (target - cumw[i - 1]) / max(ws[i], 1e-12)
                        return float(xs[i - 1] + frac * (xs[i] - xs[i - 1]))

                    rows = []
                    for c in desc_vars:
                        x, w = _wstat(c)
                        if len(x) == 0:
                            continue
                        row = {'Zmienna': c, 'Etykieta': var_labels.get(c, c)}

                        n_valid = int(_df_raw_s[c].notna().sum())
                        n_miss  = int(_df_raw_s[c].isna().sum())
                        n_w     = float(w.sum())

                        if d_n_valid:  row['N wa\u017cnych (wa\u017cone)'] = int(round(n_w))
                        if d_n_miss:   row['N brak\u00f3w']               = n_miss

                        if d_mean:
                            row['Srednia'] = _wmean(x, w)
                        if d_trimmed:
                            row['Srednia obci\u0119ta (5%)'] = float(stats.trim_mean(x, 0.05))
                        if d_median:
                            row['Mediana'] = _wquantile(x, w, 0.5)
                        if d_mode:
                            mode_res = stats.mode(x, keepdims=True)
                            row['Dominanta'] = float(mode_res.mode[0]) if len(mode_res.mode) > 0 else np.nan
                        if d_std:
                            row['Odch. std.'] = float(np.sqrt(_wvar(x, w)))
                        if d_var:
                            row['Wariancja'] = _wvar(x, w)
                        if d_se:
                            row['B\u0142\u0105d std. (SE)'] = float(np.sqrt(_wvar(x, w) / max(n_w, 1)))
                        if d_min:  row['Min'] = float(x.min())
                        if d_max:  row['Max'] = float(x.max())
                        if d_range: row['Rozst\u0119p'] = float(x.max() - x.min())
                        if d_q1:   row['Q1 (25%)'] = _wquantile(x, w, 0.25)
                        if d_q3:   row['Q3 (75%)'] = _wquantile(x, w, 0.75)
                        if d_iqr:
                            row['IQR'] = _wquantile(x, w, 0.75) - _wquantile(x, w, 0.25)
                        if d_cv:
                            mn = _wmean(x, w)
                            std = np.sqrt(_wvar(x, w))
                            row['CV (%)'] = float(std / mn * 100) if mn != 0 else np.nan
                        if d_skew:
                            mn  = _wmean(x, w)
                            std = np.sqrt(_wvar(x, w))
                            if std > 0:
                                row['Sko\u015bno\u015b\u0107'] = float(
                                    (w * ((x - mn) / std) ** 3).sum() / w.sum())
                            else:
                                row['Sko\u015bno\u015b\u0107'] = np.nan
                        if d_kurt:
                            mn  = _wmean(x, w)
                            std = np.sqrt(_wvar(x, w))
                            if std > 0:
                                row['Kurtoza'] = float(
                                    (w * ((x - mn) / std) ** 4).sum() / w.sum() - 3)
                            else:
                                row['Kurtoza'] = np.nan

                        rows.append(row)

                    if not rows:
                        continue
                    desc_df = pd.DataFrame(rows).set_index('Zmienna')
                    title_desc = f"Statystyki opisowe | {_grp_lbl}" if _grp_lbl else "Statystyki opisowe"
                    st.session_state.results['opisowe'][title_desc] = desc_df
                    num_cols_desc = [c for c in desc_df.columns if desc_df[c].dtype in [float, np.float64]]
                    with st.expander(title_desc, expanded=True):
                        _split_badge(_grp_lbl)
                        st.dataframe(
                            desc_df.style.format({c: '{:.3f}' for c in num_cols_desc}),
                            use_container_width=True
                        )
            except Exception as e:
                st.error(str(e))

    with tab_corr:
        corr_vars = st.multiselect("Zmienne do macierzy korelacji:", numeric_cols,
                                    format_func=lambda x: get_var_display_name(x, var_labels))

        c_opt1, c_opt2, c_opt3 = st.columns(3)
        with c_opt1:
            corr_method = st.selectbox("Metoda:", ["pearson", "spearman", "kendall"],
                                        key="corr_method",
                                        help="Pearson: liniowa; Spearman: rang (odporna); Kendall: rang (ma\u0142e pr\u00f3by)")
        with c_opt2:
            show_heatmap = st.checkbox("\U0001f321\ufe0f Mapa ciep\u0142a", key="corr_heatmap")
        with c_opt3:
            corr_threshold = st.slider(
                "Prog silnej korelacji (|r|):", 0.0, 1.0, 0.5, 0.05,
                key="corr_thresh",
                help="Pary o warto\u015bci bezwzgl\u0119dnej korelacji \u2265 progu zostan\u0105 wyro\u017cnione kolorem i wy\u015bwietlone jako lista."
            )

        if _tracked_button("\u25b6\ufe0f Oblicz korelacje", "analyses", "correlations", type="primary") and len(corr_vars) > 1:
            try:
                _w_full_corr = st.session_state.weights if use_weights else None
                for _grp_lbl, _df_s, _df_raw_s, _w_s in _iter_split_groups(
                        df, df_raw, var_labels, st.session_state.split_var,
                        weights=_w_full_corr):
                    corr_df, n_obs = calculate_correlations(_df_raw_s, corr_vars,
                                                             weights=_w_s,
                                                             method=corr_method)

                    # Raw numeric matrix (for styling and heatmap)
                    if corr_method == 'pearson' and _w_s is not None:
                        import re as _re
                        num_corr = corr_df.map(
                            lambda x: float(_re.sub(r'[*\n].*', '', str(x))) if str(x) not in ('1.000', 'N/A') else (1.0 if str(x) == '1.000' else np.nan)
                        )
                    else:
                        num_corr = _df_raw_s[corr_vars].corr(method=corr_method)

                    corr_df.index   = [var_labels.get(c, c) for c in corr_df.index]
                    corr_df.columns = [var_labels.get(c, c) for c in corr_df.columns]
                    num_corr.index   = [var_labels.get(c, c) for c in num_corr.index]
                    num_corr.columns = [var_labels.get(c, c) for c in num_corr.columns]

                    title_corr = f"Macierz Korelacji | {_grp_lbl}" if _grp_lbl else "Macierz Korelacji"
                    st.session_state.results['korelacje'][title_corr] = corr_df
                    with st.expander(title_corr, expanded=True):
                        _split_badge(_grp_lbl)
                        st.write(f"**Metoda:** {corr_method.title()} | **N wa\u017cnych obserwacji:** {int(round(float(n_obs)))}")

                        _color_corr_cell = _make_color_corr_cell(corr_threshold)
                        styled = corr_df.style.map(_color_corr_cell).format(
                            lambda x: x if isinstance(x, str) else f'{x:.3f}'
                        )
                        st.dataframe(styled, use_container_width=True)

                        st.markdown(
                            "\U0001f7e9 **Silna dodatnia** (r \u2265 0.70) &nbsp;&nbsp;"
                            "\U0001f7e5 **Silna ujemna** (r \u2264 \u22120.70) &nbsp;&nbsp;"
                            "\U0001f7e8 **Umiarkowana dodatnia / ujemna** (|r| \u2265 prog) &nbsp;&nbsp;"
                            "\u25fb Poni\u017cej progu"
                        )

                        strong_pairs = []
                        cols_list = list(num_corr.columns)
                        for i in range(len(cols_list)):
                            for j in range(i + 1, len(cols_list)):
                                r_val = float(num_corr.iloc[i, j])
                                if abs(r_val) >= corr_threshold:
                                    if abs(r_val) >= 0.7:
                                        strength = "silna"
                                    elif abs(r_val) >= 0.5:
                                        strength = "umiarkowana"
                                    else:
                                        strength = "s\u0142aba"
                                    direction = "dodatnia" if r_val > 0 else "ujemna"
                                    strong_pairs.append({
                                        "Zmienna A":    cols_list[i],
                                        "Zmienna B":    cols_list[j],
                                        "r":            round(r_val, 4),
                                        "|r|":          round(abs(r_val), 4),
                                        "Si\u0142a":    strength,
                                        "Kierunek":     direction,
                                    })

                        if strong_pairs:
                            strong_df = pd.DataFrame(strong_pairs).sort_values("|r|", ascending=False)
                            n_strong = len(strong_df)
                            st.markdown(f"**\U0001f517 Silnie skorelowane pary (|r| \u2265 {corr_threshold:.2f}) \u2014 {n_strong} par:**")
                            st.dataframe(
                                strong_df.style.apply(_color_pair_row, axis=1)
                                               .format({'r': '{:.4f}', '|r|': '{:.4f}'}),
                                use_container_width=True, hide_index=True
                            )
                        else:
                            st.info(f"Brak par o |r| \u2265 {corr_threshold:.2f}. Obni\u017c pr\u00f3g aby zobaczy\u0107 wi\u0119cej par.")

                        if show_heatmap:
                            fig = px.imshow(
                                num_corr, color_continuous_scale='RdBu_r',
                                zmin=-1, zmax=1,
                                title=f'Mapa ciep\u0142a korelacji ({corr_method.title()}) \u2014 {_grp_lbl}' if _grp_lbl else f'Mapa ciep\u0142a korelacji ({corr_method.title()})',
                                text_auto='.2f'
                            )
                            fig.update_layout(height=max(400, len(corr_vars) * 40 + 100))
                            st.plotly_chart(fig, use_container_width=True, key=f"pc_corr_heatmap_{corr_method}_{_grp_lbl or 'full'}")

            except Exception as e:
                st.error(str(e))

# -------------------------------------------------------------
# MODU? 4: REGRESJA
# -------------------------------------------------------------
elif menu == "\U0001f4c9 Regresja":
    _require_module_access("regression")
    _require_data()
    module_header("\U0001f4c9", "Regresja", "OLS (liniowa) i logistyczna (binarna/wielomianowa)")
    tab_ols, tab_log = st.tabs(["OLS (liniowa)", "Logistyczna"])

    with tab_ols:
        st.markdown("##### Regresja Liniowa OLS")

        with st.expander("Jak wykona\u0107 i interpretowa\u0107 regresj\u0119 -- kliknij aby rozwin\u0105\u0107", expanded=False):
            st.markdown("""
    **Kiedy u\u017cywa\u0107?** Gdy chcesz sprawdzi\u0107, kt\u00f3re zmienne (predyktory X) przewiduj\u0105 warto\u015b\u0107 innej zmiennej (Y), oraz jak silny jest ten zwi\u0105zek.

    **Jak wykona\u0107:**
    1. Wybierz **zmienn\u0105 zale\u017cn\u0105 (Y)** -- musi by\u0107 numeryczna (np. wynik testu, poziom satysfakcji).
    2. Wybierz **predyktory** w jednym lub kilku blokach.
    3. Bloki hierarchiczne (jak w SPSS): ka\u017cdy kolejny blok dodaje nowe zmienne i pokazuje przyrost R\u00b2.
    4. Kliknij **Uruchom regresj\u0119**.

    **Jak interpretowa\u0107 wyniki:**

    | Wska\u017anik | Interpretacja |
    |---|---|
    | **R\u00b2** | % wariancji Y wyja\u015bniany przez model. Wy\u017cszy = lepszy. |
    | **Skorygowane R\u00b2** | R\u00b2 poprawiony o liczb\u0119 predyktor\u00f3w -- por\u00f3wnuj mi\u0119dzy modelami. |
    | **\u0394R\u00b2** | O ile wzros\u0142o R\u00b2 po dodaniu nowego bloku predyktor\u00f3w. |
    | **F modelu / p** | Czy ca\u0142y model jest istotny statystycznie (p < 0.05 = TAK). |
    | **F zmiany / p** | Czy nowy blok predyktor\u00f3w istotnie poprawi\u0142 model. |
    | **B** | Niestandaryzowany wsp\u00f3\u0142czynnik: zmiana Y przy wzro\u015bcie X o 1 jednostk\u0119. |
    | **Beta (std.)** | Standaryzowany -- por\u00f3wnuje si\u0142\u0119 r\u00f3\u017cnych predyktor\u00f3w (niezale\u017cnie od skali). |
    | **VIF** | Wska\u017anik wsp\u00f3\u0142liniowo\u015bci: < 5 OK \u00b7 5-10 uwaga \u00b7 > 10 problem. |
    | **Tolerancja** | 1/VIF. Im ni\u017csza, tym wi\u0119kszy problem ze wsp\u00f3\u0142liniowo\u015bci\u0105. |

    **Wykresy diagnostyczne:**
    - *Reszty vs Dopasowane* -- punkty rozmieszczone losowo wok\u00f3\u0142 0 = OK (spe\u0142nione za\u0142o\u017cenie homoskedastyczno\u015bci).
    - *Q-Q plot* -- punkty blisko czerwonej linii = reszty maj\u0105 rozk\u0142ad normalny = OK.
            """)

        st.info("Dodawaj predyktory **blokami** (hierarchicznie). Ka\u017cdy blok poka\u017ce zmian\u0119 R\u00b2 i test F-zmiany.")

        dep_var = st.selectbox("Zmienna zale\u017cna (Y):", numeric_cols, format_func=lambda x: get_var_display_name(x, var_labels))

        if st.button("\u2795 Dodaj blok predyktor\u00f3w"):
            st.session_state.reg_blocks.append([])
            st.rerun()

        blocks_to_delete = None
        for b_idx in range(len(st.session_state.reg_blocks)):
            c1, c2 = st.columns([6, 1])
            with c1:
                chosen = st.multiselect(
                    f"Blok {b_idx + 1} -- predyktory:",
                    [c for c in numeric_cols if c != dep_var],
                    default=st.session_state.reg_blocks[b_idx],
                    format_func=lambda x: get_var_display_name(x, var_labels),
                    key=f"reg_block_{b_idx}"
                )
                st.session_state.reg_blocks[b_idx] = chosen
            with c2:
                if b_idx > 0 and st.button("\U0001f5d1\ufe0f", key=f"del_block_{b_idx}", help="Usu\u0144 blok"):
                    blocks_to_delete = b_idx

        if blocks_to_delete is not None:
            st.session_state.reg_blocks.pop(blocks_to_delete)
            st.rerun()

        st.divider()
        if _tracked_button("\u25b6\ufe0f Uruchom regresj\u0119", "regression", "run_ols", type="primary"):
            valid_blocks = [b for b in st.session_state.reg_blocks if b]
            if not valid_blocks:
                st.error("Dodaj co najmniej jeden predyktor.")
            else:
                with st.spinner("Obliczanie regresji OLS..."):
                    _w_full = st.session_state.weights if use_weights else None
                    for _grp_lbl, _df_s, _df_raw_s, _w_s in _iter_split_groups(
                            df, df_raw, var_labels, st.session_state.split_var, weights=_w_full):
                        _new_reg = run_regression_block(
                            _df_raw_s, dep_var, valid_blocks, weights=_w_s
                        )
                        for _r in _new_reg:
                            _r['group_label'] = _grp_lbl
                            _merge_result(st.session_state.regression_results, _r,
                                key_fn=lambda r: (r.get('dep_var',''), r.get('block_idx', 0),
                                                  r.get('group_label', '')))

        if st.session_state.regression_results:
            _rgc1, _rgc2 = st.columns([5, 1])
            _rgc1.markdown(f"**Zapisane modele regresji ({len(st.session_state.regression_results)}):**")
            with _rgc2:
                if st.button("\U0001f5d1\ufe0f Usu\u0144 wszystkie", key="del_all_reg",
                             use_container_width=True):
                    st.session_state.regression_results = []
                    st.rerun()

        for _ri, res in enumerate(list(st.session_state.regression_results)):
            if 'error' in res:
                st.error(res['error'])
                continue
            dep_label = var_labels.get(res['dep_var'], res['dep_var'])
            _rc1, _rc2 = st.columns([6, 1])
            with _rc1:
                _rexp = st.expander(
                    f"\U0001f4ca Blok {res['Blok']} -- [{res['dep_var']}] {dep_label}"
                    + (f" | {res.get('group_label','')}" if res.get('group_label') else ""),
                    expanded=True
                )
            with _rc2:
                if st.button("\U0001f5d1\ufe0f", key=f"del_reg_{_ri}",
                             help=f"Usu\u0144 model {res['dep_var']}"):
                    st.session_state.regression_results.pop(_ri)
                    st.rerun()
            with _rexp:
                _split_badge(res.get('group_label', ''))
                m1, m2, m3, m4 = st.columns(4)
                m1.metric("N", f"{int(round(float(res['N']))):,}")
                m2.metric("R\u00b2", f"{res['R2']:.4f}")
                m3.metric("Skorygowane R\u00b2", f"{res['Skor_R2']:.4f}")
                m4.metric("\u0394R\u00b2", f"{res['Delta_R2']:.4f}")
                f1, f2, f3, f4 = st.columns(4)
                f_pval = res['p (F modelu)']
                f1.metric("F modelu", f"{res['F modelu']:.3f}")
                f2.metric("p (F modelu)", f"('OK' if f_pval < 0.05 else 'NS') {f_pval:.4f}")
                fc, fcp = res['F zmiany'], res['p (F zmiany)']
                try:
                    if not np.isnan(fc):
                        f3.metric("F zmiany", f"{fc:.3f}")
                        f4.metric("p (F zmiany)", f"('OK' if fcp < 0.05 else 'NS') {fcp:.4f}")
                except: pass

                st.markdown("**Wsp\u00f3\u0142czynniki regresji:**")
                coef_df = res['coef_df'].copy()



                styled = coef_df.style \
                    .format({'B': '{:.4f}', 'B\u0142\u0105d std. B': '{:.4f}', 'Beta (std.)': '{:.4f}',
                             't': '{:.3f}', 'p-value': '{:.4f}', 'VIF': '{:.2f}', 'Tolerancja': '{:.4f}'}) \
                    .map(_style_p, subset=['p-value']) \
                    .map(_style_vif, subset=['VIF'])
                st.dataframe(styled, use_container_width=True)
                st.caption("\U0001f7e2 p < 0.05 -- istotne statystycznie \u00b7 VIF > 10 = wsp\u00f3\u0142liniowo\u015b\u0107 (czerwony) \u00b7 VIF 5-10 = uwaga (pomara\u0144czowy)")

                # Diagnostic plots
                all_pred = res['Wszystkie predyktory']
                df_diag = df_raw[[res['dep_var']] + all_pred].dropna()
                if len(df_diag) > 5:
                    X_d = sm.add_constant(df_diag[all_pred])
                    mod_d = sm.OLS(df_diag[res['dep_var']], X_d).fit()
                    ch1, ch2 = st.columns(2)
                    with ch1:
                        fig_r = px.scatter(x=mod_d.fittedvalues, y=mod_d.resid,
                                           labels={'x': 'Wartosci dopasowane', 'y': 'Reszty'},
                                           title='Reszty vs Warto\u015bci dopasowane', color_discrete_sequence=['#2E75B6'])
                        fig_r.add_hline(y=0, line_dash='dash', line_color='red')
                        _chart_key_base = f"{res.get('dep_var','ols')}_{res.get('Blok','')}_{res.get('group_label','')}_{_ri}"
                        st.plotly_chart(fig_r, use_container_width=True, key=f"pc_ols_{_chart_key_base}_resid")
                    with ch2:
                        (osm, osr), (slope, intercept, _r) = stats.probplot(mod_d.resid)
                        fig_qq = go.Figure()
                        fig_qq.add_trace(go.Scatter(x=list(osm), y=list(osr), mode='markers', name='Reszty', marker=dict(color='#2E75B6')))
                        fig_qq.add_trace(go.Scatter(x=[min(osm), max(osm)], y=[slope * min(osm) + intercept, slope * max(osm) + intercept],
                                                    mode='lines', name='Linia ref.', line=dict(color='red', dash='dash')))
                        fig_qq.update_layout(title='Wykres Q-Q (normalno\u015b\u0107 reszt)',
                                             xaxis_title='Kwantyle teoretyczne', yaxis_title='Kwantyle pr\u00f3bkowe')
                        st.plotly_chart(fig_qq, use_container_width=True, key=f"pc_ols_{_chart_key_base}_qq")

    # \u2500\u2500 TAB: Regresja Logistyczna \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
    with tab_log:
        st.markdown("##### Regresja Logistyczna")

        with st.expander("Jak wykona\u0107 i interpretowa\u0107", expanded=False):
            st.markdown("""
**Czym jest regresja logistyczna?**
Modeluje prawdopodobie\u0144stwo przynale\u017cno\u015bci do kategorii.
- **Binarna** \u2014 zmienna zale\u017cna ma 2 kategorie (np. zakup: tak/nie, 0/1)
- **Wielomianowa (Multinomial)** \u2014 zmienna zale\u017cna ma 3+ kategorie

**Kluczowe wska\u017aniki:**

| Wska\u017anik | Interpretacja |
|---|---|
| **Iloraz szans (OR)** | OR > 1: zmienna zwi\u0119ksza szans\u0119; OR < 1: zmniejsza |
| **95% CI** | Przedzia\u0142 ufno\u015bci dla OR. Je\u015bli nie zawiera 1: istotne |
| **p-value** | < 0.05: zmienna istotnie wp\u0142ywa na wynik |
| **Pseudo R\u00b2 (McFadden)** | 0.2-0.4: dobre dopasowanie modelu |
| **AIC** | Kryterium informacyjne. Mniejszy = lepszy model |
            """)

        log_type = st.radio("Typ regresji:", ["Binarna (Logit)", "Wielomianowa (MNLogit)"],
                             horizontal=True, key="log_type")
        st.divider()
        col_lg1, col_lg2 = st.columns(2)
        with col_lg1:
            log_dep = st.selectbox("Zmienna zale\u017cna:", visible_columns,
                                    format_func=lambda x: get_var_display_name(x, var_labels), key="log_dep")
            dep_vals = sorted(df[log_dep].dropna().unique())
            st.caption(f"Unikalne warto\u015bci ({len(dep_vals)}): {', '.join(str(v) for v in dep_vals[:8])}"
                       + (" ..." if len(dep_vals) > 8 else ""))
            if log_type == "Binarna (Logit)" and len(dep_vals) != 2:
                st.warning("Regresja binarna wymaga dok\u0142adnie 2 unikalnych warto\u015bci.")
        with col_lg2:
            log_indep = st.multiselect("Predyktory:",
                                        [c for c in visible_columns if c != log_dep],
                                        format_func=lambda x: get_var_display_name(x, var_labels), key="log_indep")
            log_dummy = st.checkbox("Automatycznie zakoduj zmienne kategoryczne (dummy coding)",
                                     value=True, key="log_dummy")

        if _tracked_button("\u25b6\ufe0f Uruchom regresj\u0119 logistyczn\u0105", "regression", "run_logistic", type="primary", key="log_run"):
            if not log_indep:
                st.error("Wybierz co najmniej jeden predyktor.")
            else:
                with st.spinner("Obliczanie..."):
                    _w_full = st.session_state.weights if use_weights else None
                    for _grp_lbl, _df_s, _df_raw_s, _w_s in _iter_split_groups(
                            df, df_raw, var_labels, st.session_state.split_var, weights=_w_full):
                        try:
                            df_lg = _df_s[[log_dep] + log_indep].dropna().copy()
                            n_obs = len(df_lg)
                            _w_log = None
                            if _w_s is not None:
                                _w_log = pd.Series(_w_s, index=_df_s.index
                                                   ).reindex(df_lg.index).fillna(1).clip(lower=0).values
                                n_obs = round(float(_w_log.sum()), 1)
                            if len(df_lg) < 20:
                                _grp_disp = _grp_lbl or 'pe\u0142na baza'
                                st.error(f"Za ma\u0142o obserwacji ({len(df_lg)}) dla grupy `{_grp_disp}`.")
                                continue
                            y_series = df_lg[log_dep]
                            if log_type == "Binarna (Logit)":
                                uniq = sorted(y_series.dropna().unique())
                                if len(uniq) != 2:
                                    st.error(f"Zmienna zale\u017cna musi mie\u0107 dok\u0142adnie 2 warto\u015bci (grupa `{_grp_lbl}`).")
                                    continue
                                y = (y_series == uniq[1]).astype(np.float64)
                                dep_ref = str(uniq[0]); dep_pos = str(uniq[1])
                            else:
                                y = y_series.astype(str)
                                dep_ref = str(sorted(y_series.unique())[0])
                                dep_pos = "wszystkie"

                            cat_cols = [c for c in log_indep
                                        if df_lg[c].dtype == object
                                        or (df_lg[c].nunique() <= 10 and log_dummy)]
                            num_cols_lg = [c for c in log_indep if c not in cat_cols]

                            parts = []
                            if num_cols_lg:
                                num_part = df_lg[num_cols_lg].apply(
                                    pd.to_numeric, errors='coerce').fillna(0).astype(np.float64)
                                parts.append(num_part)
                            if cat_cols and log_dummy:
                                for cc in cat_cols:
                                    dummies = pd.get_dummies(
                                        df_lg[cc].astype(str),
                                        prefix=cc, drop_first=True
                                    ).astype(np.float64)
                                    parts.append(dummies)
                            elif cat_cols:
                                for cc in cat_cols:
                                    df_lg[cc] = pd.Categorical(
                                        df_lg[cc].astype(str)).codes.astype(np.float64)
                                parts.append(df_lg[cat_cols].astype(np.float64))

                            if parts:
                                X = pd.concat(parts, axis=1)
                            else:
                                st.error("Brak predyktor\u00f3w do modelu.")
                                continue

                            X = X.astype(np.float64)
                            X_const = sm.add_constant(X, has_constant='add').astype(np.float64)
                            y_fit = y.astype(np.float64)

                            if log_type == "Binarna (Logit)":
                                # GLM(Binomial) zamiast sm.Logit: poprawnie wazy przez
                                # freq_weights (sm.Logit cicho ignoruje wagi). Wagi sumuja
                                # sie do N, wiec nobs/AIC/istotnosc pozostaja poprawne.
                                _glm_kw = {'freq_weights': _w_log} if _w_log is not None else {}
                                model = sm.GLM(y_fit, X_const,
                                               family=sm.families.Binomial(),
                                               **_glm_kw).fit(maxiter=200)
                                # GLM nie ma .prsquared/.llr_pvalue jak Logit -> licz recznie
                                try:
                                    _log_pr2 = float(model.pseudo_rsquared(kind='mcf'))
                                except Exception:
                                    _log_pr2 = np.nan
                                try:
                                    _log_lr = 2.0 * (model.llf - model.llnull)
                                    _log_llr_p = float(stats.chi2.sf(_log_lr, max(int(model.df_model), 1)))
                                except Exception:
                                    _log_llr_p = np.nan
                                params = model.params
                                pvals  = model.pvalues
                                conf   = model.conf_int()
                                or_vals = np.exp(params)
                                or_lo   = np.exp(conf.iloc[:, 0])
                                or_hi   = np.exp(conf.iloc[:, 1])
                                coef_df = pd.DataFrame({
                                    'Zmienna':      params.index.tolist(),
                                    'Wspolczynnik': params.values.round(4),
                                    'Iloraz szans': or_vals.values.round(4),
                                    'CI 95% (dol)': or_lo.values.round(4),
                                    'CI 95% (gor)': or_hi.values.round(4),
                                    'p-value':      pvals.values.round(4),
                                    'Istotny':      ['Tak' if p < 0.05 else 'Nie' for p in pvals.values]
                                })
                                result_entry = {'type': 'Binarna', 'dep_var': log_dep, 'dep_ref': dep_ref,
                                                'dep_pos': dep_pos, 'indep_vars': log_indep, 'n_obs': n_obs,
                                                'pseudo_r2': _log_pr2, 'llr_p': _log_llr_p,
                                                'aic': model.aic, 'bic': model.bic_llf, 'log_lik': model.llf,
                                                'coef_df': coef_df, 'model': model, 'error': None,
                                                'group_label': _grp_lbl}
                            else:
                                # MNLogit nie wspiera freq_weights -> wazenie przez
                                # replikacje obserwacji (skalowanie wzgledem najmniejszej
                                # wagi, by wazenie realnie wplynelo na wspolczynniki).
                                _mn_y, _mn_X = y_fit, X_const
                                _mn_wnote = None
                                if _w_log is not None:
                                    _wv = np.asarray(_w_log, dtype=float)
                                    _posm = _wv[_wv > 0]
                                    _minw = float(_posm.min()) if _posm.size else 0.0
                                    if _minw > 0:
                                        _reps = np.round(_wv / _minw).astype(int)
                                        _reps[_reps < 0] = 0
                                        _tot = int(_reps.sum())
                                        if 10 <= _tot <= 500000:
                                            _ri = np.repeat(np.arange(len(_wv)), _reps)
                                            _mn_y = y_fit.iloc[_ri].reset_index(drop=True)
                                            _mn_X = X_const.iloc[_ri].reset_index(drop=True)
                                            _mn_wnote = ("\u2696\ufe0f Wazenie wielomianowej przez replikacje: "
                                                         "wsp\u00f3\u0142czynniki s\u0105 wa\u017cone, ale b\u0142\u0119dy std., "
                                                         "p-warto\u015bci, AIC i pseudo-R\u00b2 s\u0105 przybli\u017cone "
                                                         "(oparte na pr\u00f3bie replikowanej).")
                                        else:
                                            _grp_disp = _grp_lbl or 'pe\u0142na baza'
                                            st.warning(f"Grupa `{_grp_disp}`: wagi zbyt zr\u00f3\u017cnicowane "
                                                       "do replikacji \u2014 model wielomianowy policzony bez wag.")
                                model = sm.MNLogit(_mn_y, _mn_X).fit(disp=False, maxiter=200)
                                result_entry = {'type': 'Wielomianowa', 'dep_var': log_dep, 'dep_ref': dep_ref,
                                                'indep_vars': log_indep, 'n_obs': n_obs,
                                                'pseudo_r2': model.prsquared, 'llr_p': model.llr_pvalue,
                                                'aic': model.aic, 'bic': model.bic, 'log_lik': model.llf,
                                                'coef_df': None, 'model': model, 'error': None,
                                                'weight_note': _mn_wnote,
                                                'group_label': _grp_lbl}
                            _merge_result(st.session_state.logistic_results, result_entry,
                                key_fn=lambda r: (r.get('dep_var',''), r.get('type',''),
                                                  r.get('group_label','')))
                        except Exception as _lg_err:
                            _grp_disp = _grp_lbl or 'pe\u0142na baza'
                            st.error(f"B\u0142\u0105d dla grupy `{_grp_disp}`: {_lg_err}")
                    st.success("\u2705 Regresja logistyczna obliczona!")

        if st.session_state.logistic_results:
            _lgc1, _lgc2 = st.columns([5, 1])
            _lgc1.markdown(f"**Zapisane modele logistyczne ({len(st.session_state.logistic_results)}):**")
            with _lgc2:
                if st.button("\U0001f5d1\ufe0f Usu\u0144 wszystkie", key="del_all_log",
                             use_container_width=True):
                    st.session_state.logistic_results = []
                    st.rerun()

        for _lgi, res_lg in enumerate(list(st.session_state.logistic_results)):
            if res_lg.get('error'): st.error(res_lg['error']); continue
            dep_lbl = var_labels.get(res_lg['dep_var'], res_lg['dep_var'])
            _lgec1, _lgec2 = st.columns([6, 1])
            with _lgec1:
                _lgexp = st.expander(
                    f"\U0001f4c9 {res_lg['type']}: [{res_lg['dep_var']}] {dep_lbl}"
                    + (f" | {res_lg.get('group_label','')}" if res_lg.get('group_label') else ""),
                    expanded=True
                )
            with _lgec2:
                if st.button("\U0001f5d1\ufe0f", key=f"del_log_{_lgi}",
                             help=f"Usu\u0144 {res_lg['dep_var']}"):
                    st.session_state.logistic_results.pop(_lgi)
                    st.rerun()
            with _lgexp:
                _split_badge(res_lg.get('group_label', ''))
                m1,m2,m3,m4,m5 = st.columns(5)
                m1.metric("N", f"{int(round(float(res_lg['n_obs']))):,}")
                m2.metric("Pseudo R\u00b2", f"{res_lg['pseudo_r2']:.4f}")
                m3.metric("p (LLR)", f"{res_lg['llr_p']:.4f}")
                m4.metric("AIC", f"{res_lg['aic']:.1f}")
                m5.metric("BIC", f"{res_lg['bic']:.1f}")
                if res_lg['type'] == 'Binarna' and res_lg['coef_df'] is not None:
                    st.markdown(f"**Kategoria ref.:** `{res_lg['dep_ref']}` | **Modelowana:** `{res_lg['dep_pos']}`")
                    cdf = res_lg['coef_df'].copy()
                    styled = cdf.style.apply(_color_sig, axis=1).format(
                        {'Wspolczynnik': '{:.4f}', 'Iloraz szans': '{:.4f}',
                         'CI 95% (dol)': '{:.4f}', 'CI 95% (gor)': '{:.4f}', 'p-value': '{:.4f}'})
                    st.dataframe(styled, use_container_width=True, hide_index=True)
                    plot_df = cdf[cdf['Zmienna'] != 'const'].copy()
                    if not plot_df.empty:
                        fig_or = go.Figure()
                        colors = ['#C00000' if p >= 0.05 else '#2E75B6' for p in plot_df['p-value']]
                        fig_or.add_trace(go.Scatter(x=plot_df['Iloraz szans'], y=plot_df['Zmienna'],
                                                    mode='markers', marker=dict(size=10, color=colors),
                                                    error_x=dict(type='data', symmetric=False,
                                                                 array=(plot_df['CI 95% (gor)'] - plot_df['Iloraz szans']).tolist(),
                                                                 arrayminus=(plot_df['Iloraz szans'] - plot_df['CI 95% (dol)']).tolist())))
                        fig_or.add_vline(x=1, line_dash='dash', line_color='gray')
                        fig_or.update_layout(title='Ilorazy szans (OR) z 95% CI', xaxis_title='OR',
                                              height=max(300, len(plot_df)*35+100), showlegend=False)
                        st.plotly_chart(fig_or, use_container_width=True,
                                        key=f"pc_log_{res_lg.get('dep_var','log')}_{res_lg.get('group_label','')}_{_lgi}_or")
                else:
                    if res_lg.get('weight_note'):
                        st.caption(res_lg['weight_note'])
                    st.text(res_lg['model'].summary().as_text())

        if st.session_state.logistic_results:
            if st.button("\U0001f5d1\ufe0f Wyczy\u015b\u0107 wyniki regresji logistycznej", type="secondary", key="log_clear_tab"):
                st.session_state.logistic_results = []
                st.rerun()

# -------------------------------------------------------------
# MODUL 5: ANOVA
# -------------------------------------------------------------
# =============================================================
# MODUL: TESTY NORMALNOSCI
# =============================================================
elif menu == "\U0001f4d0 Testy Normalno\u015bci":
    _require_module_access("normality")
    _require_data()
    module_header("\U0001f4d0", "Testy Normalno\u015bci",
                  "Sprawdzanie za\u0142o\u017cenia normalno\u015bci rozk\u0142adu \u2014 wymagane przed ANOVA, regresj\u0105 i innymi testami parametrycznymi")

    with st.expander("Jak przeprowadzi\u0107 i interpretowa\u0107 testy normalno\u015bci", expanded=False):
        st.markdown("""
### Czym jest test normalno\u015bci?

Test normalno\u015bci sprawdza, czy rozk\u0142ad zmiennej jest zbli\u017cony do rozk\u0142adu normalnego (Gaussa). Jest to
za\u0142o\u017cenie wielu test\u00f3w parametrycznych: **t-testu, ANOVA, regresji liniowej, korelacji Pearsona**.

---

### Dost\u0119pne testy

| Test | Kiedy stosowa\u0107 | Hipoteza zerowa (H\u2080) |
|---|---|---|
| **Shapiro-Wilk** | Ma\u0142e i \u015brednie pr\u00f3by (N \u2264 2000) \u2014 **najsilniejszy test** | Rozk\u0142ad jest normalny |
| **Kolmogorov-Smirnov** | Du\u017ce pr\u00f3by (N > 2000), por\u00f3wnanie z rozk\u0142adem teoretycznym | Rozk\u0142ad jest normalny |
| **Lilliefors** | Wariant K-S gdy \u015brednia i odch. std. s\u0105 szacowane z danych | Rozk\u0142ad jest normalny |
| **D\u2019Agostino-Pearson** | Opiera si\u0119 na sko\u015bno\u015bci i kurtozie | Rozk\u0142ad jest normalny |

---

### Jak interpretowa\u0107 wyniki?

**p-value:**
- **p > 0.05** \u2192 Brak podstaw do odrzucenia H\u2080 \u2192 dane **mog\u0105 pochodzi\u0107** z rozk\u0142adu normalnego \u2705
- **p \u2264 0.05** \u2192 Odrzucamy H\u2080 \u2192 dane **nie pochodz\u0105** z rozk\u0142adu normalnego \u274c

> \u26a0\ufe0f **Uwaga:** Dla du\u017cych pr\u00f3b (N > 200) nawet ma\u0142e, praktycznie nieistotne odchylenia od normalno\u015bci
> daj\u0105 p < 0.05. W takich przypadkach **wa\u017cniejsza jest ocena wizualna** (Q-Q plot, histogram)
> i miary sko\u015bno\u015bci / kurtozy. Regu\u0142a: |sko\u015bno\u015b\u0107| < 2 i |kurtoza| < 7 to akceptowalna normalno\u015b\u0107.

---

### Ocena wizualna

- **Histogram** \u2014 kszta\u0142t dzwonu = dobry znak; skos lub gruby ogon = odchylenie
- **Wykres Q-Q** \u2014 punkty blisko prostej = normalno\u015b\u0107; wygi\u0119cie = skos lub kurtoza

---

### Praktyczne zalecenia (jak w SPSS)

| Wielko\u015b\u0107 pr\u00f3by | Zalecany test |
|---|---|
| N < 50 | Shapiro-Wilk |
| 50 \u2264 N \u2264 2000 | Shapiro-Wilk + wizualna ocena Q-Q |
| N > 2000 | Lilliefors lub D\u2019Agostino + wizualna ocena Q-Q |
        """)

    st.divider()

    norm_vars = st.multiselect(
        "Wybierz zmienne numeryczne do testowania:",
        numeric_cols,
        format_func=lambda x: get_var_display_name(x, var_labels),
        key="norm_vars"
    )

    ncol1, ncol2 = st.columns(2)
    with ncol1:
        norm_tests = st.multiselect(
            "Wybierz testy:",
            ["Shapiro-Wilk", "Kolmogorov-Smirnov (Lilliefors)", "D\u2019Agostino-Pearson"],
            default=["Shapiro-Wilk", "D\u2019Agostino-Pearson"],
            key="norm_tests"
        )
        norm_alpha = st.select_slider(
            "Poziom istotno\u015bci (\u03b1):",
            options=[0.01, 0.05, 0.10],
            value=0.05,
            key="norm_alpha"
        )
    with ncol2:
        norm_show_qq   = st.checkbox("Wykres Q-Q", value=True, key="norm_qq")
        norm_show_hist = st.checkbox("Histogram z krzywa normaln\u0105", value=True, key="norm_hist")
        norm_show_desc = st.checkbox("Statystyki opisowe (sko\u015bno\u015b\u0107, kurtoza)", value=True, key="norm_desc")

    if _tracked_button("\u25b6\ufe0f Przeprowad\u017a testy normalno\u015bci", "normality", "run_normality", type="primary",
                       key="norm_run") and norm_vars and norm_tests:
        _w_full_norm = st.session_state.weights if use_weights else None
        for _grp_lbl, _df_s, _df_raw_s, _w_s in _iter_split_groups(
                df, df_raw, var_labels, st.session_state.split_var, weights=_w_full_norm):
            for var in norm_vars:
                if _w_s is not None:
                    mask = _df_raw_s[var].notna()
                    x_raw = _df_raw_s.loc[mask, var].values.astype(float)
                    w_raw = pd.Series(_w_s, index=_df_raw_s.index).loc[mask].values
                    # Weighted sample via repetition (SPSS approach for normality tests)
                    if len(w_raw) > 0 and w_raw.min() > 0:
                        reps = np.round(w_raw / w_raw.min()).astype(int)
                        x = np.repeat(x_raw, reps)
                    else:
                        x = x_raw
                else:
                    x = _df_raw_s[var].dropna().values.astype(float)

                n = len(x)
                n_raw = len(_df_raw_s[var].dropna())
                if n < 3:
                    _grp_disp = _grp_lbl or 'pe\u0142na baza'
                    st.warning(f"Za ma\u0142o obserwacji dla `{var}` w grupie `{_grp_disp}` (N={n}).")
                    continue

                test_rows = []
                warn_msg = None
                if "Shapiro-Wilk" in norm_tests:
                    if n > 5000:
                        warn_msg = "Shapiro-Wilk: pr\u00f3ba zbyt du\u017ca (N > 5000). U\u017cyto losowej pr\u00f3bki 5000."
                        x_sw = np.random.choice(x, 5000, replace=False)
                    else:
                        x_sw = x
                    sw_stat, sw_p = stats.shapiro(x_sw)
                    test_rows.append({
                        "Test": "Shapiro-Wilk",
                        "Statystyka": round(float(sw_stat), 4),
                        "p-value": round(float(sw_p), 4),
                        "Wynik": "\u2705 Normalny" if sw_p > norm_alpha else "\u274c Nienormalny"
                    })

                if "Kolmogorov-Smirnov (Lilliefors)" in norm_tests:
                    from statsmodels.stats.diagnostic import kstest_normal
                    lf_stat, lf_p = kstest_normal(x, dist='norm')
                    test_rows.append({
                        "Test": "Lilliefors (K-S)",
                        "Statystyka": round(float(lf_stat), 4),
                        "p-value": round(float(lf_p), 4),
                        "Wynik": "\u2705 Normalny" if lf_p > norm_alpha else "\u274c Nienormalny"
                    })

                if "D\u2019Agostino-Pearson" in norm_tests:
                    dag_stat, dag_p = stats.normaltest(x)
                    test_rows.append({
                        "Test": "D\u2019Agostino-Pearson",
                        "Statystyka": round(float(dag_stat), 4),
                        "p-value": round(float(dag_p), 4),
                        "Wynik": "\u2705 Normalny" if dag_p > norm_alpha else "\u274c Nienormalny"
                    })

                result_key = f"{var} | {_grp_lbl}" if _grp_lbl else var
                st.session_state.normality_results[result_key] = {
                    'var': var,
                    'group_label': _grp_lbl,
                    'x': x.tolist(),
                    'n': n,
                    'n_raw': n_raw,
                    'test_df': pd.DataFrame(test_rows),
                    'warn_msg': warn_msg,
                    'show_qq': bool(norm_show_qq),
                    'show_hist': bool(norm_show_hist),
                    'show_desc': bool(norm_show_desc),
                    'alpha': float(norm_alpha),
                }

    # \u2500\u2500 Display all stored results with delete UI \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
    if st.session_state.normality_results:
        st.divider()
        _nhc1, _nhc2 = st.columns([5, 1])
        _nhc1.markdown(f"**Zapisane testy normalno\u015bci ({len(st.session_state.normality_results)}):**")
        with _nhc2:
            if st.button("\U0001f5d1\ufe0f Usu\u0144 wszystkie", key="del_all_norm",
                         use_container_width=True):
                st.session_state.normality_results = {}
                st.rerun()

        for result_key, nres in list(st.session_state.normality_results.items()):
            var = nres.get('var', result_key)
            lbl = var_labels.get(var, var)
            _grp_l = nres.get('group_label', '')
            _title = f"\U0001f4ca {var} \u2014 {lbl}" + (f" | {_grp_l}" if _grp_l else "")
            _nec1, _nec2 = st.columns([6, 1])
            with _nec1:
                _nexp = st.expander(_title, expanded=True)
            with _nec2:
                if st.button("\U0001f5d1\ufe0f", key=f"del_norm_{result_key}",
                             help=f"Usu\u0144 {result_key}"):
                    st.session_state.normality_results.pop(result_key, None)
                    st.rerun()
            with _nexp:
                _split_badge(_grp_l)
                x = np.array(nres['x'])
                n = nres['n']
                st.caption(f"N = {nres['n_raw']:,} obserwacji (efektywne N do test\u00f3w: {n:,})")
                if nres.get('warn_msg'):
                    st.warning(nres['warn_msg'])

                def _style_norm_row(row):
                    color = '#E2EFDA' if '\u2705' in row['Wynik'] else '#FCE4D6'
                    return [f'background-color: {color}'] * len(row)

                st.dataframe(nres['test_df'].style.apply(_style_norm_row, axis=1)
                             .format({'Statystyka': '{:.4f}', 'p-value': '{:.4f}'}),
                             use_container_width=True, hide_index=True)

                if nres.get('show_desc'):
                    skew_v = float(stats.skew(x))
                    kurt_v = float(stats.kurtosis(x))
                    sk_interp = ("symetryczny" if abs(skew_v) < 0.5
                                 else ("lekko sko\u015bny" if abs(skew_v) < 1.0
                                       else ("umiarkowanie sko\u015bny" if abs(skew_v) < 2.0
                                             else "silnie sko\u015bny")))
                    kt_interp = ("mezokurtyczny (normalny)" if abs(kurt_v) < 1.0
                                 else ("platykurtyczny (sp\u0142aszczony)" if kurt_v < 0
                                       else "leptokurtyczny (smuk\u0142y)"))
                    desc_c1, desc_c2, desc_c3, desc_c4 = st.columns(4)
                    desc_c1.metric("Sko\u015bno\u015b\u0107", f"{skew_v:.4f}")
                    desc_c2.metric("Kurtoza", f"{kurt_v:.4f}")
                    desc_c3.metric("Ocena sko\u015bno\u015bci", sk_interp)
                    desc_c4.metric("Ocena kurtozy", kt_interp)

                plot_cols = st.columns(2 if (nres.get('show_qq') and nres.get('show_hist')) else 1)
                plot_idx = 0

                if nres.get('show_qq'):
                    (osm, osr), (slope, intercept, r) = stats.probplot(x, dist='norm')
                    fig_qq = go.Figure()
                    fig_qq.add_trace(go.Scatter(
                        x=list(osm), y=list(osr),
                        mode='markers', name='Obserwacje',
                        marker=dict(color='#2E75B6', size=4, opacity=0.6)
                    ))
                    fig_qq.add_trace(go.Scatter(
                        x=[min(osm), max(osm)],
                        y=[slope * min(osm) + intercept, slope * max(osm) + intercept],
                        mode='lines', name='Linia normalna',
                        line=dict(color='#C00000', dash='dash')
                    ))
                    fig_qq.update_layout(
                        title=f"Wykres Q-Q: {lbl}",
                        xaxis_title="Kwantyle teoretyczne",
                        yaxis_title="Kwantyle pr\u00f3bkowe",
                        height=350, showlegend=True
                    )
                    with plot_cols[plot_idx]:
                        st.plotly_chart(fig_qq, use_container_width=True, key=f"qq_{result_key}")
                    plot_idx += 1

                if nres.get('show_hist'):
                    fig_hist = go.Figure()
                    fig_hist.add_trace(go.Histogram(
                        x=x, name="Dane",
                        histnorm='probability density',
                        marker_color='#2E75B6', opacity=0.7,
                        nbinsx=min(50, max(10, n // 10))
                    ))
                    x_range = np.linspace(x.min(), x.max(), 200)
                    mu, sigma = float(x.mean()), float(x.std())
                    y_norm = stats.norm.pdf(x_range, mu, sigma)
                    fig_hist.add_trace(go.Scatter(
                        x=x_range, y=y_norm, mode='lines', name='Rozk\u0142ad normalny',
                        line=dict(color='#C00000', width=2)
                    ))
                    fig_hist.update_layout(
                        title=f"Histogram z krzywa normaln\u0105: {lbl}",
                        xaxis_title=lbl, yaxis_title="G\u0119sto\u015b\u0107",
                        height=350, showlegend=True, barmode='overlay'
                    )
                    with plot_cols[plot_idx]:
                        st.plotly_chart(fig_hist, use_container_width=True, key=f"hist_{result_key}")


elif menu == "\U0001f4ca ANOVA":
    _require_module_access("anova")
    _require_data()
    module_header("\U0001f4ca", "ANOVA", "Jednoczynnikowa analiza wariancji z testem post-hoc Tukeya")

    with st.expander("Jak wykona\u0107 i interpretowa\u0107 ANOVA -- kliknij aby rozwin\u0105\u0107", expanded=False):
        st.markdown("""
**Kiedy u\u017cywa\u0107?** Gdy chcesz por\u00f3wna\u0107 **\u015brednie 3 lub wi\u0119cej grup** jednocze\u015bnie.
Przyk\u0142ad: czy \u015brednia satysfakcja r\u00f3\u017cni si\u0119 mi\u0119dzy miastem A, B i C?
Dla 2 grup u\u017cyj testu T (modu\u0142 Analizy \u2192 \u015arednie).

**Jak wykona\u0107:**
1. Wybierz **zmienn\u0105 zale\u017cn\u0105** (ci\u0105g\u0142a, numeryczna -- np. satysfakcja 1-10).
2. Wybierz **czynnik grupuj\u0105cy** (kategoryczna -- np. miasto, wiek, segment).
3. Kliknij **Uruchom ANOVA**.

**Jak interpretowa\u0107 wyniki:**

| Wska\u017anik | Interpretacja |
|---|---|
| **F** | Im wy\u017csze, tym wi\u0119ksza r\u00f3\u017cnica mi\u0119dzy grupami wzgl\u0119dem zmienno\u015bci wewn\u0105trz grup. |
| **p-value** | p < 0.05 = grupy r\u00f3\u017cni\u0105 si\u0119 istotnie statystycznie. |
| **Eta\u00b2 (\u03b7\u00b2)** | Miara si\u0142y efektu: < 0.01 s\u0142aby \u00b7 0.01-0.06 umiarkowany \u00b7 > 0.14 du\u017cy. |
| **MS (mi\u0119dzy grupami)** | Wariancja wyja\u015bniona przez przynale\u017cno\u015b\u0107 do grupy. |
| **MS (wewn\u0105trz grup)** | Wariancja wewn\u0105trz ka\u017cdej grupy (b\u0142\u0105d). |
| **Test Levene'a** | Sprawdza jednorodnosc wariancji. p < 0.05 = wariancje niejednorodne (naruszenie za\u0142o\u017cenia). |

**Test post-hoc Tukey HSD:**
Gdy ANOVA jest istotna (p < 0.05), Tukey wskazuje **kt\u00f3re konkretnie pary grup** r\u00f3\u017cni\u0105 si\u0119 od siebie.
p < 0.05 dla danej pary = ta para jest istotnie r\u00f3\u017cna.

**Za\u0142o\u017cenia ANOVA:** normalno\u015b\u0107 rozk\u0142adu w grupach, jednorodnosc wariancji (Levene), niezale\u017cno\u015b\u0107 obserwacji.
        """)

    st.info("Por\u00f3wnaj \u015brednie zmiennej ci\u0105g\u0142ej mi\u0119dzy grupami zdefiniowanymi przez zmienn\u0105 kategoryczn\u0105.")

    col1, col2 = st.columns(2)
    with col1:
        anova_dep = st.selectbox("\U0001f3af Zmienna zale\u017cna (ci\u0105g\u0142a):", numeric_cols, format_func=lambda x: get_var_display_name(x, var_labels))
    with col2:
        anova_grp = st.selectbox("\U0001f465 Czynnik grupuj\u0105cy (kategoryczna):", visible_columns, format_func=lambda x: get_var_display_name(x, var_labels))

    if _tracked_button("\u25b6\ufe0f Uruchom ANOVA", "anova", "run_anova", type="primary"):
        with st.spinner("Obliczanie ANOVA..."):
            _w_full = st.session_state.weights if use_weights else None
            for _grp_lbl, _df_s, _df_raw_s, _w_s in _iter_split_groups(
                    df, df_raw, var_labels, st.session_state.split_var, weights=_w_full):
                result, err = run_anova(_df_raw_s, anova_dep, anova_grp, _df_s,
                                          weights=_w_s)
                if err:
                    _grp_disp = _grp_lbl or 'Pe\u0142na baza'
                    st.error(f"{_grp_disp}: {err}")
                else:
                    result['group_label'] = _grp_lbl
                    _merge_result(st.session_state.anova_results, result,
                        key_fn=lambda r: (r.get('dep_var',''), r.get('group_var',''),
                                          r.get('group_label', '')))

    if st.session_state.anova_results:
        _anc1, _anc2 = st.columns([5, 1])
        _anc1.markdown(f"**Zapisane analizy ANOVA ({len(st.session_state.anova_results)}):**")
        with _anc2:
            if st.button("\U0001f5d1\ufe0f Usu\u0144 wszystkie", key="del_all_anova",
                         use_container_width=True):
                st.session_state.anova_results = []
                st.rerun()

    for _ai, res in enumerate(list(st.session_state.anova_results)):
        dep_l = var_labels.get(res['dep_var'], res['dep_var'])
        grp_l = var_labels.get(res['group_var'], res['group_var'])
        _ac1, _ac2 = st.columns([6, 1])
        _title_anova = f"\U0001f4ca ANOVA: [{res['dep_var']}] {dep_l} \u00d7 [{res['group_var']}] {grp_l}"
        if res.get('group_label'):
            _title_anova += f" | {res['group_label']}"
        with _ac1:
            _aexp = st.expander(_title_anova, expanded=True)
        with _ac2:
            if st.button("\U0001f5d1\ufe0f", key=f"del_anova_{_ai}",
                         help=f"Usu\u0144 {res['dep_var']} x {res['group_var']}"):
                st.session_state.anova_results.pop(_ai)
                st.rerun()
        with _aexp:
            _split_badge(res.get('group_label', ''))
            sig_label = "\u2705 Istotna statystycznie (p < 0.05)" if res['p'] < 0.05 else "\u274c Brak istotno\u015bci (p \u2265 0.05)"
            m1, m2, m3, m4 = st.columns(4)
            m1.metric("F", f"{res['F']:.3f}")
            m2.metric("p-value", f"{res['p']:.4f}")
            m3.metric("Eta\u00b2 (efekt)", f"{res['eta2']:.4f}")
            m4.metric("Wynik", sig_label)

            st.markdown("**Tabela ANOVA:**")
            anova_table = pd.DataFrame({
                '\u0179r\u00f3d\u0142o': ['Mi\u0119dzy grupami', 'Wewn\u0105trz grup', 'Og\u00f3\u0142em'],
                'SS': [res['ss_between'], res['ss_within'], res['ss_total']],
                'df': [res['df_between'], res['df_within'], res['df_between'] + res['df_within']],
                'MS': [res['ms_between'], res['ms_within'], ''],
                'F': [res['F'], '', ''],
                'p': [res['p'], '', ''],
                'Eta\u00b2': [res['eta2'], '', ''],
            })
            st.dataframe(anova_table, use_container_width=True, hide_index=True)

            st.markdown(f"**Test Levene'a (jednorodnosc wariancji):** stat={res['lev_stat']:.3f}, p={res['lev_p']:.4f} -- {'Wariancje niejednorodne (p<0.05)' if res['lev_p'] < 0.05 else 'Wariancje jednorodne'}")

            st.markdown("**Statystyki opisowe wg grupy:**")
            st.dataframe(res['desc_df'].style.format({'Srednia': '{:.3f}', 'Odch. std.': '{:.3f}', 'Min': '{:.2f}', 'Max': '{:.2f}'}),
                         use_container_width=True, hide_index=True)

            # Bar chart with error bars
            fig_anova = go.Figure()
            fig_anova.add_trace(go.Bar(
                x=res['desc_df']['Grupa'].astype(str),
                y=res['desc_df']['Srednia'],
                error_y=dict(type='data', array=res['desc_df']['Odch. std.'].values, visible=True),
                marker_color='#2E75B6', name='\u015arednia \u00b1 Odch.std.'
            ))
            fig_anova.update_layout(title=f"\u015arednie wg grup -- {dep_l}", xaxis_title=grp_l, yaxis_title='Srednia', height=350)
            st.plotly_chart(fig_anova, use_container_width=True, key=f"pc_anova_bar_{res['dep_var']}_{res['group_var']}_{res.get('group_label','')}")

            if not res['posthoc_df'].empty:
                st.markdown("**Test post-hoc Tukey HSD:**")
                st.dataframe(res['posthoc_df'], use_container_width=True, hide_index=True)

    if st.session_state.anova_results:
        if st.button("\U0001f5d1\ufe0f Wyczy\u015b\u0107 wyniki ANOVA", type="secondary"):
            st.session_state.anova_results = []
            st.rerun()

# -------------------------------------------------------------
# MODU? 6: ANALIZA CZYNNIKOWA
# -------------------------------------------------------------
elif menu == "\U0001f52c Analiza Czynnikowa":
    _require_module_access("factor")
    _require_data()
    module_header("\U0001f52c", "Analiza Czynnikowa", "Eksploracyjna Analiza Czynnikowa (EFA)")

    with st.expander("Jak wykona\u0107 i interpretowa\u0107 EFA -- kliknij aby rozwin\u0105\u0107", expanded=False):
        st.markdown("""
**Kiedy u\u017cywa\u0107?** Gdy chcesz odkry\u0107 **ukryte konstrukty (czynniki)** kryj\u0105ce si\u0119 za korelacjami mi\u0119dzy wieloma zmiennymi.
Przyk\u0142ad: 15 pyta\u0144 o satysfakcj\u0119 mo\u017ce odzwierciedla\u0107 3 ukryte wymiary: satysfakcja z produktu, obs\u0142ugi i ceny.

**Jak wykona\u0107:**
1. Wybierz **min. 3 zmienne numeryczne** ze wsp\u00f3lnej baterii pyta\u0144 (np. pytania Likerta 1-5).
2. Zdecyduj o **liczbie czynnik\u00f3w** -- zacznij od 2-4, u\u017cyj wykresu osypiska jako wskaz\u00f3wki.
3. Wybierz **rotacj\u0119**: Varimax (czynniki niezale\u017cne, najcz\u0119stsza), Promax (czynniki mog\u0105 by\u0107 powi\u0105zane).
4. Wybierz **metod\u0119 ekstrakcji**: Principal (PA), MinRes, ML.
5. Kliknij **Uruchom analiz\u0119 czynnikow\u0105**.

**Jak interpretowa\u0107 wyniki:**

| Wska\u017anik | Interpretacja |
|---|---|
| **KMO** | Adekwatno\u015b\u0107 pr\u00f3by: \u2265 0.9 znakomita \u00b7 \u2265 0.8 b.dobra \u00b7 \u2265 0.7 dobra \u00b7 \u2265 0.6 umiarkowana \u00b7 < 0.5 nieodpowiednia |
| **Test Bartletta (p)** | p < 0.05 = macierz korelacji nadaje si\u0119 do EFA (zmienne s\u0105 powi\u0105zane). |
| **Warto\u015b\u0107 w\u0142asna (EV)** | Zasada Kaisera: zachowaj czynniki z EV > 1. Sprawd\u017a na wykresie osypiska. |
| **% wyja\u015bnionej wariancji** | Ile zmienno\u015bci danych wyja\u015bnia dany czynnik. \u0141\u0105cznie powinno by\u0107 \u2265 50-60%. |
| **\u0141adunek czynnikowy** | Si\u0142a i kierunek powi\u0105zania zmiennej z czynnikiem. |\u0142adunek| \u2265 0.40 = istotny (pogrubiony). |
| **Komunalno\u015b\u0107 (h\u00b2)** | % wariancji danej zmiennej wyja\u015bniony przez wszystkie czynniki. < 0.30 = zmienna s\u0142abo pasuje. |

**Wykres osypiska (Scree Plot):** Szukaj miejsca, gdzie krzywa "ugina si\u0119" (elbow). Czynniki przed tym miejscem warto zachowa\u0107.

**Wskaz\u00f3wki praktyczne:**
- Ka\u017cda zmienna powinna \u0142adowa\u0107 istotnie (|\u2265 0.40|) na **jeden g\u0142\u00f3wny czynnik** (prosta struktura).
- Zmienne z nisk\u0105 komunalno\u015bci\u0105 (< 0.30) lub \u0142adunkami krzy\u017cowymi warto usun\u0105\u0107.
- Minimalna pr\u00f3ba: N \u2265 5 \u00d7 liczba zmiennych (najlepiej N \u2265 200).
        """)

    st.info("Zidentyfikuj ukryte konstrukty (czynniki) kryj\u0105ce si\u0119 za korelacjami mi\u0119dzy zmiennymi. Wyniki analogiczne do SPSS.")

    fa_vars = st.multiselect("Zmienne do analizy (min. 3):", numeric_cols, format_func=lambda x: get_var_display_name(x, var_labels))

    if fa_vars and len(fa_vars) >= 3:
        col1, col2, col3 = st.columns(3)
        with col1:
            n_factors = st.number_input("Liczba czynnik\u00f3w:", min_value=1, max_value=min(len(fa_vars) - 1, 15), value=min(3, len(fa_vars) - 1))
        with col2:
            rotation = st.selectbox("Rotacja:", ['varimax', 'promax', 'oblimin', 'quartimax', 'none'])
        with col3:
            method = st.selectbox("Metoda ekstrakcji:", ['principal', 'minres', 'ml'])

        show_scree = st.checkbox("\U0001f4c8 Wykres osypiska (Scree Plot)")

        if _tracked_button("\u25b6\ufe0f Uruchom analiz\u0119 czynnikow\u0105", "factor", "run_factor", type="primary"):
            with st.spinner("Obliczanie analizy czynnikowej..."):
                _w_full = st.session_state.weights if use_weights else None
                for _grp_lbl, _df_s, _df_raw_s, _w_s in _iter_split_groups(
                        df, df_raw, var_labels, st.session_state.split_var, weights=_w_full):
                    result, err = run_factor_analysis(
                        _df_raw_s, fa_vars, int(n_factors), rotation, method,
                        weights=_w_s
                    )
                    if err:
                        _grp_disp = _grp_lbl or 'Pe\u0142na baza'
                        st.error(f"{_grp_disp}: {err}")
                    else:
                        result['group_label'] = _grp_lbl
                        _merge_result(st.session_state.factor_results, result,
                            key_fn=lambda r: (tuple(sorted(r.get('variables',[]))),
                                              r.get('rotation',''),
                                              r.get('group_label', '')))

        if st.session_state.factor_results:
            _fcc1, _fcc2 = st.columns([5, 1])
            _fcc1.markdown(f"**Zapisane analizy czynnikowe ({len(st.session_state.factor_results)}):**")
            with _fcc2:
                if st.button("\U0001f5d1\ufe0f Usu\u0144 wszystkie", key="del_all_factor",
                             use_container_width=True):
                    st.session_state.factor_results = []
                    st.rerun()

        for _fi, res in enumerate(list(st.session_state.factor_results)):
            _fec1, _fec2 = st.columns([6, 1])
            with _fec1:
                _fexp = st.expander(
                    f"\U0001f52c Analiza czynnikowa -- {res['rotation'].upper()} -- N={res['n']}"
                    + (f" | {res.get('group_label','')}" if res.get('group_label') else ""),
                    expanded=True
                )
            with _fec2:
                if st.button("\U0001f5d1\ufe0f", key=f"del_factor_{_fi}",
                             help="Usu\u0144 analiz\u0119"):
                    st.session_state.factor_results.pop(_fi)
                    st.rerun()
            with _fexp:
                _split_badge(res.get('group_label', ''))
                m1, m2, m3, m4 = st.columns(4)
                m1.metric("N obserwacji", f"{int(round(float(res['n']))):,}")
                m2.metric("KMO", f"{res['kmo']:.3f}", help="\u22650.7 = dobra adekwatnosc proby")
                m3.metric("Bartlett Chi\u00b2", f"{res['bartlett_chi2']:.2f}")
                m4.metric("Bartlett p", f"{res['bartlett_p']:.4f}", delta="\u2705 OK" if res['bartlett_p'] < 0.05 else "\u274c")

                kmo_interp = ("Nieodpowiednia" if res['kmo'] < 0.5 else "S\u0142aba" if res['kmo'] < 0.6 else
                              "Umiarkowana" if res['kmo'] < 0.7 else "Dobra" if res['kmo'] < 0.8 else
                              "Bardzo dobra" if res['kmo'] < 0.9 else "Znakomita")
                st.caption(f"KMO = {res['kmo']:.3f} \u2192 adekwatnosc proby: **{kmo_interp}** | Bartlett p {'< 0.05 -- nadaje sie do EFA' if res['bartlett_p'] < 0.05 else '>= 0.05 -- macierz moze byc jednostkowa'}")

                st.markdown("**Macierz \u0142adunk\u00f3w czynnikowych** (pogrubione |\u0142adunek| \u2265 0.40):")


                load_display = res['loadings'].copy()
                comm_col = res['communalities']['Komunalnosc (h2)']
                load_display['Komunalnosc (h2)'] = comm_col
                styled_load = load_display.style \
                    .format('{:.3f}') \
                    .map(_style_loading, subset=res['loadings'].columns.tolist())
                st.dataframe(styled_load, use_container_width=True)

                st.markdown("**Wyja\u015bniona wariancja:**")
                st.dataframe(res['variance'].style.format('{:.3f}'), use_container_width=True)

                if show_scree:
                    ev_vals = res['eigenvalues']['Warto\u015b\u0107 w\u0142asna'].values[:min(len(fa_vars), 15)]
                    fig_scree = go.Figure()
                    fig_scree.add_trace(go.Scatter(y=ev_vals, x=list(range(1, len(ev_vals) + 1)),
                                                   mode='lines+markers', name='Warto\u015bci w\u0142asne',
                                                   line=dict(color='#2E75B6', width=2), marker=dict(size=8)))
                    fig_scree.add_hline(y=1, line_dash='dash', line_color='red', annotation_text='Kryterium Kaisera (EV=1)')
                    fig_scree.update_layout(title='Wykres osypiska (Scree Plot)',
                                            xaxis_title='Numer czynnika', yaxis_title='Warto\u015b\u0107 w\u0142asna', height=350)
                    st.plotly_chart(fig_scree, use_container_width=True, key="pc_efa_scree")

    elif fa_vars:
        st.warning("Wybierz co najmniej 3 zmienne.")

    if st.session_state.factor_results:
        if st.button("\U0001f5d1\ufe0f Wyczy\u015b\u0107 wyniki analizy czynnikowej", type="secondary"):
            st.session_state.factor_results = []
            st.rerun()

# -------------------------------------------------------------
# MODU? 7: EKSPORT DO EXCELA
# -------------------------------------------------------------
# =============================================================
# MODUL: CONJOINT
# =============================================================
elif menu == "\U0001f4ca Conjoint":
    _require_module_access("conjoint")
    _require_data()
    module_header("\U0001f4ca", "Analiza Conjoint", "Rating-based (OLS) i CBC (Logit) \u2014 u\u017cyteczno\u015bci cz\u0105stkowe i wa\u017cno\u015b\u0107 atrybut\u00f3w")
    _weights_ignored_note(use_weights)

    with st.expander("Jak wykona\u0107 i interpretowa\u0107 -- kliknij aby rozwin\u0105\u0107", expanded=False):
        st.markdown("""
**Czym jest Conjoint?**
Conjoint (analiza l\u0105czna) mierzy, jak poszczeg\u00f3lne cechy produktu wp\u0142ywaj\u0105 na preferencje respondent\u00f3w. Wynikiem s\u0105 **u\u017cyteczno\u015bci cz\u0105stkowe** (part-worth utilities) oraz **wa\u017cno\u015b\u0107 atrybut\u00f3w**.

**Dwa dost\u0119pne warianty:**
- **Rating-based**: Respondenci oceniaj\u0105 profile produkt\u00f3w w skali (np. 1-10). Zmienna zale\u017cna = ocena.
- **CBC (Choice-Based)**: Respondenci wybieraj\u0105 mi\u0119dzy profilami. Zmienna zale\u017cna = 0/1 (czy profil zosta\u0142 wybrany).

**Jak interpretowa\u0107:**

| Wska\u017anik | Interpretacja |
|---|---|
| **Wa\u017cno\u015b\u0107 atrybutu (%)** | Im wy\u017csza, tym bardziej ten atrybut wp\u0142ywa na decyzje |
| **U\u017cyteczno\u015b\u0107 cz\u0105stkowa** | Dodatnia = preferowany poziom, ujemna = niepreferowan |
| **R\u00b2** | Odsetek wariancji ocen wyja\u015bniany przez model (rating) |
| **Pseudo R\u00b2** | Miara dopasowania modelu logit (CBC), >0.2 = dobre |

**Wymagania dotycz\u0105ce danych:**
- Rating: min. 30 respondent\u00f3w, zmienne atrybut\u00f3w kategoryczne lub liczbowe
- CBC: dane w formacie long (jeden wiersz = jeden profil-respondent)
        """)

    conj_method = st.radio("Wariant analizy:", ["Rating-based (OLS)", "CBC (Choice-Based Logit)"],
                            horizontal=True, key="conj_method")
    st.divider()

    col_a, col_b = st.columns(2)
    with col_a:
        if conj_method == "Rating-based (OLS)":
            conj_rating = st.selectbox("Zmienna zale\u017cna (ocena profilu):",
                                        numeric_cols, format_func=lambda x: get_var_display_name(x, var_labels),
                                        key="conj_rating")
            conj_attrs = st.multiselect("Atrybuty produktu (zmienne niezale\u017cne):",
                                         [c for c in visible_columns if c != conj_rating],
                                         format_func=lambda x: get_var_display_name(x, var_labels),
                                         key="conj_attrs_r")
        else:
            conj_choice = st.selectbox("Zmienna wyboru (0=nie, 1=tak):",
                                        numeric_cols, format_func=lambda x: get_var_display_name(x, var_labels),
                                        key="conj_choice")
            conj_attrs = st.multiselect("Atrybuty produkt\u00f3w:",
                                         [c for c in visible_columns if c != conj_choice],
                                         format_func=lambda x: get_var_display_name(x, var_labels),
                                         key="conj_attrs_c")
    with col_b:
        st.markdown("**Wskaz\u00f3wka:**")
        if conj_method == "Rating-based (OLS)":
            st.info("Wybierz zmienn\u0105 z ocen\u0105 profilu (np. 1-10) i zmienne opisuj\u0105ce atrybuty (kategorie lub liczby).")
        else:
            st.info("Wybierz zmienn\u0105 binarny\u0105 wyboru (1=wybrany profil) i zmienne atrybut\u00f3w. Dane musz\u0105 by\u0107 w formacie long.")

    if _tracked_button("\u25b6\ufe0f Uruchom analiz\u0119 Conjoint", "conjoint", "run_conjoint", type="primary"):
        if not conj_attrs:
            st.error("Wybierz co najmniej jeden atrybut.")
        else:
            with st.spinner("Obliczanie..."):
                for _grp_lbl, _df_s, _df_raw_s, _w_s in _iter_split_groups(
                        df, df_raw, var_labels, st.session_state.split_var, weights=None):
                    if conj_method == "Rating-based (OLS)":
                        res, err = run_conjoint_rating(_df_raw_s, conj_rating, conj_attrs)
                    else:
                        res, err = run_conjoint_cbc(_df_raw_s, conj_choice, conj_attrs)
                    if err:
                        _grp_disp = _grp_lbl or 'Pe\u0142na baza'
                        st.error(f"{_grp_disp}: {err}")
                    else:
                        res['group_label'] = _grp_lbl
                        _merge_result(st.session_state.conjoint_results, res,
                            key_fn=lambda r: (r.get('method',''),
                                              tuple(sorted(r.get('attribute_vars',[]))),
                                              r.get('group_label', '')))
                st.success("\u2705 Analiza Conjoint uko\u0144czona!")

    if st.session_state.conjoint_results:
        _cjc1, _cjc2 = st.columns([5, 1])
        _cjc1.markdown(f"**Zapisane analizy Conjoint ({len(st.session_state.conjoint_results)}):**")
        with _cjc2:
            if st.button("\U0001f5d1\ufe0f Usu\u0144 wszystkie", key="del_all_conj",
                         use_container_width=True):
                st.session_state.conjoint_results = []
                st.rerun()

    for _cji, res in enumerate(list(st.session_state.conjoint_results)):
        if res.get('error'):
            st.error(res['error']); continue
        _cjec1, _cjec2 = st.columns([6, 1])
        with _cjec1:
            _cjexp = st.expander(
                f"\U0001f4ca {res['method']} -- {len(res['attribute_vars'])} atrybut\u00f3w"
                + (f" | {res.get('group_label','')}" if res.get('group_label') else ""),
                expanded=True
            )
        with _cjec2:
            if st.button("\U0001f5d1\ufe0f", key=f"del_conj_{_cji}",
                         help="Usu\u0144 analiz\u0119"):
                st.session_state.conjoint_results.pop(_cji)
                st.rerun()
        with _cjexp:
            _split_badge(res.get('group_label', ''))
            # Summary metrics
            mc1, mc2, mc3, mc4 = st.columns(4)
            mc1.metric("N", f"{int(round(float(res['n']))):,}")
            if 'r2' in res:
                mc2.metric("R\u00b2", f"{res['r2']:.4f}")
                mc3.metric("R\u00b2 skor.", f"{res['r2_adj']:.4f}")
                mc4.metric("p (F)", f"{res['p']:.4f} {'OK' if res['p'] < 0.05 else 'NS'}")
            elif 'pseudo_r2' in res:
                mc2.metric("Pseudo R\u00b2", f"{res['pseudo_r2']:.4f}")
                mc3.metric("LLR p", f"{res['llr_pvalue']:.4f}")

            # Importance chart
            st.markdown("**Wa\u017cno\u015b\u0107 atrybut\u00f3w:**")
            imp_df = pd.DataFrame(list(res['importance'].items()), columns=['Atrybut', 'Wa\u017cno\u015b\u0107 (%)'])
            imp_df['Etykieta'] = imp_df['Atrybut'].apply(lambda x: var_labels.get(x, x))
            imp_df = imp_df.sort_values('Wa\u017cno\u015b\u0107 (%)', ascending=True)
            fig_imp = px.bar(imp_df, x='Wa\u017cno\u015b\u0107 (%)', y='Etykieta', orientation='h',
                             color='Wa\u017cno\u015b\u0107 (%)', color_continuous_scale='Blues',
                             title='Wa\u017cno\u015b\u0107 atrybut\u00f3w (%)'),
            st.plotly_chart(fig_imp[0], use_container_width=True, key=f"pc_conj_imp_{_cji}")

            # Utilities per attribute
            st.markdown("**U\u017cyteczno\u015bci cz\u0105stkowe (part-worth utilities):**")
            for attr, utils in res['utilities'].items():
                if not utils: continue
                attr_lbl = var_labels.get(attr, attr)
                # Build display labels: "attr_valueLabel" instead of "attr_code"
                _vvl = meta_orig.variable_value_labels.get(attr, {}) if is_spss else {}
                _cvl = st.session_state.custom_val_labels.get(attr, {})
                def _nice_level(raw_lvl):
                    """Convert dummy-column name like 'Q1_A1' \u2192 'Q1_zdecydowanie'."""
                    s = str(raw_lvl)
                    # Strip attribute prefix (e.g. 'Q1_') to isolate the code
                    prefix = f"{attr}_"
                    code_part = s[len(prefix):] if s.startswith(prefix) else s
                    # Try custom val labels first, then SPSS labels
                    lbl_txt = _cvl.get(code_part, _cvl.get(str(code_part), ""))
                    if not lbl_txt:
                        # Try numeric conversion for SPSS labels
                        for key_variant in (code_part, str(code_part)):
                            try:
                                fv = float(key_variant)
                                lbl_txt = _vvl.get(fv, _vvl.get(int(fv), ""))
                                if lbl_txt:
                                    break
                            except (ValueError, TypeError):
                                pass
                        if not lbl_txt:
                            lbl_txt = _vvl.get(code_part, "")
                    return f"{attr}_{lbl_txt}" if lbl_txt else s

                u_df = pd.DataFrame(list(utils.items()), columns=['Poziom', 'U\u017cyteczno\u015b\u0107'])
                u_df['Poziom'] = u_df['Poziom'].apply(_nice_level)
                u_df = u_df.sort_values('U\u017cyteczno\u015b\u0107', ascending=True)
                fig_u = px.bar(u_df, x='U\u017cyteczno\u015b\u0107', y='Poziom', orientation='h',
                               title=f"[{attr}] {attr_lbl}",
                               color='U\u017cyteczno\u015b\u0107', color_continuous_scale='RdYlGn',
                               color_continuous_midpoint=0)
                fig_u.add_vline(x=0, line_dash='dash', line_color='gray')
                st.plotly_chart(fig_u, use_container_width=True, key=f"pc_conj_util_{_cji}_{attr}")

    if st.session_state.conjoint_results:
        if st.button("\U0001f5d1\ufe0f Wyczy\u015b\u0107 wyniki Conjoint", type="secondary"):
            st.session_state.conjoint_results = []
            st.rerun()


# =============================================================
# MODUL: MAXDIFF
# =============================================================
elif menu == "\U0001f522 MaxDiff":
    _require_module_access("maxdiff")
    _require_data()
    module_header("\U0001f522", "MaxDiff", "Best-Worst Scaling \u2014 ranking wa\u017cno\u015bci element\u00f3w")
    _weights_ignored_note(use_weights)

    with st.expander("Jak wykona\u0107 i interpretowa\u0107 -- kliknij aby rozwin\u0105\u0107", expanded=False):
        st.markdown("""
**Czym jest MaxDiff?**
MaxDiff (Maximum Difference Scaling) mierzy wzgl\u0119dn\u0105 wa\u017cno\u015b\u0107/preferencj\u0119 element\u00f3w.
Respondenci wskazuj\u0105 **Najwa\u017cniejszy (Best)** i **Najmniej wa\u017cny (Worst)** spo\u015br\u00f3d ka\u017cdego zestawu pozycji.

**Format danych w pliku SPSS:**
Ka\u017cdy **zestaw** ma dwie kolumny:
- Kolumna **Best** -- kt\u00f3r\u0105 pozycj\u0119 wybrano jako najwa\u017cniejsz\u0105 (np. "Produkt A")
- Kolumna **Worst** -- kt\u00f3r\u0105 wybrano jako najmniej wa\u017cn\u0105

Przyk\u0142ad: `Zestaw1_Best`, `Zestaw1_Worst`, `Zestaw2_Best`, `Zestaw2_Worst`, ...

**Jak skonfigurowa\u0107:**
1. Podaj pary kolumn (Best / Worst) dla ka\u017cdego zestawu.
2. Podaj list\u0119 element\u00f3w (pozycji) wyst\u0119puj\u0105cych w tych kolumnach.
3. Kliknij **Uruchom MaxDiff**.

**Jak interpretowa\u0107 wyniki:**

| Wska\u017anik | Interpretacja |
|---|---|
| **B-W Score** | Liczba wybor\u00f3w Best minus Worst. Wy\u017cszy = preferowany |
| **B-W Score (%)** | B-W Score / N respondent\u00f3w * 100. Por\u00f3wnywalny mi\u0119dzy badaniami |
| **Wynik standaryzowany (0-100)** | Rescalowany do skali 0-100. Najlepszy element = 100 |
        """)

    st.divider()
    st.markdown("##### 1. Zdefiniuj pary kolumn Best/Worst")
    st.info("Dla ka\u017cdego zestawu wybierz kolumn\u0119 'Najwa\u017cniejszy' (Best) i 'Najmniej wa\u017cny' (Worst).")

    if 'maxdiff_pairs' not in st.session_state:
        st.session_state.maxdiff_pairs = [('', '')]

    pairs_to_remove = None
    new_pairs = []

    def _md_fmt(x):
        if not x:
            return '-- wybierz --'
        lbl = var_labels.get(x, '')
        return f"[{x}] {lbl}" if lbl else x

    for pi, (bc, wc) in enumerate(st.session_state.maxdiff_pairs):
        col_b, col_w, col_rm = st.columns([3, 3, 1])
        with col_b:
            sel_b = st.selectbox(f"Zestaw {pi+1} -- Najwa\u017cniejszy (Best):",
                                  [''] + list(df_raw.columns),
                                  index=list([''] + list(df_raw.columns)).index(bc) if bc in df_raw.columns else 0,
                                  format_func=_md_fmt,
                                  key=f"md_best_{pi}")
        with col_w:
            sel_w = st.selectbox(f"Zestaw {pi+1} -- Najmniej wa\u017cny (Worst):",
                                  [''] + list(df_raw.columns),
                                  index=list([''] + list(df_raw.columns)).index(wc) if wc in df_raw.columns else 0,
                                  format_func=_md_fmt,
                                  key=f"md_worst_{pi}")
        with col_rm:
            st.write("")
            if pi > 0 and st.button("\U0001f5d1\ufe0f", key=f"md_rm_{pi}"):
                pairs_to_remove = pi
        new_pairs.append((sel_b, sel_w))

    st.session_state.maxdiff_pairs = new_pairs
    if pairs_to_remove is not None:
        st.session_state.maxdiff_pairs.pop(pairs_to_remove)
        st.rerun()

    if st.button("\u2795 Dodaj zestaw", key="md_add_pair"):
        st.session_state.maxdiff_pairs.append(('', ''))
        st.rerun()

    st.divider()
    st.markdown("##### 2. Okre\u015bl pozycje (elementy) badania")

    # Auto-detect items from selected columns
    valid_pairs = [(b, w) for b, w in st.session_state.maxdiff_pairs if b and w and b in df_raw.columns and w in df_raw.columns]
    if valid_pairs:
        auto_items = set()
        for bc, wc in valid_pairs:
            auto_items.update(df_raw[bc].dropna().astype(str).unique())
            auto_items.update(df_raw[wc].dropna().astype(str).unique())
        auto_items = sorted(auto_items)
        st.caption(f"Automatycznie wykryto {len(auto_items)} unikalnych pozycji z wybranych kolumn.")
        md_items_raw = st.text_area(
            "Pozycje (jedna na wiersz):",
            value='\n'.join(auto_items),
            height=180,
            key="md_items",
            help="Mo\u017cesz edytowa\u0107 list\u0119 i zmienia\u0107 kolejno\u015b\u0107."
        )
        md_items = [x.strip() for x in md_items_raw.splitlines() if x.strip()]
    else:
        md_items_raw = st.text_area("Pozycje (jedna na wiersz):", height=150, key="md_items_manual")
        md_items = [x.strip() for x in md_items_raw.splitlines() if x.strip()]

    st.divider()
    st.markdown("##### 3. Nazwa analizy i uruchomienie")
    md_name = st.text_input("Nazwa analizy MaxDiff:", value="MaxDiff", key="md_name")

    if _tracked_button("\u25b6\ufe0f Uruchom analiz\u0119 MaxDiff", "maxdiff", "run_maxdiff", type="primary"):
        if not valid_pairs:
            st.error("Wybierz co najmniej jedn\u0105 par\u0119 kolumn Best/Worst.")
        elif len(md_items) < 2:
            st.error("Podaj co najmniej 2 pozycje.")
        else:
            with st.spinner("Obliczanie wynik\u00f3w MaxDiff..."):
                for _grp_lbl, _df_s, _df_raw_s, _w_s in _iter_split_groups(
                        df, df_raw, var_labels, st.session_state.split_var, weights=None):
                    df_scores = run_maxdiff(_df_raw_s, valid_pairs, md_items)
                    name_eff = f"{md_name} | {_grp_lbl}" if _grp_lbl else md_name
                    result_md = {
                        'name': name_eff,
                        'group_label': _grp_lbl,
                        'pairs': valid_pairs,
                        'items': md_items,
                        'n_resp': len(_df_raw_s),
                        'n_tasks': len(valid_pairs),
                        'scores': df_scores,
                    }
                    _merge_result(st.session_state.maxdiff_results, result_md,
                        key_fn=lambda r: r.get('name',''))
            st.success(f"\u2705 MaxDiff uko\u0144czony! Przeanalizowano {len(valid_pairs)} zestaw\u00f3w, {len(md_items)} pozycji.")

    if st.session_state.maxdiff_results:
        _mdc1, _mdc2 = st.columns([5, 1])
        _mdc1.markdown(f"**Zapisane analizy MaxDiff ({len(st.session_state.maxdiff_results)}):**")
        with _mdc2:
            if st.button("\U0001f5d1\ufe0f Usu\u0144 wszystkie", key="del_all_md",
                         use_container_width=True):
                st.session_state.maxdiff_results = []
                st.rerun()

    for _mdi, res in enumerate(list(st.session_state.maxdiff_results)):
        _base_md, _grp_md = _extract_split_from_title(res['name'])
        _mdec1, _mdec2 = st.columns([6, 1])
        with _mdec1:
            _mdexp = st.expander(
                f"\U0001f522 {_base_md} -- {res['n_tasks']} zestaw\u00f3w, {len(res['items'])} pozycji"
                + (f" \u2014 \U0001f500 {_grp_md}" if _grp_md else ""),
                expanded=True
            )
        with _mdec2:
            if st.button("\U0001f5d1\ufe0f", key=f"del_md_{_mdi}",
                         help=f"Usu\u0144 {res['name']}"):
                st.session_state.maxdiff_results.pop(_mdi)
                st.rerun()
        with _mdexp:
            _split_badge(_grp_md)
            df_s = res['scores']
            mc1, mc2, mc3 = st.columns(3)
            mc1.metric("N respondent\u00f3w", f"{int(round(float(res['n_resp']))):,}")
            mc2.metric("Liczba zestaw\u00f3w", res['n_tasks'])
            mc3.metric("Pozycji", len(res['items']))

            st.markdown("**Wyniki MaxDiff -- ranking wa\u017cno\u015bci:**")

            _style_md = _make_style_md(len(df_s))

            st.dataframe(
                df_s.style.apply(_style_md, axis=1)
                    .format({'Best [N]': '{:.0f}', 'Worst [N]': '{:.0f}', 'Pokazano [N]': '{:.0f}',
                             'B-W Score': '{:.0f}', 'B-W Score (%)': '{:.1f}',
                             'Wynik standaryzowany (0-100)': '{:.1f}'}),
                use_container_width=True, hide_index=True
            )

            # Bar chart
            fig_md = px.bar(
                df_s.sort_values('Wynik standaryzowany (0-100)', ascending=True),
                x='Wynik standaryzowany (0-100)', y='Item', orientation='h',
                color='Wynik standaryzowany (0-100)', color_continuous_scale='Blues',
                title=f"MaxDiff -- Ranking wa\u017cno\u015bci: {res['name']}",
                text='Wynik standaryzowany (0-100)'
            )
            fig_md.update_traces(texttemplate='%{text:.1f}', textposition='outside')
            fig_md.update_layout(height=max(300, len(res['items']) * 35 + 80),
                                  coloraxis_showscale=False)
            st.plotly_chart(fig_md, use_container_width=True, key="pc_maxdiff_bar")

    if st.session_state.maxdiff_results:
        if st.button("\U0001f5d1\ufe0f Wyczy\u015b\u0107 wyniki MaxDiff", type="secondary"):
            st.session_state.maxdiff_results = []
            st.rerun()


# =============================================================
# MODUL: CHMURA SLOW
# =============================================================
# =============================================================
# MODUL: SKUPIENIA HIERARCHICZNE
# =============================================================
elif menu == "\U0001f3af Skupienia i Segmentacja":
    _require_module_access("cluster")
    _require_data()
    module_header("\U0001f3af", "Skupienia i Segmentacja", "Skupienia hierarchiczne (dendrogram) i segmentacja K-Means")
    _weights_ignored_note(use_weights)
    tab_hc, tab_kmeans = st.tabs(["Skupienia Hierarchiczne", "Segmentacja K-Means"])

    with tab_hc:
        st.markdown("##### Analiza skupie\u0144 hierarchicznych")

        with st.expander("Jak wykona\u0107 i interpretowa\u0107", expanded=False):
            st.markdown("""
    **Czym s\u0105 skupienia hierarchiczne?**
    Metoda grupowania respondent\u00f3w bez konieczno\u015bci wcze\u015bniejszego podania liczby grup (w odst\u0119pstwie do K-Means).
    Wynikiem jest **dendrogram** \u2014 drzewo podobie\u0144stw, kt\u00f3re pomaga dobra\u0107 optym aln\u0105 liczb\u0119 skupie\u0144.

    **Etapy analizy:**
    1. Wybierz zmienne numeryczne i metod\u0119 (\u0142\u0105czenia)
    2. Odczytaj dendrogram \u2014 szukaj du\u017cych "skok\u00f3w" na osi Y (odleg\u0142o\u015b\u0107)
    3. Ustaw progowy ci\u0119cie lub liczb\u0119 skupie\u0144
    4. Dodaj zmiennn\u0105 z przypisaniem do skupie\u0144 do bazy

    | Metoda | Zastosowanie |
    |---|---|
    | **Ward** | Minimalizuje wariancj\u0119 wewn\u0105trzgrupow\u0105 \u2014 zazwyczaj najlepsza |
    | **Complete** | U\u017cywa maksymalnej odleg\u0142o\u015bci \u2014 tworzy zborne skupienia |
    | **Average** | U\u017cywa \u015bredniej odleg\u0142o\u015bci \u2014 kompromis |
    | **Single** | U\u017cywa minimalnej odleg\u0142o\u015bci \u2014 podatna na efekt \u0142a\u0144cucha |
            """)

        col_hc1, col_hc2 = st.columns([2, 1])
        with col_hc1:
            hc_vars = st.multiselect(
                "Zmienne numeryczne do analizy:",
                numeric_cols,
                format_func=lambda x: get_var_display_name(x, var_labels),
                key="hc_vars"
            )
            hc_method = st.selectbox(
                "Metoda \u0142\u0105czenia:",
                ["ward", "complete", "average", "single"],
                key="hc_method",
                help="Ward: zalecana. Complete/Average: r\u00f3wnowa\u017cne. Single: nie polecana."
            )
            hc_metric = st.selectbox(
                "Miara odleg\u0142o\u015bci:",
                ["euclidean", "cosine", "correlation"],
                key="hc_metric",
                help="Ward wymaga euclidean."
            )
            if hc_method == "ward" and hc_metric != "euclidean":
                st.warning("Metoda Ward wymaga odleg\u0142o\u015bci euklidesowej \u2014 zmieniono automatycznie.")
                hc_metric = "euclidean"

        with col_hc2:
            hc_standardize = st.checkbox("Standaryzuj zmienne (Z-score)", value=True, key="hc_std")
            hc_max_obs = st.number_input(
                "Maks. respondent\u00f3w (wydajno\u015b\u0107):",
                min_value=50, max_value=5000, value=500, step=50, key="hc_maxobs",
                help="Dendrogram dla du\u017cych baz jest nieczytelny. Losowa pr\u00f3bka."
            )
            hc_n_clusters = st.slider("Liczba skupie\u0144 do wyci\u0119cia:", 2, 15, 3, key="hc_nclust")
            hc_var_name = st.text_input(
                "Nazwa nowej zmiennej skupie\u0144:",
                value="Skupienie_H",
                key="hc_varname"
            )
            # Warn if the variable name is already used by an earlier hclust result
            _existing_hc_names = {r.get('var_name') for r in st.session_state.hclust_results}
            if hc_var_name.strip() in _existing_hc_names:
                st.warning(
                    f"\u26a0\ufe0f Nazwa `{hc_var_name.strip()}` jest ju\u017c u\u017cywana przez inn\u0105 "
                    "analiz\u0119 skupie\u0144. Po wygenerowaniu nowej, **poprzedni wynik zostanie nadpisany**. "
                    "Zmie\u0144 nazw\u0119, je\u015bli chcesz zachowa\u0107 obydwa wyniki."
                )

        if _tracked_button("\u25b6\ufe0f Generuj dendrogram i skupienia", "cluster", "run_hclust", type="primary", key="hc_run"):
            if len(hc_vars) < 2:
                st.error("Wybierz co najmniej 2 zmienne.")
            else:
                from scipy.cluster.hierarchy import linkage, dendrogram, fcluster
                from scipy.spatial.distance import pdist
                matplotlib.use('Agg')
                import matplotlib.pyplot as plt

                for _grp_lbl, _df_s, _df_raw_s, _w_s in _iter_split_groups(
                        df, df_raw, var_labels, st.session_state.split_var, weights=None):
                    # Per-group variable name suffix to avoid collisions between groups
                    grp_suffix = ""
                    if _grp_lbl:
                        _safe = _grp_lbl.replace('=', '_').replace(' ', '_')
                        grp_suffix = f"_{_safe}"
                    hc_var_name_eff = f"{hc_var_name}{grp_suffix}"

                    df_hc = _df_raw_s[hc_vars].dropna()
                    if len(df_hc) < 3:
                        _grp_disp = _grp_lbl or 'pe\u0142na baza'
                        st.warning(f"Za ma\u0142o obserwacji dla grupy `{_grp_disp}` (N={len(df_hc)}).")
                        continue
                    n_used = min(len(df_hc), int(hc_max_obs))
                    if len(df_hc) > n_used:
                        df_hc = df_hc.sample(n=n_used, random_state=42)
                        _grp_disp = _grp_lbl or 'Pe\u0142na baza'
                        st.info(f"[{_grp_disp}] Losowa pr\u00f3bka: {n_used} z {len(_df_raw_s[hc_vars].dropna())} kompletnych obserwacji.")

                    if hc_standardize:
                        from sklearn.preprocessing import StandardScaler
                        X = StandardScaler().fit_transform(df_hc.values)
                    else:
                        X = df_hc.values

                    _grp_disp_sp = _grp_lbl or 'pe\u0142na baza'
                    with st.spinner(f"Obliczanie skupie\u0144 dla `{_grp_disp_sp}`..."):
                        Z = linkage(X, method=hc_method,
                                    metric=hc_metric if hc_method != 'ward' else 'euclidean')

                    # Dendrogram
                    fig_dend, ax = plt.subplots(figsize=(14, 5))
                    dendrogram(
                        Z, ax=ax,
                        truncate_mode='lastp', p=50,
                        leaf_rotation=90, leaf_font_size=8,
                        color_threshold=0.7 * max(Z[:, 2]),
                        above_threshold_color='#888888',
                    )
                    _title_d = f"Dendrogram skupie\u0144 hierarchicznych ({hc_method.title()}, n={n_used})"
                    if _grp_lbl:
                        _title_d += f" \u2014 {_grp_lbl}"
                    ax.set_title(_title_d, fontsize=12, fontweight='bold')
                    ax.set_xlabel("Indeks obserwacji lub liczba skupionych obiekt\u00f3w")
                    ax.set_ylabel("Odleg\u0142o\u015b\u0107")
                    ax.axhline(y=Z[-int(hc_n_clusters)+1, 2], color='#C00000',
                               linestyle='--', linewidth=1.5,
                               label=f"Ci\u0119cie: {hc_n_clusters} skupie\u0144")
                    ax.legend(fontsize=9)
                    plt.tight_layout()
                    st.pyplot(fig_dend, use_container_width=True)

                    buf_d = io.BytesIO()
                    fig_dend.savefig(buf_d, format='png', dpi=150, bbox_inches='tight')
                    buf_d.seek(0)
                    plt.close(fig_dend)
                    st.download_button(
                        f"\u2b07\ufe0f Pobierz dendrogram (PNG)" + (f" \u2014 {_grp_lbl}" if _grp_lbl else ""),
                        data=buf_d.getvalue(),
                        file_name=f"dendrogram_{hc_method}{grp_suffix}.png",
                        mime="image/png",
                        key=f"hc_dl_{hc_var_name_eff}"
                    )

                    # Assign clusters to full dataset (of this group slice)
                    if hc_standardize:
                        X_full = StandardScaler().fit_transform(_df_raw_s[hc_vars].dropna().values)
                    else:
                        X_full = _df_raw_s[hc_vars].dropna().values

                    Z_full = linkage(X_full, method=hc_method,
                                     metric=hc_metric if hc_method != 'ward' else 'euclidean')
                    labels_full = fcluster(Z_full, hc_n_clusters, criterion='maxclust')

                    idx_full = _df_raw_s[hc_vars].dropna().index
                    df_raw.loc[idx_full, hc_var_name_eff] = labels_full
                    df.loc[idx_full,     hc_var_name_eff] = [f"Skupienie {c}" for c in labels_full]
                    var_labels[hc_var_name_eff] = f"Skupienia hierarchiczne ({hc_n_clusters} grup, {hc_method})" + (f" [{_grp_lbl}]" if _grp_lbl else "")

                    # Cluster sizes
                    sizes = pd.Series(labels_full).value_counts().sort_index()
                    sizes_df = pd.DataFrame({
                        "Skupienie": [f"Skupienie {i}" for i in sizes.index],
                        "N": sizes.values,
                        "%": (sizes.values / sizes.sum() * 100).round(1),
                    })

                    # Profile
                    profile_df = _df_raw_s.loc[idx_full, hc_vars + [hc_var_name_eff]].copy()
                    profile_df[hc_var_name_eff] = profile_df[hc_var_name_eff].astype(int)
                    cluster_means = profile_df.groupby(hc_var_name_eff)[hc_vars].mean().round(2)
                    cluster_means.index = [f"Skupienie {i}" for i in cluster_means.index]

                    result_entry = {
                        'method': hc_method,
                        'metric': hc_metric,
                        'n_clusters': int(hc_n_clusters),
                        'vars': hc_vars,
                        'var_name': hc_var_name_eff,
                        'group_label': _grp_lbl,
                        'n_obs': len(idx_full),
                        'sizes': sizes_df,
                        'profile': cluster_means,
                        'Z': Z_full.tolist(),
                        'standardize': hc_standardize,
                        'labels_data': {str(i): int(lbl) for i, lbl in zip(idx_full, labels_full)},
                    }
                    _merge_result(st.session_state.hclust_results, result_entry,
                        key_fn=lambda r: r.get('var_name',''))
                    _grp_disp = _grp_lbl or 'Pe\u0142na baza'
                    st.success(f"\u2705 [{_grp_disp}] Zmienna `{hc_var_name_eff}` dodana do bazy.")

        if st.session_state.hclust_results:
            _hcc1, _hcc2 = st.columns([5, 1])
            _hcc1.markdown(f"**Zapisane skupienia hierarchiczne ({len(st.session_state.hclust_results)}):**")
            with _hcc2:
                if st.button("\U0001f5d1\ufe0f Usu\u0144 wszystkie", key="del_all_hc",
                             use_container_width=True):
                    st.session_state.hclust_results = []
                    st.rerun()

        for _hci, res_hc in enumerate(list(st.session_state.hclust_results)):
            _hcec1, _hcec2 = st.columns([6, 1])
            with _hcec1:
                _hcexp = st.expander(
                    f"\U0001f333 {res_hc['var_name']} \u2014 {res_hc['n_clusters']} skupie\u0144 ({res_hc['method']})",
                    expanded=True
                )
            with _hcec2:
                if st.button("\U0001f5d1\ufe0f", key=f"del_hc_{_hci}",
                             help=f"Usu\u0144 {res_hc['var_name']}"):
                    st.session_state.hclust_results.pop(_hci)
                    st.rerun()
            with _hcexp:
                _split_badge(res_hc.get('group_label', ''))
                sc1, sc2, sc3 = st.columns(3)
                sc1.metric("N obserwacji", f"{int(round(float(res_hc['n_obs']))):,}")
                sc2.metric("Skupie\u0144", res_hc['n_clusters'])
                sc3.metric("Metoda", res_hc['method'].title())

                st.markdown("**Wielko\u015b\u0107 skupie\u0144:**")
                st.dataframe(res_hc['sizes'], use_container_width=True, hide_index=True)

                st.markdown("**Profil skupie\u0144 (\u015brednie zmiennych):**")
                st.dataframe(res_hc['profile'].style.format("{:.2f}"), use_container_width=True)

                fig_bar = px.bar(
                    res_hc['sizes'], x="Skupienie", y="N",
                    color="Skupienie", title=f"Liczebno\u015b\u0107 skupie\u0144: {res_hc['var_name']}",
                    color_discrete_sequence=px.colors.qualitative.Set2
                )
                fig_bar.update_layout(showlegend=False, height=300)
                st.plotly_chart(fig_bar, use_container_width=True, key=f"pc_hclust_{res_hc.get('var_name','hc')}")

        if st.session_state.hclust_results:
            if st.button("\U0001f5d1\ufe0f Wyczy\u015b\u0107 wyniki skupie\u0144", type="secondary", key="hc_clear"):
                st.session_state.hclust_results = []
                st.rerun()

    # \u2500\u2500 TAB: Segmentacja K-Means \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
    with tab_kmeans:
        st.markdown("##### Segmentacja K-Means")
        st.info("Pogrupuj respondent\u00f3w metod\u0105 K-Means. Nowa zmienna segmentacyjna zostanie dodana do bazy.")
        seg_vars = st.multiselect("Zmienne numeryczne:", numeric_cols,
                                   format_func=lambda x: get_var_display_name(x, var_labels),
                                   key="seg_vars_mod")
        k_clusters = st.slider("Liczba segment\u00f3w (K):", 2, 10, 3, key="k_clusters_mod")
        seg_name = st.text_input("Nazwa zmiennej segmentacyjnej:",
                                  value=f"Segmentacja_{len(st.session_state.segmentations) + 1}",
                                  key="seg_name_mod")
        if _tracked_button("\u25b6\ufe0f Wykonaj segmentacj\u0119 K-Means", "cluster", "run_kmeans", type="primary", key="seg_run_mod"):
            if len(seg_vars) < 2:
                st.error("Wybierz co najmniej 2 zmienne.")
            else:
                st.session_state.segmentations.append(
                    {'vars': seg_vars, 'k': k_clusters, 'name': seg_name})
                st.success(f"\u2705 Segmentacja `{seg_name}` utworzona.")
                st.rerun()

        if st.session_state.segmentations:
            st.divider()
            st.markdown("**Zdefiniowane segmentacje:**")
            to_del = None
            for i, seg in enumerate(st.session_state.segmentations):
                c1, c2 = st.columns([5, 1])
                c1.write(f"- `{seg['name']}` \u2014 {seg['k']} grup, bazuje na {len(seg['vars'])} zmiennych")
                if c2.button("\U0001f5d1\ufe0f", key=f"del_seg_mod_{i}"):
                    to_del = i
            if to_del is not None:
                st.session_state.segmentations.pop(to_del)
                st.rerun()

            # Show cluster profiles for each segmentation
            for seg in st.session_state.segmentations:
                if seg['name'] in df_raw.columns:
                    with st.expander(f"\U0001f3af Profil: `{seg['name']}`", expanded=False):
                        profile = df_raw.groupby(seg['name'])[seg['vars']].mean().round(2)
                        profile.index = [f"Segment {int(i)}" for i in profile.index]
                        sizes = df_raw[seg['name']].value_counts().sort_index()
                        sizes.index = [f"Segment {int(i)}" for i in sizes.index]
                        col_p1, col_p2 = st.columns([3, 1])
                        col_p1.dataframe(profile.style.format("{:.2f}"), use_container_width=True)
                        col_p2.dataframe(sizes.rename("N"), use_container_width=True)
                        fig_seg = px.bar(
                            sizes.reset_index(),
                            x='index', y=seg['name'],
                            title=f"Liczebno\u015b\u0107 segment\u00f3w: {seg['name']}",
                            color='index',
                            color_discrete_sequence=px.colors.qualitative.Set2
                        )
                        fig_seg.update_layout(showlegend=False, height=280,
                                               xaxis_title="Segment", yaxis_title="N")
                        st.plotly_chart(fig_seg, use_container_width=True, key=f"pc_seg_{seg.get('name','s')}")


elif menu == "\u2601\ufe0f Chmura S\u0142\u00f3w":
    _require_module_access("wordcloud")
    _require_data()
    module_header("\u2601\ufe0f", "Chmura S\u0142\u00f3w", "Wizualizacja odpowiedzi otwartych \u2014 eksport PNG/JPG")
    _weights_ignored_note(use_weights)

    # Check for wordcloud availability
    try:
        from wordcloud import WordCloud
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        _wc_available = True
    except ImportError:
        _wc_available = False
        st.error(
            "Biblioteka `wordcloud` nie jest zainstalowana. "
            "Uruchom: `pip install wordcloud` i restart\u01b3 aplikacj\u0119."
        )

    if _wc_available:
        # \u2500\u2500 Column selector \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
        # For SPSS: also include numeric variables that have value labels
        # (they represent coded categorical responses, not open text)
        if is_spss:
            val_label_cols = [c for c in visible_columns
                              if c in meta_orig.variable_value_labels
                              and meta_orig.variable_value_labels[c]]
            text_cols = list(dict.fromkeys(
                [c for c in visible_columns if df_raw[c].dtype == object]
                + val_label_cols
            ))
        else:
            text_cols = [c for c in visible_columns if df_raw[c].dtype == object]

        if not text_cols:
            st.warning("Brak zmiennych tekstowych w bazie danych.")
        else:
            st.info(
                "Wybierz zmienn\u0105 z odpowiedziami otwartymi, dostosuj wygl\u0105d i wygeneruj chmur\u0119 s\u0142\u00f3w. "
                "Gotow\u0105 grafik\u0119 mo\u017cesz pobra\u0107 jako PNG lub JPG."
            )

            col_cfg1, col_cfg2 = st.columns([2, 1])

            with col_cfg1:
                wc_var = st.selectbox(
                    "Pytanie otwarte (zmienna tekstowa):",
                    text_cols,
                    format_func=lambda x: get_var_display_name(x, var_labels),
                    key="wc_var"
                )

                # Stop words
                st.markdown("**Stop words** \u2014 s\u0142owa do wykluczenia:")
                default_stopwords = (
                    "i w z na do nie to a ale o jak tak si\u0119 co ten ta tego tej "
                    "jest by\u0142 by\u0142a by\u0142o s\u0105 b\u0119d\u0105 ma mam masz ma\u0107 mie\u0107 "
                    "one oni one nas nam ich im kt\u00f3ry kt\u00f3ra kt\u00f3re tego "
                    "tego tej temu tym te ten tego dla ze przy po czy "
                    "bardzo wi\u0119c jednak tylko jeszcze ju\u017c bo bo\u017c gdy "
                    "mi\u0119dzy przez mo\u017ce mo\u017cna po za przed"
                )
                stopwords_raw = st.text_area(
                    "Wpisz s\u0142owa oddzielone spacjami lub przecinkami:",
                    value=default_stopwords,
                    height=100,
                    key="wc_stopwords",
                    help="Te s\u0142owa nie pojawi\u0105 si\u0119 w chmurze. Mo\u017cesz usun\u0105\u0107 lub doda\u0107 w\u0142asne."
                )

                # Case handling
                wc_lowercase = st.checkbox(
                    "Zamie\u0144 na ma\u0142e litery (Warszawa = warszawa)",
                    value=True, key="wc_lower"
                )

                # Min word frequency
                wc_min_freq = st.slider(
                    "Minimalna cz\u0119sto\u015b\u0107 wyst\u0105pienia s\u0142owa:",
                    min_value=1, max_value=20, value=1, key="wc_minfreq"
                )

                # Max words
                wc_max_words = st.slider(
                    "Maksymalna liczba s\u0142\u00f3w w chmurze:",
                    min_value=10, max_value=300, value=100, key="wc_maxwords"
                )

            with col_cfg2:
                st.markdown("**Wygl\u0105d chmury:**")

                wc_bg = st.color_picker("Kolor t\u0142a:", value="#FFFFFF", key="wc_bg")
                wc_width  = st.number_input("Szeroko\u015b\u0107 (px):", min_value=400, max_value=3000,
                                             value=1200, step=100, key="wc_w")
                wc_height = st.number_input("Wysoko\u015b\u0107 (px):", min_value=200, max_value=2000,
                                             value=600, step=100, key="wc_h")

                PALETTES = {
                    "Niebieski (domowy)":  "Blues",
                    "Czerwony":            "Reds",
                    "Zielony":             "Greens",
                    "Fioletowy":           "Purples",
                    "Ciep\u0142e kolory":  "YlOrRd",
                    "Ch\u0142odne kolory": "cool",
                    "T\u0119czowa":        "rainbow",
                    "Czarno-bia\u0142a":   "Greys",
                    "Niebiesko-zielona":   "GnBu",
                    "Czerwono-niebieska":  "RdBu",
                }
                wc_palette_name = st.selectbox(
                    "Paleta kolor\u00f3w:",
                    list(PALETTES.keys()),
                    key="wc_palette"
                )
                wc_palette = PALETTES[wc_palette_name]

                fmt_choice = st.radio(
                    "Format pobierania:",
                    ["PNG", "JPG"],
                    horizontal=True, key="wc_fmt"
                )

            st.divider()

            if _tracked_button("\u25b6\ufe0f Generuj chmur\u0119 s\u0142\u00f3w", "wordcloud", "generate_wordcloud", type="primary",
                               use_container_width=True, key="wc_generate"):

                for _grp_lbl, _df_s, _df_raw_s, _w_s in _iter_split_groups(
                        df, df_raw, var_labels, st.session_state.split_var, weights=None):
                    if _grp_lbl:
                        st.markdown(f"### {_grp_lbl}")

                    # Build text corpus \u2014 use value labels for SPSS coded variables
                    raw_series = _df_raw_s[wc_var].dropna()

                    if is_spss:
                        spss_val_labels = {}
                        spss_val_labels.update(meta_orig.variable_value_labels.get(wc_var, {}))
                        spss_val_labels.update(st.session_state.custom_val_labels.get(wc_var, {}))

                        if spss_val_labels:
                            def _map_label(v):
                                lbl = (spss_val_labels.get(v)
                                       or spss_val_labels.get(int(v) if isinstance(v, float) and v == int(v) else v)
                                       or spss_val_labels.get(str(int(v)) if isinstance(v, float) and v == int(v) else str(v))
                                       or str(v))
                                return str(lbl)
                            texts = raw_series.map(_map_label).tolist()
                        else:
                            texts = raw_series.astype(str).tolist()
                    else:
                        texts = raw_series.astype(str).tolist()

                    corpus = " ".join(texts)
                    if wc_lowercase:
                        corpus = corpus.lower()

                    raw_sw = stopwords_raw.replace(',', ' ').split()
                    stop_set = {w.strip().lower() for w in raw_sw if w.strip()}

                    if not corpus.strip():
                        st.warning(f"Brak tekstu do analizy{' dla grupy ' + _grp_lbl if _grp_lbl else ''}.")
                        continue

                    try:
                        import matplotlib.cm as cm

                        token_re = re.compile(
                            r"[\w\u0104\u0105\u0106\u0107\u0118\u0119"
                            r"\u0141\u0142\u0143\u0144\u00d3\u00f3"
                            r"\u015a\u015b\u0179\u017a\u017b\u017c]+",
                            re.UNICODE
                        )
                        tokens = token_re.findall(corpus)

                        freq = {}
                        for tok in tokens:
                            w = tok.lower() if wc_lowercase else tok
                            if len(w) < 2:
                                continue
                            if w.lower() in stop_set:
                                continue
                            freq[w] = freq.get(w, 0) + 1

                        if wc_min_freq > 1:
                            freq = {w: c for w, c in freq.items()
                                    if c >= int(wc_min_freq)}

                        if not freq:
                            st.warning(
                                "Po zastosowaniu filtr\u00f3w nie pozosta\u0142o \u017cadne s\u0142owo"
                                + (f" dla grupy `{_grp_lbl}`" if _grp_lbl else "")
                                + ". Spr\u00f3buj zmniejszy\u0107 minimaln\u0105 cz\u0119sto\u015b\u0107 lub "
                                  "skr\u00f3ci\u0107 list\u0119 stop words."
                            )
                            continue

                        wc_obj = WordCloud(
                            width=int(wc_width),
                            height=int(wc_height),
                            background_color=wc_bg,
                            colormap=wc_palette,
                            max_words=int(wc_max_words),
                            min_font_size=8,
                            max_font_size=None,
                            min_word_length=2,
                            collocations=False,
                            relative_scaling=0.5,
                            prefer_horizontal=0.9,
                        ).generate_from_frequencies(freq)

                        fig, ax = plt.subplots(
                            figsize=(int(wc_width) / 100, int(wc_height) / 100),
                            dpi=100
                        )
                        ax.imshow(wc_obj, interpolation='bilinear')
                        ax.axis('off')
                        fig.patch.set_facecolor(wc_bg)
                        plt.tight_layout(pad=0)

                        # Save PNG + JPG bytes to session state (both formats for download)
                        _png_buf = io.BytesIO()
                        fig.savefig(_png_buf, format='png', dpi=150, bbox_inches='tight', pad_inches=0)
                        _png_buf.seek(0)
                        _jpg_buf = io.BytesIO()
                        fig.savefig(_jpg_buf, format='jpeg', dpi=150, bbox_inches='tight', pad_inches=0)
                        _jpg_buf.seek(0)
                        plt.close(fig)

                        freq_df = pd.DataFrame(
                            sorted(freq.items(), key=lambda x: x[1], reverse=True),
                            columns=["S\u0142owo", "Liczba wyst\u0105pie\u0144"]
                        )
                        freq_df.index = range(1, len(freq_df) + 1)

                        _entry = {
                            'var':          wc_var,
                            'var_label':    var_labels.get(wc_var, wc_var),
                            'group_label':  _grp_lbl,
                            'png_bytes':    _png_buf.getvalue(),
                            'jpg_bytes':    _jpg_buf.getvalue(),
                            'default_fmt':  fmt_choice.lower(),
                            'freq_df':      freq_df,
                            'n_words':      len(freq),
                        }
                        # Replace if same (var, group_label) pair already exists
                        _merge_result(st.session_state.wordcloud_results, _entry,
                            key_fn=lambda r: (r.get('var',''), r.get('group_label','')))

                    except Exception as _wc_err:
                        st.error(f"B\u0142\u0105d generowania chmury: {_wc_err}")
                        st.exception(_wc_err)

        # \u2500\u2500 Persistent display of wordcloud results with delete UI \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
        if st.session_state.wordcloud_results:
            st.divider()
            _wcc1, _wcc2 = st.columns([5, 1])
            _wcc1.markdown(f"**\u2601\ufe0f Zapisane chmury s\u0142\u00f3w ({len(st.session_state.wordcloud_results)}):**")
            with _wcc2:
                if st.button("\U0001f5d1\ufe0f Usu\u0144 wszystkie", key="del_all_wc",
                             use_container_width=True):
                    st.session_state.wordcloud_results = []
                    st.rerun()

            for _wci, _wc_entry in enumerate(list(st.session_state.wordcloud_results)):
                _wc_var   = _wc_entry.get('var', '')
                _wc_lbl   = _wc_entry.get('var_label', _wc_var)
                _wc_grp   = _wc_entry.get('group_label', '')
                _wc_n     = _wc_entry.get('n_words', 0)
                _wc_fmt   = _wc_entry.get('default_fmt', 'png')

                _title = f"\u2601\ufe0f [{_wc_var}] {_wc_lbl} \u2014 {_wc_n} s\u0142\u00f3w"
                if _wc_grp:
                    _title += f" \u2014 \U0001f500 {_wc_grp}"

                _wec1, _wec2 = st.columns([6, 1])
                with _wec1:
                    _wcexp = st.expander(_title, expanded=True)
                with _wec2:
                    if st.button("\U0001f5d1\ufe0f", key=f"del_wc_{_wci}",
                                 help=f"Usu\u0144 chmur\u0119 dla {_wc_var}"):
                        st.session_state.wordcloud_results.pop(_wci)
                        st.rerun()

                with _wcexp:
                    _split_badge(_wc_grp)

                    # Display image (use PNG for best quality in browser)
                    st.image(_wc_entry['png_bytes'], use_container_width=True)

                    # Download buttons
                    _dc1, _dc2 = st.columns(2)
                    _safe_lbl = (_wc_grp.replace('=', '_').replace(' ', '_')
                                 if _wc_grp else 'pelna')
                    with _dc1:
                        st.download_button(
                            label="\u2b07\ufe0f Pobierz (PNG)",
                            data=_wc_entry['png_bytes'],
                            file_name=f"chmura_slow_{_wc_var}_{_safe_lbl}.png",
                            mime="image/png",
                            use_container_width=True,
                            key=f"wc_dl_png_{_wci}"
                        )
                    with _dc2:
                        st.download_button(
                            label="\u2b07\ufe0f Pobierz (JPG)",
                            data=_wc_entry['jpg_bytes'],
                            file_name=f"chmura_slow_{_wc_var}_{_safe_lbl}.jpg",
                            mime="image/jpeg",
                            use_container_width=True,
                            key=f"wc_dl_jpg_{_wci}"
                        )

                    # Frequency table
                    with st.expander(
                        "\U0001f4ca Cz\u0119sto\u015b\u0107 s\u0142\u00f3w w chmurze",
                        expanded=False
                    ):
                        st.dataframe(_wc_entry['freq_df'], use_container_width=True, height=300)


# =============================================================
# MODUL: POROWNANIE FAL BADANIA (wave-over-wave)
# =============================================================
elif menu == "\U0001f30a Por\u00f3wnanie fal":
    _require_module_access("waves")
    module_header("\U0001f30a", "Por\u00f3wnanie fal",
                  "Por\u00f3wnaj wyniki cz\u0119sto\u015bci mi\u0119dzy falami badania (z zapisanych projekt\u00f3w)")

    with st.expander("Instrukcja \u2014 jak por\u00f3wnywa\u0107 fale badania", expanded=False):
        st.markdown("""
##### \U0001f3af Do czego s\u0142u\u017cy ten modu\u0142?
Pozwala zestawi\u0107 wyniki **tych samych pyta\u0144 z r\u00f3\u017cnych fal** badania (np. pomiar
kwartalny, rok do roku) w jednej tabeli, z **deltami** (r\u00f3\u017cnicami w punktach
procentowych) i **testem istotno\u015bci zmian** mi\u0119dzy kolejnymi falami.

##### \U0001f527 Jak korzysta\u0107
1. W **ka\u017cdej fali** osobno: wczytaj dane, wygeneruj **tablice cz\u0119sto\u015bci**
   (modu\u0142 *Analizy i Tabele*) i zapisz projekt do pliku `.json`
   (modu\u0142 *Projekt i S\u0142ownik* \u2014 z zaznaczon\u0105 opcj\u0105 zapisu wynik\u00f3w).
2. Tutaj wgraj **2-6 plik\u00f3w projekt\u00f3w** \u2014 po jednym na fal\u0119.
3. Nadaj **etykiety fal** (kolejno\u015b\u0107 wgrania = o\u015b czasu).
4. Wybierz **zmienne** wsp\u00f3lne dla wszystkich fal i odczytaj por\u00f3wnanie.

##### \U0001f4a1 Jak czyta\u0107 wyniki
- Kolumny fal pokazuj\u0105 **procenty** ka\u017cdej kategorii; wiersz **Baza (N)** = liczebno\u015b\u0107.
- Kolumny **\u0394** to r\u00f3\u017cnica mi\u0119dzy kolejnymi falami w **punktach procentowych (pp)**.
- **Zielony** = istotny wzrost, **czerwony** = istotny spadek (test Z dla dw\u00f3ch
  proporcji z niezale\u017cnych pr\u00f3b, poziom 95%). Gwiazdka `*` oznacza istotn\u0105 zmian\u0119.

> Por\u00f3wnywane s\u0105 wy\u0142\u0105cznie **tablice cz\u0119sto\u015bci**. Upewnij si\u0119, \u017ce w ka\u017cdej fali
> pytanie ma te same etykiety warto\u015bci \u2014 kategorie s\u0105 dopasowywane po nazwie.
        """)

    # \u2014 lokalny deserializer DataFrame z formatu projektu \u2014
    def _wave_df_from_dict(d):
        if not isinstance(d, dict) or not d.get('__df__'):
            return None
        import io as _io2
        try:
            return pd.read_json(_io2.StringIO(d['data']), orient='split')
        except Exception:
            return None

    _wv_files = st.file_uploader(
        "Wgraj pliki projekt\u00f3w (.json) \u2014 po jednym na fal\u0119 (2-6):",
        type="json", accept_multiple_files=True, key="wv_uploader")

    if _wv_files:
        # Parsuj wgrane pliki i zapisz do trwalego magazynu sesji,
        # dzieki czemu wyniki sa widoczne po powrocie z innej zakladki.
        _waves_new = []
        for _wf in _wv_files:
            try:
                _raw = json.loads(_wf.getvalue())
            except Exception:
                st.error(f"Nie mo\u017cna odczyta\u0107 pliku **{_wf.name}** (nieprawid\u0142owy JSON).")
                continue
            _cz = (_raw.get('results') or {}).get('czestosci', {}) or {}
            _freq = {}
            for _k, _v in _cz.items():
                if ' | ' in _k:   # pomijamy wyniki ze splitem (split file)
                    continue
                _df_w = _wave_df_from_dict(_v)
                if _df_w is not None and not _df_w.empty:
                    _freq[_k] = _df_w
            _meta = _raw.get('meta', {}) or {}
            _def_lbl = (_meta.get('name') or _wf.name.rsplit('.', 1)[0]).strip()
            _waves_new.append({
                'name': _wf.name, 'default': _def_lbl, 'freq': _freq,
                'cvl': _raw.get('custom_var_labels', {}) or {},
                'saved_at': _raw.get('_saved_at', ''),
            })
        _no_freq = [_w['name'] for _w in _waves_new if not _w['freq']]
        if _no_freq:
            st.warning("Pliki bez zapisanych tablic cz\u0119sto\u015bci (pomini\u0119te): "
                       + ", ".join(_no_freq)
                       + ". Wygeneruj cz\u0119sto\u015bci w danym projekcie i zapisz go z wynikami.")
        st.session_state['wv_store'] = [_w for _w in _waves_new if _w['freq']]

    # Wczytaj fale z magazynu sesji (przetrwaly nawigacje miedzy zakladkami).
    _waves = st.session_state.get('wv_store', []) or []

    if not _waves:
        st.info("Wgraj co najmniej **2 pliki projekt\u00f3w** zawieraj\u0105ce zapisane tablice "
                "cz\u0119sto\u015bci, aby zobaczy\u0107 por\u00f3wnanie fal.")
    else:
        _wv_src1, _wv_src2 = st.columns([5, 1])
        _wv_src1.caption(
            f"Wczytane fale: **{len(_waves)}** "
            + ("(z wgranych plik\u00f3w)" if _wv_files else "(z pami\u0119ci sesji)"))
        with _wv_src2:
            if st.button("\U0001f5d1\ufe0f Wyczy\u015b\u0107 fale", key="wv_clear_store",
                         use_container_width=True,
                         help="Usu\u0144 wgrane fale i wyniki z pami\u0119ci sesji."):
                st.session_state.pop('wv_store', None)
                st.session_state.pop('wv_results', None)
                st.rerun()

        if len(_waves) < 2:
            st.info("Potrzebne s\u0105 co najmniej **2 fale** z zapisanymi tablicami cz\u0119sto\u015bci.")
        else:
            st.markdown("##### Etykiety fal (kolejno\u015b\u0107 = o\u015b czasu)")
            _wv_labels_raw = []
            _lbl_cols = st.columns(len(_waves))
            for _i, _w in enumerate(_waves):
                with _lbl_cols[_i]:
                    _lab = st.text_input(f"Fala {_i + 1}", value=_w['default'],
                                         key=f"wv_lbl_{_i}")
                    _cap_dt = (_w.get('saved_at') or '')[:10]
                    if _cap_dt:
                        st.caption(f"zapis: {_cap_dt}")
                    _wv_labels_raw.append((_lab or f"Fala {_i + 1}").strip()
                                          or f"Fala {_i + 1}")
            # etykiety musza byc unikalne (build_wave_comparison indeksuje kolumny po nazwie)
            _seen_lbl, _wv_labels = {}, []
            for _l in _wv_labels_raw:
                if _l in _seen_lbl:
                    _seen_lbl[_l] += 1
                    _wv_labels.append(f"{_l} ({_seen_lbl[_l]})")
                else:
                    _seen_lbl[_l] = 1
                    _wv_labels.append(_l)

            # mapa etykiet zmiennych (globalna: pierwsza fala, ktora ma etykiete)
            _cvl_all = {}
            for _w in _waves:
                for _k, _v in (_w['cvl'] or {}).items():
                    _cvl_all.setdefault(_k, _v)

            def _wv_disp(k):
                return _cvl_all.get(k, k)

            # \u2014 fala referencyjna: z niej wybieramy liste zmiennych \u2014
            _ref_i = st.selectbox(
                "Fala referencyjna (z niej wybierasz zmienne):",
                list(range(len(_waves))),
                format_func=lambda i: _wv_labels[i], key="wv_ref")
            _ref_wave = _waves[_ref_i]
            _ref_cvl = _ref_wave['cvl'] or {}
            _ref_vars = sorted(_ref_wave['freq'].keys())

            def _ref_disp(k):
                return _ref_cvl.get(k, _wv_disp(k))

            # auto-dopasowanie zmiennej z fali ref. do innej fali:
            # 1) identyczny klucz, 2) ta sama etykieta, 3) klucz bez wzgledu na wielkosc liter
            def _wv_automatch(var_key, wave):
                _fr = wave['freq']
                if var_key in _fr:
                    return var_key
                _want = _ref_disp(var_key)
                for _k in _fr:
                    if (wave['cvl'] or {}).get(_k, _k) == _want:
                        return _k
                _low = var_key.lower()
                for _k in _fr:
                    if _k.lower() == _low:
                        return _k
                return None

            if not _ref_vars:
                st.error("Fala referencyjna nie zawiera tablic cz\u0119sto\u015bci. "
                         "Wybierz inn\u0105 fal\u0119 referencyjn\u0105.")
            else:
                _sel = st.multiselect("Zmienne do por\u00f3wnania (z fali referencyjnej):",
                                      _ref_vars, format_func=_ref_disp, key="wv_vars")
                _oc1, _oc2, _oc3 = st.columns(3)
                _wv_show_sig = _oc1.checkbox("Oznacz istotne zmiany (95%)", value=True,
                                             key="wv_sig")
                _wv_shown = _oc2.checkbox("Poka\u017c liczebno\u015bci (N)", value=True,
                                          key="wv_shown")
                _wv_chart = _oc3.checkbox("Poka\u017c wykres trendu", value=True,
                                          key="wv_chart")

                # \u2014 mapowanie nazw zmiennych miedzy falami (gdy klucze sie roznia) \u2014
                _resolved = {}
                _other_idx = [_i for _i in range(len(_waves)) if _i != _ref_i]
                if _sel:
                    _need_attn = any(
                        _wv_automatch(_v, _waves[_wi]) is None
                        for _v in _sel for _wi in _other_idx)
                    _exp_lbl = ("\U0001f517 Mapowanie nazw zmiennych mi\u0119dzy falami"
                                + (" \u2014 wymaga uwagi" if _need_attn
                                   else " (auto-dopasowane)"))
                    with st.expander(_exp_lbl, expanded=_need_attn):
                        st.caption(
                            "Dla ka\u017cdej fali wska\u017c zmienn\u0105 odpowiadaj\u0105c\u0105 tej z fali "
                            "referencyjnej. Auto-dopasowanie: identyczna nazwa \u2192 ta sama "
                            "etykieta \u2192 nazwa bez wzgl\u0119du na wielko\u015b\u0107 liter. Wybierz "
                            "\"pomi\u0144\", aby wykluczy\u0107 fal\u0119 dla danej zmiennej.")
                        for _v in _sel:
                            _res = [None] * len(_waves)
                            _res[_ref_i] = _v
                            st.markdown("**" + _ref_disp(_v) + "**")
                            _mcols = st.columns(len(_other_idx))
                            for _ci, _wi in enumerate(_other_idx):
                                _wave = _waves[_wi]
                                _opts = [None] + sorted(_wave['freq'].keys())
                                _auto = _wv_automatch(_v, _wave)
                                _didx = _opts.index(_auto) if _auto in _opts else 0
                                with _mcols[_ci]:
                                    _pick = st.selectbox(
                                        _wv_labels[_wi], _opts, index=_didx,
                                        format_func=lambda k, _wj=_wi: (
                                            "\u2014 pomi\u0144 t\u0119 fal\u0119 \u2014" if k is None
                                            else ((_waves[_wj]['cvl'] or {}).get(k, k))),
                                        key=f"wv_map_{_ref_i}_{_v}_{_wi}")
                                    if _auto is None and _pick is None:
                                        st.caption("\u26a0\ufe0f brak dopasowania")
                                    _res[_wi] = _pick
                            _resolved[_v] = _res

                # \u2014 renderer HTML tabeli porownania \u2014
                def _render_wave_html(comp, show_n=False, show_sig=True):
                    from html import escape as _esc
                    import math as _m
                    _pct = comp['pct']; _base = comp['base']
                    _n = comp['n']
                    _delta = comp['delta']; _sig = comp['sig']
                    _wvs = list(_pct.columns)
                    _prs = comp['pair_labels']
                    style = (
                        '<style>'
                        '.wv-tbl{border-collapse:collapse;font-size:12px;width:100%;}'
                        '.wv-tbl th,.wv-tbl td{border:1px solid #B0C4DE;padding:4px 8px;white-space:nowrap;}'
                        '.wv-tbl .wv-h{background:#1F4E79;color:#fff;font-weight:bold;text-align:center;}'
                        '.wv-tbl .wv-hd{background:#2E75B6;color:#fff;font-weight:bold;text-align:center;}'
                        '.wv-tbl .idx{text-align:left;background:#F2F2F2;font-weight:bold;min-width:160px;}'
                        '.wv-tbl td.dat{text-align:center;}'
                        '.wv-tbl tr.base td{background:#D6E4F0;font-weight:bold;}'
                        '.wv-tbl td.wv-up{background:#E2F0D9;color:#548235;font-weight:bold;text-align:center;}'
                        '.wv-tbl td.wv-dn{background:#FCE4E4;color:#C00000;font-weight:bold;text-align:center;}'
                        '.wv-tbl td.wv-fl{color:#999;text-align:center;}'
                        '.wv-tbl div.wv-n{font-size:10px;color:#777;font-weight:normal;}'
                        '</style>'
                    )
                    hdr = '<tr><th class="wv-h idx">Kategoria</th>'
                    for _wv in _wvs:
                        hdr += '<th class="wv-h">' + _esc(str(_wv)) + '</th>'
                    for _pr in _prs:
                        hdr += '<th class="wv-hd">\u0394 ' + _esc(str(_pr)) + '</th>'
                    hdr += '</tr>'
                    body = ''
                    for _ri, _cat in enumerate(list(_pct.index)):
                        _trc = 'even' if _ri % 2 == 0 else 'odd'
                        body += '<tr class="' + _trc + '"><td class="idx">' + _esc(str(_cat)) + '</td>'
                        for _wv in _wvs:
                            _v = _pct.loc[_cat, _wv]
                            if _v is None or (isinstance(_v, float) and _m.isnan(_v)):
                                body += '<td class="wv-fl">\u2014</td>'
                            else:
                                _cell = '%.1f%%' % float(_v)
                                if show_n:
                                    _nv = _n.loc[_cat, _wv]
                                    if not (_nv is None or (isinstance(_nv, float)
                                                            and _m.isnan(_nv))):
                                        _cell += ('<div class="wv-n">(N='
                                                  + str(int(round(float(_nv)))) + ')</div>')
                                body += '<td class="dat">' + _cell + '</td>'
                        for _pr in _prs:
                            _d = _delta.loc[_cat, _pr] if _pr in _delta.columns else None
                            if _d is None or (isinstance(_d, float) and _m.isnan(_d)):
                                body += '<td class="wv-fl">\u2014</td>'
                            else:
                                _sv = (int(_sig.loc[_cat, _pr])
                                       if (show_sig and _pr in _sig.columns) else 0)
                                _cls = 'wv-up' if _sv > 0 else ('wv-dn' if _sv < 0 else 'dat')
                                _txt = '%+.1f pp' % float(_d)
                                if _sv != 0:
                                    _txt += ' *'
                                body += '<td class="' + _cls + '">' + _txt + '</td>'
                        body += '</tr>'
                    body += '<tr class="base"><td class="idx">Baza (N)</td>'
                    for _wv in _wvs:
                        _b = _base.get(_wv)
                        if _b is None or (isinstance(_b, float) and _m.isnan(_b)):
                            body += '<td class="dat"></td>'
                        else:
                            body += '<td class="dat">' + str(int(round(float(_b)))) + '</td>'
                    for _pr in _prs:
                        body += '<td class="wv-fl"></td>'
                    body += '</tr>'
                    # wiersz "Suma %" \u2014 suma procentow w kolumnie (bez box-setow '[')
                    _nonbox = [_c for _c in _pct.index
                               if not str(_c).startswith('[')]
                    body += '<tr class="base"><td class="idx">Suma %</td>'
                    for _wv in _wvs:
                        _sp = _pct[_wv].reindex(_nonbox).dropna()
                        if _sp.empty:
                            body += '<td class="dat"></td>'
                        else:
                            body += '<td class="dat">%.1f%%</td>' % float(_sp.sum())
                    for _pr in _prs:
                        body += '<td class="wv-fl"></td>'
                    body += '</tr>'
                    return (style + '<div style="overflow-x:auto"><table class="wv-tbl">'
                            + '<thead>' + hdr + '</thead><tbody>' + body
                            + '</tbody></table></div>')

                # \u2014 generowanie: wyniki zapisujemy w session_state, dzieki czemu
                # sa widoczne po powrocie z innej zakladki (nie gina jak stan widgetow) \u2014
                _gc1, _gc2 = st.columns([3, 1])
                with _gc1:
                    if st.button("\u25b6\ufe0f Generuj por\u00f3wnania", type="primary",
                                 key="wv_generate", disabled=not _sel):
                        _wvres = dict(st.session_state.get('wv_results', {}))
                        for _var in _sel:
                            _rk = _resolved.get(_var, [None] * len(_waves))
                            _freqs = [
                                _waves[_i]['freq'].get(_rk[_i]) if _rk[_i] else None
                                for _i in range(len(_waves))
                            ]
                            try:
                                comp = build_wave_comparison(_freqs, _wv_labels, do_sig=True)
                            except Exception as _e:
                                st.error(f"B\u0142\u0105d por\u00f3wnania dla {_ref_disp(_var)}: {_e}")
                                continue
                            _wvres[_var] = {'disp': _ref_disp(_var), 'comp': comp,
                                            'labels': list(_wv_labels)}
                        st.session_state['wv_results'] = _wvres
                        st.rerun()
                with _gc2:
                    if st.session_state.get('wv_results') and st.button(
                            "\U0001f5d1\ufe0f Usu\u0144 wszystkie", key="wv_del_all",
                            use_container_width=True):
                        st.session_state.pop('wv_results', None)
                        st.rerun()

                # \u2014 wyswietlanie zapisanych wynikow (przetrwaly nawigacje) \u2014
                _wvres = st.session_state.get('wv_results', {})
                if _wvres:
                    st.divider()
                    st.markdown(f"**Zapisane por\u00f3wnania ({len(_wvres)}):**")
                    for _var, _entry in list(_wvres.items()):
                        comp = _entry['comp']
                        _disp = _entry.get('disp', _var)
                        _wc1, _wc2 = st.columns([6, 1])
                        with _wc1:
                            _wexp = st.expander(_disp, expanded=True)
                        with _wc2:
                            if st.button("\U0001f5d1\ufe0f", key=f"wv_del_{_var}",
                                         help=f"Usu\u0144 por\u00f3wnanie dla {_disp}"):
                                _wvres.pop(_var, None)
                                st.session_state['wv_results'] = _wvres
                                st.rerun()
                        with _wexp:
                            st.markdown(
                                _render_wave_html(comp, show_n=_wv_shown,
                                                  show_sig=_wv_show_sig),
                                unsafe_allow_html=True)
                            if _wv_show_sig:
                                st.caption("Zielony = istotny wzrost, czerwony = istotny spadek "
                                           "(test Z dla dw\u00f3ch proporcji, 95%). pp = punkty procentowe.")
                            if _wv_chart:
                                _pc = comp['pct']
                                _plot = _pc[~_pc.index.astype(str).str.startswith('[')]
                                _plot = _plot.dropna(how='all')
                                if not _plot.empty:
                                    _long = _plot.reset_index().melt(
                                        id_vars=_plot.index.name or 'index',
                                        var_name='Fala', value_name='Procent')
                                    _long = _long.rename(
                                        columns={_long.columns[0]: 'Kategoria'})
                                    _figw = px.line(
                                        _long, x='Fala', y='Procent', color='Kategoria',
                                        markers=True, title=_disp)
                                    _figw.update_layout(
                                        height=380, yaxis_title='%',
                                        yaxis_ticksuffix='%')
                                    st.plotly_chart(_figw, use_container_width=True,
                                                    key=f"pc_wave_{_var}")

                    # \u2014 eksport do Excela (gotowy gdy sa wyniki) \u2014
                    if _wvres:
                        st.divider()
                        _wv_buf = io.BytesIO()
                        try:
                            with pd.ExcelWriter(_wv_buf, engine='xlsxwriter') as _wvwr:
                                _wb = _wvwr.book
                                _f_hdr = _wb.add_format({'bold': True, 'bg_color': '#1F4E79',
                                    'font_color': '#FFFFFF', 'border': 1, 'align': 'center',
                                    'valign': 'vcenter', 'text_wrap': True})
                                _f_hd2 = _wb.add_format({'bold': True, 'bg_color': '#2E75B6',
                                    'font_color': '#FFFFFF', 'border': 1, 'align': 'center',
                                    'valign': 'vcenter', 'text_wrap': True})
                                _f_idx = _wb.add_format({'bold': True, 'bg_color': '#F2F2F2',
                                    'border': 1})
                                _f_pct = _wb.add_format({'border': 1, 'align': 'center',
                                    'num_format': '0.0"%"'})
                                _f_up = _wb.add_format({'border': 1, 'align': 'center',
                                    'bg_color': '#E2F0D9', 'font_color': '#548235', 'bold': True,
                                    'num_format': '+0.0" pp";-0.0" pp"'})
                                _f_dn = _wb.add_format({'border': 1, 'align': 'center',
                                    'bg_color': '#FCE4E4', 'font_color': '#C00000', 'bold': True,
                                    'num_format': '+0.0" pp";-0.0" pp"'})
                                _f_neu = _wb.add_format({'border': 1, 'align': 'center',
                                    'num_format': '+0.0" pp";-0.0" pp"'})
                                _f_base = _wb.add_format({'bold': True, 'bg_color': '#D6E4F0',
                                    'border': 1, 'align': 'center', 'num_format': '0'})
                                _f_bpct = _wb.add_format({'bold': True, 'bg_color': '#D6E4F0',
                                    'border': 1, 'align': 'center', 'num_format': '0.0"%"'})
                                _f_bidx = _wb.add_format({'bold': True, 'bg_color': '#D6E4F0',
                                    'border': 1})
                                _f_em = _wb.add_format({'border': 1, 'align': 'center'})
                                _f_n = _wb.add_format({'border': 1, 'align': 'center',
                                    'num_format': '0'})
                                _used_names = set()

                                def _wv_sheet_name(_nm):
                                    for _bad in '[]:*?/\\':
                                        _nm = _nm.replace(_bad, ' ')
                                    _nm = (_nm.strip() or 'Zmienna')[:28]
                                    _base_nm = _nm
                                    _k = 2
                                    while _nm.lower() in _used_names:
                                        _nm = f"{_base_nm[:25]} {_k}"
                                        _k += 1
                                    _used_names.add(_nm.lower())
                                    return _nm

                                for _var, _entry in _wvres.items():
                                    comp = _entry['comp']
                                    _vdisp = _entry.get('disp', _var)
                                    _ws = _wb.add_worksheet(_wv_sheet_name(_vdisp))
                                    _ws.set_column(0, 0, 28)
                                    _pct = comp['pct']; _ncomp = comp['n']
                                    _base = comp['base']
                                    _delta = comp['delta']; _sig = comp['sig']
                                    _wvs = list(_pct.columns)
                                    _prs = comp['pair_labels']
                                    _stride = 2 if _wv_shown else 1
                                    _dcol0 = 1 + len(_wvs) * _stride
                                    _last_col = _dcol0 + len(_prs) - 1
                                    _ws.set_column(1, max(1, _last_col), 13)
                                    _ws.merge_range(0, 0, 0, _last_col,
                                                    _vdisp, _f_hdr)
                                    _r0 = 1
                                    _ws.write(_r0, 0, 'Kategoria', _f_hdr)
                                    for _ci, _wv in enumerate(_wvs):
                                        _c = 1 + _ci * _stride
                                        if _wv_shown:
                                            _ws.write(_r0, _c, str(_wv) + ' [%]', _f_hdr)
                                            _ws.write(_r0, _c + 1, str(_wv) + ' [N]', _f_hdr)
                                        else:
                                            _ws.write(_r0, _c, str(_wv), _f_hdr)
                                    for _pi, _pr in enumerate(_prs):
                                        _ws.write(_r0, _dcol0 + _pi,
                                                  '\u0394 ' + str(_pr), _f_hd2)
                                    _rr = _r0 + 1
                                    for _cat in list(_pct.index):
                                        _ws.write(_rr, 0, str(_cat), _f_idx)
                                        for _ci, _wv in enumerate(_wvs):
                                            _c = 1 + _ci * _stride
                                            _v = _pct.loc[_cat, _wv]
                                            if _v is None or (isinstance(_v, float)
                                                              and _v != _v):
                                                _ws.write_blank(_rr, _c, None, _f_em)
                                            else:
                                                _ws.write_number(_rr, _c, float(_v), _f_pct)
                                            if _wv_shown:
                                                _nv = _ncomp.loc[_cat, _wv]
                                                if _nv is None or (isinstance(_nv, float)
                                                                   and _nv != _nv):
                                                    _ws.write_blank(_rr, _c + 1, None, _f_em)
                                                else:
                                                    _ws.write_number(_rr, _c + 1,
                                                                     int(round(float(_nv))),
                                                                     _f_n)
                                        for _pi, _pr in enumerate(_prs):
                                            _col = _dcol0 + _pi
                                            _d = (_delta.loc[_cat, _pr]
                                                  if _pr in _delta.columns else None)
                                            if _d is None or (isinstance(_d, float)
                                                              and _d != _d):
                                                _ws.write_blank(_rr, _col, None, _f_em)
                                            else:
                                                _sv = (int(_sig.loc[_cat, _pr])
                                                       if (_wv_show_sig
                                                           and _pr in _sig.columns) else 0)
                                                _fmt = (_f_up if _sv > 0
                                                        else (_f_dn if _sv < 0 else _f_neu))
                                                _ws.write_number(_rr, _col, float(_d), _fmt)
                                        _rr += 1
                                    # wiersz bazy
                                    _ws.write(_rr, 0, 'Baza (N)', _f_bidx)
                                    for _ci, _wv in enumerate(_wvs):
                                        _c = 1 + _ci * _stride
                                        _b = _base.get(_wv)
                                        if _b is None or (isinstance(_b, float) and _b != _b):
                                            _ws.write_blank(_rr, _c, None, _f_base)
                                        else:
                                            _ws.write_number(_rr, _c,
                                                             int(round(float(_b))), _f_base)
                                        if _wv_shown:
                                            _ws.write_blank(_rr, _c + 1, None, _f_base)
                                    for _pi in range(len(_prs)):
                                        _ws.write_blank(_rr, _dcol0 + _pi, None, _f_base)
                                    # wiersz "Suma %" (bez box-setow '[')
                                    _rr += 1
                                    _nonbox = [_c for _c in _pct.index
                                               if not str(_c).startswith('[')]
                                    _ws.write(_rr, 0, 'Suma %', _f_bidx)
                                    for _ci, _wv in enumerate(_wvs):
                                        _c = 1 + _ci * _stride
                                        _sp = _pct[_wv].reindex(_nonbox).dropna()
                                        if _sp.empty:
                                            _ws.write_blank(_rr, _c, None, _f_bpct)
                                        else:
                                            _ws.write_number(_rr, _c,
                                                             float(_sp.sum()), _f_bpct)
                                        if _wv_shown:
                                            _ws.write_blank(_rr, _c + 1, None, _f_base)
                                    for _pi in range(len(_prs)):
                                        _ws.write_blank(_rr, _dcol0 + _pi, None, _f_base)
                            st.download_button(
                                "\u2b07\ufe0f Pobierz por\u00f3wnanie fal (Excel)",
                                data=_wv_buf.getvalue(),
                                file_name="Porownanie_fal.xlsx",
                                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                                use_container_width=True, key="wv_xls_dl")
                        except Exception as _e:
                            st.error(f"B\u0142\u0105d generowania pliku Excel: {_e}")


elif menu == "\U0001f4be Eksport do Excela":
    _require_module_access("export_excel")
    _require_data()
    module_header("\U0001f4be", "Eksport do Excela", "Raport analityczny, wykresy, baza danych, spis tre\u015bci")

    # \u2500\u2500 Separate standalone DB download (still available) \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
    with st.expander("Osobny plik z baz\u0105 danych", expanded=False):
        col1, col2 = st.columns(2)
        col1.info("**Baza z etykietami** -- warto\u015bci kod\u00f3w zast\u0105pione tekstem (np. 1 \u2192 'Kobieta').")
        col2.info("**Baza surowa** -- oryginalne warto\u015bci liczbowe. Wiersz 1: nazwy, Wiersz 2: etykiety.")

        db_header_style = st.radio(
            "Nag\u0142\u00f3wki kolumn:",
            ["Nazwy zmiennych", "Etykiety zmiennych"],
            key="db_header_style", horizontal=True,
            help="Wybierz czy w pierwszym wierszu maj\u0105 znale\u017a\u0107 si\u0119 nazwy zmiennych "
                 "(np. `Q1_1`) czy ich etykiety (np. `Jak oceniasz obs\u0142ug\u0119`)."
        )

        if st.button("\U0001f4e5 Pobierz osobny plik z baz\u0105 danych", use_container_width=True):
            with st.spinner("Generowanie..."):
                db_data = export_db_to_excel(df_raw, df, var_labels,
                                              header_mode=("labels" if db_header_style == "Etykiety zmiennych" else "names"))
            fname = "Baza_Danych_Excel.xlsx" if is_excel else "Baza_Danych_SPSS.xlsx"
            st.download_button("\u2b07\ufe0f Pobierz " + fname, data=db_data,
                               file_name=fname,
                               mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                               use_container_width=True)

    st.divider()
    st.markdown("### Raport analityczny (Excel)")

    any_results = (
        any(st.session_state.results.get(g) for g in ['czestosci', 'krzyzowe', 'srednie', 'opisowe', 'korelacje']) or
        bool(st.session_state.regression_results) or
        bool(st.session_state.anova_results) or
        bool(st.session_state.factor_results) or
        bool(st.session_state.matrix_results) or
        bool(st.session_state.conjoint_results) or
        bool(st.session_state.maxdiff_results)
    )
    if not any_results:
        st.warning("Brak wynik\u00f3w do eksportu. Przejd\u017a do modu\u0142\u00f3w analitycznych i wygeneruj tabele.")
    else:
        # Summary of available results
        available = []
        for grp, name in [('czestosci', 'Cz\u0119sto\u015bci'), ('krzyzowe', 'Krzy\u017cowe'), ('srednie', '\u015arednie'),
                          ('opisowe', 'Opisowe'), ('korelacje', 'Korelacje')]:
            if st.session_state.results.get(grp):
                available.append(f"\u2705 {name} ({len(st.session_state.results[grp])} tabel)")
        if st.session_state.matrix_results:
            available.append(f"\u2705 Pytania matrycowe ({len(st.session_state.matrix_results)} pyta\u0144)")
        if st.session_state.regression_results:
            available.append(f"\u2705 Regresja OLS ({len([r for r in st.session_state.regression_results if 'error' not in r])} blok\u00f3w)")
        if st.session_state.anova_results:
            available.append(f"\u2705 ANOVA ({len(st.session_state.anova_results)} analiz)")
        if st.session_state.factor_results:
            available.append(f"\u2705 Analiza Czynnikowa ({len(st.session_state.factor_results)} analiz)")
        if st.session_state.conjoint_results:
            available.append(f"\u2705 Conjoint ({len(st.session_state.conjoint_results)} analiz)")
        if st.session_state.maxdiff_results:
            available.append(f"\u2705 MaxDiff ({len(st.session_state.maxdiff_results)} analiz)")
        st.success("**Gotowe do eksportu:**\n" + " \u00b7 ".join(available))

        st.markdown("**Opcje eksportu:**")
        opt_col1, opt_col2, opt_col3 = st.columns(3)

        # Chart option
        has_freq = bool(st.session_state.results.get('czestosci'))
        add_freq_charts = opt_col1.checkbox(
            "\U0001f4ca Wykresy do tabel cz\u0119sto\u015bci",
            value=False,
            key="export_add_charts",
            help="Wstawia natywne wykresy Excela (edytowalne) obok ka\u017cdej tabeli cz\u0119sto\u015bci.",
            disabled=not has_freq,
        )

        # DB options
        incl_db_labeled = opt_col2.checkbox(
            "\U0001f4c2 Baza danych z etykietami",
            value=False,
            key="export_db_labeled",
            help="Dodaje do pliku arkusz 'Baza z etykietami' (warto\u015bci tekstowe) zaraz po Spisie Tre\u015bci.",
        )
        incl_db_raw = opt_col3.checkbox(
            "\U0001f4cb Baza danych surowa",
            value=False,
            key="export_db_raw",
            help="Dodaje arkusz 'Baza surowa (numeryczna)' z oryginalnymi kodami liczbowymi.",
        )


        # -- ukrywanie pustych kategorii (tylko SPSS) --
        if is_spss:
            excel_drop_empty = st.checkbox(
                "\U0001f9f9 Ukryj kategorie bez odpowiedzi (N = 0)",
                value=False, key="excel_drop_empty",
                help="Dla danych SPSS: pomija w tabelach kategorie ze s\u0142ownika "
                     "warto\u015bci, kt\u00f3rych nikt nie wybra\u0142 (liczebno\u015b\u0107 0).")
        else:
            excel_drop_empty = False

        # -- szczeg\u00f3\u0142owy wyb\u00f3r wynik\u00f3w do eksportu --
        with st.expander("Szczeg\u00f3\u0142owy wyb\u00f3r wynik\u00f3w do eksportu", expanded=False):
            _xl_cats = []
            for _gk, _gt in [('czestosci', 'Tablice cz\u0119sto\u015bci'),
                             ('krzyzowe', 'Tabele krzy\u017cowe'),
                             ('srednie', 'Tabele \u015brednich'),
                             ('opisowe', 'Statystyki opisowe'),
                             ('korelacje', 'Korelacje'),
                             ('banner', 'Tabele zbiorcze (Banner)')]:
                _opts = list(st.session_state.results.get(_gk, {}).keys())
                if _opts:
                    _xl_cats.append((_gk, _gt, _opts,
                                     _gk in ('czestosci', 'krzyzowe', 'srednie', 'opisowe', 'banner')))
            if st.session_state.matrix_results:
                _xl_cats.append(('matrix', 'Tabele matrycowe',
                                 [e['name'] for e in st.session_state.matrix_results], True))
            _xl_reg0 = [r for r in st.session_state.regression_results if 'error' not in r]
            if _xl_reg0:
                _xl_cats.append(('regression', 'Regresja OLS',
                                 ['Zm. zal.: ' + r.get('dep_var', 'Wynik ' + str(i + 1))
                                  for i, r in enumerate(_xl_reg0)], False))
            if st.session_state.anova_results:
                _xl_cats.append(('anova', 'ANOVA',
                                 [e.get('dep_var', '?') + ' wg ' + e.get('group_var', '?')
                                  for e in st.session_state.anova_results], False))
            if st.session_state.factor_results:
                _xl_cats.append(('factor', 'Analiza czynnikowa',
                                 [e.get('title', 'EFA ' + str(i + 1))
                                  for i, e in enumerate(st.session_state.factor_results)], False))
            if st.session_state.conjoint_results:
                _xl_cats.append(('conjoint', 'Conjoint',
                                 [e.get('title', 'Conjoint ' + str(i + 1))
                                  for i, e in enumerate(st.session_state.conjoint_results)], False))
            if st.session_state.maxdiff_results:
                _xl_cats.append(('maxdiff', 'MaxDiff',
                                 [e.get('title', 'MaxDiff ' + str(i + 1))
                                  for i, e in enumerate(st.session_state.maxdiff_results)], False))
            render_granular_selector('excel', _xl_cats, var_labels)

        # -- kolejnosc wynikow w eksporcie --
        _xl_ro_cats = []
        for _rk, _rn in [('czestosci', 'Tablice cz\u0119sto\u015bci'),
                         ('krzyzowe',  'Tabele krzy\u017cowe'),
                         ('matrix',    'Tabele matrycowe'),
                         ('srednie',   'Tabele \u015brednich'),
                         ('opisowe',   'Statystyki opisowe'),
                         ('korelacje', 'Korelacje'),
                         ('banner',    'Tabele zbiorcze (Banner)')]:
            _conf_xl = st.session_state.get('excel_sel_confirmed', {})
            _keys_xl = list(_conf_xl.get(_rk, []))
            if not _keys_xl:
                if _rk == 'matrix':
                    _keys_xl = [e['name'] for e in st.session_state.matrix_results]
                else:
                    _keys_xl = list(st.session_state.results.get(_rk, {}).keys())
            if _keys_xl:
                _xl_ro_cats.append((_rk, _rn, _keys_xl))
        _xl_order = render_reorder_ui('excel', _xl_ro_cats, var_labels)

        if _tracked_button("\U0001f4ca Generuj pe\u0142ny raport analityczny", "export_excel", "generate_report", type="primary", use_container_width=True):
            with st.spinner("Generowanie pliku Excel... To mo\u017ce chwil\u0119 potrwa\u0107."):
                output = io.BytesIO()
                try:
                    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
                        # 1. ToC -- always first tab
                        toc_ws = writer.book.add_worksheet('Spis Tre\u015bci')
                        toc_ws.set_tab_color('#1F4E79')

                        # 2. DB sheets -- right after ToC (second/third position)
                        if incl_db_labeled:
                            write_db_sheet(writer, 'Baza z etykietami',
                                           df, var_labels, '#1F4E79')
                        if incl_db_raw:
                            write_db_sheet(writer, 'Baza surowa (numeryczna)',
                                           df_raw, var_labels, '#2E75B6')

                        # -- filtr wyboru + kolejnosc + transformacja tabel --
                        _xl_keep = export_selection_filter('excel')
                        _xl_ord  = st.session_state.get('excel_result_order', {})

                        def _xl_prep(_grp):
                            _d = st.session_state.results.get(_grp, {})
                            if _grp in _xl_ord and _xl_ord[_grp]:
                                return {k: prepare_export_table(_d[k], drop_empty_cats=excel_drop_empty)
                                        for k in _xl_ord[_grp] if k in _d}
                            return {k: prepare_export_table(v, drop_empty_cats=excel_drop_empty)
                                    for k, v in _d.items() if _xl_keep(_grp, k)}
                        _xl_res = {_g: _xl_prep(_g) for _g in
                                   ['czestosci', 'krzyzowe', 'srednie', 'opisowe', 'korelacje', 'banner']}
                        _xl_mat_src = st.session_state.matrix_results
                        if 'matrix' in _xl_ord and _xl_ord['matrix']:
                            _mord = {k: i for i, k in enumerate(_xl_ord['matrix'])}
                            _xl_mtx = sorted(
                                [dict(e, df=prepare_export_table(e['df'], drop_empty_cats=excel_drop_empty))
                                 for e in _xl_mat_src if e['name'] in _mord],
                                key=lambda e: _mord.get(e['name'], 9999))
                        else:
                            _xl_mtx = [dict(e, df=prepare_export_table(e['df'], drop_empty_cats=excel_drop_empty))
                                       for e in _xl_mat_src if _xl_keep('matrix', e['name'])]

                        # 3. Analytical results sheets
                        sheet_map = {}
                        for grp, s_name in [('czestosci', 'Cz\u0119sto\u015bci'), ('krzyzowe', 'Krzy\u017cowe'),
                                            ('srednie', '\u015arednie'), ('opisowe', 'Opisowe'),
                                            ('korelacje', 'Korelacje'), ('banner', 'Banner')]:
                            if _xl_res.get(grp):
                                _charts = add_freq_charts if grp == 'czestosci' else False
                                row_map = export_tables_to_sheet(
                                    writer, s_name, _xl_res[grp], var_labels,
                                    add_charts=_charts
                                )
                                sheet_map[s_name] = row_map

                        if _xl_mtx:
                            export_matrix_to_excel(writer, _xl_mtx, var_labels)

                        _reg_all = [r for r in st.session_state.regression_results if 'error' not in r]
                        valid_reg = [r for i, r in enumerate(_reg_all)
                                     if _xl_keep('regression', 'Zm. zal.: ' + r.get('dep_var', 'Wynik ' + str(i + 1)))]
                        if valid_reg:
                            export_regression_to_excel(writer, valid_reg, var_labels)
                        _anova_sel = [e for e in st.session_state.anova_results
                                      if _xl_keep('anova', e.get('dep_var', '?') + ' wg ' + e.get('group_var', '?'))]
                        if _anova_sel:
                            export_anova_to_excel(writer, _anova_sel, var_labels)
                        _factor_sel = [e for i, e in enumerate(st.session_state.factor_results)
                                       if _xl_keep('factor', e.get('title', 'EFA ' + str(i + 1)))]
                        if _factor_sel:
                            export_factor_to_excel(writer, _factor_sel, var_labels)
                        _conj_sel = [e for i, e in enumerate(st.session_state.conjoint_results)
                                     if not e.get('error')
                                     and _xl_keep('conjoint', e.get('title', 'Conjoint ' + str(i + 1)))]
                        if _conj_sel:
                            export_conjoint_to_excel(
                                writer, _conj_sel, var_labels,
                                meta_vvl=meta_orig.variable_value_labels if is_spss else {},
                                custom_val_labels=st.session_state.custom_val_labels
                            )
                        _md_sel = [e for i, e in enumerate(st.session_state.maxdiff_results)
                                   if _xl_keep('maxdiff', e.get('title', 'MaxDiff ' + str(i + 1)))]
                        if _md_sel:
                            export_maxdiff_to_excel(writer, _md_sel, var_labels)

                        # 4. Fill ToC content last (all sheet_maps ready)
                        export_toc_sheet(
                            writer, _xl_res, _xl_mtx,
                            var_labels, sheet_map,
                            regression_results=valid_reg,
                            anova_results=_anova_sel,
                            factor_results=_factor_sel,
                            conjoint_results=_conj_sel,
                            maxdiff_results=_md_sel,
                            pre_created_ws=toc_ws,
                        )

                    st.download_button(
                        "\u2b07\ufe0f Pobierz Raport_Analityczny.xlsx",
                        data=output.getvalue(),
                        file_name="Raport_Analityczny.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        use_container_width=True,
                    )
                except Exception as e:
                    st.error(f"B\u0142\u0105d podczas generowania pliku: {e}")
                    st.exception(e)

# =============================================================
# MODUL: EKSPORT DO POWERPOINT
# =============================================================
elif menu == "\U0001f4ca Eksport do PowerPoint":
    _require_module_access("export_pptx")
    _require_data()
    module_header("\U0001f4ca", "Eksport do PowerPoint", "Edytowalne wykresy kolumnowe z cz\u0119sto\u015bci i tabel krzy\u017cowych")
    st.info(
        "Generuje plik PowerPoint z edytowalnymi wykresami kolumnowymi. "
        "Ka\u017cdy wykres jest zagnie\u017cd\u017cony jako natywny obiekt PPT z w\u0142asn\u0105 "
        "tabel\u0105 danych \u2014 mo\u017cna go edytowa\u0107 bezpo\u015brednio w PowerPoint. "
        "Eksportowane s\u0105 wy\u0142\u0105cznie wyniki tabel cz\u0119sto\u015bci i tabel krzy\u017cowych."
    )

    freq_res  = st.session_state.results.get('czestosci', {})
    cross_res = st.session_state.results.get('krzyzowe', {})

    if not freq_res and not cross_res:
        st.warning(
            "Brak wynik\u00f3w do eksportu. Wygeneruj tabele cz\u0119sto\u015bci lub krzy\u017cowe "
            "w module **Analizy i Tabele**."
        )
    else:
        n_freq  = len(freq_res)
        n_cross = len(cross_res)

        st.markdown("**Dost\u0119pne wyniki:**")
        mc1, mc2 = st.columns(2)
        mc1.metric("\U0001f4c8 Tablice cz\u0119sto\u015bci", n_freq)
        mc2.metric("\U0001f500 Tabele krzy\u017cowe", n_cross)

        st.divider()
        st.markdown("**Opcje prezentacji:**")

        # \u2500\u2500 Szablon slajd\u00f3w \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
        pptx_template_file = st.file_uploader(
            "Szablon slajd\u00f3w (.pptx) \u2014 opcjonalny",
            type=["pptx"], key="ppt_template_file",
            help="Wgraj plik .pptx z gotowym layoutem (t\u0142o, logo, stopka). "
                 "Ka\u017cdy wykres zostanie dodany do nowego slajdu z tego szablonu."
        )
        if pptx_template_file:
            st.success("\u2705 Szablon slajd\u00f3w wczytany.")
        _crtx_style = None

        st.divider()

        # -- Szablony wykres\u00f3w: tworzone w "Projekt i S\u0142ownik", tutaj tylko wybor --
        st.session_state.setdefault('ppt_chart_templates', {})
        _DEFAULT_CHART_TEMPLATE = CHART_TPL_DEFAULT

        if not st.session_state.ppt_chart_templates:
            st.info(
                "Nie masz jeszcze \u017cadnego szablonu wykres\u00f3w. Utw\u00f3rz go w module "
                "**\U0001f4c1 Projekt i S\u0142ownik \u2192 \U0001f3a8 Szablony wykres\u00f3w**. "
                "Bez wybranego szablonu u\u017cyty zostanie wygl\u0105d domy\u015blny."
            )
        else:
            st.caption(
                "Szablony wykres\u00f3w tworzysz w module **Projekt i S\u0142ownik \u2192 "
                "Szablony wykres\u00f3w**. Poni\u017cej wybierz tylko, kt\u00f3rego u\u017cy\u0107."
            )

        #st.divider()

        # \u2500\u2500 Selekcja szablonu dla cz\u0119sto\u015bci i krzy\u017cowych \u2500\u2500\u2500\u2500\u2500
        def _ppt_tpl_opts(chart_type):
            return ["(domy\u015blny)"] + [
                k for k, v in st.session_state.ppt_chart_templates.items()
                if v.get("chart_type", chart_type) == chart_type]

        tsel_col1, tsel_col2 = st.columns(2)
        with tsel_col1:
            _freq_opts = _ppt_tpl_opts("czestosci")
            ppt_freq_tpl = st.selectbox(
                "Szablon \u2014 tablice cz\u0119sto\u015bci:",
                _freq_opts, index=0, key="ppt_freq_tpl"
            )
        with tsel_col2:
            _cross_opts = _ppt_tpl_opts("krzyzowe")
            ppt_cross_tpl = st.selectbox(
                "Szablon \u2014 tabele krzy\u017cowe:",
                _cross_opts, index=0, key="ppt_cross_tpl"
            )

        def _get_tpl(name):
            if name == "(domy\u015blny)" or name not in st.session_state.ppt_chart_templates:
                return _DEFAULT_CHART_TEMPLATE
            return st.session_state.ppt_chart_templates[name]

        freq_tpl_def  = _get_tpl(ppt_freq_tpl)
        cross_tpl_def = _get_tpl(ppt_cross_tpl)

        st.divider()
        oc1, oc2, oc3 = st.columns(3)
        with oc1:
            ppt_metric = st.radio(
                "Warto\u015b\u0107 na wykresie:",
                ["Procent [%]", "Liczebno\u015b\u0107 [N]"],
                key="ppt_metric",
                help="Cz\u0119sto\u015bci: kolumna Procent lub Liczebno\u015b\u0107."
            )
        with oc2:
            ppt_title_prefix = st.text_input(
                "Prefiks tytu\u0142u slajdu:", value="",
                key="ppt_prefix", placeholder="np. Badanie 2025 |"
            )
        with oc3:
            ppt_show_base = st.checkbox(
                "Poka\u017c baz\u0119 (N) w tytu\u0142ach",
                value=True, key="ppt_base"
            )
            if is_spss:
                ppt_drop_empty = st.checkbox(
                    "Ukryj kategorie bez odpowiedzi (N = 0)",
                    value=False, key="ppt_drop_empty",
                    help="Dla danych SPSS: pomija na wykresach kategorie bez wyborow (liczebnosc 0).")
            else:
                ppt_drop_empty = False

        # \u2500\u2500 Split-aware PPT options \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
        _has_split_results = False
        try:
            for _k in list(st.session_state.results.get('czestosci', {}).keys()) + \
                      list(st.session_state.results.get('krzyzowe', {}).keys()):
                if " | " in _k and "=" in _k.split(" | ", 1)[-1]:
                    _has_split_results = True
                    break
        except Exception:
            pass

        if _has_split_results:
            st.markdown("**\U0001f500 Opcje podzia\u0142u na podzbiory:**")
            spl1, spl2 = st.columns(2)
            with spl1:
                ppt_group_slides = st.checkbox(
                    "Grupuj slajdy wg podzia\u0142u",
                    value=True, key="ppt_group_slides",
                    help="Slajdy zostan\u0105 u\u0142o\u017cone wed\u0142ug grup (np. najpierw wszystkie Kobiety, potem wszyscy M\u0119\u017cczy\u017ani) zamiast wg zmiennych."
                )
            with spl2:
                ppt_section_dividers = st.checkbox(
                    "Dodaj slajd-przerywnik przed ka\u017cd\u0105 grup\u0105",
                    value=True, key="ppt_section_dividers",
                    help="Przed pierwszym slajdem ka\u017cdej grupy zostanie wstawiony slajd tytu\u0142owy z nazw\u0105 grupy."
                )
        else:
            ppt_group_slides = False
            ppt_section_dividers = False

        # \u2500\u2500 Motyw kolorystyczny prezentacji \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
        tc1, tc2 = st.columns([1, 2])
        with tc1:
            _title_hex = st.color_picker(
                "Motyw kolorystyczny prezentacji:",
                value="#1F4E79", key="ppt_theme_color",
                help="Kolor paska z tytu\u0142em slajdu (dotyczy tylko prezentacji bez szablonu .pptx)."
            )
        with tc2:
            st.markdown(
                f'<div style="margin-top:32px;display:inline-block;'
                f'width:40px;height:30px;background:{_title_hex};'
                f'border:1px solid #ccc;border-radius:4px;vertical-align:middle;"></div> '
                f'<span style="margin-left:8px;line-height:30px;vertical-align:middle;">'
                f'Wybrany kolor: <code>{_title_hex}</code></span>',
                unsafe_allow_html=True
            )

        # bar color still comes from freq template's first color
        _bar_hex = freq_tpl_def["colors"][0]

        st.divider()

        # -- szczegolowy wybor wynikow (multiselect + przycisk zatwierdzajacy) --
        _ppt_cats = []
        if freq_res:
            _ppt_cats.append(('czestosci', 'Tablice cz\u0119sto\u015bci',
                              list(freq_res.keys()), True))
        if cross_res:
            _ppt_cats.append(('krzyzowe', 'Tabele krzy\u017cowe',
                              list(cross_res.keys()), True))
        if _ppt_cats:
            render_granular_selector('ppt', _ppt_cats, var_labels,
                                     intro="Rozwi\u0144 kategori\u0119, aby wybra\u0107 konkretne wykresy. "
                                           "Po zmianie kliknij \u201eZatwierd\u017a wyb\u00f3r\u201c.")
        _ppt_keep = export_selection_filter('ppt')
        _ppt_ord  = st.session_state.get('ppt_result_order', {})

        # Zastosuj kolejnosc (jesli ustawiona) lub filtr wyboru
        def _ppt_ordered(src_keys, cat_key):
            if cat_key in _ppt_ord and _ppt_ord[cat_key]:
                return [k for k in _ppt_ord[cat_key] if k in set(src_keys)]
            return [k for k in src_keys if _ppt_keep(cat_key, k)]

        sel_freq_keys  = _ppt_ordered(list(freq_res.keys()),  'czestosci')
        sel_cross_keys = _ppt_ordered(list(cross_res.keys()), 'krzyzowe')

        # -- reorder panel --
        _ppt_ro_cats = []
        if freq_res:
            _ppt_ro_cats.append(('czestosci', 'Tablice cz\u0119sto\u015bci',
                                  list(freq_res.keys())))
        if cross_res:
            _ppt_ro_cats.append(('krzyzowe', 'Tabele krzy\u017cowe',
                                  list(cross_res.keys())))
        render_reorder_ui('ppt', _ppt_ro_cats, var_labels)

        n_selected = len(sel_freq_keys) + len(sel_cross_keys)
        st.caption(f"Wybrano {n_selected} wykres\u00f3w do eksportu.")

        if st.button(
            f"\U0001f4ca Generuj plik PowerPoint ({n_selected} slajd\u00f3w)",
            type="primary", use_container_width=True,
            key="ppt_generate", disabled=(n_selected == 0)
        ):
            try:
                from pptx import Presentation
                from pptx.util import Inches, Pt
                from pptx.chart.data import ChartData
                from pptx.enum.chart import XL_CHART_TYPE
                from pptx.dml.color import RGBColor
                from pptx.util import Emu
                import lxml.etree as _etree

                def _hex_to_rgb(h):
                    if isinstance(h, RGBColor):
                        return h   # already an RGBColor (from crtx parser)
                    h = h.lstrip("#")
                    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

                BAR_COLOR   = _hex_to_rgb(_bar_hex)
                TITLE_COLOR = _hex_to_rgb(_title_hex)
                WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

                import colorsys as _cs

                def _palette_series(n, base_hex):
                    # If crtx provides colors, cycle through them
                    if _crtx_style and _crtx_style['series_colors']:
                        colors = _crtx_style['series_colors']
                        return [colors[i % len(colors)] for i in range(n)]
                    bh = base_hex if isinstance(base_hex, str) else f"#{base_hex[0]:02X}{base_hex[1]:02X}{base_hex[2]:02X}"
                    bh = bh.lstrip('#')
                    r,g,b = int(bh[:2],16)/255, int(bh[2:4],16)/255, int(bh[4:],16)/255
                    h,s,v = _cs.rgb_to_hsv(r,g,b)
                    colors = []
                    for i in range(n):
                        vi = max(0.25, v - i * (v-0.25)/(max(n-1,1)))
                        si = max(0.15, s - i * (s-0.15)/(max(n-1,1)))
                        rr,gg,bb = _cs.hsv_to_rgb(h, si, vi)
                        colors.append(RGBColor(int(rr*255), int(gg*255), int(bb*255)))
                    return colors

                # \u2500\u2500 Presentation base \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
                if pptx_template_file is not None:
                    import io as _io
                    prs = Presentation(_io.BytesIO(pptx_template_file.getvalue()))
                    # Use the first slide layout that looks blank, or fall back to [6]
                    blank_layout = None
                    for lay in prs.slide_layouts:
                        if len(lay.placeholders) == 0:
                            blank_layout = lay
                            break
                    if blank_layout is None:
                        blank_layout = prs.slide_layouts[min(6, len(prs.slide_layouts)-1)]
                else:
                    prs = Presentation()
                    prs.slide_width  = Inches(13.33)
                    prs.slide_height = Inches(7.5)
                    blank_layout = prs.slide_layouts[6]

                # \u2500\u2500 Font from crtx \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
                _chart_font = (_crtx_style.get('font_name') if _crtx_style else None) or None

                # \u2500\u2500 Gridlines setting from crtx \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
                _show_gridlines = (_crtx_style.get('gridlines', False)
                                   if _crtx_style else False)
                _gridline_color = (_crtx_style.get('gridline_color')
                                   if _crtx_style else None)

                def _slide_base(title_text, subtitle=""):
                    """Create a slide with title bar and return (slide, chart_top)."""
                    slide = prs.slides.add_slide(blank_layout)

                    if pptx_template_file is None:
                        # No template \u2014 draw custom white bg + blue/crtx title bar
                        bg = slide.shapes.add_shape(1, Inches(0), Inches(0),
                                                    prs.slide_width, prs.slide_height)
                        bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
                        bg.line.fill.background()

                        tb = slide.shapes.add_shape(1, Inches(0), Inches(0),
                                                    prs.slide_width, Inches(1.1))
                        tb.fill.solid(); tb.fill.fore_color.rgb = TITLE_COLOR
                        tb.line.fill.background()

                        ttf = slide.shapes.add_textbox(Inches(0.3), Inches(0.15),
                                                        Inches(12.5), Inches(0.8))
                        ttf.text_frame.word_wrap = True
                        tp = ttf.text_frame.paragraphs[0]
                        tp.text = title_text
                        tp.font.size = Pt(20); tp.font.bold = True
                        tp.font.color.rgb = WHITE
                        if _chart_font:
                            try: tp.font.name = _chart_font
                            except Exception: pass

                        chart_top = Inches(1.2)
                        if subtitle:
                            stf = slide.shapes.add_textbox(Inches(0.3), Inches(1.1),
                                                            Inches(12.5), Inches(0.3))
                            sp = stf.text_frame.paragraphs[0]
                            sp.text = subtitle; sp.font.size = Pt(11)
                            sp.font.color.rgb = RGBColor(0x59, 0x59, 0x59)
                            if _chart_font:
                                try: sp.font.name = _chart_font
                                except Exception: pass
                            chart_top = Inches(1.45)
                    else:
                        # Template slide \u2014 populate existing placeholders
                        chart_top = Inches(1.5)
                        for ph in slide.placeholders:
                            if ph.placeholder_format.idx == 0:  # title
                                ph.text = title_text
                                if _chart_font:
                                    try: ph.text_frame.paragraphs[0].font.name = _chart_font
                                    except Exception: pass
                                try:
                                    chart_top = Inches(
                                        (ph.top + ph.height) / 914400 + 0.1
                                    )
                                except Exception:
                                    chart_top = Inches(1.5)
                                break
                        if subtitle:
                            for ph in slide.placeholders:
                                if ph.placeholder_format.idx == 1:
                                    ph.text = subtitle
                                    break

                    return slide, chart_top

                def _fmt_numfmt(chart, use_pct):
                    """Apply Y-axis and gridline settings, optionally from crtx."""
                    try:
                        vax = chart.value_axis
                        vax.has_major_gridlines = _show_gridlines
                        if _show_gridlines and _gridline_color:
                            try:
                                vax.major_gridlines.format.line.color.rgb = _hex_to_rgb(_gridline_color)
                            except Exception:
                                pass
                        # Hide axis tick labels (values shown in data labels)
                        vax.tick_labels.font.size = Pt(1)
                        vax.tick_labels.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
                        vax.format.line.color.rgb = RGBColor(0xFF,0xFF,0xFF)
                    except Exception:
                        pass

                def _apply_dls(dls, color, tpl=None):
                    """Apply data label style using template."""
                    if tpl is None:
                        tpl = _DEFAULT_CHART_TEMPLATE
                    dls.show_value = bool(tpl.get("show_data_labels", True))
                    dls.show_category_name = False
                    dls.font.size = Pt(int(tpl.get("font_size_data", 9)))
                    dls.font.bold = bool(tpl.get("data_label_bold", True))
                    dls.font.color.rgb = color
                    if _chart_font:
                        try:
                            dls.font.name = _chart_font
                        except Exception:
                            pass

                def _apply_tpl_axes(chart, tpl, is_pct):
                    """Apply axis, gridline and legend settings from template."""
                    try:
                        vax = chart.value_axis
                        vax.has_major_gridlines = bool(tpl.get("show_gridlines", False))
                        if not tpl.get("show_y_axis", False):
                            vax.tick_labels.font.size = Pt(1)
                            vax.tick_labels.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
                            vax.format.line.color.rgb = RGBColor(0xFF,0xFF,0xFF)
                        else:
                            vax.tick_labels.font.size = Pt(int(tpl.get("font_size_labels", 10)))
                        vax.tick_labels.number_format = "0.0%" if is_pct else "#,##0"
                        vax.tick_labels.number_format_is_linked = False
                    except Exception:
                        pass
                    try:
                        cax = chart.category_axis
                        if tpl.get("show_x_axis", True):
                            cax.tick_labels.font.size = Pt(int(tpl.get("font_size_labels", 10)))
                        else:
                            cax.tick_labels.font.size = Pt(1)
                            cax.tick_labels.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
                    except Exception:
                        pass
                    # Legend
                    _LEG_POS = {"bottom": 4, "top": 1, "right": 2, "left": 3}
                    lp = tpl.get("legend_position", "bottom")
                    if lp == "none":
                        chart.has_legend = False
                    else:
                        chart.has_legend = True
                        try:
                            chart.legend.position = _LEG_POS.get(lp, 4)
                            chart.legend.include_in_layout = False
                            chart.legend.font.size = Pt(int(tpl.get("font_size_labels", 10)))
                        except Exception:
                            pass

                def _add_chart_slide(prs, title_text, categories, values,
                                     subtitle="", is_pct=False, tpl=None):
                    """Single-series frequency chart."""
                    if tpl is None:
                        tpl = _DEFAULT_CHART_TEMPLATE
                    slide, chart_top = _slide_base(title_text, subtitle)

                    cd = ChartData()
                    cd.categories = [str(c)[:40] for c in categories]
                    data_vals = []
                    for v in values:
                        if v == v and v is not None:
                            data_vals.append(float(v)/100.0 if is_pct else float(v))
                        else:
                            data_vals.append(0.0)
                    cd.add_series("Wynik", tuple(data_vals))

                    chart = slide.shapes.add_chart(
                        XL_CHART_TYPE.COLUMN_CLUSTERED,
                        Inches(0.4), chart_top,
                        Inches(12.5), prs.slide_height - chart_top - Inches(0.3),
                        cd
                    ).chart

                    chart.has_title = False
                    series = chart.series[0]
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = _hex_to_rgb(tpl["colors"][0])

                    dls = series.data_labels
                    _apply_dls(dls, _hex_to_rgb(tpl.get("title_color", "#1F4E79")), tpl)
                    # Format based on tpl.data_label_format
                    _fmt = tpl.get("data_label_format", "auto")
                    if _fmt == "percent" or (_fmt == "auto" and is_pct):
                        dls.number_format = "0.0%"
                    else:
                        dls.number_format = "#,##0"
                    dls.number_format_is_linked = False

                    # For freq charts legend is single-series \u2014 usually hide
                    if tpl.get("legend_position", "bottom") == "none":
                        chart.has_legend = False
                    else:
                        chart.has_legend = False  # single series \u2014 no legend
                    _apply_tpl_axes(chart, tpl, is_pct)

                def _add_cross_chart_slide(prs, title_text, categories,
                                           series_dict, subtitle="", is_pct=False, tpl=None):
                    """
                    Grouped column chart for cross-tabs.
                    categories  = row variable values (x-axis)
                    series_dict = {series_name: [values...]} one per column category
                    """
                    if tpl is None:
                        tpl = _DEFAULT_CHART_TEMPLATE
                    slide, chart_top = _slide_base(title_text, subtitle)

                    cd = ChartData()
                    cd.categories = [str(c)[:35] for c in categories]

                    ser_names = list(series_dict.keys())
                    tpl_colors = tpl.get("colors", _DEFAULT_CHART_TEMPLATE["colors"])

                    for sname, svals in series_dict.items():
                        clean = []
                        for v in svals:
                            if v == v and v is not None:
                                clean.append(float(v)/100.0 if is_pct else float(v))
                            else:
                                clean.append(0.0)
                        cd.add_series(str(sname)[:30], tuple(clean))

                    chart = slide.shapes.add_chart(
                        XL_CHART_TYPE.COLUMN_CLUSTERED,
                        Inches(0.4), chart_top,
                        Inches(12.5), prs.slide_height - chart_top - Inches(0.3),
                        cd
                    ).chart

                    chart.has_title = False

                    for i, series in enumerate(chart.series):
                        c_hex = tpl_colors[i % len(tpl_colors)]
                        series.format.fill.solid()
                        series.format.fill.fore_color.rgb = _hex_to_rgb(c_hex)

                        dls = series.data_labels
                        _apply_dls(dls, _hex_to_rgb(tpl.get("title_color", "#1F4E79")), tpl)
                        _fmt = tpl.get("data_label_format", "auto")
                        if _fmt == "percent" or (_fmt == "auto" and is_pct):
                            dls.number_format = "0.0%"
                        else:
                            dls.number_format = "#,##0"
                        dls.number_format_is_linked = False

                    _apply_tpl_axes(chart, tpl, is_pct)

                slides_added = 0
                use_pct = (ppt_metric == "Procent [%]")
                prefix  = (ppt_title_prefix.strip() + " " if ppt_title_prefix.strip() else "")

                # Helper: add a big section divider slide with group name
                def _add_section_divider(prs, title_text):
                    layout_blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
                    sld = prs.slides.add_slide(layout_blank)
                    # Background color bar
                    from pptx.util import Inches, Pt
                    from pptx.dml.color import RGBColor
                    from pptx.enum.shapes import MSO_SHAPE
                    sw, sh = prs.slide_width, prs.slide_height
                    bg = sld.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, sh)
                    try:
                        _r, _g, _b = _hex_to_rgb(_title_hex)
                        bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(_r, _g, _b)
                        bg.line.fill.background()
                    except Exception:
                        pass
                    # Title text box centered
                    tb = sld.shapes.add_textbox(Inches(0.5), sh/2 - Inches(1),
                                                 sw - Inches(1), Inches(2))
                    tf = tb.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.alignment = 2  # PP_ALIGN.CENTER = 2
                    run = p.add_run()
                    run.text = title_text
                    run.font.size = Pt(36)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

                # Helper: extract group label from a result key (returns '' if none)
                def _extract_group(k):
                    if " | " in k and "=" in k.split(" | ", 1)[-1]:
                        return k.rsplit(" | ", 1)[1]
                    return ""

                # Reorder keys if grouping requested
                if ppt_group_slides:
                    # Stable sort by group label (empty group first = full base)
                    sel_freq_keys  = sorted(sel_freq_keys,  key=lambda k: _extract_group(k))
                    sel_cross_keys = sorted(sel_cross_keys, key=lambda k: _extract_group(k))

                # Track current group to emit section dividers
                _current_group = object()  # sentinel so first comparison differs

                # \u2500\u2500 Frequency tables (single series) \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
                for var_name_key in sel_freq_keys:
                    df_res = prepare_export_table(freq_res[var_name_key], drop_empty_cats=ppt_drop_empty)
                    try:
                        # Split off group label if present
                        _base_var = var_name_key
                        _grp_suffix = ""
                        if " | " in var_name_key and "=" in var_name_key.split(" | ", 1)[-1]:
                            _base_var, _grp_suffix = var_name_key.rsplit(" | ", 1)

                        # Section divider if group changes
                        if ppt_section_dividers and ppt_group_slides and _grp_suffix != _current_group:
                            _current_group = _grp_suffix
                            if _grp_suffix:
                                _add_section_divider(prs, f"\U0001f500 {_grp_suffix}")
                                slides_added += 1

                        df_c = df_res.copy()
                        for sr in ['Suma','Og\u00f3\u0142em (Wa\u017cne)','Og\u00f3\u0142em',
                                   'Braki danych','Brak odpowiedzi']:
                            df_c = df_c[df_c.index.astype(str) != sr]
                        df_c = df_c[~df_c.index.astype(str).str.startswith('[')]

                        pct_col = next((c for c in df_c.columns
                                        if 'Procent' in str(c) or '%' in str(c)), None)
                        n_col   = next((c for c in df_c.columns
                                        if 'Liczebno' in str(c) or c == 'N'), None)
                        val_col = (pct_col if use_pct else n_col) or df_c.columns[0]

                        vals = pd.to_numeric(
                            df_c[val_col].apply(_to_float_pct),
                            errors='coerce').tolist()
                        cats = df_c.index.tolist()
                        if not cats or all(v != v for v in vals):
                            continue

                        lbl = var_labels.get(_base_var, _base_var)
                        unit = "%" if use_pct else "N"
                        if ppt_show_base:
                            try:
                                _skip = ['Suma','Og\u00f3\u0142em (Wa\u017cne)','Og\u00f3\u0142em',
                                         'Braki danych','Brak odpowiedzi']
                                _n_col = next((c for c in df_res.columns
                                               if 'Liczebno' in str(c) or c == 'N'), df_res.columns[0])
                                _n_ser = df_res[_n_col]
                                _n_ser = _n_ser[~_n_ser.index.astype(str).isin(_skip)]
                                _n_ser = _n_ser[~_n_ser.index.astype(str).str.startswith('[')]
                                base_n = int(pd.to_numeric(_n_ser, errors='coerce').dropna().sum())
                                sub = f"Baza: N={base_n} | Warto\u015bci: {unit}"
                            except Exception:
                                sub = f"Warto\u015bci: {unit}"
                        else:
                            sub = f"Warto\u015bci: {unit}"
                        if _grp_suffix:
                            sub = f"\U0001f500 {_grp_suffix} | {sub}"

                        _title_txt = f"{prefix}{lbl}"
                        if _grp_suffix:
                            _title_txt += f" \u2014 {_grp_suffix}"
                        _add_chart_slide(prs, _title_txt[:120],
                                         cats, vals, sub, is_pct=use_pct, tpl=freq_tpl_def)
                        slides_added += 1
                    except Exception as _e:
                        st.warning(f"Pomini\u0119to '{var_name_key}': {_e}")

                # \u2500\u2500 Cross-tabs (grouped series, 1 chart per table) \u2500\u2500
                for cross_key in sel_cross_keys:
                    df_res = prepare_export_table(cross_res[cross_key], drop_empty_cats=ppt_drop_empty)
                    try:
                        df_c = df_res.copy()
                        for sr in ['Suma','Braki danych',
                                   'Braki danych (wykluczone z tabeli)']:
                            df_c = df_c[df_c.index.astype(str) != sr]
                        df_c = df_c[~df_c.index.astype(str).str.startswith('[')]

                        # Select value columns
                        if use_pct:
                            val_cols = [c for c in df_c.columns
                                        if '%' in str(c) or 'Procent' in str(c)]
                        else:
                            val_cols = [c for c in df_c.columns
                                        if '[N]' in str(c) and 'Suma' not in str(c)]
                        if not val_cols:
                            val_cols = [c for c in df_c.columns if c != 'Suma']

                        cats = df_c.index.tolist()
                        if not cats:
                            continue

                        # Build series dict: clean column name -> numeric values
                        series_dict = {}
                        for col in val_cols:
                            col_label = (str(col).replace('[%]','')
                                                  .replace('[N]','')
                                                  .replace('[% Kolumnowe]','')
                                                  .replace('[% Wierszowe]','')
                                                  .strip())
                            svals = pd.to_numeric(
                                df_c[col].apply(_to_float_pct),
                                errors='coerce').tolist()
                            if any(v == v for v in svals):  # at least one non-NaN
                                series_dict[col_label] = svals

                        if not series_dict:
                            continue

                        unit = "%" if use_pct else "N"
                        # Split off group label from cross_key if present
                        _base_cross = cross_key
                        _grp_suffix = ""
                        if " | " in cross_key and "=" in cross_key.split(" | ", 1)[-1]:
                            _base_cross, _grp_suffix = cross_key.rsplit(" | ", 1)

                        # Section divider if group changes
                        if ppt_section_dividers and ppt_group_slides and _grp_suffix != _current_group:
                            _current_group = _grp_suffix
                            if _grp_suffix:
                                _add_section_divider(prs, f"\U0001f500 {_grp_suffix}")
                                slides_added += 1

                        _title_s = f"{prefix}{_base_cross}"
                        if _grp_suffix:
                            _title_s += f" \u2014 {_grp_suffix}"
                        title_s = _title_s[:120]
                        sub = f"Warto\u015bci: {unit} | Serie = kategorie zmiennej w kolumnach"
                        if _grp_suffix:
                            sub = f"\U0001f500 {_grp_suffix} | {sub}"

                        _add_cross_chart_slide(prs, title_s, cats,
                                               series_dict, sub, is_pct=use_pct, tpl=cross_tpl_def)
                        slides_added += 1
                    except Exception as _e:
                        st.warning(f"Pomini\u0119to '{cross_key}': {_e}")

                if slides_added == 0:
                    st.error("Nie uda\u0142o si\u0119 wygenerowa\u0107 \u017cadnego wykresu.")
                else:
                    ppt_buf = io.BytesIO()
                    prs.save(ppt_buf)
                    ppt_buf.seek(0)
                    st.success(f"\u2705 Wygenerowano {slides_added} slajd\u00f3w.")
                    st.download_button(
                        label=f"\u2b07\ufe0f Pobierz prezentacj\u0119 ({slides_added} slajd\u00f3w)",
                        data=ppt_buf.getvalue(),
                        file_name="Wykresy_Analiz.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True
                    )

            except ImportError:
                st.error("Biblioteka `python-pptx` nie jest zainstalowana. Uruchom: `pip install python-pptx`")
            except Exception as _ppt_err:
                st.error(f"B\u0142\u0105d generowania PowerPoint: {_ppt_err}")
                st.exception(_ppt_err)


# =============================================================
# EKSPORT DO WORDA
# =============================================================
elif menu == "\U0001f4c4 Eksport do Worda":
    _require_module_access("export_word")
    _require_data()
    module_header("\U0001f4c4", "Eksport do Worda",
                  "Raport .docx z tabelami analitycznymi i interaktywnymi wykresami HTML")

    st.session_state.setdefault('ppt_chart_templates', {})

    # ---- helpers -------------------------------------------------------
    def _wd_hex_bgr(hex_color):
        """Convert #RRGGBB to OOXML shading hex string (same order for Word)."""
        h = hex_color.lstrip("#")
        return h.upper()

    def _wd_shade_cell(cell, hex_color):
        """Set background shading of a docx table cell."""
        from docx.oxml.ns import qn as _qn
        from docx.oxml import OxmlElement as _OE
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        shd = _OE('w:shd')
        shd.set(_qn('w:val'), 'clear')
        shd.set(_qn('w:color'), 'auto')
        shd.set(_qn('w:fill'), _wd_hex_bgr(hex_color))
        tcPr.append(shd)

    def _wd_set_cell_font(cell, bold=False, color_hex=None, size_pt=9, center=False):
        from docx.shared import Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        for para in cell.paragraphs:
            if center:
                para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in para.runs:
                run.bold = bold
                run.font.size = Pt(size_pt)
                if color_hex:
                    h = color_hex.lstrip('#')
                    run.font.color.rgb = RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

    def _wd_add_hyperlink(para, url, text, color_hex="#1F4E79"):
        """Insert a clickable hyperlink into an existing paragraph."""
        from docx.oxml.ns import qn as _qn
        from docx.oxml import OxmlElement as _OE
        from docx.shared import RGBColor, Pt
        r_id = para.part.relate_to(
            url,
            'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink',
            is_external=True
        )
        hyperlink = _OE('w:hyperlink')
        hyperlink.set(_qn('r:id'), r_id)
        r = _OE('w:r')
        rPr = _OE('w:rPr')
        rStyle = _OE('w:rStyle')
        rStyle.set(_qn('w:val'), 'Hyperlink')
        rPr.append(rStyle)
        u = _OE('w:u')
        u.set(_qn('w:val'), 'single')
        rPr.append(u)
        c = _OE('w:color')
        h = color_hex.lstrip('#')
        c.set(_qn('w:val'), h.upper())
        rPr.append(c)
        sz = _OE('w:sz')
        sz.set(_qn('w:val'), '18')
        rPr.append(sz)
        r.append(rPr)
        t = _OE('w:t')
        t.text = text
        r.append(t)
        hyperlink.append(r)
        para._p.append(hyperlink)

    def _wd_apply_template_to_fig(fig, tpl):
        """Apply a ppt_chart_templates entry to a Plotly figure in-place."""
        if not tpl:
            return fig
        colors   = tpl.get('colors', ['#2E75B6'])
        title_c  = tpl.get('title_color', '#1F4E79')
        fs_title = tpl.get('font_size_title', 14)
        fs_lbl   = tpl.get('font_size_labels', 10)
        fs_data  = tpl.get('font_size_data', 9)
        show_dl  = tpl.get('show_data_labels', True)
        dl_fmt   = tpl.get('data_label_format', 'auto')
        dl_bold  = tpl.get('data_label_bold', True)
        leg_pos  = tpl.get('legend_position', 'bottom')
        show_grd = tpl.get('show_gridlines', False)
        show_y   = tpl.get('show_y_axis', False)
        show_x   = tpl.get('show_x_axis', True)

        fig.update_layout(
            title_font=dict(size=fs_title, color=title_c),
            legend=dict(
                orientation='h' if leg_pos in ('bottom','top') else 'v',
                yanchor='bottom' if leg_pos == 'bottom' else ('top' if leg_pos == 'top' else 'middle'),
                y=-0.25 if leg_pos == 'bottom' else (1.02 if leg_pos == 'top' else 0.5),
                xanchor='center' if leg_pos in ('bottom','top') else ('left' if leg_pos == 'right' else 'right'),
                x=0.5 if leg_pos in ('bottom','top') else (1.02 if leg_pos == 'right' else -0.02),
            ) if leg_pos != 'none' else dict(visible=False),
            xaxis=dict(visible=show_x, tickfont=dict(size=fs_lbl)),
            yaxis=dict(visible=show_y, showgrid=show_grd, tickfont=dict(size=fs_lbl)),
            plot_bgcolor='white',
            paper_bgcolor='white',
            margin=dict(l=20, r=20, t=50, b=60),
        )
        # Apply color palette to traces
        for i, trace in enumerate(fig.data):
            trace.marker.color = colors[i % len(colors)]
        # Data labels
        if show_dl:
            for trace in fig.data:
                if dl_fmt == 'percent':
                    trace.texttemplate = '%{y:.1f}%'
                elif dl_fmt == 'number':
                    trace.texttemplate = '%{y:.0f}'
                else:
                    # auto: dodaj % gdy kolumna wygladna na procentowa
                    _tn = str(getattr(trace, 'name', '') or '')
                    _ty = str(getattr(trace, 'hovertemplate', '') or '')
                    _is_pct = ('%' in _tn or '%' in _ty
                               or '[%]' in _tn or 'Procent' in _tn)
                    trace.texttemplate = '%{y:.1f}%' if _is_pct else '%{y:.1f}'
                trace.textposition = 'outside'
                trace.textfont = dict(size=fs_data, color=title_c)
        else:
            for trace in fig.data:
                trace.mode = 'markers' if hasattr(trace, 'mode') else trace.mode if hasattr(trace, 'mode') else None
                trace.text = None
        return fig

    def _wd_chart_html(fig, title_str, tpl):
        """Return standalone Plotly HTML string for a chart, template applied."""
        import plotly.io as _pio
        _wd_apply_template_to_fig(fig, tpl)
        fig.update_layout(title_text=title_str, height=480)
        return _pio.to_html(fig, full_html=True, include_plotlyjs='cdn',
                            config={'displayModeBar': True, 'responsive': True})

    def _wd_add_df_table(doc, df, theme_hex='#1F4E79', highlight_sum=True, size_pt=8.5, banner_blocks=None):
        """Add a styled DataFrame as a docx table. Returns the table object."""
        from docx.shared import Pt
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        n_rows, n_cols = df.shape
        extra_hdr = 1 if banner_blocks else 0
        table = doc.add_table(rows=n_rows + 1 + extra_hdr, cols=n_cols + 1)
        table.style = 'Table Grid'
        _SUM_IDX = {
            'Suma', 'Baza (N)', 'Baza (N) / Suma (%)',
            'Braki danych', 'Braki danych (wykluczone z tabeli)',
            'Og\u00f3\u0142em (Wa\u017cne)',
        }
        if banner_blocks:
            # Wiersz 0: scalone naglowki blokow (ciemny niebieski)
            blk_cells = table.rows[0].cells
            blk_cells[0].text = ''
            _wd_shade_cell(blk_cells[0], theme_hex)
            _wd_set_cell_font(blk_cells[0], bold=True, color_hex='#FFFFFF', size_pt=size_pt)
            col_cursor = 1
            for _blk_lbl, _blk_cols in banner_blocks:
                _n = len(_blk_cols)
                if _n > 1:
                    _mc = blk_cells[col_cursor].merge(blk_cells[col_cursor + _n - 1])
                else:
                    _mc = blk_cells[col_cursor]
                _mc.text = _blk_lbl
                _wd_shade_cell(_mc, theme_hex)
                _wd_set_cell_font(_mc, bold=True, color_hex='#FFFFFF', size_pt=size_pt, center=True)
                for para in _mc.paragraphs:
                    para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                col_cursor += _n
        # Wiersz naglowkowy kategorii
        hdr_cells = table.rows[extra_hdr].cells
        hdr_cells[0].text = ''
        if banner_blocks:
            _wd_shade_cell(hdr_cells[0], '#D6E4F0')
            _wd_set_cell_font(hdr_cells[0], bold=True, color_hex='#1F4E79', size_pt=size_pt)
        else:
            _wd_shade_cell(hdr_cells[0], theme_hex)
            _wd_set_cell_font(hdr_cells[0], bold=True, color_hex='#FFFFFF', size_pt=size_pt)
        for j, col in enumerate(df.columns):
            c = hdr_cells[j + 1]
            col_s = str(col)
            if banner_blocks:
                if '=' in col_s:
                    cat_disp = col_s.split('=', 1)[1]
                elif ' [' in col_s:
                    _pfx = col_s.split(' [')[0]
                    cat_disp = col_s[len(_pfx):].lstrip()
                else:
                    cat_disp = col_s
                c.text = cat_disp
                _wd_shade_cell(c, '#D6E4F0')
                _wd_set_cell_font(c, bold=True, color_hex='#1F4E79', size_pt=size_pt, center=True)
            else:
                c.text = col_s
                _wd_shade_cell(c, theme_hex)
                _wd_set_cell_font(c, bold=True, color_hex='#FFFFFF', size_pt=size_pt, center=True)
            for para in c.paragraphs:
                para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        # Data rows
        import math as _mwd
        for i, (idx, row) in enumerate(df.iterrows()):
            cells = table.rows[i + 1 + extra_hdr].cells
            # Banner: wiersz-naglowek pytania (cala linia pusta) -> scalony pasek
            if banner_blocks and all(
                    (v is None) or (hasattr(v, '__float__') and not isinstance(v, str)
                                    and _mwd.isnan(float(v)))
                    for v in row):
                _merged = cells[0]
                for _k in range(1, len(cells)):
                    _merged = _merged.merge(cells[_k])
                _merged.text = str(idx)
                _wd_shade_cell(_merged, '#2E75B6')
                _wd_set_cell_font(_merged, bold=True, color_hex='#FFFFFF', size_pt=size_pt)
                continue
            is_sum = str(idx) in _SUM_IDX
            bg = '#D6E4F0' if (is_sum and highlight_sum) else ('#F2F2F2' if i % 2 == 0 else '#FFFFFF')
            cells[0].text = str(idx)
            _wd_shade_cell(cells[0], bg)
            _wd_set_cell_font(cells[0], bold=is_sum and highlight_sum, size_pt=size_pt)
            _is_base_row = str(idx) == 'Baza (N) / Suma (%)'
            for j, val in enumerate(row):
                c = cells[j + 1]
                _col_s = str(df.columns[j]).lower() if j < len(df.columns) else ''
                _is_pct_col = ('%' in _col_s or 'procent' in _col_s)
                _is_n_col = (
                    '[n]' in _col_s
                    or _col_s in ('n', 'liczebnosc [n]', 'liczebno\u015b\u0107 [n]')
                    or _col_s.startswith('liczebno')
                    or _col_s.startswith('liczba')
                ) and not _is_pct_col
                if hasattr(val, '__float__') and not isinstance(val, str):
                    try:
                        import math as _m
                        fv = float(val)
                        if _m.isnan(fv):
                            c.text = ''
                        elif _is_n_col:
                            c.text = str(int(round(fv)))
                        elif _is_base_row and _is_pct_col:
                            # wiersz podsumowania, kolumna %: suma procentow
                            c.text = f'{fv:.0f}%'
                        elif fv == round(fv) and abs(fv) >= 1000:
                            c.text = f'{fv:.0f}'
                        else:
                            c.text = f'{fv:.1f}'
                    except (ValueError, TypeError):
                        c.text = str(val) if val is not None else ''
                else:
                    c.text = '' if (val is None or (isinstance(val, float) and __import__('math').isnan(val))) else str(val)
                _wd_shade_cell(c, bg)
                _wd_set_cell_font(c, bold=is_sum and highlight_sum, size_pt=size_pt, center=True)
                for para in c.paragraphs:
                    para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        return table

    def _wd_section_heading(doc, text, level=1):
        from docx.shared import Pt, RGBColor
        h = doc.add_heading(text, level=level)
        for run in h.runs:
            run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
        return h

    def _wd_build_zip(cfg):
        """Build ZIP bytes: raport.docx + wykresy/*.html."""
        import io as _io, zipfile as _zf, datetime as _dt
        from docx import Document
        from docx.shared import Pt, Inches, Cm, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn as _qn
        from docx.oxml import OxmlElement as _OE

        ss   = st.session_state
        res  = ss.results

        def _wd_resolve_tpl(key):
            name = cfg.get(key, '(domy\u015blny)')
            if name and name not in ('(domy\u015blny)', '') \
                    and name in ss.ppt_chart_templates:
                return ss.ppt_chart_templates[name]
            return {}

        tpl       = _wd_resolve_tpl('freq_template')
        cross_tpl = _wd_resolve_tpl('cross_template')

        doc = Document()
        # Page orientation
        if cfg.get('landscape'):
            from docx.oxml import OxmlElement as _OE2
            from docx.oxml.ns import qn as _qn2
            section = doc.sections[0]
            new_w, new_h = section.page_height, section.page_width
            section.page_width  = new_w
            section.page_height = new_h
            section.left_margin   = Cm(1.5)
            section.right_margin  = Cm(1.5)
            section.top_margin    = Cm(2)
            section.bottom_margin = Cm(2)
        else:
            section = doc.sections[0]
            section.left_margin   = Cm(2.5)
            section.right_margin  = Cm(2.5)
            section.top_margin    = Cm(2.5)
            section.bottom_margin = Cm(2.5)

        theme_hex = '#1F4E79'

        # ---- Title page ------------------------------------------------
        doc.add_paragraph()
        title_para = doc.add_paragraph()
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        title_run = title_para.add_run(cfg.get('title') or 'Raport analityczny')
        title_run.bold = True
        title_run.font.size = Pt(24)
        title_run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
        doc.add_paragraph()
        if cfg.get('author'):
            ap = doc.add_paragraph()
            ap.alignment = WD_ALIGN_PARAGRAPH.CENTER
            ar = ap.add_run(cfg['author'])
            ar.font.size = Pt(13)
        dp = doc.add_paragraph()
        dp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        dr = dp.add_run(_dt.datetime.now().strftime('%d.%m.%Y'))
        dr.font.size = Pt(11)
        dr.font.color.rgb = RGBColor(0x70, 0x70, 0x70)
        if cfg.get('desc'):
            doc.add_paragraph()
            descp = doc.add_paragraph()
            descp.alignment = WD_ALIGN_PARAGRAPH.CENTER
            descr = descp.add_run(cfg['desc'])
            descr.font.size = Pt(10)
            descr.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
        # N= info
        n_resp = len(ss.get('df_raw_orig', ss.df_raw)) if hasattr(ss, 'df_raw') else 0
        try:
            n_resp = len(ss.df_raw)
        except Exception:
            pass
        if n_resp:
            np_ = doc.add_paragraph()
            np_.alignment = WD_ALIGN_PARAGRAPH.CENTER
            npr = np_.add_run(f'N = {n_resp:,}')
            npr.font.size = Pt(11)
            npr.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
        doc.add_page_break()

        # ---- TOC placeholder (Word refreshes on open) ------------------
        _wd_section_heading(doc, 'Spis tre\u015bci', level=1)
        toc_para = doc.add_paragraph()
        fldChar1 = _OE('w:fldChar')
        fldChar1.set(_qn('w:fldCharType'), 'begin')
        instrText = _OE('w:instrText')
        instrText.set(_qn('xml:space'), 'preserve')
        instrText.text = ' TOC \\o "1-2" \\h \\z \\u '
        fldChar2 = _OE('w:fldChar')
        fldChar2.set(_qn('w:fldCharType'), 'separate')
        fldChar_end = _OE('w:fldChar')
        fldChar_end.set(_qn('w:fldCharType'), 'end')
        toc_run = toc_para.add_run()
        toc_run._r.append(fldChar1)
        toc_run = toc_para.add_run()
        toc_run._r.append(instrText)
        toc_run = toc_para.add_run()
        toc_run._r.append(fldChar2)
        toc_run = toc_para.add_run()
        toc_run._r.append(fldChar_end)
        nop = doc.add_paragraph()
        nop.add_run('(Naci\u015bnij Ctrl+A, nast\u0119pnie F9 aby od\u015bwie\u017cy\u0107 spis tre\u015bci po otwarciu dokumentu)') \
            .font.size = Pt(8)
        if nop.runs:
            nop.runs[0].font.color.rgb = RGBColor(0x99,0x99,0x99)
        doc.add_page_break()

        charts_html = {}
        chart_counter = [0]
        _chart_mode = cfg.get('chart_mode', 'html')

        def _add_chart_link(section_title, df_for_chart, chart_type='bar', use_tpl=None):
            """Build Plotly figure; embed PNG or HTML link depending on chart_mode."""
            if not cfg.get('include_charts'):
                return
            _active_tpl = use_tpl if use_tpl is not None else tpl
            try:
                import plotly.express as _px
                import plotly.io as _pio
                import io as _chio
                plot_df = df_for_chart.copy()
                _skip = {'Suma','Baza (N) / Suma (%)','Braki danych',
                         'Braki danych (wykluczone z tabeli)','Og\u00f3\u0142em (Wa\u017cne)'}
                plot_df = plot_df[~plot_df.index.astype(str).isin(_skip)]
                plot_df = plot_df[~plot_df.index.astype(str).str.startswith('[')]
                if plot_df.empty:
                    return
                pct_cols = [c for c in plot_df.columns if '%' in str(c)]
                n_cols   = [c for c in plot_df.columns if '[N]' in str(c) and '%' not in str(c)]
                val_col  = pct_cols[0] if pct_cols else (n_cols[0] if n_cols else plot_df.columns[0])
                try:
                    plot_df[val_col] = pd.to_numeric(
                        plot_df[val_col].apply(_to_float_pct), errors='coerce')
                    plot_df = plot_df.dropna(subset=[val_col])
                except Exception:
                    return
                plot_df.index.name = plot_df.index.name or "Kategoria"
                fig = _px.bar(
                    plot_df.reset_index(), x=plot_df.index.name,
                    y=val_col, title=section_title,
                    color_discrete_sequence=_active_tpl.get('colors', ['#2E75B6']) if _active_tpl else ['#2E75B6']
                )
                _wd_apply_template_to_fig(fig, _active_tpl)
                fig.update_layout(title_text=section_title, height=420)
                chart_counter[0] += 1
                if _chart_mode == 'image':
                    from docx.shared import Cm as _ChCm
                    _png = _pio.to_image(fig, format='png', width=960, height=420, scale=2)
                    doc.add_picture(_chio.BytesIO(_png), width=_ChCm(16))
                    doc.add_paragraph()
                else:
                    fname = f'chart_{chart_counter[0]:03d}.html'
                    charts_html[fname] = _wd_chart_html(fig, section_title, tpl)
                    link_para = doc.add_paragraph()
                    _wd_add_hyperlink(link_para, f'wykresy/{fname}',
                                      '\U0001f4ca Otw\u00f3rz wykres interaktywny \u2192 ' + fname,
                                      color_hex='#1F4E79')
            except Exception:
                pass

        inc = cfg.get('include', {})
        sel = cfg.get('selected', {})
        _wd_drop_e = cfg.get('drop_empty', False)
        _order = cfg.get('order', {})

        def _apply_order_wd(d, cat_key):
            if cat_key in _order and _order[cat_key]:
                return {k: d[k] for k in _order[cat_key] if k in d}
            if cat_key in sel:
                _s = set(sel[cat_key])
                return {k: v for k, v in d.items() if k in _s}
            return dict(d)

        # ---- Tablice cz\u0119sto\u015bci ----------------------------------------
        czestosci = _apply_order_wd(res.get('czestosci', {}), 'czestosci')
        if inc.get('czestosci') and czestosci:
            _wd_section_heading(doc, 'Tablice cz\u0119sto\u015bci', level=1)
            for title, df_r in czestosci.items():
                df_r = prepare_export_table(df_r, drop_empty_cats=_wd_drop_e)
                _base, _grp = _extract_split_from_title(title)
                heading_txt = _base + (' \u2014 ' + _grp if _grp else '')
                _wd_section_heading(doc, heading_txt, level=2)
                try:
                    _wd_add_df_table(doc, df_r, theme_hex=theme_hex,
                                     highlight_sum=cfg.get('highlight_sum', True))
                    _add_chart_link(heading_txt, df_r)
                except Exception as _e:
                    doc.add_paragraph(f'[B\u0142\u0105d tabeli: {_e}]')
            doc.add_page_break()

        # ---- Tabele krzy\u017cowe ------------------------------------------
        krzyzowe = _apply_order_wd(res.get('krzyzowe', {}), 'krzyzowe')
        if inc.get('krzyzowe') and krzyzowe:
            _wd_section_heading(doc, 'Tabele krzy\u017cowe', level=1)
            for title, df_r in krzyzowe.items():
                df_r = prepare_export_table(df_r, drop_empty_cats=_wd_drop_e)
                _base, _grp = _extract_split_from_title(title)
                heading_txt = _base + (' \u2014 ' + _grp if _grp else '')
                _wd_section_heading(doc, heading_txt, level=2)
                try:
                    _wd_add_df_table(doc, df_r, theme_hex=theme_hex,
                                     highlight_sum=cfg.get('highlight_sum', True))
                    _add_chart_link(heading_txt, df_r, use_tpl=cross_tpl)
                except Exception as _e:
                    doc.add_paragraph(f'[B\u0142\u0105d tabeli: {_e}]')
            doc.add_page_break()

        # ---- Tabele zbiorcze / Banner ---------------------------------
        banner = _apply_order_wd(res.get('banner', {}), 'banner')
        if inc.get('banner') and banner:
            _wd_section_heading(doc, 'Tabele zbiorcze (Banner)', level=1)
            for title, df_r in banner.items():
                df_r = prepare_export_table(df_r, drop_empty_cats=_wd_drop_e)
                _base, _grp = _extract_split_from_title(title)
                heading_txt = _base + (' \u2014 ' + _grp if _grp else '')
                _wd_section_heading(doc, heading_txt, level=2)
                try:
                    _wd_add_df_table(doc, df_r, theme_hex=theme_hex,
                                     highlight_sum=cfg.get('highlight_sum', True),
                                     banner_blocks=parse_banner_blocks(df_r.columns))
                except Exception as _e:
                    doc.add_paragraph(f'[B\u0142\u0105d tabeli: {_e}]')
            doc.add_page_break()

        # ---- Tabele matrycowe -----------------------------------------
        matrix_results = ss.matrix_results
        if 'matrix' in _order and _order['matrix']:
            _mo = {k: i for i, k in enumerate(_order['matrix'])}
            matrix_results = sorted(
                [e for e in matrix_results if e['name'] in _mo],
                key=lambda e: _mo.get(e['name'], 9999))
        elif 'matrix' in sel:
            _msk = set(sel['matrix'])
            matrix_results = [e for e in matrix_results if e['name'] in _msk]
        if inc.get('matrix') and matrix_results:
            _wd_section_heading(doc, 'Tabele matrycowe', level=1)
            for entry in matrix_results:
                _base, _grp = _extract_split_from_title(entry['name'])
                heading_txt = _base + (' \u2014 ' + _grp if _grp else '')
                _wd_section_heading(doc, heading_txt, level=2)
                try:
                    _mdf = prepare_export_table(entry['df'], drop_empty_cats=_wd_drop_e)
                    _wd_add_df_table(doc, _mdf, theme_hex=theme_hex,
                                     highlight_sum=cfg.get('highlight_sum', True))
                    _add_chart_link(heading_txt, _mdf)
                except Exception as _e:
                    doc.add_paragraph(f'[B\u0142\u0105d tabeli: {_e}]')
            doc.add_page_break()

        # ---- Tabele \u015brednich ------------------------------------------
        srednie = _apply_order_wd(res.get('srednie', {}), 'srednie')
        if inc.get('srednie') and srednie:
            _wd_section_heading(doc, 'Tabele \u015brednich', level=1)
            for title, df_r in srednie.items():
                df_r = prepare_export_table(df_r, drop_empty_cats=_wd_drop_e)
                _base, _grp = _extract_split_from_title(title)
                heading_txt = _base + (' \u2014 ' + _grp if _grp else '')
                _wd_section_heading(doc, heading_txt, level=2)
                try:
                    _wd_add_df_table(doc, df_r, theme_hex=theme_hex,
                                     highlight_sum=cfg.get('highlight_sum', True))
                except Exception as _e:
                    doc.add_paragraph(f'[B\u0142\u0105d tabeli: {_e}]')
            doc.add_page_break()

        # ---- Statystyki opisowe ---------------------------------------
        opisowe = _apply_order_wd(res.get('opisowe', {}), 'opisowe')
        if inc.get('opisowe') and opisowe:
            _wd_section_heading(doc, 'Statystyki opisowe', level=1)
            for title, df_r in opisowe.items():
                df_r = prepare_export_table(df_r, drop_empty_cats=_wd_drop_e)
                _base, _grp = _extract_split_from_title(title)
                heading_txt = _base + (' \u2014 ' + _grp if _grp else '')
                _wd_section_heading(doc, heading_txt, level=2)
                try:
                    _wd_add_df_table(doc, df_r, theme_hex=theme_hex,
                                     highlight_sum=cfg.get('highlight_sum', True))
                except Exception as _e:
                    doc.add_paragraph(f'[B\u0142\u0105d tabeli: {_e}]')
            doc.add_page_break()

        # ---- Korelacje -----------------------------------------------
        korelacje = _apply_order_wd(res.get('korelacje', {}), 'korelacje')
        if inc.get('korelacje') and korelacje:
            _wd_section_heading(doc, 'Korelacje', level=1)
            for title, df_r in korelacje.items():
                _wd_section_heading(doc, title, level=2)
                try:
                    _wd_add_df_table(doc, df_r, theme_hex=theme_hex, highlight_sum=False)
                except Exception as _e:
                    doc.add_paragraph(f'[B\u0142\u0105d tabeli: {_e}]')
            doc.add_page_break()

        # ---- Regresja OLS --------------------------------------------
        _reg_all = ss.regression_results
        if 'regression' in sel:
            _regsk = set(sel['regression'])
            reg_results = [e for i, e in enumerate(_reg_all)
                           if ('Zm. zal.: ' + e.get('dep_var', 'Wynik ' + str(i+1))) in _regsk]
        else:
            reg_results = _reg_all
        if inc.get('regression') and reg_results:
            _wd_section_heading(doc, 'Regresja OLS', level=1)
            for res_r in reg_results:
                dep = var_labels.get(res_r.get('dep_var',''), res_r.get('dep_var',''))
                _wd_section_heading(doc, f'Zmienna zale\u017cna: {dep}', level=2)
                coef_df = res_r.get('coef_df')
                if coef_df is not None:
                    try:
                        _wd_add_df_table(doc, coef_df, theme_hex=theme_hex, highlight_sum=False)
                    except Exception as _e:
                        doc.add_paragraph(f'[B\u0142\u0105d: {_e}]')
                fit = res_r.get('fit', {})
                if fit:
                    fp = doc.add_paragraph()
                    r2   = fit.get('r2',    float('nan'))
                    r2a  = fit.get('r2_adj',float('nan'))
                    fst  = fit.get('f_stat',float('nan'))
                    fp_v = fit.get('f_p',   float('nan'))
                    fp.add_run(
                        f'R\u00b2={r2:.3f}  R\u00b2 adj={r2a:.3f}  F={fst:.2f}  p={fp_v:.3f}'
                    ).font.size = __import__('docx.shared', fromlist=['Pt']).Pt(9)
            doc.add_page_break()

        # ---- Regresja logistyczna ------------------------------------
        _log_all = ss.logistic_results
        if 'logistic' in sel:
            _logsk = set(sel['logistic'])
            log_results = [e for i, e in enumerate(_log_all)
                           if ('Zm. zal.: ' + e.get('dep_var', 'Wynik ' + str(i+1))) in _logsk]
        else:
            log_results = _log_all
        if inc.get('logistic') and log_results:
            _wd_section_heading(doc, 'Regresja logistyczna', level=1)
            for res_l in log_results:
                dep = var_labels.get(res_l.get('dep_var',''), res_l.get('dep_var',''))
                _wd_section_heading(doc, f'Zmienna zale\u017cna: {dep}', level=2)
                coef_df = res_l.get('coef_df')
                if coef_df is not None:
                    try:
                        _wd_add_df_table(doc, coef_df, theme_hex=theme_hex, highlight_sum=False)
                    except Exception as _e:
                        doc.add_paragraph(f'[B\u0142\u0105d: {_e}]')
            doc.add_page_break()

        # ---- ANOVA ----------------------------------------------------
        _anova_all = ss.anova_results
        if 'anova' in sel:
            _ansk = set(sel['anova'])
            anova_results = [e for e in _anova_all
                             if (e.get('dep_var','?') + ' wg ' + e.get('group_var','?')) in _ansk]
        else:
            anova_results = _anova_all
        if inc.get('anova') and anova_results:
            _wd_section_heading(doc, 'ANOVA', level=1)
            for res_a in anova_results:
                dep = var_labels.get(res_a.get('dep_var',''), res_a.get('dep_var',''))
                grp = var_labels.get(res_a.get('group_var',''), res_a.get('group_var',''))
                _wd_section_heading(doc, f'{dep} wg {grp}', level=2)
                for key in ('anova_table', 'means_df', 'tukey_df'):
                    df_part = res_a.get(key)
                    if df_part is not None and isinstance(df_part, __import__('pandas').DataFrame):
                        try:
                            _wd_add_df_table(doc, df_part, theme_hex=theme_hex, highlight_sum=False)
                            doc.add_paragraph()
                        except Exception:
                            pass
            doc.add_page_break()

        # ---- Analiza czynnikowa --------------------------------------
        _fa_all = ss.factor_results
        if 'factor' in sel:
            _fask = set(sel['factor'])
            fa_results = [e for i, e in enumerate(_fa_all)
                          if e.get('title', 'EFA ' + str(i+1)) in _fask]
        else:
            fa_results = _fa_all
        if inc.get('factor') and fa_results:
            _wd_section_heading(doc, 'Analiza czynnikowa', level=1)
            for res_f in fa_results:
                _wd_section_heading(doc, res_f.get('title', 'EFA'), level=2)
                loads = res_f.get('loadings_df')
                if loads is not None:
                    try:
                        _wd_add_df_table(doc, loads, theme_hex=theme_hex, highlight_sum=False)
                    except Exception:
                        pass
            doc.add_page_break()

        # ---- Conjoint -------------------------------------------------
        _conj_all = ss.conjoint_results
        if 'conjoint' in sel:
            _conjsk = set(sel['conjoint'])
            conj_results = [e for i, e in enumerate(_conj_all)
                            if e.get('title', 'Conjoint ' + str(i+1)) in _conjsk]
        else:
            conj_results = _conj_all
        if inc.get('conjoint') and conj_results:
            _wd_section_heading(doc, 'Conjoint', level=1)
            for res_c in conj_results:
                _wd_section_heading(doc, res_c.get('title', 'Conjoint'), level=2)
                for key in ('utils_df', 'importance_df'):
                    df_c = res_c.get(key)
                    if df_c is not None:
                        try:
                            _wd_add_df_table(doc, df_c, theme_hex=theme_hex, highlight_sum=False)
                            doc.add_paragraph()
                        except Exception:
                            pass
            doc.add_page_break()

        # ---- MaxDiff --------------------------------------------------
        _md_all = ss.maxdiff_results
        if 'maxdiff' in sel:
            _mdsk = set(sel['maxdiff'])
            md_results = [e for i, e in enumerate(_md_all)
                          if e.get('title', 'MaxDiff ' + str(i+1)) in _mdsk]
        else:
            md_results = _md_all
        if inc.get('maxdiff') and md_results:
            _wd_section_heading(doc, 'MaxDiff', level=1)
            for res_m in md_results:
                _wd_section_heading(doc, res_m.get('title', 'MaxDiff'), level=2)
                scores_df = res_m.get('scores_df')
                if scores_df is not None:
                    try:
                        _wd_add_df_table(doc, scores_df, theme_hex=theme_hex, highlight_sum=False)
                        _add_chart_link(res_m.get('title', 'MaxDiff'), scores_df)
                    except Exception:
                        pass
            doc.add_page_break()

        # ---- Metadata footer ------------------------------------------
        _wd_section_heading(doc, 'Metodologia', level=1)
        meta_p = doc.add_paragraph()
        meta_lines = [
            f'Data wygenerowania: {_dt.datetime.now().strftime("%d.%m.%Y %H:%M")}',
            f'N (liczba respondent\u00f3w): {n_resp:,}' if n_resp else '',
            f'Wagi aktywne: {"Tak" if ss.weights is not None else "Nie"}',
            f'Podzia\u0142 na podzbiory: {ss.split_var or "Brak"}',
        ]
        for line in meta_lines:
            if line:
                p = doc.add_paragraph(line, style='List Bullet')
                p.runs[0].font.size = __import__('docx.shared', fromlist=['Pt']).Pt(9)

        # ---- Package --------------------------------------------------
        doc_buf = _io.BytesIO()
        doc.save(doc_buf)
        if _chart_mode == 'image':
            return doc_buf.getvalue()
        zip_buf = _io.BytesIO()
        with _zf.ZipFile(zip_buf, 'w', _zf.ZIP_DEFLATED) as zf:
            zf.writestr('raport.docx', doc_buf.getvalue())
            for fname, html in charts_html.items():
                zf.writestr(f'wykresy/{fname}', html.encode('utf-8'))
        return zip_buf.getvalue()

    # ---- END helpers ---------------------------------------------------

    tab_cfg, tab_gen = st.tabs([
        "\u2699\ufe0f Zawarto\u015b\u0107",
        "\U0001f4c4 Generuj",
    ])

    # ================================================================
    # TAB 1 \u2014 Zawarto\u015b\u0107
    # ================================================================
    with tab_cfg:
        st.markdown("#### Konfiguracja dokumentu")
        wc1, wc2 = st.columns(2)
        with wc1:
            wd_title  = st.text_input("Tytu\u0142 dokumentu:", key="wd_title",
                                       placeholder="np. Raport z badania satysfakcji 2025")
            wd_author = st.text_input("Autor:", key="wd_author",
                                       placeholder="np. Jan Kowalski")
        with wc2:
            wd_desc   = st.text_area("Opis / notatki:", key="wd_desc", height=68)
            wd_orient = st.radio("Orientacja stron:", ["Pionowa", "Pozioma"],
                                 horizontal=True, key="wd_orient")

        # Kategorie sa wlaczane automatycznie, jesli maja wyniki; o ostatecznej
        # zawartosci raportu decyduje sekcja "Szczegolowy wybor wynikow" ponizej.
        wd_inc_freq = wd_inc_cross = wd_inc_matrix = wd_inc_means = True
        wd_inc_desc = wd_inc_corr = wd_inc_reg = wd_inc_log = True
        wd_inc_anova = wd_inc_fa = wd_inc_conj = wd_inc_md = True
        wd_inc_banner = True

        # ---- Granular result selector -----------------------------------
        _wd_selected = {}
        _wd_has_any = any([
            wd_inc_freq   and bool(st.session_state.results.get('czestosci')),
            wd_inc_cross  and bool(st.session_state.results.get('krzyzowe')),
            wd_inc_banner and bool(st.session_state.results.get('banner')),
            wd_inc_matrix and bool(st.session_state.matrix_results),
            wd_inc_means  and bool(st.session_state.results.get('srednie')),
            wd_inc_desc   and bool(st.session_state.results.get('opisowe')),
            wd_inc_corr   and bool(st.session_state.results.get('korelacje')),
            wd_inc_reg    and bool(st.session_state.regression_results),
            wd_inc_log    and bool(st.session_state.logistic_results),
            wd_inc_anova  and bool(st.session_state.anova_results),
            wd_inc_fa     and bool(st.session_state.factor_results),
            wd_inc_conj   and bool(st.session_state.conjoint_results),
            wd_inc_md     and bool(st.session_state.maxdiff_results),
        ])
        if _wd_has_any:
            st.markdown("---")
            st.markdown("#### Szczeg\u00f3\u0142owy wyb\u00f3r wynik\u00f3w")
            st.caption("Rozwi\u0144 kategori\u0119 i odznacz to, czego NIE chcesz eksportowa\u0107, "
                       "a nast\u0119pnie kliknij **Zatwierd\u017a wyb\u00f3r**. Domy\u015blnie zaznaczone s\u0105 wszystkie.")

            def _wd_key_label(key):
                s = str(key)
                _m = re.search(r'\s*\[[^\]]*\]\s*$', s)
                if _m:
                    s = s[:_m.start()]
                _base, _grp = _extract_split_from_title(s)
                if ' x ' in _base:
                    _lbl = ' \u00d7 '.join(
                        get_var_display_name(_p.strip(), var_labels)
                        for _p in _base.split(' x '))
                else:
                    _lbl = get_var_display_name(_base, var_labels)
                return _lbl + (' \u2014 ' + _grp if _grp else '')

            # Wczytaj zapisany wybor jako domyslny (persystuje przez nawigacje)
            _wd_conf_for_sel = st.session_state.get('wd_sel_confirmed', {})

            def _wd_default(cat_key, opts):
                return [k for k in _wd_conf_for_sel.get(cat_key, opts)
                        if k in set(opts)]

            if wd_inc_freq and st.session_state.results.get('czestosci'):
                _wsopts = list(st.session_state.results['czestosci'].keys())
                with st.expander(("Tablice cz\u0119sto\u015bci (" + (str(len(_wd_default('czestosci', _wsopts))) + "/" + str(len(_wsopts)) if len(_wd_default('czestosci', _wsopts)) < len(_wsopts) else str(len(_wsopts))) + ")"), expanded=False):
                    _wd_selected['czestosci'] = st.multiselect(
                        "Wybierz:", _wsopts, default=_wd_default('czestosci', _wsopts), key="wd_sel_freq",
                        format_func=_wd_key_label, label_visibility="collapsed")
            if wd_inc_cross and st.session_state.results.get('krzyzowe'):
                _wsopts = list(st.session_state.results['krzyzowe'].keys())
                with st.expander(("Tabele krzy\u017cowe (" + (str(len(_wd_default('krzyzowe', _wsopts))) + "/" + str(len(_wsopts)) if len(_wd_default('krzyzowe', _wsopts)) < len(_wsopts) else str(len(_wsopts))) + ")"), expanded=False):
                    _wd_selected['krzyzowe'] = st.multiselect(
                        "Wybierz:", _wsopts, default=_wd_default('krzyzowe', _wsopts), key="wd_sel_cross",
                        format_func=_wd_key_label, label_visibility="collapsed")
            if wd_inc_banner and st.session_state.results.get('banner'):
                _wsopts = list(st.session_state.results['banner'].keys())
                with st.expander(("Tabele zbiorcze / Banner (" + (str(len(_wd_default('banner', _wsopts))) + "/" + str(len(_wsopts)) if len(_wd_default('banner', _wsopts)) < len(_wsopts) else str(len(_wsopts))) + ")"), expanded=False):
                    _wd_selected['banner'] = st.multiselect(
                        "Wybierz:", _wsopts, default=_wd_default('banner', _wsopts), key="wd_sel_banner",
                        format_func=_wd_key_label, label_visibility="collapsed")
            if wd_inc_matrix and st.session_state.matrix_results:
                _wsopts = [e['name'] for e in st.session_state.matrix_results]
                with st.expander(("Tabele matrycowe (" + (str(len(_wd_default('matrix', _wsopts))) + "/" + str(len(_wsopts)) if len(_wd_default('matrix', _wsopts)) < len(_wsopts) else str(len(_wsopts))) + ")"), expanded=False):
                    _wd_selected['matrix'] = st.multiselect(
                        "Wybierz:", _wsopts, default=_wd_default('matrix', _wsopts), key="wd_sel_matrix",
                        format_func=_wd_key_label, label_visibility="collapsed")
            if wd_inc_means and st.session_state.results.get('srednie'):
                _wsopts = list(st.session_state.results['srednie'].keys())
                with st.expander(("Tabele \u015brednich (" + (str(len(_wd_default('srednie', _wsopts))) + "/" + str(len(_wsopts)) if len(_wd_default('srednie', _wsopts)) < len(_wsopts) else str(len(_wsopts))) + ")"), expanded=False):
                    _wd_selected['srednie'] = st.multiselect(
                        "Wybierz:", _wsopts, default=_wd_default('srednie', _wsopts), key="wd_sel_means",
                        format_func=_wd_key_label, label_visibility="collapsed")
            if wd_inc_desc and st.session_state.results.get('opisowe'):
                _wsopts = list(st.session_state.results['opisowe'].keys())
                with st.expander(("Statystyki opisowe (" + (str(len(_wd_default('opisowe', _wsopts))) + "/" + str(len(_wsopts)) if len(_wd_default('opisowe', _wsopts)) < len(_wsopts) else str(len(_wsopts))) + ")"), expanded=False):
                    _wd_selected['opisowe'] = st.multiselect(
                        "Wybierz:", _wsopts, default=_wd_default('opisowe', _wsopts), key="wd_sel_desc",
                        format_func=_wd_key_label, label_visibility="collapsed")
            if wd_inc_corr and st.session_state.results.get('korelacje'):
                _wsopts = list(st.session_state.results['korelacje'].keys())
                with st.expander(("Korelacje (" + (str(len(_wd_default('korelacje', _wsopts))) + "/" + str(len(_wsopts)) if len(_wd_default('korelacje', _wsopts)) < len(_wsopts) else str(len(_wsopts))) + ")"), expanded=False):
                    _wd_selected['korelacje'] = st.multiselect(
                        "Wybierz:", _wsopts, default=_wd_default('korelacje', _wsopts), key="wd_sel_corr",
                        format_func=_wd_key_label, label_visibility="collapsed")
            if wd_inc_reg and st.session_state.regression_results:
                _wsopts = ['Zm. zal.: ' + e.get('dep_var', 'Wynik ' + str(i+1))
                           for i, e in enumerate(st.session_state.regression_results)]
                with st.expander(("Regresja OLS (" + (str(len(_wd_default('regression', _wsopts))) + "/" + str(len(_wsopts)) if len(_wd_default('regression', _wsopts)) < len(_wsopts) else str(len(_wsopts))) + ")"), expanded=False):
                    _wd_selected['regression'] = st.multiselect(
                        "Wybierz:", _wsopts, default=_wd_default('regression', _wsopts), key="wd_sel_reg",
                        label_visibility="collapsed")
            if wd_inc_log and st.session_state.logistic_results:
                _wsopts = ['Zm. zal.: ' + e.get('dep_var', 'Wynik ' + str(i+1))
                           for i, e in enumerate(st.session_state.logistic_results)]
                with st.expander(("Regresja logistyczna (" + (str(len(_wd_default('logistic', _wsopts))) + "/" + str(len(_wsopts)) if len(_wd_default('logistic', _wsopts)) < len(_wsopts) else str(len(_wsopts))) + ")"), expanded=False):
                    _wd_selected['logistic'] = st.multiselect(
                        "Wybierz:", _wsopts, default=_wd_default('logistic', _wsopts), key="wd_sel_log",
                        label_visibility="collapsed")
            if wd_inc_anova and st.session_state.anova_results:
                _wsopts = [e.get('dep_var', '?') + ' wg ' + e.get('group_var', '?')
                           for e in st.session_state.anova_results]
                with st.expander(("ANOVA (" + (str(len(_wd_default('anova', _wsopts))) + "/" + str(len(_wsopts)) if len(_wd_default('anova', _wsopts)) < len(_wsopts) else str(len(_wsopts))) + ")"), expanded=False):
                    _wd_selected['anova'] = st.multiselect(
                        "Wybierz:", _wsopts, default=_wd_default('anova', _wsopts), key="wd_sel_anova",
                        label_visibility="collapsed")
            if wd_inc_fa and st.session_state.factor_results:
                _wsopts = [e.get('title', 'EFA ' + str(i+1))
                           for i, e in enumerate(st.session_state.factor_results)]
                with st.expander(("Analiza czynnikowa (" + (str(len(_wd_default('factor', _wsopts))) + "/" + str(len(_wsopts)) if len(_wd_default('factor', _wsopts)) < len(_wsopts) else str(len(_wsopts))) + ")"), expanded=False):
                    _wd_selected['factor'] = st.multiselect(
                        "Wybierz:", _wsopts, default=_wd_default('factor', _wsopts), key="wd_sel_fa",
                        label_visibility="collapsed")
            if wd_inc_conj and st.session_state.conjoint_results:
                _wsopts = [e.get('title', 'Conjoint ' + str(i+1))
                           for i, e in enumerate(st.session_state.conjoint_results)]
                with st.expander(("Conjoint (" + (str(len(_wd_default('conjoint', _wsopts))) + "/" + str(len(_wsopts)) if len(_wd_default('conjoint', _wsopts)) < len(_wsopts) else str(len(_wsopts))) + ")"), expanded=False):
                    _wd_selected['conjoint'] = st.multiselect(
                        "Wybierz:", _wsopts, default=_wd_default('conjoint', _wsopts), key="wd_sel_conj",
                        label_visibility="collapsed")
            if wd_inc_md and st.session_state.maxdiff_results:
                _wsopts = [e.get('title', 'MaxDiff ' + str(i+1))
                           for i, e in enumerate(st.session_state.maxdiff_results)]
                with st.expander(("MaxDiff (" + (str(len(_wd_default('maxdiff', _wsopts))) + "/" + str(len(_wsopts)) if len(_wd_default('maxdiff', _wsopts)) < len(_wsopts) else str(len(_wsopts))) + ")"), expanded=False):
                    _wd_selected['maxdiff'] = st.multiselect(
                        "Wybierz:", _wsopts, default=_wd_default('maxdiff', _wsopts), key="wd_sel_md",
                        label_visibility="collapsed")

            st.markdown("")
            if st.button("\u2705 Zatwierd\u017a wyb\u00f3r wynik\u00f3w do eksportu",
                         type="primary", key="wd_confirm_sel"):
                st.session_state['wd_sel_confirmed'] = {
                    _k: list(_v) for _k, _v in _wd_selected.items()}
                st.success("Wyb\u00f3r zatwierdzony \u2014 zostanie u\u017cyty przy generowaniu raportu.")
            _wd_conf = st.session_state.get('wd_sel_confirmed')
            if _wd_conf is not None:
                _conf_n = sum(len(_v) for _v in _wd_conf.values())
                st.caption("\u2705 Zatwierdzono do eksportu: " + str(_conf_n) + " wynik\u00f3w. "
                           "Po zmianie zaznacze\u0144 kliknij ponownie **Zatwierd\u017a wyb\u00f3r**.")
            else:
                st.caption("\u2139\ufe0f Wyb\u00f3r niezatwierdzony \u2014 wyeksportowane zostan\u0105 wszystkie "
                           "zaznaczone wyniki.")

        # -- panel kolejnosci (reorder) --
        _wd_active = st.session_state.get('wd_sel_confirmed', _wd_selected)
        _wd_ro_cats = []
        for _rk, _rn in [('czestosci', 'Tablice cz\u0119sto\u015bci'),
                          ('krzyzowe',  'Tabele krzy\u017cowe'),
                          ('banner',    'Tabele zbiorcze (Banner)'),
                          ('matrix',    'Tabele matrycowe'),
                          ('srednie',   'Tabele \u015brednich'),
                          ('opisowe',   'Statystyki opisowe'),
                          ('korelacje', 'Korelacje')]:
            _rkeys = list(_wd_active.get(_rk, []))
            if _rkeys:
                _wd_ro_cats.append((_rk, _rn, _rkeys))
        for _rk, _rn, _rkeys in [
            ('regression', 'Regresja OLS',
             ['Zm. zal.: ' + e.get('dep_var', 'Wynik ' + str(i+1))
              for i, e in enumerate(st.session_state.regression_results)]),
            ('logistic',   'Regresja logistyczna',
             ['Zm. zal.: ' + e.get('dep_var', 'Wynik ' + str(i+1))
              for i, e in enumerate(st.session_state.logistic_results)]),
            ('anova',      'ANOVA',
             [e.get('dep_var','?') + ' wg ' + e.get('group_var','?')
              for e in st.session_state.anova_results]),
            ('factor',     'Analiza czynnikowa',
             [e.get('title','EFA ' + str(i+1))
              for i, e in enumerate(st.session_state.factor_results)]),
            ('conjoint',   'Conjoint',
             [e.get('title','Conjoint ' + str(i+1))
              for i, e in enumerate(st.session_state.conjoint_results)]),
            ('maxdiff',    'MaxDiff',
             [e.get('title','MaxDiff ' + str(i+1))
              for i, e in enumerate(st.session_state.maxdiff_results)]),
        ]:
            _sel_here = list(_wd_active.get(_rk, _rkeys))
            _filtered = [k for k in _sel_here if k in set(_rkeys)]
            if len(_filtered) >= 2:
                _wd_ro_cats.append((_rk, _rn, _filtered))
        _wd_order = render_reorder_ui('wd', _wd_ro_cats, var_labels)

        st.markdown("---")
        st.markdown("#### Opcje tabel i wykres\u00f3w")
        to1, to2 = st.columns(2)
        with to1:
            wd_highlight = st.checkbox("Wyr\u00f3\u017cnij wiersz Suma / Baza",
                                        value=True, key="wd_highlight")
            if is_spss:
                wd_drop_empty = st.checkbox(
                    "Ukryj kategorie bez odpowiedzi (N = 0)",
                    value=False, key="wd_drop_empty",
                    help="Dla danych SPSS: pomija w tabelach kategorie ze s\u0142ownika "
                         "warto\u015bci, kt\u00f3rych nikt nie wybra\u0142 (liczebno\u015b\u0107 0).")
            else:
                wd_drop_empty = False
        with to2:
            wd_charts = st.checkbox('\U0001f4ca Do\u0142\u0105cz wykresy (obrazy w dokumencie)',
                                     value=True, key='wd_charts')
            if wd_charts:
                wd_chart_mode = 'image'
                _wd_freq_opts = ['(domy\u015blny)'] + [
                    k for k, v in st.session_state.ppt_chart_templates.items()
                    if v.get('chart_type', 'czestosci') == 'czestosci']
                _wd_cross_opts = ['(domy\u015blny)'] + [
                    k for k, v in st.session_state.ppt_chart_templates.items()
                    if v.get('chart_type', 'krzyzowe') == 'krzyzowe']
                wd_freq_tpl_sel  = st.selectbox('Szablon \u2014 tablice cz\u0119sto\u015bci:', _wd_freq_opts, key='wd_freq_tpl_sel')
                wd_cross_tpl_sel = st.selectbox('Szablon \u2014 tabele krzy\u017cowe:', _wd_cross_opts, key='wd_cross_tpl_sel')
            else:
                wd_chart_mode = 'image'
                wd_freq_tpl_sel  = '(domy\u015blny)'
                wd_cross_tpl_sel = '(domy\u015blny)'

    # ================================================================
    # TAB 3 \u2014 Generuj
    # ================================================================
    with tab_gen:
        # Summary of what will be included -- liczymy wg zatwierdzonego wyboru
        _wd_conf_g = st.session_state.get('wd_sel_confirmed', {})
        _wd_order_g = st.session_state.get('wd_result_order', {})

        def _wd_sel_count(cat_key, all_items):
            """Liczba wynikow wybranych do eksportu (wg potwierdzenia lub kolejnosci)."""
            if cat_key in _wd_order_g and _wd_order_g[cat_key]:
                return len(_wd_order_g[cat_key])
            if cat_key in _wd_conf_g:
                return len(_wd_conf_g[cat_key])
            return len(all_items)

        _wd_sections = []
        if wd_inc_freq and st.session_state.results.get('czestosci'):
            _n = _wd_sel_count('czestosci', st.session_state.results['czestosci'])
            _wd_sections.append("Tablice cz\u0119sto\u015bci (" + str(_n) + ")")
        if wd_inc_cross and st.session_state.results.get('krzyzowe'):
            _n = _wd_sel_count('krzyzowe', st.session_state.results['krzyzowe'])
            _wd_sections.append("Tabele krzy\u017cowe (" + str(_n) + ")")
        if wd_inc_banner and st.session_state.results.get('banner'):
            _n = _wd_sel_count('banner', st.session_state.results['banner'])
            _wd_sections.append("Tabele zbiorcze / Banner (" + str(_n) + ")")
        if wd_inc_matrix and st.session_state.matrix_results:
            _n = _wd_sel_count('matrix', st.session_state.matrix_results)
            _wd_sections.append("Tabele matrycowe (" + str(_n) + ")")
        if wd_inc_means and st.session_state.results.get('srednie'):
            _n = _wd_sel_count('srednie', st.session_state.results['srednie'])
            _wd_sections.append("Tabele \u015brednich (" + str(_n) + ")")
        if wd_inc_desc and st.session_state.results.get('opisowe'):
            _n = _wd_sel_count('opisowe', st.session_state.results['opisowe'])
            _wd_sections.append("Statystyki opisowe (" + str(_n) + ")")
        if wd_inc_corr and st.session_state.results.get('korelacje'):
            _n = _wd_sel_count('korelacje', st.session_state.results['korelacje'])
            _wd_sections.append("Korelacje (" + str(_n) + ")")
        if wd_inc_reg and st.session_state.regression_results:
            _n = _wd_sel_count('regression', st.session_state.regression_results)
            _wd_sections.append("Regresja OLS (" + str(_n) + ")")
        if wd_inc_log and st.session_state.logistic_results:
            _n = _wd_sel_count('logistic', st.session_state.logistic_results)
            _wd_sections.append("Regresja logistyczna (" + str(_n) + ")")
        if wd_inc_anova and st.session_state.anova_results:
            _n = _wd_sel_count('anova', st.session_state.anova_results)
            _wd_sections.append("ANOVA (" + str(_n) + ")")
        if wd_inc_fa and st.session_state.factor_results:
            _n = _wd_sel_count('factor', st.session_state.factor_results)
            _wd_sections.append("Analiza czynnikowa (" + str(_n) + ")")
        if wd_inc_conj and st.session_state.conjoint_results:
            _n = _wd_sel_count('conjoint', st.session_state.conjoint_results)
            _wd_sections.append("Conjoint (" + str(_n) + ")")
        if wd_inc_md and st.session_state.maxdiff_results:
            _n = _wd_sel_count('maxdiff', st.session_state.maxdiff_results)
            _wd_sections.append("MaxDiff (" + str(_n) + ")")

        if _wd_sections:
            st.markdown("**Zawarto\u015b\u0107 raportu:**")
            for _s in _wd_sections:
                st.markdown(f"- \u2705 {_s}")
            if wd_charts:
                if wd_chart_mode == 'image':
                    st.info('\U0001f5bc\ufe0f Wykresy b\u0119d\u0105 osadzone jako obrazy PNG bezpo\u015brednio w dokumencie .docx.')
                else:
                    st.info('\U0001f4ca Wykresy b\u0119d\u0105 do\u0142\u0105czone jako interaktywne pliki HTML w folderze `wykresy/` wewn\u0105trz archiwum ZIP.')
        else:
            st.warning("\u26a0\ufe0f Brak wynik\u00f3w do eksportu. Wykonaj analizy w modu\u0142ach analitycznych i wr\u00f3\u0107 tu.")

        st.markdown("---")
        _wd_gen_disabled = (len(_wd_sections) == 0)
        if st.button(
            f"\U0001f4c4 Generuj raport Word ({len(_wd_sections)} sekcji)",
            type="primary", use_container_width=True,
            key="wd_generate", disabled=_wd_gen_disabled
        ):
            with st.spinner("Generowanie raportu Word..."):
                try:
                    from docx import Document as _DocTest  # noqa: check import
                    _wd_cfg = {
                        'title':         wd_title,
                        'author':        wd_author,
                        'desc':          wd_desc,
                        'landscape':     (wd_orient == "Pozioma"),
                        'include': {
                            'czestosci': wd_inc_freq,
                            'krzyzowe':  wd_inc_cross,
                            'banner':    wd_inc_banner,
                            'matrix':    wd_inc_matrix,
                            'srednie':   wd_inc_means,
                            'opisowe':   wd_inc_desc,
                            'korelacje': wd_inc_corr,
                            'regression':wd_inc_reg,
                            'logistic':  wd_inc_log,
                            'anova':     wd_inc_anova,
                            'factor':    wd_inc_fa,
                            'conjoint':  wd_inc_conj,
                            'maxdiff':   wd_inc_md,
                        },
                        'highlight_sum':  wd_highlight,
                        'include_charts': wd_charts,
                        'chart_mode':     wd_chart_mode,
                        'freq_template':   wd_freq_tpl_sel,
                        'cross_template':  wd_cross_tpl_sel,
                        'selected':       st.session_state.get('wd_sel_confirmed', _wd_selected),
                        'drop_empty':     wd_drop_empty,
                        'order':          st.session_state.get('wd_result_order', {}),
                    }
                    _wd_result = _wd_build_zip(_wd_cfg)
                    _doc_name = (wd_title or 'Raport').replace(' ', '_')[:40]
                    _is_img_mode = (_wd_cfg.get('chart_mode') == 'image')
                    st.download_button(
                        '\u2b07\ufe0f Pobierz raport Word (.docx)' if _is_img_mode else '\u2b07\ufe0f Pobierz raport Word (.zip)',
                        data=_wd_result,
                        file_name=f'{_doc_name}.docx' if _is_img_mode else f'{_doc_name}_word.zip',
                        mime='application/vnd.openxmlformats-officedocument.wordprocessingml.document' if _is_img_mode else 'application/zip',
                        type='primary',
                        use_container_width=True,
                        key='wd_download',
                    )
                    if _is_img_mode:
                        st.success('\u2705 Raport wygenerowany! Otw\u00f3rz plik .docx i naci\u015bnij Ctrl+A \u2192 F9 aby od\u015bwie\u017cy\u0107 spis tre\u015bci.')
                    else:
                        st.success('\u2705 Raport wygenerowany! Po pobraniu archiwum ZIP: wypakuj do folderu, otw\u00f3rz `raport.docx` i naci\u015bnij Ctrl+A \u2192 F9 aby od\u015bwie\u017cy\u0107 spis tre\u015bci.')
                except ImportError:
                    st.error("Biblioteka `python-docx` nie jest zainstalowana. Uruchom: `pip install python-docx`")
                except Exception as _wd_err:
                    st.error(f"B\u0142\u0105d generowania: {_wd_err}")
                    st.exception(_wd_err)


# =============================================================
# PANEL ADMINA
# =============================================================
elif menu == "\U0001f512 Panel admina":
    _require_module_access("admin")
    module_header("\U0001f512", "Panel admina",
                  "Zarz\u0105dzanie u\u017cytkownikami, uprawnieniami i aktywno\u015bci\u0105")

    _adm_me = st.session_state.get("current_user_id")

    # -- przycisk odswiezenia danych --
    _adm_ref_c1, _adm_ref_c2 = st.columns([8, 2])
    with _adm_ref_c2:
        import datetime as _adm_dt
        _adm_now_str = _adm_dt.datetime.now().strftime("%H:%M:%S")
        _adm_ref_inner = st.columns([2, 1])
        _adm_ref_inner[0].caption("Dane z: " + _adm_now_str)
        if _adm_ref_inner[1].button(
            "\U0001f504",
            key="adm_refresh_btn",
            help="Od\u015bwie\u017c wszystkie dane panelu (sesje, aktywno\u015b\u0107, itp.)",
            use_container_width=True
        ):
            st.rerun()

    def _adm_gen_pw(length=12):
        import random
        _chars = string.ascii_letters + string.digits + "!@#%^&*"
        return "".join(random.SystemRandom().choices(_chars, k=length))

    (
        _adm_tab_users, _adm_tab_perms, _adm_tab_sess,
        _adm_tab_hist, _adm_tab_act, _adm_tab_stats, _adm_tab_cfg
    ) = st.tabs([
        "\U0001f464 U\u017cytkownicy",
        "\U0001f510 Uprawnienia",
        "\U0001f4f6 Aktywne sesje",
        "\U0001f4cb Historia logowa\u0144",
        "\U0001f4c8 Aktywno\u015b\u0107",
        "\U0001f4ca Statystyki",
        "\u2699\ufe0f Ustawienia",
    ])

    # ============================================================
    # TAB 1 -- UZYTKOWNICY
    # ============================================================
    with _adm_tab_users:
        st.markdown("### Dodaj nowego u\u017cytkownika")
        with st.form("adm_add_user_form", clear_on_submit=True):
            _au_c1, _au_c2, _au_c3 = st.columns(3)
            _au_username = _au_c1.text_input("Nazwa u\u017cytkownika *")
            _au_role     = _au_c2.selectbox("Rola", ["user", "admin"])
            _au_expires  = _au_c3.date_input(
                "Wygasa (puste = bez limitu)", value=None)
            _au_c4, _au_c5 = st.columns(2)
            _au_pw_mode   = _au_c4.radio(
                "Has\u0142o", ["Generuj losowe", "Wpisz r\u0119cznie"], horizontal=True)
            _au_pw_custom = _au_c5.text_input(
                "Has\u0142o (je\u015bli r\u0119czne)", type="password")
            _au_mcp = st.checkbox(
                "Wymu\u015b zmian\u0119 has\u0142a przy pierwszym logowaniu", value=True)
            _au_sub = st.form_submit_button(
                "\u2795 Utw\u00f3rz u\u017cytkownika", type="primary")

        if _au_sub:
            _au_uname_clean = (_au_username or "").strip().lower()
            if not _au_uname_clean:
                st.error("Podaj nazw\u0119 u\u017cytkownika.")
            elif _get_user_by_name(_au_uname_clean):
                st.error("U\u017cytkownik '" + _au_uname_clean + "' ju\u017c istnieje.")
            else:
                _au_pw_final = _adm_gen_pw() if _au_pw_mode == "Generuj losowe" else _au_pw_custom
                _au_pw_err = None if _au_mcp else _validate_password_policy(_au_pw_final, _au_uname_clean)
                if not _au_pw_final:
                    st.error("Podaj has\u0142o lub wybierz generowanie losowe.")
                elif _au_pw_err:
                    st.error("\u274c " + _au_pw_err)
                else:
                    _au_h, _au_s = _hash_password(_au_pw_final)
                    _au_exp_str = _au_expires.isoformat() if _au_expires else None
                    _au_new_id = get_db().execute(
                        "INSERT INTO users(username,password_hash,password_salt,role,"
                        "created_at,created_by,is_active,expires_at,must_change_password)"
                        " VALUES(?,?,?,?,?,?,1,?,?)",
                        (_au_uname_clean, _au_h, _au_s, _au_role,
                         _now_iso(), _adm_me, _au_exp_str, int(_au_mcp))
                    ).lastrowid
                    if _au_role == "admin":
                        for _mk in _MODULE_KEYS.keys():
                            get_db().execute(
                                "INSERT OR IGNORE INTO module_permissions"
                                "(user_id,module_key,granted) VALUES(?,?,1)",
                                (_au_new_id, _mk))
                    get_db().execute(
                        "INSERT INTO audit_log(actor_user_id,target_user_id,"
                        "event_type,details_json,created_at) VALUES(?,?,?,?,?)",
                        (_adm_me, _au_new_id, "create_user",
                         json.dumps({"username": _au_uname_clean, "role": _au_role}),
                         _now_iso()))
                    _au_ok_msg = "\u2705 Utworzono u\u017cytkownika **" + _au_uname_clean + "**"
                    if _au_pw_mode == "Generuj losowe":
                        _au_ok_msg += "\n\n**Has\u0142o (skopiuj teraz!):** `" + _au_pw_final + "`"
                    st.success(_au_ok_msg)
                    st.rerun()

        st.markdown("---")
        st.markdown("### Lista u\u017cytkownik\u00f3w")
        _au_all = get_db().execute(
            "SELECT id,username,role,is_active,expires_at,last_login_at,"
            "failed_login_count,locked_until,must_change_password FROM users"
            " ORDER BY role DESC,username"
        ).fetchall()

        for _au_row in _au_all:
            _au_uid   = _au_row["id"]
            _au_ud    = _au_row["username"]
            _au_act   = bool(_au_row["is_active"])
            _au_lkd   = _is_locked(_au_row)
            _au_exp_v = _au_row["expires_at"] or "bez limitu"
            _au_ll    = (_au_row["last_login_at"] or "\u2014")[:16].replace("T", " ")
            _au_ic    = "\U0001f7e2" if (_au_act and not _au_lkd) else ("\U0001f534" if _au_lkd else "\u26ab")
            _au_rb    = " [admin]" if _au_row["role"] == "admin" else " [user]"
            _au_mb    = " \U0001f511" if _au_row["must_change_password"] else ""
            with st.expander(
                _au_ic + " " + _au_ud + _au_rb + _au_mb
                + "   ostatnie log.: " + _au_ll, expanded=False
            ):
                _au_a1, _au_a2, _au_a3, _au_a4 = st.columns(4)
                # Aktywacja/deaktywacja
                if _au_act:
                    if _au_a1.button("\u26ab Dezaktywuj", key="adm_deact_" + str(_au_uid)):
                        get_db().execute("UPDATE users SET is_active=0 WHERE id=?", (_au_uid,))
                        get_db().execute(
                            "INSERT INTO audit_log(actor_user_id,target_user_id,event_type,"
                            "details_json,created_at) VALUES(?,?,?,?,?)",
                            (_adm_me, _au_uid, "deactivate",
                             json.dumps({"username": _au_ud}), _now_iso()))
                        st.rerun()
                else:
                    if _au_a1.button("\U0001f7e2 Aktywuj", key="adm_act_" + str(_au_uid)):
                        get_db().execute("UPDATE users SET is_active=1 WHERE id=?", (_au_uid,))
                        get_db().execute(
                            "INSERT INTO audit_log(actor_user_id,target_user_id,event_type,"
                            "details_json,created_at) VALUES(?,?,?,?,?)",
                            (_adm_me, _au_uid, "activate",
                             json.dumps({"username": _au_ud}), _now_iso()))
                        st.rerun()
                # Odblokowanie
                if _au_lkd:
                    if _au_a2.button("\U0001f513 Odblokuj", key="adm_unlock_" + str(_au_uid)):
                        get_db().execute(
                            "UPDATE users SET failed_login_count=0,locked_until=NULL WHERE id=?",
                            (_au_uid,))
                        get_db().execute(
                            "INSERT INTO audit_log(actor_user_id,target_user_id,event_type,"
                            "details_json,created_at) VALUES(?,?,?,?,?)",
                            (_adm_me, _au_uid, "unlock",
                             json.dumps({"username": _au_ud}), _now_iso()))
                        st.rerun()
                # Reset hasla
                if _au_a3.button("\U0001f504 Reset has\u0142a", key="adm_rpw_" + str(_au_uid)):
                    _au_npw = _adm_gen_pw()
                    _au_nh, _au_ns = _hash_password(_au_npw)
                    get_db().execute(
                        "UPDATE users SET password_hash=?,password_salt=?,"
                        "must_change_password=1 WHERE id=?",
                        (_au_nh, _au_ns, _au_uid))
                    get_db().execute(
                        "INSERT INTO audit_log(actor_user_id,target_user_id,event_type,"
                        "details_json,created_at) VALUES(?,?,?,?,?)",
                        (_adm_me, _au_uid, "reset_password",
                         json.dumps({"username": _au_ud}), _now_iso()))
                    st.success("Nowe has\u0142o **" + _au_ud + "**: `" + _au_npw + "`")
                # Usuwanie (tylko innych)
                if _au_uid != _adm_me:
                    if _au_a4.button("\U0001f5d1\ufe0f Usu\u0144",
                                     key="adm_del_" + str(_au_uid)):
                        get_db().execute("DELETE FROM users WHERE id=?", (_au_uid,))
                        get_db().execute(
                            "DELETE FROM module_permissions WHERE user_id=?", (_au_uid,))
                        get_db().execute(
                            "INSERT INTO audit_log(actor_user_id,target_user_id,event_type,"
                            "details_json,created_at) VALUES(?,?,?,?,?)",
                            (_adm_me, _au_uid, "delete_user",
                             json.dumps({"username": _au_ud}), _now_iso()))
                        st.success("U\u017cytkownik **" + _au_ud + "** usuni\u0119ty.")
                        st.rerun()
                # Edycja daty wygasniecia
                st.caption("Wygasa: " + str(_au_exp_v))
                _au_ne, _au_ec1, _au_ec2 = st.columns([2, 1, 1])
                _au_new_exp = _au_ne.date_input(
                    "Nowa data wygasni\u0119cia", value=None,
                    key="adm_exp_" + str(_au_uid))
                if _au_ec1.button("Ustaw", key="adm_setexp_" + str(_au_uid)):
                    _au_ev = _au_new_exp.isoformat() if _au_new_exp else None
                    get_db().execute(
                        "UPDATE users SET expires_at=? WHERE id=?", (_au_ev, _au_uid))
                    get_db().execute(
                        "INSERT INTO audit_log(actor_user_id,target_user_id,event_type,"
                        "details_json,created_at) VALUES(?,?,?,?,?)",
                        (_adm_me, _au_uid, "set_expires",
                         json.dumps({"username": _au_ud, "expires_at": _au_ev}),
                         _now_iso()))
                    st.success("Zaktualizowano dat\u0119 wygasni\u0119cia.")
                    st.rerun()
                if _au_ec2.button("Usu\u0144 limit", key="adm_delexp_" + str(_au_uid)):
                    get_db().execute(
                        "UPDATE users SET expires_at=NULL WHERE id=?", (_au_uid,))
                    get_db().execute(
                        "INSERT INTO audit_log(actor_user_id,target_user_id,event_type,"
                        "details_json,created_at) VALUES(?,?,?,?,?)",
                        (_adm_me, _au_uid, "remove_expires",
                         json.dumps({"username": _au_ud}), _now_iso()))
                    st.success("Usuni\u0119to limit czasu.")
                    st.rerun()

    # ============================================================
    # TAB 2 -- UPRAWNIENIA
    # ============================================================
    with _adm_tab_perms:
        if st.session_state.pop("_adm_perms_saved", False):
            st.success("\u2705 Uprawnienia zosta\u0142y zapisane.")
        st.markdown("### Matryca uprawnie\u0144 (bez kont admin \u2014 maj\u0105 pe\u0142ny dost\u0119p)")
        _prm_users = get_db().execute(
            "SELECT id,username FROM users WHERE role='user' AND is_active=1"
            " ORDER BY username"
        ).fetchall()
        if not _prm_users:
            st.info("Brak aktywnych u\u017cytkownik\u00f3w z rol\u0105 'user'.")
        else:
            _prm_mod_keys = [k for k in _MODULE_KEYS.keys() if k not in ("admin", "import")]
            _prm_mod_short = {
                "dashboard":    "Dashboard", "project":    "Projekt",
                "prep":         "Przygot.",  "analyses":   "Analizy",
                "waves":        "Fale",
                "regression":   "Regresja",  "anova":      "ANOVA",
                "normality":    "Normalno\u015b\u0107", "factor": "Czynnikowa",
                "cluster":      "Skupienia", "conjoint":   "Conjoint",
                "maxdiff":      "MaxDiff",   "wordcloud":  "Chmura S\u0142\u00f3w",
                "export_excel": "Excel",     "export_pptx":"PPT",
                "export_word":  "Word",
            }
            _prm_rows = []
            _prm_uid_map = {}
            for _pu in _prm_users:
                _pu_prms = {
                    r["module_key"]: bool(r["granted"])
                    for r in get_db().execute(
                        "SELECT module_key,granted FROM module_permissions WHERE user_id=?",
                        (_pu["id"],)).fetchall()
                }
                _prm_row = {_prm_mod_short.get(k, k): _pu_prms.get(k, False)
                            for k in _prm_mod_keys}
                _prm_row["__user__"] = _pu["username"]
                _prm_rows.append(_prm_row)
                _prm_uid_map[_pu["username"]] = _pu["id"]

            # --- szybkie przypisanie uprawnien per uzytkownik ---
            st.caption("Szybkie przypisanie \u2014 nadaj lub cofnij wszystkie uprawnienia dla u\u017cytkownika jednym klikni\u0119ciem:")
            _prm_quick_cols = st.columns(min(len(_prm_users), 5))
            for _pqi, _pqu in enumerate(_prm_users):
                with _prm_quick_cols[_pqi % min(len(_prm_users), 5)]:
                    st.markdown("**" + _pqu["username"] + "**")
                    _pq_c1, _pq_c2 = st.columns(2)
                    if _pq_c1.button(
                        "\u2705 Wszystkie",
                        key="adm_grant_all_" + str(_pqu["id"]),
                        use_container_width=True,
                        help="Przyznaj wszystkie uprawnienia u\u017cytkownikowi " + _pqu["username"]
                    ):
                        for _gmk in _prm_mod_keys:
                            _set_module_perm(_pqu["id"], _gmk, True, _adm_me)
                        get_db().execute(
                            "INSERT INTO audit_log(actor_user_id,target_user_id,"
                            "event_type,details_json,created_at) VALUES(?,?,?,?,?)",
                            (_adm_me, _pqu["id"], "grant_all_perms",
                             json.dumps({"username": _pqu["username"]}), _now_iso())
                        )
                        st.rerun()
                    if _pq_c2.button(
                        "\u274c \u017badne",
                        key="adm_revoke_all_" + str(_pqu["id"]),
                        use_container_width=True,
                        help="Cofnij wszystkie uprawnienia u\u017cytkownikowi " + _pqu["username"]
                    ):
                        for _rmk in _prm_mod_keys:
                            _set_module_perm(_pqu["id"], _rmk, False, _adm_me)
                        get_db().execute(
                            "INSERT INTO audit_log(actor_user_id,target_user_id,"
                            "event_type,details_json,created_at) VALUES(?,?,?,?,?)",
                            (_adm_me, _pqu["id"], "revoke_all_perms",
                             json.dumps({"username": _pqu["username"]}), _now_iso())
                        )
                        st.rerun()
            st.markdown("---")
            # --- matryca uprawnien (edycja szczegolowa) ---
            _prm_df = pd.DataFrame(_prm_rows).set_index("__user__")
            _prm_edited = st.data_editor(
                _prm_df, use_container_width=True, key="adm_perm_editor",
                column_config={c: st.column_config.CheckboxColumn(c)
                               for c in _prm_df.columns},
            )
            if st.button("\U0001f4be Zapisz uprawnienia", type="primary",
                         key="adm_save_perms"):
                for _pu_name, _pu_row in _prm_edited.iterrows():
                    _pu_id = _prm_uid_map.get(_pu_name)
                    if not _pu_id:
                        continue
                    for _sh_key, _granted in _pu_row.items():
                        _mk = next(
                            (k for k, v in _prm_mod_short.items() if v == _sh_key), None)
                        if _mk:
                            _set_module_perm(_pu_id, _mk, bool(_granted), _adm_me)
                    # Odswierz uprawnienia w session_state jesli to biezacy user
                    if _pu_id == st.session_state.get("current_user_id"):
                        st.session_state.current_user_perms = _load_user_perms(
                            _pu_id, "user")
                st.session_state._adm_perms_saved = True
                st.rerun()

    # ============================================================
    # TAB 3 -- AKTYWNE SESJE
    # ============================================================
    with _adm_tab_sess:
        st.markdown("### Aktywne sesje")
        _sess_idle = _get_setting_int("idle_timeout_minutes", 60)
        _sess_cutoff = (
            datetime.datetime.utcnow() - datetime.timedelta(minutes=_sess_idle)
        ).replace(microsecond=0).isoformat()
        _sess_rows = get_db().execute(
            "SELECT s.id,s.session_token,s.ip_address,s.geo_country,s.geo_city,"
            "s.started_at,s.last_seen_at,u.username"
            " FROM sessions s JOIN users u ON s.user_id=u.id"
            " WHERE s.ended_at IS NULL AND s.login_success=1"
            " AND s.last_seen_at > ?"
            " ORDER BY s.last_seen_at DESC",
            (_sess_cutoff,)
        ).fetchall()
        if not _sess_rows:
            st.info("Brak aktywnych sesji.")
        else:
            for _sr in _sess_rows:
                _sr_geo = (((_sr["geo_city"] or "") + ", " + (_sr["geo_country"] or "")).strip(", ") or "?")
                _sr_start = (_sr["started_at"] or "")[:16].replace("T", " ")
                _sr_last  = (_sr["last_seen_at"] or "")[:16].replace("T", " ")
                _sc1, _sc2 = st.columns([4, 1])
                _sc1.markdown(
                    "**" + _sr["username"] + "** \u2014 IP: `"
                    + (_sr["ip_address"] or "?") + "` \u2014 Geo: "
                    + _sr_geo + "  \n"
                    "\U0001f4c5 Start: " + _sr_start
                    + "  \u23f0 Ostatnia aktywno\u015b\u0107: " + _sr_last
                )
                if _sc2.button("\u21a9\ufe0f Wyloguj",
                               key="adm_fs_" + str(_sr["id"])):
                    _end_session(_sr["id"], "admin_force_logout")
                    get_db().execute(
                        "INSERT INTO audit_log(actor_user_id,target_user_id,event_type,"
                        "details_json,created_at) VALUES(?,?,?,?,?)",
                        (_adm_me, None, "force_logout",
                         json.dumps({"username": _sr["username"],
                                     "session_id": _sr["id"]}),
                         _now_iso()))
                    st.success("Sesja u\u017cytkownika **" + _sr["username"] + "** zako\u0144czona.")
                    st.rerun()
                st.markdown("---")

    # ============================================================
    # TAB 4 -- HISTORIA LOGOWAN
    # ============================================================
    with _adm_tab_hist:
        st.markdown("### Historia logowa\u0144")
        _hst_c1, _hst_c2, _hst_c3 = st.columns(3)
        _hst_user_filter = _hst_c1.text_input(
            "Filtruj u\u017cytkownika", placeholder="wszystkie", key="adm_hf_user")
        _hst_only_fail = _hst_c2.checkbox("Tylko nieudane", key="adm_hf_fail")
        _hst_limit = _hst_c3.selectbox("Max wierszy", [50, 100, 250, 500], key="adm_hf_lim")
        _hst_where = "WHERE 1=1"
        _hst_params = []
        if _hst_user_filter:
            _hst_where += " AND (u.username LIKE ? OR s.attempted_username LIKE ?)"
            _hst_params += ["%" + _hst_user_filter + "%", "%" + _hst_user_filter + "%"]
        if _hst_only_fail:
            _hst_where += " AND s.login_success=0"
        _hst_q = (
            "SELECT s.started_at,s.login_success,s.ip_address,s.geo_country,s.geo_city,"
            "s.logout_reason,s.attempted_username,u.username"
            " FROM sessions s LEFT JOIN users u ON s.user_id=u.id " + _hst_where
            + " ORDER BY s.started_at DESC LIMIT ?"
        )
        _hst_rows = get_db().execute(_hst_q, _hst_params + [_hst_limit]).fetchall()
        if not _hst_rows:
            st.info("Brak rekord\u00f3w spe\u0142niaj\u0105cych kryteria.")
        else:
            _hst_data = []
            for _hr in _hst_rows:
                _hst_uname = _hr["username"] or _hr["attempted_username"] or "?"
                _hst_geo = (((_hr["geo_city"] or "") + " " + (_hr["geo_country"] or "")).strip() or "?")
                _hst_data.append({
                    "Czas (UTC)":   (_hr["started_at"] or "")[:16].replace("T", " "),
                    "U\u017cytkownik": _hst_uname,
                    "Wynik":        "\u2705 OK" if _hr["login_success"] else "\u274c B\u0142\u0105d",
                    "IP":           _hr["ip_address"] or "?",
                    "Lokalizacja":  _hst_geo,
                    "Pow\u00f3d":   _hr["logout_reason"] or "",
                })
            st.dataframe(pd.DataFrame(_hst_data), use_container_width=True, hide_index=True)

    # ============================================================
    # TAB 5 -- AKTYWNOSC
    # ============================================================
    with _adm_tab_act:
        st.markdown("### Log aktywno\u015bci analiz")
        _act_c1, _act_c2, _act_c3 = st.columns(3)
        _act_uf  = _act_c1.text_input(
            "U\u017cytkownik", placeholder="wszyscy", key="adm_af_user")
        _act_mf  = _act_c2.selectbox(
            "Modu\u0142", ["-- wszystkie --"] + list(_MODULE_KEYS.keys()),
            key="adm_af_mod")
        _act_lim = _act_c3.selectbox(
            "Max wierszy", [100, 250, 500, 1000], key="adm_af_lim")
        _act_where = "WHERE 1=1"
        _act_params = []
        if _act_uf:
            _act_where += " AND u.username LIKE ?"
            _act_params.append("%" + _act_uf + "%")
        if _act_mf != "-- wszystkie --":
            _act_where += " AND a.module=?"
            _act_params.append(_act_mf)
        _act_rows = get_db().execute(
            "SELECT a.created_at,a.module,a.action,a.ip_address,a.metadata_json,"
            "u.username FROM activity_log a"
            " LEFT JOIN users u ON a.user_id=u.id "
            + _act_where + " ORDER BY a.created_at DESC LIMIT ?",
            _act_params + [_act_lim]
        ).fetchall()
        if not _act_rows:
            st.info("Brak rekord\u00f3w.")
        else:
            _act_data = []
            for _ar in _act_rows:
                _act_data.append({
                    "Czas (UTC)":   (_ar["created_at"] or "")[:16].replace("T", " "),
                    "U\u017cytkownik": _ar["username"] or "?",
                    "Modu\u0142":   _ar["module"] or "",
                    "Akcja":        _ar["action"] or "",
                    "IP":           _ar["ip_address"] or "",
                    "Metadane":     str(_ar["metadata_json"] or "")[:80],
                })
            _act_df = pd.DataFrame(_act_data)
            st.dataframe(_act_df, use_container_width=True, hide_index=True)
            # Eksport do Excela
            if st.button("\U0001f4be Eksportuj log do Excela", key="adm_act_export"):
                _act_buf = io.BytesIO()
                with pd.ExcelWriter(_act_buf, engine="xlsxwriter") as _act_wr:
                    _act_df.to_excel(_act_wr, sheet_name="Aktywnosc", index=False)
                st.download_button(
                    "\u2b07\ufe0f Pobierz log_aktywnosci.xlsx",
                    data=_act_buf.getvalue(),
                    file_name="log_aktywnosci.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )
        # Podsumowanie: top uzytkownicy
        st.markdown("---")
        st.markdown("**Top u\u017cytkownicy wg liczby analiz**")
        _top_rows = get_db().execute(
            "SELECT u.username, COUNT(a.id) as cnt"
            " FROM activity_log a JOIN users u ON a.user_id=u.id"
            " WHERE a.module != 'system'"
            " GROUP BY u.id ORDER BY cnt DESC LIMIT 10"
        ).fetchall()
        if _top_rows:
            _top_df = pd.DataFrame([{"U\u017cytkownik": r["username"],
                                      "Liczba analiz": r["cnt"]} for r in _top_rows])
            st.bar_chart(_top_df.set_index("U\u017cytkownik"))

    # ============================================================
    # TAB 6 -- STATYSTYKI
    # ============================================================
    with _adm_tab_stats:
        st.markdown("### Statystyki u\u017cytkowania")
        _st_c1, _st_c2 = st.columns(2)
        # Top modulow
        _st_mod_rows = get_db().execute(
            "SELECT module, COUNT(*) as cnt FROM activity_log"
            " WHERE module != 'system' GROUP BY module ORDER BY cnt DESC LIMIT 10"
        ).fetchall()
        if _st_mod_rows:
            with _st_c1:
                st.markdown("**Top modu\u0142\u00f3w**")
                _st_mod_df = pd.DataFrame([{"Modu\u0142": r["module"],
                                             "Analizy": r["cnt"]} for r in _st_mod_rows])
                st.bar_chart(_st_mod_df.set_index("Modu\u0142"))
        # Analizy per dzien (30 dni)
        _st_day_rows = get_db().execute(
            "SELECT substr(created_at,1,10) as day, COUNT(*) as cnt"
            " FROM activity_log WHERE module != 'system'"
            " AND created_at >= date('now','-30 days')"
            " GROUP BY day ORDER BY day"
        ).fetchall()
        if _st_day_rows:
            with _st_c2:
                st.markdown("**Analizy / dzie\u0144 (ostatnie 30 dni)**")
                _st_day_df = pd.DataFrame([{"Dzie\u0144": r["day"],
                                             "Analizy": r["cnt"]} for r in _st_day_rows])
                st.bar_chart(_st_day_df.set_index("Dzie\u0144"))
        # Logowania per kraj
        _st_geo_rows = get_db().execute(
            "SELECT geo_country, COUNT(*) as cnt FROM sessions"
            " WHERE login_success=1 AND geo_country IS NOT NULL"
            " GROUP BY geo_country ORDER BY cnt DESC"
        ).fetchall()
        if _st_geo_rows:
            st.markdown("**Logowania wg kraju**")
            _st_geo_df = pd.DataFrame([{"Kraj": r["geo_country"],
                                         "Logowania": r["cnt"]} for r in _st_geo_rows])
            try:
                _st_fig = px.bar(_st_geo_df, x="Kraj", y="Logowania",
                                 title="Logowania wg kraju")
                st.plotly_chart(_st_fig, use_container_width=True)
            except Exception:
                st.dataframe(_st_geo_df, use_container_width=True, hide_index=True)

    # ============================================================
    # TAB 7 -- USTAWIENIA
    # ============================================================
    with _adm_tab_cfg:
        st.markdown("### Ustawienia systemowe")
        _cfg_rows = get_db().execute("SELECT key,value FROM settings ORDER BY key").fetchall()
        _cfg_labels = {
            "idle_timeout_minutes": "Timeout bezczynno\u015bci (minuty)",
            "lockout_minutes":      "Czas blokady konta po b\u0142\u0119dach (minuty)",
            "max_fail_attempts":    "Max pr\u00f3b logowania przed blokad\u0105",
            "min_pw_length":        "Minimalna d\u0142ugo\u015b\u0107 has\u0142a (znaki)",
            "rate_limit_window":    "Okno rate-limit (minuty)",
        }
        _cfg_form_vals = {r["key"]: r["value"] for r in _cfg_rows}
        with st.form("adm_settings_form"):
            _cfg_inputs = {}
            for _ck, _cv in _cfg_form_vals.items():
                _cfg_inputs[_ck] = st.number_input(
                    _cfg_labels.get(_ck, _ck),
                    value=int(_cv),
                    min_value=1,
                    key="adm_cfg_" + _ck,
                )
            _cfg_sub = st.form_submit_button(
                "\U0001f4be Zapisz ustawienia", type="primary")
        if _cfg_sub:
            for _ck, _cv in _cfg_inputs.items():
                _set_setting(_ck, str(int(_cv)))
            st.success("\u2705 Ustawienia zapisane.")
        st.markdown("---")
        st.markdown("**Audit log (ostatnie 50 zdarze\u0144)**")
        _aud_rows = get_db().execute(
            "SELECT a.created_at,a.event_type,a.details_json,"
            "ua.username as actor,ut.username as target"
            " FROM audit_log a"
            " LEFT JOIN users ua ON a.actor_user_id=ua.id"
            " LEFT JOIN users ut ON a.target_user_id=ut.id"
            " ORDER BY a.created_at DESC LIMIT 50"
        ).fetchall()
        if _aud_rows:
            _aud_data = []
            for _ar in _aud_rows:
                _aud_data.append({
                    "Czas":      (_ar["created_at"] or "")[:16].replace("T", " "),
                    "Zdarzenie": _ar["event_type"] or "",
                    "Aktor":     _ar["actor"] or "system",
                    "Cel":       _ar["target"] or "",
                    "Szczeg\u00f3\u0142y": str(_ar["details_json"] or "")[:100],
                })
            st.dataframe(pd.DataFrame(_aud_data), use_container_width=True,
                         hide_index=True)

else:
    st.info("\U0001f448 Wybierz modu\u0142 z menu bocznego.")

