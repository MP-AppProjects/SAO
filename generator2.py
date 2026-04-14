# -*- coding: utf-8 -*-
import streamlit as st
import pandas as pd
import numpy as np
import pyreadstat
import tempfile
import io
import json
import re
import copy
import string
import datetime
import plotly.express as px
import plotly.graph_objects as go
import scipy.stats as stats
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.cm as cm
from collections import defaultdict
from sklearn.cluster import KMeans
from sklearn.preprocessing import StandardScaler
from scipy.cluster.hierarchy import linkage, dendrogram, fcluster
import statsmodels.api as sm
from statsmodels.stats.outliers_influence import variance_inflation_factor
from factor_analyzer import FactorAnalyzer
from factor_analyzer.factor_analyzer import calculate_kmo, calculate_bartlett_sphericity

st.set_page_config(page_title="System Analiz Openfield (SAO)", layout="wide", page_icon="\U0001f4ca")

# =============================================================
# LOGOWANIE
# =============================================================
# -------------------------------------------------------------
# CSS -- minimalistyczny, profesjonalny styl
# -------------------------------------------------------------
st.markdown("""
<style>
    /* \u2500\u2500 Extra top padding so content doesn\u2019t hide under Streamlit\u2019s fixed header \u2500\u2500 */
    .block-container { padding-top: 3.5rem !important; }

    /* \u2500\u2500 Sticky tab bar that stays below the Streamlit header (58px) \u2500\u2500 */
    .stTabs [data-baseweb="tab-list"] {
        position: sticky;
        top: 58px;
        z-index: 99;
        background: white;
        padding-top: 4px;
        padding-bottom: 3px;          /* room for the active-tab underline */
        overflow-x: auto;
        overflow-y: visible;          /* must be visible so underline isn\u2019t clipped */
        flex-wrap: nowrap;
        scrollbar-width: thin;
        scrollbar-color: #2E75B6 #f0f0f0;
        box-shadow: 0 1px 0 0 #e6e6e6; /* replicate the baseline separator */
    }

    /* \u2500\u2500 Compact tabs: smaller font + tight padding \u2500\u2500 */
    .stTabs [data-baseweb="tab"] {
        font-size: 0.78rem;
        padding: 5px 10px;
        white-space: nowrap;
    }

    /* Chrome / Edge / Safari scrollbar */
    .stTabs [data-baseweb="tab-list"]::-webkit-scrollbar { height: 4px; }
    .stTabs [data-baseweb="tab-list"]::-webkit-scrollbar-track {
        background: #f0f0f0; border-radius: 2px;
    }
    .stTabs [data-baseweb="tab-list"]::-webkit-scrollbar-thumb {
        background: #2E75B6; border-radius: 2px;
    }

    .metric-card {
        background: #f8f9fa; border-radius: 8px; padding: 14px 18px;
        border-left: 4px solid #2E75B6; margin-bottom: 8px;
    }
    .metric-card b { color: #2E75B6; font-size: 1.1rem; }
    .sig-green { color: #006100; font-weight: bold; }
    .sig-red   { color: #C00000; }
    .sig-orange{ color: #E36C09; }
    .section-header {
        background: linear-gradient(90deg, #2E75B6, #1F4E79);
        color: white; padding: 8px 16px; border-radius: 6px;
        font-size: 0.95rem; font-weight: bold; margin-bottom: 12px;
    }
    div[data-testid="stSidebarNav"] { display: none; }
    .stAlert { border-radius: 8px; }

</style>
""", unsafe_allow_html=True)

# -------------------------------------------------------------
# FUNKCJE POMOCNICZE
# -------------------------------------------------------------



@st.cache_data
def load_spss_data(uploaded_file):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".sav") as tmp:
        tmp.write(uploaded_file.getvalue())
        tmp_path = tmp.name
    df, meta = pyreadstat.read_sav(tmp_path)
    df_labeled = df.copy()
    for col, labels in meta.variable_value_labels.items():
        if col in df_labeled.columns:
            ordered_cats = [labels[k] for k in sorted(labels.keys())]
            df_labeled[col] = df_labeled[col].map(labels).fillna(df_labeled[col])
            existing = df_labeled[col].dropna().unique()
            extra = [e for e in existing if e not in ordered_cats]
            df_labeled[col] = pd.Categorical(df_labeled[col], categories=ordered_cats + extra, ordered=True)
    return df, df_labeled, meta


# ------------------------------------------------------------------
# Excel compatibility layer -- mimics pyreadstat meta interface
# ------------------------------------------------------------------
class ExcelMeta:
    """Thin compatibility wrapper so all meta_orig references work for Excel data."""
    def __init__(self, columns, col_labels=None):
        col_labels = col_labels or {}
        self.column_names_to_labels = {c: col_labels.get(c, c) for c in columns}
        self.variable_value_labels  = {}   # no SPSS value labels for Excel

    def get(self, key, default=None):
        return self.column_names_to_labels.get(key, default)


@st.cache_data
def load_excel_data(uploaded_file, sheet_name, col_type_overrides_json="{}", custom_missing_json="{}"):
    """
    Load an Excel file.
    Row 1 = variable names, row 2+ = data.
    col_type_overrides_json: JSON string of {col: 'numeric'|'categorical'}
    custom_missing_json:     JSON string of {col: [val, ...]} missing values per column
    (JSON strings are hashable, enabling proper Streamlit cache keying)
    Returns (df_raw, df_labeled, ExcelMeta)
    """
    col_type_overrides = json.loads(col_type_overrides_json)
    custom_missing     = json.loads(custom_missing_json)
    df = pd.read_excel(uploaded_file, sheet_name=sheet_name, header=0)

    # Clean column names
    df.columns = [str(c).strip() for c in df.columns]

    # Auto-detect and apply types
    df_raw = df.copy()
    df_labeled = df.copy()

    # Collect text\u2192numeric mappings to store as value labels
    _text_to_num_maps = {}  # {col: {1: 'Kobieta', 2: 'M\u0119\u017cczyzna', ...}}

    def _missing_str_set(col):
        """Build set of string representations of missing values for a column."""
        m_vals = custom_missing.get(col, [])
        s = set()
        for v in m_vals:
            s.add(str(v))
            try:
                s.add(str(int(float(v))))
                s.add(str(float(v)))
            except (ValueError, TypeError):
                pass
        return s

    for col in df.columns:
        override = col_type_overrides.get(col)
        if override == 'numeric':
            # Try direct numeric conversion first
            numeric_attempt = pd.to_numeric(df_raw[col], errors='coerce')
            non_null = df_raw[col].dropna()
            survived = numeric_attempt.notna().sum()
            if len(non_null) == 0 or survived / len(non_null) >= 0.9:
                df_raw[col]    = numeric_attempt
                df_labeled[col] = numeric_attempt.copy()
            else:
                # Column is text \u2014 encode as consecutive integers 1, 2, 3...
                # Skip values that are defined as missing
                miss_strs = _missing_str_set(col)
                unique_vals = sorted(
                    [v for v in df_raw[col].dropna().unique()
                     if str(v) not in miss_strs],
                    key=lambda x: str(x)
                )
                code_map  = {v: i + 1 for i, v in enumerate(unique_vals)}
                label_map = {i + 1: str(v) for i, v in enumerate(unique_vals)}
                df_raw[col]    = df_raw[col].map(code_map)   # missing vals \u2192 NaN
                df_labeled[col] = df_raw[col].copy()
                _text_to_num_maps[col] = label_map
        elif override == 'categorical':
            df_raw[col]    = df_raw[col].astype(str).where(df_raw[col].notna(), np.nan)
            df_labeled[col] = df_raw[col].copy()
        else:
            # Auto: try numeric first
            numeric_attempt = pd.to_numeric(df_raw[col], errors='coerce')
            non_null = df_raw[col].dropna()
            if len(non_null) > 0:
                numeric_rate = numeric_attempt.notna().sum() / len(non_null)
                if numeric_rate >= 0.9:
                    df_raw[col]    = numeric_attempt
                    df_labeled[col] = numeric_attempt.copy()
                # else: leave as object (categorical)

    meta = ExcelMeta(df.columns.tolist())
    meta._text_to_num_maps = _text_to_num_maps
    return df_raw, df_labeled, meta

def auto_detect_mrs(df_raw):
    binary_cols = [c for c in df_raw.columns if set(df_raw[c].dropna().unique()).issubset({0, 1}) and len(set(df_raw[c].dropna().unique())) > 0]
    mrs_candidates = defaultdict(list)
    for col in binary_cols:
        prefix = col.rsplit('_', 1)[0] if '_' in col else col[:-1]
        mrs_candidates[prefix].append(col)
    return {k: v for k, v in mrs_candidates.items() if len(v) > 1}

def auto_detect_matrix(df_raw):
    """Detect matrix/battery questions by shared prefix -- works for numeric AND text columns."""
    binary_cols = set(c for c in df_raw.columns if set(df_raw[c].dropna().unique()).issubset({0, 1}))
    # All non-binary columns (numeric and text)
    candidates = [c for c in df_raw.columns if c not in binary_cols]
    matrix_candidates = defaultdict(list)
    for col in candidates:
        if '_' in col:
            prefix = col.rsplit('_', 1)[0]
        elif len(col) > 1 and col[-1].isdigit():
            prefix = col[:-1]
        else:
            continue
        matrix_candidates[prefix].append(col)
    return {k: sorted(v) for k, v in matrix_candidates.items() if len(v) >= 2}

def build_matrix_table(df, df_raw, matrix_cols, var_labels, weights, meta_vvl, custom_val_labels=None):
    """
    Build a matrix frequency table with TRANSPOSED layout:
      Rows    = scale values / categories  (e.g. 1, 2, 3 ... or 'Tak','Nie')
      Columns = subquestions (variable names / labels)
      Cells   = N and % per subquestion \u00d7 value combination

    Also appends a combined 'Baza (N) / Suma (%)' summary row with N and % side by side.
    custom_val_labels: {var_name: {str(code): new_label}} overrides display labels.

    Returns: df_out, all_cats (display labels), sub_labels
    """
    w = weights
    if custom_val_labels is None:
        custom_val_labels = {}

    # -- 1. Collect all unique raw category values across the battery --
    raw_cats_set = []
    for col in matrix_cols:
        series = (df[col] if col in df.columns else df_raw[col]).dropna()
        for cat in series.unique():
            cat_str = str(cat)
            if cat_str not in raw_cats_set:
                raw_cats_set.append(cat_str)
    try:
        raw_cats_sorted = sorted(raw_cats_set, key=lambda x: float(x))
    except (ValueError, TypeError):
        raw_cats_sorted = sorted(raw_cats_set)

    # -- 2. Build display label map for categories (apply custom_val_labels) --
    # Since value labels may differ per column, build a unified best-effort map
    # using the first column's custom/SPSS labels as reference, then merge.
    cat_display = {}   # raw_str ? display_str
    for raw_str in raw_cats_sorted:
        cat_display[raw_str] = raw_str   # default: show raw value

    # Apply SPSS value labels from the first column that has them
    for col in matrix_cols:
        spss_vvl = meta_vvl.get(col, {})
        col_custom = custom_val_labels.get(col, {})
        for raw_str in raw_cats_sorted:
            # Custom label overrides SPSS label
            if raw_str in col_custom:
                cat_display[raw_str] = col_custom[raw_str]
            elif raw_str not in col_custom:
                # Try SPSS numeric key
                try:
                    num_key = float(raw_str)
                    if num_key in spss_vvl and cat_display[raw_str] == raw_str:
                        cat_display[raw_str] = spss_vvl[num_key]
                except (ValueError, TypeError):
                    if raw_str in spss_vvl and cat_display[raw_str] == raw_str:
                        cat_display[raw_str] = spss_vvl[raw_str]

    display_cats = [cat_display[r] for r in raw_cats_sorted]   # display labels, in order

    # -- 3. For each subquestion: count N and % per category --
    sub_labels = []
    data_n   = {}
    data_pct = {}
    data_base = {}

    for col in matrix_cols:
        sub_lbl = var_labels.get(col, col)
        if sub_lbl in sub_labels:
            sub_lbl = f"{sub_lbl} [{col}]"
        sub_labels.append(sub_lbl)

        series = (df[col] if col in df.columns else df_raw[col])
        missing_mask = series.isna()
        base_w = float(w[~missing_mask].sum())
        data_base[sub_lbl] = base_w

        counts = {}
        for raw_str, disp_str in zip(raw_cats_sorted, display_cats):
            mask = (series.astype(str) == raw_str) & (~missing_mask)
            counts[disp_str] = float(w[mask].sum())
        data_n[sub_lbl]   = counts
        data_pct[sub_lbl] = {disp: (v / base_w * 100 if base_w > 0 else 0.0)
                             for disp, v in counts.items()}

    # -- 4. Build output DataFrame --
    # Columns interleaved: SubA [N], SubA [%], SubB [N], SubB [%], ...
    interleaved_cols = []
    for lbl in sub_labels:
        interleaved_cols.append(f"{lbl} [N]")
        interleaved_cols.append(f"{lbl} [%]")

    # Rows: one per display category + one combined summary row "Baza (N) / Suma (%)"
    SUMMARY_ROW = "Baza (N) / Suma (%)"
    all_rows = display_cats + [SUMMARY_ROW]
    df_out = pd.DataFrame(index=all_rows, columns=interleaved_cols, dtype=object)

    for sub_lbl in sub_labels:
        col_sum_pct = 0.0
        for disp_str in display_cats:
            n_val   = data_n[sub_lbl][disp_str]
            pct_val = data_pct[sub_lbl][disp_str]
            df_out.loc[disp_str, f"{sub_lbl} [N]"] = n_val
            df_out.loc[disp_str, f"{sub_lbl} [%]"] = pct_val
            col_sum_pct += pct_val
        # Summary row: N = base respondents, % = sum of percentages (100%)
        df_out.loc[SUMMARY_ROW, f"{sub_lbl} [N]"] = data_base[sub_lbl]
        df_out.loc[SUMMARY_ROW, f"{sub_lbl} [%]"] = round(col_sum_pct, 1)

    return df_out, display_cats, sub_labels

def apply_segmentations(df_raw, df, meta_labels, segmentations_list):
    for seg in segmentations_list:
        cols, k, name = seg['vars'], seg['k'], seg['name']
        X = df_raw[cols].copy()
        for c in cols:
            m_val = X[c].mean()
            X[c] = X[c].fillna(m_val if pd.notna(m_val) else 0)
        scaler = StandardScaler()
        X_scaled = scaler.fit_transform(X)
        kmeans = KMeans(n_clusters=k, random_state=42, n_init=10)
        clusters = kmeans.fit_predict(X_scaled) + 1
        df_raw[name] = clusters
        df[name] = [f"Segment {c}" for c in clusters]
        df[name] = pd.Categorical(df[name], categories=[f"Segment {i}" for i in range(1, k + 1)], ordered=True)
        meta_labels[name] = f"Segmentacja K-Means ({k} grup)"

def apply_recodings(df_raw, df, var_labels, recodings_list):
    """Apply variable recodings stored in session state.
    Supports both numeric and text source variables.
    Output type is determined by the mapped values (numeric if all parseable, else text).
    """
    for rec in recodings_list:
        src = rec['source']
        new_name = rec['new_name']
        mapping = rec['mapping']   # {old_val_str: new_val_str}
        label = rec.get('label', new_name)
        output_type = rec.get('output_type', 'auto')  # 'numeric', 'text', 'auto'
        if src not in df_raw.columns:
            continue

        src_series = df_raw[src].copy().astype(str).str.strip()

        # Build a normalised lookup: str(old_val) -> new_val_str
        lookup = {str(k).strip(): str(v) for k, v in mapping.items()}

        new_col = src_series.map(lookup)   # unmapped ? NaN

        # Decide output type
        if output_type == 'numeric':
            df_raw[new_name] = pd.to_numeric(new_col, errors='coerce')
            df[new_name] = df_raw[new_name].copy()
        elif output_type == 'text':
            df_raw[new_name] = new_col
            df[new_name] = new_col
        else:  # auto: numeric if all non-null values are parseable
            numeric_attempt = pd.to_numeric(new_col, errors='coerce')
            if new_col.dropna().empty or numeric_attempt.notna().sum() == new_col.notna().sum():
                df_raw[new_name] = numeric_attempt
                df[new_name] = numeric_attempt.copy()
            else:
                df_raw[new_name] = new_col
                df[new_name] = new_col

        var_labels[new_name] = label

def apply_cleaning_ops(df_raw, df, cleaning_ops_list):
    """
    Apply stored cleaning operations in-place to the original columns.
    Each entry in cleaning_ops_list:
      {
        'cols': [col1, col2, ...],   # columns to clean
        'ops':  {                    # which operations to apply
            'strip': bool,
            'dbl_sp': bool,
            'tabs': bool,
            'newlines': bool,
            'quotes': bool,
            'case': 'none'|'upper'|'lower'|'title',
            'special': bool,
        }
      }
    """
    QUOTES_MAP = [
        ('\u201c', '"'), ('\u201d', '"'), ('\u201e', '"'),
        ('\u2018', "'"), ('\u2019', "'"), ('\u201a', "'"),
    ]
    for entry in cleaning_ops_list:
        ops  = entry.get('ops', {})
        cols = entry.get('cols', [])
        for col in cols:
            if col not in df_raw.columns:
                continue
            was_null = df_raw[col].isna()
            series = df_raw[col].astype(str).copy()

            if ops.get('strip'):
                series = series.str.strip()
            if ops.get('dbl_sp'):
                series = series.str.replace(r' {2,}', ' ', regex=True)
            if ops.get('tabs'):
                series = series.str.replace('\t', ' ', regex=False)
            if ops.get('newlines'):
                series = series.str.replace(r'[\n\r]', ' ', regex=True)
            if ops.get('quotes'):
                for old_q, new_q in QUOTES_MAP:
                    series = series.str.replace(old_q, new_q, regex=False)
            case = ops.get('case', 'none')
            if   case == 'upper': series = series.str.upper()
            elif case == 'lower': series = series.str.lower()
            elif case == 'title': series = series.str.title()
            if ops.get('special'):
                series = series.str.replace(r'[^\w\s]', '', regex=True)
                series = series.str.replace('_', '', regex=False)

            # Write back, preserving original NaN positions.
            # Convert Categorical columns to object first \u2014 cleaned values
            # may not exist in the original category set.
            if hasattr(df_raw[col], 'cat'):
                df_raw[col] = df_raw[col].astype(object)
            if hasattr(df[col], 'cat'):
                df[col] = df[col].astype(object)

            df_raw.loc[~was_null, col] = series[~was_null]
            df.loc[~was_null, col]     = series[~was_null]


def get_var_display_name(var_name, meta):
    # MRS / matrix set virtual names
    if 'mrs_sets' in st.session_state and var_name in st.session_state.mrs_sets:
        return f"[{var_name}] Zestaw Wielokrotnych Odpowiedzi"
    if 'matrix_sets' in st.session_state and var_name in st.session_state.matrix_sets:
        return f"[{var_name}] Pytanie matrycowe"

    # Resolve label
    label = meta.get(var_name, var_name) if isinstance(meta, dict) else meta.column_names_to_labels.get(var_name, var_name)
    if len(label) > 60:
        label = label[:57] + "..."

    # Mark derived (added during session) with a visible prefix
    # 'original_cols' is set at load time; fall back gracefully if not yet defined
    try:
        is_derived = var_name not in original_cols
    except NameError:
        is_derived = False

    prefix = "[+] " if is_derived else ""
    return f"{prefix}[{var_name}] {label}"

def get_weighted_stats(x, w):
    mask = ~np.isnan(x)
    x_valid, w_valid = x[mask], w[mask]
    sum_w = w_valid.sum()
    if sum_w == 0:
        return np.nan, np.nan, 0
    mean = (x_valid * w_valid).sum() / sum_w
    var = (w_valid * (x_valid - mean) ** 2).sum() / sum_w
    ess = (sum_w ** 2) / (w_valid ** 2).sum() if (w_valid ** 2).sum() > 0 else 0
    return mean, var, ess

def apply_means_sig_testing(df_means, df_vars, df_ess):
    cols = df_means.columns
    letters = list(string.ascii_uppercase)
    col_letters = {c: letters[i % 26] * (i // 26 + 1) for i, c in enumerate(cols)}
    sig_df = pd.DataFrame("", index=df_means.index, columns=cols)
    for r in df_means.index:
        for i, c1 in enumerate(cols):
            for j, c2 in enumerate(cols):
                if i >= j: continue
                m1, m2 = df_means.loc[r, c1], df_means.loc[r, c2]
                v1, v2 = df_vars.loc[r, c1], df_vars.loc[r, c2]
                n1, n2 = df_ess.loc[r, c1], df_ess.loc[r, c2]
                if pd.isna(m1) or pd.isna(m2) or n1 < 2 or n2 < 2 or (v1 == 0 and v2 == 0): continue
                se = np.sqrt(v1 / n1 + v2 / n2)
                if se == 0: continue
                t_stat = (m1 - m2) / se
                if t_stat > 1.96:   sig_df.loc[r, c1] += " " + col_letters[c2]
                elif t_stat < -1.96: sig_df.loc[r, c2] += " " + col_letters[c1]
    return sig_df, col_letters

def apply_sig_testing(df_pct, df_n):
    bases = df_n.loc['Suma']
    cols = [c for c in df_pct.columns if c != 'Suma']
    letters = list(string.ascii_uppercase)
    col_letters = {c: letters[i % 26] * (i // 26 + 1) for i, c in enumerate(cols)}
    sig_df = pd.DataFrame("", index=df_pct.index, columns=df_pct.columns)
    for r in df_pct.index:
        if r in ['Suma', 'Braki danych', 'Braki danych (wykluczone z tabeli)']: continue
        for i, c1 in enumerate(cols):
            for j, c2 in enumerate(cols):
                if i >= j: continue
                p1 = df_pct.loc[r, c1] / 100.0 if pd.notna(df_pct.loc[r, c1]) else np.nan
                p2 = df_pct.loc[r, c2] / 100.0 if pd.notna(df_pct.loc[r, c2]) else np.nan
                n1, n2 = bases[c1], bases[c2]
                if pd.isna(p1) or pd.isna(p2) or n1 == 0 or n2 == 0: continue
                p_pool = ((p1 * n1) + (p2 * n2)) / (n1 + n2)
                if p_pool == 0 or p_pool == 1: continue
                se = np.sqrt(p_pool * (1 - p_pool) * (1 / n1 + 1 / n2))
                if se == 0: continue
                z = (p1 - p2) / se
                if z > 1.96:   sig_df.loc[r, c1] += " " + col_letters[c2]
                elif z < -1.96: sig_df.loc[r, c2] += " " + col_letters[c1]
    return sig_df, col_letters

def module_header(icon, title, subtitle=""):
    """Render a blue gradient banner header identical to the Dashboard banner."""
    # Back button \u2014 above the banner, left-aligned
    if st.button("\u2190 Powr\u00f3\u0107 do Dashboardu", key=f"back_dash_{title}"):
        st.session_state.nav_to = "\U0001f3e0 Dashboard"
        st.rerun()

    sub_html = (f'<p style="margin:6px 0 0;opacity:.85;font-size:0.95rem;">{subtitle}</p>'
                if subtitle else "")
    st.markdown(f"""
<div style="background:linear-gradient(90deg,#1F4E79,#2E75B6);
     padding:22px 32px;border-radius:10px;margin-bottom:20px;color:white;">
  <h2 style="margin:0;font-size:1.55rem;">{icon} {title}</h2>
  {sub_html}
</div>
""", unsafe_allow_html=True)


def get_streamlit_format(df):
    format_dict = {}
    for col in df.columns:
        if df[col].apply(lambda x: isinstance(x, str) and ("%" in x or any(c.isalpha() for c in x if c not in ['N', 'a', 'A']))).any(): continue
        if "%" in str(col).lower() or "procent" in str(col).lower():
            format_dict[col] = lambda x: f"{x:.1f}%" if pd.notnull(x) and not isinstance(x, str) else str(x)
        else:
            format_dict[col] = lambda x: f"{x:.2f}" if pd.notnull(x) and not isinstance(x, str) and abs(x - round(x)) > 0.01 else (f"{x:.0f}" if pd.notnull(x) and not isinstance(x, str) else str(x))
    return format_dict

def safe_style(df):
    """
    Apply get_streamlit_format styling safely.
    Pandas Styler crashes on non-unique index or columns \u2014 deduplicate before styling.
    """
    d = df.copy()
    # Deduplicate index
    if not d.index.is_unique:
        seen = {}
        new_idx = []
        for v in d.index:
            k = str(v)
            if k in seen:
                seen[k] += 1
                new_idx.append(f"{v} ({seen[k]})")
            else:
                seen[k] = 0
                new_idx.append(v)
        d.index = new_idx
    # Deduplicate columns
    if not d.columns.is_unique:
        seen = {}
        new_cols = []
        for v in d.columns:
            k = str(v)
            if k in seen:
                seen[k] += 1
                new_cols.append(f"{v} ({seen[k]})")
            else:
                seen[k] = 0
                new_cols.append(v)
        d.columns = new_cols
    return d.style.format(get_streamlit_format(d))

# ---------------------------------------------------------------------------
# Module-level style/format helpers (defined once, reused across all reruns)
# ---------------------------------------------------------------------------

def _to_float_pct(x):
    """Strip % and letter suffixes, return float or original value."""
    if isinstance(x, str):
        clean = x.replace('%', '').strip().split()[0] if x.strip() else ''
        try:
            return float(clean)
        except (ValueError, TypeError):
            return x
    return x

def _fmt_cell(v):
    """Format a matrix table cell for display."""
    if v == "" or (isinstance(v, float) and np.isnan(v)):
        return ""
    if isinstance(v, float):
        return f"{v:.1f}"
    return str(v)

def _style_p(val):
    """Color p-value cells: green if significant, red if not."""
    try:
        v = float(val)
        if v < 0.001: return 'color:#006100;font-weight:bold'
        if v < 0.05:  return 'color:#006100'
        return 'color:#C00000'
    except (ValueError, TypeError):
        return ''

def _style_vif(val):
    """Color VIF cells: red > 10, orange > 5, green otherwise."""
    try:
        v = float(val)
        if v > 10: return 'color:#C00000;font-weight:bold'
        if v > 5:  return 'color:#E36C09'
        return 'color:#006100'
    except (ValueError, TypeError):
        return ''

def _style_loading(val):
    """Bold + green background for factor loadings >= 0.4."""
    try:
        if abs(float(val)) >= 0.4:
            return 'font-weight:bold;background-color:#E2EFDA'
        return ''
    except (ValueError, TypeError):
        return ''

def _color_pair_row(row):
    """Color correlation pair rows by strength."""
    abs_r = abs(row['r'])
    if abs_r >= 0.7:
        c = '#E2EFDA' if row['r'] > 0 else '#FCE4D6'
    else:
        c = '#FFFACD'
    return [f'background-color: {c}'] * len(row)

def _color_sig(row):
    """Green background for significant logistic regression rows."""
    color = '#E2EFDA' if row.get('Istotny') == 'Tak' else ''
    return [f'background-color: {color}'] * len(row)

def _make_style_matrix_row(sumrow_label):
    """Factory: returns a styler for matrix table rows."""
    def _style(row):
        if row.name == sumrow_label:
            return ['background-color:#E2EFDA; font-weight:bold'] * len(row)
        return [''] * len(row)
    return _style

def _make_color_corr_cell(threshold):
    """Factory: returns a cell styler for correlation matrix given threshold."""
    def _style(val):
        try:
            v = float(str(val).split()[0])
        except (ValueError, TypeError):
            return ''
        if abs(v) >= 1.0:
            return ''
        abs_v = abs(v)
        if abs_v >= 0.7:
            bg = '#C00000' if v < 0 else '#375623'
            return f'background-color: {bg}; color: white; font-weight: bold'
        elif abs_v >= threshold:
            bg = '#FCE4D6' if v < 0 else '#E2EFDA'
            return f'background-color: {bg}; font-weight: bold'
        return ''
    return _style

def _make_style_md(n_rows):
    """Factory: returns a MaxDiff row styler knowing total row count."""
    def _style(row):
        if row['Ranking'] == 1:
            return ['background-color:#E2EFDA; font-weight:bold'] * len(row)
        if row['Ranking'] == n_rows:
            return ['background-color:#FCE4D6'] * len(row)
        return [''] * len(row)
    return _style

def calculate_rim_weights(df, target_dict, max_iterations=50):
    weights = np.ones(len(df))
    for iteration in range(max_iterations):
        max_error = 0
        for var, targets in target_dict.items():
            for cat, target_pct in targets.items():
                mask = (df[var] == cat)
                if mask.sum() == 0: continue
                current_pct = weights[mask].sum() / weights.sum()
                if current_pct > 0:
                    adjustment = target_pct / current_pct
                    weights[mask] *= adjustment
                    max_error = max(max_error, abs(target_pct - current_pct))
        if max_error < 0.001: break
    return weights

def calculate_correlations(df, cols, weights=None, method='pearson'):
    """
    Compute correlation matrix with optional case weights.
    For Pearson: uses weighted covariance matrix (consistent with SPSS WLS).
    For Spearman/Kendall: weights are approximated via replication (SPSS approach).
    Returns (corr_matrix_with_stars, n_effective).
    """
    df_clean = df[cols].dropna()

    if weights is not None:
        w = pd.Series(weights, index=df.index).reindex(df_clean.index).fillna(0)
        w = w.clip(lower=0)
    else:
        w = pd.Series(np.ones(len(df_clean)), index=df_clean.index)

    n = int(w.sum())
    corr_matrix = pd.DataFrame(index=cols, columns=cols)

    if method == 'pearson':
        # Weighted Pearson: r = cov_w(x,y) / (sd_w(x) * sd_w(y))
        w_sum = w.sum()
        w_arr = w.values

        def _wcov(x, y, w_arr, w_sum):
            mx = (x * w_arr).sum() / w_sum
            my = (y * w_arr).sum() / w_sum
            return (w_arr * (x - mx) * (y - my)).sum() / w_sum

        for c1 in cols:
            for c2 in cols:
                if c1 == c2:
                    corr_matrix.loc[c1, c2] = "1.000"
                else:
                    try:
                        x = df_clean[c1].values.astype(float)
                        y = df_clean[c2].values.astype(float)
                        cov_xy = _wcov(x, y, w_arr, w_sum)
                        cov_xx = _wcov(x, x, w_arr, w_sum)
                        cov_yy = _wcov(y, y, w_arr, w_sum)
                        denom = np.sqrt(cov_xx * cov_yy)
                        r = cov_xy / denom if denom > 0 else 0.0
                        r = max(-1.0, min(1.0, r))
                        # t-stat for significance
                        df_t = max(n - 2, 1)
                        t = r * np.sqrt(df_t / max(1 - r**2, 1e-12))
                        p = 2 * stats.t.sf(abs(t), df_t)
                        stars = "**" if p < 0.01 else ("*" if p < 0.05 else "")
                        corr_matrix.loc[c1, c2] = f"{r:.3f}{stars}\n(p={p:.3f})"
                    except Exception:
                        corr_matrix.loc[c1, c2] = "N/A"
    else:
        # Spearman / Kendall \u2014 use unweighted (methodologically standard;
        # SPSS does not weight rank-based correlations either)
        df_c2 = df_clean.copy()
        for c1 in cols:
            for c2 in cols:
                if c1 == c2:
                    corr_matrix.loc[c1, c2] = "1.000"
                else:
                    try:
                        if method == 'spearman':
                            r, p = stats.spearmanr(df_c2[c1], df_c2[c2])
                        else:
                            r, p = stats.kendalltau(df_c2[c1], df_c2[c2])
                        stars = "**" if p < 0.01 else ("*" if p < 0.05 else "")
                        corr_matrix.loc[c1, c2] = f"{r:.3f}{stars}\n(p={p:.3f})"
                    except Exception:
                        corr_matrix.loc[c1, c2] = "N/A"

    return corr_matrix, n

# -------------------------------------------------------------
# REGRESJA
# -------------------------------------------------------------

def run_regression_block(df_data, dep_var, indep_vars_blocks, weights=None):
    results_list = []
    prev_r2 = 0.0
    cumulative_vars = []
    for block_idx, block_vars in enumerate(indep_vars_blocks):
        cumulative_vars = cumulative_vars + block_vars
        df_reg = df_data[[dep_var] + cumulative_vars].dropna()
        if len(df_reg) < len(cumulative_vars) + 2:
            results_list.append({'error': f"Blok {block_idx + 1}: Za ma\u0142o obserwacji ({len(df_reg)})."})
            continue
        y = df_reg[dep_var]
        X = sm.add_constant(df_reg[cumulative_vars])
        try:
            if weights is not None:
                w_reg = pd.Series(weights, index=df_data.index).reindex(df_reg.index).fillna(1).clip(lower=0)
                model = sm.WLS(y, X, weights=w_reg).fit()
            else:
                model = sm.OLS(y, X).fit()
        except Exception as e:
            results_list.append({'error': str(e)})
            continue
        r2 = model.rsquared
        r2_adj = model.rsquared_adj
        delta_r2 = r2 - prev_r2
        f_stat = model.fvalue
        f_pval = model.f_pvalue
        if block_idx == 0:
            f_change, f_change_p, df1_change = f_stat, f_pval, len(cumulative_vars)
        else:
            df1_change = len(block_vars)
            df2_change = len(df_reg) - len(cumulative_vars) - 1
            if df2_change > 0 and (1 - r2) > 0:
                f_change = (delta_r2 / df1_change) / ((1 - r2) / df2_change)
                f_change_p = 1 - stats.f.cdf(f_change, df1_change, df2_change)
            else:
                f_change, f_change_p = np.nan, np.nan
        vif_dict = {}
        if len(cumulative_vars) > 1:
            X_vif = df_reg[cumulative_vars].astype(float)
            for i, v in enumerate(cumulative_vars):
                try:
                    vif_dict[v] = variance_inflation_factor(X_vif.values, i)
                except:
                    vif_dict[v] = np.nan
        else:
            vif_dict[cumulative_vars[0]] = np.nan
        std_y = y.std()
        beta_dict = {}
        for v in cumulative_vars:
            std_x = df_reg[v].std()
            beta_dict[v] = model.params[v] * std_x / std_y if std_y > 0 and std_x > 0 else np.nan
        coef_rows = []
        for v in cumulative_vars:
            vif_val = vif_dict.get(v, np.nan)
            coef_rows.append({
                'Zmienna': v,
                'B': model.params.get(v, np.nan),
                'B\u0142\u0105d std. B': model.bse.get(v, np.nan),
                'Beta (std.)': beta_dict.get(v, np.nan),
                't': model.tvalues.get(v, np.nan),
                'p-value': model.pvalues.get(v, np.nan),
                'VIF': vif_val,
                'Tolerancja': 1 / vif_val if pd.notna(vif_val) and vif_val > 0 else np.nan,
            })
        results_list.append({
            'Blok': block_idx + 1,
            'Zmienne w bloku': ', '.join(block_vars),
            'Wszystkie predyktory': cumulative_vars[:],
            'dep_var': dep_var,
            'N': len(df_reg),
            'R': np.sqrt(r2),
            'R2': r2,
            'Skor_R2': r2_adj,
            'Delta_R2': delta_r2,
            'F modelu': f_stat,
            'p (F modelu)': f_pval,
            'F zmiany': f_change,
            'p (F zmiany)': f_change_p,
            'df1 (F zmiany)': df1_change,
            'df2 (F zmiany)': len(df_reg) - len(cumulative_vars) - 1,
            'coef_df': pd.DataFrame(coef_rows),
            '_model': model,
            '_df_reg': df_reg,
        })
        prev_r2 = r2
    return results_list

# -------------------------------------------------------------
# ANOVA
# -------------------------------------------------------------

def run_anova(df_raw, dep_var, group_var, df_labeled, weights=None):
    """One-way ANOVA with post-hoc Tukey HSD. Supports case weights (SPSS-compatible)."""
    tmp = pd.DataFrame({
        'dep': df_raw[dep_var].values,
        'grp': df_labeled[group_var].values
    }, index=df_raw.index)
    if weights is not None:
        tmp['w'] = pd.Series(weights, index=df_raw.index).reindex(tmp.index).fillna(0).clip(lower=0)
    else:
        tmp['w'] = 1.0
    tmp = tmp.dropna(subset=['dep', 'grp'])
    tmp = tmp[tmp['w'] > 0]

    groups = tmp['grp'].unique()
    group_data = [tmp.loc[tmp['grp'] == g, 'dep'].values for g in groups]
    group_w    = [tmp.loc[tmp['grp'] == g, 'w'].values  for g in groups]
    group_data = [(d, w) for d, w in zip(group_data, group_w) if len(d) >= 2]
    if len(group_data) < 2:
        return None, "Za ma\u0142o grup z wystarczaj\u0105c\u0105 liczb\u0105 obserwacji."

    # Weighted grand mean
    total_w  = tmp['w'].sum()
    grand_mean = (tmp['dep'] * tmp['w']).sum() / total_w

    # Weighted group stats
    desc_rows = []
    for g in groups:
        sub = tmp[tmp['grp'] == g]
        n_g   = sub['w'].sum()
        mean_g = (sub['dep'] * sub['w']).sum() / n_g if n_g > 0 else np.nan
        var_g  = (sub['w'] * (sub['dep'] - mean_g) ** 2).sum() / max(n_g - 1, 1)
        std_g  = np.sqrt(var_g)
        desc_rows.append({
            'Grupa': g,
            'N (wa\u017cone)': round(n_g, 2),
            'Srednia': round(mean_g, 4),
            'Odch. std.': round(std_g, 4),
            'Min': sub['dep'].min(),
            'Max': sub['dep'].max()
        })
    desc_df = pd.DataFrame(desc_rows)

    # Weighted SS
    ss_between = sum(
        (tmp[tmp['grp'] == g]['w'].sum()) *
        ((tmp[tmp['grp'] == g]['dep'] * tmp[tmp['grp'] == g]['w']).sum() /
         tmp[tmp['grp'] == g]['w'].sum() - grand_mean) ** 2
        for g in groups
    )
    ss_total   = (tmp['w'] * (tmp['dep'] - grand_mean) ** 2).sum()
    ss_within  = ss_total - ss_between

    df_between = len(groups) - 1
    df_within  = total_w - len(groups)   # effective df (sum of weights - k)
    ms_between = ss_between / df_between if df_between > 0 else np.nan
    ms_within  = ss_within  / df_within  if df_within  > 0 else np.nan
    f_stat     = ms_between / ms_within  if ms_within and ms_within > 0 else np.nan
    p_val      = stats.f.sf(f_stat, df_between, df_within) if not np.isnan(f_stat) else np.nan

    # Levene's test (unweighted \u2014 Levene is robust; SPSS uses unweighted too)
    lev_stat, lev_p = stats.levene(*[d for d, _ in group_data])

    eta2 = ss_between / ss_total if ss_total > 0 else np.nan

    # Tukey HSD post-hoc (weighted)
    from itertools import combinations
    posthoc_rows = []
    for g1, g2 in combinations(groups, 2):
        s1 = tmp[tmp['grp'] == g1]; s2 = tmp[tmp['grp'] == g2]
        n1  = s1['w'].sum(); n2  = s2['w'].sum()
        m1  = (s1['dep'] * s1['w']).sum() / n1 if n1 > 0 else np.nan
        m2  = (s2['dep'] * s2['w']).sum() / n2 if n2 > 0 else np.nan
        if n1 < 2 or n2 < 2 or pd.isna(ms_within): continue
        diff    = m1 - m2
        se_tukey = np.sqrt(ms_within * (1 / n1 + 1 / n2) / 2)
        q = abs(diff) / se_tukey if se_tukey > 0 else np.nan
        try:
            from scipy.stats import studentized_range
            p_tukey = studentized_range.sf(q * np.sqrt(2), len(groups), df_within)
        except Exception:
            p_tukey = np.nan
        posthoc_rows.append({
            'Grupa A': g1, 'Grupa B': g2,
            'R\u00f3\u017cnica \u015brednich (A-B)': round(diff, 4),
            'p-value (Tukey)': round(p_tukey, 4) if not np.isnan(p_tukey) else np.nan,
            'Istotna (p<0.05)': '\u2705' if (not np.isnan(p_tukey) and p_tukey < 0.05) else '\u274c'
        })

    posthoc_df = pd.DataFrame(posthoc_rows) if posthoc_rows else pd.DataFrame()

    result = {
        'dep_var': dep_var, 'group_var': group_var,
        'F': f_stat, 'p': p_val,
        'df_between': df_between, 'df_within': df_within,
        'eta2': eta2, 'lev_stat': lev_stat, 'lev_p': lev_p,
        'desc_df': desc_df, 'posthoc_df': posthoc_df,
        'ss_between': ss_between, 'ss_within': ss_within, 'ss_total': ss_total,
        'ms_between': ms_between, 'ms_within': ms_within,
    }
    return result, None

# -------------------------------------------------------------
# ANALIZA CZYNNIKOWA
# -------------------------------------------------------------

def run_factor_analysis(df_raw, cols, n_factors, rotation='varimax', method='principal', weights=None):
    df_fa = df_raw[cols].dropna()
    if len(df_fa) < len(cols) + 5:
        return None, f"Za ma\u0142o obserwacji ({len(df_fa)}). Potrzeba co najmniej {len(cols)+5}."
    if n_factors >= len(cols):
        return None, f"Liczba czynnik\u00f3w ({n_factors}) musi by\u0107 mniejsza ni\u017c liczba zmiennych ({len(cols)})."
    try:
        # When weights are provided, build weighted correlation matrix and pass to FA.
        # This is consistent with SPSS FACTOR (WLS/GLS approach via cov_matrix).
        if weights is not None:
            w = pd.Series(weights, index=df_raw.index).reindex(df_fa.index).fillna(0).clip(lower=0)
            w_arr = w.values
            w_sum = w_arr.sum()
            X = df_fa.values.astype(float)
            # Weighted means
            means = (w_arr[:, None] * X).sum(axis=0) / w_sum
            Xc = X - means
            # Weighted covariance matrix
            cov_w = (w_arr[:, None] * Xc).T @ Xc / (w_sum - 1)
            # Convert to correlation matrix
            std_w = np.sqrt(np.diag(cov_w))
            corr_w = cov_w / np.outer(std_w, std_w)
            np.fill_diagonal(corr_w, 1.0)
            corr_df_w = pd.DataFrame(corr_w, index=cols, columns=cols)
            # KMO and Bartlett on weighted corr
            try:
                kmo_all, kmo_model = calculate_kmo(corr_df_w)
                bart_chi2, bart_p = calculate_bartlett_sphericity(corr_df_w)
            except Exception:
                kmo_all, kmo_model = calculate_kmo(df_fa)
                bart_chi2, bart_p = calculate_bartlett_sphericity(df_fa)
            fa = FactorAnalyzer(n_factors=n_factors, rotation=rotation, method=method,
                                is_corr_matrix=True)
            fa.fit(corr_df_w)
        else:
            kmo_all, kmo_model = calculate_kmo(df_fa)
            bart_chi2, bart_p = calculate_bartlett_sphericity(df_fa)
            fa = FactorAnalyzer(n_factors=n_factors, rotation=rotation, method=method)
            fa.fit(df_fa)

        loadings = pd.DataFrame(fa.loadings_, index=cols,
                                columns=[f"Czynnik {i + 1}" for i in range(n_factors)])
        communalities = pd.DataFrame({'Komunalnosc (h2)': fa.get_communalities()}, index=cols)
        ev, v = fa.get_eigenvalues()
        eigenvalues_df = pd.DataFrame({'Warto\u015b\u0107 w\u0142asna': ev, '% wariancji': ev / len(cols) * 100,
                                       'Skumulowany %': np.cumsum(ev / len(cols) * 100)},
                                      index=[f"Czynnik {i + 1}" for i in range(len(ev))])
        var_explained = fa.get_factor_variance()
        var_df = pd.DataFrame({'SS \u0141adunk\u00f3w': var_explained[0], '% wariancji': var_explained[1] * 100,
                               'Skumulowany %': var_explained[2] * 100},
                              index=[f"Czynnik {i + 1}" for i in range(n_factors)])
        return {
            'loadings': loadings, 'communalities': communalities,
            'eigenvalues': eigenvalues_df, 'variance': var_df,
            'kmo': kmo_model, 'kmo_all': kmo_all,
            'bartlett_chi2': bart_chi2, 'bartlett_p': bart_p,
            'n': len(df_fa), 'cols': cols, 'rotation': rotation,
        }, None
    except Exception as e:
        return None, str(e)

# -------------------------------------------------------------
# EKSPORT DO EXCELA -- TABLICE WYNIKOWE (NAPRAWIONY)
# -------------------------------------------------------------

def safe_excel_val(val):
    """Convert value to Excel-safe type."""
    if val is None:
        return ""
    if isinstance(val, float) and np.isnan(val):
        return ""
    if isinstance(val, (np.integer,)):
        return int(val)
    if isinstance(val, (np.floating,)):
        return float(val)
    if isinstance(val, pd.CategoricalDtype):
        return str(val)
    return val

def export_toc_sheet(writer, results, matrix_results, var_labels, sheet_map,
                     regression_results=None, anova_results=None, factor_results=None,
                     conjoint_results=None, maxdiff_results=None,
                     pre_created_ws=None):
    """
    Create a 'Spis Tre\u015bci' sheet with clickable hyperlinks to each table.
    If pre_created_ws is provided (worksheet created before data sheets), use it directly
    so the ToC appears as the first tab in Excel.
    sheet_map: dict {sheet_name: {title: excel_row}} -- row index where each table starts.
    """
    workbook = writer.book
    if pre_created_ws is not None:
        worksheet = pre_created_ws
    else:
        worksheet = workbook.add_worksheet('\U0001f4cb Spis Tre\u015bci')
        worksheet.set_tab_color('#1F4E79')
    worksheet.activate()

    fmt_title   = workbook.add_format({'bold': True, 'font_size': 14, 'font_color': '#1F4E79',
                                        'bottom': 2, 'bottom_color': '#1F4E79'})
    fmt_section = workbook.add_format({'bold': True, 'bg_color': '#1F4E79', 'font_color': 'white',
                                        'border': 1, 'align': 'left', 'font_size': 11})
    fmt_link    = workbook.add_format({'font_color': '#0563C1', 'underline': True, 'border': 1, 'align': 'left'})
    fmt_sub     = workbook.add_format({'italic': True, 'font_color': '#595959', 'border': 1, 'align': 'left'})
    fmt_hdr     = workbook.add_format({'bold': True, 'bg_color': '#D6E4F0', 'border': 1})
    fmt_num     = workbook.add_format({'border': 1, 'align': 'right', 'num_format': '#,##0'})
    fmt_empty   = workbook.add_format({'border': 1})

    worksheet.set_column(0, 0,  6)   # #
    worksheet.set_column(1, 1, 55)   # Nazwa tabeli
    worksheet.set_column(2, 2, 22)   # Arkusz
    worksheet.set_column(3, 3, 12)   # Typ

    row = 0
    worksheet.merge_range(row, 0, row, 3, '\U0001f4cb Spis Tre\u015bci Raportu Analitycznego', fmt_title)
    row += 2

    counter = 1

    # -- Tablice cz?sto?ci --
    czestosci = results.get('czestosci', {})
    if czestosci:
        worksheet.merge_range(row, 0, row, 3, '\U0001f4c8 Tablice Cz\u0119sto\u015bci', fmt_section)
        row += 1
        worksheet.write(row, 0, '#',          fmt_hdr)
        worksheet.write(row, 1, 'Zmienna',    fmt_hdr)
        worksheet.write(row, 2, 'Arkusz',     fmt_hdr)
        worksheet.write(row, 3, 'Typ',        fmt_hdr)
        row += 1
        for title in czestosci:
            display = f"[{title}] {var_labels.get(title, title)}"
            target_row = sheet_map.get('Cz\u0119sto\u015bci', {}).get(title, 1)
            cell_addr = f"A{target_row + 1}"
            try:
                worksheet.write_url(row, 1, f"internal:'Cz\u0119sto\u015bci'!{cell_addr}", fmt_link, display)
            except Exception:
                worksheet.write(row, 1, display, fmt_link)
            worksheet.write(row, 0, counter,     fmt_num)
            worksheet.write(row, 2, 'Cz\u0119sto\u015bci', fmt_empty)
            worksheet.write(row, 3, 'Cz\u0119sto\u015bci', fmt_empty)
            counter += 1
            row += 1
        row += 1

    # -- Tablice krzy?owe --
    krzyzowe = results.get('krzyzowe', {})
    if krzyzowe:
        worksheet.merge_range(row, 0, row, 3, '\U0001f500 Tablice Krzy\u017cowe', fmt_section)
        row += 1
        worksheet.write(row, 0, '#',         fmt_hdr)
        worksheet.write(row, 1, 'Tabela',    fmt_hdr)
        worksheet.write(row, 2, 'Arkusz',    fmt_hdr)
        worksheet.write(row, 3, 'Typ',       fmt_hdr)
        row += 1
        for title in krzyzowe:
            if ' x ' in title:
                rv, cv = title.split(' x ', 1)
                display = (f"Wiersz: [{rv}] {var_labels.get(rv, rv)}  \u00d7  "
                           f"Kolumna: [{cv}] {var_labels.get(cv, cv)}")
            else:
                display = title
            target_row = sheet_map.get('Krzy\u017cowe', {}).get(title, 1)
            cell_addr = f"A{target_row + 1}"
            try:
                worksheet.write_url(row, 1, f"internal:'Krzy\u017cowe'!{cell_addr}", fmt_link, display)
            except Exception:
                worksheet.write(row, 1, display, fmt_link)
            worksheet.write(row, 0, counter,    fmt_num)
            worksheet.write(row, 2, 'Krzy\u017cowe', fmt_empty)
            worksheet.write(row, 3, 'Krzy\u017cowe', fmt_empty)
            counter += 1
            row += 1
        row += 1

    # -- Pytania matrycowe --
    if matrix_results:
        worksheet.merge_range(row, 0, row, 3, '\U0001f522 Pytania Matrycowe', fmt_section)
        row += 1
        worksheet.write(row, 0, '#',       fmt_hdr)
        worksheet.write(row, 1, 'Pytanie', fmt_hdr)
        worksheet.write(row, 2, 'Arkusz',  fmt_hdr)
        worksheet.write(row, 3, 'Typ',     fmt_hdr)
        row += 1
        for entry in matrix_results:
            try:
                worksheet.write_url(row, 1, "internal:'Pytania Matrycowe'!A1", fmt_link, entry['name'])
            except Exception:
                worksheet.write(row, 1, entry['name'], fmt_link)
            worksheet.write(row, 0, counter,              fmt_num)
            worksheet.write(row, 2, 'Pytania Matrycowe',  fmt_empty)
            worksheet.write(row, 3, 'Matryca',            fmt_empty)
            counter += 1
            row += 1
        row += 1

    # -- Other sheets --
    other_sheets = [('\u015arednie', 'srednie', '\U0001f4ca'), ('Opisowe', 'opisowe', '\U0001f522'),
                    ('Korelacje', 'korelacje', '\U0001f517')]
    for sheet_name, key, icon in other_sheets:
        if results.get(key):
            worksheet.write(row, 0, '',  fmt_empty)
            try:
                worksheet.write_url(row, 1, f"internal:'{sheet_name}'!A1", fmt_link, f"{icon} {sheet_name}")
            except Exception:
                worksheet.write(row, 1, f"{icon} {sheet_name}", fmt_link)
            worksheet.write(row, 2, sheet_name, fmt_empty)
            worksheet.write(row, 3, '',          fmt_empty)
            row += 1

    # -- Optional analytical sheets (only shown if results exist) --
    valid_reg = [r for r in (regression_results or []) if 'error' not in r]
    has_anova  = bool(anova_results)
    has_fa     = bool(factor_results)
    valid_conj = [r for r in (conjoint_results or []) if not r.get('error')]
    has_md     = bool(maxdiff_results)

    for sheet_name, has_data, icon in [
        ('Regresja',        bool(valid_reg), '\U0001f4c9'),
        ('ANOVA',           has_anova,       '\U0001f4ca'),
        ('Anal. Czynnikowa', has_fa,         '\U0001f52c'),
        ('Conjoint',        bool(valid_conj),'\U0001f4ca'),
        ('MaxDiff',         has_md,          '\U0001f522'),
    ]:
        if not has_data:
            continue
        try:
            worksheet.write_url(row, 1, f"internal:'{sheet_name}'!A1", fmt_link, f"{icon} {sheet_name}")
            worksheet.write(row, 0, '', fmt_empty)
            worksheet.write(row, 2, sheet_name, fmt_empty)
            worksheet.write(row, 3, '', fmt_empty)
            row += 1
        except Exception:
            pass


def export_tables_to_sheet(writer, s_name, results_dict, var_labels, add_charts=False):
    workbook = writer.book

    # Limit sheet name to 31 chars (Excel limit)
    sheet_name = s_name[:31]
    worksheet = workbook.add_worksheet(sheet_name)

    fmt_title     = workbook.add_format({'bold': True, 'align': 'center', 'valign': 'vcenter', 'border': 1, 'bg_color': '#1F4E79', 'font_color': 'white', 'font_size': 10})
    fmt_header    = workbook.add_format({'bold': True, 'align': 'center', 'valign': 'vcenter', 'border': 1, 'bg_color': '#D6E4F0', 'text_wrap': True})
    fmt_index_b   = workbook.add_format({'bold': True, 'border': 1, 'align': 'left', 'bg_color': '#F2F2F2'})
    fmt_index_n   = workbook.add_format({'border': 1, 'align': 'left'})
    fmt_n         = workbook.add_format({'num_format': '#,##0', 'border': 1, 'align': 'right'})
    fmt_pct       = workbook.add_format({'num_format': '0.0"%"', 'border': 1, 'align': 'right'})
    fmt_float     = workbook.add_format({'num_format': '#,##0.00', 'border': 1, 'align': 'right'})
    fmt_str       = workbook.add_format({'border': 1, 'align': 'center', 'text_wrap': True})
    fmt_empty     = workbook.add_format({})

    worksheet.set_column(0, 0, 42)

    sr = 0
    title_row_map = {}   # title -> excel row number (for ToC hyperlinks)

    # Summary rows that should be excluded from charts
    _chart_exclude = {
        'Suma', 'Braki danych', 'Braki danych (wykluczone z tabeli)',
        'Og\u00f3\u0142em (Wa\u017cne)', 'Baza (N) / Suma (%)',
    }

    for title, df_res in results_dict.items():
        title_row_map[title] = sr   # record where this table starts
        df_export = df_res.copy()

        # Convert percentage columns to float for proper formatting
        if s_name != 'Korelacje':
            for col in df_export.columns:
                col_str = str(col).lower()
                if "%" in col_str or "procent" in col_str:
                    df_export[col] = df_export[col].apply(_to_float_pct)

        num_cols = len(df_export.columns)

        # Build display title
        if s_name == 'Cz\u0119sto\u015bci':
            display_title = f"[{title}] {var_labels.get(title, title)}"
            chart_title   = var_labels.get(title, title)   # label only (no [code])
        elif s_name in ['Krzy\u017cowe', '\u015arednie']:
            if ' x ' in title:
                r_v, c_v = title.split(' x ', 1)
                display_title = f"Wiersz: [{r_v}] {var_labels.get(r_v, r_v)}  \u00d7  Kolumna: [{c_v}] {var_labels.get(c_v, c_v)}"
            else:
                display_title = title
            chart_title = display_title
        else:
            display_title = title
            chart_title   = title

        # Title row
        if num_cols > 1:
            worksheet.merge_range(sr, 0, sr, num_cols, display_title, fmt_title)
        else:
            worksheet.write(sr, 0, display_title, fmt_title)
        sr += 1

        # Header row
        worksheet.write(sr, 0, "Kategorie / Statystyki", fmt_header)
        for c_idx, col_name in enumerate(df_export.columns):
            col_w = 16 if s_name not in ['Korelacje'] else 22
            worksheet.set_column(c_idx + 1, c_idx + 1, col_w)
            worksheet.write(sr, c_idx + 1, str(col_name), fmt_header)
        sr += 1

        # Data rows
        bold_rows = {'Suma', 'Braki danych', 'Braki danych (wykluczone z tabeli)',
                     'Og\u00f3\u0142em (Wa\u017cne)', 'Srednia', 'Odchylenie Std.', 'Baza (N)'}
        for r_idx, row_name in enumerate(df_export.index):
            is_bold = str(row_name) in bold_rows
            idx_fmt = fmt_index_b if is_bold else fmt_index_n
            row_name_str = "" if (pd.isna(row_name) if not isinstance(row_name, str) else False) else str(row_name)
            worksheet.write(sr + r_idx, 0, row_name_str, idx_fmt)

            for c_idx, col_name in enumerate(df_export.columns):
                raw_val = df_export.iloc[r_idx, c_idx]
                col_str = str(col_name).lower()
                is_pct_col = "%" in col_str or "procent" in col_str

                try:
                    is_empty = pd.isna(raw_val)
                except:
                    is_empty = False
                if is_empty or str(raw_val).strip() in ('', 'nan', 'None'):
                    worksheet.write(sr + r_idx, c_idx + 1, "", fmt_empty)
                    continue

                if isinstance(raw_val, str):
                    worksheet.write(sr + r_idx, c_idx + 1, raw_val, fmt_str)
                elif s_name == 'Korelacje':
                    worksheet.write(sr + r_idx, c_idx + 1, str(raw_val), fmt_str)
                elif s_name == 'Opisowe':
                    worksheet.write(sr + r_idx, c_idx + 1, float(raw_val), fmt_float)
                elif s_name == '\u015arednie' and str(row_name) == 'Baza (N)':
                    worksheet.write(sr + r_idx, c_idx + 1, float(raw_val), fmt_n)
                elif s_name == '\u015arednie':
                    worksheet.write(sr + r_idx, c_idx + 1, float(raw_val), fmt_float)
                elif is_pct_col:
                    try:
                        worksheet.write(sr + r_idx, c_idx + 1, float(raw_val), fmt_pct)
                    except:
                        worksheet.write(sr + r_idx, c_idx + 1, str(raw_val), fmt_str)
                else:
                    try:
                        worksheet.write(sr + r_idx, c_idx + 1, float(raw_val), fmt_n)
                    except:
                        worksheet.write(sr + r_idx, c_idx + 1, str(raw_val), fmt_str)

        # \u2500\u2500 Native Excel chart (frequency tables only) \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
        if add_charts and s_name == 'Cz\u0119sto\u015bci':
            # Identify the % column index (1-based, col 0 = category labels)
            pct_col_idx = None
            for ci, col_name in enumerate(df_export.columns):
                cn_low = str(col_name).lower()
                if 'procent' in cn_low or '%' in cn_low:
                    pct_col_idx = ci + 1   # +1 because col 0 is the index label
                    break

            if pct_col_idx is not None:
                # Collect rows that should appear in the chart (exclude summary rows)
                chart_data_rows = []   # 0-based excel row numbers
                for r_idx, row_name in enumerate(df_export.index):
                    rn = str(row_name)
                    if rn in _chart_exclude:
                        continue
                    # Skip Box rows like [Top 2 Box]
                    if rn.startswith('[') and rn.endswith(']'):
                        continue
                    chart_data_rows.append(sr + r_idx)   # sr already advanced past headers

                if len(chart_data_rows) >= 2:
                    first_row = chart_data_rows[0]
                    last_row  = chart_data_rows[-1]

                    chart = workbook.add_chart({'type': 'bar'})   # horizontal bars
                    chart.add_series({
                        'name':       chart_title[:60],
                        'categories': [sheet_name, first_row, 0,          last_row, 0],
                        'values':     [sheet_name, first_row, pct_col_idx, last_row, pct_col_idx],
                        'fill':       {'color': '#2E75B6'},
                        'border':     {'color': '#1F4E79'},
                        'gap':        60,
                        'data_labels': {
                            'value':      True,
                            'num_format': '0.0"%"',
                            'font':       {'size': 9},
                            'position':   'outside_end',
                        },
                    })
                    chart.set_title({
                        'name':    chart_title[:80],
                        'overlay': False,
                    })
                    # X axis: no title, no tick labels, no gridlines
                    chart.set_x_axis({
                        'name':         '',
                        'min':           0,
                        'num_font':     {'size': 1, 'color': '#FFFFFF'},  # invisible labels
                        'major_gridlines': {'visible': False},
                        'minor_gridlines': {'visible': False},
                        'major_tick_mark': 'none',
                        'minor_tick_mark': 'none',
                        'line':          {'none': True},
                    })
                    # Y axis: categories top-to-bottom, no gridlines
                    chart.set_y_axis({
                        'reverse':         True,
                        'num_font':        {'size': 9},
                        'major_gridlines': {'visible': False},
                        'minor_gridlines': {'visible': False},
                        'major_tick_mark': 'none',
                        'minor_tick_mark': 'none',
                        'line':            {'none': True},
                    })
                    chart.set_legend({'none': True})
                    chart.set_plotarea({'border': {'none': True}})
                    chart.set_chartarea({'border': {'color': '#D6E4F0'}})

                    # Height: match the table exactly.
                    # Excel default row height = 15pt = 20px.
                    # Table occupies: 1 title row + 1 header row + len(df_export) data rows + 1 blank = len+3 rows
                    # We use 20px per row as a close approximation.
                    table_rows   = len(df_export) + 2   # title + header + data rows
                    c_height     = max(180, table_rows * 20)
                    chart.set_size({'width': 480, 'height': c_height})

                    # Insert aligned with table title row, to the right
                    insert_col = num_cols + 2
                    title_row  = title_row_map[title]
                    worksheet.insert_chart(title_row, insert_col, chart,
                                           {'x_offset': 5, 'y_offset': 0})

        # \u2500\u2500 Advance row pointer \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
        sr += len(df_export) + 1

        # Chi-square note
        if s_name == 'Krzy\u017cowe' and title in st.session_state.chi_results:
            worksheet.write(sr, 0, st.session_state.chi_results[title],
                            workbook.add_format({'italic': True, 'font_color': '#595959'}))
            sr += 1
        sr += 2  # blank rows between tables

    return title_row_map   # {title: starting excel row} for ToC hyperlinks


def export_regression_to_excel(writer, regression_results, var_labels):
    workbook = writer.book
    worksheet = workbook.add_worksheet('Regresja')

    fmt_title   = workbook.add_format({'bold': True, 'bg_color': '#1F4E79', 'font_color': 'white', 'align': 'center', 'valign': 'vcenter', 'border': 1, 'font_size': 11})
    fmt_section = workbook.add_format({'bold': True, 'bg_color': '#D6E4F0', 'border': 1, 'align': 'left'})
    fmt_header  = workbook.add_format({'bold': True, 'bg_color': '#F2F2F2', 'border': 1, 'align': 'center', 'text_wrap': True})
    fmt_label   = workbook.add_format({'bold': True, 'border': 1})
    fmt_val     = workbook.add_format({'border': 1, 'align': 'right', 'num_format': '#,##0.000'})
    fmt_int     = workbook.add_format({'border': 1, 'align': 'right', 'num_format': '#,##0'})
    fmt_warn    = workbook.add_format({'border': 1, 'align': 'right', 'num_format': '#,##0.000', 'font_color': '#C00000'})
    fmt_ok      = workbook.add_format({'border': 1, 'align': 'right', 'num_format': '#,##0.000', 'font_color': '#006100'})
    fmt_dash    = workbook.add_format({'border': 1, 'align': 'center'})

    worksheet.set_column(0, 0, 38)
    for i in range(1, 10): worksheet.set_column(i, i, 16)

    row = 0
    for res in regression_results:
        if 'error' in res:
            worksheet.write(row, 0, f"B\u0141\u0104D: {res['error']}", fmt_section)
            row += 2
            continue
        dep_label = var_labels.get(res['dep_var'], res['dep_var'])
        worksheet.merge_range(row, 0, row, 8, f"REGRESJA OLS -- Zmienna zale\u017cna: [{res['dep_var']}] {dep_label}  |  Blok {res['Blok']}", fmt_title)
        row += 1
        worksheet.write(row, 0, "Podsumowanie Modelu", fmt_section)
        row += 1
        for field, val, fmt in [
            ('N (obserwacje)', res['N'], fmt_int),
            ('R', res['R'], fmt_val),
            ('R2', res['R2'], fmt_val),
            ('Skor_R2', res['Skor_R2'], fmt_val),
            ('\u0394R\u00b2 (zmiana R\u00b2)', res['Delta_R2'], fmt_val),
            ('F modelu', res['F modelu'], fmt_val),
            ('p (F modelu)', res['p (F modelu)'], fmt_val),
            ('F zmiany', res['F zmiany'], fmt_val),
            ('p (F zmiany)', res['p (F zmiany)'], fmt_val),
        ]:
            worksheet.write(row, 0, field, fmt_label)
            try:
                v = float(val)
                worksheet.write(row, 1, v, fmt)
            except:
                worksheet.write(row, 1, '--', fmt_dash)
            row += 1
        row += 1
        worksheet.write(row, 0, "Wsp\u00f3\u0142czynniki Regresji", fmt_section)
        row += 1
        for ci, h in enumerate(['Zmienna', 'B', 'B\u0142\u0105d std. B', 'Beta (std.)', 't', 'p-value', 'VIF', 'Tolerancja']):
            worksheet.write(row, ci, h, fmt_header)
        row += 1
        for _, r_data in res['coef_df'].iterrows():
            vn = r_data['Zmienna']
            worksheet.write(row, 0, f"[{vn}] {var_labels.get(vn, vn)}", fmt_label)
            worksheet.write(row, 1, float(r_data['B']), fmt_val)
            worksheet.write(row, 2, float(r_data['B\u0142\u0105d std. B']), fmt_val)
            try: worksheet.write(row, 3, float(r_data['Beta (std.)']), fmt_val)
            except: worksheet.write(row, 3, '--', fmt_dash)
            worksheet.write(row, 4, float(r_data['t']), fmt_val)
            p = r_data['p-value']
            try:
                pf = float(p)
                worksheet.write(row, 5, pf, fmt_ok if pf < 0.05 else fmt_val)
            except: worksheet.write(row, 5, '--', fmt_dash)
            vif = r_data['VIF']
            try:
                vf = float(vif)
                vif_fmt = fmt_warn if vf > 10 else fmt_val
                worksheet.write(row, 6, vf, vif_fmt)
                worksheet.write(row, 7, float(r_data['Tolerancja']), vif_fmt)
            except:
                worksheet.write(row, 6, '--', fmt_dash)
                worksheet.write(row, 7, '--', fmt_dash)
            row += 1
        worksheet.write(row, 0, "VIF > 10 = problem ze wsp\u00f3\u0142liniowo\u015bci\u0105  |  p < 0.05 = istotne statystycznie",
                        workbook.add_format({'italic': True, 'font_color': '#595959'}))
        row += 3


def export_anova_to_excel(writer, anova_results, var_labels):
    workbook = writer.book
    worksheet = workbook.add_worksheet('ANOVA')
    fmt_title   = workbook.add_format({'bold': True, 'bg_color': '#1F4E79', 'font_color': 'white', 'border': 1, 'align': 'center', 'font_size': 11})
    fmt_section = workbook.add_format({'bold': True, 'bg_color': '#D6E4F0', 'border': 1})
    fmt_header  = workbook.add_format({'bold': True, 'bg_color': '#F2F2F2', 'border': 1, 'align': 'center'})
    fmt_label   = workbook.add_format({'bold': True, 'border': 1})
    fmt_val     = workbook.add_format({'border': 1, 'align': 'right', 'num_format': '#,##0.000'})
    fmt_int     = workbook.add_format({'border': 1, 'align': 'right', 'num_format': '#,##0'})
    fmt_str     = workbook.add_format({'border': 1, 'align': 'center'})
    worksheet.set_column(0, 0, 30)
    for i in range(1, 8): worksheet.set_column(i, i, 16)
    row = 0
    for res in anova_results:
        dep_l = var_labels.get(res['dep_var'], res['dep_var'])
        grp_l = var_labels.get(res['group_var'], res['group_var'])
        worksheet.merge_range(row, 0, row, 6, f"ANOVA -- Zmienna zale\u017cna: {dep_l}  |  Czynnik: {grp_l}", fmt_title)
        row += 2
        worksheet.write(row, 0, "Tabela ANOVA", fmt_section)
        row += 1
        for ci, h in enumerate(['\u0179r\u00f3d\u0142o', 'SS', 'df', 'MS', 'F', 'p-value', 'Eta\u00b2']):
            worksheet.write(row, ci, h, fmt_header)
        row += 1
        worksheet.write(row, 0, "Mi\u0119dzy grupami", fmt_label)
        worksheet.write(row, 1, res['ss_between'], fmt_val)
        worksheet.write(row, 2, res['df_between'], fmt_int)
        worksheet.write(row, 3, res['ms_between'], fmt_val)
        worksheet.write(row, 4, res['F'], fmt_val)
        worksheet.write(row, 5, res['p'], fmt_val)
        worksheet.write(row, 6, res['eta2'], fmt_val)
        row += 1
        worksheet.write(row, 0, "Wewn\u0105trz grup", fmt_label)
        worksheet.write(row, 1, res['ss_within'], fmt_val)
        worksheet.write(row, 2, res['df_within'], fmt_int)
        worksheet.write(row, 3, res['ms_within'], fmt_val)
        row += 2
        # Descriptives
        worksheet.write(row, 0, "Statystyki opisowe wg grupy", fmt_section)
        row += 1
        for ci, h in enumerate(res['desc_df'].columns):
            worksheet.write(row, ci, h, fmt_header)
        row += 1
        for _, r_d in res['desc_df'].iterrows():
            for ci, v in enumerate(r_d):
                try: worksheet.write(row, ci, float(v), fmt_val)
                except: worksheet.write(row, ci, str(v), fmt_str)
            row += 1
        row += 1
        # Post-hoc
        if not res['posthoc_df'].empty:
            worksheet.write(row, 0, "Test post-hoc: Tukey HSD", fmt_section)
            row += 1
            for ci, h in enumerate(res['posthoc_df'].columns):
                worksheet.write(row, ci, h, fmt_header)
            row += 1
            for _, r_d in res['posthoc_df'].iterrows():
                for ci, v in enumerate(r_d):
                    try: worksheet.write(row, ci, float(v), fmt_val)
                    except: worksheet.write(row, ci, str(v), fmt_str)
                row += 1
        row += 3


def export_factor_to_excel(writer, factor_results, var_labels):
    workbook = writer.book
    worksheet = workbook.add_worksheet('Anal. Czynnikowa')
    fmt_title   = workbook.add_format({'bold': True, 'bg_color': '#1F4E79', 'font_color': 'white', 'border': 1, 'align': 'center', 'font_size': 11})
    fmt_section = workbook.add_format({'bold': True, 'bg_color': '#D6E4F0', 'border': 1})
    fmt_header  = workbook.add_format({'bold': True, 'bg_color': '#F2F2F2', 'border': 1, 'align': 'center'})
    fmt_label   = workbook.add_format({'bold': True, 'border': 1})
    fmt_val     = workbook.add_format({'border': 1, 'align': 'right', 'num_format': '#,##0.000'})
    fmt_hi      = workbook.add_format({'border': 1, 'align': 'right', 'num_format': '#,##0.000', 'bold': True, 'bg_color': '#E2EFDA'})
    worksheet.set_column(0, 0, 38)
    for i in range(1, 12): worksheet.set_column(i, i, 14)
    row = 0
    for res in factor_results:
        worksheet.merge_range(row, 0, row, res['loadings'].shape[1], f"ANALIZA CZYNNIKOWA -- Rotacja: {res['rotation'].upper()}  |  N={res['n']}", fmt_title)
        row += 1
        # KMO and Bartlett
        worksheet.write(row, 0, "Adekwatno\u015b\u0107 pr\u00f3by KMO", fmt_label)
        worksheet.write(row, 1, res['kmo'], fmt_val)
        row += 1
        worksheet.write(row, 0, "Test sferyczno\u015bci Bartletta (Chi\u00b2)", fmt_label)
        worksheet.write(row, 1, res['bartlett_chi2'], fmt_val)
        row += 1
        worksheet.write(row, 0, "Test sferyczno\u015bci Bartletta (p)", fmt_label)
        worksheet.write(row, 1, res['bartlett_p'], fmt_val)
        row += 2
        # Loadings
        worksheet.write(row, 0, "Macierz \u0141adunk\u00f3w Czynnikowych", fmt_section)
        row += 1
        worksheet.write(row, 0, "Zmienna", fmt_header)
        for ci, col in enumerate(res['loadings'].columns):
            worksheet.write(row, ci + 1, col, fmt_header)
        worksheet.write(row, len(res['loadings'].columns) + 1, "Komunalno\u015b\u0107 (h\u00b2)", fmt_header)
        row += 1
        for var in res['loadings'].index:
            worksheet.write(row, 0, f"[{var}] {var_labels.get(var, var)}", fmt_label)
            for ci, col in enumerate(res['loadings'].columns):
                val = res['loadings'].loc[var, col]
                fmt_use = fmt_hi if abs(val) >= 0.4 else fmt_val
                worksheet.write(row, ci + 1, float(val), fmt_use)
            worksheet.write(row, len(res['loadings'].columns) + 1,
                            float(res['communalities'].loc[var, 'Komunalnosc (h2)']), fmt_val)
            row += 1
        row += 2
        # Variance explained
        worksheet.write(row, 0, "Wyja\u015bniona Wariancja", fmt_section)
        row += 1
        for ci, col in enumerate(['', 'SS \u0141adunk\u00f3w', '% wariancji', 'Skumulowany %']):
            worksheet.write(row, ci, col, fmt_header)
        row += 1
        for idx, r_d in res['variance'].iterrows():
            worksheet.write(row, 0, str(idx), fmt_label)
            for ci, v in enumerate(r_d):
                worksheet.write(row, ci + 1, float(v), fmt_val)
            row += 1
        row += 3


def export_matrix_to_excel(writer, matrix_results, var_labels):
    """
    Export matrix/battery frequency tables.
    Layout: Rows = scale values, Columns = subquestions (N | %).
    """
    workbook  = writer.book
    worksheet = workbook.add_worksheet('Pytania Matrycowe')

    fmt_title    = workbook.add_format({'bold': True, 'bg_color': '#1F4E79', 'font_color': 'white',
                                         'align': 'center', 'valign': 'vcenter', 'border': 1, 'font_size': 11})
    fmt_sub_hdr  = workbook.add_format({'bold': True, 'bg_color': '#2E75B6', 'font_color': 'white',
                                         'align': 'center', 'border': 1, 'text_wrap': True})
    fmt_np_hdr   = workbook.add_format({'bold': True, 'bg_color': '#D6E4F0',
                                         'align': 'center', 'border': 1})
    fmt_val_lbl  = workbook.add_format({'bold': True, 'border': 1, 'align': 'left', 'bg_color': '#F2F2F2'})
    fmt_base_lbl = workbook.add_format({'bold': True, 'border': 1, 'align': 'left', 'bg_color': '#E2EFDA'})
    fmt_suma_lbl = workbook.add_format({'bold': True, 'border': 1, 'align': 'left', 'bg_color': '#D6E4F0', 'italic': True})
    fmt_n        = workbook.add_format({'num_format': '#,##0',  'border': 1, 'align': 'right'})
    fmt_pct      = workbook.add_format({'num_format': '0.0"%"', 'border': 1, 'align': 'right'})
    fmt_base_n   = workbook.add_format({'num_format': '#,##0',  'border': 1, 'align': 'right',
                                         'bold': True, 'bg_color': '#E2EFDA'})
    fmt_suma_pct = workbook.add_format({'num_format': '0.0"%"', 'border': 1, 'align': 'right',
                                         'bold': True, 'bg_color': '#D6E4F0', 'italic': True})
    fmt_empty    = workbook.add_format({'border': 1})

    worksheet.set_column(0, 0, 14)   # row-index column (scale values)

    row = 0
    for entry in matrix_results:
        name         = entry['name']
        df_matrix    = entry['df']
        all_cats     = entry['cats']
        sub_labels   = entry['sub_labels']
        display_mode = entry.get('display_mode', 'N + %')   # default: show both

        n_subs = len(sub_labels)
        # Determine how many columns per subquestion based on display mode
        cols_per_sub = 1 if display_mode in ('Tylko N', 'Tylko %') else 2
        total_data_cols = n_subs * cols_per_sub
        total_cols = total_data_cols

        # -- Title --
        if total_cols > 0:
            worksheet.merge_range(row, 0, row, total_cols, f"Pytanie matrycowe: {name}", fmt_title)
        else:
            worksheet.write(row, 0, f"Pytanie matrycowe: {name}", fmt_title)
        row += 1

        # -- Row 1: subquestion labels --
        worksheet.write(row, 0, "Warto\u015b\u0107 \\ Subpytanie", fmt_sub_hdr)
        col_cur = 1
        for sub_lbl in sub_labels:
            disp = sub_lbl if len(sub_lbl) <= 40 else sub_lbl[:37] + "..."
            if display_mode == 'N + %':
                worksheet.merge_range(row, col_cur, row, col_cur + 1, disp, fmt_sub_hdr)
                worksheet.set_column(col_cur,     col_cur,     11)
                worksheet.set_column(col_cur + 1, col_cur + 1, 9)
                col_cur += 2
            else:
                worksheet.write(row, col_cur, disp, fmt_sub_hdr)
                worksheet.set_column(col_cur, col_cur, 12)
                col_cur += 1
        row += 1

        # -- Row 2: N / % sub-headers --
        worksheet.write(row, 0, "", fmt_np_hdr)
        col_cur = 1
        for _ in sub_labels:
            if display_mode == 'N + %':
                worksheet.write(row, col_cur,     "N",  fmt_np_hdr)
                worksheet.write(row, col_cur + 1, "%",  fmt_np_hdr)
                col_cur += 2
            elif display_mode == 'Tylko N':
                worksheet.write(row, col_cur, "N", fmt_np_hdr)
                col_cur += 1
            else:
                worksheet.write(row, col_cur, "%", fmt_np_hdr)
                col_cur += 1
        row += 1

        # -- Data rows --
        for cat_val in all_cats:
            worksheet.write(row, 0, str(cat_val), fmt_val_lbl)
            col_cur = 1
            for sub_lbl in sub_labels:
                if display_mode == 'N + %':
                    n_val   = df_matrix.loc[cat_val, f"{sub_lbl} [N]"]
                    pct_val = df_matrix.loc[cat_val, f"{sub_lbl} [%]"]
                    try: worksheet.write(row, col_cur,     float(n_val),   fmt_n)
                    except: worksheet.write(row, col_cur,     "", fmt_empty)
                    try: worksheet.write(row, col_cur + 1, float(pct_val), fmt_pct)
                    except: worksheet.write(row, col_cur + 1, "", fmt_empty)
                    col_cur += 2
                elif display_mode == 'Tylko N':
                    n_val = df_matrix.loc[cat_val, f"{sub_lbl} [N]"]
                    try: worksheet.write(row, col_cur, float(n_val), fmt_n)
                    except: worksheet.write(row, col_cur, "", fmt_empty)
                    col_cur += 1
                else:  # Tylko %
                    pct_val = df_matrix.loc[cat_val, f"{sub_lbl} [%]"]
                    try: worksheet.write(row, col_cur, float(pct_val), fmt_pct)
                    except: worksheet.write(row, col_cur, "", fmt_empty)
                    col_cur += 1
            row += 1

        # -- Single combined summary row: "Baza (N) / Suma (%)" --
        # N and % sit side by side, matching the frequency table style
        SUMROW = "Baza (N) / Suma (%)"
        worksheet.write(row, 0, SUMROW, fmt_suma_lbl)
        col_cur = 1
        for sub_lbl in sub_labels:
            base_val = df_matrix.loc[SUMROW, f"{sub_lbl} [N]"]
            suma_val = df_matrix.loc[SUMROW, f"{sub_lbl} [%]"]
            if display_mode == 'N + %':
                try: worksheet.write(row, col_cur,     float(base_val), fmt_base_n)
                except: worksheet.write(row, col_cur,     "", fmt_empty)
                try: worksheet.write(row, col_cur + 1, float(suma_val), fmt_suma_pct)
                except: worksheet.write(row, col_cur + 1, "", fmt_empty)
                col_cur += 2
            elif display_mode == 'Tylko N':
                try: worksheet.write(row, col_cur, float(base_val), fmt_base_n)
                except: worksheet.write(row, col_cur, "", fmt_empty)
                col_cur += 1
            else:  # Tylko %
                try: worksheet.write(row, col_cur, float(suma_val), fmt_suma_pct)
                except: worksheet.write(row, col_cur, "", fmt_empty)
                col_cur += 1
        row += 3   # gap between batteries


def write_db_sheet(writer, sheet_label, data_df, var_labels, hdr_color='#1F4E79'):
    """Write a single database sheet into an already-open ExcelWriter.
    Row 0 = column names (header), data from row 1 onwards.
    """
    workbook = writer.book
    ws = workbook.add_worksheet(sheet_label[:31])
    fmt_h = workbook.add_format({
        'bold': True, 'bg_color': hdr_color, 'font_color': 'white',
        'border': 1, 'align': 'center',
    })
    for ci, col in enumerate(data_df.columns):
        ws.write(0, ci, col, fmt_h)
        ws.set_column(ci, ci, 16)
    for ri, (_, row_data) in enumerate(data_df.iterrows()):
        for ci, val in enumerate(row_data):
            try:
                is_na = pd.isna(val)
            except Exception:
                is_na = False
            if is_na:
                ws.write(ri + 1, ci, '')
            elif isinstance(val, (int, float, np.integer, np.floating)):
                ws.write(ri + 1, ci, float(val))
            else:
                s = str(val)
                ws.write(ri + 1, ci, '' if s in ('nan', 'None', '<NA>') else s)


def export_db_to_excel(df_raw, df_labeled, var_labels):
    """Standalone download: both sheets in one file (kept for backward compat)."""
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
        write_db_sheet(writer, 'Baza z etykietami',        df_labeled, var_labels, '#1F4E79')
        write_db_sheet(writer, 'Baza surowa (numeryczna)', df_raw,     var_labels, '#2E75B6')
    return output.getvalue()

# -------------------------------------------------------------
# =============================================================
# CONJOINT ANALYSIS
# =============================================================

def run_conjoint_rating(df_raw, rating_var, attribute_vars):
    """Rating-based Conjoint via OLS. Returns part-worth utilities and importance."""
    from sklearn.preprocessing import LabelEncoder
    df_c = df_raw[[rating_var] + attribute_vars].dropna()
    if len(df_c) < 10:
        return None, "Za ma\u0142o obserwacji."
    y = df_c[rating_var].astype(float)
    # Dummy-encode categorical attributes; numeric attributes treated as linear
    X_parts = []
    attr_info = {}
    for attr in attribute_vars:
        col = df_c[attr]
        if col.dtype == object or col.nunique() <= 8:
            dummies = pd.get_dummies(col.astype(str), prefix=attr, drop_first=False)
            X_parts.append(dummies)
            attr_info[attr] = {'type': 'categorical', 'levels': list(dummies.columns)}
        else:
            X_parts.append(col.rename(attr).to_frame())
            attr_info[attr] = {'type': 'numeric', 'levels': [attr]}
    X = pd.concat(X_parts, axis=1).astype(float)
    X_const = sm.add_constant(X)
    model = sm.OLS(y, X_const).fit()
    # Part-worth utilities
    utilities = {}
    for attr, info in attr_info.items():
        lvl_utils = {}
        for lv in info['levels']:
            lvl_utils[lv] = model.params.get(lv, 0.0)
        # Zero-center (sum-to-zero coding adjustment)
        mean_u = np.mean(list(lvl_utils.values()))
        lvl_utils = {k: v - mean_u for k, v in lvl_utils.items()}
        utilities[attr] = lvl_utils
    # Relative importance: range of utilities per attribute
    ranges = {attr: max(u.values()) - min(u.values()) for attr, u in utilities.items()}
    total_range = sum(ranges.values())
    importance = {attr: (r / total_range * 100) if total_range > 0 else 0
                  for attr, r in ranges.items()}
    return {
        'method': 'Rating-based (OLS)',
        'rating_var': rating_var,
        'attribute_vars': attribute_vars,
        'n': len(df_c),
        'r2': model.rsquared,
        'r2_adj': model.rsquared_adj,
        'f': model.fvalue,
        'p': model.f_pvalue,
        'utilities': utilities,
        'importance': importance,
        'model': model,
        'attr_info': attr_info,
    }, None


def run_conjoint_cbc(df_raw, choice_var, attribute_vars):
    """Choice-Based Conjoint via logistic regression. choice_var = 0/1."""
    from sklearn.linear_model import LogisticRegression
    df_c = df_raw[[choice_var] + attribute_vars].dropna()
    if len(df_c) < 20:
        return None, "Za ma\u0142o obserwacji (min. 20)."
    y = df_c[choice_var].astype(int)
    X_parts = []
    attr_info = {}
    for attr in attribute_vars:
        col = df_c[attr]
        if col.dtype == object or col.nunique() <= 8:
            dummies = pd.get_dummies(col.astype(str), prefix=attr, drop_first=True)
            X_parts.append(dummies)
            attr_info[attr] = {'type': 'categorical', 'levels': list(dummies.columns)}
        else:
            X_parts.append(col.rename(attr).to_frame())
            attr_info[attr] = {'type': 'numeric', 'levels': [attr]}
    X = pd.concat(X_parts, axis=1).astype(float)
    X_const = sm.add_constant(X)
    try:
        model = sm.Logit(y, X_const).fit(disp=False)
    except Exception as e:
        return None, str(e)
    utilities = {}
    for attr, info in attr_info.items():
        lvl_utils = {}
        for lv in info['levels']:
            lvl_utils[lv] = model.params.get(lv, 0.0)
        mean_u = np.mean(list(lvl_utils.values())) if lvl_utils else 0
        utilities[attr] = {k: v - mean_u for k, v in lvl_utils.items()}
    ranges = {attr: max(u.values()) - min(u.values()) if u else 0 for attr, u in utilities.items()}
    total_range = sum(ranges.values())
    importance = {attr: (r / total_range * 100) if total_range > 0 else 0
                  for attr, r in ranges.items()}
    return {
        'method': 'CBC (Logit)',
        'choice_var': choice_var,
        'attribute_vars': attribute_vars,
        'n': len(df_c),
        'llr': model.llr,
        'llr_pvalue': model.llr_pvalue,
        'pseudo_r2': model.prsquared,
        'utilities': utilities,
        'importance': importance,
        'model': model,
        'attr_info': attr_info,
    }, None


def export_conjoint_to_excel(writer, conjoint_results, var_labels):
    workbook = writer.book
    ws = workbook.add_worksheet('Conjoint')
    fmt_t  = workbook.add_format({'bold': True, 'bg_color': '#1F4E79', 'font_color': 'white',
                                   'border': 1, 'align': 'center', 'font_size': 11})
    fmt_s  = workbook.add_format({'bold': True, 'bg_color': '#D6E4F0', 'border': 1})
    fmt_h  = workbook.add_format({'bold': True, 'bg_color': '#F2F2F2', 'border': 1, 'align': 'center'})
    fmt_lbl= workbook.add_format({'border': 1, 'bold': True})
    fmt_val= workbook.add_format({'border': 1, 'align': 'right', 'num_format': '#,##0.000'})
    fmt_pct= workbook.add_format({'border': 1, 'align': 'right', 'num_format': '0.0"%"'})
    fmt_str= workbook.add_format({'border': 1})
    ws.set_column(0, 0, 35)
    ws.set_column(1, 5, 16)
    row = 0
    for res in conjoint_results:
        if res.get('error'):
            ws.write(row, 0, f"B\u0141\u0104D: {res['error']}", fmt_s); row += 2; continue
        ws.merge_range(row, 0, row, 4, f"CONJOINT -- {res['method']}", fmt_t); row += 1
        # Model summary
        ws.write(row, 0, "N", fmt_lbl)
        ws.write(row, 1, res['n'], workbook.add_format({'border':1,'align':'right','num_format':'#,##0'}))
        row += 1
        if 'r2' in res:
            ws.write(row, 0, "R\u00b2", fmt_lbl); ws.write(row, 1, res['r2'], fmt_val); row += 1
            ws.write(row, 0, "R\u00b2 skor.", fmt_lbl); ws.write(row, 1, res['r2_adj'], fmt_val); row += 1
            ws.write(row, 0, "F / p", fmt_lbl)
            ws.write(row, 1, res['f'], fmt_val); ws.write(row, 2, res['p'], fmt_val); row += 1
        if 'pseudo_r2' in res:
            ws.write(row, 0, "Pseudo R\u00b2 (McFadden)", fmt_lbl); ws.write(row, 1, res['pseudo_r2'], fmt_val); row += 1
        row += 1
        # Importance
        ws.write(row, 0, "Wa\u017cno\u015b\u0107 atrybut\u00f3w (%)", fmt_s); row += 1
        for attr, imp in sorted(res['importance'].items(), key=lambda x: -x[1]):
            ws.write(row, 0, f"[{attr}] {var_labels.get(attr, attr)}", fmt_lbl)
            ws.write(row, 1, imp, fmt_pct); row += 1
        row += 1
        # Utilities
        ws.write(row, 0, "U\u017cyteczno\u015bci cz\u0105stkowe (part-worth utilities)", fmt_s); row += 1
        ws.write(row, 0, "Atrybut / Poziom", fmt_h)
        ws.write(row, 1, "U\u017cyteczno\u015b\u0107", fmt_h); row += 1
        for attr, utils in res['utilities'].items():
            ws.write(row, 0, f"[{attr}] {var_labels.get(attr, attr)}", fmt_lbl)
            ws.write(row, 1, "", fmt_str); row += 1
            for level, util in sorted(utils.items(), key=lambda x: -x[1]):
                ws.write(row, 0, f"  {level}", fmt_str)
                ws.write(row, 1, util, fmt_val); row += 1
        row += 3


# =============================================================
# MAXDIFF ANALYSIS
# =============================================================

def run_maxdiff(df_raw, task_pairs, item_values):
    """
    MaxDiff scoring from paired Best/Worst columns.
    task_pairs: list of (best_col, worst_col) tuples
    item_values: list of unique item labels (strings) that appear in those columns
    Returns: DataFrame with item scores and ranks.
    """
    n_resp = len(df_raw)
    counts = {item: {'best': 0, 'worst': 0, 'shown': 0} for item in item_values}
    for best_col, worst_col in task_pairs:
        if best_col not in df_raw.columns or worst_col not in df_raw.columns:
            continue
        best_series  = df_raw[best_col].dropna().astype(str)
        worst_series = df_raw[worst_col].dropna().astype(str)
        for item in item_values:
            counts[item]['best']  += (best_series  == str(item)).sum()
            counts[item]['worst'] += (worst_series == str(item)).sum()
            counts[item]['shown'] += ((best_series == str(item)) | (worst_series == str(item))).sum()
    rows = []
    for item in item_values:
        b = counts[item]['best']
        w = counts[item]['worst']
        shown = counts[item]['shown']
        bw_score = b - w
        bw_pct   = bw_score / n_resp * 100 if n_resp > 0 else 0
        rows.append({'Item': item, 'Best [N]': b, 'Worst [N]': w,
                     'B-W Score': bw_score, 'B-W Score (%)': round(bw_pct, 2),
                     'Pokazano [N]': shown})
    df_scores = pd.DataFrame(rows).sort_values('B-W Score', ascending=False).reset_index(drop=True)
    df_scores.insert(0, 'Ranking', range(1, len(df_scores) + 1))
    # Rescale to 0-100 (most positive = 100)
    mn, mx = df_scores['B-W Score'].min(), df_scores['B-W Score'].max()
    if mx > mn:
        df_scores['Wynik standaryzowany (0-100)'] = ((df_scores['B-W Score'] - mn) / (mx - mn) * 100).round(1)
    else:
        df_scores['Wynik standaryzowany (0-100)'] = 50.0
    return df_scores


def export_maxdiff_to_excel(writer, maxdiff_results, var_labels):
    workbook = writer.book
    ws = workbook.add_worksheet('MaxDiff')
    fmt_t   = workbook.add_format({'bold': True, 'bg_color': '#1F4E79', 'font_color': 'white',
                                    'border': 1, 'align': 'center', 'font_size': 11})
    fmt_h   = workbook.add_format({'bold': True, 'bg_color': '#F2F2F2', 'border': 1, 'align': 'center'})
    fmt_lbl = workbook.add_format({'border': 1})
    fmt_n   = workbook.add_format({'border': 1, 'align': 'right', 'num_format': '#,##0'})
    fmt_val = workbook.add_format({'border': 1, 'align': 'right', 'num_format': '#,##0.00'})
    fmt_pct = workbook.add_format({'border': 1, 'align': 'right', 'num_format': '0.0'})
    ws.set_column(0, 0, 8); ws.set_column(1, 1, 35)
    ws.set_column(2, 7, 18)
    row = 0
    for res in maxdiff_results:
        ws.merge_range(row, 0, row, 6, f"MaxDiff -- {res['name']}", fmt_t); row += 1
        ws.write(row, 0, f"N respondent\u00f3w: {res['n_resp']}  |  Liczba zestaw\u00f3w: {res['n_tasks']}",
                 workbook.add_format({'italic': True})); row += 2
        cols_out = ['Ranking', 'Item', 'Best [N]', 'Worst [N]', 'B-W Score', 'B-W Score (%)', 'Wynik standaryzowany (0-100)']
        for ci, col in enumerate(cols_out):
            ws.write(row, ci, col, fmt_h)
        row += 1
        df_s = res['scores']
        for _, r in df_s.iterrows():
            ws.write(row, 0, int(r['Ranking']), fmt_n)
            ws.write(row, 1, str(r['Item']), fmt_lbl)
            ws.write(row, 2, int(r['Best [N]']), fmt_n)
            ws.write(row, 3, int(r['Worst [N]']), fmt_n)
            ws.write(row, 4, float(r['B-W Score']), fmt_val)
            ws.write(row, 5, float(r['B-W Score (%)']), fmt_pct)
            ws.write(row, 6, float(r['Wynik standaryzowany (0-100)']), fmt_pct)
            row += 1
        row += 3


# SESSION STATE
# -------------------------------------------------------------
if 'mrs_sets'            not in st.session_state: st.session_state.mrs_sets = {}
if 'matrix_sets'         not in st.session_state: st.session_state.matrix_sets = {}
if 'matrix_results'      not in st.session_state: st.session_state.matrix_results = []
if 'custom_var_labels'   not in st.session_state: st.session_state.custom_var_labels = {}
if 'custom_val_labels'   not in st.session_state: st.session_state.custom_val_labels = {}
if 'user_cleared_val_labels' not in st.session_state: st.session_state.user_cleared_val_labels = set()
if 'box_sets'            not in st.session_state: st.session_state.box_sets = defaultdict(dict)
if 'segmentations'       not in st.session_state: st.session_state.segmentations = []
if 'hclust_results'      not in st.session_state: st.session_state.hclust_results = []
if 'logistic_results'    not in st.session_state: st.session_state.logistic_results = []
if 'recodings'           not in st.session_state: st.session_state.recodings = []
if 'cleaning_ops'        not in st.session_state: st.session_state.cleaning_ops = []  # [{cols, ops}]
if 'results'             not in st.session_state: st.session_state.results = {'czestosci': {}, 'krzyzowe': {}, 'srednie': {}, 'opisowe': {}, 'korelacje': {}}
if 'chi_results'         not in st.session_state: st.session_state.chi_results = {}
if 'custom_missing'      not in st.session_state: st.session_state.custom_missing = {}
if 'weights'             not in st.session_state: st.session_state.weights = None
if 'weight_targets'      not in st.session_state: st.session_state.weight_targets = {}
if 'treat_empty_as_miss' not in st.session_state: st.session_state.treat_empty_as_miss = False
if 'regression_results'  not in st.session_state: st.session_state.regression_results = []
if 'anova_results'       not in st.session_state: st.session_state.anova_results = []
if 'factor_results'      not in st.session_state: st.session_state.factor_results = []
if 'reg_blocks'          not in st.session_state: st.session_state.reg_blocks = [[]]
if 'conjoint_results'    not in st.session_state: st.session_state.conjoint_results = []
if 'maxdiff_results'     not in st.session_state: st.session_state.maxdiff_results = []
if 'maxdiff_pairs'       not in st.session_state: st.session_state.maxdiff_pairs = [('', '')]
if 'data_source'         not in st.session_state: st.session_state.data_source = 'spss'
if 'excel_col_types'     not in st.session_state: st.session_state.excel_col_types = {}
if 'excel_sheet'         not in st.session_state: st.session_state.excel_sheet = None

# -------------------------------------------------------------
# ?? Source selection ??????????????????????????????????????????
st.sidebar.markdown("## \U0001f4c1 \u0179r\u00f3d\u0142o danych")
data_source = st.sidebar.radio(
    "Wybierz format pliku:",
    ["\U0001f4ca SPSS (.sav)", "\U0001f4c8 Excel (.xlsx)"],
    key="data_source_radio",
    horizontal=False,
)
is_spss  = (data_source == "\U0001f4ca SPSS (.sav)")
is_excel = not is_spss

st.sidebar.markdown("---")

if is_spss:
    uploaded_file = st.sidebar.file_uploader("Plik SPSS (.sav)", type="sav", label_visibility="collapsed")
    excel_file = None
else:
    uploaded_file = None
    excel_file = st.sidebar.file_uploader("Plik Excel (.xlsx)", type=["xlsx", "xls"], label_visibility="collapsed")

# ?? Stop if no file ???????????????????????????????????????????
if is_spss and uploaded_file is None:
    st.info("\U0001f448 Wczytaj plik SPSS (.sav) lub Excel z paska bocznego, aby rozpocz\u0105\u0107 prac\u0119.")
    st.stop()
if is_excel and excel_file is None:
    st.info("\U0001f448 Wczytaj plik SPSS (.sav) lub Excel z paska bocznego, aby rozpocz\u0105\u0107 prac\u0119.")
    st.stop()

# ?? Excel: sheet selector (shown inline, above spinner) ???????
if is_excel:
    try:
        xf = pd.ExcelFile(excel_file)
        sheet_names = xf.sheet_names
    except Exception as _xe:
        st.error(f"Nie mo\u017cna otworzy\u0107 pliku Excel: {_xe}")
        st.stop()

    if len(sheet_names) == 1:
        selected_sheet = sheet_names[0]
    else:
        selected_sheet = st.sidebar.selectbox(
            "Wybierz arkusz:",
            sheet_names,
            index=sheet_names.index(st.session_state.excel_sheet)
                  if st.session_state.excel_sheet in sheet_names else 0,
            key="sheet_selector"
        )
        st.session_state.excel_sheet = selected_sheet

# ?? Load data ?????????????????????????????????????????????????
with st.spinner("Wczytywanie i strukturyzowanie bazy..."):
    if is_spss:
        df_orig_raw, df_orig, meta_orig = load_spss_data(uploaded_file)
        loaded_name = uploaded_file.name
    else:
        _overrides_json = json.dumps(st.session_state.excel_col_types, sort_keys=True)
        _missing_json   = json.dumps(st.session_state.custom_missing, sort_keys=True)
        df_orig_raw, df_orig, meta_orig = load_excel_data(
            excel_file, selected_sheet,
            col_type_overrides_json=_overrides_json,
            custom_missing_json=_missing_json
        )
        loaded_name = excel_file.name
        # Apply text\u2192numeric encoding maps as value labels (only for newly encoded cols)
        _tnm = getattr(meta_orig, '_text_to_num_maps', {})
        for _col, _lmap in _tnm.items():
            if (_col not in st.session_state.custom_val_labels
                    and _col not in st.session_state.user_cleared_val_labels):
                st.session_state.custom_val_labels[_col] = _lmap

df_raw = df_orig_raw.copy()
df     = df_orig.copy()
var_labels = meta_orig.column_names_to_labels.copy()

# Track which columns are original (from file) vs derived (added in-session)
original_cols = set(df_orig_raw.columns)

# Apply user-edited variable labels (from S\u0142ownik tab)
for _col, _lbl in st.session_state.custom_var_labels.items():
    var_labels[_col] = _lbl

# Apply segmentations, recodings, cleaning only when the relevant state has changed.
# We fingerprint the inputs; if fingerprint matches the cached version, skip re-apply.
_pipeline_key = json.dumps({
    'seg':      len(st.session_state.segmentations),
    'rec':      len(st.session_state.recodings),
    'clean':    len(st.session_state.cleaning_ops),
    'miss':     len(st.session_state.custom_missing),
    'empty':    st.session_state.treat_empty_as_miss,
}, sort_keys=True)

# Always apply \u2014 these are fast dict/pandas ops and run on copies not originals
apply_segmentations(df_raw, df, var_labels, st.session_state.segmentations)
apply_recodings(df_raw, df, var_labels, st.session_state.recodings)

# Apply in-place text cleaning
apply_cleaning_ops(df_raw, df, st.session_state.cleaning_ops)

# Apply empty-as-missing
if st.session_state.treat_empty_as_miss:
    for c in df_raw.columns:
        if df_raw[c].dtype == object:
            df_raw[c] = df_raw[c].replace(r'^\s*$', np.nan, regex=True)
            df[c]     = df[c].replace(r'^\s*$', np.nan, regex=True)

# Apply custom missing values
for c, m_vals in st.session_state.custom_missing.items():
    if c in df_raw.columns:
        # Build a comprehensive replace list:
        # - original values (numeric)
        # - string versions: "999", "999.0"  (for object/text columns)
        # - int versions: 999 (for columns stored as int64)
        replace_vals = []
        for v in m_vals:
            replace_vals.append(v)
            try:
                iv = int(v)
                replace_vals.append(iv)
                replace_vals.append(str(iv))          # "999"
            except (ValueError, TypeError):
                pass
            try:
                replace_vals.append(str(v))           # "999.0"
                replace_vals.append(str(float(v)))
            except (ValueError, TypeError):
                pass
        replace_vals = list(dict.fromkeys(replace_vals))  # deduplicate, preserve order

        df_raw[c] = df_raw[c].replace(replace_vals, np.nan)
        if is_spss:
            # Also replace value-labelled versions in df
            label_vals = [
                meta_orig.variable_value_labels.get(c, {}).get(v)
                for v in m_vals
                if v in meta_orig.variable_value_labels.get(c, {})
            ]
            df[c] = df[c].replace(label_vals + replace_vals, np.nan)
        else:
            df[c] = df[c].replace(replace_vals, np.nan)

hidden_cols = set()
for set_data in st.session_state.mrs_sets.values():
    cols = set_data if isinstance(set_data, list) else set_data.get('cols', [])
    hidden_cols.update(cols)
visible_columns = [c for c in df.columns if c not in hidden_cols]
all_options     = visible_columns + list(st.session_state.mrs_sets.keys()) + list(st.session_state.matrix_sets.keys())
numeric_cols_raw = df_raw.select_dtypes(include=[np.number]).columns.tolist()
numeric_cols     = [c for c in numeric_cols_raw if c not in hidden_cols and c in visible_columns]

# Sidebar status
st.sidebar.markdown("---")
n_rows, n_cols = len(df_raw), len(df_raw.columns)
src_icon = "\U0001f4ca" if is_spss else "\U0001f4c8"
st.sidebar.success(f"{src_icon} **{loaded_name}**\n\n{n_rows:,} respondent\u00f3w \u00b7 {n_cols} zmiennych")

if st.session_state.weights is not None:
    st.sidebar.markdown("---")
    use_weights = st.sidebar.checkbox("\u2696\ufe0f Zastosuj wagi w analizach", value=True)
else:
    use_weights = False

st.sidebar.markdown("---")
st.sidebar.markdown("## \U0001f4cc Nawigacja")

_MENU_ITEMS = [
    "\U0001f3e0 Dashboard",
    "\U0001f4c1 Projekt i S\u0142ownik",
    "\U0001f6e0\ufe0f Przygotowanie Danych",
    "\U0001f4c8 Analizy i Tabele",
    "\U0001f4c9 Regresja",
    "\U0001f4ca ANOVA",
    "\U0001f4d0 Testy Normalno\u015bci",
    "\U0001f52c Analiza Czynnikowa",
    "\U0001f3af Skupienia i Segmentacja",
    "\U0001f4ca Conjoint",
    "\U0001f522 MaxDiff",
    "\u2601\ufe0f Chmura S\u0142\u00f3w",
    "\U0001f4be Eksport do Excela",
    "\U0001f4ca Eksport do PowerPoint",
]

# Allow tile clicks to navigate by writing to session state
if 'nav_to' not in st.session_state:
    st.session_state.nav_to = None
if 'current_menu' not in st.session_state:
    st.session_state.current_menu = _MENU_ITEMS[0]

# If a tile was clicked, update current_menu then clear nav_to
if st.session_state.nav_to and st.session_state.nav_to in _MENU_ITEMS:
    st.session_state.current_menu = st.session_state.nav_to
    st.session_state.nav_to = None

# Use key="current_menu" so Streamlit syncs radio <-> session_state automatically.
# This avoids the double-click bug caused by manually setting index= every rerun.
menu = st.sidebar.radio("", _MENU_ITEMS,
                         key="current_menu",
                         label_visibility="collapsed")

# =============================================================
# DASHBOARD
# =============================================================
if menu == "\U0001f3e0 Dashboard":
    n_rows, n_cols_db = len(df_raw), len(df_raw.columns)

    st.markdown(f"""
<div style="background:linear-gradient(90deg,#1F4E79,#2E75B6);
     padding:28px 36px;border-radius:12px;margin-bottom:24px;color:white;">
  <h2 style="margin:0;font-size:1.8rem;">System Analiz Openfield (SAO)</h2>
  <p style="margin:6px 0 0;opacity:.85;font-size:1rem;">
    \U0001f4c1 <b>{loaded_name}</b> &nbsp;\u00b7&nbsp;
    \U0001f465 {n_rows:,} respondent\u00f3w &nbsp;\u00b7&nbsp;
    \U0001f4ca {n_cols_db} zmiennych
  </p>
</div>
""", unsafe_allow_html=True)

    # Tile definitions: (emoji, label, description, menu_key)
    _TILES = [
        ("\U0001f4c1", "Projekt i S\u0142ownik",
         "Zapisz/wczytaj projekt, przegl\u0105daj s\u0142ownik zmiennych i edytuj etykiety",
         "\U0001f4c1 Projekt i S\u0142ownik"),
        ("\U0001f6e0\ufe0f", "Przygotowanie Danych",
         "Braki, rekodowanie, czyszczenie, zestawy MRS, pytania matrycowe, wa\u017cenie",
         "\U0001f6e0\ufe0f Przygotowanie Danych"),
        ("\U0001f4c8", "Analizy i Tabele",
         "Tablice cz\u0119sto\u015bci, tabele krzy\u017cowe, \u015brednie, statystyki opisowe, korelacje",
         "\U0001f4c8 Analizy i Tabele"),
        ("\U0001f4d0", "Testy Normalno\u015bci",
         "Shapiro-Wilk, Kolmogorov-Smirnov, Lilliefors, D\u2019Agostino \u2014 wykresy Q-Q i histogramy",
         "\U0001f4d0 Testy Normalno\u015bci"),
        ("\U0001f4c9", "Regresja",
         "OLS (liniowa) i logistyczna (binarna/wielomianowa) w jednym module",
         "\U0001f4c9 Regresja"),
        ("\U0001f4ca", "ANOVA",
         "Jednoczynnikowa ANOVA, test Levene\u2019a, eta\u00b2, Tukey HSD post-hoc",
         "\U0001f4ca ANOVA"),
        ("\U0001f52c", "Analiza Czynnikowa",
         "EFA, KMO, test Bartletta, wykres osypiska, macierz \u0142adunk\u00f3w",
         "\U0001f52c Analiza Czynnikowa"),
        ("\U0001f3af", "Skupienia i Segmentacja",
         "Hierarchiczne (dendrogram) + K-Means: dobr\u00f3r skupie\u0144, profile grup",
         "\U0001f3af Skupienia i Segmentacja"),
        ("\U0001f4ca", "Conjoint",
         "Rating-based (OLS) i CBC (Logit) \u2014 u\u017cyteczno\u015bci, wa\u017cno\u015b\u0107 atrybut\u00f3w",
         "\U0001f4ca Conjoint"),
        ("\U0001f522", "MaxDiff",
         "Best-Worst Scaling, B-W scores, ranking wa\u017cno\u015bci, wynik standaryzowany 0-100",
         "\U0001f522 MaxDiff"),
        ("\u2601\ufe0f", "Chmura S\u0142\u00f3w",
         "Wizualizacja pyta\u0144 otwartych, stop words, palety kolor\u00f3w, eksport PNG/JPG",
         "\u2601\ufe0f Chmura S\u0142\u00f3w"),
        ("\U0001f4be", "Eksport do Excela",
         "Raport analityczny, wykresy, baza danych, spis tre\u015bci z hiperlink\u00f3w",
         "\U0001f4be Eksport do Excela"),
        ("\U0001f4ca", "Eksport do PowerPoint",
         "Edytowalne wykresy kolumnowe z cz\u0119sto\u015bci i tabel krzy\u017cowych \u2014 ka\u017cdy na osobnym slajdzie",
         "\U0001f4ca Eksport do PowerPoint"),
    ]

    # Status badges for tiles that have results
    def _has_results(menu_key):
        if "Cz\u0119sto\u015bci" in menu_key or "Analizy" in menu_key:
            return any(st.session_state.results.get(g)
                       for g in ['czestosci','krzyzowe','srednie','opisowe','korelacje']) \
                   or bool(st.session_state.matrix_results)
        if "Regresja" in menu_key and "Logistyczna" not in menu_key:
            return bool(st.session_state.regression_results)
        if "Regresja" in menu_key and "Logistyczna" in menu_key:
            return bool(st.session_state.logistic_results)
        if "ANOVA" in menu_key:
            return bool(st.session_state.anova_results)
        if "Czynnikowa" in menu_key:
            return bool(st.session_state.factor_results)
        if "Skupienia" in menu_key or "Segmentacja" in menu_key:
            return bool(st.session_state.hclust_results)
        if "Conjoint" in menu_key:
            return bool(st.session_state.conjoint_results)
        if "MaxDiff" in menu_key:
            return bool(st.session_state.maxdiff_results)
        return False

    st.markdown("### \U0001f4cc Modu\u0142y")

    # Render tiles in rows of 3
    for row_start in range(0, len(_TILES), 3):
        cols = st.columns(3, gap="medium")
        for ci, tile in enumerate(_TILES[row_start:row_start+3]):
            icon, title, desc, key = tile
            has_res = _has_results(key)
            badge = " \u2705" if has_res else ""
            with cols[ci]:
                st.markdown(f"""
<div style="
    background:#fff;border-radius:10px;
    border:1.5px solid {'#2E75B6' if has_res else '#e0e0e0'};
    padding:20px 22px 14px;height:160px;
    box-shadow:0 2px 8px rgba(0,0,0,0.06);
    transition:border-color .2s;
    display:flex;flex-direction:column;justify-content:space-between;
">
  <div>
    <span style="font-size:1.6rem">{icon}</span>
    <span style="font-size:1rem;font-weight:700;color:#1F4E79;margin-left:8px">{title}{badge}</span>
    <p style="font-size:0.78rem;color:#595959;margin:8px 0 0;line-height:1.4">{desc}</p>
  </div>
</div>
""", unsafe_allow_html=True)
                if st.button(f"Przejd\u017a \u2192", key=f"tile_{key}",
                              use_container_width=True):
                    st.session_state.nav_to = key
                    st.rerun()
        st.write("")  # vertical gap between rows

    # Quick stats row
    st.divider()
    st.markdown("### \U0001f4cb Podsumowanie sesji")
    qs = st.columns(6)
    qs[0].metric("\U0001f5c2\ufe0f Zestawy MRS",      len(st.session_state.mrs_sets))
    qs[1].metric("\U0001f4cb Pyt. matrycowe",          len(st.session_state.matrix_sets))
    qs[2].metric("\U0001f504 Rekodowania",             len(st.session_state.recodings))
    qs[3].metric("\U0001f9f9 Regu\u0142y czyszczenia", len(st.session_state.cleaning_ops))
    n_analyses = (
        sum(1 for g in ['czestosci','krzyzowe','srednie','opisowe','korelacje']
            if st.session_state.results.get(g))
        + bool(st.session_state.regression_results)
        + bool(st.session_state.anova_results)
        + bool(st.session_state.factor_results)
        + bool(st.session_state.conjoint_results)
        + bool(st.session_state.maxdiff_results)
    )
    qs[4].metric("\U0001f4ca Wykonane analizy",  n_analyses)
    qs[5].metric("\u2696\ufe0f Wagi",
                 "\u2705 Tak" if st.session_state.weights is not None else "\u274c Nie")

# -------------------------------------------------------------
# MODUL 1: PROJEKT I SLOWNIK
# -------------------------------------------------------------
elif menu == "\U0001f4c1 Projekt i S\u0142ownik":
    module_header("\U0001f4c1", "Projekt i S\u0142ownik")
    tab_proj, tab_summary, tab_dict = st.tabs(["\u2699\ufe0f Projekt", "\U0001f4cb Podsumowanie Bazy", "\U0001f4d6 S\u0142ownik Zmiennych"])

    with tab_proj:
        st.markdown("#### Zarz\u0105dzanie projektem")

        # \u2500\u2500 Helpers for serialising / deserialising DataFrames \u2500
        def _df_to_dict(df):
            """Serialize a DataFrame to a JSON-safe dict."""
            if df is None or not isinstance(df, pd.DataFrame):
                return None
            return {'__df__': True, 'data': df.to_json(orient='split')}

        def _dict_to_df(d):
            """Deserialize a DataFrame from saved dict."""
            if not isinstance(d, dict) or not d.get('__df__'):
                return None
            return pd.read_json(d['data'], orient='split')

        def _ser_results(res_dict):
            """Convert {title: DataFrame} dict to JSON-safe format."""
            return {k: _df_to_dict(v) for k, v in res_dict.items()}

        def _deser_results(raw):
            """Restore {title: DataFrame} from saved format."""
            if not isinstance(raw, dict):
                return {}
            return {k: _dict_to_df(v) for k, v in raw.items()
                    if _dict_to_df(v) is not None}

        def _ser_matrix_results(lst):
            """Serialize matrix_results list."""
            out = []
            for entry in lst:
                e = {}
                for k, v in entry.items():
                    e[k] = _df_to_dict(v) if isinstance(v, pd.DataFrame) else v
                out.append(e)
            return out

        def _deser_matrix_results(lst):
            if not isinstance(lst, list):
                return []
            out = []
            for entry in lst:
                e = {}
                for k, v in entry.items():
                    e[k] = _dict_to_df(v) if isinstance(v, dict) and v.get('__df__') else v
                out.append(e)
            return out

        def _safe_val(v):
            """Convert any value to a JSON-safe type, silently dropping non-serializable objects."""
            if v is None or isinstance(v, (bool, int, float, str)):
                return v
            if isinstance(v, pd.DataFrame):
                return _df_to_dict(v)
            if isinstance(v, np.ndarray):
                return v.tolist()
            if isinstance(v, (np.integer,)):
                return int(v)
            if isinstance(v, (np.floating,)):
                return float(v)
            if isinstance(v, (np.bool_,)):
                return bool(v)
            if isinstance(v, dict):
                return {kk: _safe_val(vv) for kk, vv in v.items()}
            if isinstance(v, (list, tuple)):
                return [_safe_val(i) for i in v]
            # Try basic JSON round-trip; if it fails, drop the value
            try:
                import json as _json
                _json.dumps(v)
                return v
            except (TypeError, ValueError):
                return None   # drop silently

        def _ser_reg_results(lst):
            """Serialize regression results \u2014 drop unserialisable model objects."""
            out = []
            for res in lst:
                e = {}
                for k, v in res.items():
                    if k == 'model':
                        continue   # statsmodels object \u2014 always skip
                    e[k] = _safe_val(v)
                out.append(e)
            return out

        def _deser_reg_results(lst):
            if not isinstance(lst, list):
                return []
            out = []
            for entry in lst:
                e = {}
                for k, v in entry.items():
                    if isinstance(v, dict) and v.get('__df__'):
                        e[k] = _dict_to_df(v)
                    else:
                        e[k] = v
                e['model'] = None   # model not restored \u2014 display from coef_df
                out.append(e)
            return out

        def _ser_factor_results(lst):
            out = []
            for res in lst:
                e = {k: _safe_val(v) for k, v in res.items()}
                out.append(e)
            return out

        def _deser_factor_results(lst):
            if not isinstance(lst, list):
                return []
            out = []
            for entry in lst:
                e = {}
                for k, v in entry.items():
                    e[k] = _dict_to_df(v) if isinstance(v, dict) and v.get('__df__') else v
                out.append(e)
            return out

        def _ser_hclust(lst):
            out = []
            for res in lst:
                e = {k: _safe_val(v) for k, v in res.items()}
                out.append(e)
            return out

        def _deser_hclust(lst):
            if not isinstance(lst, list):
                return []
            out = []
            for entry in lst:
                e = {}
                for k, v in entry.items():
                    if isinstance(v, dict) and v.get('__df__'):
                        e[k] = _dict_to_df(v)
                    elif isinstance(v, list) and k == 'Z':
                        e[k] = v   # keep as list
                    else:
                        e[k] = v
                out.append(e)
            return out

        def _build_project_data(include_results=True):
            data = {
                "_version": "3.0",
                "_saved_at": datetime.datetime.now().isoformat(timespec="seconds"),
                "_source": "spss" if is_spss else "excel",
                "meta": {
                    "name":   st.session_state.get("proj_name_inp", ""),
                    "author": st.session_state.get("proj_author_inp", ""),
                    "desc":   st.session_state.get("proj_desc_inp", ""),
                },
                # Config
                "treat_empty_as_miss": st.session_state.treat_empty_as_miss,
                "custom_missing":      st.session_state.custom_missing,
                "excel_col_types":     st.session_state.excel_col_types,
                "mrs_sets":            st.session_state.mrs_sets,
                "matrix_sets":         st.session_state.matrix_sets,
                "box_sets":            dict(st.session_state.box_sets),
                "custom_var_labels":   st.session_state.custom_var_labels,
                "custom_val_labels":   st.session_state.custom_val_labels,
                "user_cleared_val_labels": list(st.session_state.user_cleared_val_labels),
                "recodings":           st.session_state.recodings,
                "cleaning_ops":        st.session_state.cleaning_ops,
                "segmentations":       st.session_state.segmentations,
                "weight_targets":      st.session_state.weight_targets,
                "weights":             list(st.session_state.weights)
                                       if st.session_state.weights is not None else None,
                "maxdiff_pairs":       st.session_state.maxdiff_pairs,
                "reg_blocks":          st.session_state.reg_blocks,
            }
            if include_results:
                data["results"] = {
                    "czestosci": _ser_results(st.session_state.results.get('czestosci', {})),
                    "krzyzowe":  _ser_results(st.session_state.results.get('krzyzowe', {})),
                    "srednie":   _ser_results(st.session_state.results.get('srednie', {})),
                    "opisowe":   _ser_results(st.session_state.results.get('opisowe', {})),
                    "korelacje": _ser_results(st.session_state.results.get('korelacje', {})),
                }
                data["chi_results"]       = st.session_state.chi_results
                data["matrix_results"]    = _ser_matrix_results(st.session_state.matrix_results)
                data["regression_results"]= _ser_reg_results(st.session_state.regression_results)
                data["logistic_results"]  = _ser_reg_results(st.session_state.logistic_results)
                data["anova_results"]     = _ser_reg_results(st.session_state.anova_results)
                data["factor_results"]    = _ser_factor_results(st.session_state.factor_results)
                data["conjoint_results"]  = _ser_factor_results(st.session_state.conjoint_results)
                data["maxdiff_results"]   = _ser_factor_results(st.session_state.maxdiff_results)
                data["hclust_results"]    = _ser_hclust(st.session_state.hclust_results)
            return data

        col1, col2 = st.columns(2)

        # \u2500\u2500 SAVE \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
        with col1:
            st.markdown("**\U0001f4e5 Zapisz projekt**")

            proj_name   = st.text_input("Nazwa projektu:", key="proj_name_inp",
                                         placeholder="np. Badanie satysfakcji klient\u00f3w 2025")
            proj_author = st.text_input("Autor:", key="proj_author_inp",
                                         placeholder="np. Jan Kowalski")
            proj_desc   = st.text_area("Opis / notatki:", key="proj_desc_inp",
                                        height=70,
                                        placeholder="Opcjonalny opis badania, wersji danych itp.")

            include_res = st.checkbox(
                "\U0001f4ca Do\u0142\u0105cz wyniki analiz do pliku",
                value=True, key="proj_save_results",
                help="Tablice cz\u0119sto\u015bci, tabele krzy\u017cowe, regresje, ANOVA, EFA itp. "
                     "Plik b\u0119dzie wi\u0119kszy, ale po wczytaniu nie trzeba b\u0119dzie "
                     "ponownie wyklikywa\u0107 analiz."
            )

            try:
                proj_json = json.dumps(
                    _build_project_data(include_results=include_res),
                    ensure_ascii=False, indent=2
                )
            except Exception as _se:
                proj_json = None
                st.error(f"B\u0142\u0105d serializacji: {_se}")

            safe_name = (proj_name or "Projekt").replace(" ", "_")[:40]
            if proj_json:
                st.download_button(
                    "\U0001f4e5 Zapisz projekt (.json)",
                    data=proj_json,
                    file_name=f"{safe_name}.json",
                    mime="application/json",
                    type="primary",
                    use_container_width=True
                )

            n_analyses = (
                sum(1 for g in ['czestosci','krzyzowe','srednie','opisowe','korelacje']
                    if st.session_state.results.get(g)) +
                bool(st.session_state.regression_results) +
                bool(st.session_state.logistic_results) +
                bool(st.session_state.anova_results) +
                bool(st.session_state.factor_results) +
                bool(st.session_state.conjoint_results) +
                bool(st.session_state.maxdiff_results)
            )
            n_config = (len(st.session_state.mrs_sets) + len(st.session_state.matrix_sets) +
                        len(st.session_state.recodings) + len(st.session_state.segmentations) +
                        len(st.session_state.custom_var_labels) + len(st.session_state.custom_missing))
            st.caption(f"Konfiguracja: {n_config} element\u00f3w | Analizy: {n_analyses} modu\u0142\u00f3w z wynikami")

        # \u2500\u2500 LOAD \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
        with col2:
            st.markdown("**\U0001f504 Wczytaj projekt**")
            proj_file = st.file_uploader("Wgraj plik projektu (.json)", type="json",
                                          key="proj_uploader")
            if proj_file is not None:
                try:
                    raw_data = json.loads(proj_file.getvalue())
                except Exception:
                    st.error("Nieprawid\u0142owy plik JSON.")
                    raw_data = None

                if raw_data is not None:
                    ver      = raw_data.get("_version", "1.0")
                    meta     = raw_data.get("meta", {})
                    saved_at = raw_data.get("_saved_at", "")
                    src_lbl  = raw_data.get("_source", "?")
                    has_res  = "results" in raw_data

                    with st.expander("\U0001f4cb Podgl\u0105d wczytywanego projektu", expanded=True):
                        if meta.get("name"):
                            st.markdown(f"**Nazwa:** {meta['name']}")
                        if meta.get("author"):
                            st.markdown(f"**Autor:** {meta['author']}")
                        if meta.get("desc"):
                            st.markdown(f"**Opis:** {meta['desc']}")
                        st.caption(
                            f"Wersja: {ver} \u00b7 "
                            f"\u0179r\u00f3d\u0142o: {'SPSS' if src_lbl=='spss' else 'Excel' if src_lbl=='excel' else src_lbl} \u00b7 "
                            f"Zapisano: {saved_at or 'nieznane'}"
                        )
                        if has_res:
                            n_saved = sum(
                                len(raw_data.get("results", {}).get(g, {}))
                                for g in ['czestosci','krzyzowe','srednie','opisowe','korelacje']
                            )
                            st.success(f"\U0001f4ca Plik zawiera wyniki analiz ({n_saved} tabel).")
                        else:
                            st.info("Plik zawiera tylko konfiguracj\u0119 (bez wynik\u00f3w analiz).")

                        summary_parts = []
                        for key, label in [("mrs_sets","zestawy MRS"),("matrix_sets","pytania matrycowe"),
                                           ("recodings","rekodowania"),("segmentations","segmentacje")]:
                            n = len(raw_data.get(key, {}))
                            if n: summary_parts.append(f"{n} {label}")
                        if summary_parts:
                            st.info(", ".join(summary_parts))

                    if st.button("\u2705 Przywr\u00f3\u0107 z pliku",
                                  type="primary", use_container_width=True):
                        # Config
                        st.session_state.mrs_sets          = raw_data.get("mrs_sets", {})
                        st.session_state.matrix_sets       = raw_data.get("matrix_sets", {})
                        st.session_state.custom_var_labels = raw_data.get("custom_var_labels", {})
                        st.session_state.custom_val_labels = raw_data.get("custom_val_labels", {})
                        st.session_state.box_sets          = defaultdict(dict, raw_data.get("box_sets", {}))
                        st.session_state.segmentations     = raw_data.get("segmentations", [])
                        st.session_state.recodings         = raw_data.get("recodings", [])
                        st.session_state.cleaning_ops      = raw_data.get("cleaning_ops", [])
                        st.session_state.custom_missing    = raw_data.get("custom_missing", {})
                        st.session_state.user_cleared_val_labels = set(raw_data.get("user_cleared_val_labels", []))
                        st.session_state.weight_targets    = raw_data.get("weight_targets", {})
                        st.session_state.treat_empty_as_miss = raw_data.get("treat_empty_as_miss", False)
                        st.session_state.excel_col_types   = raw_data.get("excel_col_types", {})
                        st.session_state.maxdiff_pairs     = raw_data.get("maxdiff_pairs", [('', '')])
                        st.session_state.reg_blocks        = raw_data.get("reg_blocks", [[]])
                        w = raw_data.get("weights")
                        st.session_state.weights = np.array(w) if w else None
                        # Results
                        if has_res:
                            raw_res = raw_data.get("results", {})
                            st.session_state.results = {
                                'czestosci': _deser_results(raw_res.get('czestosci', {})),
                                'krzyzowe':  _deser_results(raw_res.get('krzyzowe', {})),
                                'srednie':   _deser_results(raw_res.get('srednie', {})),
                                'opisowe':   _deser_results(raw_res.get('opisowe', {})),
                                'korelacje': _deser_results(raw_res.get('korelacje', {})),
                            }
                            st.session_state.chi_results        = raw_data.get("chi_results", {})
                            st.session_state.matrix_results     = _deser_matrix_results(raw_data.get("matrix_results", []))
                            st.session_state.regression_results = _deser_reg_results(raw_data.get("regression_results", []))
                            st.session_state.logistic_results   = _deser_reg_results(raw_data.get("logistic_results", []))
                            st.session_state.anova_results      = _deser_reg_results(raw_data.get("anova_results", []))
                            st.session_state.factor_results     = _deser_factor_results(raw_data.get("factor_results", []))
                            st.session_state.conjoint_results   = _deser_factor_results(raw_data.get("conjoint_results", []))
                            st.session_state.maxdiff_results    = _deser_factor_results(raw_data.get("maxdiff_results", []))
                            st.session_state.hclust_results     = _deser_hclust(raw_data.get("hclust_results", []))
                        st.success("\u2705 Projekt wczytany pomy\u015blnie!")
                        st.rerun()

    with tab_summary:
        st.markdown("#### Podsumowanie bazy danych")
        total_rows = len(df_raw)
        total_cols_n = len(df_raw.columns)
        num_c = len(df_raw.select_dtypes(include=[np.number]).columns)
        cat_c = total_cols_n - num_c
        complete = int(df_raw.dropna().shape[0])
        miss_cells = int(df_raw.isna().sum().sum())
        total_cells = total_rows * total_cols_n
        miss_pct = miss_cells / total_cells * 100 if total_cells > 0 else 0

        m1, m2, m3, m4 = st.columns(4)
        m1.metric("\U0001f465 Respondenci", f"{total_rows:,}")
        m2.metric("\U0001f4ca Zmienne", f"{total_cols_n:,}")
        m3.metric("\U0001f522 Numeryczne", f"{num_c:,}")
        m4.metric("\U0001f524 Kategoryczne", f"{cat_c:,}")
        m5, m6, m7, m8 = st.columns(4)
        m5.metric("\u2705 Kompletne wiersze", f"{complete:,}")
        m6.metric("\u26a0\ufe0f Wiersze z brakami", f"{total_rows - complete:,}")
        m7.metric("\U0001f573\ufe0f Kom\u00f3rki z NaN", f"{miss_cells:,}")
        m8.metric("\U0001f4c9 % brak\u00f3w w bazie", f"{miss_pct:.1f}%")

        st.divider()
        st.markdown("**Braki danych i statystyki wg zmiennej:**")
        summ_rows = []
        for c in df_raw.columns:
            n_miss = df_raw[c].isna().sum()
            summ_rows.append({
                'Zmienna': c, 'Etykieta': var_labels.get(c, ''),
                'Typ danych': str(df_raw[c].dtype),
                'Braki [N]': n_miss,
                'Braki [%]': round(n_miss / total_rows * 100, 1),
                'Unikalne warto\u015bci': df_raw[c].nunique(),
                'Min': df_raw[c].min() if pd.api.types.is_numeric_dtype(df_raw[c]) else '--',
                'Max': df_raw[c].max() if pd.api.types.is_numeric_dtype(df_raw[c]) else '--',
            })
        st.dataframe(pd.DataFrame(summ_rows), use_container_width=True, height=400)

    with tab_dict:
        st.markdown("#### S\u0142ownik zmiennych")
        st.info("\U0001f3f7\ufe0f Edycja etykiet zmiennych i warto\u015bci dost\u0119pna w module "
                "**Przygotowanie Danych \u2192 Etykiety**.")
        dict_rows = []
        for col in df_raw.columns:
            orig_lbl = meta_orig.column_names_to_labels.get(col, "")
            curr_lbl = var_labels.get(col, orig_lbl)
            has_vl   = (bool(meta_orig.variable_value_labels.get(col))
                        or bool(st.session_state.custom_val_labels.get(col)))
            is_custom_var = col in st.session_state.custom_var_labels
            is_custom_val = col in st.session_state.custom_val_labels
            dict_rows.append({
                "Zmienna": col,
                "Etykieta bie\u017c\u0105ca": curr_lbl,
                "Etykieta oryginalna": orig_lbl,
                "Typ": str(df_raw[col].dtype),
                "Et. warto\u015bci": "\u2705" if has_vl else "--",
                "Zmodyfikowana": ("\U0001f3f7\ufe0f+\u270f\ufe0f" if is_custom_var and is_custom_val
                                  else "\U0001f3f7\ufe0f" if is_custom_val
                                  else "\u270f\ufe0f" if is_custom_var else ""),
            })
        st.dataframe(pd.DataFrame(dict_rows), use_container_width=True, height=420)
        st.caption("\u270f\ufe0f = zmieniona etykieta zmiennej | \U0001f3f7\ufe0f = zmienione etykiety warto\u015bci")


# -------------------------------------------------------------
# MODU? 2: PRZYGOTOWANIE DANYCH
# -------------------------------------------------------------
elif menu == "\U0001f6e0\ufe0f Przygotowanie Danych":
    module_header("\U0001f6e0\ufe0f", "Przygotowanie Danych")
    # For Excel: add an extra tab for type overrides
    if is_excel:
        tab_miss, tab_labels, tab_types, tab_clean, tab_mrs, tab_matrix, tab_weight, tab_recode, tab_box = st.tabs([
            "\u26a0\ufe0f Braki", "\U0001f3f7\ufe0f Etykiety", "\U0001f522 Typy", "\U0001f9f9 Czyszczenie",
            "\U0001f5f9 Wielokrotne odp.", "\U0001f4cb Matrycowe",
            "\u2696\ufe0f Wa\u017cenie", "\U0001f504 Rekodowanie", "\U0001f4e6 Top/Bottom Box"
        ])
    else:
        tab_miss, tab_labels, tab_clean, tab_mrs, tab_matrix, tab_weight, tab_recode, tab_box = st.tabs([
            "\u26a0\ufe0f Braki", "\U0001f3f7\ufe0f Etykiety", "\U0001f9f9 Czyszczenie",
            "\U0001f5f9 Wielokrotne odp.", "\U0001f4cb Matrycowe",
            "\u2696\ufe0f Wa\u017cenie", "\U0001f504 Rekodowanie", "\U0001f4e6 Top/Bottom Box"
        ])
        tab_types = None

    # -- BRAKI DANYCH ---------------------------------
    with tab_miss:
        st.markdown("#### Definiowanie brak\u00f3w danych")
        st.session_state.treat_empty_as_miss = st.checkbox(
            "\U0001f532 Traktuj puste ci\u0105gi tekstowe ('', ' ') jako braki danych (NaN)",
            value=st.session_state.treat_empty_as_miss,
            help="Zaznacz, je\u015bli puste kom\u00f3rki tekstowe maj\u0105 by\u0107 traktowane jako brak odpowiedzi."
        )
        st.divider()

        # SPSS-specific: default missing values from file
        if is_spss:
            use_spss_missing = st.checkbox("U\u017cywaj domy\u015blnych brak\u00f3w danych z pliku SPSS (zalecane)", value=True)
        else:
            use_spss_missing = False
            st.info("\U0001f4c8 Dane Excel nie maj\u0105 wbudowanych brak\u00f3w danych. Zdefiniuj warto\u015bci brak\u00f3w poni\u017cej.")

        if not use_spss_missing:
            col1, col2 = st.columns(2)
            with col1:
                st.markdown("**Globalne braki danych**")
                global_missing_str = st.text_input("Warto\u015bci oddzielone przecinkiem (np. 98, 99):")
                if st.button("Zastosuj do wszystkich zmiennych", use_container_width=True):
                    try:
                        vals = [float(x.strip()) for x in global_missing_str.split(',') if x.strip()]
                        for c in df_raw.columns:
                            st.session_state.custom_missing[c] = vals
                        st.success(f"Zastosowano braki {vals} do wszystkich zmiennych.")
                    except Exception:
                        st.error("Wprowad\u017a poprawne liczby.")
            with col2:
                st.markdown("**Braki dla konkretnej zmiennej**")
                var_missing = st.selectbox("Zmienna:", df_raw.columns,
                    format_func=lambda x: get_var_display_name(x, var_labels),
                    key="missing_var_select")
                var_missing_vals = st.multiselect("Warto\u015bci traktowane jako braki:",
                    sorted(df_raw[var_missing].dropna().unique(), key=lambda x: str(x)))
                if st.button("Ustaw braki dla tej zmiennej", use_container_width=True):
                    st.session_state.custom_missing[var_missing] = var_missing_vals
                    st.success("Zapisano!")

            if st.session_state.custom_missing:
                st.divider()
                st.markdown(f"**\U0001f4cb Zapisane braki danych ({len(st.session_state.custom_missing)} zmiennych):**")

                for _col, _vals in list(st.session_state.custom_missing.items()):
                    mc1, mc2, mc3, mc4 = st.columns([3, 3, 2, 1])
                    mc1.markdown(f"`{_col}` {var_labels.get(_col, '')}")
                    new_miss_str = mc2.text_input(
                        "", value=", ".join(str(v) for v in _vals),
                        key=f"miss_edit_{_col}", label_visibility="collapsed"
                    )
                    with mc3:
                        if st.button("\U0001f4be Zapisz", key=f"miss_upd_{_col}",
                                     use_container_width=True):
                            try:
                                new_vals = [float(x.strip())
                                            for x in new_miss_str.split(',') if x.strip()]
                                st.session_state.custom_missing[_col] = new_vals
                                st.rerun()
                            except Exception:
                                st.error("Nieprawid\u0142owe warto\u015bci.")
                    with mc4:
                        if st.button("\U0001f5d1\ufe0f", key=f"miss_del_{_col}",
                                     help=f"Usu\u0144 braki dla {_col}"):
                            st.session_state.custom_missing.pop(_col, None)
                            st.rerun()

                if st.button("\U0001f5d1\ufe0f Usu\u0144 wszystkie braki danych", type="secondary",
                              use_container_width=True, key="miss_clear_all"):
                    st.session_state.custom_missing = {}
                    st.rerun()

    # -- REKODOWANIE -----------------------------------
    with tab_recode:
        st.markdown("#### Rekodowanie zmiennych")
        st.info("Utw\u00f3rz now\u0105 zmienn\u0105 przez przypisanie nowych warto\u015bci do istniej\u0105cych kod\u00f3w. Dzia\u0142a zar\u00f3wno dla zmiennych **numerycznych**, jak i **tekstowych**. Nowa zmienna zostanie dodana do bazy.")

        col_r1, col_r2 = st.columns(2)
        with col_r1:
            # All visible columns -- numeric AND text
            all_recode_candidates = visible_columns
            src_var = st.selectbox(
                "Zmienna \u017ar\u00f3d\u0142owa:",
                all_recode_candidates,
                format_func=lambda x: get_var_display_name(x, var_labels),
                key="rec_src"
            )
            new_var_name = st.text_input("Nazwa nowej zmiennej:", value=f"{src_var}_r", key="rec_name")
            new_var_label = st.text_input("Etykieta nowej zmiennej:", key="rec_label")

            # Detect source type
            is_numeric_src = src_var in numeric_cols
            src_type_label = "numeryczna" if is_numeric_src else "tekstowa"
            st.caption(f"Typ zmiennej \u017ar\u00f3d\u0142owej: **{src_type_label}**")

            output_type = st.radio(
                "Typ warto\u015bci wyj\u015bciowych:",
                ["auto", "numeric", "text"],
                index=0,
                horizontal=True,
                key="rec_output_type",
                help="'auto' = automatyczne wykrycie (liczby je\u015bli mo\u017cliwe, tekst w p.p.)\n'numeric' = wymu\u015b liczby\n'text' = wymu\u015b tekst"
            )

        with col_r2:
            # Use df_raw for numeric, df (labeled) for categorical/text
            if is_numeric_src:
                unique_vals = sorted(df_raw[src_var].dropna().unique()) if src_var in df_raw.columns else []
            else:
                unique_vals = sorted(df[src_var].dropna().unique().astype(str)) if src_var in df.columns else []

            st.markdown(f"**Mapowanie warto\u015bci** (`{src_var}` -- {len(unique_vals)} unikalnych warto\u015bci):")

            mapping = {}
            if unique_vals:
                if len(unique_vals) > 50:
                    st.warning(f"Zmienna ma {len(unique_vals)} warto\u015bci -- wy\u015bwietlono pierwsze 50.")
                    unique_vals = unique_vals[:50]

                for val in unique_vals:
                    val_str = str(val)
                    # Show SPSS label in the field name if available (numeric vars)
                    if is_numeric_src:
                        try:
                            orig_label = meta_orig.variable_value_labels.get(src_var, {}).get(float(val_str), '')
                        except (ValueError, TypeError):
                            orig_label = ''
                    else:
                        orig_label = ''
                    display = f"{val_str}" + (f"  ({orig_label})" if orig_label else "")

                    # Always use text_input so both numeric and text targets are possible
                    new_val = st.text_input(
                        f"{display}  \u2192",
                        value=val_str,
                        key=f"rec_{src_var}_{val_str}"
                    )
                    mapping[val_str] = new_val
            else:
                st.info("Brak warto\u015bci do rekodowania (zmienna pusta lub brak danych).")

        if unique_vals and st.button("\u2705 Utw\u00f3rz rekodowan\u0105 zmienn\u0105", type="primary", use_container_width=True):
            if not new_var_name.strip():
                st.error("Podaj nazw\u0119 nowej zmiennej.")
            elif new_var_name.strip() in df_raw.columns:
                st.error(f"Zmienna `{new_var_name.strip()}` ju\u017c istnieje. Wybierz inn\u0105 nazw\u0119.")
            else:
                rec_entry = {
                    'source': src_var,
                    'new_name': new_var_name.strip(),
                    'label': new_var_label.strip() or f"Rekodowanie: {src_var}",
                    'mapping': mapping,
                    'output_type': output_type,
                }
                st.session_state.recodings.append(rec_entry)
                st.success(f"\u2705 Zmienna `{new_var_name.strip()}` zostanie dodana po od\u015bwie\u017ceniu.")
                st.rerun()

        if st.session_state.recodings:
            st.divider()
            st.markdown("**Zapisane rekodowania:**")
            to_del = None
            for i, rec in enumerate(st.session_state.recodings):
                c1, c2 = st.columns([5, 1])
                c1.markdown(f"- `{rec['new_name']}` \u2190 `{rec['source']}` \u00b7 *{rec['label']}* \u00b7 typ wyj\u015bcia: `{rec.get('output_type','auto')}`")
                if c2.button("\U0001f5d1\ufe0f", key=f"del_rec_{i}"):
                    to_del = i
            if to_del is not None:
                st.session_state.recodings.pop(to_del)
                st.rerun()

    # -- ETYKIETY ZMIENNYCH I WARTOSCI ----------------
    with tab_labels:
        st.markdown("#### \U0001f3f7\ufe0f Etykiety zmiennych i warto\u015bci")

        lab_sub_var, lab_sub_val, lab_sub_all = st.tabs([
            "\u270f\ufe0f Etykieta zmiennej",
            "\U0001f3f7\ufe0f Etykiety warto\u015bci (kody)",
            "\U0001f4cb Wszystkie moje zmiany"
        ])

        # \u2500\u2500 Sub-tab 1: Variable label \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
        with lab_sub_var:
            col_ev1, col_ev2 = st.columns(2)
            with col_ev1:
                edit_var = st.selectbox("Zmienna:", df_raw.columns,
                    format_func=lambda x: f"[{x}] {var_labels.get(x, x)}",
                    key="lab_edit_var")
                new_lbl = st.text_input("Nowa etykieta:",
                    value=var_labels.get(edit_var, ''), key="lab_new_lbl")
                bcol1, bcol2 = st.columns(2)
                with bcol1:
                    if st.button("\U0001f4be Zapisz", key="lab_save_var_lbl",
                                 use_container_width=True, type="primary"):
                        st.session_state.custom_var_labels[edit_var] = new_lbl
                        var_labels[edit_var] = new_lbl
                        st.success(f"Zapisano etykiet\u0119 `{edit_var}`.")
                        st.rerun()
                with bcol2:
                    if edit_var in st.session_state.custom_var_labels:
                        if st.button("\U0001f5d1\ufe0f Usu\u0144", key="lab_del_var_lbl",
                                     use_container_width=True):
                            st.session_state.custom_var_labels.pop(edit_var, None)
                            st.rerun()
            with col_ev2:
                st.markdown("**Etykiety warto\u015bci tej zmiennej:**")
                spss_vvl_preview = meta_orig.variable_value_labels.get(edit_var, {})
                cust_vl_preview  = st.session_state.custom_val_labels.get(edit_var, {})
                if spss_vvl_preview or cust_vl_preview:
                    preview_rows = []
                    for k, v in sorted(spss_vvl_preview.items()):
                        preview_rows.append({'Kod': k, 'Etykieta \u017ar\u00f3d\u0142owa': v,
                            'Niestandardowa': cust_vl_preview.get(str(k), '--')})
                    for rk, cv in cust_vl_preview.items():
                        if not any(str(r['Kod']) == rk for r in preview_rows):
                            preview_rows.append({'Kod': rk, 'Etykieta \u017ar\u00f3d\u0142owa': '--',
                                'Niestandardowa': cv})
                    st.dataframe(pd.DataFrame(preview_rows), use_container_width=True, hide_index=True)
                else:
                    uniq = sorted(df_raw[edit_var].dropna().unique())
                    st.info(f"Brak etykiet warto\u015bci. Unikalne warto\u015bci: "
                            f"{', '.join(str(v) for v in uniq[:20])}"
                            + (" ..." if len(uniq) > 20 else ""))

        # \u2500\u2500 Sub-tab 2: Value labels editor \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
        with lab_sub_val:
            st.caption("Zmie\u0144 wy\u015bwietlane etykiety warto\u015bci dla wybranej zmiennej.")
            col_vv1, col_vv2 = st.columns([1, 1])
            with col_vv1:
                val_edit_var = st.selectbox("Zmienna:", df_raw.columns,
                    format_func=lambda x: f"[{x}] {var_labels.get(x, x)}",
                    key="lab_val_edit_var")
                spss_codes = meta_orig.variable_value_labels.get(val_edit_var, {})
                raw_vals   = sorted(df_raw[val_edit_var].dropna().unique())
                st.caption(f"{len(raw_vals)} unikalnych warto\u015bci \u00b7 {len(spss_codes)} etykiet \u017ar\u00f3d\u0142owych")
                display_items = {}
                for v in raw_vals:
                    raw_str = str(v)
                    spss_lbl = spss_codes.get(v, spss_codes.get(raw_str, ''))
                    display_items[raw_str] = spss_lbl if spss_lbl else raw_str
                if len(display_items) > 50:
                    st.warning(f"Wy\u015bwietlono pierwsze 50 z {len(display_items)} warto\u015bci.")
                    display_items = dict(list(display_items.items())[:50])
            with col_vv2:
                st.markdown("**Nowe etykiety wy\u015bwietlania:**")
                existing_custom = st.session_state.custom_val_labels.get(val_edit_var, {})
                new_val_map = {}
                for raw_str, spss_lbl in display_items.items():
                    current = existing_custom.get(raw_str, spss_lbl)
                    hint = f" [{spss_lbl}]" if spss_lbl and spss_lbl != raw_str else ""
                    new_label = st.text_input(f"Kod {raw_str}{hint}:", value=current,
                        key=f"lab_vl_{val_edit_var}_{raw_str}")
                    new_val_map[raw_str] = new_label
                col_bsave, col_bclear = st.columns(2)
                with col_bsave:
                    if st.button("\U0001f4be Zapisz etykiety warto\u015bci", key="lab_save_val_lbls",
                                 use_container_width=True, type="primary"):
                        filtered = {k: v for k, v in new_val_map.items()
                                    if v.strip() and v.strip() != k}
                        st.session_state.custom_val_labels[val_edit_var] = filtered
                        st.session_state.user_cleared_val_labels.discard(val_edit_var)
                        st.success(f"Zapisano {len(filtered)} etykiet dla `{val_edit_var}`.")
                        st.rerun()
                with col_bclear:
                    if existing_custom and st.button("\U0001f5d1\ufe0f Usu\u0144",
                                                     key="lab_clear_val_lbls",
                                                     use_container_width=True):
                        st.session_state.custom_val_labels.pop(val_edit_var, None)
                        st.session_state.user_cleared_val_labels.add(val_edit_var)
                        st.rerun()

            # \u2500\u2500 Saved labels for current variable \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
            all_saved = st.session_state.custom_val_labels
            if all_saved:
                st.divider()
                st.markdown("**\U0001f4cb Zapisane etykiety warto\u015bci (wszystkie zmienne):**")
                for sv_col, sv_map in list(all_saved.items()):
                    with st.expander(
                        f"`{sv_col}` \u2014 {var_labels.get(sv_col, sv_col)} ({len(sv_map)} etykiet)",
                        expanded=(sv_col == val_edit_var)
                    ):
                        sv_edited = {}
                        for code, lbl in sorted(sv_map.items(), key=lambda x: x[0]):
                            ec1, ec2, ec3 = st.columns([1, 4, 1])
                            ec1.markdown(f"**{code}**")
                            new_v = ec2.text_input("", value=str(lbl),
                                key=f"sv_val_{sv_col}_{code}",
                                label_visibility="collapsed")
                            sv_edited[code] = new_v
                            with ec3:
                                if st.button("\U0001f5d1\ufe0f", key=f"sv_del_row_{sv_col}_{code}"):
                                    upd = {k: v for k, v in sv_map.items() if str(k) != str(code)}
                                    if upd:
                                        st.session_state.custom_val_labels[sv_col] = upd
                                    else:
                                        st.session_state.custom_val_labels.pop(sv_col, None)
                                        st.session_state.user_cleared_val_labels.add(sv_col)
                                    st.rerun()
                        sc1, sc2 = st.columns(2)
                        with sc1:
                            if st.button("\U0001f4be Zapisz", key=f"sv_save_{sv_col}",
                                         use_container_width=True):
                                st.session_state.custom_val_labels[sv_col] = sv_edited
                                st.success("Zapisano.")
                                st.rerun()
                        with sc2:
                            if st.button("\U0001f5d1\ufe0f Usu\u0144 ca\u0142\u0105 zmienn\u0105",
                                         key=f"sv_del_{sv_col}", use_container_width=True):
                                st.session_state.custom_val_labels.pop(sv_col, None)
                                st.session_state.user_cleared_val_labels.add(sv_col)
                                st.rerun()

        # \u2500\u2500 Sub-tab 3: All custom labels overview \u2500\u2500\u2500\u2500\u2500
        with lab_sub_all:
            has_any = st.session_state.custom_var_labels or st.session_state.custom_val_labels
            if not has_any:
                st.info("Nie wprowadzono \u017cadnych niestandardowych etykiet.")
            else:
                if st.session_state.custom_var_labels:
                    st.markdown("**\u270f\ufe0f Zmienione etykiety zmiennych:**")
                    changed_var = False
                    for col_cv, lbl_cv in list(st.session_state.custom_var_labels.items()):
                        rc1, rc2, rc3 = st.columns([3, 3, 1])
                        rc1.markdown(f"`{col_cv}`")
                        new_cv = rc2.text_input("", value=lbl_cv,
                            key=f"all_var_lbl_{col_cv}", label_visibility="collapsed")
                        with rc3:
                            if st.button("\U0001f5d1\ufe0f", key=f"all_del_var_{col_cv}"):
                                st.session_state.custom_var_labels.pop(col_cv, None)
                                st.rerun()
                        if new_cv != lbl_cv:
                            st.session_state.custom_var_labels[col_cv] = new_cv
                            changed_var = True
                    st.divider()

                if st.session_state.custom_val_labels:
                    st.markdown("**\U0001f3f7\ufe0f Zmienione etykiety warto\u015bci:**")
                    for col_cv, val_map in list(st.session_state.custom_val_labels.items()):
                        with st.expander(
                            f"`{col_cv}` \u2014 {var_labels.get(col_cv, col_cv)}"
                            f" ({len(val_map)} etykiet)", expanded=False
                        ):
                            edited_vmap = {}
                            for code, lbl in sorted(val_map.items(), key=lambda x: x[0]):
                                ec1, ec2 = st.columns([1, 3])
                                ec1.markdown(f"Kod **{code}**")
                                new_v = ec2.text_input("", value=str(lbl),
                                    key=f"all_val_{col_cv}_{code}",
                                    label_visibility="collapsed")
                                edited_vmap[code] = new_v
                            sc1, sc2 = st.columns(2)
                            with sc1:
                                if st.button("\U0001f4be Zapisz", key=f"all_save_val_{col_cv}",
                                             use_container_width=True):
                                    st.session_state.custom_val_labels[col_cv] = edited_vmap
                                    st.session_state.user_cleared_val_labels.discard(col_cv)
                                    st.success("Zapisano.")
                                    st.rerun()
                            with sc2:
                                if st.button("\U0001f5d1\ufe0f Usu\u0144",
                                             key=f"all_del_val_{col_cv}",
                                             use_container_width=True):
                                    st.session_state.custom_val_labels.pop(col_cv, None)
                                    st.session_state.user_cleared_val_labels.add(col_cv)
                                    st.rerun()

            st.divider()
            if st.button("\U0001f5d1\ufe0f Usu\u0144 WSZYSTKIE moje etykiety",
                         type="secondary", use_container_width=True, key="all_del_all_labels"):
                _all_val_cols = set(st.session_state.custom_val_labels.keys())
                st.session_state.custom_var_labels = {}
                st.session_state.custom_val_labels = {}
                st.session_state.user_cleared_val_labels = _all_val_cols
                st.rerun()

    # -- CZYSZCZENIE DANYCH -------------------------
    with tab_clean:
        st.markdown("#### \U0001f9f9 Czyszczenie danych tekstowych")
        st.info(
            "Operacje czyszczenia s\u0105 stosowane tymczasowo w tej sesji i zapisywane jako **rekodowania**. "
            "Mo\u017cesz stosowa\u0107 je globalnie (do wszystkich zmiennych tekstowych) lub tylko do wybranych zmiennych."
        )

        # \u2500\u2500 Helper: text columns \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
        text_cols = [c for c in visible_columns if df_raw[c].dtype == object]

        if not text_cols:
            st.warning("Brak zmiennych tekstowych w bazie danych.")
        else:
            # \u2500\u2500 SECTION A: Globally on all text columns \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
            st.markdown("### \U0001f310 Globalnie \u2014 wszystkie zmienne tekstowe")
            st.caption(f"Dzia\u0142a na {len(text_cols)} zmiennych tekstowych jednocze\u015bnie.")

            clean_col1, clean_col2 = st.columns(2)
            with clean_col1:
                st.markdown("**Czyszczenie spacji i bia\u0142ych znak\u00f3w:**")
                do_strip    = st.checkbox("Usu\u0144 spacje na pocz\u0105tku i ko\u0144cu", key="clean_g_strip", value=True)
                do_dbl_sp   = st.checkbox("Usu\u0144 podw\u00f3jne spacje",               key="clean_g_dbl")
                do_tabs     = st.checkbox("Usu\u0144 tabulatory (\\t)",                   key="clean_g_tabs")
                do_newlines = st.checkbox("Usu\u0144 znaki nowej linii (\\n, \\r)",       key="clean_g_nl")

                st.markdown("**Standaryzacja cudzys\u0142ow\u00f3w:**")
                do_quotes   = st.checkbox("Zamie\u0144 cudzys\u0142owy na standardowe ASCII (\u201c\u201d\u2018\u2019 \u2192 \"')", key="clean_g_quotes")

            with clean_col2:
                st.markdown("**Wielko\u015b\u0107 liter:**")
                case_mode = st.radio(
                    "Zastosuj:",
                    ["Bez zmian", "WIELKIE LITERY", "ma\u0142e litery", "Pierwsza Wielka (Title Case)"],
                    key="clean_g_case", horizontal=False
                )

                st.markdown("**Znaki specjalne:**")
                do_special = st.checkbox(
                    "Usu\u0144 znaki specjalne (pozostaw litery, cyfry i spacje)",
                    key="clean_g_special",
                    help="Usuwa wszystko poza literami (w tym polskimi), cyframi i spacjami."
                )

            if st.button("\u25b6\ufe0f Zastosuj globalnie do wszystkich zmiennych tekstowych",
                         type="primary", key="clean_global_apply", use_container_width=True):
                ops = {
                    'strip':    do_strip,
                    'dbl_sp':   do_dbl_sp,
                    'tabs':     do_tabs,
                    'newlines': do_newlines,
                    'quotes':   do_quotes,
                    'case':     {'Bez zmian': 'none', 'WIELKIE LITERY': 'upper',
                                 'ma\u0142e litery': 'lower', 'Pierwsza Wielka (Title Case)': 'title'}.get(case_mode, 'none'),
                    'special':  do_special,
                }
                if any(v for k, v in ops.items() if k != 'case') or ops['case'] != 'none':
                    st.session_state.cleaning_ops.append({'cols': list(text_cols), 'ops': ops})
                    st.success(f"\u2705 Czyszczenie zapisane dla {len(text_cols)} zmiennych i b\u0119dzie stosowane przy ka\u017cdym wczytaniu danych.")
                    st.rerun()
                else:
                    st.warning("Nie wybrano \u017cadnej operacji czyszczenia.")

            st.divider()

            # \u2500\u2500 SECTION B: Per-column \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
            st.markdown("### \U0001f527 Dla wybranych zmiennych")

            sel_cols = st.multiselect(
                "Wybierz zmienne do czyszczenia:",
                text_cols,
                format_func=lambda x: get_var_display_name(x, var_labels),
                key="clean_sel_cols"
            )

            if sel_cols:
                c1, c2 = st.columns(2)
                with c1:
                    st.markdown("**Spacje i bia\u0142e znaki:**")
                    p_strip    = st.checkbox("Usu\u0144 spacje na pocz\u0105tku i ko\u0144cu", key="clean_p_strip",  value=True)
                    p_dbl      = st.checkbox("Usu\u0144 podw\u00f3jne spacje",               key="clean_p_dbl")
                    p_tabs     = st.checkbox("Usu\u0144 tabulatory",                         key="clean_p_tabs")
                    p_nl       = st.checkbox("Usu\u0144 znaki nowej linii",                  key="clean_p_nl")
                    p_quotes   = st.checkbox("Zamie\u0144 cudzys\u0142owy na ASCII",         key="clean_p_quotes")
                with c2:
                    st.markdown("**Wielko\u015b\u0107 liter:**")
                    p_case = st.radio(
                        "Zastosuj:",
                        ["Bez zmian", "WIELKIE LITERY", "ma\u0142e litery", "Pierwsza Wielka"],
                        key="clean_p_case", horizontal=False
                    )
                    st.markdown("**Znaki specjalne:**")
                    p_special = st.checkbox(
                        "Usu\u0144 znaki specjalne",
                        key="clean_p_special"
                    )

                # Preview
                if sel_cols:
                    st.markdown("**Podgl\u0105d pierwszych 5 wierszy (przed / po):**")
                    preview_col = sel_cols[0]
                    orig = df_raw[preview_col].dropna().astype(str).head(5)
                    prev = orig.copy()
                    if p_strip:   prev = prev.str.strip()
                    if p_dbl:     prev = prev.str.replace(r' {2,}', ' ', regex=True)
                    if p_tabs:    prev = prev.str.replace('\t', ' ', regex=False)
                    if p_nl:      prev = prev.str.replace(r'[\n\r]', ' ', regex=True)
                    if p_quotes:
                        for old_q, new_q in [('\u201c','"'), ('\u201d','"'), ('\u201e','"'),
                                              ('\u2018',"'"), ('\u2019',"'"), ('\u201a',"'")]:
                            prev = prev.str.replace(old_q, new_q, regex=False)
                    if p_case == "WIELKIE LITERY":        prev = prev.str.upper()
                    elif p_case == "ma\u0142e litery":    prev = prev.str.lower()
                    elif p_case == "Pierwsza Wielka":     prev = prev.str.title()
                    if p_special:
                        prev = prev.str.replace(r'[^\w\s]', '', regex=True).str.replace('_','', regex=False)

                    preview_df = pd.DataFrame({'Przed': orig.values, 'Po': prev.values})
                    st.dataframe(preview_df, use_container_width=True, hide_index=True)

                if st.button("\u25b6\ufe0f Zastosuj dla wybranych zmiennych",
                             type="primary", key="clean_per_apply", use_container_width=True):
                    ops = {
                        'strip':    p_strip,
                        'dbl_sp':   p_dbl,
                        'tabs':     p_tabs,
                        'newlines': p_nl,
                        'quotes':   p_quotes,
                        'case':     {'Bez zmian': 'none', 'WIELKIE LITERY': 'upper',
                                     'ma\u0142e litery': 'lower', 'Pierwsza Wielka': 'title'}.get(p_case, 'none'),
                        'special':  p_special,
                    }
                    if any(v for k, v in ops.items() if k != 'case') or ops['case'] != 'none':
                        st.session_state.cleaning_ops.append({'cols': list(sel_cols), 'ops': ops})
                        st.success(f"\u2705 Czyszczenie zapisane dla {len(sel_cols)} zmiennych: {', '.join(sel_cols[:3])}{'...' if len(sel_cols) > 3 else ''}.")
                        st.rerun()
                    else:
                        st.warning("Nie wybrano \u017cadnej operacji czyszczenia.")

            # \u2500\u2500 Active cleaning rules panel \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
            if st.session_state.cleaning_ops:
                st.divider()
                st.markdown("### \U0001f9f9 Aktywne regu\u0142y czyszczenia")
                st.caption("Regu\u0142y s\u0105 stosowane automatycznie przy ka\u017cdym wczytaniu danych. Kolejno\u015b\u0107 ma znaczenie.")
                to_remove = None
                for i, entry in enumerate(st.session_state.cleaning_ops):
                    ops  = entry['ops']
                    cols = entry['cols']
                    # Build human-readable summary
                    op_labels = []
                    if ops.get('strip'):    op_labels.append("trim spacji")
                    if ops.get('dbl_sp'):   op_labels.append("podw\u00f3jne spacje")
                    if ops.get('tabs'):     op_labels.append("tabulatory")
                    if ops.get('newlines'): op_labels.append("nowe linie")
                    if ops.get('quotes'):   op_labels.append("cudzys\u0142owy")
                    case = ops.get('case', 'none')
                    if case != 'none': op_labels.append({'upper': 'WIELKIE', 'lower': 'ma\u0142e', 'title': 'Pierwsze Wielkie'}.get(case, case))
                    if ops.get('special'):  op_labels.append("znaki specjalne")
                    col_desc = f"{len(cols)} zmiennych" if len(cols) > 3 else ", ".join(cols)
                    row_c1, row_c2 = st.columns([5, 1])
                    row_c1.markdown(f"**{i+1}.** `{col_desc}` \u2192 {', '.join(op_labels)}")
                    if row_c2.button("\U0001f5d1\ufe0f", key=f"del_clean_{i}", help="Usu\u0144 t\u0119 regu\u0142\u0119"):
                        to_remove = i
                if to_remove is not None:
                    st.session_state.cleaning_ops.pop(to_remove)
                    st.rerun()
                if st.button("\U0001f5d1\ufe0f Usu\u0144 wszystkie regu\u0142y czyszczenia", type="secondary", key="clean_clear_all"):
                    st.session_state.cleaning_ops = []
                    st.rerun()

    # -- TYPY ZMIENNYCH (tylko Excel) ------------------
    if is_excel and tab_types is not None:
        with tab_types:
            st.markdown("#### Typy zmiennych (auto-detekcja + korekta)")
            st.info(
                "Program automatycznie wykrywa typy zmiennych. "
                "Je\u015bli kolumna ma \u226510 unikalnych warto\u015bci i same liczby \u2192 **numeryczna**; "
                "w przeciwnym razie \u2192 **kategoryczna**. "
                "Mo\u017cesz r\u0119cznie poprawi\u0107 dowolny typ poni\u017cej."
            )

            type_rows = []
            _type_label = {"numeric": "numeryczna", "categorical": "kategoryczna", "auto": "auto"}
            for col in df_orig_raw.columns:
                auto_type = "numeryczna" if col in numeric_cols_raw else "kategoryczna"
                override  = st.session_state.excel_col_types.get(col, "auto")
                effective = auto_type if override == "auto" else _type_label.get(override, override)
                type_rows.append({
                    "Zmienna": col,
                    "Etykieta": var_labels.get(col, col),
                    "Auto-detekcja": auto_type,
                    "Korekta u\u017cytkownika": _type_label.get(override, override),
                    "Efektywny typ": effective,
                    "Unikalnych warto\u015bci": df_orig_raw[col].nunique(),
                    "Brak\u00f3w [N]": df_orig_raw[col].isna().sum(),
                })
            st.dataframe(pd.DataFrame(type_rows), use_container_width=True, height=280)

            st.divider()
            st.markdown("**Zmie\u0144 typ wybranej zmiennej:**")
            col_t1, col_t2, col_t3 = st.columns([3, 2, 2])
            with col_t1:
                type_edit_var = st.selectbox(
                    "Zmienna:", df_orig_raw.columns,
                    format_func=lambda x: f"[{x}] {var_labels.get(x, x)}",
                    key="type_edit_var"
                )
            with col_t2:
                cur_override = st.session_state.excel_col_types.get(type_edit_var, "auto")
                new_type = st.selectbox(
                    "Typ:", ["auto", "numeric", "categorical"],
                    format_func=lambda x: {"auto": "auto", "numeric": "numeryczna", "categorical": "kategoryczna"}[x],
                    index=["auto", "numeric", "categorical"].index(cur_override) if cur_override in ["auto", "numeric", "categorical"] else 0,
                    key="type_edit_val"
                )
            with col_t3:
                st.write("")
                st.write("")
                if st.button("\U0001f4be Zastosuj", key="type_apply", use_container_width=True):
                    if new_type == "auto":
                        st.session_state.excel_col_types.pop(type_edit_var, None)
                    else:
                        st.session_state.excel_col_types[type_edit_var] = new_type
                    load_excel_data.clear()
                    st.success(f"Typ `{type_edit_var}` zmieniony na **{new_type}**. Strona zostanie od\u015bwie\u017cona.")
                    st.rerun()

            if st.session_state.excel_col_types:
                if st.button("\U0001f5d1\ufe0f Zresetuj wszystkie korekty typ\u00f3w", type="secondary", key="type_reset"):
                    st.session_state.excel_col_types = {}
                    load_excel_data.clear()
                    st.rerun()

            # \u2500\u2500 Summary of converted variables \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
            numeric_overrides = {c: v for c, v in st.session_state.excel_col_types.items()
                                 if v == "numeric"}
            if numeric_overrides:
                st.divider()
                st.markdown("**\u2705 Zmienne zmienione na numeryczne:**")
                for col_n in numeric_overrides:
                    val_lbl = st.session_state.custom_val_labels.get(col_n, {})
                    tag = " \u2014 zakodowana tekstowo" if val_lbl else " \u2014 by\u0142a ju\u017c liczbowa"
                    with st.expander(f"**{var_labels.get(col_n, col_n)}** (`{col_n}`){tag}", expanded=False):
                        if val_lbl:
                            st.markdown("**Etykiety warto\u015bci** (mo\u017cesz edytowa\u0107):")
                            edited_labels = {}
                            for code, lbl in sorted(val_lbl.items(), key=lambda x: x[0]):
                                new_lbl = st.text_input(
                                    f"Kod {code}:",
                                    value=str(lbl),
                                    key=f"lbl_edit_{col_n}_{code}"
                                )
                                edited_labels[code] = new_lbl
                            if st.button("\U0001f4be Zapisz etykiety", key=f"lbl_save_{col_n}",
                                         use_container_width=True):
                                st.session_state.custom_val_labels[col_n] = edited_labels
                                st.success("Etykiety zapisane.")
                                st.rerun()
                        else:
                            st.info("Warto\u015bci by\u0142y ju\u017c liczbowe \u2014 brak etykiet do edycji.")
                        if st.button("\u21a9\ufe0f Cofnij zmian\u0119 typu", key=f"type_revert_{col_n}",
                                     use_container_width=True):
                            st.session_state.excel_col_types.pop(col_n, None)
                            st.session_state.custom_val_labels.pop(col_n, None)
                            load_excel_data.clear()
                            st.rerun()

    # -- MRS ------------------------------------------
    with tab_mrs:
        col1, col2 = st.columns([1, 1])
        with col1:
            st.markdown("#### Dodaj zestaw wielokrotnego wyboru")

            # Numeric columns only
            numeric_mrs_cols = [c for c in df_raw.columns
                                if pd.api.types.is_numeric_dtype(df_raw[c])]

            if st.button("\U0001f50d Autowykrywanie (zmienne 0/1)", use_container_width=True):
                detected = auto_detect_mrs(df_raw)
                added = 0
                for name, cols in detected.items():
                    key = f"Auto_{name}"
                    if key not in st.session_state.mrs_sets:
                        # Store as dict with count_val
                        st.session_state.mrs_sets[key] = {'cols': cols, 'count_val': 1}
                        added += 1
                st.success(f"Wykryto i dodano {added} zestaw\u00f3w wielokrotnych odpowiedzi.")
                st.rerun()

            mrs_name = st.text_input("Nazwa nowego zestawu:", key="mrs_new_name")
            mrs_cols = st.multiselect(
                "Zmienne (numeryczne):",
                options=numeric_mrs_cols,
                format_func=lambda x: get_var_display_name(x, var_labels),
                key="mrs_new_cols"
            )

            # Show possible values from selected columns
            if mrs_cols:
                all_vals = sorted(set(
                    v for c in mrs_cols
                    for v in df_raw[c].dropna().unique()
                ))
                val_options = [int(v) if float(v) == int(v) else float(v) for v in all_vals]
            else:
                val_options = [1]

            mrs_count_val = st.selectbox(
                "Warto\u015b\u0107 do zliczania:",
                options=val_options if val_options else [1],
                index=val_options.index(1) if 1 in val_options else 0,
                key="mrs_new_val",
                help="Kt\u00f3ra warto\u015b\u0107 oznacza zaznaczenie odpowiedzi? Najcz\u0119\u015bciej: 1."
            )

            if st.button("\u2795 Dodaj zestaw", use_container_width=True):
                if mrs_name and mrs_cols:
                    st.session_state.mrs_sets[mrs_name] = {
                        'cols': mrs_cols,
                        'count_val': mrs_count_val
                    }
                    st.rerun()
                else:
                    st.warning("Podaj nazw\u0119 i wybierz co najmniej jedn\u0105 zmienn\u0105.")

        with col2:
            st.markdown("#### Zdefiniowane zestawy -- kliknij aby edytowa\u0107")
            if not st.session_state.mrs_sets:
                st.info("Brak zdefiniowanych zestaw\u00f3w wielokrotnych odpowiedzi.")
            for set_name, set_data in list(st.session_state.mrs_sets.items()):
                # Support both old format (list) and new format (dict)
                if isinstance(set_data, list):
                    set_cols = set_data
                    set_val  = 1
                else:
                    set_cols = set_data.get('cols', [])
                    set_val  = set_data.get('count_val', 1)

                with st.expander(f"**{set_name}** ({len(set_cols)} zmiennych, zliczana: {set_val})", expanded=False):
                    new_mrs_name = st.text_input("Nazwa zestawu:", value=set_name, key=f"mrs_rename_{set_name}")
                    st.caption("Wybierz zmienne numeryczne zestawu:")
                    edited_cols = st.multiselect(
                        "Zmienne zestawu:",
                        options=numeric_mrs_cols,
                        default=[c for c in set_cols if c in numeric_mrs_cols],
                        format_func=lambda x: get_var_display_name(x, var_labels),
                        key=f"mrs_edit_{set_name}"
                    )
                    if edited_cols:
                        edit_vals = sorted(set(
                            v for c in edited_cols
                            for v in df_raw[c].dropna().unique()
                        ))
                        edit_val_options = [int(v) if float(v) == int(v) else float(v) for v in edit_vals]
                    else:
                        edit_val_options = [1]

                    edited_val = st.selectbox(
                        "Warto\u015b\u0107 do zliczania:",
                        options=edit_val_options if edit_val_options else [1],
                        index=edit_val_options.index(set_val) if set_val in edit_val_options else 0,
                        key=f"mrs_val_{set_name}"
                    )

                    col_save, col_del = st.columns([3, 1])
                    with col_save:
                        if st.button("\U0001f4be Zapisz zmiany", key=f"save_mrs_{set_name}", use_container_width=True):
                            if not edited_cols:
                                st.error("Zestaw nie mo\u017ce by\u0107 pusty.")
                            else:
                                del st.session_state.mrs_sets[set_name]
                                final_name = new_mrs_name.strip() or set_name
                                st.session_state.mrs_sets[final_name] = {
                                    'cols': edited_cols,
                                    'count_val': edited_val
                                }
                                st.success(f"Zapisano '{final_name}'.")
                                st.rerun()
                    with col_del:
                        if st.button("\U0001f5d1\ufe0f Usu\u0144", key=f"del_mrs_{set_name}", use_container_width=True):
                            del st.session_state.mrs_sets[set_name]
                            st.rerun()

    # -- BATERIE MATRYCOWE -----------------------------
    with tab_matrix:
        st.markdown("#### Pytania matrycowe (pytania Likerta)")
        st.info("Grupuj subpytania z tej samej baterii (np. Q5_1, Q5_2, Q5_3) w jeden zestaw. B\u0119d\u0105 dost\u0119pne jako jedna tabela matrycowa w analizach i eksporcie.")

        col1, col2 = st.columns([1, 1])
        with col1:
            st.markdown("**Dodaj bateri\u0119**")
            if st.button("\U0001f50d Autowykrywanie (wsp\u00f3lny prefiks)", use_container_width=True, key="auto_matrix"):
                detected_m = auto_detect_matrix(df_raw)
                added_m = 0
                for name, cols in detected_m.items():
                    key = f"Auto_{name}"
                    if key not in st.session_state.matrix_sets:
                        st.session_state.matrix_sets[key] = cols
                        added_m += 1
                if added_m:
                    st.success(f"Wykryto i dodano {added_m} pyta\u0144 matrycowych.")
                else:
                    st.info("Nie znaleziono nowych baterii lub wszystkie ju\u017c istniej\u0105.")
                st.rerun()

            matrix_name = st.text_input("Nazwa baterii (np. 'Satysfakcja z produktu'):", key="matrix_new_name")
            matrix_cols_sel = st.multiselect(
                "Subpytania (zmienne numeryczne lub tekstowe):",
                options=[c for c in visible_columns if c not in hidden_cols],
                format_func=lambda x: get_var_display_name(x, var_labels),
                key="matrix_new_cols"
            )
            if st.button("\u2795 Dodaj bateri\u0119", use_container_width=True, key="add_matrix"):
                if matrix_name and len(matrix_cols_sel) >= 2:
                    st.session_state.matrix_sets[matrix_name] = matrix_cols_sel
                    st.success(f"Dodano bateri\u0119 '{matrix_name}'.")
                    st.rerun()
                else:
                    st.warning("Podaj nazw\u0119 i wybierz co najmniej 2 subpytania.")

        with col2:
            st.markdown("**Zdefiniowane pytania -- kliknij aby edytowa\u0107**")
            if not st.session_state.matrix_sets:
                st.info("Brak zdefiniowanych pyta\u0144 matrycowych.")
            to_del_mat = None
            for mat_name, mat_cols in list(st.session_state.matrix_sets.items()):
                with st.expander(f"**{mat_name}** ({len(mat_cols)} subpyta\u0144)", expanded=False):
                    new_mat_name = st.text_input("Nazwa pytania matrycowego:", value=mat_name, key=f"matrix_rename_{mat_name}")
                    edited_mat_cols = st.multiselect(
                        "Subpytania:",
                        options=[c for c in visible_columns if c not in hidden_cols],
                        default=[c for c in mat_cols if c in visible_columns],
                        format_func=lambda x: get_var_display_name(x, var_labels),
                        key=f"matrix_edit_{mat_name}"
                    )
                    for c in edited_mat_cols:
                        st.caption(f"  `{c}` -- {var_labels.get(c, '')}")
                    col_ms, col_md = st.columns([3, 1])
                    with col_ms:
                        if st.button("\U0001f4be Zapisz zmiany", key=f"save_mat_{mat_name}", use_container_width=True):
                            if len(edited_mat_cols) >= 2:
                                final_name = new_mat_name.strip() or mat_name
                                del st.session_state.matrix_sets[mat_name]
                                st.session_state.matrix_sets[final_name] = edited_mat_cols
                                st.success(f"Zapisano '{final_name}'.")
                                st.rerun()
                            else:
                                st.error("Pytanie musi mie\u0107 co najmniej 2 subpytania.")
                    with col_md:
                        if st.button("\U0001f5d1\ufe0f Usu\u0144", key=f"del_mat_{mat_name}", use_container_width=True):
                            del st.session_state.matrix_sets[mat_name]
                            st.rerun()

    # -- TOP/BOTTOM BOX --------------------------------
    with tab_box:
        col1, col2 = st.columns([1, 1])
        with col1:
            st.markdown("#### Dodaj grup\u0119 (Top/Bottom Box)")
            box_var = st.selectbox("Zmienna:", visible_columns, format_func=lambda x: get_var_display_name(x, var_labels), key="box_var_select")
            if box_var:
                box_name = st.text_input("Nazwa grupy (np. Top 2 Box):")
                box_cats = st.multiselect("Odpowiedzi w grupie:", df[box_var].dropna().unique())
                if st.button("\u2795 Dodaj grup\u0119", use_container_width=True):
                    if box_name and box_cats:
                        st.session_state.box_sets[box_var][f"[{box_name}]"] = box_cats
                        st.success(f"Dodano: {box_name}")
                        st.rerun()
        with col2:
            st.markdown("#### Zdefiniowane grupy")
            if not st.session_state.box_sets:
                st.info("Brak grup Box.")
            to_del_v, to_del_b = None, None
            for var, boxes in st.session_state.box_sets.items():
                st.markdown(f"**{var}**")
                for b_name, b_cats in boxes.items():
                    c1, c2 = st.columns([4, 1])
                    c1.write(f"- {b_name}: {', '.join(str(c) for c in b_cats)}")
                    if c2.button("\U0001f5d1\ufe0f", key=f"del_box_{var}_{b_name}"): to_del_v, to_del_b = var, b_name
            if to_del_v:
                del st.session_state.box_sets[to_del_v][to_del_b]
                if not st.session_state.box_sets[to_del_v]: del st.session_state.box_sets[to_del_v]
                st.rerun()

    # -- SEGMENTACJA -----------------------------------
    # -- WAZENIE ---------------------------------------
    with tab_weight:
        st.markdown("#### Wa\u017cenie RIM (iteracyjne dopasowanie proporcjonalne)")
        weight_vars = st.multiselect("Zmienne do wa\u017cenia:", visible_columns, format_func=lambda x: get_var_display_name(x, var_labels))
        if weight_vars:
            st.write("Wprowad\u017a docelowe odsetki (suma musi wynosi\u0107 100%):")
            targets, valid_targets = {}, True
            for w_var in weight_vars:
                st.markdown(f"**{get_var_display_name(w_var, var_labels)}**")
                categories = df[w_var].dropna().unique()
                targets[w_var] = {}
                cols = st.columns(min(len(categories), 4))
                sum_pct = 0
                for i, cat in enumerate(categories):
                    val = cols[i % 4].number_input(f"{cat}", 0.0, 100.0, 0.0, key=f"w_{w_var}_{cat}")
                    targets[w_var][cat] = val / 100.0
                    sum_pct += val
                if not np.isclose(sum_pct, 100.0, atol=0.1):
                    st.error(f"Suma = {sum_pct:.1f}%. Musi wynosi\u0107 100%!")
                    valid_targets = False
            if valid_targets and st.button("\u2696\ufe0f Oblicz wagi", type="primary"):
                st.session_state.weights = calculate_rim_weights(df, targets)
                st.session_state.weight_targets = targets
                st.success("\u2705 Wagi obliczone!")

# -------------------------------------------------------------
# MODU? 3: ANALIZY I TABELE
# -------------------------------------------------------------
elif menu == "\U0001f4c8 Analizy i Tabele":
    module_header("\U0001f4c8", "Analizy i Tabele")
    tab_freq, tab_matrix_an, tab_cross, tab_means, tab_desc, tab_corr = st.tabs([
        "\U0001f4c8 Cz\u0119sto\u015bci", "\U0001f522 Pytania Matrycowe", "\U0001f500 Tabele Krzy\u017cowe", "\U0001f4ca \u015arednie (T-Test)", "\U0001f522 Statystyki Opisowe", "\U0001f517 Korelacje"
    ])

    with tab_freq:
        freq_vars = st.multiselect("Wybierz zmienne:", all_options, format_func=lambda x: get_var_display_name(x, var_labels))
        show_charts_freq = st.checkbox("\U0001f4ca Wy\u015bwietlaj wykresy", key="charts_freq")
        if st.button("\u25b6\ufe0f Generuj tablice cz\u0119sto\u015bci", type="primary") and freq_vars:
            st.session_state.results['czestosci'] = {}
            w = st.session_state.weights if use_weights else np.ones(len(df_raw))
            for freq_var in freq_vars:
                if freq_var in st.session_state.matrix_sets:
                    # Matrix variable -- show as transposed freq table
                    mat_cols = st.session_state.matrix_sets[freq_var]
                    try:
                        mat_df, cats, sub_lbls = build_matrix_table(df, df_raw, mat_cols, var_labels, w, meta_orig.variable_value_labels, st.session_state.custom_val_labels)
                        st.session_state.results['czestosci'][freq_var] = mat_df
                        with st.expander(f"[Pytanie matrycowe] {freq_var}", expanded=True):
                            pct_cols_m = [f"{s} [%]" for s in sub_lbls]
                            n_cols_m   = [f"{s} [N]"  for s in sub_lbls]
                            st.dataframe(
                                mat_df.style
                                    .format(lambda x: f"{x:.0f}" if isinstance(x, (int, float)) else x, subset=n_cols_m)
                                    .format(lambda x: f"{x:.1f}%" if isinstance(x, (int, float)) else x, subset=pct_cols_m),
                                use_container_width=True
                            )
                    except Exception as e:
                        st.error(f"B\u0142\u0105d dla baterii {freq_var}: {e}")
                    continue

                if freq_var in st.session_state.mrs_sets:
                    set_data = st.session_state.mrs_sets[freq_var]
                    # Support both old format (list) and new format (dict)
                    if isinstance(set_data, list):
                        cols = set_data
                        count_val = 1
                    else:
                        cols = set_data.get('cols', [])
                        count_val = set_data.get('count_val', 1)
                    missing_mask = df_raw[cols].isna().all(axis=1)
                    missing_count = w[missing_mask].sum()
                    counts = pd.Series({var_labels.get(c, c): w[(df_raw[c] == count_val)].sum() for c in cols})
                    total_respondents = w[~missing_mask].sum()
                    pcts = (counts / total_respondents) * 100 if total_respondents > 0 else counts * 0
                    res_df = pd.DataFrame({'Liczebnosc [N]': counts, 'Procent [%]': pcts})
                    res_df.loc['Og\u00f3\u0142em (Wa\u017cne)'] = [total_respondents, pcts.sum()]
                    res_df.loc['Braki danych'] = [missing_count, np.nan]
                else:
                    missing_mask = df[freq_var].isna()
                    missing_count = w[missing_mask].sum()
                    valid_df = pd.DataFrame({'val': df[freq_var], 'w': w}).dropna()
                    counts = valid_df.groupby('val', observed=False)['w'].sum()
                    if freq_var in st.session_state.box_sets:
                        for box_name, b_cats in st.session_state.box_sets[freq_var].items():
                            box_val = counts[counts.index.isin(b_cats)].sum()
                            counts.loc[box_name] = box_val
                    pcts = (counts / valid_df['w'].sum()) * 100 if valid_df['w'].sum() > 0 else counts * 0
                    res_df = pd.DataFrame({'Liczebnosc [N]': counts, 'Procent [%]': pcts})
                    sum_n = counts[~counts.index.astype(str).str.startswith('[')].sum()
                    res_df.loc['Suma'] = [sum_n, 100.0 if sum_n > 0 else 0]
                    res_df.loc['Braki danych'] = [missing_count, np.nan]
                st.session_state.results['czestosci'][freq_var] = res_df
                with st.expander(get_var_display_name(freq_var, var_labels), expanded=True):
                    st.dataframe(res_df.style.format(get_streamlit_format(res_df)), use_container_width=True)
                    if show_charts_freq:
                        plot_df = res_df.drop(index=['Suma', 'Og\u00f3\u0142em (Wa\u017cne)', 'Braki danych'], errors='ignore')
                        plot_df = plot_df[~plot_df.index.astype(str).str.startswith('[')]
                        if not plot_df.empty and 'Procent [%]' in plot_df.columns:
                            plot_df = plot_df.dropna(subset=['Procent [%]'])
                        if not plot_df.empty and 'Procent [%]' in plot_df.columns:
                            fig = px.bar(plot_df, x='Procent [%]', y=plot_df.index, orientation='h',
                                         title=var_labels.get(freq_var, freq_var), color_discrete_sequence=['#2E75B6'])
                            fig.update_layout(yaxis={'categoryorder': 'total ascending'}, height=max(300, len(plot_df) * 30 + 80))
                            st.plotly_chart(fig, use_container_width=True, key=f"pc_freq_gen_{freq_var}")
            st.success("\u2705 Tablice cz\u0119sto\u015bci wygenerowane!")

        # \u2500\u2500 Persistent display of stored results \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
        if st.session_state.results.get('czestosci'):
            st.divider()
            st.markdown(f"**Zapisane wyniki ({len(st.session_state.results['czestosci'])} tablic):**")
            for freq_var, res_df in st.session_state.results['czestosci'].items():
                with st.expander(get_var_display_name(freq_var, var_labels), expanded=False):
                    st.dataframe(res_df.style.format(get_streamlit_format(res_df)),
                                 use_container_width=True)
                    if show_charts_freq:
                        plot_df = res_df.drop(
                            index=['Suma','Og\u00f3\u0142em (Wa\u017cne)','Braki danych'], errors='ignore')
                        plot_df = plot_df[~plot_df.index.astype(str).str.startswith('[')]
                        if not plot_df.empty and 'Procent [%]' in plot_df.columns:
                            fig = px.bar(plot_df, x='Procent [%]', y=plot_df.index,
                                         orientation='h',
                                         title=var_labels.get(freq_var, freq_var),
                                         color_discrete_sequence=['#2E75B6'])
                            fig.update_layout(yaxis={'categoryorder':'total ascending'}, height=350)
                            st.plotly_chart(fig, use_container_width=True, key=f"pc_freq_saved_{freq_var}")

    # -- BATERIE MATRYCOWE -- dedykowana zak?adka analityczna --
    with tab_matrix_an:
        st.markdown("#### Tabele matrycowe (pytania matrycowe (Likert))")
        st.info("Wiersze = warto\u015bci skali, Kolumny = subpytania. Dost\u0119pne tryby: tylko N, tylko %, lub N + %.")

        if not st.session_state.matrix_sets:
            st.warning("Brak zdefiniowanych pyta\u0144 matrycowych. Przejd\u017a do **Przygotowanie Danych \u2192 Pytania Matrycowe** i dodaj baterie.")
        else:
            mat_sel = st.multiselect(
                "Wybierz pytania matrycowe:",
                list(st.session_state.matrix_sets.keys()),
                default=list(st.session_state.matrix_sets.keys()),
                key="matrix_an_sel"
            )

            col_mode, col_chart = st.columns([2, 1])
            with col_mode:
                mat_display_mode = st.radio(
                    "Prezentuj warto\u015bci:",
                    ["N + %", "Tylko N", "Tylko %"],
                    index=0, horizontal=True, key="mat_display_mode"
                )
            with col_chart:
                show_chart_mat = st.checkbox("\U0001f4ca Wy\u015bwietl wykres", key="chart_mat")

            if st.button("\u25b6\ufe0f Generuj tabele matrycowe", type="primary", key="gen_matrix"):
                w = st.session_state.weights if use_weights else np.ones(len(df_raw))
                st.session_state.matrix_results = []
                for mat_name in mat_sel:
                    mat_cols = st.session_state.matrix_sets[mat_name]
                    try:
                        mat_df, cats, sub_lbls = build_matrix_table(df, df_raw, mat_cols, var_labels, w, meta_orig.variable_value_labels, st.session_state.custom_val_labels)
                        st.session_state.matrix_results.append({
                            'name': mat_name, 'df': mat_df, 'cats': cats,
                            'sub_labels': sub_lbls, 'cols': mat_cols,
                            'display_mode': mat_display_mode,
                        })
                    except Exception as e:
                        st.error(f"B\u0142\u0105d dla '{mat_name}': {e}")
                st.success(f"\u2705 Wygenerowano {len(st.session_state.matrix_results)} tabel matrycowych.")

            for entry in st.session_state.matrix_results:
                with st.expander(f"\U0001f522 {entry['name']}", expanded=True):
                    mat_df   = entry['df']
                    cats     = entry['cats']
                    sub_lbls = entry.get('sub_labels', [])

                    pct_cols_m = [f"{s} [%]" for s in sub_lbls]
                    n_cols_m   = [f"{s} [N]"  for s in sub_lbls]

                    # Build view according to display mode
                    SUMROW = "Baza (N) / Suma (%)"
                    if mat_display_mode == "Tylko N":
                        view_cols = n_cols_m
                        view_df   = mat_df[view_cols].copy()
                    elif mat_display_mode == "Tylko %":
                        view_cols = pct_cols_m
                        view_df   = mat_df[view_cols].copy()
                    else:  # N + %
                        view_cols = []
                        for s in sub_lbls:
                            view_cols += [f"{s} [N]", f"{s} [%]"]
                        view_df = mat_df[view_cols].copy()

                    _style_matrix_row = _make_style_matrix_row(SUMROW)


                    styled = view_df.style.apply(_style_matrix_row, axis=1).format(_fmt_cell)
                    st.dataframe(styled, use_container_width=True)

                    if show_chart_mat and cats and sub_lbls:
                        chart_rows = [c for c in cats if c in mat_df.index]
                        chart_data = pd.DataFrame(index=chart_rows, columns=sub_lbls, dtype=float)
                        for sub_l in sub_lbls:
                            col_key = f"{sub_l} [%]"
                            if col_key in mat_df.columns:
                                for cat_v in chart_rows:
                                    try:
                                        chart_data.loc[cat_v, sub_l] = float(mat_df.loc[cat_v, col_key])
                                    except Exception:
                                        chart_data.loc[cat_v, sub_l] = 0
                        chart_data = chart_data.reset_index().rename(columns={'index': 'Wartosc'})
                        chart_long = chart_data.melt(id_vars='Wartosc', var_name='Subpytanie', value_name='%')
                        fig_mat = px.bar(
                            chart_long, x='%', y='Wartosc', color='Subpytanie',
                            barmode='group', orientation='h', title=entry['name'],
                            color_discrete_sequence=px.colors.qualitative.Set2
                        )
                        fig_mat.update_layout(yaxis={'categoryorder': 'category ascending'},
                                              height=max(300, len(cats) * 40 + 80))
                        st.plotly_chart(fig_mat, use_container_width=True, key=f"pc_mat_{entry.get('name','m')}")

    with tab_cross:
        col1, col2 = st.columns(2)
        with col1: row_vars = st.multiselect("Zmienne w wierszach:", all_options, format_func=lambda x: get_var_display_name(x, var_labels))
        with col2: col_vars = st.multiselect("Zmienne w kolumnach:", all_options, format_func=lambda x: get_var_display_name(x, var_labels))
        pct_type = st.radio("Spos\u00f3b prezentacji:", ["Liczebno\u015bci", "Kolumnowe (%)", "Wierszowe (%)", "Liczebno\u015bci + Kolumnowe (%)", "Liczebno\u015bci + Wierszowe (%)"], horizontal=True)
        c1, c2, c3, c4 = st.columns(4)
        do_sig_test = c1.checkbox("\U0001f520 Testy Z (95%)")
        do_chi_square = c2.checkbox("\U0001f9ee Chi-kwadrat")
        do_cramer = c3.checkbox("\U0001f4cf V Kramera", help="Si\u0142a zwi\u0105zku: 0=brak, 0.1=s\u0142aby, 0.3=umiarkowany, 0.5+=silny")
        show_charts_cross = c4.checkbox("\U0001f4ca Wykresy")

        if st.button("\u25b6\ufe0f Generuj tabele krzy\u017cowe", type="primary") and row_vars and col_vars:
            st.session_state.results['krzyzowe'] = {}
            st.session_state.chi_results = {}
            w = st.session_state.weights if use_weights else np.ones(len(df_raw))
            for row_var in row_vars:
                for col_var in col_vars:
                    is_row_mrs = row_var in st.session_state.mrs_sets
                    is_col_mrs = col_var in st.session_state.mrs_sets
                    try:
                        tmp_df = pd.DataFrame({'w': w})
                        if is_row_mrs: tmp_df['R_miss'] = df_raw[st.session_state.mrs_sets[row_var]].isna().all(axis=1)
                        else: tmp_df['R_miss'] = df[row_var].isna()
                        if is_col_mrs: tmp_df['C_miss'] = df_raw[st.session_state.mrs_sets[col_var]].isna().all(axis=1)
                        else: tmp_df['C_miss'] = df[col_var].isna()
                        missing_count = tmp_df.loc[tmp_df['R_miss'] | tmp_df['C_miss'], 'w'].sum()

                        if not is_row_mrs and not is_col_mrs:
                            df_n = pd.crosstab(df[row_var], df[col_var], values=w, aggfunc='sum', dropna=False).fillna(0)
                        elif is_row_mrs and not is_col_mrs:
                            cols = st.session_state.mrs_sets[row_var]
                            mrs_w = df_raw[cols].replace(np.nan, 0).multiply(w, axis=0)
                            df_n = mrs_w.groupby(df[col_var], observed=False).sum().T
                            df_n.index = [var_labels.get(c, c) for c in df_n.index]
                        elif not is_row_mrs and is_col_mrs:
                            cols = st.session_state.mrs_sets[col_var]
                            mrs_w = df_raw[cols].replace(np.nan, 0).multiply(w, axis=0)
                            df_n = mrs_w.groupby(df[row_var], observed=False).sum()
                            df_n.columns = [var_labels.get(c, c) for c in df_n.columns]

                        if row_var in st.session_state.box_sets and not is_row_mrs:
                            for box_name, b_cats in st.session_state.box_sets[row_var].items():
                                df_n.loc[box_name] = df_n.loc[df_n.index.intersection(b_cats)].sum(axis=0)

                        if (do_chi_square or do_cramer) and not is_row_mrs and not is_col_mrs:
                            obs = df_n.loc[~df_n.index.astype(str).str.startswith('[')]
                            obs = obs[[c for c in obs.columns if not str(c).startswith('[') and c != 'Suma']]
                            if obs.shape[0] > 1 and obs.shape[1] > 1 and obs.sum().sum() > 0:
                                chi2, p, dof, ex = stats.chi2_contingency(obs)
                                n_total = obs.sum().sum()
                                k = min(obs.shape[0], obs.shape[1])
                                cramer_v = float(np.sqrt(chi2 / (n_total * (k - 1)))) if n_total > 0 and k > 1 else 0.0

                                # Interpret Cramer's V
                                if cramer_v < 0.1:
                                    v_interp = "brak/zaniedbywalny"
                                elif cramer_v < 0.3:
                                    v_interp = "s\u0142aby"
                                elif cramer_v < 0.5:
                                    v_interp = "umiarkowany"
                                else:
                                    v_interp = "silny"

                                parts = []
                                if do_chi_square:
                                    parts.append(f"Chi\u00b2={chi2:.2f}, df={dof}, p={p:.3f}")
                                if do_cramer:
                                    parts.append(f"V Kramera={cramer_v:.3f} ({v_interp})")
                                st.session_state.chi_results[f"{row_var} x {col_var}"] = " | ".join(parts)

                        if not is_row_mrs and not is_col_mrs:
                            df_n['Suma'] = df_n.loc[~df_n.index.astype(str).str.startswith('[')].sum(axis=1)
                            df_n.loc['Suma'] = df_n.loc[~df_n.index.astype(str).str.startswith('[')].sum(axis=0)
                        elif is_row_mrs and not is_col_mrs:
                            df_n['Suma'] = mrs_w.sum().values
                            df_n.loc['Suma'] = df_n.sum(axis=0)
                        elif not is_row_mrs and is_col_mrs:
                            df_n['Suma'] = df_n.loc[~df_n.index.astype(str).str.startswith('[')].sum(axis=1)
                            df_n.loc['Suma'] = mrs_w.sum().values.tolist() + [mrs_w.sum().sum()]

                        df_pct = pd.DataFrame(np.nan, index=df_n.index, columns=df_n.columns)
                        if "Kolumnowe" in pct_type:
                            if not is_row_mrs and not is_col_mrs: df_pct = df_n.div(df_n.loc['Suma'].replace(0, np.nan), axis=1) * 100
                            elif is_row_mrs and not is_col_mrs:
                                base = tmp_df.loc[~tmp_df['C_miss']].groupby(df[col_var], observed=False)['w'].sum()
                                base['Suma'] = base.sum()
                                df_pct = df_n.div(base.replace(0, np.nan), axis=1) * 100
                            elif not is_row_mrs and is_col_mrs:
                                df_pct = df_n.div(df_n.loc['Suma'].replace(0, np.nan), axis=1) * 100
                        elif "Wierszowe" in pct_type:
                            if not is_row_mrs and not is_col_mrs: df_pct = df_n.div(df_n['Suma'].replace(0, np.nan), axis=0) * 100
                            elif is_row_mrs and not is_col_mrs: df_pct = df_n.div(df_n['Suma'].replace(0, np.nan), axis=0) * 100
                            elif not is_row_mrs and is_col_mrs:
                                base = tmp_df.loc[~tmp_df['R_miss']].groupby(df[row_var], observed=False)['w'].sum()
                                for box_name, b_cats in st.session_state.box_sets.get(row_var, {}).items():
                                    base[box_name] = base[base.index.intersection(b_cats)].sum()
                                base['Suma'] = base.loc[~base.index.astype(str).str.startswith('[')].sum()
                                df_pct = df_n.div(base.replace(0, np.nan), axis=0) * 100

                        if "Kolumnowe" in pct_type or "Wierszowe" in pct_type:
                            df_pct = df_pct.fillna(0)

                        if do_sig_test and "Kolumnowe" in pct_type:
                            sig_df, col_letters = apply_sig_testing(df_pct, df_n)
                            rename_dict = {c: f"{c} [{col_letters.get(c, '')}]" for c in df_n.columns if c != 'Suma'}
                            df_n.rename(columns=rename_dict, inplace=True)
                            df_pct.rename(columns=rename_dict, inplace=True)
                            sig_df.rename(columns=rename_dict, inplace=True)
                            df_pct_str = pd.DataFrame("", index=df_pct.index, columns=df_pct.columns)
                            for c in df_pct.columns:
                                df_pct_str[c] = df_pct[c].apply(lambda x: f"{x:.1f}%" if pd.notna(x) else "")
                            df_pct = df_pct_str + sig_df

                        if pct_type == "Liczebno\u015bci":
                            cross_df = df_n.add_suffix(' [N]')
                        elif pct_type in ["Kolumnowe (%)", "Wierszowe (%)"]:
                            cross_df = df_pct.add_suffix(' [%]')
                        else:
                            p_lbl = "[% Kolumnowe]" if "Kolumnowe" in pct_type else "[% Wierszowe]"
                            cross_cols = []
                            for c in df_n.columns:
                                cross_cols.extend([f"{c} [N]", f"{c} {p_lbl}"])
                            cross_df = pd.DataFrame(index=df_n.index, columns=cross_cols)
                            for c in df_n.columns:
                                cross_df[f"{c} [N]"] = df_n[c]
                                cross_df[f"{c} {p_lbl}"] = df_pct[c]

                        cross_df.loc['Braki danych (wykluczone z tabeli)'] = [missing_count] + [np.nan] * (len(cross_df.columns) - 1)
                        title = f"{row_var} x {col_var}"
                        st.session_state.results['krzyzowe'][title] = cross_df
                        with st.expander(title):
                            st.dataframe(safe_style(cross_df), use_container_width=True)
                            if title in st.session_state.chi_results:
                                st.caption(f"\U0001f9ee {st.session_state.chi_results[title]}")
                            if show_charts_cross:
                                plot_df = cross_df.drop(index=['Suma', 'Braki danych', 'Braki danych (wykluczone z tabeli)'], errors='ignore')
                                plot_df = plot_df[~plot_df.index.astype(str).str.startswith('[')]
                                if "Kolumnowe" in pct_type or "Wierszowe" in pct_type:
                                    p_cols = [c for c in plot_df.columns if '[%]' in c or '[% Kolumnowe]' in c or '[% Wierszowe]' in c]
                                else:
                                    p_cols = [c for c in plot_df.columns if '[N]' in c]
                                p_cols = [c for c in p_cols if 'Suma' not in c]
                                if p_cols and not plot_df.empty:
                                    temp_plot = plot_df[p_cols].copy()
                                    for col in temp_plot.columns:
                                        temp_plot[col] = pd.to_numeric(temp_plot[col].apply(lambda x: str(x).split('%')[0].strip() if isinstance(x, str) else x), errors='coerce')
                                    temp_plot.columns = [c.split(' [')[0] for c in temp_plot.columns]
                                    fig = px.bar(temp_plot, barmode='group', orientation='h',
                                                 title=f"{var_labels.get(row_var, row_var)} wg {var_labels.get(col_var, col_var)}")
                                    fig.update_layout(yaxis={'categoryorder': 'total ascending'}, height=350)
                                    st.plotly_chart(fig, use_container_width=True, key=f"pc_cross_{title}")
                    except Exception as e:
                        st.error(f"B\u0142\u0105d dla {row_var} \u00d7 {col_var}: {e}")
            st.success("\u2705 Tabele krzy\u017cowe wygenerowane!")

        # \u2500\u2500 Persistent display of stored cross-tab results \u2500\u2500\u2500\u2500
        if st.session_state.results.get('krzyzowe'):
            st.divider()
            st.markdown(f"**Zapisane wyniki ({len(st.session_state.results['krzyzowe'])} tabel):**")
            for title, cross_df in st.session_state.results['krzyzowe'].items():
                with st.expander(title, expanded=False):
                    st.dataframe(safe_style(cross_df), use_container_width=True)
                    if title in st.session_state.chi_results:
                        st.caption(f"\U0001f9ee {st.session_state.chi_results[title]}")

    with tab_means:
        col1, col2 = st.columns(2)
        with col1: mean_rows = st.multiselect("Zmienne ci\u0105g\u0142e (wiersze):", numeric_cols, format_func=lambda x: get_var_display_name(x, var_labels))
        with col2: mean_cols_sel = st.multiselect("Metryczka (kolumny):", all_options, format_func=lambda x: get_var_display_name(x, var_labels))
        do_means_sig = st.checkbox("\U0001f520 Oznacz istotne r\u00f3\u017cnice \u015brednich (T-Test 95%)")
        if st.button("\u25b6\ufe0f Generuj tabele \u015brednich", type="primary") and mean_rows and mean_cols_sel:
            st.session_state.results['srednie'] = {}
            w = st.session_state.weights if use_weights else np.ones(len(df_raw))
            for row_var in mean_rows:
                for col_var in mean_cols_sel:
                    try:
                        x = pd.to_numeric(df_raw[row_var], errors='coerce')
                        c_series = df[col_var]
                        cats = c_series.dropna().unique()
                        df_means = pd.DataFrame(index=['Srednia', 'Odchylenie Std.', 'Baza (N)'], columns=cats)
                        df_vars = pd.DataFrame(index=['Srednia'], columns=cats)
                        df_ess = pd.DataFrame(index=['Srednia'], columns=cats)
                        for cat in cats:
                            mask = (c_series == cat)
                            mean, var, ess = get_weighted_stats(x[mask].values, w[mask])
                            std = np.sqrt(var) if pd.notna(var) and var >= 0 else np.nan
                            df_means.loc['Srednia', cat] = mean
                            df_means.loc['Odchylenie Std.', cat] = std
                            df_means.loc['Baza (N)', cat] = w[mask & ~np.isnan(x)].sum()
                            df_vars.loc['Srednia', cat] = var
                            df_ess.loc['Srednia', cat] = ess
                        mean, var, ess = get_weighted_stats(x.values, w)
                        df_means['Og\u00f3\u0142em'] = [mean, np.sqrt(var) if pd.notna(var) and var >= 0 else np.nan, w[~np.isnan(x)].sum()]
                        df_vars['Og\u00f3\u0142em'] = var
                        df_ess['Og\u00f3\u0142em'] = ess
                        if do_means_sig:
                            sig_df, col_letters = apply_means_sig_testing(df_means.loc[['Srednia']], df_vars.loc[['Srednia']], df_ess.loc[['Srednia']])
                            rename_dict = {c: f"{c} [{col_letters.get(c, '')}]" for c in cats}
                            df_means.rename(columns=rename_dict, inplace=True)
                            sig_df.rename(columns=rename_dict, inplace=True)
                            df_str = pd.DataFrame("", index=df_means.index, columns=df_means.columns)
                            for c in df_means.columns:
                                df_str.loc['Srednia', c] = f"{df_means.loc['Srednia', c]:.2f}" if pd.notna(df_means.loc['Srednia', c]) else ""
                                df_str.loc['Odchylenie Std.', c] = df_means.loc['Odchylenie Std.', c]
                                df_str.loc['Baza (N)', c] = df_means.loc['Baza (N)', c]
                            df_str.loc['Srednia'] = df_str.loc['Srednia'] + sig_df.loc['Srednia']
                            df_means = df_str
                        title = f"{row_var} x {col_var}"
                        st.session_state.results['srednie'][title] = df_means
                        with st.expander(title):
                            st.dataframe(df_means.style.format(get_streamlit_format(df_means)), use_container_width=True)
                    except Exception as e:
                        st.error(f"B\u0142\u0105d: {e}")
            st.success("\u2705 Tabele \u015brednich wygenerowane!")

    with tab_desc:
        desc_vars = st.multiselect("Zmienne numeryczne:", numeric_cols,
                                    format_func=lambda x: get_var_display_name(x, var_labels))

        st.markdown("**Wybierz statystyki do prezentacji:**")

        # Group checkboxes into 3 columns for a clean layout
        d1, d2, d3 = st.columns(3)
        with d1:
            st.markdown("*Tendencja centralna*")
            d_mean     = st.checkbox("Srednia",              value=True,  key="ds_mean")
            d_median   = st.checkbox("Mediana (Q2 / 50%)",   value=True,  key="ds_median")
            d_mode     = st.checkbox("Dominanta (moda)",     value=False, key="ds_mode")
            d_trimmed  = st.checkbox("Srednia obci\u0119ta (5%)", value=False, key="ds_trimmed",
                                      help="Srednia po odci\u0119ciu 5% obserwacji z ka\u017cdego ko\u0144ca")
        with d2:
            st.markdown("*Rozrzut*")
            d_std      = st.checkbox("Odchylenie std.",      value=True,  key="ds_std")
            d_var      = st.checkbox("Wariancja",            value=False, key="ds_var")
            d_se       = st.checkbox("B\u0142\u0105d std. sredniej (SE)", value=False, key="ds_se")
            d_range    = st.checkbox("Rozst\u0119p (max-min)", value=False, key="ds_range")
            d_iqr      = st.checkbox("IQR (Q3-Q1)",          value=False, key="ds_iqr")
            d_cv       = st.checkbox("Wsp. zmienno\u015bci (%)", value=False, key="ds_cv",
                                      help="CV = (Odch. std / Srednia) * 100%")
        with d3:
            st.markdown("*Kszta\u0142t rozk\u0142adu*")
            d_skew     = st.checkbox("Sko\u015bno\u015b\u0107",  value=False, key="ds_skew")
            d_kurt     = st.checkbox("Kurtoza",               value=False, key="ds_kurt")
            d_min      = st.checkbox("Min",                   value=True,  key="ds_min")
            d_max      = st.checkbox("Max",                   value=True,  key="ds_max")
            d_q1       = st.checkbox("Q1 (25. percentyl)",    value=False, key="ds_q1")
            d_q3       = st.checkbox("Q3 (75. percentyl)",    value=False, key="ds_q3")
            st.markdown("*Obserwacje*")
            d_n_valid  = st.checkbox("N wa\u017cnych",         value=True,  key="ds_nvalid")
            d_n_miss   = st.checkbox("N brak\u00f3w",          value=True,  key="ds_nmiss")

        if st.button("\u25b6\ufe0f Generuj statystyki opisowe", type="primary") and desc_vars:
            try:
                _w_desc = (st.session_state.weights if use_weights and st.session_state.weights is not None
                           else None)

                def _wstat(col):
                    """Return (values, weights) arrays after dropping NaN."""
                    mask = df_raw[col].notna()
                    x = df_raw.loc[mask, col].values.astype(float)
                    if _w_desc is not None:
                        w = pd.Series(_w_desc, index=df_raw.index).loc[mask].values.clip(min=0)
                    else:
                        w = np.ones(len(x))
                    return x, w

                def _wmean(x, w):
                    return float((x * w).sum() / w.sum()) if w.sum() > 0 else np.nan

                def _wvar(x, w):
                    m = _wmean(x, w)
                    n_eff = w.sum()
                    return float((w * (x - m) ** 2).sum() / max(n_eff - 1, 1))

                def _wquantile(x, w, q):
                    """Weighted quantile (linear interpolation, consistent with SPSS)."""
                    idx = np.argsort(x)
                    xs, ws = x[idx], w[idx]
                    cumw = np.cumsum(ws)
                    total = cumw[-1]
                    target = q * total
                    # Find where cumulative weight crosses target
                    i = np.searchsorted(cumw, target)
                    if i == 0:
                        return float(xs[0])
                    if i >= len(xs):
                        return float(xs[-1])
                    # Linear interpolation
                    frac = (target - cumw[i - 1]) / max(ws[i], 1e-12)
                    return float(xs[i - 1] + frac * (xs[i] - xs[i - 1]))

                rows = []
                for c in desc_vars:
                    x, w = _wstat(c)
                    row = {'Zmienna': c, 'Etykieta': var_labels.get(c, c)}

                    n_valid = int(df_raw[c].notna().sum())
                    n_miss  = int(df_raw[c].isna().sum())
                    n_w     = float(w.sum())   # effective weighted N

                    if d_n_valid:  row['N wa\u017cnych (wa\u017cone)'] = round(n_w, 2)
                    if d_n_miss:   row['N brak\u00f3w']               = n_miss

                    if d_mean:
                        row['Srednia'] = _wmean(x, w)
                    if d_trimmed:
                        # Trimmed mean: use unweighted (SPSS behaviour)
                        row['Srednia obci\u0119ta (5%)'] = float(stats.trim_mean(x, 0.05))
                    if d_median:
                        row['Mediana'] = _wquantile(x, w, 0.5)
                    if d_mode:
                        mode_res = stats.mode(x, keepdims=True)
                        row['Dominanta'] = float(mode_res.mode[0]) if len(mode_res.mode) > 0 else np.nan
                    if d_std:
                        row['Odch. std.'] = float(np.sqrt(_wvar(x, w)))
                    if d_var:
                        row['Wariancja'] = _wvar(x, w)
                    if d_se:
                        row['B\u0142\u0105d std. (SE)'] = float(np.sqrt(_wvar(x, w) / max(n_w, 1)))
                    if d_min:  row['Min'] = float(x.min())
                    if d_max:  row['Max'] = float(x.max())
                    if d_range: row['Rozst\u0119p'] = float(x.max() - x.min())
                    if d_q1:   row['Q1 (25%)'] = _wquantile(x, w, 0.25)
                    if d_q3:   row['Q3 (75%)'] = _wquantile(x, w, 0.75)
                    if d_iqr:
                        row['IQR'] = _wquantile(x, w, 0.75) - _wquantile(x, w, 0.25)
                    if d_cv:
                        mn = _wmean(x, w)
                        std = np.sqrt(_wvar(x, w))
                        row['CV (%)'] = float(std / mn * 100) if mn != 0 else np.nan
                    if d_skew:
                        # Weighted skewness (SPSS formula)
                        mn  = _wmean(x, w)
                        std = np.sqrt(_wvar(x, w))
                        if std > 0:
                            row['Sko\u015bno\u015b\u0107'] = float(
                                (w * ((x - mn) / std) ** 3).sum() / w.sum())
                        else:
                            row['Sko\u015bno\u015b\u0107'] = np.nan
                    if d_kurt:
                        mn  = _wmean(x, w)
                        std = np.sqrt(_wvar(x, w))
                        if std > 0:
                            row['Kurtoza'] = float(
                                (w * ((x - mn) / std) ** 4).sum() / w.sum() - 3)
                        else:
                            row['Kurtoza'] = np.nan

                    rows.append(row)

                desc_df = pd.DataFrame(rows).set_index('Zmienna')
                st.session_state.results['opisowe']['Statystyki opisowe'] = desc_df
                num_cols_desc = [c for c in desc_df.columns if desc_df[c].dtype in [float, np.float64]]
                st.dataframe(
                    desc_df.style.format({c: '{:.3f}' for c in num_cols_desc}),
                    use_container_width=True
                )
            except Exception as e:
                st.error(str(e))

    with tab_corr:
        corr_vars = st.multiselect("Zmienne do macierzy korelacji:", numeric_cols,
                                    format_func=lambda x: get_var_display_name(x, var_labels))

        c_opt1, c_opt2, c_opt3 = st.columns(3)
        with c_opt1:
            corr_method = st.selectbox("Metoda:", ["pearson", "spearman", "kendall"],
                                        key="corr_method",
                                        help="Pearson: liniowa; Spearman: rang (odporna); Kendall: rang (ma\u0142e pr\u00f3by)")
        with c_opt2:
            show_heatmap = st.checkbox("\U0001f321\ufe0f Mapa ciep\u0142a", key="corr_heatmap")
        with c_opt3:
            corr_threshold = st.slider(
                "Prog silnej korelacji (|r|):", 0.0, 1.0, 0.5, 0.05,
                key="corr_thresh",
                help="Pary o warto\u015bci bezwzgl\u0119dnej korelacji \u2265 progu zostan\u0105 wyro\u017cnione kolorem i wy\u015bwietlone jako lista."
            )

        if st.button("\u25b6\ufe0f Oblicz korelacje", type="primary") and len(corr_vars) > 1:
            try:
                _w_corr = st.session_state.weights if use_weights else None
                corr_df, n_obs = calculate_correlations(df_raw, corr_vars,
                                                         weights=_w_corr,
                                                         method=corr_method)

                # Raw numeric matrix (for styling and heatmap)
                if corr_method == 'pearson' and _w_corr is not None:
                    # Reconstruct numeric matrix from the starred corr_df
                    import re as _re
                    num_corr = corr_df.map(
                        lambda x: float(_re.sub(r'[*\n].*', '', str(x))) if str(x) not in ('1.000', 'N/A') else (1.0 if str(x) == '1.000' else np.nan)
                    )
                else:
                    num_corr = df_raw[corr_vars].corr(method=corr_method)

                corr_df.index   = [var_labels.get(c, c) for c in corr_df.index]
                corr_df.columns = [var_labels.get(c, c) for c in corr_df.columns]
                num_corr.index   = [var_labels.get(c, c) for c in num_corr.index]
                num_corr.columns = [var_labels.get(c, c) for c in num_corr.columns]

                st.session_state.results['korelacje']['Macierz Korelacji'] = corr_df
                st.write(f"**Metoda:** {corr_method.title()} | **N wa\u017cnych obserwacji:** {n_obs}")

                # \u2500\u2500 Color-coded table \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
                _color_corr_cell = _make_color_corr_cell(corr_threshold)

                styled = corr_df.style.map(_color_corr_cell).format(
                    lambda x: x if isinstance(x, str) else f'{x:.3f}'
                )
                st.dataframe(styled, use_container_width=True)

                # \u2500\u2500 Legend \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
                st.markdown(
                    "\U0001f7e9 **Silna dodatnia** (r \u2265 0.70) &nbsp;&nbsp;"
                    "\U0001f7e5 **Silna ujemna** (r \u2264 \u22120.70) &nbsp;&nbsp;"
                    "\U0001f7e8 **Umiarkowana dodatnia / ujemna** (|r| \u2265 prog) &nbsp;&nbsp;"
                    "\u25fb Poni\u017cej progu"
                )

                # \u2500\u2500 Strong pairs list \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
                strong_pairs = []
                cols_list = list(num_corr.columns)
                for i in range(len(cols_list)):
                    for j in range(i + 1, len(cols_list)):
                        r_val = float(num_corr.iloc[i, j])
                        if abs(r_val) >= corr_threshold:
                            if abs(r_val) >= 0.7:
                                strength = "silna"
                            elif abs(r_val) >= 0.5:
                                strength = "umiarkowana"
                            else:
                                strength = "s\u0142aba"
                            direction = "dodatnia" if r_val > 0 else "ujemna"
                            strong_pairs.append({
                                "Zmienna A":    cols_list[i],
                                "Zmienna B":    cols_list[j],
                                "r":            round(r_val, 4),
                                "|r|":          round(abs(r_val), 4),
                                "Si\u0142a":    strength,
                                "Kierunek":     direction,
                            })

                if strong_pairs:
                    strong_df = pd.DataFrame(strong_pairs).sort_values("|r|", ascending=False)
                    n_strong = len(strong_df)
                    st.markdown(f"**\U0001f517 Silnie skorelowane pary (|r| \u2265 {corr_threshold:.2f}) \u2014 {n_strong} par:**")


                    st.dataframe(
                        strong_df.style.apply(_color_pair_row, axis=1)
                                       .format({'r': '{:.4f}', '|r|': '{:.4f}'}),
                        use_container_width=True, hide_index=True
                    )
                else:
                    st.info(f"Brak par o |r| \u2265 {corr_threshold:.2f}. Obni\u017c pr\u00f3g aby zobaczy\u0107 wi\u0119cej par.")

                # \u2500\u2500 Heatmap \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
                if show_heatmap:
                    fig = px.imshow(
                        num_corr, color_continuous_scale='RdBu_r',
                        zmin=-1, zmax=1,
                        title=f'Mapa ciep\u0142a korelacji ({corr_method.title()})',
                        text_auto='.2f'
                    )
                    fig.update_layout(height=max(400, len(corr_vars) * 40 + 100))
                    st.plotly_chart(fig, use_container_width=True, key=f"pc_ols_{res.get('dep_var','ols')}_diag")

            except Exception as e:
                st.error(str(e))

# -------------------------------------------------------------
# MODU? 4: REGRESJA
# -------------------------------------------------------------
elif menu == "\U0001f4c9 Regresja":
    module_header("\U0001f4c9", "Regresja", "OLS (liniowa) i logistyczna (binarna/wielomianowa)")
    tab_ols, tab_log = st.tabs(["\U0001f4c9 OLS (liniowa)", "\U0001f4c9 Logistyczna"])

    with tab_ols:
        st.markdown("##### Regresja Liniowa OLS -- wyniki w stylu SPSS")

        with st.expander("\U0001f4d6 Jak wykona\u0107 i interpretowa\u0107 regresj\u0119 -- kliknij aby rozwin\u0105\u0107", expanded=False):
            st.markdown("""
    **Kiedy u\u017cywa\u0107?** Gdy chcesz sprawdzi\u0107, kt\u00f3re zmienne (predyktory X) przewiduj\u0105 warto\u015b\u0107 innej zmiennej (Y), oraz jak silny jest ten zwi\u0105zek.

    **Jak wykona\u0107:**
    1. Wybierz **zmienn\u0105 zale\u017cn\u0105 (Y)** -- musi by\u0107 numeryczna (np. wynik testu, poziom satysfakcji).
    2. Wybierz **predyktory** w jednym lub kilku blokach.
    3. Bloki hierarchiczne (jak w SPSS): ka\u017cdy kolejny blok dodaje nowe zmienne i pokazuje przyrost R\u00b2.
    4. Kliknij **Uruchom regresj\u0119**.

    **Jak interpretowa\u0107 wyniki:**

    | Wska\u017anik | Interpretacja |
    |---|---|
    | **R\u00b2** | % wariancji Y wyja\u015bniany przez model. Wy\u017cszy = lepszy. |
    | **Skorygowane R\u00b2** | R\u00b2 poprawiony o liczb\u0119 predyktor\u00f3w -- por\u00f3wnuj mi\u0119dzy modelami. |
    | **\u0394R\u00b2** | O ile wzros\u0142o R\u00b2 po dodaniu nowego bloku predyktor\u00f3w. |
    | **F modelu / p** | Czy ca\u0142y model jest istotny statystycznie (p < 0.05 = TAK). |
    | **F zmiany / p** | Czy nowy blok predyktor\u00f3w istotnie poprawi\u0142 model. |
    | **B** | Niestandaryzowany wsp\u00f3\u0142czynnik: zmiana Y przy wzro\u015bcie X o 1 jednostk\u0119. |
    | **Beta (std.)** | Standaryzowany -- por\u00f3wnuje si\u0142\u0119 r\u00f3\u017cnych predyktor\u00f3w (niezale\u017cnie od skali). |
    | **VIF** | Wska\u017anik wsp\u00f3\u0142liniowo\u015bci: < 5 OK \u00b7 5-10 uwaga \u00b7 > 10 problem. |
    | **Tolerancja** | 1/VIF. Im ni\u017csza, tym wi\u0119kszy problem ze wsp\u00f3\u0142liniowo\u015bci\u0105. |

    **Wykresy diagnostyczne:**
    - *Reszty vs Dopasowane* -- punkty rozmieszczone losowo wok\u00f3\u0142 0 = OK (spe\u0142nione za\u0142o\u017cenie homoskedastyczno\u015bci).
    - *Q-Q plot* -- punkty blisko czerwonej linii = reszty maj\u0105 rozk\u0142ad normalny = OK.
            """)

        st.info("Dodawaj predyktory **blokami** (hierarchicznie). Ka\u017cdy blok poka\u017ce zmian\u0119 R\u00b2 i test F-zmiany.")

        dep_var = st.selectbox("\U0001f3af Zmienna zale\u017cna (Y):", numeric_cols, format_func=lambda x: get_var_display_name(x, var_labels))

        if st.button("\u2795 Dodaj blok predyktor\u00f3w"):
            st.session_state.reg_blocks.append([])
            st.rerun()

        blocks_to_delete = None
        for b_idx in range(len(st.session_state.reg_blocks)):
            c1, c2 = st.columns([6, 1])
            with c1:
                chosen = st.multiselect(
                    f"Blok {b_idx + 1} -- predyktory:",
                    [c for c in numeric_cols if c != dep_var],
                    default=st.session_state.reg_blocks[b_idx],
                    format_func=lambda x: get_var_display_name(x, var_labels),
                    key=f"reg_block_{b_idx}"
                )
                st.session_state.reg_blocks[b_idx] = chosen
            with c2:
                if b_idx > 0 and st.button("\U0001f5d1\ufe0f", key=f"del_block_{b_idx}", help="Usu\u0144 blok"):
                    blocks_to_delete = b_idx

        if blocks_to_delete is not None:
            st.session_state.reg_blocks.pop(blocks_to_delete)
            st.rerun()

        st.divider()
        if st.button("\u25b6\ufe0f Uruchom regresj\u0119", type="primary"):
            valid_blocks = [b for b in st.session_state.reg_blocks if b]
            if not valid_blocks:
                st.error("Dodaj co najmniej jeden predyktor.")
            else:
                with st.spinner("Obliczanie regresji OLS..."):
                    st.session_state.regression_results = run_regression_block(
                        df_raw, dep_var, valid_blocks,
                        weights=st.session_state.weights if use_weights else None
                    )

        for res in st.session_state.regression_results:
            if 'error' in res:
                st.error(res['error'])
                continue
            dep_label = var_labels.get(res['dep_var'], res['dep_var'])
            with st.expander(f"\U0001f4ca Blok {res['Blok']} -- [{res['dep_var']}] {dep_label}", expanded=True):
                m1, m2, m3, m4 = st.columns(4)
                m1.metric("N", f"{res['N']:,}")
                m2.metric("R\u00b2", f"{res['R2']:.4f}")
                m3.metric("Skorygowane R\u00b2", f"{res['Skor_R2']:.4f}")
                m4.metric("\u0394R\u00b2", f"{res['Delta_R2']:.4f}")
                f1, f2, f3, f4 = st.columns(4)
                f_pval = res['p (F modelu)']
                f1.metric("F modelu", f"{res['F modelu']:.3f}")
                f2.metric("p (F modelu)", f"('OK' if f_pval < 0.05 else 'NS') {f_pval:.4f}")
                fc, fcp = res['F zmiany'], res['p (F zmiany)']
                try:
                    if not np.isnan(fc):
                        f3.metric("F zmiany", f"{fc:.3f}")
                        f4.metric("p (F zmiany)", f"('OK' if fcp < 0.05 else 'NS') {fcp:.4f}")
                except: pass

                st.markdown("**Wsp\u00f3\u0142czynniki regresji:**")
                coef_df = res['coef_df'].copy()



                styled = coef_df.style \
                    .format({'B': '{:.4f}', 'B\u0142\u0105d std. B': '{:.4f}', 'Beta (std.)': '{:.4f}',
                             't': '{:.3f}', 'p-value': '{:.4f}', 'VIF': '{:.2f}', 'Tolerancja': '{:.4f}'}) \
                    .map(_style_p, subset=['p-value']) \
                    .map(_style_vif, subset=['VIF'])
                st.dataframe(styled, use_container_width=True)
                st.caption("\U0001f7e2 p < 0.05 -- istotne statystycznie \u00b7 VIF > 10 = wsp\u00f3\u0142liniowo\u015b\u0107 (czerwony) \u00b7 VIF 5-10 = uwaga (pomara\u0144czowy)")

                # Diagnostic plots
                all_pred = res['Wszystkie predyktory']
                df_diag = df_raw[[res['dep_var']] + all_pred].dropna()
                if len(df_diag) > 5:
                    X_d = sm.add_constant(df_diag[all_pred])
                    mod_d = sm.OLS(df_diag[res['dep_var']], X_d).fit()
                    ch1, ch2 = st.columns(2)
                    with ch1:
                        fig_r = px.scatter(x=mod_d.fittedvalues, y=mod_d.resid,
                                           labels={'x': 'Wartosci dopasowane', 'y': 'Reszty'},
                                           title='Reszty vs Warto\u015bci dopasowane', color_discrete_sequence=['#2E75B6'])
                        fig_r.add_hline(y=0, line_dash='dash', line_color='red')
                        st.plotly_chart(fig_r, use_container_width=True, key=f"pc_ols_{res.get('dep_var','ols')}_resid")
                    with ch2:
                        (osm, osr), (slope, intercept, _r) = stats.probplot(mod_d.resid)
                        fig_qq = go.Figure()
                        fig_qq.add_trace(go.Scatter(x=list(osm), y=list(osr), mode='markers', name='Reszty', marker=dict(color='#2E75B6')))
                        fig_qq.add_trace(go.Scatter(x=[min(osm), max(osm)], y=[slope * min(osm) + intercept, slope * max(osm) + intercept],
                                                    mode='lines', name='Linia ref.', line=dict(color='red', dash='dash')))
                        fig_qq.update_layout(title='Wykres Q-Q (normalno\u015b\u0107 reszt)',
                                             xaxis_title='Kwantyle teoretyczne', yaxis_title='Kwantyle pr\u00f3bkowe')
                        st.plotly_chart(fig_qq, use_container_width=True, key=f"pc_ols_{res.get('dep_var','ols')}_qq")

    # \u2500\u2500 TAB: Regresja Logistyczna \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
    with tab_log:
        st.markdown("##### Regresja Logistyczna")

        with st.expander("\U0001f4d6 Jak wykona\u0107 i interpretowa\u0107", expanded=False):
            st.markdown("""
**Czym jest regresja logistyczna?**
Modeluje prawdopodobie\u0144stwo przynale\u017cno\u015bci do kategorii.
- **Binarna** \u2014 zmienna zale\u017cna ma 2 kategorie (np. zakup: tak/nie, 0/1)
- **Wielomianowa (Multinomial)** \u2014 zmienna zale\u017cna ma 3+ kategorie

**Kluczowe wska\u017aniki:**

| Wska\u017anik | Interpretacja |
|---|---|
| **Iloraz szans (OR)** | OR > 1: zmienna zwi\u0119ksza szans\u0119; OR < 1: zmniejsza |
| **95% CI** | Przedzia\u0142 ufno\u015bci dla OR. Je\u015bli nie zawiera 1: istotne |
| **p-value** | < 0.05: zmienna istotnie wp\u0142ywa na wynik |
| **Pseudo R\u00b2 (McFadden)** | 0.2-0.4: dobre dopasowanie modelu |
| **AIC** | Kryterium informacyjne. Mniejszy = lepszy model |
            """)

        log_type = st.radio("Typ regresji:", ["Binarna (Logit)", "Wielomianowa (MNLogit)"],
                             horizontal=True, key="log_type")
        st.divider()
        col_lg1, col_lg2 = st.columns(2)
        with col_lg1:
            log_dep = st.selectbox("\U0001f3af Zmienna zale\u017cna:", visible_columns,
                                    format_func=lambda x: get_var_display_name(x, var_labels), key="log_dep")
            dep_vals = sorted(df[log_dep].dropna().unique())
            st.caption(f"Unikalne warto\u015bci ({len(dep_vals)}): {', '.join(str(v) for v in dep_vals[:8])}"
                       + (" ..." if len(dep_vals) > 8 else ""))
            if log_type == "Binarna (Logit)" and len(dep_vals) != 2:
                st.warning("Regresja binarna wymaga dok\u0142adnie 2 unikalnych warto\u015bci.")
        with col_lg2:
            log_indep = st.multiselect("\U0001f4e6 Predyktory:",
                                        [c for c in visible_columns if c != log_dep],
                                        format_func=lambda x: get_var_display_name(x, var_labels), key="log_indep")
            log_dummy = st.checkbox("Automatycznie zakoduj zmienne kategoryczne (dummy coding)",
                                     value=True, key="log_dummy")

        if st.button("\u25b6\ufe0f Uruchom regresj\u0119 logistyczn\u0105", type="primary", key="log_run"):
            if not log_indep:
                st.error("Wybierz co najmniej jeden predyktor.")
            else:
                with st.spinner("Obliczanie..."):
                    try:
                        df_lg = df[[log_dep] + log_indep].dropna().copy()
                        n_obs = len(df_lg)
                        # Weights for logistic
                        _w_log = None
                        if use_weights and st.session_state.weights is not None:
                            _w_log = pd.Series(st.session_state.weights, index=df.index
                                               ).reindex(df_lg.index).fillna(1).clip(lower=0).values
                            n_obs = round(float(_w_log.sum()), 1)
                        if len(df_lg) < 20:
                            st.error(f"Za ma\u0142o obserwacji ({len(df_lg)}).")
                        else:
                            y_series = df_lg[log_dep]
                            if log_type == "Binarna (Logit)":
                                uniq = sorted(y_series.dropna().unique())
                                if len(uniq) != 2:
                                    st.error("Zmienna zale\u017cna musi mie\u0107 dok\u0142adnie 2 warto\u015bci.")
                                    st.stop()
                                y = (y_series == uniq[1]).astype(np.float64)
                                dep_ref = str(uniq[0]); dep_pos = str(uniq[1])
                            else:
                                y = y_series.astype(str)
                                dep_ref = str(sorted(y_series.unique())[0])
                                dep_pos = "wszystkie"

                            # \u2500\u2500 Build X matrix robustly \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
                            # Separate categorical (need dummies) from numeric
                            cat_cols = [c for c in log_indep
                                        if df_lg[c].dtype == object
                                        or (df_lg[c].nunique() <= 10 and log_dummy)]
                            num_cols_lg = [c for c in log_indep if c not in cat_cols]

                            parts = []
                            if num_cols_lg:
                                num_part = df_lg[num_cols_lg].apply(
                                    pd.to_numeric, errors='coerce').fillna(0).astype(np.float64)
                                parts.append(num_part)
                            if cat_cols and log_dummy:
                                for cc in cat_cols:
                                    dummies = pd.get_dummies(
                                        df_lg[cc].astype(str),
                                        prefix=cc, drop_first=True
                                    ).astype(np.float64)
                                    parts.append(dummies)
                            elif cat_cols:
                                for cc in cat_cols:
                                    df_lg[cc] = pd.Categorical(
                                        df_lg[cc].astype(str)).codes.astype(np.float64)
                                parts.append(df_lg[cat_cols].astype(np.float64))

                            if parts:
                                X = pd.concat(parts, axis=1)
                            else:
                                st.error("Brak predyktor\u00f3w do modelu.")
                                st.stop()

                            # Ensure all float64 \u2014 no object columns
                            X = X.astype(np.float64)
                            # Keep as DataFrame so model.params has named index
                            X_const = sm.add_constant(X, has_constant='add').astype(np.float64)
                            y_fit = y.astype(np.float64)

                            if log_type == "Binarna (Logit)":
                                if _w_log is not None:
                                    model = sm.Logit(y_fit, X_const,
                                                     freq_weights=_w_log).fit(disp=False, maxiter=200)
                                else:
                                    model = sm.Logit(y_fit, X_const).fit(disp=False, maxiter=200)
                                params = model.params
                                pvals  = model.pvalues
                                conf   = model.conf_int()
                                or_vals = np.exp(params)
                                or_lo   = np.exp(conf.iloc[:, 0])
                                or_hi   = np.exp(conf.iloc[:, 1])
                                coef_df = pd.DataFrame({
                                    'Zmienna':      params.index.tolist(),
                                    'Wspolczynnik': params.values.round(4),
                                    'Iloraz szans': or_vals.values.round(4),
                                    'CI 95% (dol)': or_lo.values.round(4),
                                    'CI 95% (gor)': or_hi.values.round(4),
                                    'p-value':      pvals.values.round(4),
                                    'Istotny':      ['Tak' if p < 0.05 else 'Nie' for p in pvals.values]
                                })
                                result_entry = {'type': 'Binarna', 'dep_var': log_dep, 'dep_ref': dep_ref,
                                                'dep_pos': dep_pos, 'indep_vars': log_indep, 'n_obs': n_obs,
                                                'pseudo_r2': model.prsquared, 'llr_p': model.llr_pvalue,
                                                'aic': model.aic, 'bic': model.bic, 'log_lik': model.llf,
                                                'coef_df': coef_df, 'model': model, 'error': None}
                            else:
                                model = sm.MNLogit(y_fit, X_const).fit(disp=False, maxiter=200)
                                result_entry = {'type': 'Wielomianowa', 'dep_var': log_dep, 'dep_ref': dep_ref,
                                                'indep_vars': log_indep, 'n_obs': n_obs,
                                                'pseudo_r2': model.prsquared, 'llr_p': model.llr_pvalue,
                                                'aic': model.aic, 'bic': model.bic, 'log_lik': model.llf,
                                                'coef_df': None, 'model': model, 'error': None}
                            st.session_state.logistic_results.append(result_entry)
                            st.success("\u2705 Regresja logistyczna obliczona!")
                    except Exception as _lg_err:
                        st.error(f"B\u0142\u0105d: {_lg_err}")

        for res_lg in st.session_state.logistic_results:
            if res_lg.get('error'): st.error(res_lg['error']); continue
            dep_lbl = var_labels.get(res_lg['dep_var'], res_lg['dep_var'])
            with st.expander(f"\U0001f4c9 {res_lg['type']}: [{res_lg['dep_var']}] {dep_lbl}", expanded=True):
                m1,m2,m3,m4,m5 = st.columns(5)
                m1.metric("N", f"{res_lg['n_obs']:,}")
                m2.metric("Pseudo R\u00b2", f"{res_lg['pseudo_r2']:.4f}")
                m3.metric("p (LLR)", f"{res_lg['llr_p']:.4f}")
                m4.metric("AIC", f"{res_lg['aic']:.1f}")
                m5.metric("BIC", f"{res_lg['bic']:.1f}")
                if res_lg['type'] == 'Binarna' and res_lg['coef_df'] is not None:
                    st.markdown(f"**Kategoria ref.:** `{res_lg['dep_ref']}` | **Modelowana:** `{res_lg['dep_pos']}`")
                    cdf = res_lg['coef_df'].copy()
                    styled = cdf.style.apply(_color_sig, axis=1).format(
                        {'Wspolczynnik': '{:.4f}', 'Iloraz szans': '{:.4f}',
                         'CI 95% (dol)': '{:.4f}', 'CI 95% (gor)': '{:.4f}', 'p-value': '{:.4f}'})
                    st.dataframe(styled, use_container_width=True, hide_index=True)
                    plot_df = cdf[cdf['Zmienna'] != 'const'].copy()
                    if not plot_df.empty:
                        fig_or = go.Figure()
                        colors = ['#C00000' if p >= 0.05 else '#2E75B6' for p in plot_df['p-value']]
                        fig_or.add_trace(go.Scatter(x=plot_df['Iloraz szans'], y=plot_df['Zmienna'],
                                                    mode='markers', marker=dict(size=10, color=colors),
                                                    error_x=dict(type='data', symmetric=False,
                                                                 array=(plot_df['CI 95% (gor)'] - plot_df['Iloraz szans']).tolist(),
                                                                 arrayminus=(plot_df['Iloraz szans'] - plot_df['CI 95% (dol)']).tolist())))
                        fig_or.add_vline(x=1, line_dash='dash', line_color='gray')
                        fig_or.update_layout(title='Ilorazy szans (OR) z 95% CI', xaxis_title='OR',
                                              height=max(300, len(plot_df)*35+100), showlegend=False)
                        st.plotly_chart(fig_or, use_container_width=True, key=f"pc_log_{res_lg.get('dep_var','log')}_or")
                else:
                    st.text(res_lg['model'].summary().as_text())

        if st.session_state.logistic_results:
            if st.button("\U0001f5d1\ufe0f Wyczy\u015b\u0107 wyniki regresji logistycznej", type="secondary", key="log_clear_tab"):
                st.session_state.logistic_results = []
                st.rerun()

# -------------------------------------------------------------
# MODUL 5: ANOVA
# -------------------------------------------------------------
# =============================================================
# MODUL: TESTY NORMALNOSCI
# =============================================================
elif menu == "\U0001f4d0 Testy Normalno\u015bci":
    module_header("\U0001f4d0", "Testy Normalno\u015bci",
                  "Sprawdzanie za\u0142o\u017cenia normalno\u015bci rozk\u0142adu \u2014 wymagane przed ANOVA, regresj\u0105 i innymi testami parametrycznymi")

    with st.expander("\U0001f4d6 Jak przeprowadzi\u0107 i interpretowa\u0107 testy normalno\u015bci", expanded=False):
        st.markdown("""
### Czym jest test normalno\u015bci?

Test normalno\u015bci sprawdza, czy rozk\u0142ad zmiennej jest zbli\u017cony do rozk\u0142adu normalnego (Gaussa). Jest to
za\u0142o\u017cenie wielu test\u00f3w parametrycznych: **t-testu, ANOVA, regresji liniowej, korelacji Pearsona**.

---

### Dost\u0119pne testy

| Test | Kiedy stosowa\u0107 | Hipoteza zerowa (H\u2080) |
|---|---|---|
| **Shapiro-Wilk** | Ma\u0142e i \u015brednie pr\u00f3by (N \u2264 2000) \u2014 **najsilniejszy test** | Rozk\u0142ad jest normalny |
| **Kolmogorov-Smirnov** | Du\u017ce pr\u00f3by (N > 2000), por\u00f3wnanie z rozk\u0142adem teoretycznym | Rozk\u0142ad jest normalny |
| **Lilliefors** | Wariant K-S gdy \u015brednia i odch. std. s\u0105 szacowane z danych | Rozk\u0142ad jest normalny |
| **D\u2019Agostino-Pearson** | Opiera si\u0119 na sko\u015bno\u015bci i kurtozie | Rozk\u0142ad jest normalny |

---

### Jak interpretowa\u0107 wyniki?

**p-value:**
- **p > 0.05** \u2192 Brak podstaw do odrzucenia H\u2080 \u2192 dane **mog\u0105 pochodzi\u0107** z rozk\u0142adu normalnego \u2705
- **p \u2264 0.05** \u2192 Odrzucamy H\u2080 \u2192 dane **nie pochodz\u0105** z rozk\u0142adu normalnego \u274c

> \u26a0\ufe0f **Uwaga:** Dla du\u017cych pr\u00f3b (N > 200) nawet ma\u0142e, praktycznie nieistotne odchylenia od normalno\u015bci
> daj\u0105 p < 0.05. W takich przypadkach **wa\u017cniejsza jest ocena wizualna** (Q-Q plot, histogram)
> i miary sko\u015bno\u015bci / kurtozy. Regu\u0142a: |sko\u015bno\u015b\u0107| < 2 i |kurtoza| < 7 to akceptowalna normalno\u015b\u0107.

---

### Ocena wizualna

- **Histogram** \u2014 kszta\u0142t dzwonu = dobry znak; skos lub gruby ogon = odchylenie
- **Wykres Q-Q** \u2014 punkty blisko prostej = normalno\u015b\u0107; wygi\u0119cie = skos lub kurtoza

---

### Praktyczne zalecenia (jak w SPSS)

| Wielko\u015b\u0107 pr\u00f3by | Zalecany test |
|---|---|
| N < 50 | Shapiro-Wilk |
| 50 \u2264 N \u2264 2000 | Shapiro-Wilk + wizualna ocena Q-Q |
| N > 2000 | Lilliefors lub D\u2019Agostino + wizualna ocena Q-Q |
        """)

    st.divider()

    norm_vars = st.multiselect(
        "Wybierz zmienne numeryczne do testowania:",
        numeric_cols,
        format_func=lambda x: get_var_display_name(x, var_labels),
        key="norm_vars"
    )

    ncol1, ncol2 = st.columns(2)
    with ncol1:
        norm_tests = st.multiselect(
            "Wybierz testy:",
            ["Shapiro-Wilk", "Kolmogorov-Smirnov (Lilliefors)", "D\u2019Agostino-Pearson"],
            default=["Shapiro-Wilk", "D\u2019Agostino-Pearson"],
            key="norm_tests"
        )
        norm_alpha = st.select_slider(
            "Poziom istotno\u015bci (\u03b1):",
            options=[0.01, 0.05, 0.10],
            value=0.05,
            key="norm_alpha"
        )
    with ncol2:
        norm_show_qq   = st.checkbox("Wykres Q-Q", value=True, key="norm_qq")
        norm_show_hist = st.checkbox("Histogram z krzywa normaln\u0105", value=True, key="norm_hist")
        norm_show_desc = st.checkbox("Statystyki opisowe (sko\u015bno\u015b\u0107, kurtoza)", value=True, key="norm_desc")

    if st.button("\u25b6\ufe0f Przeprowad\u017a testy normalno\u015bci", type="primary",
                 key="norm_run") and norm_vars and norm_tests:
        for var in norm_vars:
            _w_norm = st.session_state.weights if use_weights else None
            if _w_norm is not None:
                mask = df_raw[var].notna()
                x_raw = df_raw.loc[mask, var].values.astype(float)
                w_raw = pd.Series(_w_norm, index=df_raw.index).loc[mask].values
                # Weighted sample via repetition (SPSS approach for normality tests)
                reps = np.round(w_raw / w_raw.min()).astype(int)
                x = np.repeat(x_raw, reps)
            else:
                x = df_raw[var].dropna().values.astype(float)

            n = len(x)
            lbl = var_labels.get(var, var)

            st.markdown(f"---\n#### \U0001f4ca `{var}` \u2014 {lbl}")
            st.caption(f"N = {len(df_raw[var].dropna()):,} obserwacji (efektywne N do test\u00f3w: {n:,})")

            # \u2500\u2500 Normality tests \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
            test_rows = []
            if "Shapiro-Wilk" in norm_tests:
                if n > 5000:
                    st.warning("Shapiro-Wilk: pr\u00f3ba zbyt du\u017ca (N > 5000). U\u017cyto losowej pr\u00f3bki 5000.")
                    x_sw = np.random.choice(x, 5000, replace=False)
                else:
                    x_sw = x
                sw_stat, sw_p = stats.shapiro(x_sw)
                test_rows.append({
                    "Test": "Shapiro-Wilk",
                    "Statystyka": round(float(sw_stat), 4),
                    "p-value": round(float(sw_p), 4),
                    "Wynik": "\u2705 Normalny" if sw_p > norm_alpha else "\u274c Nienormalny"
                })

            if "Kolmogorov-Smirnov (Lilliefors)" in norm_tests:
                from statsmodels.stats.diagnostic import kstest_normal
                lf_stat, lf_p = kstest_normal(x, dist='norm')
                test_rows.append({
                    "Test": "Lilliefors (K-S)",
                    "Statystyka": round(float(lf_stat), 4),
                    "p-value": round(float(lf_p), 4),
                    "Wynik": "\u2705 Normalny" if lf_p > norm_alpha else "\u274c Nienormalny"
                })

            if "D\u2019Agostino-Pearson" in norm_tests:
                dag_stat, dag_p = stats.normaltest(x)
                test_rows.append({
                    "Test": "D\u2019Agostino-Pearson",
                    "Statystyka": round(float(dag_stat), 4),
                    "p-value": round(float(dag_p), 4),
                    "Wynik": "\u2705 Normalny" if dag_p > norm_alpha else "\u274c Nienormalny"
                })

            test_df = pd.DataFrame(test_rows)

            def _style_norm_row(row):
                color = '#E2EFDA' if '\u2705' in row['Wynik'] else '#FCE4D6'
                return [f'background-color: {color}'] * len(row)

            st.dataframe(test_df.style.apply(_style_norm_row, axis=1)
                         .format({'Statystyka': '{:.4f}', 'p-value': '{:.4f}'}),
                         use_container_width=True, hide_index=True)

            # \u2500\u2500 Descriptives \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
            if norm_show_desc:
                skew_v = float(stats.skew(x))
                kurt_v = float(stats.kurtosis(x))
                sk_interp = ("symetryczny" if abs(skew_v) < 0.5
                             else ("lekko sko\u015bny" if abs(skew_v) < 1.0
                                   else ("umiarkowanie sko\u015bny" if abs(skew_v) < 2.0
                                         else "silnie sko\u015bny")))
                kt_interp = ("mezokurtyczny (normalny)" if abs(kurt_v) < 1.0
                             else ("platykurtyczny (sp\u0142aszczony)" if kurt_v < 0
                                   else "leptokurtyczny (smuk\u0142y)"))
                desc_c1, desc_c2, desc_c3, desc_c4 = st.columns(4)
                desc_c1.metric("Sko\u015bno\u015b\u0107", f"{skew_v:.4f}")
                desc_c2.metric("Kurtoza", f"{kurt_v:.4f}")
                desc_c3.metric("Ocena sko\u015bno\u015bci", sk_interp)
                desc_c4.metric("Ocena kurtozy", kt_interp)

            # \u2500\u2500 Plots \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
            plot_cols = st.columns(2 if (norm_show_qq and norm_show_hist) else 1)
            plot_idx = 0

            if norm_show_qq:
                (osm, osr), (slope, intercept, r) = stats.probplot(x, dist='norm')
                fig_qq = go.Figure()
                fig_qq.add_trace(go.Scatter(
                    x=list(osm), y=list(osr),
                    mode='markers', name='Obserwacje',
                    marker=dict(color='#2E75B6', size=4, opacity=0.6)
                ))
                fig_qq.add_trace(go.Scatter(
                    x=[min(osm), max(osm)],
                    y=[slope * min(osm) + intercept, slope * max(osm) + intercept],
                    mode='lines', name='Linia normalna',
                    line=dict(color='#C00000', dash='dash')
                ))
                fig_qq.update_layout(
                    title=f"Wykres Q-Q: {lbl}",
                    xaxis_title="Kwantyle teoretyczne",
                    yaxis_title="Kwantyle pr\u00f3bkowe",
                    height=350, showlegend=True
                )
                with plot_cols[plot_idx]:
                    st.plotly_chart(fig_qq, use_container_width=True, key=f"qq_{var}")
                plot_idx += 1

            if norm_show_hist:
                fig_hist = go.Figure()
                fig_hist.add_trace(go.Histogram(
                    x=x, name="Dane",
                    histnorm='probability density',
                    marker_color='#2E75B6', opacity=0.7,
                    nbinsx=min(50, max(10, n // 10))
                ))
                # Overlay normal curve
                x_range = np.linspace(x.min(), x.max(), 200)
                mu, sigma = float(x.mean()), float(x.std())
                y_norm = stats.norm.pdf(x_range, mu, sigma)
                fig_hist.add_trace(go.Scatter(
                    x=x_range, y=y_norm, mode='lines', name='Rozk\u0142ad normalny',
                    line=dict(color='#C00000', width=2)
                ))
                fig_hist.update_layout(
                    title=f"Histogram z krzywa normaln\u0105: {lbl}",
                    xaxis_title=lbl, yaxis_title="G\u0119sto\u015b\u0107",
                    height=350, showlegend=True, barmode='overlay'
                )
                with plot_cols[plot_idx]:
                    st.plotly_chart(fig_hist, use_container_width=True, key=f"hist_{var}")


elif menu == "\U0001f4ca ANOVA":
    module_header("\U0001f4ca", "ANOVA", "Jednoczynnikowa analiza wariancji z testem post-hoc Tukeya")

    with st.expander("\U0001f4d6 Jak wykona\u0107 i interpretowa\u0107 ANOVA -- kliknij aby rozwin\u0105\u0107", expanded=False):
        st.markdown("""
**Kiedy u\u017cywa\u0107?** Gdy chcesz por\u00f3wna\u0107 **\u015brednie 3 lub wi\u0119cej grup** jednocze\u015bnie.
Przyk\u0142ad: czy \u015brednia satysfakcja r\u00f3\u017cni si\u0119 mi\u0119dzy miastem A, B i C?
Dla 2 grup u\u017cyj testu T (modu\u0142 Analizy \u2192 \u015arednie).

**Jak wykona\u0107:**
1. Wybierz **zmienn\u0105 zale\u017cn\u0105** (ci\u0105g\u0142a, numeryczna -- np. satysfakcja 1-10).
2. Wybierz **czynnik grupuj\u0105cy** (kategoryczna -- np. miasto, wiek, segment).
3. Kliknij **Uruchom ANOVA**.

**Jak interpretowa\u0107 wyniki:**

| Wska\u017anik | Interpretacja |
|---|---|
| **F** | Im wy\u017csze, tym wi\u0119ksza r\u00f3\u017cnica mi\u0119dzy grupami wzgl\u0119dem zmienno\u015bci wewn\u0105trz grup. |
| **p-value** | p < 0.05 = grupy r\u00f3\u017cni\u0105 si\u0119 istotnie statystycznie. |
| **Eta\u00b2 (\u03b7\u00b2)** | Miara si\u0142y efektu: < 0.01 s\u0142aby \u00b7 0.01-0.06 umiarkowany \u00b7 > 0.14 du\u017cy. |
| **MS (mi\u0119dzy grupami)** | Wariancja wyja\u015bniona przez przynale\u017cno\u015b\u0107 do grupy. |
| **MS (wewn\u0105trz grup)** | Wariancja wewn\u0105trz ka\u017cdej grupy (b\u0142\u0105d). |
| **Test Levene'a** | Sprawdza jednorodnosc wariancji. p < 0.05 = wariancje niejednorodne (naruszenie za\u0142o\u017cenia). |

**Test post-hoc Tukey HSD:**
Gdy ANOVA jest istotna (p < 0.05), Tukey wskazuje **kt\u00f3re konkretnie pary grup** r\u00f3\u017cni\u0105 si\u0119 od siebie.
p < 0.05 dla danej pary = ta para jest istotnie r\u00f3\u017cna.

**Za\u0142o\u017cenia ANOVA:** normalno\u015b\u0107 rozk\u0142adu w grupach, jednorodnosc wariancji (Levene), niezale\u017cno\u015b\u0107 obserwacji.
        """)

    st.info("Por\u00f3wnaj \u015brednie zmiennej ci\u0105g\u0142ej mi\u0119dzy grupami zdefiniowanymi przez zmienn\u0105 kategoryczn\u0105.")

    col1, col2 = st.columns(2)
    with col1:
        anova_dep = st.selectbox("\U0001f3af Zmienna zale\u017cna (ci\u0105g\u0142a):", numeric_cols, format_func=lambda x: get_var_display_name(x, var_labels))
    with col2:
        anova_grp = st.selectbox("\U0001f465 Czynnik grupuj\u0105cy (kategoryczna):", visible_columns, format_func=lambda x: get_var_display_name(x, var_labels))

    if st.button("\u25b6\ufe0f Uruchom ANOVA", type="primary"):
        with st.spinner("Obliczanie ANOVA..."):
            result, err = run_anova(df_raw, anova_dep, anova_grp, df,
                                      weights=st.session_state.weights if use_weights else None)
            if err:
                st.error(err)
            else:
                st.session_state.anova_results.append(result)

    for res in st.session_state.anova_results:
        dep_l = var_labels.get(res['dep_var'], res['dep_var'])
        grp_l = var_labels.get(res['group_var'], res['group_var'])
        with st.expander(f"\U0001f4ca ANOVA: [{res['dep_var']}] {dep_l} \u00d7 [{res['group_var']}] {grp_l}", expanded=True):
            sig_label = "\u2705 Istotna statystycznie (p < 0.05)" if res['p'] < 0.05 else "\u274c Brak istotno\u015bci (p \u2265 0.05)"
            m1, m2, m3, m4 = st.columns(4)
            m1.metric("F", f"{res['F']:.3f}")
            m2.metric("p-value", f"{res['p']:.4f}")
            m3.metric("Eta\u00b2 (efekt)", f"{res['eta2']:.4f}")
            m4.metric("Wynik", sig_label)

            st.markdown("**Tabela ANOVA:**")
            anova_table = pd.DataFrame({
                '\u0179r\u00f3d\u0142o': ['Mi\u0119dzy grupami', 'Wewn\u0105trz grup', 'Og\u00f3\u0142em'],
                'SS': [res['ss_between'], res['ss_within'], res['ss_total']],
                'df': [res['df_between'], res['df_within'], res['df_between'] + res['df_within']],
                'MS': [res['ms_between'], res['ms_within'], ''],
                'F': [res['F'], '', ''],
                'p': [res['p'], '', ''],
                'Eta\u00b2': [res['eta2'], '', ''],
            })
            st.dataframe(anova_table, use_container_width=True, hide_index=True)

            st.markdown(f"**Test Levene'a (jednorodnosc wariancji):** stat={res['lev_stat']:.3f}, p={res['lev_p']:.4f} -- {'Wariancje niejednorodne (p<0.05)' if res['lev_p'] < 0.05 else 'Wariancje jednorodne'}")

            st.markdown("**Statystyki opisowe wg grupy:**")
            st.dataframe(res['desc_df'].style.format({'Srednia': '{:.3f}', 'Odch. std.': '{:.3f}', 'Min': '{:.2f}', 'Max': '{:.2f}'}),
                         use_container_width=True, hide_index=True)

            # Bar chart with error bars
            fig_anova = go.Figure()
            fig_anova.add_trace(go.Bar(
                x=res['desc_df']['Grupa'].astype(str),
                y=res['desc_df']['Srednia'],
                error_y=dict(type='data', array=res['desc_df']['Odch. std.'].values, visible=True),
                marker_color='#2E75B6', name='\u015arednia \u00b1 Odch.std.'
            ))
            fig_anova.update_layout(title=f"\u015arednie wg grup -- {dep_l}", xaxis_title=grp_l, yaxis_title='Srednia', height=350)
            st.plotly_chart(fig_anova, use_container_width=True, key="pc_anova_bar")

            if not res['posthoc_df'].empty:
                st.markdown("**Test post-hoc Tukey HSD:**")
                st.dataframe(res['posthoc_df'], use_container_width=True, hide_index=True)

    if st.session_state.anova_results:
        if st.button("\U0001f5d1\ufe0f Wyczy\u015b\u0107 wyniki ANOVA", type="secondary"):
            st.session_state.anova_results = []
            st.rerun()

# -------------------------------------------------------------
# MODU? 6: ANALIZA CZYNNIKOWA
# -------------------------------------------------------------
elif menu == "\U0001f52c Analiza Czynnikowa":
    module_header("\U0001f52c", "Analiza Czynnikowa", "Eksploracyjna Analiza Czynnikowa (EFA)")

    with st.expander("\U0001f4d6 Jak wykona\u0107 i interpretowa\u0107 EFA -- kliknij aby rozwin\u0105\u0107", expanded=False):
        st.markdown("""
**Kiedy u\u017cywa\u0107?** Gdy chcesz odkry\u0107 **ukryte konstrukty (czynniki)** kryj\u0105ce si\u0119 za korelacjami mi\u0119dzy wieloma zmiennymi.
Przyk\u0142ad: 15 pyta\u0144 o satysfakcj\u0119 mo\u017ce odzwierciedla\u0107 3 ukryte wymiary: satysfakcja z produktu, obs\u0142ugi i ceny.

**Jak wykona\u0107:**
1. Wybierz **min. 3 zmienne numeryczne** ze wsp\u00f3lnej baterii pyta\u0144 (np. pytania Likerta 1-5).
2. Zdecyduj o **liczbie czynnik\u00f3w** -- zacznij od 2-4, u\u017cyj wykresu osypiska jako wskaz\u00f3wki.
3. Wybierz **rotacj\u0119**: Varimax (czynniki niezale\u017cne, najcz\u0119stsza), Promax (czynniki mog\u0105 by\u0107 powi\u0105zane).
4. Wybierz **metod\u0119 ekstrakcji**: Principal (PA), MinRes, ML.
5. Kliknij **Uruchom analiz\u0119 czynnikow\u0105**.

**Jak interpretowa\u0107 wyniki:**

| Wska\u017anik | Interpretacja |
|---|---|
| **KMO** | Adekwatno\u015b\u0107 pr\u00f3by: \u2265 0.9 znakomita \u00b7 \u2265 0.8 b.dobra \u00b7 \u2265 0.7 dobra \u00b7 \u2265 0.6 umiarkowana \u00b7 < 0.5 nieodpowiednia |
| **Test Bartletta (p)** | p < 0.05 = macierz korelacji nadaje si\u0119 do EFA (zmienne s\u0105 powi\u0105zane). |
| **Warto\u015b\u0107 w\u0142asna (EV)** | Zasada Kaisera: zachowaj czynniki z EV > 1. Sprawd\u017a na wykresie osypiska. |
| **% wyja\u015bnionej wariancji** | Ile zmienno\u015bci danych wyja\u015bnia dany czynnik. \u0141\u0105cznie powinno by\u0107 \u2265 50-60%. |
| **\u0141adunek czynnikowy** | Si\u0142a i kierunek powi\u0105zania zmiennej z czynnikiem. |\u0142adunek| \u2265 0.40 = istotny (pogrubiony). |
| **Komunalno\u015b\u0107 (h\u00b2)** | % wariancji danej zmiennej wyja\u015bniony przez wszystkie czynniki. < 0.30 = zmienna s\u0142abo pasuje. |

**Wykres osypiska (Scree Plot):** Szukaj miejsca, gdzie krzywa "ugina si\u0119" (elbow). Czynniki przed tym miejscem warto zachowa\u0107.

**Wskaz\u00f3wki praktyczne:**
- Ka\u017cda zmienna powinna \u0142adowa\u0107 istotnie (|\u2265 0.40|) na **jeden g\u0142\u00f3wny czynnik** (prosta struktura).
- Zmienne z nisk\u0105 komunalno\u015bci\u0105 (< 0.30) lub \u0142adunkami krzy\u017cowymi warto usun\u0105\u0107.
- Minimalna pr\u00f3ba: N \u2265 5 \u00d7 liczba zmiennych (najlepiej N \u2265 200).
        """)

    st.info("Zidentyfikuj ukryte konstrukty (czynniki) kryj\u0105ce si\u0119 za korelacjami mi\u0119dzy zmiennymi. Wyniki analogiczne do SPSS.")

    fa_vars = st.multiselect("Zmienne do analizy (min. 3):", numeric_cols, format_func=lambda x: get_var_display_name(x, var_labels))

    if fa_vars and len(fa_vars) >= 3:
        col1, col2, col3 = st.columns(3)
        with col1:
            n_factors = st.number_input("Liczba czynnik\u00f3w:", min_value=1, max_value=min(len(fa_vars) - 1, 15), value=min(3, len(fa_vars) - 1))
        with col2:
            rotation = st.selectbox("Rotacja:", ['varimax', 'promax', 'oblimin', 'quartimax', 'none'])
        with col3:
            method = st.selectbox("Metoda ekstrakcji:", ['principal', 'minres', 'ml'])

        show_scree = st.checkbox("\U0001f4c8 Wykres osypiska (Scree Plot)")

        if st.button("\u25b6\ufe0f Uruchom analiz\u0119 czynnikow\u0105", type="primary"):
            with st.spinner("Obliczanie analizy czynnikowej..."):
                result, err = run_factor_analysis(
                    df_raw, fa_vars, int(n_factors), rotation, method,
                    weights=st.session_state.weights if use_weights else None
                )
                if err:
                    st.error(err)
                else:
                    st.session_state.factor_results.append(result)

        for res in st.session_state.factor_results:
            with st.expander(f"\U0001f52c Analiza czynnikowa -- {res['rotation'].upper()} -- N={res['n']}", expanded=True):
                m1, m2, m3, m4 = st.columns(4)
                m1.metric("N obserwacji", f"{res['n']:,}")
                m2.metric("KMO", f"{res['kmo']:.3f}", help="\u22650.7 = dobra adekwatnosc proby")
                m3.metric("Bartlett Chi\u00b2", f"{res['bartlett_chi2']:.2f}")
                m4.metric("Bartlett p", f"{res['bartlett_p']:.4f}", delta="\u2705 OK" if res['bartlett_p'] < 0.05 else "\u274c")

                kmo_interp = ("Nieodpowiednia" if res['kmo'] < 0.5 else "S\u0142aba" if res['kmo'] < 0.6 else
                              "Umiarkowana" if res['kmo'] < 0.7 else "Dobra" if res['kmo'] < 0.8 else
                              "Bardzo dobra" if res['kmo'] < 0.9 else "Znakomita")
                st.caption(f"KMO = {res['kmo']:.3f} \u2192 adekwatnosc proby: **{kmo_interp}** | Bartlett p {'< 0.05 -- nadaje sie do EFA' if res['bartlett_p'] < 0.05 else '>= 0.05 -- macierz moze byc jednostkowa'}")

                st.markdown("**Macierz \u0142adunk\u00f3w czynnikowych** (pogrubione |\u0142adunek| \u2265 0.40):")


                load_display = res['loadings'].copy()
                comm_col = res['communalities']['Komunalnosc (h2)']
                load_display['Komunalnosc (h2)'] = comm_col
                styled_load = load_display.style \
                    .format('{:.3f}') \
                    .map(_style_loading, subset=res['loadings'].columns.tolist())
                st.dataframe(styled_load, use_container_width=True)

                st.markdown("**Wyja\u015bniona wariancja:**")
                st.dataframe(res['variance'].style.format('{:.3f}'), use_container_width=True)

                if show_scree:
                    ev_vals = res['eigenvalues']['Warto\u015b\u0107 w\u0142asna'].values[:min(len(fa_vars), 15)]
                    fig_scree = go.Figure()
                    fig_scree.add_trace(go.Scatter(y=ev_vals, x=list(range(1, len(ev_vals) + 1)),
                                                   mode='lines+markers', name='Warto\u015bci w\u0142asne',
                                                   line=dict(color='#2E75B6', width=2), marker=dict(size=8)))
                    fig_scree.add_hline(y=1, line_dash='dash', line_color='red', annotation_text='Kryterium Kaisera (EV=1)')
                    fig_scree.update_layout(title='Wykres osypiska (Scree Plot)',
                                            xaxis_title='Numer czynnika', yaxis_title='Warto\u015b\u0107 w\u0142asna', height=350)
                    st.plotly_chart(fig_scree, use_container_width=True, key="pc_efa_scree")

    elif fa_vars:
        st.warning("Wybierz co najmniej 3 zmienne.")

    if st.session_state.factor_results:
        if st.button("\U0001f5d1\ufe0f Wyczy\u015b\u0107 wyniki analizy czynnikowej", type="secondary"):
            st.session_state.factor_results = []
            st.rerun()

# -------------------------------------------------------------
# MODU? 7: EKSPORT DO EXCELA
# -------------------------------------------------------------
# =============================================================
# MODUL: CONJOINT
# =============================================================
elif menu == "\U0001f4ca Conjoint":
    module_header("\U0001f4ca", "Analiza Conjoint", "Rating-based (OLS) i CBC (Logit) \u2014 u\u017cyteczno\u015bci cz\u0105stkowe i wa\u017cno\u015b\u0107 atrybut\u00f3w")

    with st.expander("\U0001f4d6 Jak wykona\u0107 i interpretowa\u0107 -- kliknij aby rozwin\u0105\u0107", expanded=False):
        st.markdown("""
**Czym jest Conjoint?**
Conjoint (analiza l\u0105czna) mierzy, jak poszczeg\u00f3lne cechy produktu wp\u0142ywaj\u0105 na preferencje respondent\u00f3w. Wynikiem s\u0105 **u\u017cyteczno\u015bci cz\u0105stkowe** (part-worth utilities) oraz **wa\u017cno\u015b\u0107 atrybut\u00f3w**.

**Dwa dost\u0119pne warianty:**
- **Rating-based**: Respondenci oceniaj\u0105 profile produkt\u00f3w w skali (np. 1-10). Zmienna zale\u017cna = ocena.
- **CBC (Choice-Based)**: Respondenci wybieraj\u0105 mi\u0119dzy profilami. Zmienna zale\u017cna = 0/1 (czy profil zosta\u0142 wybrany).

**Jak interpretowa\u0107:**

| Wska\u017anik | Interpretacja |
|---|---|
| **Wa\u017cno\u015b\u0107 atrybutu (%)** | Im wy\u017csza, tym bardziej ten atrybut wp\u0142ywa na decyzje |
| **U\u017cyteczno\u015b\u0107 cz\u0105stkowa** | Dodatnia = preferowany poziom, ujemna = niepreferowan |
| **R\u00b2** | Odsetek wariancji ocen wyja\u015bniany przez model (rating) |
| **Pseudo R\u00b2** | Miara dopasowania modelu logit (CBC), >0.2 = dobre |

**Wymagania dotycz\u0105ce danych:**
- Rating: min. 30 respondent\u00f3w, zmienne atrybut\u00f3w kategoryczne lub liczbowe
- CBC: dane w formacie long (jeden wiersz = jeden profil-respondent)
        """)

    conj_method = st.radio("Wariant analizy:", ["Rating-based (OLS)", "CBC (Choice-Based Logit)"],
                            horizontal=True, key="conj_method")
    st.divider()

    col_a, col_b = st.columns(2)
    with col_a:
        if conj_method == "Rating-based (OLS)":
            conj_rating = st.selectbox("\U0001f3af Zmienna zale\u017cna (ocena profilu):",
                                        numeric_cols, format_func=lambda x: get_var_display_name(x, var_labels),
                                        key="conj_rating")
            conj_attrs = st.multiselect("\U0001f4e6 Atrybuty produktu (zmienne niezale\u017cne):",
                                         [c for c in visible_columns if c != conj_rating],
                                         format_func=lambda x: get_var_display_name(x, var_labels),
                                         key="conj_attrs_r")
        else:
            conj_choice = st.selectbox("\U0001f3af Zmienna wyboru (0=nie, 1=tak):",
                                        numeric_cols, format_func=lambda x: get_var_display_name(x, var_labels),
                                        key="conj_choice")
            conj_attrs = st.multiselect("\U0001f4e6 Atrybuty produkt\u00f3w:",
                                         [c for c in visible_columns if c != conj_choice],
                                         format_func=lambda x: get_var_display_name(x, var_labels),
                                         key="conj_attrs_c")
    with col_b:
        st.markdown("**Wskaz\u00f3wka:**")
        if conj_method == "Rating-based (OLS)":
            st.info("Wybierz zmienn\u0105 z ocen\u0105 profilu (np. 1-10) i zmienne opisuj\u0105ce atrybuty (kategorie lub liczby).")
        else:
            st.info("Wybierz zmienn\u0105 binarny\u0105 wyboru (1=wybrany profil) i zmienne atrybut\u00f3w. Dane musz\u0105 by\u0107 w formacie long.")

    if st.button("\u25b6\ufe0f Uruchom analiz\u0119 Conjoint", type="primary"):
        if not conj_attrs:
            st.error("Wybierz co najmniej jeden atrybut.")
        else:
            with st.spinner("Obliczanie..."):
                if conj_method == "Rating-based (OLS)":
                    res, err = run_conjoint_rating(df_raw, conj_rating, conj_attrs)
                else:
                    res, err = run_conjoint_cbc(df_raw, conj_choice, conj_attrs)
            if err:
                st.error(err)
            else:
                st.session_state.conjoint_results.append(res)
                st.success("\u2705 Analiza Conjoint uko\u0144czona!")

    for res in st.session_state.conjoint_results:
        if res.get('error'):
            st.error(res['error']); continue
        with st.expander(f"\U0001f4ca {res['method']} -- {len(res['attribute_vars'])} atrybut\u00f3w", expanded=True):
            # Summary metrics
            mc1, mc2, mc3, mc4 = st.columns(4)
            mc1.metric("N", f"{res['n']:,}")
            if 'r2' in res:
                mc2.metric("R\u00b2", f"{res['r2']:.4f}")
                mc3.metric("R\u00b2 skor.", f"{res['r2_adj']:.4f}")
                mc4.metric("p (F)", f"{res['p']:.4f} {'OK' if res['p'] < 0.05 else 'NS'}")
            elif 'pseudo_r2' in res:
                mc2.metric("Pseudo R\u00b2", f"{res['pseudo_r2']:.4f}")
                mc3.metric("LLR p", f"{res['llr_pvalue']:.4f}")

            # Importance chart
            st.markdown("**Wa\u017cno\u015b\u0107 atrybut\u00f3w:**")
            imp_df = pd.DataFrame(list(res['importance'].items()), columns=['Atrybut', 'Wa\u017cno\u015b\u0107 (%)'])
            imp_df['Etykieta'] = imp_df['Atrybut'].apply(lambda x: var_labels.get(x, x))
            imp_df = imp_df.sort_values('Wa\u017cno\u015b\u0107 (%)', ascending=True)
            fig_imp = px.bar(imp_df, x='Wa\u017cno\u015b\u0107 (%)', y='Etykieta', orientation='h',
                             color='Wa\u017cno\u015b\u0107 (%)', color_continuous_scale='Blues',
                             title='Wa\u017cno\u015b\u0107 atrybut\u00f3w (%)'),
            st.plotly_chart(fig_imp[0], use_container_width=True, key="pc_conj_imp")

            # Utilities per attribute
            st.markdown("**U\u017cyteczno\u015bci cz\u0105stkowe (part-worth utilities):**")
            for attr, utils in res['utilities'].items():
                if not utils: continue
                attr_lbl = var_labels.get(attr, attr)
                u_df = pd.DataFrame(list(utils.items()), columns=['Poziom', 'U\u017cyteczno\u015b\u0107'])
                u_df = u_df.sort_values('U\u017cyteczno\u015b\u0107', ascending=True)
                fig_u = px.bar(u_df, x='U\u017cyteczno\u015b\u0107', y='Poziom', orientation='h',
                               title=f"[{attr}] {attr_lbl}",
                               color='U\u017cyteczno\u015b\u0107', color_continuous_scale='RdYlGn',
                               color_continuous_midpoint=0)
                fig_u.add_vline(x=0, line_dash='dash', line_color='gray')
                st.plotly_chart(fig_u, use_container_width=True, key=f"pc_conj_util_{attr}")

    if st.session_state.conjoint_results:
        if st.button("\U0001f5d1\ufe0f Wyczy\u015b\u0107 wyniki Conjoint", type="secondary"):
            st.session_state.conjoint_results = []
            st.rerun()


# =============================================================
# MODUL: MAXDIFF
# =============================================================
elif menu == "\U0001f522 MaxDiff":
    module_header("\U0001f522", "MaxDiff", "Best-Worst Scaling \u2014 ranking wa\u017cno\u015bci element\u00f3w")

    with st.expander("\U0001f4d6 Jak wykona\u0107 i interpretowa\u0107 -- kliknij aby rozwin\u0105\u0107", expanded=False):
        st.markdown("""
**Czym jest MaxDiff?**
MaxDiff (Maximum Difference Scaling) mierzy wzgl\u0119dn\u0105 wa\u017cno\u015b\u0107/preferencj\u0119 element\u00f3w.
Respondenci wskazuj\u0105 **Najwa\u017cniejszy (Best)** i **Najmniej wa\u017cny (Worst)** spo\u015br\u00f3d ka\u017cdego zestawu pozycji.

**Format danych w pliku SPSS:**
Ka\u017cdy **zestaw** ma dwie kolumny:
- Kolumna **Best** -- kt\u00f3r\u0105 pozycj\u0119 wybrano jako najwa\u017cniejsz\u0105 (np. "Produkt A")
- Kolumna **Worst** -- kt\u00f3r\u0105 wybrano jako najmniej wa\u017cn\u0105

Przyk\u0142ad: `Zestaw1_Best`, `Zestaw1_Worst`, `Zestaw2_Best`, `Zestaw2_Worst`, ...

**Jak skonfigurowa\u0107:**
1. Podaj pary kolumn (Best / Worst) dla ka\u017cdego zestawu.
2. Podaj list\u0119 element\u00f3w (pozycji) wyst\u0119puj\u0105cych w tych kolumnach.
3. Kliknij **Uruchom MaxDiff**.

**Jak interpretowa\u0107 wyniki:**

| Wska\u017anik | Interpretacja |
|---|---|
| **B-W Score** | Liczba wybor\u00f3w Best minus Worst. Wy\u017cszy = preferowany |
| **B-W Score (%)** | B-W Score / N respondent\u00f3w * 100. Por\u00f3wnywalny mi\u0119dzy badaniami |
| **Wynik standaryzowany (0-100)** | Rescalowany do skali 0-100. Najlepszy element = 100 |
        """)

    st.divider()
    st.markdown("##### 1. Zdefiniuj pary kolumn Best/Worst")
    st.info("Dla ka\u017cdego zestawu wybierz kolumn\u0119 'Najwa\u017cniejszy' (Best) i 'Najmniej wa\u017cny' (Worst).")

    if 'maxdiff_pairs' not in st.session_state:
        st.session_state.maxdiff_pairs = [('', '')]

    pairs_to_remove = None
    new_pairs = []

    def _md_fmt(x):
        if not x:
            return '-- wybierz --'
        lbl = var_labels.get(x, '')
        return f"[{x}] {lbl}" if lbl else x

    for pi, (bc, wc) in enumerate(st.session_state.maxdiff_pairs):
        col_b, col_w, col_rm = st.columns([3, 3, 1])
        with col_b:
            sel_b = st.selectbox(f"Zestaw {pi+1} -- Najwa\u017cniejszy (Best):",
                                  [''] + list(df_raw.columns),
                                  index=list([''] + list(df_raw.columns)).index(bc) if bc in df_raw.columns else 0,
                                  format_func=_md_fmt,
                                  key=f"md_best_{pi}")
        with col_w:
            sel_w = st.selectbox(f"Zestaw {pi+1} -- Najmniej wa\u017cny (Worst):",
                                  [''] + list(df_raw.columns),
                                  index=list([''] + list(df_raw.columns)).index(wc) if wc in df_raw.columns else 0,
                                  format_func=_md_fmt,
                                  key=f"md_worst_{pi}")
        with col_rm:
            st.write("")
            if pi > 0 and st.button("\U0001f5d1\ufe0f", key=f"md_rm_{pi}"):
                pairs_to_remove = pi
        new_pairs.append((sel_b, sel_w))

    st.session_state.maxdiff_pairs = new_pairs
    if pairs_to_remove is not None:
        st.session_state.maxdiff_pairs.pop(pairs_to_remove)
        st.rerun()

    if st.button("\u2795 Dodaj zestaw", key="md_add_pair"):
        st.session_state.maxdiff_pairs.append(('', ''))
        st.rerun()

    st.divider()
    st.markdown("##### 2. Okre\u015bl pozycje (elementy) badania")

    # Auto-detect items from selected columns
    valid_pairs = [(b, w) for b, w in st.session_state.maxdiff_pairs if b and w and b in df_raw.columns and w in df_raw.columns]
    if valid_pairs:
        auto_items = set()
        for bc, wc in valid_pairs:
            auto_items.update(df_raw[bc].dropna().astype(str).unique())
            auto_items.update(df_raw[wc].dropna().astype(str).unique())
        auto_items = sorted(auto_items)
        st.caption(f"Automatycznie wykryto {len(auto_items)} unikalnych pozycji z wybranych kolumn.")
        md_items_raw = st.text_area(
            "Pozycje (jedna na wiersz):",
            value='\n'.join(auto_items),
            height=180,
            key="md_items",
            help="Mo\u017cesz edytowa\u0107 list\u0119 i zmienia\u0107 kolejno\u015b\u0107."
        )
        md_items = [x.strip() for x in md_items_raw.splitlines() if x.strip()]
    else:
        md_items_raw = st.text_area("Pozycje (jedna na wiersz):", height=150, key="md_items_manual")
        md_items = [x.strip() for x in md_items_raw.splitlines() if x.strip()]

    st.divider()
    st.markdown("##### 3. Nazwa analizy i uruchomienie")
    md_name = st.text_input("Nazwa analizy MaxDiff:", value="MaxDiff", key="md_name")

    if st.button("\u25b6\ufe0f Uruchom analiz\u0119 MaxDiff", type="primary"):
        if not valid_pairs:
            st.error("Wybierz co najmniej jedn\u0105 par\u0119 kolumn Best/Worst.")
        elif len(md_items) < 2:
            st.error("Podaj co najmniej 2 pozycje.")
        else:
            with st.spinner("Obliczanie wynik\u00f3w MaxDiff..."):
                df_scores = run_maxdiff(df_raw, valid_pairs, md_items)
            result_md = {
                'name': md_name,
                'pairs': valid_pairs,
                'items': md_items,
                'n_resp': len(df_raw),
                'n_tasks': len(valid_pairs),
                'scores': df_scores,
            }
            st.session_state.maxdiff_results.append(result_md)
            st.success(f"\u2705 MaxDiff uko\u0144czony! Przeanalizowano {len(valid_pairs)} zestaw\u00f3w, {len(md_items)} pozycji.")

    for res in st.session_state.maxdiff_results:
        with st.expander(f"\U0001f522 {res['name']} -- {res['n_tasks']} zestaw\u00f3w, {len(res['items'])} pozycji", expanded=True):
            df_s = res['scores']
            mc1, mc2, mc3 = st.columns(3)
            mc1.metric("N respondent\u00f3w", f"{res['n_resp']:,}")
            mc2.metric("Liczba zestaw\u00f3w", res['n_tasks'])
            mc3.metric("Pozycji", len(res['items']))

            st.markdown("**Wyniki MaxDiff -- ranking wa\u017cno\u015bci:**")

            _style_md = _make_style_md(len(df_s))

            st.dataframe(
                df_s.style.apply(_style_md, axis=1)
                    .format({'Best [N]': '{:.0f}', 'Worst [N]': '{:.0f}', 'Pokazano [N]': '{:.0f}',
                             'B-W Score': '{:.0f}', 'B-W Score (%)': '{:.1f}',
                             'Wynik standaryzowany (0-100)': '{:.1f}'}),
                use_container_width=True, hide_index=True
            )

            # Bar chart
            fig_md = px.bar(
                df_s.sort_values('Wynik standaryzowany (0-100)', ascending=True),
                x='Wynik standaryzowany (0-100)', y='Item', orientation='h',
                color='Wynik standaryzowany (0-100)', color_continuous_scale='Blues',
                title=f"MaxDiff -- Ranking wa\u017cno\u015bci: {res['name']}",
                text='Wynik standaryzowany (0-100)'
            )
            fig_md.update_traces(texttemplate='%{text:.1f}', textposition='outside')
            fig_md.update_layout(height=max(300, len(res['items']) * 35 + 80),
                                  coloraxis_showscale=False)
            st.plotly_chart(fig_md, use_container_width=True, key="pc_maxdiff_bar")

    if st.session_state.maxdiff_results:
        if st.button("\U0001f5d1\ufe0f Wyczy\u015b\u0107 wyniki MaxDiff", type="secondary"):
            st.session_state.maxdiff_results = []
            st.rerun()


# =============================================================
# MODUL: CHMURA SLOW
# =============================================================
# =============================================================
# MODUL: SKUPIENIA HIERARCHICZNE
# =============================================================
elif menu == "\U0001f3af Skupienia i Segmentacja":
    module_header("\U0001f3af", "Skupienia i Segmentacja", "Skupienia hierarchiczne (dendrogram) i segmentacja K-Means")
    tab_hc, tab_kmeans = st.tabs(["\U0001f333 Skupienia Hierarchiczne", "\U0001f3af Segmentacja K-Means"])

    with tab_hc:
        st.markdown("##### Analiza skupie\u0144 hierarchicznych")

        with st.expander("\U0001f4d6 Jak wykona\u0107 i interpretowa\u0107", expanded=False):
            st.markdown("""
    **Czym s\u0105 skupienia hierarchiczne?**
    Metoda grupowania respondent\u00f3w bez konieczno\u015bci wcze\u015bniejszego podania liczby grup (w odst\u0119pstwie do K-Means).
    Wynikiem jest **dendrogram** \u2014 drzewo podobie\u0144stw, kt\u00f3re pomaga dobra\u0107 optym aln\u0105 liczb\u0119 skupie\u0144.

    **Etapy analizy:**
    1. Wybierz zmienne numeryczne i metod\u0119 (\u0142\u0105czenia)
    2. Odczytaj dendrogram \u2014 szukaj du\u017cych "skok\u00f3w" na osi Y (odleg\u0142o\u015b\u0107)
    3. Ustaw progowy ci\u0119cie lub liczb\u0119 skupie\u0144
    4. Dodaj zmiennn\u0105 z przypisaniem do skupie\u0144 do bazy

    | Metoda | Zastosowanie |
    |---|---|
    | **Ward** | Minimalizuje wariancj\u0119 wewn\u0105trzgrupow\u0105 \u2014 zazwyczaj najlepsza |
    | **Complete** | U\u017cywa maksymalnej odleg\u0142o\u015bci \u2014 tworzy zborne skupienia |
    | **Average** | U\u017cywa \u015bredniej odleg\u0142o\u015bci \u2014 kompromis |
    | **Single** | U\u017cywa minimalnej odleg\u0142o\u015bci \u2014 podatna na efekt \u0142a\u0144cucha |
            """)

        col_hc1, col_hc2 = st.columns([2, 1])
        with col_hc1:
            hc_vars = st.multiselect(
                "Zmienne numeryczne do analizy:",
                numeric_cols,
                format_func=lambda x: get_var_display_name(x, var_labels),
                key="hc_vars"
            )
            hc_method = st.selectbox(
                "Metoda \u0142\u0105czenia:",
                ["ward", "complete", "average", "single"],
                key="hc_method",
                help="Ward: zalecana. Complete/Average: r\u00f3wnowa\u017cne. Single: nie polecana."
            )
            hc_metric = st.selectbox(
                "Miara odleg\u0142o\u015bci:",
                ["euclidean", "cosine", "correlation"],
                key="hc_metric",
                help="Ward wymaga euclidean."
            )
            if hc_method == "ward" and hc_metric != "euclidean":
                st.warning("Metoda Ward wymaga odleg\u0142o\u015bci euklidesowej \u2014 zmieniono automatycznie.")
                hc_metric = "euclidean"

        with col_hc2:
            hc_standardize = st.checkbox("Standaryzuj zmienne (Z-score)", value=True, key="hc_std")
            hc_max_obs = st.number_input(
                "Maks. respondent\u00f3w (wydajno\u015b\u0107):",
                min_value=50, max_value=5000, value=500, step=50, key="hc_maxobs",
                help="Dendrogram dla du\u017cych baz jest nieczytelny. Losowa pr\u00f3bka."
            )
            hc_n_clusters = st.slider("Liczba skupie\u0144 do wyci\u0119cia:", 2, 15, 3, key="hc_nclust")
            hc_var_name = st.text_input(
                "Nazwa nowej zmiennej skupie\u0144:",
                value="Skupienie_H",
                key="hc_varname"
            )

        if st.button("\u25b6\ufe0f Generuj dendrogram i skupienia", type="primary", key="hc_run"):
            if len(hc_vars) < 2:
                st.error("Wybierz co najmniej 2 zmienne.")
            else:
                from scipy.cluster.hierarchy import linkage, dendrogram, fcluster
                from scipy.spatial.distance import pdist
                matplotlib.use('Agg')
                import matplotlib.pyplot as plt

                df_hc = df_raw[hc_vars].dropna()
                n_used = min(len(df_hc), int(hc_max_obs))
                if len(df_hc) > n_used:
                    df_hc = df_hc.sample(n=n_used, random_state=42)
                    st.info(f"Losowa pr\u00f3bka: {n_used} z {len(df_raw[hc_vars].dropna())} kompletnych obserwacji.")

                if hc_standardize:
                    from sklearn.preprocessing import StandardScaler
                    X = StandardScaler().fit_transform(df_hc.values)
                else:
                    X = df_hc.values

                with st.spinner("Obliczanie skupie\u0144..."):
                    Z = linkage(X, method=hc_method,
                                metric=hc_metric if hc_method != 'ward' else 'euclidean')

                # \u2500\u2500 Dendrogram \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
                fig_dend, ax = plt.subplots(figsize=(14, 5))
                dendrogram(
                    Z, ax=ax,
                    truncate_mode='lastp', p=50,
                    leaf_rotation=90, leaf_font_size=8,
                    color_threshold=0.7 * max(Z[:, 2]),
                    above_threshold_color='#888888',
                )
                ax.set_title(f"Dendrogram skupie\u0144 hierarchicznych ({hc_method.title()}, n={n_used})",
                             fontsize=12, fontweight='bold')
                ax.set_xlabel("Indeks obserwacji lub liczba skupionych obiekt\u00f3w")
                ax.set_ylabel("Odleg\u0142o\u015b\u0107")
                ax.axhline(y=Z[-int(hc_n_clusters)+1, 2], color='#C00000',
                           linestyle='--', linewidth=1.5,
                           label=f"Ci\u0119cie: {hc_n_clusters} skupie\u0144")
                ax.legend(fontsize=9)
                plt.tight_layout()
                st.pyplot(fig_dend, use_container_width=True)

                # Download dendrogram
                buf_d = io.BytesIO()
                fig_dend.savefig(buf_d, format='png', dpi=150, bbox_inches='tight')
                buf_d.seek(0)
                plt.close(fig_dend)
                st.download_button(
                    "\u2b07\ufe0f Pobierz dendrogram (PNG)",
                    data=buf_d.getvalue(),
                    file_name=f"dendrogram_{hc_method}.png",
                    mime="image/png"
                )

                # \u2500\u2500 Assign clusters to full dataset \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
                if hc_standardize:
                    X_full = StandardScaler().fit_transform(df_raw[hc_vars].dropna().values)
                else:
                    X_full = df_raw[hc_vars].dropna().values

                Z_full = linkage(X_full, method=hc_method,
                                 metric=hc_metric if hc_method != 'ward' else 'euclidean')
                labels_full = fcluster(Z_full, hc_n_clusters, criterion='maxclust')

                idx_full = df_raw[hc_vars].dropna().index
                df_raw.loc[idx_full, hc_var_name] = labels_full
                df.loc[idx_full,     hc_var_name] = [f"Skupienie {c}" for c in labels_full]
                var_labels[hc_var_name] = f"Skupienia hierarchiczne ({hc_n_clusters} grup, {hc_method})"

                # \u2500\u2500 Cluster sizes \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
                sizes = pd.Series(labels_full).value_counts().sort_index()
                sizes_df = pd.DataFrame({
                    "Skupienie": [f"Skupienie {i}" for i in sizes.index],
                    "N": sizes.values,
                    "%": (sizes.values / sizes.sum() * 100).round(1),
                })

                # Profile: cluster means per variable
                profile_df = df_raw.loc[idx_full, hc_vars + [hc_var_name]].copy()
                profile_df[hc_var_name] = profile_df[hc_var_name].astype(int)
                cluster_means = profile_df.groupby(hc_var_name)[hc_vars].mean().round(2)
                cluster_means.index = [f"Skupienie {i}" for i in cluster_means.index]

                result_entry = {
                    'method': hc_method,
                    'metric': hc_metric,
                    'n_clusters': int(hc_n_clusters),
                    'vars': hc_vars,
                    'var_name': hc_var_name,
                    'n_obs': len(idx_full),
                    'sizes': sizes_df,
                    'profile': cluster_means,
                    'Z': Z_full.tolist(),
                }
                st.session_state.hclust_results.append(result_entry)
                st.success(f"\u2705 Skupienia hierarchiczne obliczone! Zmienna `{hc_var_name}` dodana do bazy.")

        for res_hc in st.session_state.hclust_results:
            with st.expander(
                f"\U0001f333 {res_hc['var_name']} \u2014 {res_hc['n_clusters']} skupie\u0144 ({res_hc['method']})",
                expanded=True
            ):
                sc1, sc2, sc3 = st.columns(3)
                sc1.metric("N obserwacji", f"{res_hc['n_obs']:,}")
                sc2.metric("Skupie\u0144", res_hc['n_clusters'])
                sc3.metric("Metoda", res_hc['method'].title())

                st.markdown("**Wielko\u015b\u0107 skupie\u0144:**")
                st.dataframe(res_hc['sizes'], use_container_width=True, hide_index=True)

                st.markdown("**Profil skupie\u0144 (\u015brednie zmiennych):**")
                st.dataframe(res_hc['profile'].style.format("{:.2f}"), use_container_width=True)

                fig_bar = px.bar(
                    res_hc['sizes'], x="Skupienie", y="N",
                    color="Skupienie", title=f"Liczebno\u015b\u0107 skupie\u0144: {res_hc['var_name']}",
                    color_discrete_sequence=px.colors.qualitative.Set2
                )
                fig_bar.update_layout(showlegend=False, height=300)
                st.plotly_chart(fig_bar, use_container_width=True, key=f"pc_hclust_{res_hc.get('var_name','hc')}")

        if st.session_state.hclust_results:
            if st.button("\U0001f5d1\ufe0f Wyczy\u015b\u0107 wyniki skupie\u0144", type="secondary", key="hc_clear"):
                st.session_state.hclust_results = []
                st.rerun()

    # \u2500\u2500 TAB: Segmentacja K-Means \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
    with tab_kmeans:
        st.markdown("##### Segmentacja K-Means")
        st.info("Pogrupuj respondent\u00f3w metod\u0105 K-Means. Nowa zmienna segmentacyjna zostanie dodana do bazy.")
        seg_vars = st.multiselect("Zmienne numeryczne:", numeric_cols,
                                   format_func=lambda x: get_var_display_name(x, var_labels),
                                   key="seg_vars_mod")
        k_clusters = st.slider("Liczba segment\u00f3w (K):", 2, 10, 3, key="k_clusters_mod")
        seg_name = st.text_input("Nazwa zmiennej segmentacyjnej:",
                                  value=f"Segmentacja_{len(st.session_state.segmentations) + 1}",
                                  key="seg_name_mod")
        if st.button("\u25b6\ufe0f Wykonaj segmentacj\u0119 K-Means", type="primary", key="seg_run_mod"):
            if len(seg_vars) < 2:
                st.error("Wybierz co najmniej 2 zmienne.")
            else:
                st.session_state.segmentations.append(
                    {'vars': seg_vars, 'k': k_clusters, 'name': seg_name})
                st.success(f"\u2705 Segmentacja `{seg_name}` utworzona.")
                st.rerun()

        if st.session_state.segmentations:
            st.divider()
            st.markdown("**Zdefiniowane segmentacje:**")
            to_del = None
            for i, seg in enumerate(st.session_state.segmentations):
                c1, c2 = st.columns([5, 1])
                c1.write(f"- `{seg['name']}` \u2014 {seg['k']} grup, bazuje na {len(seg['vars'])} zmiennych")
                if c2.button("\U0001f5d1\ufe0f", key=f"del_seg_mod_{i}"):
                    to_del = i
            if to_del is not None:
                st.session_state.segmentations.pop(to_del)
                st.rerun()

            # Show cluster profiles for each segmentation
            for seg in st.session_state.segmentations:
                if seg['name'] in df_raw.columns:
                    with st.expander(f"\U0001f3af Profil: `{seg['name']}`", expanded=False):
                        profile = df_raw.groupby(seg['name'])[seg['vars']].mean().round(2)
                        profile.index = [f"Segment {int(i)}" for i in profile.index]
                        sizes = df_raw[seg['name']].value_counts().sort_index()
                        sizes.index = [f"Segment {int(i)}" for i in sizes.index]
                        col_p1, col_p2 = st.columns([3, 1])
                        col_p1.dataframe(profile.style.format("{:.2f}"), use_container_width=True)
                        col_p2.dataframe(sizes.rename("N"), use_container_width=True)
                        fig_seg = px.bar(
                            sizes.reset_index(),
                            x='index', y=seg['name'],
                            title=f"Liczebno\u015b\u0107 segment\u00f3w: {seg['name']}",
                            color='index',
                            color_discrete_sequence=px.colors.qualitative.Set2
                        )
                        fig_seg.update_layout(showlegend=False, height=280,
                                               xaxis_title="Segment", yaxis_title="N")
                        st.plotly_chart(fig_seg, use_container_width=True, key=f"pc_seg_{seg.get('name','s')}")


elif menu == "\u2601\ufe0f Chmura S\u0142\u00f3w":
    module_header("\u2601\ufe0f", "Chmura S\u0142\u00f3w", "Wizualizacja odpowiedzi otwartych \u2014 eksport PNG/JPG")

    # Check for wordcloud availability
    try:
        from wordcloud import WordCloud
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        _wc_available = True
    except ImportError:
        _wc_available = False
        st.error(
            "Biblioteka `wordcloud` nie jest zainstalowana. "
            "Uruchom: `pip install wordcloud` i restart\u01b3 aplikacj\u0119."
        )

    if _wc_available:
        # \u2500\u2500 Column selector \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
        # For SPSS: also include numeric variables that have value labels
        # (they represent coded categorical responses, not open text)
        if is_spss:
            val_label_cols = [c for c in visible_columns
                              if c in meta_orig.variable_value_labels
                              and meta_orig.variable_value_labels[c]]
            text_cols = list(dict.fromkeys(
                [c for c in visible_columns if df_raw[c].dtype == object]
                + val_label_cols
            ))
        else:
            text_cols = [c for c in visible_columns if df_raw[c].dtype == object]

        if not text_cols:
            st.warning("Brak zmiennych tekstowych w bazie danych.")
        else:
            st.info(
                "Wybierz zmienn\u0105 z odpowiedziami otwartymi, dostosuj wygl\u0105d i wygeneruj chmur\u0119 s\u0142\u00f3w. "
                "Gotow\u0105 grafik\u0119 mo\u017cesz pobra\u0107 jako PNG lub JPG."
            )

            col_cfg1, col_cfg2 = st.columns([2, 1])

            with col_cfg1:
                wc_var = st.selectbox(
                    "Pytanie otwarte (zmienna tekstowa):",
                    text_cols,
                    format_func=lambda x: get_var_display_name(x, var_labels),
                    key="wc_var"
                )

                # Stop words
                st.markdown("**Stop words** \u2014 s\u0142owa do wykluczenia:")
                default_stopwords = (
                    "i w z na do nie to a ale o jak tak si\u0119 co ten ta tego tej "
                    "jest by\u0142 by\u0142a by\u0142o s\u0105 b\u0119d\u0105 ma mam masz ma\u0107 mie\u0107 "
                    "one oni one nas nam ich im kt\u00f3ry kt\u00f3ra kt\u00f3re tego "
                    "tego tej temu tym te ten tego dla ze przy po czy "
                    "bardzo wi\u0119c jednak tylko jeszcze ju\u017c bo bo\u017c gdy "
                    "mi\u0119dzy przez mo\u017ce mo\u017cna po za przed"
                )
                stopwords_raw = st.text_area(
                    "Wpisz s\u0142owa oddzielone spacjami lub przecinkami:",
                    value=default_stopwords,
                    height=100,
                    key="wc_stopwords",
                    help="Te s\u0142owa nie pojawi\u0105 si\u0119 w chmurze. Mo\u017cesz usun\u0105\u0107 lub doda\u0107 w\u0142asne."
                )

                # Case handling
                wc_lowercase = st.checkbox(
                    "Zamie\u0144 na ma\u0142e litery (Warszawa = warszawa)",
                    value=True, key="wc_lower"
                )

                # Min word frequency
                wc_min_freq = st.slider(
                    "Minimalna cz\u0119sto\u015b\u0107 wyst\u0105pienia s\u0142owa:",
                    min_value=1, max_value=20, value=1, key="wc_minfreq"
                )

                # Max words
                wc_max_words = st.slider(
                    "Maksymalna liczba s\u0142\u00f3w w chmurze:",
                    min_value=10, max_value=300, value=100, key="wc_maxwords"
                )

            with col_cfg2:
                st.markdown("**Wygl\u0105d chmury:**")

                wc_bg = st.color_picker("Kolor t\u0142a:", value="#FFFFFF", key="wc_bg")
                wc_width  = st.number_input("Szeroko\u015b\u0107 (px):", min_value=400, max_value=3000,
                                             value=1200, step=100, key="wc_w")
                wc_height = st.number_input("Wysoko\u015b\u0107 (px):", min_value=200, max_value=2000,
                                             value=600, step=100, key="wc_h")

                PALETTES = {
                    "Niebieski (domowy)":  "Blues",
                    "Czerwony":            "Reds",
                    "Zielony":             "Greens",
                    "Fioletowy":           "Purples",
                    "Ciep\u0142e kolory":  "YlOrRd",
                    "Ch\u0142odne kolory": "cool",
                    "T\u0119czowa":        "rainbow",
                    "Czarno-bia\u0142a":   "Greys",
                    "Niebiesko-zielona":   "GnBu",
                    "Czerwono-niebieska":  "RdBu",
                }
                wc_palette_name = st.selectbox(
                    "Paleta kolor\u00f3w:",
                    list(PALETTES.keys()),
                    key="wc_palette"
                )
                wc_palette = PALETTES[wc_palette_name]

                fmt_choice = st.radio(
                    "Format pobierania:",
                    ["PNG", "JPG"],
                    horizontal=True, key="wc_fmt"
                )

            st.divider()

            if st.button("\u25b6\ufe0f Generuj chmur\u0119 s\u0142\u00f3w", type="primary",
                         use_container_width=True, key="wc_generate"):

                # Build text corpus \u2014 use value labels for SPSS coded variables
                raw_series = df_raw[wc_var].dropna()

                if is_spss:
                    # Check if this variable has value labels (coded numeric/categorical)
                    spss_val_labels = {}
                    spss_val_labels.update(meta_orig.variable_value_labels.get(wc_var, {}))
                    spss_val_labels.update(st.session_state.custom_val_labels.get(wc_var, {}))

                    if spss_val_labels:
                        # Map each code to its label; fall back to str(code) if no label
                        def _map_label(v):
                            # SPSS codes may be float (1.0) or int
                            lbl = (spss_val_labels.get(v)
                                   or spss_val_labels.get(int(v) if isinstance(v, float) and v == int(v) else v)
                                   or spss_val_labels.get(str(int(v)) if isinstance(v, float) and v == int(v) else str(v))
                                   or str(v))
                            return str(lbl)
                        texts = raw_series.map(_map_label).tolist()
                    else:
                        # Open-ended text variable \u2014 use as-is
                        texts = raw_series.astype(str).tolist()
                else:
                    texts = raw_series.astype(str).tolist()

                corpus = " ".join(texts)

                if wc_lowercase:
                    corpus = corpus.lower()

                # Build stop words set
                raw_sw = stopwords_raw.replace(',', ' ').split()
                stop_set = {w.strip().lower() for w in raw_sw if w.strip()}

                if not corpus.strip():
                    st.error("Wybrana zmienna jest pusta \u2014 brak tekstu do analizy.")
                else:
                    try:
                        import matplotlib.cm as cm

                        # \u2500\u2500 1. Tokenize corpus \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
                        token_re = re.compile(
                            r"[\w\u0104\u0105\u0106\u0107\u0118\u0119"
                            r"\u0141\u0142\u0143\u0144\u00d3\u00f3"
                            r"\u015a\u015b\u0179\u017a\u017b\u017c]+",
                            re.UNICODE
                        )
                        tokens = token_re.findall(corpus)

                        # \u2500\u2500 2. Count frequencies (excluding stop words) \u2500
                        freq = {}
                        for tok in tokens:
                            w = tok.lower() if wc_lowercase else tok
                            if len(w) < 2:
                                continue
                            if w.lower() in stop_set:
                                continue
                            freq[w] = freq.get(w, 0) + 1

                        # \u2500\u2500 3. Apply minimum frequency filter \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
                        if wc_min_freq > 1:
                            freq = {w: c for w, c in freq.items()
                                    if c >= int(wc_min_freq)}

                        if not freq:
                            st.warning(
                                "Po zastosowaniu filtr\u00f3w nie pozosta\u0142o \u017cadne s\u0142owo. "
                                "Spr\u00f3buj zmniejszy\u0107 minimaln\u0105 cz\u0119sto\u015b\u0107 lub "
                                "skr\u00f3ci\u0107 list\u0119 stop words."
                            )
                        else:
                            # \u2500\u2500 4. Build WordCloud from frequency dict \u2500\u2500
                            wc_obj = WordCloud(
                                width=int(wc_width),
                                height=int(wc_height),
                                background_color=wc_bg,
                                colormap=wc_palette,
                                max_words=int(wc_max_words),
                                min_font_size=8,
                                max_font_size=None,
                                min_word_length=2,
                                collocations=False,  # already filtered
                                relative_scaling=0.5,
                                prefer_horizontal=0.9,
                            ).generate_from_frequencies(freq)

                            # \u2500\u2500 5. Render with matplotlib \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
                            fig, ax = plt.subplots(
                                figsize=(int(wc_width) / 100, int(wc_height) / 100),
                                dpi=100
                            )
                            ax.imshow(wc_obj, interpolation='bilinear')
                            ax.axis('off')
                            fig.patch.set_facecolor(wc_bg)
                            plt.tight_layout(pad=0)

                            # Show in app
                            st.pyplot(fig, use_container_width=True)

                            # \u2500\u2500 6. Download \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
                            buf = io.BytesIO()
                            save_fmt = fmt_choice.lower()
                            save_kwargs = {
                                'format':       'jpeg' if save_fmt == 'jpg' else save_fmt,
                                'dpi':          150,
                                'bbox_inches':  'tight',
                                'pad_inches':   0,
                            }
                            fig.savefig(buf, **save_kwargs)
                            buf.seek(0)
                            plt.close(fig)

                            fname = f"chmura_slow_{wc_var}.{save_fmt}"
                            mime  = "image/jpeg" if save_fmt == "jpg" else "image/png"
                            st.download_button(
                                label=f"\u2b07\ufe0f Pobierz chmur\u0119 ({fmt_choice})",
                                data=buf.getvalue(),
                                file_name=fname,
                                mime=mime,
                                use_container_width=True,
                            )

                            # \u2500\u2500 7. Frequency table \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
                            with st.expander(
                                "\U0001f4ca Cz\u0119sto\u015b\u0107 s\u0142\u00f3w w chmurze",
                                expanded=False
                            ):
                                freq_df = pd.DataFrame(
                                    sorted(freq.items(), key=lambda x: x[1], reverse=True),
                                    columns=["S\u0142owo", "Liczba wyst\u0105pie\u0144"]
                                )
                                freq_df.index = range(1, len(freq_df) + 1)
                                st.dataframe(freq_df, use_container_width=True, height=300)

                    except Exception as _wc_err:
                        st.error(f"B\u0142\u0105d generowania chmury: {_wc_err}")
                        st.exception(_wc_err)


elif menu == "\U0001f4be Eksport do Excela":
    module_header("\U0001f4be", "Eksport do Excela", "Raport analityczny, wykresy, baza danych, spis tre\u015bci")

    # \u2500\u2500 Separate standalone DB download (still available) \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
    with st.expander("\U0001f4c2 Osobny plik z baz\u0105 danych", expanded=False):
        col1, col2 = st.columns(2)
        col1.info("**Baza z etykietami** -- warto\u015bci kod\u00f3w zast\u0105pione tekstem (np. 1 \u2192 'Kobieta').")
        col2.info("**Baza surowa** -- oryginalne warto\u015bci liczbowe. Wiersz 1: nazwy, Wiersz 2: etykiety.")
        if st.button("\U0001f4e5 Pobierz osobny plik z baz\u0105 danych", use_container_width=True):
            with st.spinner("Generowanie..."):
                db_data = export_db_to_excel(df_raw, df, var_labels)
            fname = "Baza_Danych_Excel.xlsx" if is_excel else "Baza_Danych_SPSS.xlsx"
            st.download_button("\u2b07\ufe0f Pobierz " + fname, data=db_data,
                               file_name=fname,
                               mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                               use_container_width=True)

    st.divider()
    st.markdown("### \U0001f4cb Raport analityczny (Excel)")

    any_results = (
        any(st.session_state.results.get(g) for g in ['czestosci', 'krzyzowe', 'srednie', 'opisowe', 'korelacje']) or
        bool(st.session_state.regression_results) or
        bool(st.session_state.anova_results) or
        bool(st.session_state.factor_results) or
        bool(st.session_state.matrix_results) or
        bool(st.session_state.conjoint_results) or
        bool(st.session_state.maxdiff_results)
    )
    if not any_results:
        st.warning("Brak wynik\u00f3w do eksportu. Przejd\u017a do modu\u0142\u00f3w analitycznych i wygeneruj tabele.")
    else:
        # Summary of available results
        available = []
        for grp, name in [('czestosci', 'Cz\u0119sto\u015bci'), ('krzyzowe', 'Krzy\u017cowe'), ('srednie', '\u015arednie'),
                          ('opisowe', 'Opisowe'), ('korelacje', 'Korelacje')]:
            if st.session_state.results.get(grp):
                available.append(f"\u2705 {name} ({len(st.session_state.results[grp])} tabel)")
        if st.session_state.matrix_results:
            available.append(f"\u2705 Pytania matrycowe ({len(st.session_state.matrix_results)} pyta\u0144)")
        if st.session_state.regression_results:
            available.append(f"\u2705 Regresja OLS ({len([r for r in st.session_state.regression_results if 'error' not in r])} blok\u00f3w)")
        if st.session_state.anova_results:
            available.append(f"\u2705 ANOVA ({len(st.session_state.anova_results)} analiz)")
        if st.session_state.factor_results:
            available.append(f"\u2705 Analiza Czynnikowa ({len(st.session_state.factor_results)} analiz)")
        if st.session_state.conjoint_results:
            available.append(f"\u2705 Conjoint ({len(st.session_state.conjoint_results)} analiz)")
        if st.session_state.maxdiff_results:
            available.append(f"\u2705 MaxDiff ({len(st.session_state.maxdiff_results)} analiz)")
        st.success("**Gotowe do eksportu:**\n" + " \u00b7 ".join(available))

        st.markdown("**Opcje eksportu:**")
        opt_col1, opt_col2, opt_col3 = st.columns(3)

        # Chart option
        has_freq = bool(st.session_state.results.get('czestosci'))
        add_freq_charts = opt_col1.checkbox(
            "\U0001f4ca Wykresy do tabel cz\u0119sto\u015bci",
            value=False,
            key="export_add_charts",
            help="Wstawia natywne wykresy Excela (edytowalne) obok ka\u017cdej tabeli cz\u0119sto\u015bci.",
            disabled=not has_freq,
        )

        # DB options
        incl_db_labeled = opt_col2.checkbox(
            "\U0001f4c2 Baza danych z etykietami",
            value=False,
            key="export_db_labeled",
            help="Dodaje do pliku arkusz 'Baza z etykietami' (warto\u015bci tekstowe) zaraz po Spisie Tre\u015bci.",
        )
        incl_db_raw = opt_col3.checkbox(
            "\U0001f4cb Baza danych surowa",
            value=False,
            key="export_db_raw",
            help="Dodaje arkusz 'Baza surowa (numeryczna)' z oryginalnymi kodami liczbowymi.",
        )

        if st.button("\U0001f4ca Generuj pe\u0142ny raport analityczny", type="primary", use_container_width=True):
            with st.spinner("Generowanie pliku Excel... To mo\u017ce chwil\u0119 potrwa\u0107."):
                output = io.BytesIO()
                try:
                    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
                        # 1. ToC -- always first tab
                        toc_ws = writer.book.add_worksheet('\U0001f4cb Spis Tre\u015bci')
                        toc_ws.set_tab_color('#1F4E79')

                        # 2. DB sheets -- right after ToC (second/third position)
                        if incl_db_labeled:
                            write_db_sheet(writer, 'Baza z etykietami',
                                           df, var_labels, '#1F4E79')
                        if incl_db_raw:
                            write_db_sheet(writer, 'Baza surowa (numeryczna)',
                                           df_raw, var_labels, '#2E75B6')

                        # 3. Analytical results sheets
                        sheet_map = {}
                        for grp, s_name in [('czestosci', 'Cz\u0119sto\u015bci'), ('krzyzowe', 'Krzy\u017cowe'),
                                            ('srednie', '\u015arednie'), ('opisowe', 'Opisowe'), ('korelacje', 'Korelacje')]:
                            if st.session_state.results.get(grp):
                                _charts = add_freq_charts if grp == 'czestosci' else False
                                row_map = export_tables_to_sheet(
                                    writer, s_name, st.session_state.results[grp], var_labels,
                                    add_charts=_charts
                                )
                                sheet_map[s_name] = row_map

                        if st.session_state.matrix_results:
                            export_matrix_to_excel(writer, st.session_state.matrix_results, var_labels)

                        valid_reg = [r for r in st.session_state.regression_results if 'error' not in r]
                        if valid_reg:
                            export_regression_to_excel(writer, valid_reg, var_labels)
                        if st.session_state.anova_results:
                            export_anova_to_excel(writer, st.session_state.anova_results, var_labels)
                        if st.session_state.factor_results:
                            export_factor_to_excel(writer, st.session_state.factor_results, var_labels)
                        if st.session_state.conjoint_results:
                            valid_conj = [r for r in st.session_state.conjoint_results if not r.get('error')]
                            if valid_conj:
                                export_conjoint_to_excel(writer, valid_conj, var_labels)
                        if st.session_state.maxdiff_results:
                            export_maxdiff_to_excel(writer, st.session_state.maxdiff_results, var_labels)

                        # 4. Fill ToC content last (all sheet_maps ready)
                        export_toc_sheet(
                            writer, st.session_state.results, st.session_state.matrix_results,
                            var_labels, sheet_map,
                            regression_results=st.session_state.regression_results,
                            anova_results=st.session_state.anova_results,
                            factor_results=st.session_state.factor_results,
                            conjoint_results=st.session_state.conjoint_results,
                            maxdiff_results=st.session_state.maxdiff_results,
                            pre_created_ws=toc_ws,
                        )

                    st.download_button(
                        "\u2b07\ufe0f Pobierz Raport_Analityczny.xlsx",
                        data=output.getvalue(),
                        file_name="Raport_Analityczny.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        use_container_width=True,
                    )
                except Exception as e:
                    st.error(f"B\u0142\u0105d podczas generowania pliku: {e}")
                    st.exception(e)

# =============================================================
# MODUL: EKSPORT DO POWERPOINT
# =============================================================
elif menu == "\U0001f4ca Eksport do PowerPoint":
    module_header("\U0001f4ca", "Eksport do PowerPoint", "Edytowalne wykresy kolumnowe z cz\u0119sto\u015bci i tabel krzy\u017cowych")
    st.info(
        "Generuje plik PowerPoint z edytowalnymi wykresami kolumnowymi. "
        "Ka\u017cdy wykres jest zagnie\u017cd\u017cony jako natywny obiekt PPT z w\u0142asn\u0105 "
        "tabel\u0105 danych \u2014 mo\u017cna go edytowa\u0107 bezpo\u015brednio w PowerPoint. "
        "Eksportowane s\u0105 wy\u0142\u0105cznie wyniki tabel cz\u0119sto\u015bci i tabel krzy\u017cowych."
    )

    freq_res  = st.session_state.results.get('czestosci', {})
    cross_res = st.session_state.results.get('krzyzowe', {})

    if not freq_res and not cross_res:
        st.warning(
            "Brak wynik\u00f3w do eksportu. Wygeneruj tabele cz\u0119sto\u015bci lub krzy\u017cowe "
            "w module **Analizy i Tabele**."
        )
    else:
        n_freq  = len(freq_res)
        n_cross = len(cross_res)

        st.markdown("**Dost\u0119pne wyniki:**")
        mc1, mc2 = st.columns(2)
        mc1.metric("\U0001f4c8 Tablice cz\u0119sto\u015bci", n_freq)
        mc2.metric("\U0001f500 Tabele krzy\u017cowe", n_cross)

        st.divider()
        st.markdown("**Opcje prezentacji:**")

        # \u2500\u2500 Szablon slajd\u00f3w \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
        pptx_template_file = st.file_uploader(
            "\U0001f4cb Szablon slajd\u00f3w (.pptx) \u2014 opcjonalny",
            type=["pptx"], key="ppt_template_file",
            help="Wgraj plik .pptx z gotowym layoutem (t\u0142o, logo, stopka). "
                 "Ka\u017cdy wykres zostanie dodany do nowego slajdu z tego szablonu."
        )
        if pptx_template_file:
            st.success("\u2705 Szablon slajd\u00f3w wczytany.")
        _crtx_style = None  # crtx not used

        st.divider()
        oc1, oc2, oc3, oc4 = st.columns(4)
        with oc1:
            ppt_metric = st.radio(
                "Warto\u015b\u0107 na wykresie:",
                ["Procent [%]", "Liczebno\u015b\u0107 [N]"],
                key="ppt_metric",
                help="Cz\u0119sto\u015bci: kolumna Procent lub Liczebno\u015b\u0107."
            )
        with oc2:
            ppt_title_prefix = st.text_input(
                "Prefiks tytu\u0142u slajdu:", value="",
                key="ppt_prefix", placeholder="np. Badanie 2025 |"
            )
        with oc3:
            ppt_show_base = st.checkbox(
                "Poka\u017c baz\u0119 (N) w tytu\u0142ach",
                value=True, key="ppt_base"
            )
        with oc4:
            _PALETTES = {
                "Niebieski (domy\u015blny)": ("#2E75B6", "#1F4E79"),
                "Zielony":                   ("#375623", "#1E3A12"),
                "Czerwony":                  ("#C00000", "#7B0000"),
                "Pomara\u0144czowy":         ("#E36C09", "#8E3F00"),
                "Fioletowy":                 ("#7030A0", "#3D1960"),
                "Szary (korporacyjny)":      ("#595959", "#262626"),
                "Granatowy":                 ("#003087", "#001B55"),
                "Z\u0142oty":               ("#C09000", "#7A5A00"),
            }
            ppt_palette_name = st.selectbox(
                "Paleta kolor\u00f3w:",
                list(_PALETTES.keys()), key="ppt_palette"
            )
            _bar_hex, _title_hex = _PALETTES[ppt_palette_name]

        st.divider()
        sel_freq_keys  = []
        sel_cross_keys = []

        if freq_res:
            st.markdown("**\U0001f4c8 Tablice cz\u0119sto\u015bci \u2014 wybierz do eksportu:**")
            col_a, col_b = st.columns(2)
            for i, key in enumerate(freq_res.keys()):
                col = col_a if i % 2 == 0 else col_b
                lbl = var_labels.get(key, key)
                display = f"{key} \u2014 {lbl[:50]}" if lbl != key else key[:60]
                if col.checkbox(display, value=True, key=f"ppt_fsel_{key}"):
                    sel_freq_keys.append(key)

        if cross_res:
            st.markdown("**\U0001f500 Tabele krzy\u017cowe \u2014 wybierz do eksportu:**")
            col_c, col_d = st.columns(2)
            for i, key in enumerate(cross_res.keys()):
                col = col_c if i % 2 == 0 else col_d
                if col.checkbox(key[:70], value=True, key=f"ppt_csel_{key}"):
                    sel_cross_keys.append(key)

        n_selected = len(sel_freq_keys) + len(sel_cross_keys)
        st.caption(f"Wybrano {n_selected} wykres\u00f3w do eksportu.")

        if st.button(
            f"\U0001f4ca Generuj plik PowerPoint ({n_selected} slajd\u00f3w)",
            type="primary", use_container_width=True,
            key="ppt_generate", disabled=(n_selected == 0)
        ):
            try:
                from pptx import Presentation
                from pptx.util import Inches, Pt
                from pptx.chart.data import ChartData
                from pptx.enum.chart import XL_CHART_TYPE
                from pptx.dml.color import RGBColor
                from pptx.util import Emu
                import lxml.etree as _etree

                def _hex_to_rgb(h):
                    if isinstance(h, RGBColor):
                        return h   # already an RGBColor (from crtx parser)
                    h = h.lstrip("#")
                    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

                BAR_COLOR   = _hex_to_rgb(_bar_hex)
                TITLE_COLOR = _hex_to_rgb(_title_hex)
                WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

                import colorsys as _cs

                def _palette_series(n, base_hex):
                    # If crtx provides colors, cycle through them
                    if _crtx_style and _crtx_style['series_colors']:
                        colors = _crtx_style['series_colors']
                        return [colors[i % len(colors)] for i in range(n)]
                    bh = base_hex if isinstance(base_hex, str) else f"#{base_hex[0]:02X}{base_hex[1]:02X}{base_hex[2]:02X}"
                    bh = bh.lstrip('#')
                    r,g,b = int(bh[:2],16)/255, int(bh[2:4],16)/255, int(bh[4:],16)/255
                    h,s,v = _cs.rgb_to_hsv(r,g,b)
                    colors = []
                    for i in range(n):
                        vi = max(0.25, v - i * (v-0.25)/(max(n-1,1)))
                        si = max(0.15, s - i * (s-0.15)/(max(n-1,1)))
                        rr,gg,bb = _cs.hsv_to_rgb(h, si, vi)
                        colors.append(RGBColor(int(rr*255), int(gg*255), int(bb*255)))
                    return colors

                # \u2500\u2500 Presentation base \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
                if pptx_template_file is not None:
                    import io as _io
                    prs = Presentation(_io.BytesIO(pptx_template_file.getvalue()))
                    # Use the first slide layout that looks blank, or fall back to [6]
                    blank_layout = None
                    for lay in prs.slide_layouts:
                        if len(lay.placeholders) == 0:
                            blank_layout = lay
                            break
                    if blank_layout is None:
                        blank_layout = prs.slide_layouts[min(6, len(prs.slide_layouts)-1)]
                else:
                    prs = Presentation()
                    prs.slide_width  = Inches(13.33)
                    prs.slide_height = Inches(7.5)
                    blank_layout = prs.slide_layouts[6]

                # \u2500\u2500 Font from crtx \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
                _chart_font = (_crtx_style.get('font_name') if _crtx_style else None) or None

                # \u2500\u2500 Gridlines setting from crtx \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
                _show_gridlines = (_crtx_style.get('gridlines', False)
                                   if _crtx_style else False)
                _gridline_color = (_crtx_style.get('gridline_color')
                                   if _crtx_style else None)

                def _slide_base(title_text, subtitle=""):
                    """Create a slide with title bar and return (slide, chart_top)."""
                    slide = prs.slides.add_slide(blank_layout)

                    if pptx_template_file is None:
                        # No template \u2014 draw custom white bg + blue/crtx title bar
                        bg = slide.shapes.add_shape(1, Inches(0), Inches(0),
                                                    prs.slide_width, prs.slide_height)
                        bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
                        bg.line.fill.background()

                        tb = slide.shapes.add_shape(1, Inches(0), Inches(0),
                                                    prs.slide_width, Inches(1.1))
                        tb.fill.solid(); tb.fill.fore_color.rgb = TITLE_COLOR
                        tb.line.fill.background()

                        ttf = slide.shapes.add_textbox(Inches(0.3), Inches(0.15),
                                                        Inches(12.5), Inches(0.8))
                        ttf.text_frame.word_wrap = True
                        tp = ttf.text_frame.paragraphs[0]
                        tp.text = title_text
                        tp.font.size = Pt(20); tp.font.bold = True
                        tp.font.color.rgb = WHITE
                        if _chart_font:
                            try: tp.font.name = _chart_font
                            except Exception: pass

                        chart_top = Inches(1.2)
                        if subtitle:
                            stf = slide.shapes.add_textbox(Inches(0.3), Inches(1.1),
                                                            Inches(12.5), Inches(0.3))
                            sp = stf.text_frame.paragraphs[0]
                            sp.text = subtitle; sp.font.size = Pt(11)
                            sp.font.color.rgb = RGBColor(0x59, 0x59, 0x59)
                            if _chart_font:
                                try: sp.font.name = _chart_font
                                except Exception: pass
                            chart_top = Inches(1.45)
                    else:
                        # Template slide \u2014 populate existing placeholders
                        chart_top = Inches(1.5)
                        for ph in slide.placeholders:
                            if ph.placeholder_format.idx == 0:  # title
                                ph.text = title_text
                                if _chart_font:
                                    try: ph.text_frame.paragraphs[0].font.name = _chart_font
                                    except Exception: pass
                                try:
                                    chart_top = Inches(
                                        (ph.top + ph.height) / 914400 + 0.1
                                    )
                                except Exception:
                                    chart_top = Inches(1.5)
                                break
                        if subtitle:
                            for ph in slide.placeholders:
                                if ph.placeholder_format.idx == 1:
                                    ph.text = subtitle
                                    break

                    return slide, chart_top

                def _fmt_numfmt(chart, use_pct):
                    """Apply Y-axis and gridline settings, optionally from crtx."""
                    try:
                        vax = chart.value_axis
                        vax.has_major_gridlines = _show_gridlines
                        if _show_gridlines and _gridline_color:
                            try:
                                vax.major_gridlines.format.line.color.rgb = _hex_to_rgb(_gridline_color)
                            except Exception:
                                pass
                        # Hide axis tick labels (values shown in data labels)
                        vax.tick_labels.font.size = Pt(1)
                        vax.tick_labels.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
                        vax.format.line.color.rgb = RGBColor(0xFF,0xFF,0xFF)
                    except Exception:
                        pass

                def _apply_dls(dls, color):
                    """Apply data label style, respecting crtx font if available."""
                    dls.show_value = True; dls.show_category_name = False
                    dls.font.size = Pt(11); dls.font.bold = True
                    dls.font.color.rgb = color
                    if _chart_font:
                        try:
                            dls.font.name = _chart_font
                        except Exception:
                            pass

                def _add_chart_slide(prs, title_text, categories, values,
                                     subtitle="", is_pct=False):
                    """Single-series frequency chart."""
                    slide, chart_top = _slide_base(title_text, subtitle)

                    cd = ChartData()
                    cd.categories = [str(c)[:40] for c in categories]
                    data_vals = []
                    for v in values:
                        if v == v and v is not None:
                            data_vals.append(float(v)/100.0 if is_pct else float(v))
                        else:
                            data_vals.append(0.0)
                    cd.add_series("Wynik", tuple(data_vals))

                    chart = slide.shapes.add_chart(
                        XL_CHART_TYPE.COLUMN_CLUSTERED,
                        Inches(0.4), chart_top,
                        Inches(12.5), prs.slide_height - chart_top - Inches(0.3),
                        cd
                    ).chart

                    chart.has_legend = False; chart.has_title = False
                    series = chart.series[0]
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = BAR_COLOR

                    dls = series.data_labels
                    _apply_dls(dls, TITLE_COLOR)
                    dls.number_format = "0.0%" if is_pct else "#,##0"
                    dls.number_format_is_linked = False

                    _fmt_numfmt(chart, is_pct)
                    chart.category_axis.tick_labels.font.size = Pt(10)

                def _add_cross_chart_slide(prs, title_text, categories,
                                           series_dict, subtitle="", is_pct=False):
                    """
                    Grouped column chart for cross-tabs.
                    categories  = row variable values (x-axis)
                    series_dict = {series_name: [values...]} one per column category
                    """
                    slide, chart_top = _slide_base(title_text, subtitle)

                    cd = ChartData()
                    cd.categories = [str(c)[:35] for c in categories]

                    ser_names = list(series_dict.keys())
                    palette = _palette_series(len(ser_names), _bar_hex)

                    for sname, svals in series_dict.items():
                        clean = []
                        for v in svals:
                            if v == v and v is not None:
                                clean.append(float(v)/100.0 if is_pct else float(v))
                            else:
                                clean.append(0.0)
                        cd.add_series(str(sname)[:30], tuple(clean))

                    chart = slide.shapes.add_chart(
                        XL_CHART_TYPE.COLUMN_CLUSTERED,
                        Inches(0.4), chart_top,
                        Inches(12.5), prs.slide_height - chart_top - Inches(0.3),
                        cd
                    ).chart

                    chart.has_legend = True
                    chart.has_title  = False
                    try:
                        chart.legend.position = 4   # BOTTOM
                        chart.legend.include_in_layout = False
                        chart.legend.font.size = Pt(10)
                    except Exception:
                        pass

                    for i, series in enumerate(chart.series):
                        c = palette[i % len(palette)]
                        series.format.fill.solid()
                        series.format.fill.fore_color.rgb = c

                        dls = series.data_labels
                        _apply_dls(dls, TITLE_COLOR)
                        dls.font.size = Pt(9)
                        dls.number_format = "0.0%" if is_pct else "#,##0"
                        dls.number_format_is_linked = False

                    _fmt_numfmt(chart, is_pct)
                    chart.category_axis.tick_labels.font.size = Pt(10)

                slides_added = 0
                use_pct = (ppt_metric == "Procent [%]")
                prefix  = (ppt_title_prefix.strip() + " " if ppt_title_prefix.strip() else "")

                # \u2500\u2500 Frequency tables (single series) \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500
                for var_name_key in sel_freq_keys:
                    df_res = freq_res[var_name_key]
                    try:
                        df_c = df_res.copy()
                        for sr in ['Suma','Og\u00f3\u0142em (Wa\u017cne)','Og\u00f3\u0142em',
                                   'Braki danych','Brak odpowiedzi']:
                            df_c = df_c[df_c.index.astype(str) != sr]
                        df_c = df_c[~df_c.index.astype(str).str.startswith('[')]

                        pct_col = next((c for c in df_c.columns
                                        if 'Procent' in str(c) or '%' in str(c)), None)
                        n_col   = next((c for c in df_c.columns
                                        if 'Liczebno' in str(c) or c == 'N'), None)
                        val_col = (pct_col if use_pct else n_col) or df_c.columns[0]

                        vals = pd.to_numeric(
                            df_c[val_col].astype(str).str.replace('%','').str.strip(),
                            errors='coerce').tolist()
                        cats = df_c.index.tolist()
                        if not cats or all(v != v for v in vals):
                            continue

                        lbl = var_labels.get(var_name_key, var_name_key)
                        unit = "%" if use_pct else "N"
                        if ppt_show_base:
                            try:
                                _skip = ['Suma','Og\u00f3\u0142em (Wa\u017cne)','Og\u00f3\u0142em',
                                         'Braki danych','Brak odpowiedzi']
                                _n_col = next((c for c in df_res.columns
                                               if 'Liczebno' in str(c) or c == 'N'), df_res.columns[0])
                                _n_ser = df_res[_n_col]
                                _n_ser = _n_ser[~_n_ser.index.astype(str).isin(_skip)]
                                _n_ser = _n_ser[~_n_ser.index.astype(str).str.startswith('[')]
                                base_n = int(pd.to_numeric(_n_ser, errors='coerce').dropna().sum())
                                sub = f"Baza: N={base_n} | Warto\u015bci: {unit}"
                            except Exception:
                                sub = f"Warto\u015bci: {unit}"
                        else:
                            sub = f"Warto\u015bci: {unit}"

                        _add_chart_slide(prs, f"{prefix}{lbl}"[:120],
                                         cats, vals, sub, is_pct=use_pct)
                        slides_added += 1
                    except Exception as _e:
                        st.warning(f"Pomini\u0119to '{var_name_key}': {_e}")

                # \u2500\u2500 Cross-tabs (grouped series, 1 chart per table) \u2500\u2500
                for cross_key in sel_cross_keys:
                    df_res = cross_res[cross_key]
                    try:
                        df_c = df_res.copy()
                        for sr in ['Suma','Braki danych',
                                   'Braki danych (wykluczone z tabeli)']:
                            df_c = df_c[df_c.index.astype(str) != sr]
                        df_c = df_c[~df_c.index.astype(str).str.startswith('[')]

                        # Select value columns
                        if use_pct:
                            val_cols = [c for c in df_c.columns
                                        if '%' in str(c) or 'Procent' in str(c)]
                        else:
                            val_cols = [c for c in df_c.columns
                                        if '[N]' in str(c) and 'Suma' not in str(c)]
                        if not val_cols:
                            val_cols = [c for c in df_c.columns if c != 'Suma']

                        cats = df_c.index.tolist()
                        if not cats:
                            continue

                        # Build series dict: clean column name -> numeric values
                        series_dict = {}
                        for col in val_cols:
                            col_label = (str(col).replace('[%]','')
                                                  .replace('[N]','')
                                                  .replace('[% Kolumnowe]','')
                                                  .replace('[% Wierszowe]','')
                                                  .strip())
                            svals = pd.to_numeric(
                                df_c[col].astype(str).str.replace('%','').str.strip(),
                                errors='coerce').tolist()
                            if any(v == v for v in svals):  # at least one non-NaN
                                series_dict[col_label] = svals

                        if not series_dict:
                            continue

                        unit = "%" if use_pct else "N"
                        title_s = f"{prefix}{cross_key}"[:120]
                        sub = f"Warto\u015bci: {unit} | Serie = kategorie zmiennej w kolumnach"

                        _add_cross_chart_slide(prs, title_s, cats,
                                               series_dict, sub, is_pct=use_pct)
                        slides_added += 1
                    except Exception as _e:
                        st.warning(f"Pomini\u0119to '{cross_key}': {_e}")

                if slides_added == 0:
                    st.error("Nie uda\u0142o si\u0119 wygenerowa\u0107 \u017cadnego wykresu.")
                else:
                    ppt_buf = io.BytesIO()
                    prs.save(ppt_buf)
                    ppt_buf.seek(0)
                    st.success(f"\u2705 Wygenerowano {slides_added} slajd\u00f3w.")
                    st.download_button(
                        label=f"\u2b07\ufe0f Pobierz prezentacj\u0119 ({slides_added} slajd\u00f3w)",
                        data=ppt_buf.getvalue(),
                        file_name="Wykresy_Analiz.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True
                    )

            except ImportError:
                st.error("Biblioteka `python-pptx` nie jest zainstalowana. Uruchom: `pip install python-pptx`")
            except Exception as _ppt_err:
                st.error(f"B\u0142\u0105d generowania PowerPoint: {_ppt_err}")
                st.exception(_ppt_err)

else:
    st.info("\U0001f448 Wybierz modu\u0142 z menu bocznego.")
